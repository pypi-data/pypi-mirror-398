"""
office.py - DOCX/PPTX/XLSX extraction
"""

import logging
from typing import Dict, Any
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import io

logger = logging.getLogger(__name__)


class OfficeExtractor:
    """Extract text from Office documents"""

    async def extract(self, response) -> Dict[str, Any]:
        """Extract text from Office document"""
        result = {"text": "", "type": None}

        content_type = getattr(response, "headers", {}).get("content-type", "").lower()
        url = getattr(response, "url", "")

        try:
            if "wordprocessingml" in content_type or url.endswith(".docx"):
                result["type"] = "docx"
                result.update(await self._extract_docx(response.content))
            elif "presentationml" in content_type or url.endswith(".pptx"):
                result["type"] = "pptx"
                result.update(await self._extract_pptx(response.content))
            elif "spreadsheetml" in content_type or url.endswith(".xlsx"):
                result["type"] = "xlsx"
                result.update(await self._extract_xlsx(response.content))
        except Exception as e:
            logger.error(f"Error extracting Office document: {e}")

        return result

    async def _extract_docx(self, content: bytes) -> Dict[str, Any]:
        """Extract from DOCX"""
        doc = Document(io.BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs]
        return {"text": "\n".join(paragraphs)}

    async def _extract_pptx(self, content: bytes) -> Dict[str, Any]:
        """Extract from PPTX"""
        prs = Presentation(io.BytesIO(content))
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return {"text": "\n".join(text_runs)}

    async def _extract_xlsx(self, content: bytes) -> Dict[str, Any]:
        """Extract from XLSX"""
        wb = load_workbook(io.BytesIO(content))
        text_runs = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell)
                if row_text:
                    text_runs.append(row_text)
        return {"text": "\n".join(text_runs)}
