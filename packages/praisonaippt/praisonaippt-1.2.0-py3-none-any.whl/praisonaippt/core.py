"""
Core presentation creation logic for Bible verses PowerPoint generator.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from .utils import split_long_text, sanitize_filename
from .pdf_converter import PDFOptions, convert_pptx_to_pdf


def add_title_slide(prs, title, subtitle=""):
    """
    Add a title slide to the presentation.
    
    Args:
        prs: Presentation object
        title (str): Title text
        subtitle (str): Subtitle text (optional)
    
    Returns:
        Slide object
    """
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
    
    return slide


def add_section_slide(prs, section_name):
    """
    Add a section title slide to the presentation.
    
    Args:
        prs: Presentation object
        section_name (str): Section title text
    
    Returns:
        Slide object
    """
    section_slide_layout = prs.slide_layouts[1]
    section_slide = prs.slides.add_slide(section_slide_layout)
    section_title = section_slide.shapes.title
    section_title.text = section_name
    
    # Style section title
    section_title.text_frame.paragraphs[0].font.size = Pt(44)
    section_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    return section_slide


def _apply_highlights(paragraph, text, highlights, large_text=None):
    """
    Apply highlighting and/or large text formatting to specific words or phrases in a paragraph.
    
    Args:
        paragraph: Paragraph object from text frame
        text (str): The full text to display
        highlights (list): List of words/phrases to highlight (bold + orange)
        large_text (dict): Dictionary mapping words/phrases to font sizes, e.g. {"word": 200}
    
    Returns:
        None (modifies paragraph in place)
    """
    # Create a case-insensitive search pattern for each highlight and large_text
    import re
    
    # Build a list of (start, end, matched_text, format_type, font_size) tuples for all matches
    matches = []
    
    # Add highlights
    if highlights:
        for highlight in highlights:
            # Escape special regex characters in the highlight text
            pattern = re.escape(highlight)
            # Find all occurrences (case-insensitive)
            for match in re.finditer(pattern, text, re.IGNORECASE):
                matches.append((match.start(), match.end(), match.group(), 'highlight', None))
    
    # Add large_text entries
    if large_text:
        for word, font_size in large_text.items():
            pattern = re.escape(word)
            for match in re.finditer(pattern, text, re.IGNORECASE):
                matches.append((match.start(), match.end(), match.group(), 'large', font_size))
    
    # Sort matches by start position
    matches.sort(key=lambda x: x[0])
    
    # Remove overlapping matches (keep first occurrence)
    filtered_matches = []
    last_end = -1
    for start, end, matched_text, format_type, font_size in matches:
        if start >= last_end:
            filtered_matches.append((start, end, matched_text, format_type, font_size))
            last_end = end
    
    # Build the paragraph with highlighted sections
    if not filtered_matches:
        # No matches found, just add plain text
        paragraph.text = text
        paragraph.font.size = Pt(32)
        paragraph.font.color.rgb = RGBColor(0, 0, 0)
        return
    
    # Add text segments with appropriate formatting
    current_pos = 0
    first_run = True
    
    for start, end, matched_text, format_type, font_size in filtered_matches:
        # Add text before the formatted text
        if start > current_pos:
            if first_run:
                paragraph.text = text[current_pos:start]
                run = paragraph.runs[0]
                first_run = False
            else:
                run = paragraph.add_run()
                run.text = text[current_pos:start]
            run.font.size = Pt(32)
            run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Add the formatted text
        if first_run:
            paragraph.text = matched_text
            run = paragraph.runs[0]
            first_run = False
        else:
            run = paragraph.add_run()
            run.text = matched_text
        
        # Apply formatting based on type
        if format_type == 'highlight':
            # Apply highlight formatting (bold + orange color)
            run.font.size = Pt(32)
            run.font.bold = True
            run.font.color.rgb = RGBColor(255, 140, 0)  # Orange color for highlights
        elif format_type == 'large':
            # Apply large text formatting (custom font size, keep black color)
            run.font.size = Pt(font_size)
            run.font.color.rgb = RGBColor(0, 0, 0)  # Keep black color
        
        current_pos = end
    
    # Add any remaining text after the last highlight
    if current_pos < len(text):
        run = paragraph.add_run()
        run.text = text[current_pos:]
        run.font.size = Pt(32)
        run.font.color.rgb = RGBColor(0, 0, 0)


def add_verse_slide(prs, verse_text, reference, part_num=None, highlights=None, large_text=None):
    """
    Add a verse slide to the presentation.
    
    Args:
        prs: Presentation object
        verse_text (str): The verse text
        reference (str): The verse reference
        part_num (int): Part number if verse is split (optional)
        highlights (list): List of words/phrases to highlight (optional)
        large_text (dict): Dictionary mapping words to custom font sizes (optional)
    
    Returns:
        Slide object
    """
    verse_slide_layout = prs.slide_layouts[6]  # Blank layout
    verse_slide = prs.slides.add_slide(verse_slide_layout)
    
    # Add text box for verse
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(3)
    
    textbox = verse_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add verse text with highlighting
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    if (highlights and len(highlights) > 0) or (large_text and len(large_text) > 0):
        # Apply highlighting and/or large text formatting
        _apply_highlights(p, verse_text, highlights, large_text)
    else:
        # No formatting, just add plain text
        p.text = verse_text
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0, 0, 0)
    
    # Add reference
    reference_text = reference
    if part_num is not None:
        reference_text += f" (Part {part_num})"
    
    ref_left = Inches(1)
    ref_top = Inches(5.5)
    ref_width = Inches(8)
    ref_height = Inches(1)
    
    ref_textbox = verse_slide.shapes.add_textbox(ref_left, ref_top, ref_width, ref_height)
    ref_text_frame = ref_textbox.text_frame
    
    ref_p = ref_text_frame.paragraphs[0]
    ref_p.text = reference_text
    ref_p.alignment = PP_ALIGN.CENTER
    ref_p.font.size = Pt(22)
    ref_p.font.color.rgb = RGBColor(100, 100, 100)
    ref_p.font.italic = True
    
    return verse_slide


def create_presentation(data, output_file=None, custom_title=None, 
                        convert_to_pdf=False, pdf_options=None, pdf_backend='auto'):
    """
    Create a PowerPoint presentation from Bible verses data.
    
    Args:
        data (dict): Verses data dictionary with structure:
                     {
                         "presentation_title": "Title",
                         "presentation_subtitle": "Subtitle",
                         "sections": [
                             {
                                 "section": "Section Name",
                                 "verses": [
                                     {"reference": "Book 1:1", "text": "Verse text"}
                                 ]
                             }
                         ]
                     }
        output_file (str): Output filename (optional, auto-generated if not provided)
        custom_title (str): Custom presentation title (optional, overrides JSON title)
        convert_to_pdf (bool): Whether to also convert to PDF (default: False)
        pdf_options (PDFOptions): PDF conversion options (optional)
        pdf_backend (str): PDF conversion backend ('aspose', 'libreoffice', 'auto')
    
    Returns:
        str or dict: Path to the created presentation file, or dict with both PPTX and PDF paths
                    if convert_to_pdf is True, or None if error
    """
    if not data:
        print("Error: No data provided")
        return None
    
    # Create presentation
    prs = Presentation()
    
    # Get verses data from JSON
    verses_data = data.get("sections", [])
    
    # Add title slide
    if custom_title:
        title = custom_title
        subtitle = ""
    else:
        title = data.get("presentation_title", "Bible Verses Collection")
        subtitle = data.get("presentation_subtitle", "Selected Scriptures")
    
    add_title_slide(prs, title, subtitle)
    
    # Add slides for each verse with section slides
    for section_data in verses_data:
        # Add section title slide if section name exists (skip if custom title is provided)
        if section_data.get("section") and not custom_title:
            add_section_slide(prs, section_data["section"])
        
        # Add verse slides if there are any verses
        if section_data.get("verses") and len(section_data["verses"]) > 0:
            for verse in section_data["verses"]:
                # Split long verses into multiple parts
                verse_parts = split_long_text(verse["text"])
                
                # Get highlights and large_text if specified
                highlights = verse.get("highlights", None)
                large_text = verse.get("large_text", None)
                
                for i, part in enumerate(verse_parts):
                    part_num = i + 1 if len(verse_parts) > 1 else None
                    add_verse_slide(prs, part, verse["reference"], part_num, highlights, large_text)
    
    # Generate output filename if not provided
    if not output_file:
        if custom_title:
            base_name = sanitize_filename(custom_title)
        else:
            base_name = sanitize_filename(data.get("presentation_title", "presentation"))
        output_file = f"{base_name}.pptx"
    
    # Ensure .pptx extension
    if not output_file.endswith('.pptx'):
        output_file += '.pptx'
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✓ Presentation created successfully: {output_file}")
        
        # Convert to PDF if requested
        if convert_to_pdf:
            try:
                # Use default PDF options if none provided
                if pdf_options is None:
                    pdf_options = PDFOptions()
                
                # Generate PDF filename
                from pathlib import Path
                pdf_file = str(Path(output_file).with_suffix('.pdf'))
                
                # Convert to PDF
                pdf_result = convert_pptx_to_pdf(
                    output_file, 
                    pdf_file, 
                    backend=pdf_backend, 
                    options=pdf_options
                )
                
                print(f"✓ PDF created successfully: {pdf_result}")
                
                # Return both files
                return {
                    'pptx': output_file,
                    'pdf': pdf_result
                }
                
            except Exception as e:
                print(f"Warning: PDF conversion failed: {e}")
                print("Presentation was created successfully at:", output_file)
                return output_file
        
        return output_file
        
    except Exception as e:
        print(f"Error saving presentation: {e}")
        return None
