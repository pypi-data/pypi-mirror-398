import keyring as kr
from pydantic import BaseModel, Field, ConfigDict, field_validator
from pathlib import Path
import psycopg2
import keyring as kr
import json
from io import StringIO
from pptx import Presentation  # Install using: pip install python-pptx
from docx import Document  # Install using: pip install python-docx
import openai
from datetime import datetime
from typing import List, Dict, Optional
import logging
from tenacity import retry, stop_after_attempt, wait_fixed
from contextlib import contextmanager
import time
import threading
import hashlib
from sentence_transformers import SentenceTransformer, util
from langchain.text_splitter import RecursiveCharacterTextSplitter
import torch
from krazy.ai.tools import Config

class RAGPostgres:
    def __init__(self, config:Config, doc_table='rag_documents'):
        self.config = config
        self.conn = self.connect_to_postgres_aws()
        self.schema = 'ai'
        self.doc_table = doc_table
        self.chat_table = 'chat_history'
        self.ingested_files_table = 'ingested_files'
        self.watcher_thread = None
        self.stop_watcher_flag = threading.Event()

        self.client = openai.AzureOpenAI(
            api_key=self.config.pxv_open_ai_api_key,
            azure_endpoint=self.config.azure_endpoint,
            api_version="2023-07-01-preview"
        )
        self.sentence_encoder = SentenceTransformer("all-mpnet-base-v2")

    def connect_to_postgres_aws(self):
        return psycopg2.connect(
            dbname="postgres",
            user="root",
            password=kr.get_password('aws-posgtres', 'root'),
            host="finance-kartik-poc.cq4isu0u5wom.ap-southeast-1.rds.amazonaws.com"
        )

    @contextmanager
    def get_cursor(self):
        with self.conn.cursor() as cur:
            yield cur

    def setup_tables(self):
        with self.get_cursor() as cur:
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {self.schema}.{self.doc_table} (
                    id SERIAL PRIMARY KEY,
                    content TEXT,
                    metadata JSONB,
                    embedding VECTOR(1536),
                    tsv TSVECTOR
                );
            """)
            cur.execute(f"""
                CREATE INDEX IF NOT EXISTS idx_embedding_vector ON {self.schema}.{self.doc_table} USING ivfflat (embedding vector_cosine_ops);
            """)
            cur.execute(f"""
                CREATE INDEX IF NOT EXISTS idx_metadata ON {self.schema}.{self.doc_table} USING GIN (metadata);
            """)
            cur.execute(f"""
                CREATE INDEX IF NOT EXISTS idx_tsv ON {self.schema}.{self.doc_table} USING GIN(tsv);
            """)
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {self.schema}.{self.chat_table} (
                    id SERIAL PRIMARY KEY,
                    chat_agent TEXT,
                    agent_name TEXT,
                    user_input TEXT,
                    assistant_response TEXT,
                    timestamp TIMESTAMPTZ DEFAULT NOW(),
                    session_id UUID DEFAULT gen_random_uuid(),
                    embedding VECTOR(1536)
                );
            """)
            cur.execute(f"""
                CREATE TABLE IF NOT EXISTS {self.schema}.{self.ingested_files_table} (
                    id SERIAL PRIMARY KEY,
                    file_name TEXT UNIQUE,
                    checksum TEXT
                );
            """)
        self.conn.commit()

    def calculate_checksum(self, file_path: str) -> str:
        hasher = hashlib.sha256()
        with open(file_path, 'rb') as f:
            while chunk := f.read(8192):
                hasher.update(chunk)
        return hasher.hexdigest()

    def is_file_already_ingested(self, file_path: str) -> bool:
        checksum = self.calculate_checksum(file_path)
        with self.get_cursor() as cur:
            cur.execute(
                f"SELECT 1 FROM {self.schema}.{self.ingested_files_table} WHERE file_name = %s AND checksum = %s",
                (Path(file_path).name, checksum)
            )
            return cur.fetchone() is not None

    def record_file_ingestion(self, file_path: str):
        checksum = self.calculate_checksum(file_path)
        with self.get_cursor() as cur:
            cur.execute(
                f"INSERT INTO {self.schema}.{self.ingested_files_table} (file_name, checksum) VALUES (%s, %s) ON CONFLICT DO NOTHING",
                (Path(file_path).name, checksum)
            )
        self.conn.commit()

    @retry(stop=stop_after_attempt(3), wait=wait_fixed(2))
    def get_embedding(self, text: str) -> List[float]:
        response = self.client.embeddings.create(
            input=[text],
            model="text-embedding-ada-002"
        )
        return response.data[0].embedding

    def chunk_text(self, text: str, max_chunk_tokens=400, overlap=40) -> List[str]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=max_chunk_tokens,
            chunk_overlap=overlap,
            separators=["\n\n", "\n", ".", " "],
            length_function=lambda x: len(x.split())
        )
        return splitter.split_text(text)

    def extract_text_from_pdf(self, file_path: str) -> str:
        with fitz.open(file_path) as pdf:
            return "\n".join(page.get_text() for page in pdf).strip()

    def extract_text_from_word(self, file_path: str) -> str:
        doc = Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

    def extract_text_from_ppt(self, file_path: str) -> str:
        prs = Presentation(file_path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

    def extract_text_from_txt_or_md(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8") as file:
            return file.read().strip()

    def ingest_file(self, file_path: str, metadata: Optional[Dict] = None):
        if self.is_file_already_ingested(file_path):
            print(f"🔁 Skipping already ingested file: {file_path}")
            return

        ext = Path(file_path).suffix.lower()
        metadata = metadata or {"source": str(file_path)}
        try:
            if ext == ".pdf":
                text = self.extract_text_from_pdf(file_path)
            elif ext == ".docx":
                text = self.extract_text_from_word(file_path)
            elif ext == ".pptx":
                text = self.extract_text_from_ppt(file_path)
            elif ext in [".txt", ".md"]:
                text = self.extract_text_from_txt_or_md(file_path)
            else:
                print(f"❌ Unsupported file type for ingestion: {ext}")
                return

            chunks = self.chunk_text(text)
            self.insert_document_chunks(chunks, metadata)
            self.record_file_ingestion(file_path)
        except Exception as e:
            print(f"❌ Failed to ingest file {file_path}: {e}")

    def insert_document_chunks(self, chunks: List[str], metadata: Dict):
        with self.get_cursor() as cur:
            for chunk in chunks:
                embedding = self.get_embedding(chunk)
                cur.execute(
                    f"""
                    INSERT INTO {self.schema}.{self.doc_table} (content, metadata, embedding, tsv)
                    VALUES (%s, %s, %s, to_tsvector(%s))
                    """,
                    (chunk, json.dumps(metadata), embedding, chunk)
                )
        self.conn.commit()

    def watch_folder(self, folder_path: str, polling_interval: int = 10):
        print(f"👀 Watching folder: {folder_path}")
        folder = Path(folder_path)

        def watch():
            while not self.stop_watcher_flag.is_set():
                for file_path in folder.glob("*"):
                    if file_path.is_file():
                        self.ingest_file(str(file_path))
                time.sleep(polling_interval)

        self.watcher_thread = threading.Thread(target=watch, daemon=True)
        self.watcher_thread.start()

    def stop_watcher(self):
        if self.watcher_thread:
            self.stop_watcher_flag.set()
            self.watcher_thread.join()
            print("🛑 Watcher stopped.")

    def highlight_relevant_passage(self, query: str, chunk: str) -> str:
        sentences = [s.strip() for s in chunk.split('.') if len(s.strip()) > 0]
        if not sentences:
            return chunk
        query_embedding = self.sentence_encoder.encode(query, convert_to_tensor=True)
        sentence_embeddings = self.sentence_encoder.encode(sentences, convert_to_tensor=True)
        scores = util.pytorch_cos_sim(query_embedding, sentence_embeddings)[0]
        top_idx = scores.argmax().item()
        return sentences[top_idx]

    def search_similar_documents(self, query: str, top_k: int = 5, metadata_filter: Optional[Dict] = None, highlight=True):
        embedding = self.get_embedding(query)
        sql = f"SELECT content, metadata FROM {self.schema}.{self.doc_table}"
        params = []

        where_clauses = []
        if metadata_filter:
            where_clauses.append("metadata @> %s::jsonb")
            params.append(json.dumps(metadata_filter))
        where_clause = " WHERE " + " AND ".join(where_clauses) if where_clauses else ""

        sql += where_clause
        sql += f" ORDER BY embedding <-> %s::vector LIMIT %s"
        params.extend([embedding, top_k * 3])

        with self.get_cursor() as cur:
            cur.execute(sql, params)
            results = cur.fetchall()

        ranked = self.rerank_results(query, results, final_k=top_k)
        return ranked if highlight else results

    def rerank_results(self, query: str, chunks: List[tuple], final_k: int = 5):
        contents = [c[0] for c in chunks]
        meta = [c[1] for c in chunks]
        query_embedding = self.sentence_encoder.encode(query, convert_to_tensor=True)
        doc_embeddings = self.sentence_encoder.encode(contents, convert_to_tensor=True)
        scores = util.pytorch_cos_sim(query_embedding, doc_embeddings)[0]
        top_results = torch.topk(scores, final_k)
        return [
            {"highlight": contents[i], "metadata": meta[i]}
            for i in top_results.indices.tolist()
        ]
