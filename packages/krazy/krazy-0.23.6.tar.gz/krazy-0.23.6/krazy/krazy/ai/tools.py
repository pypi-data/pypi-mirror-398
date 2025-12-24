
import keyring as kr
from tavily import TavilyClient
import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from pydantic import BaseModel, Field, ConfigDict, field_validator, PrivateAttr
import os
import requests
from bs4 import BeautifulSoup
from crewai.tools import BaseTool
from typing import Type
import pandas as pd
from pathlib import Path
import psycopg2
import keyring as kr
import json
from io import StringIO
from pptx import Presentation  # Install using: pip install python-pptx
from docx import Document  # Install using: pip install python-docx
import fitz  # PyMuPDF
import openai
from datetime import datetime
from typing import List, Dict, Optional
from tenacity import retry, stop_after_attempt, wait_fixed
from contextlib import contextmanager

class Config:
    """
    Configuration handler for Azure OpenAI, Tavily, PostgreSQL, and embedding settings.
    Expects a dictionary containing all required credentials and parameters.
    """

    def __init__(self, params: Dict[str, str]):
        # ✅ Azure OpenAI Credentials
        self.api_key = params.get('api_key')
        self.deployment_name = params.get('deployment_name')
        self.endpoint = params.get('endpoint')

        # ✅ Model Identifiers
        self.model = params.get('model', 'gpt-4o')
        self.model_32k = params.get('model_32k', 'gpt-4-32k')
        self.summary_model = params.get('summary_model', 'gpt-35-turbo-16k')
        self.embedding_model = params.get('embedding_model', 'text-embedding-ada-002')

        # ✅ Tavily API Key
        self.tavily_api_key = params.get('tavily_api_key')

        # ✅ PostgreSQL Credentials
        self.postgres_dbname = params.get('postgres_dbname', 'postgres')
        self.postgres_user = params.get('postgres_user', 'root')
        self.postgres_password = params.get('postgres_password')
        self.postgres_host = params.get('postgres_host')
        self.postgres_port = params.get('postgres_port', '5432')

        # ✅ Azure Embedding Endpoint for ChromaDB usage
        self.endpoint_url = f"{self.endpoint}/openai/deployments/{self.embedding_model}/completions?api-version={self.deployment_name}"

        # ✅ ChromaDB Embedding Configuration
        self.embedder_config = {
            "embedder": {
                "provider": "azure_openai",
                "config": {
                    "deployment_name": self.embedding_model,
                    "api_key": self.api_key,
                    "api_base": self.endpoint_url,
                },
            },
        }

    @staticmethod
    def sample_input_param() -> Dict[str, str]:
        """
        Returns a sample dictionary of required configuration parameters.
        """
        return {
            # Azure OpenAI
            "api_key": "your-azure-openai-api-key",
            "deployment_name": "2023-07-01-preview",
            "endpoint": "https://your-resource-name.openai.azure.com",
            "model": "gpt-4o",
            "model_32k": "gpt-4-32k",
            "summary_model": "gpt-35-turbo-16k",
            "embedding_model": "text-embedding-ada-002",

            # Tavily
            "tavily_api_key": "your-tavily-api-key",

            # PostgreSQL
            "postgres_dbname": "postgres",
            "postgres_user": "root",
            "postgres_password": "your-db-password",
            "postgres_host": "your-db-host",
            "postgres_port": "5432"
        }
    


class WebSearchToolInput(BaseModel):
    """Input model for the WebSearchTool."""
    query: str = Field(..., description="The search query to fetch results from the web.")
    num_results: int = Field(default=3, description="Number of search results to retrieve.")

class WebSearchTool(BaseTool):
    """A CrewAI-compatible tool that searches the web and extracts content from search results."""
    
    name: str = "web_search"
    description: str = "Searches the web using DuckDuckGo and extracts content from results."
    model_config = ConfigDict(arbitrary_types_allowed=True)  # ✅ Allow custom types
    args_schema:Type[BaseModel] = WebSearchToolInput  # ✅ Use the input schema

    def _run(self, query: str, num_results: int = 3) -> str:
        """
        Searches the web and extracts titles, links, and content.

        Args:
            query (str): The search query.
            num_results (int): Number of results to return.

        Returns:
            str: Formatted search results.
        """
        search_url = f"https://duckduckgo.com/html/?q={query.replace(' ', '+')}"
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }

        response = requests.get(search_url, headers=headers)
        if response.status_code != 200:
            return f"❌ Error: Unable to fetch search results. Status Code: {response.status_code}"

        soup = BeautifulSoup(response.text, "html.parser")
        extracted_results = []

        # ✅ Extract search results
        for result in soup.find_all("a", class_="result__a", limit=num_results):
            title = result.get_text()
            link = result["href"]
            content = self.fetch_page_content(link)

            extracted_results.append(f"🔹 **{title}**\n🔗 {link}\n📜 {content}\n" + "-" * 50)

        return "\n\n".join(extracted_results) if extracted_results else "❌ No results found."

    def fetch_page_content(self, url: str) -> str:
        """
        Fetches and extracts the main content of a webpage.

        Args:
            url (str): The webpage URL.

        Returns:
            str: Extracted text content.
        """
        try:
            headers = {
                "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
            }

            response = requests.get(url, headers=headers, timeout=5)
            if response.status_code != 200:
                return f"❌ Error: Unable to fetch page. Status Code: {response.status_code}"

            soup = BeautifulSoup(response.text, "html.parser")
            paragraphs = soup.find_all("p")
            page_content = "\n".join([p.get_text() for p in paragraphs])

            return page_content[:1000] + "..." if len(page_content) > 1000 else page_content  # Limit to 1000 chars

        except Exception as e:
            return f"❌ Error fetching content: {e}"

class WebSearchToolTavilyInput[BaseModel]:
    query: str = Field(..., description="The search query to fetch results from the web.")

class WebSearchToolTavily(BaseTool):
    """
    A CrewAI-compatible web search tool using Tavily's API.

    Attributes:
        api_key (str): API key for Tavily. Must be provided during initialization.
    """

    name: str = "web_search"
    description: str = "Searches the web and returns results."
    model_config = ConfigDict(arbitrary_types_allowed=True)
    args_schema: Type[BaseModel] = WebSearchToolTavilyInput

    def __init__(self, api_key: str):
        super().__init__()
        if not api_key:
            raise ValueError("❌ Tavily API key must be provided.")
        self.api_key = api_key

    def _run(self, query: str) -> str:
        """
        Searches the web for the given query using Tavily.

        Args:
            query (str): The search query.

        Returns:
            str: The search results as a formatted string.
        """
        try:
            client = TavilyClient(api_key=self.api_key)
            response = client.search(query, max_results=2)

            if "results" in response and response["results"]:
                formatted_response = "\n".join([f"- {r}" for r in response["results"]])
                return f"Here are the top results for your query:\n{formatted_response}"
            else:
                return "No results found for your query. Please try a different query."

        except Exception as e:
            return f"❌ An error occurred during the web search: {e}"
        

class PDFReaderToolInput(BaseModel):
    """Input schema for the PDF Reader Tool."""
    file_path: str = Field(..., description="Path to the PDF file.")

class PDFReaderTool(BaseTool):
    name: str = "pdf_reader"
    description: str = "Reads and extracts text from a PDF file."
    model_config = ConfigDict(arbitrary_types_allowed=True)  # Use ConfigDict for Pydantic v2
    args_schema:Type[BaseModel]=PDFReaderToolInput

    def _run(self, file_path:str) -> str:

        """
        Extracts text from all pages of a PDF.

        Args:
            file_path (str): Path to the PDF file.

        Returns:
            str: Full text content of the PDF.
        """
        if not Path(file_path).exists():
            raise FileNotFoundError(f"❌ File does not exist: {file_path}")

        try:
            text = ""
            with fitz.open(file_path) as pdf:
                for page_num in range(len(pdf)):
                    page = pdf[page_num]
                    text += page.get_text()
            return text.strip()
        
        except Exception as e:
            return f"❌ Error reading PDF: {e}"

class CSVCustomReaderToolInput(BaseModel):
    """Input schema for the CSV Reader Tool."""
    code: str = Field(..., description="Python code for extracting information from the CSV.")
    file_path: str = Field(..., description="Path to the CSV file.")

class CSVCustomReaderTool(BaseTool):
    """Tool for reading and executing agent-generated Python code on a CSV file."""
    
    name: str = "csv_reader"
    description: str = "Executes Python code on a CSV file and returns processed results. Include result = ... in the code. result is what will be returned."
    args_schema:Type[BaseModel] = CSVCustomReaderToolInput

    def _run(self, code: str, file_path: str=None) -> str:
        """Executes the agent-provided Python code on the given CSV file."""
        try:
            if file_path is None:
                return "Provide path to CSV file"
            if Path(file_path).exists:
                pass
            else:
                return 'File does not exist'
            
            df = pd.DataFrame()
            exec_globals = {"df": df}
            exec(code, exec_globals)
            
            # The agent should always return the result in a variable named `result`
            result = exec_globals.get("result", "❌ No valid result found.")

            return result
        
        except Exception as e:
            return f"❌ Error executing code: {e}"


class PostgresQueryInput(BaseModel):
    """Schema for Postgres Query Tool."""
    query: str = Field(..., description="Select SQL query to execute in PostgreSQL.")

class PostgresQueryTool(BaseTool):
    """Tool for querying a PostgreSQL database."""

    name: str = "postgres_query"
    description: str = "Executes SELECT SQL queries on a PostgreSQL database."
    args_schema: Type[BaseModel] = PostgresQueryInput

    _conn_params: Dict[str, str] = PrivateAttr()

    def __init__(self, conn_params: Dict[str, str]):
        super().__init__()
        self._conn_params = conn_params

    def _run(self, query: str) -> str:
        """Executes a SQL query and returns the result as JSON."""
        if not query.lower().strip().startswith("select"):
            return "Only SELECT queries are allowed."

        try:
            conn = psycopg2.connect(
                dbname=self._conn_params.get("dbname", "postgres"),
                user=self._conn_params.get("user", "root"),
                password=self._conn_params["password"],
                host=self._conn_params["host"],
                port=self._conn_params.get("port", "5432")
            )
            cursor = conn.cursor()
            cursor.execute(query)
            records = cursor.fetchall()
            col_names = [desc[0] for desc in cursor.description]
            cursor.close()
            conn.close()

            return json.dumps([dict(zip(col_names, row)) for row in records])
        except Exception as e:
            return f"❌ Error executing query: {e}"

    def sample_conn_param(self):
        """Sample connection parameters for PostgreSQL."""
        return {
            "dbname": "your_database_name",
            "user": "your_username",
            "password": "your_password",
            "host": "localhost",
            "port": "5432"
        }
        


class ExcelImportToolInput(BaseModel):
    """Input schema for importing and processing an Excel file."""
    file_path: str = Field(..., description="Path to the Excel file to import.")
    code: str = Field(..., description="Python code to execute on the DataFrame.")

class ExcelImportTool(BaseTool):
    """Tool for importing an Excel file, executing LLM-generated code on it, and returning results."""
    
    name: str = "excel_import"
    description: str = "Imports an Excel file, processes it using the provided Python code, and returns the result."
    args_schema: Type[BaseModel] = ExcelImportToolInput

    def _run(self, file_path: str, code: str):
        """Loads an Excel file into a DataFrame, executes the given Python code, and returns the output."""
        try:
            if not os.path.exists(file_path):
                return f"❌ Error: The file '{file_path}' does not exist."

            xl = pd.ExcelFile(file_path)  # Loads the excel file
            exec_globals = {"xl": xl}  # Define execution environment

            exec(code, exec_globals)  # Execute LLM-provided code

            # Extract result from the execution
            result = exec_globals.get("result", "❌ No valid result found. Ensure LLM sets 'result' variable.")

            return result  # Return processed result

        except Exception as e:
            return f"❌ Error processing Excel file: {e}"

class CSVorExcelExportToolInput(BaseModel):
    """Input schema for exporting data to CSV."""
    file_path: str = Field(..., description="Path to save the exported CSV file.")
    dataframe_json: str = Field(..., description="DataFrame in JSON format to export.")
    export_as: str = Field("csv", description="Export format (default: csv). Mention xlsx for excel export.")

class CSVorExcelExportTool(BaseTool):
    """Tool for exporting a pandas DataFrame to a CSV file."""
    
    name: str = "csv_export"
    description: str = "Exports data (from a DataFrame) to a CSV file."
    args_schema: Type[BaseModel] = CSVorExcelExportToolInput

    def _run(self, file_path: str, dataframe_json: str, export_as: str = "csv"):
        """Exports the provided JSON data to a CSV file."""
        try:
            if export_as.lower() == "csv":
                df = pd.read_json(StringIO(dataframe_json))   # Convert JSON to DataFrame
                df.to_csv(file_path, index=False)  # Save as CSV
            if export_as.lower() == "xlsx":
                df = pd.read_json(StringIO(dataframe_json))
                df.to_excel(file_path, index=False)

            return f"✅ Data exported successfully to {file_path}"

        except Exception as e:
            return f"❌ Error exporting CSV: {e}"


class WordFileReaderToolInput(BaseModel):
    """Input schema for reading a Word (.docx) file."""
    file_path: str = Field(..., description="Path to the Word file.")

class WordFileReaderTool(BaseTool):
    """Tool to read and extract text from a Microsoft Word (.docx) file."""
    
    name: str = "word_reader"
    description: str = "Reads and extracts text from a Word document (.docx)."
    args_schema: Type[BaseModel] = WordFileReaderToolInput

    def _run(self, file_path: str):
        """Reads and extracts text from a Word file."""
        try:
            if not os.path.exists(file_path):
                return f"❌ Error: The file '{file_path}' does not exist."

            doc = Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])

            return text.strip() if text else "❌ No text found in the document."

        except Exception as e:
            return f"❌ Error reading Word file: {e}"


class PowerPointFileReaderToolInput(BaseModel):
    """Input schema for reading a PowerPoint (.pptx) file."""
    file_path: str = Field(..., description="Path to the PowerPoint file.")

class PowerPointFileReaderTool(BaseTool):
    """Tool to read and extract text from a PowerPoint (.pptx) file."""
    
    name: str = "pptx_reader"
    description: str = "Reads and extracts text from a PowerPoint presentation (.pptx)."
    args_schema: Type[BaseModel] = PowerPointFileReaderToolInput

    def _run(self, file_path: str):
        """Reads and extracts text from a PowerPoint file."""
        try:
            if not os.path.exists(file_path):
                return f"❌ Error: The file '{file_path}' does not exist."

            ppt = Presentation(file_path)
            slides_text = []

            for slide in ppt.slides:
                slide_text = "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                slides_text.append(slide_text)

            full_text = "\n\n".join(slides_text)

            return full_text.strip() if full_text else "❌ No text found in the presentation."

        except Exception as e:
            return f"❌ Error reading PowerPoint file: {e}"

class FolderListFilesToolInput(BaseModel):
    """Input schema for listing files in a folder."""
    folder_path: str = Field(..., description="Path to the folder whose files need to be listed.")

class FolderListFilesTool(BaseTool):
    """Tool to list all files in a specified folder."""
    
    name: str = "folder_list_files"
    description: str = "Lists all files in a given folder."
    args_schema: Type[BaseModel] = FolderListFilesToolInput

    def _run(self, folder_path: str):
        """Lists all files in a folder."""
        try:
            if not os.path.exists(folder_path):
                return f"❌ Error: The folder '{folder_path}' does not exist."

            if not os.path.isdir(folder_path):
                return f"❌ Error: The path '{folder_path}' is not a directory."

            files = os.listdir(folder_path)
            file_list = [f for f in files if os.path.isfile(os.path.join(folder_path, f))]

            return file_list if file_list else "📁 No files found in the folder."

        except Exception as e:
            return f"❌ Error listing files: {e}"


