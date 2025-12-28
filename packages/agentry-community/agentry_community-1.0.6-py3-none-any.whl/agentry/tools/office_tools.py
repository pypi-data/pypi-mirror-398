from typing import Optional, List, Literal
from pydantic import BaseModel, Field
import os
from .base import BaseTool, ToolResult

# --- PPTX ---
from pptx import Presentation
from pptx.util import Inches, Pt

class EditPPTXParams(BaseModel):
    file_path: str = Field(..., description='Path to PPTX file.')
    slide_index: Optional[int] = Field(None, description='Slide index (0-indexed). If None, searches all slides.')
    old_text: str = Field(..., description='Text to find and replace.')
    new_text: str = Field(..., description='Replacement text.')
    replace_all: bool = Field(True, description='Replace all occurrences.')

class CreatePPTXParams(BaseModel):
    file_path: str = Field(..., description='Path for new PPTX file.')
    title: str = Field(..., description='Title of the presentation.')
    subtitle: Optional[str] = Field(None, description='Subtitle/Author name.')
    slides: Optional[List[dict]] = Field(None, description='List of slides to add. Format: [{"title": "Slide 1", "content": "Bullet points..."}]')

class AppendSlideParams(BaseModel):
    file_path: str = Field(..., description='Path to PPTX file.')
    title: str = Field(..., description='Slide title.')
    content: str = Field(..., description='Slide content (text/bullet points).')
    layout_index: int = Field(1, description='Layout index (0=Title, 1=Content, etc.)')

class EditPPTXTool(BaseTool):
    name = "edit_pptx"
    description = "Replace text in a PowerPoint presentation."
    args_schema = EditPPTXParams

    def run(self, file_path: str, old_text: str, new_text: str, slide_index: int = None, replace_all: bool = True) -> ToolResult:
        try:
            if not os.path.exists(file_path): return ToolResult(success=False, error="File not found")
            
            prs = Presentation(file_path)
            replacements = 0
            
            slides_to_check = [prs.slides[slide_index]] if slide_index is not None and slide_index < len(prs.slides) else prs.slides
            
            for slide in slides_to_check:
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
                                replacements += 1
                                if not replace_all: break
            
            if replacements > 0:
                prs.save(file_path)
                return ToolResult(success=True, content=f"Replaced {replacements} occurrences of '{old_text}'.")
            else:
                return ToolResult(success=False, error=f"Text '{old_text}' not found.")
        except Exception as e:
            return ToolResult(success=False, error=f"PPTX Error: {e}")

class CreatePPTXTool(BaseTool):
    name = "create_pptx"
    description = "Create a new PowerPoint presentation."
    args_schema = CreatePPTXParams
    
    def run(self, file_path: str, title: str, subtitle: str = None, slides: List[dict] = None) -> ToolResult:
        try:
            prs = Presentation()
            
            # Title Slide
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            slide.shapes.title.text = title
            if subtitle:
                slide.placeholders[1].text = subtitle
                
            # Content Slides
            if slides:
                for s_data in slides:
                    s_title = s_data.get('title', 'Slide')
                    s_content = s_data.get('content', '')
                    
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = s_title
                    slide.placeholders[1].text = s_content
            
            prs.save(file_path)
            return ToolResult(success=True, content=f"Created presentation: {file_path}")
        except Exception as e:
            return ToolResult(success=False, error=f"PPTX Create Error: {e}")

class AppendSlideTool(BaseTool):
    name = "append_slide"
    description = "Add a new slide to an existing PPTX."
    args_schema = AppendSlideParams
    
    def run(self, file_path: str, title: str, content: str, layout_index: int = 1) -> ToolResult:
        try:
            if not os.path.exists(file_path): return ToolResult(success=False, error="File not found")
            prs = Presentation(file_path)
            
            if layout_index >= len(prs.slide_layouts): layout_index = 1
            
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            if slide.shapes.title: slide.shapes.title.text = title
            
            # Try to place content in body
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = content
            
            prs.save(file_path)
            return ToolResult(success=True, content=f"Appended slide '{title}' to {file_path}")
        except Exception as e:
            return ToolResult(success=False, error=f"Append Slide Error: {e}")


# --- DOCX ---
from docx import Document

class EditDOCXParams(BaseModel):
    file_path: str = Field(..., description='Path to DOCX file.')
    old_text: str = Field(..., description='Text to replace.')
    new_text: str = Field(..., description='Replacement text.')

class CreateDOCXParams(BaseModel):
    file_path: str = Field(..., description='Path for new DOCX.')
    content: str = Field(..., description='Initial text content.')

class EditDOCXTool(BaseTool):
    name = "edit_docx"
    description = "Replace text in a Word document."
    args_schema = EditDOCXParams
    
    def run(self, file_path: str, old_text: str, new_text: str) -> ToolResult:
        try:
            if not os.path.exists(file_path): return ToolResult(success=False, error="File not found")
            doc = Document(file_path)
            replaced = 0
            
            for p in doc.paragraphs:
                if old_text in p.text:
                    p.text = p.text.replace(old_text, new_text)
                    replaced += 1
            
            # Also check tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                             if old_text in p.text:
                                p.text = p.text.replace(old_text, new_text)
                                replaced += 1

            if replaced > 0:
                doc.save(file_path)
                return ToolResult(success=True, content=f"Replaced {replaced} occurrences.")
            else:
                return ToolResult(success=False, error="Text not found.")
        except Exception as e:
            return ToolResult(success=False, error=f"DOCX Error: {e}")

class CreateDOCXTool(BaseTool):
    name = "create_docx"
    description = "Create a new Word document."
    args_schema = CreateDOCXParams
    
    def run(self, file_path: str, content: str) -> ToolResult:
        try:
            doc = Document()
            doc.add_paragraph(content)
            doc.save(file_path)
            return ToolResult(success=True, content=f"Created {file_path}")
        except Exception as e:
             return ToolResult(success=False, error=f"DOCX Create Error: {e}")


# --- Excel ---
import openpyxl

class EditExcelParams(BaseModel):
    file_path: str = Field(..., description='Path to XLSX file.')
    sheet_name: str = Field("Sheet1", description='Target sheet.')
    cell: str = Field(..., description='Cell coordinate (e.g., A1).')
    value: str = Field(..., description='New value.')

class CreateExcelParams(BaseModel):
    file_path: str = Field(..., description='Path for new XLSX.')
    sheet_name: str = Field("Sheet1", description='Initial sheet name.')

class EditExcelTool(BaseTool):
    name = "edit_excel"
    description = "Write a value to a specific cell in an Excel sheet."
    args_schema = EditExcelParams
    
    def run(self, file_path: str, cell: str, value: str, sheet_name: str = "Sheet1") -> ToolResult:
        try:
            if not os.path.exists(file_path): return ToolResult(success=False, error="File not found")
            wb = openpyxl.load_workbook(file_path)
            
            if sheet_name not in wb.sheetnames:
                 wb.create_sheet(sheet_name)
            
            ws = wb[sheet_name]
            ws[cell] = value
            wb.save(file_path)
            return ToolResult(success=True, content=f"Updated {cell} in {sheet_name}.")
        except Exception as e:
             return ToolResult(success=False, error=f"Excel Error: {e}")

class CreateExcelTool(BaseTool):
    name = "create_excel"
    description = "Create a new Excel workbook."
    args_schema = CreateExcelParams
    
    def run(self, file_path: str, sheet_name: str = "Sheet1") -> ToolResult:
        try:
            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = sheet_name
            wb.save(file_path)
            return ToolResult(success=True, content=f"Created {file_path}")
        except Exception as e:
             return ToolResult(success=False, error=f"Excel Create Error: {e}")
