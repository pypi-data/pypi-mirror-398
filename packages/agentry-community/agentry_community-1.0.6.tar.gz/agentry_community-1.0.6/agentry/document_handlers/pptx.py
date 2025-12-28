from typing import Dict, Any
from .base import BaseDocumentHandler

class PPTXHandler(BaseDocumentHandler):
    """Handler for PowerPoint presentations (.pptx) using python-pptx with Ollama Vision fallback."""

    def __init__(self, file_path: str):
        super().__init__(file_path)
        self._prs = None
        self._text = ""
        self._metadata = {}

    def load(self) -> None:
        """Load and parse the PPTX file using a hybrid approach."""
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            raise RuntimeError("python-pptx is not installed.")

        try:
            self._prs = Presentation(self.file_path)
            
            text_content = []
            
            # Check availability
            try:
                from ..services.ollama_vision import OllamaVisionService
                VISION_READY = True
                from io import BytesIO
                from PIL import Image
            except ImportError:
                VISION_READY = False
            
            for i, slide in enumerate(self._prs.slides):
                slide_text_parts = []
                image_shapes = []
                
                # Check shapes for text and images
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text_parts.append(shape.text)
                    
                    if hasattr(shape, "shape_type") and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_shapes.append(shape)
                
                slide_raw_text = "\n".join(slide_text_parts)
                
                # Heuristic: If slide has images but very little text (< 50 chars), try OCR on the images
                if VISION_READY and image_shapes and len(slide_raw_text.strip()) < 50:
                    ocr_segments = []
                    for img_shape in image_shapes:
                        try:
                            image_blob = img_shape.image.blob
                            image = Image.open(BytesIO(image_blob))
                            text = OllamaVisionService.get_text_from_pil_image(image)
                            if text.strip():
                                ocr_segments.append(text)
                        except Exception:
                            pass
                    
                    if ocr_segments:
                        combined_ocr = "\n".join(ocr_segments)
                        text_content.append(f"--- Slide {i+1} (OllamaVision) ---\n{combined_ocr}")
                        self._metadata[f"slide_{i+1}_engine"] = "ollama_vision"
                        continue # Skip standard text append if we used OCR
                
                # Standard Text
                text_content.append(f"--- Slide {i+1} ---\n{slide_raw_text}")
            
            self._text = "\n\n".join(text_content)
            
            # Extract core properties
            core_props = self._prs.core_properties
            self._metadata.update({
                "author": core_props.author,
                "created": str(core_props.created),
                "modified": str(core_props.modified),
                "title": core_props.title,
                "subject": core_props.subject,
                "keywords": core_props.keywords,
                "last_modified_by": core_props.last_modified_by,
                "slide_count": len(self._prs.slides)
            })
            # Remove None values
            self._metadata = {k: v for k, v in self._metadata.items() if v is not None}

        except Exception as e:
            raise RuntimeError(f"Failed to load PPTX file {self.file_path}: {e}")

    def get_text(self) -> str:
        if self._prs is None:
            self.load()
        return self._text

    def get_metadata(self) -> Dict[str, Any]:
        if self._prs is None:
            self.load()
        return self._metadata
