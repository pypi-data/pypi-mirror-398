import json
import warnings  # Correct module for warnings, including PendingDeprecationWarning
import os
from pathlib import Path

from docx import Document
from PyPDF2 import PdfReader
from PIL import Image
import pytesseract
import uuid

from langchain.chat_models import init_chat_model
from langchain_text_splitters import RecursiveCharacterTextSplitter, CharacterTextSplitter
from openai import OpenAI
from openai import OpenAI
from .llm_wrapper import LLMClientWrapper  # Relative import
from dotenv import load_dotenv
import tempfile

from ..config.splitter_config import SplitterConfig, LLM_Structured_Output_Type
from ..rlchunker.inference import RLChunker
from ..services.propositional_index import PropositionalIndexer

load_dotenv(override=True)
# Initialize OpenAI client
client =  OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
from django.core.files.uploadedfile import UploadedFile

from enum import Enum
import pandas as pd
from pptx import Presentation
class SplitType(Enum):
    RECURSIVE = "RecursiveCharacterTextSplitter"
    CHARACTER = "CharacterTextSplitter"
    STANDARD = "standard"
    R_PRETRAINED_PROPOSITION = "RLBasedTextSplitterWithProposition"
    R_PRETRAINED = "RLBasedTextSplitter"

def extract_content_agnostic(file, filename=None):
    """
    Extract text content from a file.

    Supports:
    - PDF (.pdf)
    - Word (.docx)
    - Text (.txt)
    - Images (.png, .jpg, .jpeg, .tiff, .bmp)

    Parameters:
    - file: either a file path (str) or bytes (binary content)
    - filename: required if `file` is bytes, to determine extension
    """
    # Determine if file is path or binary
    if isinstance(file, str):
        filepath = file
        ext = os.path.splitext(filepath)[1].lower()
    elif isinstance(file, Path):
        filepath = str(file)
        ext = os.path.splitext(filepath)[1].lower()
    elif isinstance(file, bytes):
        if not filename:
            raise ValueError("filename must be provided if passing binary content")
        ext = os.path.splitext(filename)[1].lower()
        # Write bytes to a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(file)
            filepath = tmp.name
    else:
        raise TypeError("file must be a string path or bytes")

    # At this point, `filepath` is a valid file path on disk
    # TODO: implement your extraction logic based on `ext` and `filepath`
    # Example:
    # if ext == ".pdf":
    #     content = extract_pdf(filepath)
    # elif ext == ".docx":
    #     content = extract_docx(filepath)
    # ...

    text = load_file_by_type(ext, filepath)

    # If we created a temporary file, optionally delete it
    if isinstance(file, UploadedFile) and not hasattr(file, 'temporary_file_path'):
        try:
            os.remove(filepath)
        except Exception:
            pass

    return text.strip()


def extract_content(file):
    """
    Extract text content from a file.
    Supports:
    - PDF (.pdf)
    - Word (.docx)
    - Text (.txt)
    - Images (.png, .jpg, .jpeg, .tiff, .bmp)

    file: either a file path (str) or a Django UploadedFile object (request.FILES['file'])
    """
    # Determine if input is file path or UploadedFile
    if isinstance(file, UploadedFile):
        filename = file.name
        ext = os.path.splitext(filename)[1].lower()

        # Check if file is already on disk
        if hasattr(file, 'temporary_file_path'):
            filepath = file.temporary_file_path()
        else:
            # Save in-memory file to temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
                for chunk in file.chunks():
                    tmp.write(chunk)
                filepath = tmp.name
    else:
        # It's a file path
        filepath = file
        ext = os.path.splitext(filepath)[1].lower()

    text = load_file_by_type(ext, filepath)

    # If we created a temporary file, optionally delete it
    if isinstance(file, UploadedFile) and not hasattr(file, 'temporary_file_path'):
        try:
            os.remove(filepath)
        except Exception:
            pass

    return text.strip()


def load_file_by_type(ext, filepath):
    text = ""
    if ext == ".pdf":
        reader = PdfReader(filepath)
        text = "\n".join([p.extract_text() or "" for p in reader.pages])

    elif ext == ".docx":
        doc = Document(filepath)
        text = "\n".join([p.text for p in doc.paragraphs])

    elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img)

    elif ext == ".txt":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
    elif ext == ".json":
        with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
            data = json.load(f)
            # Convert JSON to text (pretty print or flatten)
            text = json.dumps(data, ensure_ascii=False, indent=2)
    # -------------------------
    # PPTX (PowerPoint)
    # -------------------------
    elif ext in [".pptx", ".ppt"]:
        pres = Presentation(filepath)
        slides_text = []
        for slide in pres.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            slides_text.append("\n".join(slide_text))
        text = "\n\n---- Slide Break ----\n\n".join(slides_text)

    # -------------------------
    # Excel (XLS / XLSX)
    # -------------------------
    elif ext in [".xlsx", ".xls"]:
        # Using pandas for convenience
        try:
            df_dict = pd.read_excel(filepath, sheet_name=None)
            all_sheets = []
            for sheet, df in df_dict.items():
                sheet_text = f"=== Sheet: {sheet} ===\n"
                sheet_text += df.to_string(index=False)
                all_sheets.append(sheet_text)
            text = "\n\n".join(all_sheets)
        except Exception as e:
            raise ValueError(f"Failed to read Excel file: {e}")
    else:
        raise ValueError(f"Unsupported file type: {ext}")
    return text


def split_text_by_config(text, splitter_config:SplitterConfig=None, binary_data=None):
    """Split text into chunks of N words."""
    if splitter_config is None:
        splitter_config =  SplitterConfig(chunk_size= 200, chunk_overlap= 0)
        words = text.split()
        return [" ".join(words[i:i + splitter_config.chunk_size]) for i in
                range(0, len(words), splitter_config.chunk_size)]
    else:
        """Split text into chunks of N words."""
        if splitter_config.split_type == SplitType.STANDARD.value:
            words = text.split()
            return [" ".join(words[i:i + splitter_config.chunk_size]) for i in
                    range(0, len(words), splitter_config.chunk_size)]

        elif splitter_config.split_type == SplitType.RECURSIVE.value:
            """Split text into chunks of N characters."""
            text_splitter = RecursiveCharacterTextSplitter(
                separators=splitter_config.separators,
                chunk_size=splitter_config.chunk_size,
                chunk_overlap=splitter_config.chunk_overlap,
            )
            chunked_content = text_splitter.split_text(text)
            return chunked_content
        elif splitter_config.split_type == SplitType.CHARACTER.value:
            """Split text into chunks of N characters."""
            text_splitter = CharacterTextSplitter(
                separators=splitter_config.separators,
                chunk_size=splitter_config.chunk_size,
                chunk_overlap=splitter_config.chunk_overlap,
            )
            chunked_content = text_splitter.split_text(text)
            return chunked_content

        elif splitter_config.split_type == SplitType.R_PRETRAINED_PROPOSITION.value:
            indexer = PropositionalIndexer(model_name="gpt-4o-mini")

            # Index directly from file
            sentences = indexer.index_file_content(text, "propositional_index.txt")

            # ✅ Combine all sentences into one big text
            combined_text = " ".join(sentences)
            # Initialize chunker once
            chunker = RLChunker(device="cpu", embedding_dim=384)

            # Chunk a single text

            chunked_content = chunker.chunk_text(combined_text,min_len=splitter_config.min_rl_chunk_size,max_len=splitter_config.max_rl_chunk_size)

            return chunked_content
        elif splitter_config.split_type == SplitType.R_PRETRAINED.value:

            # Initialize chunker once
            chunker = RLChunker(device="cpu", embedding_dim=384)

            # Chunk a single text

            chunked_content = chunker.chunk_text(text,min_len=splitter_config.min_rl_chunk_size,max_len=splitter_config.max_rl_chunk_size)

            return chunked_content
        else:
            words = text.split()
            return [" ".join(words[i:i + splitter_config.chunk_size]) for i in
                    range(0, len(words), splitter_config.chunk_size)]

import json
from langchain.schema import HumanMessage
import uuid

def process_with_llm(chunk, instructions=None, xclient=None):
    """
    Send a chunk to LLM and return structured JSON array.
    Expected format: [{"id": ..., "title": ..., "text": ...}, ...]
    """
    instructions = instructions or "Extract sections"

    # Combine chunk + instructions into one prompt
    prompt_text = f"""
    You are a helpful assistant that structures text into JSON sections.
    Take the following text and split it into sections based on the most important category headings.
     return a JSON array of objects with the following keys:
      - "id" (a UUID you generate)
      - "title" (the most important heading)
      - "text" (the remaining text under that heading)
    Return ONLY valid JSON, without extra text or backtick or markdown formatting.

    Text:
    {chunk}

    Instructions: {instructions}
    """

    # Use provided client or create new wrapper
    if xclient is None:
        xclient = init_chat_model(
            model="gpt-4o-mini",
            model_provider="openai",  # you can later swap to "anthropic", "google", etc.
            api_key=os.getenv("OPENAI_API_KEY")
        )


    # Call the LLM
    response_msg = xclient.predict_messages([HumanMessage(content=prompt_text)])
    response_text = response_msg.content

    # Parse JSON safely
    try:
        structured_data = json.loads(response_text)
        if isinstance(structured_data, str):
            # Sometimes LLM returns a JSON string inside quotes
            structured_data = json.loads(structured_data)
    except json.JSONDecodeError as e:
        print("LLM returned invalid JSON:", response_text)
        raise e

    for item in structured_data:
        if isinstance(item, dict) and "id" not in item:
            item["id"] = str(uuid.uuid4())

    return structured_data

def process_large_text(text, instructions,splitter_config:SplitterConfig=None,client=None):
    """Main function: split -> send to LLM -> collect results."""
    chunks = split_text_by_config(text, splitter_config=splitter_config)
    all_results = []
    if splitter_config.enableLLMTouchUp:
        if splitter_config.llm_structured_output_type == LLM_Structured_Output_Type.STANDARD:
            warnings.warn("bypassing LLM touch up for standard structured output")
            return chunks
        elif splitter_config.llm_structured_output_type == LLM_Structured_Output_Type.STRUCTURED_WITH_VECTOR_DB_ID_GENERATED:
            for chunk in chunks:
                structured = process_with_llm(chunk,instructions,client)
                # Ensure UUIDs exist
                for obj in structured:
                    if "id" not in obj:
                        obj["id"] = str(uuid.uuid4())
                all_results.extend(structured)

            return all_results
    else:
        return chunks



def prepare_chunked_text(file_path,file_name,instructions,chunk_size=200,splitter_config:SplitterConfig=None,client=None):
    content =extract_content_agnostic(file_path,file_name)
    results=process_large_text(content,instructions, splitter_config=splitter_config,client=client)
    print (results)
    return results

#this function takes a django file and extracts filename and byte content
def extract_file_details(uploaded_file):
    # 1. Get the filename
    filename = uploaded_file.name

    # 2. Get the file content as bytes
    file_bytes = uploaded_file.read()  # reads entire file into memory

    # Now you can call your extract_content function
    return filename, file_bytes




















































