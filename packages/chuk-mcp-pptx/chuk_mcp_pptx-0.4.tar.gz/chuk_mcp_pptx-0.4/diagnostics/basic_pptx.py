#!/usr/bin/env python3
"""
Test basic python-pptx functionality without our server
"""

from pptx import Presentation


def test_basic():
    """Create the simplest possible presentation"""
    print("Creating basic presentation with python-pptx directly...")

    # Create presentation
    prs = Presentation()

    # Add title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Test Presentation"
    subtitle.text = "Direct python-pptx test"

    # Save
    prs.save("test_basic.pptx")
    print("Saved to test_basic.pptx")

    # Try to read it back
    try:
        prs2 = Presentation("test_basic.pptx")
        print(f"Successfully read back: {len(prs2.slides)} slide(s)")
        return True
    except Exception as e:
        print(f"Error reading back: {e}")
        return False


if __name__ == "__main__":
    success = test_basic()
    if success:
        print("\n✅ Basic test passed - python-pptx is working")
        print("Try opening test_basic.pptx in PowerPoint")
    else:
        print("\n❌ Basic test failed - python-pptx issue")
