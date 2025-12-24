#!/usr/bin/env python3
"""
Test if files created by our server are valid
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "src"))

from pptx import Presentation
import zipfile


def check_pptx_file(filename):
    """Check if a PowerPoint file is valid"""
    print(f"\nChecking {filename}...")

    # Check if it's a valid ZIP file (PPTX files are ZIP archives)
    try:
        with zipfile.ZipFile(filename, "r") as z:
            print("  ✅ Valid ZIP structure")
            # Check for required PowerPoint files
            namelist = z.namelist()
            required_files = [
                "[Content_Types].xml",
                "_rels/.rels",
                "ppt/presentation.xml",
                "ppt/_rels/presentation.xml.rels",
            ]
            for req_file in required_files:
                if req_file in namelist:
                    print(f"  ✅ Found {req_file}")
                else:
                    print(f"  ❌ Missing {req_file}")
    except zipfile.BadZipFile:
        print("  ❌ Invalid ZIP structure - file is corrupted")
        return False
    except Exception as e:
        print(f"  ❌ Error checking ZIP: {e}")
        return False

    # Try to open with python-pptx
    try:
        prs = Presentation(filename)
        print("  ✅ python-pptx can read it")
        print(f"  📊 Contains {len(prs.slides)} slide(s)")

        # Check for charts
        chart_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_chart") and shape.has_chart:
                    chart_count += 1
        if chart_count > 0:
            print(f"  📈 Contains {chart_count} chart(s)")

        return True
    except Exception as e:
        print(f"  ❌ python-pptx error: {e}")
        return False


def test_server_creation():
    """Test creating a file with our server"""
    from chuk_mcp_pptx.sync_server import (
        pptx_create,
        pptx_add_title_slide,
        pptx_add_slide,
        pptx_save,
    )

    print("\n" + "=" * 60)
    print("Creating test file with server...")

    # Create presentation
    result = pptx_create(name="validity_test")
    print(f"  {result}")

    # Add title slide
    result = pptx_add_title_slide(title="Validity Test", subtitle="Testing file format")
    print(f"  {result}")

    # Add content slide
    result = pptx_add_slide(title="Test Content", content=["Point 1", "Point 2", "Point 3"])
    print(f"  {result}")

    # Save
    result = pptx_save(path="validity_test.pptx")
    print(f"  {result}")

    return "validity_test.pptx"


if __name__ == "__main__":
    print("PowerPoint File Validity Checker")
    print("=" * 60)

    # Test existing files
    test_files = [
        "test_basic.pptx",
        "async_test.pptx",
        "sales_analytics_2024.pptx",
        "chart_gallery.pptx",
    ]

    for filename in test_files:
        if os.path.exists(filename):
            check_pptx_file(filename)
        else:
            print(f"\n{filename} not found, skipping...")

    # Create and test a new file
    new_file = test_server_creation()
    check_pptx_file(new_file)

    print("\n" + "=" * 60)
    print("✅ Validity check complete")
