#!/usr/bin/env python3
"""
Test compatibility - create presentations with different features
to identify what might be causing the format error
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def create_simple_no_chart():
    """Create simplest presentation without charts"""
    prs = Presentation()

    # Just a title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Simple Test"
    slide.placeholders[1].text = "No charts"

    prs.save("test_1_simple.pptx")
    print("Created test_1_simple.pptx - basic slide only")


def create_with_text():
    """Create presentation with text content"""
    prs = Presentation()

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Text Test"

    # Content slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Content"
    content = slide.placeholders[1].text_frame
    content.text = "First bullet"
    p = content.add_paragraph()
    p.text = "Second bullet"

    prs.save("test_2_text.pptx")
    print("Created test_2_text.pptx - with text content")


def create_with_simple_chart():
    """Create presentation with a basic chart"""
    prs = Presentation()

    # Add slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank

    # Create simple chart data
    chart_data = CategoryChartData()
    chart_data.categories = ["A", "B", "C"]
    chart_data.add_series("Series 1", (1, 2, 3))

    # Add chart
    x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

    prs.save("test_3_chart.pptx")
    print("Created test_3_chart.pptx - with simple chart")


def create_with_complex_chart():
    """Create presentation with more complex chart"""
    prs = Presentation()

    # Add slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("2023", (10.5, 12.3, 14.7, 16.2))
    chart_data.add_series("2024", (11.2, 13.8, 15.9, 17.5))

    # Add chart
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(8), Inches(5), chart_data
    )

    # Add title to chart
    chart.chart.has_title = True
    chart.chart.chart_title.text_frame.text = "Quarterly Sales"

    prs.save("test_4_complex_chart.pptx")
    print("Created test_4_complex_chart.pptx - with complex chart")


def create_mixed_content():
    """Create presentation with mixed content"""
    prs = Presentation()

    # Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Mixed Content Test"

    # Text slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Agenda"
    content = slide2.placeholders[1].text_frame
    content.text = "Introduction"
    content.add_paragraph().text = "Data Analysis"
    content.add_paragraph().text = "Conclusions"

    # Chart slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = ["Jan", "Feb", "Mar"]
    chart_data.add_series("Sales", (100, 120, 140))
    slide3.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(1), Inches(1), Inches(8), Inches(5), chart_data
    )

    prs.save("test_5_mixed.pptx")
    print("Created test_5_mixed.pptx - with mixed content")


if __name__ == "__main__":
    print("Creating Compatibility Test Files")
    print("=" * 60)
    print("\nThis will create 5 test files with increasing complexity.")
    print("Try opening each in PowerPoint to identify which features cause issues.\n")

    create_simple_no_chart()
    create_with_text()
    create_with_simple_chart()
    create_with_complex_chart()
    create_mixed_content()

    print("\n" + "=" * 60)
    print("Test files created:")
    print("  1. test_1_simple.pptx - Just a title slide")
    print("  2. test_2_text.pptx - Title and content slides")
    print("  3. test_3_chart.pptx - Simple chart")
    print("  4. test_4_complex_chart.pptx - Chart with title and formatting")
    print("  5. test_5_mixed.pptx - Mixed content including charts")
    print("\nTry opening each file in PowerPoint to see which ones work.")
    print("This will help identify what's causing the format error.")
