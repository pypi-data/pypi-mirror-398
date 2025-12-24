#!/usr/bin/env python3
"""
Test chart_utils directly
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "src"))

from pptx import Presentation
from chuk_mcp_pptx.utilities.chart_utils import add_chart, add_pie_chart


def test_chart_utils():
    """Test our chart utilities"""
    print("Testing chart_utils functions...")

    # Create presentation
    prs = Presentation()

    # Add a blank slide
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)

    # Test add_chart
    print("\n1. Testing add_chart (column)...")
    try:
        chart_shape = add_chart(
            slide1,
            chart_type="column",
            left=1.0,
            top=1.0,
            width=8.0,
            height=4.5,
            categories=["Q1", "Q2", "Q3", "Q4"],
            series_data={"2023": [100, 120, 140, 160], "2024": [110, 130, 150, 170]},
            title="Quarterly Sales",
        )
        print(f"   ✅ Chart added: {chart_shape}")
    except Exception as e:
        print(f"   ❌ Error: {e}")
        import traceback

        traceback.print_exc()

    # Add another slide for pie chart
    slide2 = prs.slides.add_slide(slide_layout)

    # Test add_pie_chart
    print("\n2. Testing add_pie_chart...")
    try:
        pie_shape = add_pie_chart(
            slide2,
            left=2.0,
            top=1.0,
            width=6.0,
            height=4.5,
            categories=["Product A", "Product B", "Product C", "Product D"],
            values=[35, 25, 25, 15],
            title="Market Share",
        )
        print(f"   ✅ Pie chart added: {pie_shape}")
    except Exception as e:
        print(f"   ❌ Error: {e}")
        import traceback

        traceback.print_exc()

    # Save and check
    prs.save("test_charts.pptx")
    print("\n✅ Saved to test_charts.pptx")

    # Verify charts
    print("\nSlide 1 shapes:")
    for i, shape in enumerate(slide1.shapes):
        print(f"  Shape {i}: {type(shape).__name__}")
        if hasattr(shape, "has_chart") and shape.has_chart:
            print(f"    - Chart type: {shape.chart.chart_type}")

    print("\nSlide 2 shapes:")
    for i, shape in enumerate(slide2.shapes):
        print(f"  Shape {i}: {type(shape).__name__}")
        if hasattr(shape, "has_chart") and shape.has_chart:
            print(f"    - Chart type: {shape.chart.chart_type}")


if __name__ == "__main__":
    test_chart_utils()
