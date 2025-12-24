#!/usr/bin/env python3
"""
Minimal test to reproduce the PowerPoint format issue
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "src"))

from pptx import Presentation


def test_with_python_pptx_direct():
    """Create file directly with python-pptx"""
    print("Test 1: Direct python-pptx creation")
    prs = Presentation()

    # Add title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Direct Test"
    slide.placeholders[1].text = "Created with python-pptx directly"

    # Save
    prs.save("test_direct.pptx")
    print("  ✅ Saved test_direct.pptx")
    return "test_direct.pptx"


def test_with_our_server():
    """Create file using our server"""
    print("\nTest 2: Using our server")

    # Force sync mode to eliminate async issues
    os.environ["PPTX_ASYNC"] = "false"

    from chuk_mcp_pptx.sync_server import pptx_create, pptx_add_title_slide, pptx_save

    # Create presentation
    result = pptx_create(name="server_test")
    print(f"  {result}")

    # Add title slide
    result = pptx_add_title_slide(title="Server Test", subtitle="Created with our server")
    print(f"  {result}")

    # Save
    result = pptx_save(path="test_server.pptx")
    print(f"  {result}")

    return "test_server.pptx"


def test_manager_direct():
    """Test using PresentationManager directly"""
    print("\nTest 3: Using PresentationManager directly")

    from chuk_mcp_pptx.presentation_manager import PresentationManager

    manager = PresentationManager()

    # Create presentation
    result = manager.create("manager_test")
    print(f"  {result}")

    # Get presentation and add slide
    prs = manager.get("manager_test")
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Manager Test"
    slide.placeholders[1].text = "Created with PresentationManager"

    # Save using prs.save directly
    prs.save("test_manager.pptx")
    print("  ✅ Saved test_manager.pptx")

    return "test_manager.pptx"


def compare_files(file1, file2):
    """Compare two PPTX files"""
    import zipfile

    print(f"\nComparing {file1} vs {file2}:")

    # Check file sizes
    size1 = os.path.getsize(file1)
    size2 = os.path.getsize(file2)
    print(f"  Size: {size1} vs {size2} bytes")

    # Compare ZIP contents
    with zipfile.ZipFile(file1, "r") as z1, zipfile.ZipFile(file2, "r") as z2:
        names1 = set(z1.namelist())
        names2 = set(z2.namelist())

        if names1 == names2:
            print(f"  ✅ Same file structure ({len(names1)} files)")
        else:
            print("  ❌ Different file structure")
            only_in_1 = names1 - names2
            only_in_2 = names2 - names1
            if only_in_1:
                print(f"    Only in {file1}: {only_in_1}")
            if only_in_2:
                print(f"    Only in {file2}: {only_in_2}")


if __name__ == "__main__":
    print("=" * 60)
    print("Minimal Issue Reproduction Test")
    print("=" * 60)

    # Create test files
    file1 = test_with_python_pptx_direct()
    file2 = test_with_our_server()
    file3 = test_manager_direct()

    # Compare files
    print("\n" + "=" * 60)
    print("File Comparison")
    print("=" * 60)

    compare_files(file1, file2)
    compare_files(file1, file3)
    compare_files(file2, file3)

    print("\n" + "=" * 60)
    print("Test complete. Try opening each file:")
    print("  1. test_direct.pptx - Should work (direct python-pptx)")
    print("  2. test_server.pptx - Test if this works")
    print("  3. test_manager.pptx - Test if this works")
    print("\nIf any fail, we've isolated the issue.")
    print("=" * 60)
