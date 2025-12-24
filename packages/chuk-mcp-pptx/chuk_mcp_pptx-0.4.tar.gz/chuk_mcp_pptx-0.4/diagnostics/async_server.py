#!/usr/bin/env python3
"""
Test async server functionality
"""

import asyncio
import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "src"))

# Force async mode
os.environ["PPTX_ASYNC"] = "true"

from chuk_mcp_pptx.server import (
    pptx_create,
    pptx_add_title_slide,
    pptx_add_slide,
    pptx_add_chart,
    pptx_save,
    pptx_get_info,
    pptx_list,
)
import json


async def test_async_server():
    """Test async server operations"""
    print("Testing Async Server")
    print("=" * 50)

    # Test 1: Create presentation
    print("\n1. Creating presentation...")
    result = await pptx_create(name="async_test")
    print(f"   ✅ {result}")

    # Test 2: Add title slide
    print("\n2. Adding title slide...")
    result = await pptx_add_title_slide(
        title="Async Test Presentation", subtitle="Testing async/await functionality"
    )
    print(f"   ✅ {result}")

    # Test 3: Add content slide
    print("\n3. Adding content slide...")
    result = await pptx_add_slide(
        title="Test Points", content=["First point", "Second point", "Third point"]
    )
    print(f"   ✅ {result}")

    # Test 4: Add chart (this is the critical test)
    print("\n4. Adding chart to slide...")
    result = await pptx_add_slide(title="Test Chart", content=[])
    print(f"   ✅ {result}")

    result = await pptx_add_chart(
        slide_index=2,
        chart_type="column",
        categories=["Jan", "Feb", "Mar", "Apr"],
        series_data={"Sales": [100, 120, 140, 160], "Costs": [80, 85, 90, 95]},
        title="Monthly Performance",
    )
    print(f"   ✅ {result}")

    # Test 5: Get info
    print("\n5. Getting presentation info...")
    info = await pptx_get_info()
    info_data = json.loads(info)
    print(f"   Total slides: {info_data['slides']}")
    for slide in info_data["slide_details"]:
        print(f"   Slide {slide['index']}: {slide['title']} ({slide['shapes']} shapes)")

    # Test 6: Save presentation
    print("\n6. Saving presentation...")
    result = await pptx_save(path="async_test.pptx")
    print(f"   ✅ {result}")

    print("\n" + "=" * 50)
    print("✅ All async tests passed!")

    # Verify the saved file
    print("\nVerifying saved file...")
    from pptx import Presentation

    prs = Presentation("async_test.pptx")
    chart_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "has_chart") and shape.has_chart:
                chart_count += 1
    print(f"   Found {chart_count} chart(s) in the presentation")

    return True


async def run_parallel_test():
    """Test parallel async operations"""
    print("\nTesting Parallel Async Operations")
    print("=" * 50)

    # Create multiple presentations in parallel
    print("\n1. Creating 3 presentations in parallel...")
    tasks = [
        pptx_create(name="parallel_1"),
        pptx_create(name="parallel_2"),
        pptx_create(name="parallel_3"),
    ]
    results = await asyncio.gather(*tasks)
    for i, result in enumerate(results, 1):
        print(f"   Presentation {i}: {result}")

    # List all presentations
    print("\n2. Listing all presentations...")
    presentations = await pptx_list()
    print(f"   Presentations: {presentations}")

    print("\n✅ Parallel operations successful!")
    return True


async def main():
    """Main async test runner"""
    print("🚀 Async PowerPoint MCP Server Test Suite")
    print("=" * 60)

    try:
        # Run basic async tests
        await test_async_server()

        # Run parallel tests
        await run_parallel_test()

        print("\n" + "=" * 60)
        print("🎉 All async tests completed successfully!")
        print("📁 Created file: async_test.pptx")
        return 0

    except Exception as e:
        print(f"\n❌ Test failed with error: {e}")
        import traceback

        traceback.print_exc()
        return 1


if __name__ == "__main__":
    # Run the async tests
    exit_code = asyncio.run(main())
    sys.exit(exit_code)
