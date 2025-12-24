#!/usr/bin/env python3
"""
Debug script to test chart creation directly
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def test_direct_chart():
    """Test creating a chart directly with python-pptx"""
    print("Testing direct chart creation...")

    # Create presentation
    prs = Presentation()

    # Add a slide with title and content layout
    slide_layout = prs.slide_layouts[5]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Define chart data
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("2023", (10.2, 12.5, 15.2, 18.5))
    chart_data.add_series("2024", (15.1, 18.4, 21.3, 22.8))

    # Add chart to slide
    x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4.5)

    try:
        chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)
        print("✅ Chart added successfully!")

        # Add title
        chart.chart.has_title = True
        chart.chart.chart_title.text_frame.text = "Quarterly Revenue"
        print("✅ Chart title added!")

    except Exception as e:
        print(f"❌ Error adding chart: {e}")
        import traceback

        traceback.print_exc()

    # Save presentation
    prs.save("test_chart_debug.pptx")
    print("✅ Saved to test_chart_debug.pptx")

    # Check what's in the slide
    print(f"Slide has {len(slide.shapes)} shapes")
    for i, shape in enumerate(slide.shapes):
        print(f"  Shape {i}: {type(shape).__name__}")
        if hasattr(shape, "has_chart") and shape.has_chart:
            print(f"    - Has chart: {shape.chart.chart_type}")


if __name__ == "__main__":
    test_direct_chart()
