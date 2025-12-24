#!/usr/bin/env python3
"""
Platform Chat Showcase - Realistic platform-specific messaging UIs.

Demonstrates chat interfaces that look like real platforms:
- iMessage (iOS)
- Android Messages (Material Design)
- ChatGPT (OpenAI)
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.chat import (
    iMessageConversation,
    AndroidConversation,
    WhatsAppConversation,
    ChatGPTConversation,
    SlackConversation,
    TeamsConversation,
    FacebookMessengerConversation,
    AIMConversation,
    MSNConversation,
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_imessage_slide(prs, theme):
    """Showcase iMessage-style conversation."""
    print("  • Creating iMessage demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "iMessage (iOS)"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {"text": "Hey! Are you free for lunch?", "variant": "received", "timestamp": "11:30 AM"},
        {"text": "Yes! Where do you want to go?", "variant": "sent", "timestamp": "11:31 AM"},
        {
            "text": "How about that new Italian place?",
            "variant": "received",
            "timestamp": "11:32 AM",
        },
        {"text": "Perfect! See you at noon", "variant": "sent", "timestamp": "11:33 AM"},
        {"text": "👍", "variant": "received"},
    ]

    conversation = iMessageConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_android_messages_slide(prs, theme):
    """Showcase Android Messages-style conversation."""
    print("  • Creating Android Messages demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Android Messages"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {
            "text": "Meeting in 10 minutes!",
            "sender": "Sarah",
            "variant": "received",
            "timestamp": "2:50 PM",
        },
        {"text": "On my way", "variant": "sent", "timestamp": "2:51 PM"},
        {"text": "Should I bring my laptop?", "variant": "sent", "timestamp": "2:51 PM"},
        {
            "text": "Yes, we'll need it for the presentation",
            "sender": "Sarah",
            "variant": "received",
            "timestamp": "2:52 PM",
        },
        {"text": "Got it 👍", "variant": "sent", "timestamp": "2:52 PM"},
    ]

    conversation = AndroidConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_chatgpt_slide(prs, theme):
    """Showcase ChatGPT-style conversation."""
    print("  • Creating ChatGPT demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "ChatGPT Interface"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation using ChatGPTConversation component
    messages = [
        {"text": "Explain recursion in simple terms", "variant": "user"},
        {
            "text": "Recursion is when a function calls itself to solve a problem by breaking it into smaller, similar problems. Think of it like Russian nesting dolls - each doll contains a smaller version of itself until you reach the smallest one.",
            "variant": "assistant",
        },
        {"text": "Can you show a simple example in Python?", "variant": "user"},
        {
            "text": "Here's a classic example - calculating factorial:\n\ndef factorial(n):\n    if n == 0:\n        return 1\n    return n * factorial(n-1)",
            "variant": "assistant",
        },
    ]

    conversation = ChatGPTConversation(messages, spacing=0.0, theme=theme.__dict__)
    conversation.render(slide, left=0.5, top=1.7, width=9.0)


def create_whatsapp_slide(prs, theme):
    """Showcase WhatsApp-style conversation."""
    print("  • Creating WhatsApp demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "WhatsApp"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {"text": "Hey! Are we still meeting today?", "variant": "received", "timestamp": "10:15"},
        {"text": "Yes! 3pm at the coffee shop", "variant": "sent", "timestamp": "10:16"},
        {"text": "Perfect, see you then!", "variant": "received", "timestamp": "10:17"},
        {
            "text": "Bringing my laptop for the presentation",
            "variant": "sent",
            "timestamp": "10:18",
        },
        {"text": "👍 Great idea", "variant": "received", "timestamp": "10:19"},
    ]

    conversation = WhatsAppConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_whatsapp_group_chat(prs, theme):
    """WhatsApp group conversation."""
    print("  • Creating WhatsApp group chat...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "WhatsApp - Group Chat"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Group conversation with sender names
    messages = [
        {
            "text": "Who's bringing the snacks?",
            "sender": "Alice",
            "variant": "received",
            "timestamp": "14:00",
        },
        {"text": "I can bring chips and drinks!", "variant": "sent", "timestamp": "14:01"},
        {
            "text": "I'll bring pizza 🍕",
            "sender": "Bob",
            "variant": "received",
            "timestamp": "14:02",
        },
        {
            "text": "Perfect! I'll bring dessert",
            "sender": "Carol",
            "variant": "received",
            "timestamp": "14:03",
        },
        {"text": "This is going to be great!", "variant": "sent", "timestamp": "14:04"},
    ]

    conversation = WhatsAppConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_comparison_slide(prs, theme):
    """Show mobile platforms side by side."""
    print("  • Creating mobile platform comparison...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Mobile Platforms Comparison"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Row 1: iMessage and Android
    imsg_messages = [{"text": "iMessage (iOS)", "variant": "sent"}]
    imsg_conv = iMessageConversation(imsg_messages, theme=theme.__dict__)
    imsg_conv.render(slide, left=0.5, top=2.0, width=4.5)

    android_messages = [{"text": "Android Messages", "variant": "sent"}]
    android_conv = AndroidConversation(android_messages, theme=theme.__dict__)
    android_conv.render(slide, left=5.0, top=2.0, width=4.5)

    # Row 2: WhatsApp and Facebook Messenger
    whatsapp_messages = [{"text": "WhatsApp", "variant": "sent", "timestamp": "10:30"}]
    whatsapp_conv = WhatsAppConversation(whatsapp_messages, theme=theme.__dict__)
    whatsapp_conv.render(slide, left=0.5, top=3.3, width=4.5)

    fb_messages = [{"text": "Facebook Messenger", "variant": "sent"}]
    fb_conv = FacebookMessengerConversation(fb_messages, theme=theme.__dict__)
    fb_conv.render(slide, left=5.0, top=3.3, width=4.5)


def create_slack_slide(prs, theme):
    """Slack workspace conversation."""
    print("  • Creating Slack demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Slack - Workplace Chat"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {
            "text": "Quick standup in 5 minutes!",
            "sender": "Sarah",
            "timestamp": "9:00 AM",
            "avatar_text": "SM",
        },
        {
            "text": "On my way",
            "sender": "John",
            "timestamp": "9:01 AM",
            "avatar_text": "JD",
            "reactions": ["👍 2"],
        },
        {
            "text": "Can someone share the Zoom link?",
            "sender": "Mike",
            "timestamp": "9:02 AM",
            "avatar_text": "MC",
        },
        {
            "text": "Just posted it in the calendar event",
            "sender": "Sarah",
            "timestamp": "9:03 AM",
            "avatar_text": "SM",
            "reactions": ["✅ 3"],
        },
    ]

    conversation = SlackConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_chatgpt_code_example(prs, theme):
    """ChatGPT with code examples."""
    print("  • Creating ChatGPT code example...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "ChatGPT - Code Assistant"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    messages = [
        {"text": "How do I reverse a string in Python?", "variant": "user"},
        {
            "text": "There are several ways! The simplest is using slicing:\n\ntext = 'hello'\nreversed_text = text[::-1]\nprint(reversed_text)  # 'olleh'",
            "variant": "assistant",
        },
        {"text": "What about reversing a list?", "variant": "user"},
        {
            "text": "For lists, you can use .reverse() method or reversed():\n\nmy_list = [1, 2, 3]\nmy_list.reverse()  # [3, 2, 1]",
            "variant": "assistant",
        },
    ]

    conversation = ChatGPTConversation(messages, spacing=0.0, theme=theme.__dict__)
    conversation.render(slide, left=0.5, top=1.7, width=9.0)


def create_teams_slide(prs, theme):
    """Teams enterprise conversation."""
    print("  • Creating Teams demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Microsoft Teams - Enterprise Chat"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {
            "text": "Meeting notes have been shared",
            "sender": "Sarah Johnson",
            "timestamp": "Today at 10:30 AM",
            "avatar_text": "SJ",
        },
        {
            "text": "Great presentation today!",
            "sender": "Mike Chen",
            "timestamp": "Today at 10:35 AM",
            "avatar_text": "MC",
            "reactions": ["👍 5", "❤️ 2"],
        },
        {
            "text": "Thanks everyone! Let's sync up tomorrow",
            "sender": "Sarah Johnson",
            "timestamp": "Today at 10:40 AM",
            "avatar_text": "SJ",
            "reply_count": 3,
        },
    ]

    conversation = TeamsConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_facebook_slide(prs, theme):
    """Facebook Messenger conversation."""
    print("  • Creating Facebook Messenger demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Facebook Messenger"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {"text": "Hey! Want to grab coffee?", "variant": "received", "avatar_text": "JS"},
        {"text": "Sure! When works for you?", "variant": "sent"},
        {"text": "How about tomorrow at 3pm?", "variant": "received", "avatar_text": "JS"},
        {"text": "Perfect! See you then", "variant": "sent"},
    ]

    conversation = FacebookMessengerConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_aol_slide(prs, theme):
    """AOL Instant Messenger nostalgic conversation."""
    print("  • Creating AOL IM demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "AOL Instant Messenger (Nostalgic)"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {
            "text": "Hey! What's up?",
            "screen_name": "sk8rgrl2004",
            "variant": "received",
            "timestamp": "5:30 PM",
        },
        {
            "text": "Not much, you?",
            "screen_name": "xXCoolDude2003Xx",
            "variant": "sent",
            "timestamp": "5:31 PM",
        },
        {
            "text": "Want to hang out later?",
            "screen_name": "sk8rgrl2004",
            "variant": "received",
            "timestamp": "5:32 PM",
        },
        {
            "text": "Yeah! Let's go to the mall",
            "screen_name": "xXCoolDude2003Xx",
            "variant": "sent",
            "timestamp": "5:33 PM",
        },
    ]

    conversation = AIMConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_msn_slide(prs, theme):
    """MSN Messenger nostalgic conversation."""
    print("  • Creating MSN Messenger demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "MSN Messenger (Nostalgic)"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Conversation
    messages = [
        {
            "text": "Hey! What are you up to?",
            "display_name": "AwesomeGirl",
            "variant": "received",
            "timestamp": "21:00",
            "emoticon": "😊",
        },
        {
            "text": "Not much, just listening to music",
            "display_name": "CoolGuy",
            "variant": "sent",
            "timestamp": "21:01",
            "emoticon": "🎵",
        },
        {
            "text": "Nice! What are you listening to?",
            "display_name": "AwesomeGirl",
            "variant": "received",
            "timestamp": "21:02",
        },
        {
            "text": "The new album that just came out",
            "display_name": "CoolGuy",
            "variant": "sent",
            "timestamp": "21:03",
            "emoticon": "🎧",
        },
    ]

    conversation = MSNConversation(messages, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def main():
    """Generate platform chat showcase presentation."""
    print("\n💬 Creating Platform Chat Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    # Mobile platforms
    create_imessage_slide(prs, theme)
    create_android_messages_slide(prs, theme)
    create_whatsapp_slide(prs, theme)
    create_whatsapp_group_chat(prs, theme)
    create_facebook_slide(prs, theme)

    # AI & Workplace platforms
    create_chatgpt_slide(prs, theme)
    create_chatgpt_code_example(prs, theme)
    create_slack_slide(prs, theme)
    create_teams_slide(prs, theme)

    # Nostalgic/Legacy platforms
    create_aol_slide(prs, theme)
    create_msn_slide(prs, theme)

    # Comparison
    create_comparison_slide(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "platform_chat_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n💬 Platform Showcase Features:")
    print("  📱 Mobile Platforms:")
    print("     • iMessage (iOS) - Blue/gray bubbles, iOS design")
    print("     • Android Messages - RCS blue, Material Design")
    print("     • WhatsApp - Green bubbles with checkmarks")
    print("     • Facebook Messenger - Blue bubbles with avatars")
    print("  💼 AI & Workplace:")
    print("     • ChatGPT - AI assistant interface")
    print("     • Slack - Workspace chat with reactions")
    print("     • Microsoft Teams - Enterprise messaging")
    print("  🕰️  Nostalgic/Legacy:")
    print("     • AOL Instant Messenger - Classic screen names")
    print("     • MSN Messenger - 'says:' format with emoticons")
    print("\n💡 Demonstrates:")
    print("  • Platform-specific styling and authentic colors")
    print("  • Automatic spacing with conversation components")
    print("  • Real-world messaging scenarios")
    print("  • Modular platform organization (one file per platform)")
    print("  • No hardcoded positioning - all calculated automatically")


if __name__ == "__main__":
    main()
