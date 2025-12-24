#!/usr/bin/env python3
"""
Themes Showcase - Comprehensive demonstration of the theme system.
Shows all available themes with components, charts, and styling.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.card import Card, MetricCard
from chuk_mcp_pptx.components.core.badge import Badge
from chuk_mcp_pptx.components.core.button import Button
from chuk_mcp_pptx.components.charts import ColumnChart, PieChart
from chuk_mcp_pptx.components.core import Container, Grid, Stack
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_theme_overview_slide(prs, theme, is_first=False):
    """Create overview slide for a theme."""
    layout = prs.slide_layouts[0] if is_first else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    theme.apply_to_slide(slide)

    # Title
    if is_first:
        slide.shapes.title.text = "Theme System Showcase"
        slide.placeholders[1].text = f"Demonstrating {theme.name} and all available themes"
    else:
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = f"{theme.name.title()} Theme"
            title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color(
                "foreground.DEFAULT"
            )

    # Theme info card
    container = Container(size="md", padding="md", center=True)
    bounds = container.render(slide, top=2.5 if not is_first else 3.0)

    card = Card(variant="elevated", theme=theme.__dict__)
    card.add_child(Card.Title(f"{theme.name.title()}"))
    card.add_child(
        Card.Description(
            f"Mode: {theme.mode} | Primary: {theme.primary_hue} | Font: {theme.font_family}"
        )
    )
    card.render(slide, left=bounds["left"], top=bounds["top"], width=bounds["width"])


def create_theme_components_slide(prs, theme):
    """Show all core components in this theme."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = f"{theme.name.title()} - Components"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Layout
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    # Badges
    badges = [
        ("Primary", "default", 0),
        ("Success", "success", 2),
        ("Warning", "warning", 4),
        ("Error", "destructive", 6),
        ("Outline", "outline", 8),
    ]

    for text, variant, col_start in badges:
        pos = grid.get_cell(col_span=2, col_start=col_start, row_start=0)
        Badge(text=text, variant=variant, theme=theme.__dict__).render(
            slide, left=pos["left"] + 0.1, top=pos["top"] + 0.1
        )

    # Buttons
    buttons = [
        Button("Default", variant="default", size="sm", theme=theme.__dict__),
        Button("Secondary", variant="secondary", size="sm", theme=theme.__dict__),
        Button("Outline", variant="outline", size="sm", theme=theme.__dict__),
        Button("Ghost", variant="ghost", size="sm", theme=theme.__dict__),
    ]

    stack = Stack(direction="horizontal", gap="md", align="start")
    stack.render_children(
        slide,
        buttons,
        left=bounds["left"],
        top=bounds["top"] + 1.0,
        item_width=2.0,
        item_height=0.4,
    )

    # Cards
    card_variants = [
        ("default", "Default", 0),
        ("outlined", "Outlined", 4),
        ("elevated", "Elevated", 8),
    ]

    for variant, title, col_start in card_variants:
        pos = grid.get_cell(col_span=4, col_start=col_start, row_start=2)
        card = Card(variant=variant, theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description(f"{variant} card variant"))
        card.render(slide, **pos)


def create_theme_dashboard_slide(prs, theme):
    """Create a dashboard example with this theme."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = f"{theme.name.title()} - Dashboard"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Layout
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    # Metrics
    metrics = [
        ("Revenue", "$1.2M", "+15%", "up", 0),
        ("Users", "45K", "+8%", "up", 4),
        ("Growth", "23%", "+3%", "up", 8),
    ]

    for label, value, change, trend, col_start in metrics:
        pos = grid.get_cell(col_span=4, col_start=col_start, row_start=0)
        MetricCard(
            label=label, value=value, change=change, trend=trend, theme=theme.__dict__
        ).render(slide, **pos)

    # Main content card
    main_pos = grid.get_cell(col_span=8, col_start=0, row_start=1)
    main_card = Card(variant="elevated", theme=theme.__dict__)
    main_card.add_child(Card.Title("Analytics Dashboard"))
    main_card.add_child(Card.Description("Key performance indicators and metrics at a glance"))
    main_card.render(slide, **main_pos)

    # Sidebar
    sidebar_pos = grid.get_cell(col_span=4, col_start=8, row_start=1)
    sidebar_card = Card(variant="outlined", theme=theme.__dict__)
    sidebar_card.add_child(Card.Title("Actions"))
    sidebar_card.add_child(Card.Description("Quick links"))
    sidebar_card.render(slide, **sidebar_pos)


async def create_theme_charts_slide(prs, theme):
    """Show chart examples in this theme."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = f"{theme.name.title()} - Charts"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Column Chart
    column_chart = ColumnChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        series={"Sales": [100, 120, 140, 160], "Profit": [20, 25, 30, 35]},
        title="Quarterly Performance",
        theme=theme.__dict__,
    )
    await column_chart.render(slide, left=0.5, top=2.0, width=4.5, height=3.5)

    # Pie Chart
    pie_chart = PieChart(
        categories=["Product A", "Product B", "Product C"],
        values=[45, 30, 25],
        title="Market Share",
        theme=theme.__dict__,
    )
    await pie_chart.render(slide, left=5.2, top=2.0, width=4.0, height=3.5)


def create_theme_comparison_slide(prs, theme_manager):
    """Create a slide comparing all themes."""
    print("  • Creating theme comparison slide...")

    # Get all themes
    all_themes = [
        "dark",
        "dark-blue",
        "dark-violet",
        "dark-green",
        "light",
        "corporate",
        "light-warm",
    ]

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    main_theme = theme_manager.get_theme("dark-violet")
    main_theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "All Available Themes"
        title_shape.text_frame.paragraphs[0].font.color.rgb = main_theme.get_color(
            "foreground.DEFAULT"
        )

    # Grid layout for theme cards
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)
    grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)

    # Display theme cards
    for idx, theme_name in enumerate(all_themes):
        row = idx // 4
        col = (idx % 4) * 3

        if col < 12:  # Ensure we don't exceed grid
            pos = grid.get_cell(col_span=3, col_start=col, row_start=row)

            theme = theme_manager.get_theme(theme_name)
            card = Card(variant="outlined", theme=theme.__dict__)
            card.add_child(Card.Title(theme_name.title()))
            card.add_child(Card.Description(f"{theme.mode} mode"))
            card.render(slide, **pos)


async def main():
    """Generate comprehensive themes showcase."""
    print("\n🎨 Creating Themes Showcase")
    print("=" * 70)

    # Initialize
    theme_manager = ThemeManager()

    # Featured themes to showcase in detail (diverse visual variety)
    featured_themes = [
        "dark-violet",  # Modern, vibrant purple
        "dark-blue",  # Professional blue
        "dark-green",  # Nature, calming green
        "light",  # Clean, minimal light
        "corporate",  # Business, professional light
    ]

    for theme_name in featured_themes:
        print(f"\n📊 Creating {theme_name} theme showcase...")

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        theme = theme_manager.get_theme(theme_name)

        # Create slides for this theme
        is_first = theme_name == featured_themes[0]
        create_theme_overview_slide(prs, theme, is_first=is_first)
        create_theme_components_slide(prs, theme)
        create_theme_dashboard_slide(prs, theme)
        await create_theme_charts_slide(prs, theme)

        # Save individual theme showcase
        output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
        os.makedirs(output_dir, exist_ok=True)
        output_path = os.path.join(output_dir, f"theme_{theme_name}.pptx")
        prs.save(output_path)

        print(f"   ✅ Created {output_path}")
        print(f"      Slides: {len(prs.slides)} | Theme: {theme.name}")

    # Create comprehensive comparison presentation
    print("\n📊 Creating comprehensive themes comparison...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    create_theme_comparison_slide(prs, theme_manager)

    # Add detailed slides for each theme
    all_themes = [
        "dark",
        "dark-blue",
        "dark-violet",
        "dark-green",
        "light",
        "corporate",
        "light-warm",
    ]

    for theme_name in all_themes:
        theme = theme_manager.get_theme(theme_name)
        create_theme_overview_slide(prs, theme)
        create_theme_components_slide(prs, theme)

    output_path = os.path.join(output_dir, "themes_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print("\n🎨 Themes Showcased:")
    print("  • Dark themes: dark, dark-blue, dark-violet, dark-green")
    print("  • Light themes: light, corporate, light-warm")
    print("\n💡 Each theme demonstrates:")
    print("  • Color system and tokens")
    print("  • Component variants")
    print("  • Dashboard layouts")
    print("  • Chart styling")
    print("  • Consistent design language")


if __name__ == "__main__":
    import asyncio

    asyncio.run(main())
