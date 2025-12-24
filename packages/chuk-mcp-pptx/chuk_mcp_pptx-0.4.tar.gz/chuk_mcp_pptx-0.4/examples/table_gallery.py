#!/usr/bin/env python3
"""
Table Gallery - Showcases table components with different themes and variants.
Demonstrates various table types for business presentations.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.table import Table
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_table_gallery(theme_name: str = "corporate"):
    """Create a table gallery with specified theme."""

    print(f"\n📊 Creating Table Gallery ({theme_name} theme)")
    print("=" * 50)

    # Initialize presentation and theme
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    theme_manager = ThemeManager()
    theme = theme_manager.get_theme(theme_name)

    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    theme.apply_to_slide(slide1)

    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = "Table Gallery"
    subtitle.text = f"Professional tables with {theme_name} theme\nData presentation components"

    # ==========================================================================
    # SLIDE 2: FINANCIAL DATA
    # ==========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide2)

    title_shape = slide2.shapes.title
    title_shape.text = "Financial Report"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    financial_table = Table(
        headers=["Metric", "Q1", "Q2", "Q3", "Q4", "YoY"],
        data=[
            ["Revenue", "$12.4M", "$13.8M", "$14.2M", "$15.6M", "+28%"],
            ["Gross Profit", "$4.8M", "$5.4M", "$5.8M", "$6.2M", "+31%"],
            ["EBITDA", "$1.6M", "$2.0M", "$2.3M", "$2.6M", "+45%"],
            ["Net Income", "$1.2M", "$1.5M", "$1.8M", "$2.1M", "+52%"],
        ],
        variant="default",
        size="md",
        theme=theme,
    )
    financial_table.render(slide2, left=1.5, top=2.0, width=7.0, height=3.0)

    # ==========================================================================
    # SLIDE 3: PRODUCT COMPARISON
    # ==========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide3)

    title_shape = slide3.shapes.title
    title_shape.text = "Product Comparison"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    comparison_table = Table(
        headers=["Feature", "Basic", "Pro", "Enterprise"],
        data=[
            ["Users", "5", "50", "Unlimited"],
            ["Storage", "10 GB", "100 GB", "Unlimited"],
            ["Support", "Email", "Priority", "24/7 Phone"],
            ["API Access", "No", "Yes", "Yes"],
            ["Price/Month", "$9.99", "$49.99", "$199.99"],
        ],
        variant="bordered",
        size="md",
        theme=theme,
    )
    comparison_table.render(slide3, left=2.0, top=2.0, width=6.0, height=3.5)

    # ==========================================================================
    # SLIDE 4: PROJECT TIMELINE
    # ==========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide4)

    title_shape = slide4.shapes.title
    title_shape.text = "Project Timeline"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    timeline_table = Table(
        headers=["Phase", "Start", "End", "Duration", "Status"],
        data=[
            ["Planning", "Jan 1", "Jan 31", "4 weeks", "Complete"],
            ["Design", "Feb 1", "Feb 28", "4 weeks", "Complete"],
            ["Development", "Mar 1", "May 31", "12 weeks", "In Progress"],
            ["Testing", "Jun 1", "Jun 30", "4 weeks", "Planned"],
            ["Launch", "Jul 1", "Jul 15", "2 weeks", "Planned"],
        ],
        variant="striped",
        size="md",
        theme=theme,
    )
    timeline_table.render(slide4, left=1.0, top=2.0, width=8.0, height=3.5)

    # ==========================================================================
    # SLIDE 5: PERFORMANCE METRICS
    # ==========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide5)

    title_shape = slide5.shapes.title
    title_shape.text = "Performance Metrics"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    metrics_table = Table(
        headers=["Department", "Target", "Actual", "Achievement"],
        data=[
            ["Sales", "500 units", "547 units", "109%"],
            ["Marketing", "10K leads", "12.3K leads", "123%"],
            ["Support", "< 2hr", "1.5hr avg", "133%"],
            ["Engineering", "95% uptime", "99.2%", "104%"],
            ["Finance", "< 5% var", "3.2% var", "136%"],
        ],
        variant="default",
        size="md",
        theme=theme,
    )
    metrics_table.render(slide5, left=1.5, top=2.0, width=7.0, height=3.5)

    # ==========================================================================
    # SLIDE 6: BUDGET ALLOCATION
    # ==========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide6)

    title_shape = slide6.shapes.title
    title_shape.text = "Budget Allocation"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    budget_table = Table(
        headers=["Category", "Q1", "Q2", "Q3", "Q4", "Total"],
        data=[
            ["R&D", "$2.5M", "$2.8M", "$3.0M", "$3.2M", "$11.5M"],
            ["Sales", "$1.8M", "$2.0M", "$2.2M", "$2.0M", "$8.0M"],
            ["Operations", "$2.0M", "$2.0M", "$2.0M", "$2.0M", "$8.0M"],
            ["Support", "$1.0M", "$1.0M", "$1.1M", "$1.1M", "$4.2M"],
            ["Total", "$7.3M", "$7.8M", "$8.3M", "$8.3M", "$31.7M"],
        ],
        variant="bordered",
        size="md",
        theme=theme,
    )
    budget_table.render(slide6, left=1.0, top=2.0, width=8.0, height=3.5)

    # ==========================================================================
    # SLIDE 7: TABLE VARIANTS
    # ==========================================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide7)

    title_shape = slide7.shapes.title
    title_shape.text = "Table Variants"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Default variant (small)
    default_table = Table(
        headers=["Type", "Style"],
        data=[["Default", "Standard"], ["Data", "Format"]],
        variant="default",
        size="sm",
        theme=theme,
    )
    default_table.render(slide7, left=0.5, top=2.0, width=4.5, height=1.5)

    # Bordered variant (small)
    bordered_table = Table(
        headers=["Type", "Style"],
        data=[["Bordered", "Bold"], ["Data", "Format"]],
        variant="bordered",
        size="sm",
        theme=theme,
    )
    bordered_table.render(slide7, left=5.5, top=2.0, width=4.0, height=1.5)

    # Striped variant (small)
    striped_table = Table(
        headers=["Type", "Style"],
        data=[["Striped", "Alternating"], ["Data", "Rows"]],
        variant="striped",
        size="sm",
        theme=theme,
    )
    striped_table.render(slide7, left=0.5, top=4.0, width=4.5, height=1.5)

    # Minimal variant (small)
    minimal_table = Table(
        headers=["Type", "Style"],
        data=[["Minimal", "Clean"], ["Data", "Simple"]],
        variant="minimal",
        size="sm",
        theme=theme,
    )
    minimal_table.render(slide7, left=5.5, top=4.0, width=4.0, height=1.5)

    return prs


def main():
    """Create table galleries for key themes."""
    print("\n🎨 Creating Table Galleries")
    print("=" * 70)

    themes_to_create = ["corporate", "dark", "ocean"]

    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)

    created_files = []

    for theme_name in themes_to_create:
        print(f"\nCreating {theme_name} table gallery...")

        try:
            prs = create_table_gallery(theme_name)

            filename = f"table_gallery_{theme_name}.pptx"
            output_path = os.path.join(output_dir, filename)
            prs.save(output_path)

            created_files.append((theme_name, filename))
            print(f"  ✅ Created {filename}")

        except Exception as e:
            print(f"  ❌ Error creating {theme_name}: {e}")

    # Summary
    print(f"\n🎉 Created {len(created_files)} table galleries!")
    print("\n📊 Files created:")
    for theme_name, filename in created_files:
        print(f"  • {theme_name}: {filename}")

    print("\n📁 All files saved to: outputs/")
    print("\n💡 Table types demonstrated:")
    print("    • Financial reports")
    print("    • Product comparisons")
    print("    • Project timelines")
    print("    • Performance metrics")
    print("    • Budget allocations")
    print("    • Multiple variants (default, bordered, striped, minimal)")


if __name__ == "__main__":
    main()
