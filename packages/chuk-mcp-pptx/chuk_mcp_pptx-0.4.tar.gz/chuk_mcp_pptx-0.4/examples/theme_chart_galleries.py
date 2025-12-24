#!/usr/bin/env python3
"""
Theme Chart Galleries - Creates individual chart galleries for each theme.
Demonstrates how the same chart components look with different themes.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

# Import chart components
from chuk_mcp_pptx.components.charts import (
    ColumnChart,
    BarChart,
    LineChart,
    AreaChart,
    PieChart,
    DoughnutChart,
    ScatterChart,
    BubbleChart,
    RadarChart,
    GaugeChart,
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_chart_gallery_for_theme(theme_name: str, theme_obj):
    """Create a comprehensive chart gallery for a specific theme."""

    print(f"\n📊 Creating {theme_name} Chart Gallery")
    print("=" * 50)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ==========================================================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    theme_obj.apply_to_slide(slide1)

    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]

    title.text = f"{theme_obj.name} Chart Gallery"
    subtitle.text = f"Beautiful charts with {theme_name} theme\nComponent-based design system"

    # ==========================================================================
    # SLIDE 2: COLUMN & BAR CHARTS
    # ==========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide2)

    title_shape = slide2.shapes.title
    title_shape.text = "Column & Bar Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")

    # Clustered column chart (larger)
    column_chart = ColumnChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        series={"Revenue": [125, 148, 172, 195], "Profit": [28, 35, 42, 48]},
        variant="clustered",
        title="Quarterly Performance",
        theme=theme_obj,
    )
    column_chart.render(slide2, left=0.5, top=1.8, width=4.2, height=4.5)

    # Horizontal bar chart (larger)
    bar_chart = BarChart(
        categories=["AI/ML", "Cloud", "Security", "IoT"],
        series={"Adoption": [85, 92, 78, 65], "Investment": [70, 88, 82, 58]},
        variant="clustered",
        title="Technology Trends",
        theme=theme_obj,
    )
    bar_chart.render(slide2, left=5.3, top=1.8, width=4.2, height=4.5)

    # ==========================================================================
    # SLIDE 3: LINE & AREA CHARTS
    # ==========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide3)

    title_shape = slide3.shapes.title
    title_shape.text = "Line & Area Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")

    # Multi-series line chart (simplified)
    line_chart = LineChart(
        categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
        series={
            "Users": [1200, 1350, 1480, 1620, 1850, 2100],
            "Sessions": [800, 920, 1050, 1180, 1350, 1520],
        },
        variant="smooth",
        title="Platform Growth",
        theme=theme_obj,
    )
    line_chart.render(slide3, left=0.5, top=1.8, width=4.2, height=4.5)

    # Stacked area chart (simplified)
    area_chart = AreaChart(
        categories=["2020", "2021", "2022", "2023", "2024"],
        series={"Enterprise": [45, 52, 68, 82, 95], "SMB": [32, 38, 45, 58, 72]},
        variant="stacked",
        title="Customer Segments",
        theme=theme_obj,
    )
    area_chart.render(slide3, left=5.3, top=1.8, width=4.2, height=4.5)

    # ==========================================================================
    # SLIDE 4: PIE & DOUGHNUT CHARTS
    # ==========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide4)

    title_shape = slide4.shapes.title
    title_shape.text = "Pie & Doughnut Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")

    # Pie chart with market share (simplified)
    pie_chart = PieChart(
        categories=["SaaS", "Enterprise", "Mobile", "Other"],
        values=[40, 28, 20, 12],
        explode_slice=0,
        title="Revenue by Product",
        theme=theme_obj,
    )
    pie_chart.render(slide4, left=0.5, top=1.8, width=4.2, height=4.5)

    # Doughnut chart (simplified)
    doughnut = DoughnutChart(
        categories=["Americas", "Europe", "Asia", "Other"],
        values=[42, 28, 24, 6],
        hole_size=0.55,
        title="Global Distribution",
        theme=theme_obj,
    )
    doughnut.render(slide4, left=5.3, top=1.8, width=4.2, height=4.5)

    # ==========================================================================
    # SLIDE 5: SCATTER & BUBBLE CHARTS
    # ==========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide5)

    title_shape = slide5.shapes.title
    title_shape.text = "Scatter & Bubble Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")

    # Scatter plot (simplified)
    scatter_chart = ScatterChart(
        series_data=[
            {
                "name": "Products",
                "x_values": [20, 30, 40, 50, 60, 70, 80],
                "y_values": [25, 35, 42, 50, 58, 68, 78],
            }
        ],
        marker_size=10,
        title="Performance Correlation",
        theme=theme_obj,
    )
    scatter_chart.render(slide5, left=0.5, top=1.8, width=4.2, height=4.5)

    # Bubble chart (simplified)
    bubble_chart = BubbleChart(
        series_data=[
            {
                "name": "Portfolio",
                "points": [[30, 40, 20], [50, 60, 30], [40, 50, 25], [60, 70, 35], [70, 80, 40]],
            }
        ],
        size_scale=1.8,
        transparency=25,
        title="Market Positioning",
        theme=theme_obj,
    )
    bubble_chart.render(slide5, left=5.3, top=1.8, width=4.2, height=4.5)

    # ==========================================================================
    # SLIDE 6: RADAR & GAUGE CHARTS
    # ==========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide6)

    title_shape = slide6.shapes.title
    title_shape.text = "Radar & Gauge Charts"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")

    # Radar chart (simplified)
    radar_chart = RadarChart(
        categories=["Speed", "Quality", "Cost", "Support"],
        series={"Product A": [8.5, 9.2, 7.5, 8.8], "Product B": [7.2, 8.8, 8.8, 7.5]},
        variant="filled",
        max_value=10,
        title="Product Comparison",
        theme=theme_obj,
    )
    radar_chart.render(slide6, left=0.5, top=1.8, width=4.2, height=4.5)

    # Two larger gauge charts side by side
    gauge1 = GaugeChart(value=87, min_value=0, max_value=100, title="Satisfaction", theme=theme_obj)
    gauge1.render(slide6, left=5.3, top=1.8, width=2.0, height=2.0)

    gauge2 = GaugeChart(value=94, min_value=0, max_value=100, title="Uptime", theme=theme_obj)
    gauge2.render(slide6, left=7.5, top=1.8, width=2.0, height=2.0)

    # ==========================================================================
    # SLIDE 7: SAMPLE CHARTS
    # ==========================================================================
    slide7 = prs.slides.add_slide(prs.slide_layouts[5])
    theme_obj.apply_to_slide(slide7)

    title_shape = slide7.shapes.title
    title_shape.text = f"{theme_obj.name} Theme Colors"
    title_shape.text_frame.paragraphs[0].font.color.rgb = theme_obj.get_color("foreground.DEFAULT")

    # Larger sample charts to showcase theme colors
    sample_column = ColumnChart(
        categories=["A", "B", "C", "D"],
        series={"Series 1": [10, 15, 12, 18], "Series 2": [8, 12, 10, 14]},
        title="Multi-Series Column",
        theme=theme_obj,
    )
    sample_column.render(slide7, left=0.5, top=1.8, width=4.2, height=4.5)

    sample_pie = PieChart(
        categories=["X", "Y", "Z", "W"],
        values=[40, 25, 20, 15],
        title="Color Palette",
        theme=theme_obj,
    )
    sample_pie.render(slide7, left=5.3, top=1.8, width=4.2, height=4.5)

    return prs


def create_all_theme_galleries():
    """Create chart galleries for all themes."""

    print("\n🎨 Creating Individual Theme Chart Galleries")
    print("=" * 70)

    theme_manager = ThemeManager()

    # Define theme groups
    theme_groups = {
        "Dark Themes": ["dark", "dark-blue", "dark-violet", "dark-green", "dark-purple"],
        "Light Themes": ["light", "corporate", "light-warm"],
        "Special Themes": ["cyberpunk", "sunset", "ocean", "minimal"],
    }

    created_files = []

    for group_name, theme_names in theme_groups.items():
        print(f"\n📁 {group_name}")
        print("-" * 30)

        for theme_name in theme_names:
            theme_obj = theme_manager.get_theme(theme_name)
            if theme_obj:
                print(f"  Creating {theme_name} gallery...")

                try:
                    prs = create_chart_gallery_for_theme(theme_name, theme_obj)

                    # Save presentation
                    output_dir = os.path.join(
                        os.path.dirname(__file__), "..", "outputs", "theme_galleries"
                    )
                    os.makedirs(output_dir, exist_ok=True)

                    filename = f"chart_gallery_{theme_name.replace('-', '_')}.pptx"
                    output_path = os.path.join(output_dir, filename)
                    prs.save(output_path)

                    created_files.append((theme_name, output_path))
                    print(f"    ✅ Created {filename}")

                except Exception as e:
                    print(f"    ❌ Error creating {theme_name}: {e}")
            else:
                print(f"    ⚠️  Theme '{theme_name}' not found")

    # Summary
    print(f"\n🎉 Created {len(created_files)} theme galleries!")
    print("\n📊 Files created:")
    for theme_name, path in created_files:
        print(f"  • {theme_name}: {os.path.basename(path)}")

    print("\n📁 All files saved to: outputs/theme_galleries/")
    print("\n💡 Each gallery shows the same charts with different themes,")
    print("    demonstrating the power of the component-based design system!")

    return created_files


if __name__ == "__main__":
    create_all_theme_galleries()
