#!/usr/bin/env python3
"""
Tokens Showcase - Comprehensive demonstration of design tokens system.
Shows colors, typography, spacing, borders, shadows, and semantic color usage.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.card import Card, MetricCard
from chuk_mcp_pptx.components.core.badge import Badge
from chuk_mcp_pptx.components.core.button import Button
from chuk_mcp_pptx.components.core import Container, Grid, Stack
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_color_tokens_slide(prs, theme):
    """Showcase color token system."""
    print("  • Creating Color Tokens showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Color Tokens"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Container for organized layout
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    # Semantic colors in badges
    semantic_colors = [
        ("Primary", "default", 0),
        ("Secondary", "secondary", 3),
        ("Success", "success", 6),
        ("Warning", "warning", 9),
    ]

    for label, variant, col_start in semantic_colors:
        pos = grid.get_cell(col_span=3, col_start=col_start, row_start=0)
        Badge(text=label, variant=variant, theme=theme.__dict__).render(
            slide, left=pos["left"] + 0.3, top=pos["top"] + 0.1
        )

    # Additional semantic colors
    more_colors = [
        ("Destructive", "destructive", 0),
        ("Outline", "outline", 3),
        ("Muted", "secondary", 6),
    ]

    for label, variant, col_start in more_colors:
        pos = grid.get_cell(col_span=3, col_start=col_start, row_start=1)
        Badge(text=label, variant=variant, theme=theme.__dict__).render(
            slide, left=pos["left"] + 0.3, top=pos["top"] + 0.1
        )

    # Color cards showing token usage
    color_cards = [
        ("Background", "Card backgrounds and surfaces", 0),
        ("Foreground", "Text and icon colors", 4),
        ("Border", "Dividers and outlines", 8),
    ]

    for title, desc, col_start in color_cards:
        pos = grid.get_cell(col_span=4, col_start=col_start, row_start=2)
        card = Card(variant="outlined", theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description(desc))
        card.render(slide, **pos)


def create_typography_tokens_slide(prs, theme):
    """Showcase typography token system."""
    print("  • Creating Typography Tokens showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Typography Tokens"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Typography scale demonstration
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)

    # Create cards with different text sizes
    stack = Stack(direction="vertical", gap="sm", align="start")

    type_scale = [
        ("Heading", "Large titles and headers"),
        ("Subheading", "Section titles and emphasis"),
        ("Body", "Main content and descriptions"),
        ("Caption", "Small text and metadata"),
    ]

    cards = []
    for title, desc in type_scale:
        card = Card(variant="default", theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description(desc))
        cards.append(card)

    stack.render_children(
        slide,
        cards,
        left=bounds["left"],
        top=bounds["top"],
        item_width=bounds["width"],
        item_height=1.1,  # Give each card more height to prevent cutoff
    )


def create_spacing_tokens_slide(prs, theme):
    """Showcase spacing token system."""
    print("  • Creating Spacing Tokens showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Spacing Tokens"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Spacing scale with visual examples
    spacing_scale = [
        ("xs", 'Extra Small (0.1")', 2.2),
        ("sm", 'Small (0.2")', 3.0),
        ("md", 'Medium (0.3")', 3.8),
        ("lg", 'Large (0.5")', 4.6),
        ("xl", 'Extra Large (0.8")', 5.4),
    ]

    for size, label, top in spacing_scale:
        # Label
        Badge(text=label, variant="outline", theme=theme.__dict__).render(
            slide, left=0.5, top=top - 0.05
        )

        # Visual representation
        stack = Stack(direction="horizontal", gap=size, align="start")
        positions = stack.distribute(
            num_items=4, item_width=0.3, item_height=0.3, left=4.0, top=top
        )

        for pos in positions:
            Badge(text="•", variant="default", theme=theme.__dict__).render(
                slide, left=pos["left"], top=pos["top"]
            )


def create_component_composition_slide(prs, theme):
    """Showcase how tokens work together in components."""
    print("  • Creating Component Composition showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Tokens in Action"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Dashboard showing all tokens working together
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    # Metrics using tokens
    metrics = [
        ("Users", "12.5K", "+8%", "up", 0),
        ("Revenue", "$245K", "+15%", "up", 4),
        ("Growth", "23%", "+3%", "up", 8),
    ]

    for label, value, change, trend, col_start in metrics:
        pos = grid.get_cell(col_span=4, col_start=col_start, row_start=0)
        MetricCard(
            label=label, value=value, change=change, trend=trend, theme=theme.__dict__
        ).render(slide, **pos)

    # Card showing composition
    main_pos = grid.get_cell(col_span=8, col_start=0, row_start=1)
    main_card = Card(variant="elevated", theme=theme.__dict__)
    main_card.add_child(Card.Title("Design Tokens"))
    main_card.add_child(
        Card.Description("Colors, spacing, typography, and borders working together")
    )
    main_card.render(slide, **main_pos)

    # Sidebar with actions
    sidebar_pos = grid.get_cell(col_span=4, col_start=8, row_start=1)

    buttons = [
        Button("Primary", variant="default", size="sm", theme=theme.__dict__),
        Button("Secondary", variant="secondary", size="sm", theme=theme.__dict__),
        Button("Outline", variant="outline", size="sm", theme=theme.__dict__),
    ]

    stack = Stack(direction="vertical", gap="sm", align="start")
    stack.render_children(
        slide,
        buttons,
        left=sidebar_pos["left"] + 0.1,
        top=sidebar_pos["top"] + 0.1,
        item_width=sidebar_pos["width"] - 0.2,
        item_height=0.4,
    )


def create_variant_system_slide(prs, theme):
    """Showcase the variant system powered by tokens."""
    print("  • Creating Variant System showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Variant System"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Card variants showing token flexibility
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    card_variants = [
        ("default", "Default", "Standard styling", 0),
        ("outlined", "Outlined", "Border emphasis", 4),
        ("elevated", "Elevated", "Shadow depth", 8),
    ]

    for variant, title, desc, col_start in card_variants:
        pos = grid.get_cell(col_span=4, col_start=col_start, row_start=0)
        card = Card(variant=variant, theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description(desc))
        card.render(slide, **pos)

    # Button variants
    button_positions = [
        ("Default", "default", 0),
        ("Secondary", "secondary", 2.5),
        ("Outline", "outline", 5.0),
        ("Ghost", "ghost", 7.5),
    ]

    for text, variant, left in button_positions:
        btn = Button(text=text, variant=variant, size="md", theme=theme.__dict__)
        btn.render(slide, left=bounds["left"] + left, top=bounds["top"] + 3.0, width=2.0)


def create_semantic_colors_slide(prs, theme):
    """Showcase semantic color usage."""
    print("  • Creating Semantic Colors showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Semantic Colors"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Semantic color examples with meaning - use grid for better fit
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)

    grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)

    semantic_examples = [
        ("Success", "Completed tasks and positive actions", "success", 0, 0),
        ("Warning", "Caution and important notices", "warning", 6, 0),
        ("Destructive", "Errors and critical actions", "destructive", 0, 1),
        ("Info", "Informational messages", "default", 6, 1),
    ]

    # Render cards and badges in 2x2 grid
    for title, desc, variant, col_start, row_start in semantic_examples:
        pos = grid.get_cell(col_span=6, col_start=col_start, row_start=row_start)

        card = Card(variant="outlined", theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description(desc))
        card.render(slide, **pos)

        # Render badge next to card title
        Badge(text=variant.upper(), variant=variant, theme=theme.__dict__).render(
            slide, left=pos["left"] + pos["width"] - 1.5, top=pos["top"] + 0.15
        )


def main():
    """Generate tokens showcase presentation."""
    print("\n🎨 Creating Tokens Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    create_color_tokens_slide(prs, theme)
    create_typography_tokens_slide(prs, theme)
    create_spacing_tokens_slide(prs, theme)
    create_component_composition_slide(prs, theme)
    create_variant_system_slide(prs, theme)
    create_semantic_colors_slide(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "tokens_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n🎨 Design Tokens Demonstrated:")
    print("  • Color Tokens (semantic colors, backgrounds, foregrounds)")
    print("  • Typography Tokens (heading, body, caption scales)")
    print("  • Spacing Tokens (xs, sm, md, lg, xl)")
    print("  • Component Composition (tokens working together)")
    print("  • Variant System (flexible styling)")
    print("  • Semantic Colors (meaningful color usage)")
    print("\n💡 Shows how design tokens enable:")
    print("  • Consistent styling across components")
    print("  • Easy theme switching")
    print("  • Composable design patterns")
    print("  • Maintainable design systems")


if __name__ == "__main__":
    main()
