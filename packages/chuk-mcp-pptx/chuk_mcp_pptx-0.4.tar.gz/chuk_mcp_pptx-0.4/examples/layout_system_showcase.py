#!/usr/bin/env python3
"""
Layout System Showcase - Demonstrates the grid/layout system.
Clean, organized examples with proper spacing and alignment.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core import Container, Grid, Stack, Divider
from chuk_mcp_pptx.components.core.card import Card, MetricCard
from chuk_mcp_pptx.components.core.badge import Badge
from chuk_mcp_pptx.components.core.button import Button
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_grid_demo(prs, theme):
    """Demonstrate 12-column grid system with clean layout."""
    print("  • Creating Grid System demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "12-Column Grid System"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    grid = Grid(columns=12, gap="md")

    # Row 1: Full width example
    pos = grid.get_span(col_span=12, col_start=0, left=0.5, top=1.8, width=9.0, height=0.6)
    Badge(text="12 Columns - Full Width", variant="outline", theme=theme.__dict__).render(
        slide, left=pos["left"] + 3.5, top=pos["top"] + 0.1
    )

    # Row 2: Two equal columns (6 + 6)
    pos1 = grid.get_span(col_span=6, col_start=0, left=0.5, top=2.6, width=9.0, height=0.6)
    Badge(text="6 Columns", variant="default", theme=theme.__dict__).render(
        slide, left=pos1["left"] + 1.5, top=pos1["top"] + 0.1
    )

    pos2 = grid.get_span(col_span=6, col_start=6, left=0.5, top=2.6, width=9.0, height=0.6)
    Badge(text="6 Columns", variant="default", theme=theme.__dict__).render(
        slide, left=pos2["left"] + 1.5, top=pos2["top"] + 0.1
    )

    # Row 3: Three equal columns (4 + 4 + 4)
    for i in range(3):
        pos = grid.get_span(col_span=4, col_start=i * 4, left=0.5, top=3.4, width=9.0, height=0.6)
        Badge(text="4 Cols", variant="secondary", theme=theme.__dict__).render(
            slide, left=pos["left"] + 0.8, top=pos["top"] + 0.1
        )

    # Row 4: Asymmetric layout (8 + 4) - Main + Sidebar pattern
    main = grid.get_span(col_span=8, col_start=0, left=0.5, top=4.2, width=9.0, height=1.8)
    card_main = Card(variant="elevated", theme=theme.__dict__)
    card_main.add_child(Card.Title("Main Content (8 cols)"))
    card_main.add_child(Card.Description("Primary content area"))
    card_main.render(slide, **main)

    sidebar = grid.get_span(col_span=4, col_start=8, left=0.5, top=4.2, width=9.0, height=1.8)
    card_sidebar = Card(variant="outlined", theme=theme.__dict__)
    card_sidebar.add_child(Card.Title("Sidebar (4)"))
    card_sidebar.add_child(Card.Description("Secondary info"))
    card_sidebar.render(slide, **sidebar)

    # Row 5: Four equal columns
    for i in range(4):
        pos = grid.get_span(col_span=3, col_start=i * 3, left=0.5, top=6.2, width=9.0, height=0.6)
        Badge(text="3", variant="success", theme=theme.__dict__).render(
            slide, left=pos["left"] + 0.8, top=pos["top"] + 0.1
        )


def create_container_demo(prs, theme):
    """Demonstrate Container system with clean examples."""
    print("  • Creating Container demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Container System"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Small container
    container_sm = Container(size="sm", padding="md", center=True)
    bounds_sm = container_sm.render(slide, top=2.0)
    card_sm = Card(variant="elevated", theme=theme.__dict__)
    card_sm.add_child(Card.Title('Small (8")'))
    card_sm.add_child(Card.Description("Focused content"))
    card_sm.render(
        slide, left=bounds_sm["left"], top=bounds_sm["top"], width=bounds_sm["width"]
    )  # Auto-height

    # Medium container
    container_md = Container(size="md", padding="md", center=True)
    bounds_md = container_md.render(slide, top=3.5)
    card_md = Card(variant="elevated", theme=theme.__dict__)
    card_md.add_child(Card.Title('Medium (9")'))
    card_md.add_child(Card.Description("Balanced width"))
    card_md.render(
        slide, left=bounds_md["left"], top=bounds_md["top"], width=bounds_md["width"]
    )  # Auto-height

    # Large container
    container_lg = Container(size="lg", padding="md", center=True)
    bounds_lg = container_lg.render(slide, top=5.0)
    card_lg = Card(variant="elevated", theme=theme.__dict__)
    card_lg.add_child(Card.Title('Large (10")'))
    card_lg.add_child(Card.Description("Standard slide width"))
    card_lg.render(
        slide, left=bounds_lg["left"], top=bounds_lg["top"], width=bounds_lg["width"]
    )  # Auto-height

    # Visual indicators of centering
    Divider(
        orientation="vertical", thickness=1, color="border.DEFAULT", theme=theme.__dict__
    ).render(slide, left=5.0, top=1.8, height=4.6)


def create_stack_demo(prs, theme):
    """Demonstrate Stack layouts."""
    print("  • Creating Stack demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Stack Layouts"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Vertical stack on left
    Badge(text="Vertical Stack", variant="outline", theme=theme.__dict__).render(
        slide, left=0.5, top=1.8
    )

    # Create cards
    v_cards = []
    for i in range(4):
        card = Card(variant="default", theme=theme.__dict__)
        card.add_child(Card.Title(f"Item {i + 1}"))
        v_cards.append(card)

    # Stack handles positioning automatically
    v_stack = Stack(direction="vertical", gap="md", align="start")
    v_stack.render_children(slide, v_cards, left=0.5, top=2.2, item_width=4.0)

    # Vertical divider
    Divider(orientation="vertical", thickness=1, theme=theme.__dict__).render(
        slide, left=4.8, top=1.8, height=4.5
    )

    # Horizontal stack on right
    Badge(text="Horizontal Stack", variant="outline", theme=theme.__dict__).render(
        slide, left=5.2, top=1.8
    )

    # Create cards
    h_cards = []
    for i in range(3):
        card = Card(variant="outlined", theme=theme.__dict__)
        card.add_child(Card.Title(f"{i + 1}"))
        h_cards.append(card)

    # Stack handles positioning automatically
    h_stack = Stack(direction="horizontal", gap="md", align="start")
    h_stack.render_children(slide, h_cards, left=5.2, top=2.2, item_width=1.3, item_height=2.0)


def create_spacing_demo(prs, theme):
    """Demonstrate spacing scale."""
    print("  • Creating Spacing demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Spacing Scale"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    gaps = [
        ("xs", "Extra Small"),
        ("sm", "Small"),
        ("md", "Medium"),
        ("lg", "Large"),
        ("xl", "Extra Large"),
    ]
    top = 2.2

    for gap_size, label in gaps:
        # Label
        Badge(text=label, variant="outline", theme=theme.__dict__).render(
            slide, left=0.5, top=top - 0.05
        )

        # Visual representation with dots
        stack = Stack(direction="horizontal", gap=gap_size, align="start")
        positions = stack.distribute(
            num_items=5, item_width=0.2, item_height=0.2, left=2.5, top=top
        )

        for pos in positions:
            Badge(text="•", variant="default", theme=theme.__dict__).render(
                slide, left=pos["left"], top=pos["top"]
            )

        # Gap size indicator
        Badge(text=f"gap: {gap_size}", variant="secondary", theme=theme.__dict__).render(
            slide, left=7.5, top=top - 0.05
        )

        top += 1.0


def create_responsive_dashboard(prs, theme):
    """Create a clean responsive dashboard demonstrating 12-column grid system."""
    print("  • Creating Responsive Dashboard...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Responsive Dashboard - 12 Column Grid"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Use Container to establish content boundaries
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)

    # Grid uses container bounds - no need to pass them to each get_cell call
    grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)

    # Row 0 - Metrics (3 cells of 4 columns each)
    metrics = [
        ("Revenue", "$1.2M", "+15%", "up", 0),
        ("Users", "45K", "+8%", "up", 4),
        ("Growth", "23%", "+3%", "up", 8),
    ]

    for label, value, change, trend, col_start in metrics:
        # Grid knows its bounds - just specify cell position
        pos = grid.get_cell(col_span=4, col_start=col_start, row_start=0)
        MetricCard(
            label=label, value=value, change=change, trend=trend, theme=theme.__dict__
        ).render(slide, **pos)

    # Row 1 - Main content (8 cols)
    main_pos = grid.get_cell(col_span=8, col_start=0, row_start=1)
    main_card = Card(variant="elevated", theme=theme.__dict__)
    main_card.add_child(Card.Title("Main Content (8/12 cols)"))
    main_card.add_child(Card.Description("Primary content area scales with grid system"))
    main_card.render(slide, **main_pos)

    # Row 1 - Sidebar (4 cols) - simplified, no manual calculations
    sidebar_pos = grid.get_cell(col_span=4, col_start=8, row_start=1)

    # Just a label badge for the sidebar section
    Badge(text="Actions (4/12)", variant="outline", theme=theme.__dict__).render(
        slide, left=sidebar_pos["left"] + 0.1, top=sidebar_pos["top"] + 0.1
    )

    # Buttons stacked in sidebar - grid provides boundaries
    buttons = [
        Button("Export", variant="outline", size="sm", theme=theme.__dict__),
        Button("Refresh", variant="secondary", size="sm", theme=theme.__dict__),
        Button("Settings", variant="ghost", size="sm", theme=theme.__dict__),
    ]

    stack = Stack(direction="vertical", gap="sm", align="start")
    stack.render_children(
        slide,
        buttons,
        left=sidebar_pos["left"] + 0.1,
        top=sidebar_pos["top"] + 0.5,
        item_width=sidebar_pos["width"] - 0.2,
        item_height=0.4,
    )


def create_divider_demo(prs, theme):
    """Demonstrate dividers with clean, visually appealing layout."""
    print("  • Creating Divider demo...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Dividers & Separators"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Left side: Badges with dividers demonstrating separator usage
    sections = [
        ("Introduction", "default"),
        ("Analysis", "success"),
        ("Recommendation", "warning"),
        ("Conclusion", "outline"),
    ]

    # Create badges
    section_badges = [
        Badge(text=label, variant=variant, theme=theme.__dict__) for label, variant in sections
    ]

    # Stack badges with proper gap
    left_stack = Stack(direction="vertical", gap="lg")

    # Get positions for both badges and dividers
    positions = left_stack.distribute(
        num_items=len(sections), item_width=4.0, item_height=0.3, left=0.5, top=2.0
    )

    # Render badges at calculated positions
    for badge, pos in zip(section_badges, positions):
        badge.render(slide, left=pos["left"], top=pos["top"])

        # Add divider below each badge (except last)
        if badge != section_badges[-1]:
            Divider(orientation="horizontal", thickness=1, theme=theme.__dict__).render(
                slide, left=pos["left"], top=pos["top"] + 0.4, width=4.0
            )

    # Vertical divider in center - taller and more prominent
    Divider(orientation="vertical", thickness=2, theme=theme.__dict__).render(
        slide, left=5.0, top=1.8, height=4.6
    )

    # Right side: Content cards with proper spacing
    right_stack = Stack(direction="vertical", gap="md", align="start")

    # Get positions for cards
    card_positions = right_stack.distribute(
        num_items=2,
        item_width=4.0,
        item_height=1.5,  # Fixed height for each card
        left=5.5,
        top=2.0,
    )

    # Render cards at calculated positions
    right_cards = [
        ("Content Section A", "First content area with consistent styling"),
        ("Content Section B", "Second content area with matching layout"),
    ]

    for (title, description), pos in zip(right_cards, card_positions):
        card = Card(variant="elevated", theme=theme.__dict__)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description(description))
        card.render(slide, left=pos["left"], top=pos["top"], width=pos["width"])


def main():
    """Generate layout system showcase."""
    print("\n🎨 Creating Layout System Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    create_grid_demo(prs, theme)
    create_container_demo(prs, theme)
    create_stack_demo(prs, theme)
    create_spacing_demo(prs, theme)
    create_responsive_dashboard(prs, theme)
    create_divider_demo(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "layout_system_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n🎨 Layout Features Demonstrated:")
    print("  • 12-Column Grid System")
    print("  • Container Sizes (sm, md, lg)")
    print("  • Stack Layouts (vertical & horizontal)")
    print("  • Spacing Scale (xs → xl)")
    print("  • Responsive Dashboard")
    print("  • Dividers & Separators")
    print("\n💡 Clean, organized layouts with:")
    print("  • Consistent spacing")
    print("  • Clear visual hierarchy")
    print("  • Proper alignment")
    print("  • Readable labels")


if __name__ == "__main__":
    main()
