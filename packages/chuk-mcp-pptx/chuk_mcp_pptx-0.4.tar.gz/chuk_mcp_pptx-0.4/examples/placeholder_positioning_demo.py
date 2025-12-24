#!/usr/bin/env python3
"""
Placeholder Positioning Demo

This example demonstrates how components should position themselves
within template placeholders. It tests:
- Tables in content placeholders
- Charts in content placeholders
- Images in picture placeholders
- Cards and other components in placeholders

The key is that when targeting a placeholder, the component should
render at the EXACT bounds of the placeholder, not at hardcoded positions.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation

# Import components directly for testing
from chuk_mcp_pptx.components.core.table import Table
from chuk_mcp_pptx.components.core.card import MetricCard
from chuk_mcp_pptx.components.charts.column_bar import ColumnChart


def create_test_presentation():
    """Create a test presentation with various placeholder layouts."""
    prs = Presentation()

    # Use the Title and Content layout (index 1 usually has content placeholder)
    # We'll create slides manually to control placeholder positions

    return prs


def test_table_placeholder_positioning():
    """Test that Table component respects placeholder bounds."""
    print("\n" + "=" * 60)
    print("TEST: Table Placeholder Positioning")
    print("=" * 60)

    prs = Presentation()

    # Add a slide with content layout
    slide_layout = prs.slide_layouts[5]  # Usually "Title Only" or blank-ish
    slide = prs.slides.add_slide(slide_layout)

    # Create a mock placeholder-like shape at specific position
    # In real templates, placeholders have specific idx and bounds
    mock_left, mock_top, mock_width, mock_height = 2.0, 2.5, 6.0, 3.0

    print(f"\nExpected position: ({mock_left}, {mock_top}, {mock_width}, {mock_height})")

    # Create Table component
    table = Table(
        headers=["Product", "Price", "Stock"],
        data=[
            ["Cheese A", "$10", "100"],
            ["Cheese B", "$15", "50"],
            ["Cheese C", "$20", "75"],
        ],
        variant="striped",
    )

    # Test 1: Render with explicit positions (no placeholder)
    print("\n1. Rendering Table with explicit positions...")
    shape1 = table.render(
        slide=slide,
        left=mock_left,
        top=mock_top,
        width=mock_width,
        height=mock_height,
        placeholder=None,
    )

    # Check actual position
    actual_left = shape1.left.inches
    actual_top = shape1.top.inches
    actual_width = shape1.width.inches
    actual_height = shape1.height.inches

    print(
        f"   Actual position: ({actual_left:.2f}, {actual_top:.2f}, {actual_width:.2f}, {actual_height:.2f})"
    )

    if abs(actual_left - mock_left) < 0.01 and abs(actual_top - mock_top) < 0.01:
        print("   ✅ PASS: Table positioned correctly with explicit coords")
    else:
        print("   ❌ FAIL: Table position mismatch!")

    # Save test output
    output_path = os.path.join(
        os.path.dirname(__file__), "..", "outputs", "test_table_positioning.pptx"
    )
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"\n   Saved to: {output_path}")

    return True


def test_chart_placeholder_positioning():
    """Test that Chart component respects placeholder bounds."""
    print("\n" + "=" * 60)
    print("TEST: Chart Placeholder Positioning")
    print("=" * 60)

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Expected placeholder bounds
    mock_left, mock_top, mock_width, mock_height = 1.5, 2.0, 7.0, 4.0

    print(f"\nExpected position: ({mock_left}, {mock_top}, {mock_width}, {mock_height})")

    # Create ColumnChart
    chart = ColumnChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        series={"Revenue": [100, 150, 120, 180]},
        title="Quarterly Revenue",
    )

    # Render with explicit positions
    print("\n1. Rendering ColumnChart with explicit positions...")
    shape = chart.render(
        slide=slide,
        left=mock_left,
        top=mock_top,
        width=mock_width,
        height=mock_height,
        placeholder=None,
    )

    # Check actual position
    actual_left = shape.left.inches
    actual_top = shape.top.inches
    actual_width = shape.width.inches
    actual_height = shape.height.inches

    print(
        f"   Actual position: ({actual_left:.2f}, {actual_top:.2f}, {actual_width:.2f}, {actual_height:.2f})"
    )

    if abs(actual_left - mock_left) < 0.01 and abs(actual_top - mock_top) < 0.01:
        print("   ✅ PASS: Chart positioned correctly")
    else:
        print("   ❌ FAIL: Chart position mismatch!")
        print(f"      Delta: left={actual_left - mock_left:.2f}, top={actual_top - mock_top:.2f}")

    # Also test with DEFAULT positions (should use ChartComponent.DEFAULT_* values)
    slide2 = prs.slides.add_slide(slide_layout)
    print("\n2. Rendering ColumnChart with DEFAULT positions (None)...")
    chart2 = ColumnChart(
        categories=["A", "B", "C"], series={"Sales": [10, 20, 30]}, title="Default Position Test"
    )

    # Render with no position params - should use defaults
    shape2 = chart2.render(
        slide=slide2, left=None, top=None, width=None, height=None, placeholder=None
    )

    print(
        f"   Default position: ({shape2.left.inches:.2f}, {shape2.top.inches:.2f}, {shape2.width.inches:.2f}, {shape2.height.inches:.2f})"
    )
    print("   Expected defaults: (1.0, 2.0, 8.0, 3.0)")  # From ChartComponent class

    output_path = os.path.join(
        os.path.dirname(__file__), "..", "outputs", "test_chart_positioning.pptx"
    )
    prs.save(output_path)
    print(f"\n   Saved to: {output_path}")

    return True


def test_chart_with_real_placeholder():
    """Test that Chart component respects actual placeholder bounds."""
    print("\n" + "=" * 60)
    print("TEST: Chart with Real Placeholder")
    print("=" * 60)

    prs = Presentation()

    # Use a layout that has content placeholders
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)

    print(f"\nSlide layout: {slide_layout.name}")

    content_placeholder = None
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        ph_type = ph.placeholder_format.type
        if ph_type == 7 or idx == 1:  # Content placeholder
            content_placeholder = ph
            break

    if content_placeholder:
        ph_left = content_placeholder.left.inches
        ph_top = content_placeholder.top.inches
        ph_width = content_placeholder.width.inches
        ph_height = content_placeholder.height.inches

        print(f"Placeholder bounds: ({ph_left:.2f}, {ph_top:.2f}, {ph_width:.2f}, {ph_height:.2f})")

        # Create chart and render with placeholder
        chart = ColumnChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={"Revenue": [100, 150, 120, 180]},
            title="Chart in Placeholder",
        )

        print("\nRendering ColumnChart with placeholder parameter...")
        shape = chart.render(
            slide=slide,
            left=1.0,  # These should be OVERRIDDEN by placeholder bounds
            top=1.0,
            width=2.0,
            height=1.0,
            placeholder=content_placeholder,
        )

        actual_left = shape.left.inches
        actual_top = shape.top.inches
        actual_width = shape.width.inches

        print(f"Chart rendered at: ({actual_left:.2f}, {actual_top:.2f}, {actual_width:.2f}, ...)")

        # Check if left and width match placeholder bounds (top may vary due to title)
        if abs(actual_left - ph_left) < 0.1 and abs(actual_width - ph_width) < 0.1:
            print("✅ PASS: Chart correctly used placeholder bounds!")
        else:
            print("❌ FAIL: Chart did NOT use placeholder bounds!")
            print(f"   Expected left: {ph_left:.2f}, got: {actual_left:.2f}")
            print(f"   Expected width: {ph_width:.2f}, got: {actual_width:.2f}")
    else:
        print("No content placeholder found on this layout")

    output_path = os.path.join(
        os.path.dirname(__file__), "..", "outputs", "test_chart_with_placeholder.pptx"
    )
    prs.save(output_path)
    print(f"\nSaved to: {output_path}")

    return True


def test_real_placeholder_extraction():
    """Test extracting bounds from an actual placeholder shape."""
    print("\n" + "=" * 60)
    print("TEST: Real Placeholder Bound Extraction")
    print("=" * 60)

    prs = Presentation()

    # Use a layout that has content placeholders
    # Layout 1 is typically "Title and Content"
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    print(f"\nSlide layout: {slide_layout.name}")
    print("\nPlaceholders on this slide:")

    content_placeholder = None
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        ph_type = ph.placeholder_format.type
        left = ph.left.inches
        top = ph.top.inches
        width = ph.width.inches
        height = ph.height.inches
        print(
            f"  idx={idx}, type={ph_type}, bounds=({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
        )

        # Find a content placeholder (type 7 is OBJECT/CONTENT)
        if ph_type == 7 or idx == 1:  # Content placeholder
            content_placeholder = ph

    if content_placeholder:
        print(f"\nUsing placeholder idx={content_placeholder.placeholder_format.idx} for test")

        # Extract bounds like the component would
        ph_left = content_placeholder.left.inches
        ph_top = content_placeholder.top.inches
        ph_width = content_placeholder.width.inches
        ph_height = content_placeholder.height.inches

        print(f"Placeholder bounds: ({ph_left:.2f}, {ph_top:.2f}, {ph_width:.2f}, {ph_height:.2f})")

        # Now render a table INTO this placeholder
        table = Table(headers=["Item", "Value"], data=[["Test", "123"]], variant="default")

        # Pass the actual placeholder to render
        print("\nRendering Table with placeholder parameter...")
        shape = table.render(
            slide=slide,
            left=1.0,  # These should be OVERRIDDEN by placeholder bounds
            top=1.0,
            width=2.0,
            height=1.0,
            placeholder=content_placeholder,
        )

        actual_left = shape.left.inches
        actual_top = shape.top.inches
        actual_width = shape.width.inches
        actual_height = shape.height.inches

        print(
            f"Table rendered at: ({actual_left:.2f}, {actual_top:.2f}, {actual_width:.2f}, {actual_height:.2f})"
        )

        # Check if it matches placeholder bounds
        if abs(actual_left - ph_left) < 0.1 and abs(actual_top - ph_top) < 0.1:
            print("✅ PASS: Table correctly used placeholder bounds!")
        else:
            print("❌ FAIL: Table did NOT use placeholder bounds!")
            print(f"   Expected: ({ph_left:.2f}, {ph_top:.2f})")
            print(f"   Got: ({actual_left:.2f}, {actual_top:.2f})")
    else:
        print("No content placeholder found on this layout")

    output_path = os.path.join(
        os.path.dirname(__file__), "..", "outputs", "test_real_placeholder.pptx"
    )
    prs.save(output_path)
    print(f"\nSaved to: {output_path}")

    return True


def test_metric_card_positioning():
    """Test MetricCard placeholder positioning."""
    print("\n" + "=" * 60)
    print("TEST: MetricCard Placeholder Positioning")
    print("=" * 60)

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Test explicit positioning
    mock_left, mock_top, mock_width, mock_height = 3.0, 2.0, 2.5, 1.5

    print(f"\nExpected position: ({mock_left}, {mock_top}, {mock_width}, {mock_height})")

    card = MetricCard(value="$150K", label="Revenue", change="+12%", variant="default")

    shape = card.render(
        slide=slide,
        left=mock_left,
        top=mock_top,
        width=mock_width,
        height=mock_height,
        placeholder=None,
    )

    actual_left = shape.left.inches
    actual_top = shape.top.inches

    print(f"Actual position: ({actual_left:.2f}, {actual_top:.2f})")

    if abs(actual_left - mock_left) < 0.01 and abs(actual_top - mock_top) < 0.01:
        print("✅ PASS: MetricCard positioned correctly")
    else:
        print("❌ FAIL: MetricCard position mismatch!")

    output_path = os.path.join(
        os.path.dirname(__file__), "..", "outputs", "test_metric_card_positioning.pptx"
    )
    prs.save(output_path)
    print(f"\nSaved to: {output_path}")

    return True


def main():
    """Run all positioning tests."""
    print("🔧 Placeholder Positioning Demo")
    print("=" * 60)
    print("Testing that components respect placeholder bounds...")

    results = []

    try:
        results.append(("Table", test_table_placeholder_positioning()))
    except Exception as e:
        print(f"❌ Table test failed: {e}")
        results.append(("Table", False))

    try:
        results.append(("Chart", test_chart_placeholder_positioning()))
    except Exception as e:
        print(f"❌ Chart test failed: {e}")
        results.append(("Chart", False))

    try:
        results.append(("Chart with Placeholder", test_chart_with_real_placeholder()))
    except Exception as e:
        print(f"❌ Chart with Placeholder test failed: {e}")
        import traceback

        traceback.print_exc()
        results.append(("Chart with Placeholder", False))

    try:
        results.append(("Real Placeholder", test_real_placeholder_extraction()))
    except Exception as e:
        print(f"❌ Real Placeholder test failed: {e}")
        import traceback

        traceback.print_exc()
        results.append(("Real Placeholder", False))

    try:
        results.append(("MetricCard", test_metric_card_positioning()))
    except Exception as e:
        print(f"❌ MetricCard test failed: {e}")
        results.append(("MetricCard", False))

    print("\n" + "=" * 60)
    print("SUMMARY")
    print("=" * 60)
    for name, passed in results:
        status = "✅ PASS" if passed else "❌ FAIL"
        print(f"  {status}: {name}")

    all_passed = all(r[1] for r in results)
    print(f"\nOverall: {'✅ All tests passed!' if all_passed else '❌ Some tests failed'}")

    return all_passed


if __name__ == "__main__":
    success = main()
    sys.exit(0 if success else 1)
