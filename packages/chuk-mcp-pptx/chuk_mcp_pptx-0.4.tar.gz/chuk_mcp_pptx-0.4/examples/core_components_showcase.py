#!/usr/bin/env python3
"""
Core Components Showcase - Comprehensive demonstration of Button, Badge, Alert, and Card.
Shows all variants, sizes, composition patterns, and real-world usage examples.
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core import (
    Button,
    IconButton,
    ButtonGroup,
    Badge,
    DotBadge,
    CountBadge,
    Alert,
    Card,
    MetricCard,
    ProgressBar,
    Icon,
    IconList,
    Timeline,
    IconTile,
    ValueTile,
    Avatar,
    AvatarWithLabel,
    AvatarGroup,
    Shape,
    Connector,
    ProcessFlow,
    CycleDiagram,
    HierarchyDiagram,
    Image,
    TextBox,
    BulletList,
    Table,
)
from chuk_mcp_pptx.components.core import Container, Grid, Stack
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_button_showcase(prs, theme):
    """Showcase all button variants and types."""
    print("  • Creating Button Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Button Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Button variants - organized with better spacing
    button_variants = [
        ("Default", "default", 0.5),
        ("Secondary", "secondary", 2.3),
        ("Outline", "outline", 4.1),
        ("Ghost", "ghost", 5.7),
        ("Destructive", "destructive", 7.0),
    ]

    for text, variant, left in button_variants:
        btn = Button(text=text, variant=variant, size="md", theme=theme)
        btn.render(slide, left=left, top=2.0, width=1.6)

    # Button sizes
    size_buttons = [
        ("Small", "sm", 0.5, 1.3),
        ("Medium", "md", 2.2, 2.0),
        ("Large", "lg", 4.6, 2.8),
    ]

    for text, size, left, width in size_buttons:
        btn = Button(text=text, variant="default", size=size, theme=theme)
        btn.render(slide, left=left, top=3.2, width=width)

    # Icon buttons - compact row
    icon_buttons = [
        ("play", 0.5),
        ("pause", 1.2),
        ("settings", 1.9),
        ("star", 2.6),
        ("heart", 3.3),
        ("search", 4.0),
    ]

    for icon, left in icon_buttons:
        btn = IconButton(icon=icon, variant="ghost", size="md", theme=theme)
        btn.render(slide, left=left, top=4.5)

    # Button groups
    buttons_config = [
        {"text": "Save", "variant": "default", "size": "md"},
        {"text": "Cancel", "variant": "ghost", "size": "md"},
    ]
    group = ButtonGroup(buttons=buttons_config, orientation="horizontal", spacing=0.2, theme=theme)
    group.render(slide, left=0.5, top=5.8)


def create_badge_showcase(prs, theme):
    """Showcase all badge variants and types."""
    print("  • Creating Badge Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Badge Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Badge variants - organized row
    badge_variants = [
        ("Default", "default", 0.5),
        ("Secondary", "secondary", 1.8),
        ("Success", "success", 3.1),
        ("Warning", "warning", 4.4),
        ("Destructive", "destructive", 5.7),
        ("Outline", "outline", 7.2),
    ]

    for text, variant, left in badge_variants:
        badge = Badge(text=text, variant=variant, theme=theme)
        badge.render(slide, left=left, top=2.0)

    # Dot badges - small indicators
    dot_badges = [
        ("default", 0.5),
        ("success", 1.5),
        ("warning", 2.5),
        ("destructive", 3.5),
    ]

    for variant, left in dot_badges:
        dot = DotBadge(variant=variant, theme=theme)
        dot.render(slide, left=left, top=3.0)

    # Count badges - notification style
    count_badges = [
        (1, 5.0),
        (12, 5.8),
        (99, 6.6),
        (150, 7.4),  # Shows "99+"
    ]

    for count, left in count_badges:
        badge = CountBadge(count=count, variant="destructive", theme=theme)
        badge.render(slide, left=left, top=3.0)

    # Badge use cases - combined with text
    use_cases = [
        ("New", "default", 0.5, 4.2),
        ("Beta", "warning", 2.0, 4.2),
        ("Active", "success", 3.5, 4.2),
        ("Deprecated", "destructive", 5.2, 4.2),
    ]

    for text, variant, left, top in use_cases:
        badge = Badge(text=text, variant=variant, theme=theme)
        badge.render(slide, left=left, top=top)


def create_alert_showcase(prs, theme):
    """Showcase all alert variants."""
    print("  • Creating Alert Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Alert Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Alert variants - stacked vertically for readability
    alerts = [
        ("info", "Information", "This is an informational message.", 2.0),
        ("success", "Success!", "Your changes have been saved successfully.", 3.2),
        ("warning", "Warning", "Please review before proceeding.", 4.4),
        ("error", "Error", "An error occurred while processing.", 5.6),
    ]

    for variant, title, description, top in alerts:
        alert = Alert(variant=variant, title=title, description=description, theme=theme)
        alert.render(slide, left=0.5, top=top, width=9.0, height=0.9)


def create_alert_composition_showcase(prs, theme):
    """Showcase alert composition patterns."""
    print("  • Creating Alert Composition examples...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Alert Composition Patterns"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Composed alerts - using children
    alert1 = Alert(variant="success", theme=theme)
    alert1.add_child(Alert.Title("Deployment Complete"))
    alert1.add_child(Alert.Description("Your application has been deployed to production."))
    alert1.render(slide, left=0.5, top=2.0, width=9.0, height=1.0)

    alert2 = Alert(variant="warning", show_icon=True, theme=theme)
    alert2.add_child(Alert.Title("Quota Warning"))
    alert2.add_child(Alert.Description("You've used 90% of your monthly quota."))
    alert2.render(slide, left=0.5, top=3.5, width=9.0, height=1.0)

    # Alert without icon
    alert3 = Alert(variant="info", show_icon=False, theme=theme)
    alert3.add_child(Alert.Title("System Maintenance"))
    alert3.add_child(Alert.Description("Scheduled maintenance on Sunday 2:00 AM - 4:00 AM."))
    alert3.render(slide, left=0.5, top=5.0, width=9.0, height=1.0)


def create_card_showcase(prs, theme):
    """Showcase all card variants."""
    print("  • Creating Card Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Card Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Use Container → Grid pattern for card variants
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=1.8)

    # 12-column grid for card variants (4 cols each = 3 cards)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    card_variants = [
        ("default", "Default Card", 0),
        ("outlined", "Outlined Card", 4),
        ("elevated", "Elevated Card", 8),
    ]

    for variant, title, col_start in card_variants:
        pos = grid.get_cell(col_span=4, col_start=col_start)
        card = Card(variant=variant, theme=theme)
        card.add_child(Card.Title(title))
        card.add_child(Card.Description("Card with composition pattern"))
        card.render(slide, **pos)

    # Metric cards using 12-column grid (3 cols each = 4 cards)
    metrics = [
        ("Revenue", "$1.2M", "+12%", "up", 0),
        ("Users", "45.2K", "+8%", "up", 3),
        ("Retention", "92%", "-2%", "down", 6),
        ("NPS", "4.8", "0%", "neutral", 9),
    ]

    for label, value, change, trend, col_start in metrics:
        pos = grid.get_cell(col_span=3, col_start=col_start, row_start=1)
        metric = MetricCard(label=label, value=value, change=change, trend=trend, theme=theme)
        metric.render(slide, **pos)


def create_progress_icon_timeline_showcase(prs, theme):
    """Showcase ProgressBar, Icon, and Timeline components."""
    print("  • Creating ProgressBar, Icon & Timeline showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Progress, Icons & Timeline"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Progress bars
    ProgressBar(
        value=75, label="Project Progress", show_percentage=True, variant="success", theme=theme
    ).render(slide, left=0.5, top=1.9, width=4.0)

    ProgressBar(value=60, segments=10, style="segmented", label="Milestones", theme=theme).render(
        slide, left=0.5, top=2.8, width=4.0
    )

    # Icons row
    icons = [("check", "success"), ("star", "warning"), ("rocket", "primary"), ("target", "error")]
    left = 5.0
    for icon, variant in icons:
        Icon(icon, variant=variant, size="lg", theme=theme).render(slide, left=left, top=2.0)
        left += 0.7

    # Icon list
    features = [
        ("check", "Fast & Reliable"),
        ("check", "Easy to Use"),
        ("rocket", "High Performance"),
    ]
    IconList(features, variant="success", icon_size="sm", theme=theme).render(
        slide, left=5.0, top=2.8, width=4.0
    )

    # Timeline
    events = [
        {"date": "Q1", "title": "Plan"},
        {"date": "Q2", "title": "Build"},
        {"date": "Q3", "title": "Launch", "highlight": True},
    ]
    Timeline(events, style="arrow", theme=theme).render(slide, left=0.5, top=4.5, width=8.5)


def create_tile_avatar_showcase(prs, theme):
    """Showcase Tile and Avatar components."""
    print("  • Creating Tile & Avatar showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Tiles & Avatars"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Tiles - different variants in a row
    tiles = [
        (IconTile("rocket", label="Fast", variant="filled", color_variant="primary"), 0.5),
        (ValueTile("42", label="Tasks", variant="outlined"), 2.6),
        (IconTile("check", label="Done", variant="filled", color_variant="success"), 4.7),
        (ValueTile("98%", label="Score", variant="default"), 6.8),
    ]

    for tile, left in tiles:
        tile.theme = theme.__dict__
        tile.render(slide, left=left, top=2.0)

    # Avatars - different sizes and variants
    avatars = [
        (Avatar(text="JD", variant="filled", color_variant="primary", size="sm"), 0.5, 4.2),
        (Avatar(text="AS", variant="outlined", color_variant="success", size="md"), 1.5, 4.0),
        (Avatar(icon="user", variant="default", size="lg"), 3.0, 3.8),
        (Avatar(text="BM", variant="filled", color_variant="warning", size="md"), 5.0, 4.0),
    ]

    for avatar, left, top in avatars:
        avatar.theme = theme.__dict__
        avatar.render(slide, left=left, top=top)

    # Avatar with label - horizontal
    avatar_label = AvatarWithLabel(
        text="JD",
        label="John Doe",
        sublabel="Product Designer",
        variant="filled",
        color_variant="primary",
        orientation="horizontal",
        theme=theme,
    )
    avatar_label.render(slide, left=0.5, top=5.5, width=3.5)

    # Avatar group
    members = [
        {"text": "JD", "color_variant": "primary"},
        {"text": "AS", "color_variant": "success"},
        {"text": "BM", "color_variant": "warning"},
        {"text": "KL", "color_variant": "destructive"},
        {"text": "MN", "color_variant": "default"},
    ]
    group = AvatarGroup(members, max_display=3, overlap=True, size="sm", theme=theme)
    group.render(slide, left=5.0, top=5.6)


def create_combined_dashboard(prs, theme):
    """Create a realistic dashboard combining all components."""
    print("  • Creating combined components dashboard...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Dashboard Example"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Status badges at top
    Badge(text="Live", variant="success", theme=theme).render(slide, left=8.5, top=0.3)

    # Action buttons using Stack
    buttons = [
        Button(text="Refresh", variant="outline", size="sm", theme=theme),
        Button(text="Export", variant="ghost", size="sm", theme=theme),
    ]
    stack = Stack(direction="horizontal", gap="sm")
    stack.render_children(slide, buttons, left=0.5, top=1.8, item_width=1.2, item_height=0.4)

    # Key metrics using Grid
    container = Container(size="lg", padding="sm", center=True)
    bounds = container.render(slide, top=2.3)
    grid = Grid(columns=12, gap="md", bounds=bounds)

    metrics = [
        ("Total Sales", "$245K", "+18%", "up", 0),
        ("Conversion", "3.2%", "+0.5%", "up", 4),
        ("Bounce Rate", "42%", "-5%", "down", 8),
    ]

    for label, value, change, trend, col_start in metrics:
        pos = grid.get_cell(col_span=4, col_start=col_start)
        MetricCard(label=label, value=value, change=change, trend=trend, theme=theme).render(
            slide, **pos
        )

    # Alert notification
    alert = Alert(variant="info", theme=theme)
    alert.add_child(Alert.Title("New Features Available"))
    alert.add_child(Alert.Description("Check out the latest updates in the changelog."))
    alert.render(slide, left=0.5, top=4.5, width=9.0, height=0.9)

    # Status indicators - simple horizontal placement
    DotBadge(variant="success", theme=theme).render(slide, left=0.5, top=5.8)
    DotBadge(variant="warning", theme=theme).render(slide, left=0.8, top=5.8)
    DotBadge(variant="destructive", theme=theme).render(slide, left=1.1, top=5.8)


def create_shapes_showcase(prs, theme):
    """Showcase basic shape components."""
    print("  • Creating Shapes showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Shape Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Grid of different shapes
    shapes_demo = [
        ("rectangle", "Rectangle", 0.5, 1.8),
        ("rounded_rectangle", "Rounded", 2.2, 1.8),
        ("oval", "Circle", 3.9, 1.8),
        ("diamond", "Diamond", 5.6, 1.8),
        ("star", "Star", 7.3, 1.8),
        ("triangle", "Triangle", 0.5, 3.8),
        ("hexagon", "Hexagon", 2.2, 3.8),
        ("heart", "Heart", 3.9, 3.8),
        ("lightning", "Lightning", 5.6, 3.8),
        ("cloud", "Cloud", 7.3, 3.8),
    ]

    colors = [
        "primary.DEFAULT",
        "success.DEFAULT",
        "warning.DEFAULT",
        "destructive.DEFAULT",
        "accent.DEFAULT",
    ]

    for idx, (shape_type, label, left, top) in enumerate(shapes_demo):
        fill_color = colors[idx % len(colors)]
        shape = Shape(
            shape_type=shape_type, text=label, fill_color=fill_color, line_width=0, theme=theme
        )
        shape.render(slide, left=left, top=top, width=1.5, height=1.5)


def create_connectors_showcase(prs, theme):
    """Showcase connector and arrow components."""
    print("  • Creating Connectors showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Connectors & Arrows"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Create shapes to connect
    shape1 = Shape(
        shape_type="rounded_rectangle", text="Start", fill_color="primary.DEFAULT", theme=theme
    )
    shape1.render(slide, left=1.5, top=2, width=2, height=1)

    shape2 = Shape(
        shape_type="rounded_rectangle", text="Process", fill_color="secondary.DEFAULT", theme=theme
    )
    shape2.render(slide, left=4.5, top=2, width=2, height=1)

    shape3 = Shape(
        shape_type="rounded_rectangle", text="End", fill_color="success.DEFAULT", theme=theme
    )
    shape3.render(slide, left=7.5, top=2, width=2, height=1)

    # Connectors
    Connector(
        3.5, 2.5, 4.5, 2.5, "straight", "primary.DEFAULT", 3, arrow_end=True, theme=theme
    ).render(slide)
    Connector(
        6.5, 2.5, 7.5, 2.5, "straight", "primary.DEFAULT", 3, arrow_end=True, theme=theme
    ).render(slide)

    # Show connector types
    y = 4.5
    # Straight
    Shape(shape_type="oval", text="A", fill_color="accent.DEFAULT", theme=theme).render(
        slide, 1, y, 0.8, 0.8
    )
    Shape(shape_type="oval", text="B", fill_color="accent.DEFAULT", theme=theme).render(
        slide, 2.5, y, 0.8, 0.8
    )
    Connector(
        1.8, y + 0.4, 2.5, y + 0.4, "straight", "muted.foreground", 2, arrow_end=True, theme=theme
    ).render(slide)

    # Elbow
    Shape(shape_type="oval", text="C", fill_color="warning.DEFAULT", theme=theme).render(
        slide, 4, y, 0.8, 0.8
    )
    Shape(shape_type="oval", text="D", fill_color="warning.DEFAULT", theme=theme).render(
        slide, 5.5, y + 1, 0.8, 0.8
    )
    Connector(
        4.8, y + 0.4, 5.5, y + 1.4, "elbow", "muted.foreground", 2, arrow_end=True, theme=theme
    ).render(slide)

    # Curved
    Shape(shape_type="oval", text="E", fill_color="destructive.DEFAULT", theme=theme).render(
        slide, 7, y, 0.8, 0.8
    )
    Shape(shape_type="oval", text="F", fill_color="destructive.DEFAULT", theme=theme).render(
        slide, 8.5, y + 1, 0.8, 0.8
    )
    Connector(
        7.8, y + 0.4, 8.5, y + 1.4, "curved", "muted.foreground", 2, arrow_end=True, theme=theme
    ).render(slide)


def create_smartart_showcase(prs, theme):
    """Showcase SmartArt diagram components."""
    print("  • Creating SmartArt Diagrams showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "SmartArt Diagrams"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # ProcessFlow
    process = ProcessFlow(items=["Plan", "Design", "Build", "Test"], theme=theme)
    process.render(slide, left=0.5, top=1.8, width=9, height=1.3)

    # CycleDiagram
    cycle = CycleDiagram(items=["Plan", "Do", "Check", "Act"], theme=theme)
    cycle.render(slide, left=0.5, top=3.5, width=4, height=3.5)

    # HierarchyDiagram
    hierarchy = HierarchyDiagram(items=["CEO", "CTO", "CFO", "COO"], theme=theme)
    hierarchy.render(slide, left=5, top=3.5, width=4.5, height=3.5)


def create_table_showcase(prs, theme):
    """Showcase Table components with different variants."""
    print("  • Creating Table Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Table Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Default table - Top left
    table1 = Table(
        headers=["Product", "Q1", "Q2", "Q3"],
        data=[
            ["Laptops", "$100K", "$120K", "$110K"],
            ["Phones", "$80K", "$90K", "$95K"],
            ["Tablets", "$60K", "$65K", "$70K"],
        ],
        variant="default",
        size="sm",
        theme=theme,
    )
    table1.render(slide, left=0.5, top=1.8, width=4.5, height=2.2)

    # Label for default table
    label1 = TextBox(
        text="Default Variant", font_size=10, bold=True, alignment="center", theme=theme
    )
    label1.render(slide, left=0.5, top=1.6, width=4.5, height=0.3)

    # Bordered table - Top right
    table2 = Table(
        headers=["Feature", "Basic", "Pro"],
        data=[["Storage", "10GB", "100GB"], ["Users", "1", "10"], ["Support", "Email", "24/7"]],
        variant="bordered",
        size="sm",
        theme=theme,
    )
    table2.render(slide, left=5.5, top=1.8, width=4.0, height=2.2)

    # Label for bordered table
    label2 = TextBox(
        text="Bordered Variant", font_size=10, bold=True, alignment="center", theme=theme
    )
    label2.render(slide, left=5.5, top=1.6, width=4.0, height=0.3)

    # Striped table - Bottom left
    table3 = Table(
        headers=["Month", "Revenue", "Expenses", "Profit"],
        data=[
            ["January", "$50K", "$30K", "$20K"],
            ["February", "$55K", "$32K", "$23K"],
            ["March", "$60K", "$35K", "$25K"],
            ["April", "$65K", "$38K", "$27K"],
        ],
        variant="striped",
        size="md",
        theme=theme,
    )
    table3.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

    # Label for striped table
    label3 = TextBox(
        text="Striped Variant", font_size=10, bold=True, alignment="center", theme=theme
    )
    label3.render(slide, left=0.5, top=4.3, width=4.5, height=0.3)

    # Minimal table - Bottom right
    table4 = Table(
        headers=["Name", "Score"],
        data=[["Alice", "95"], ["Bob", "87"], ["Charlie", "92"]],
        variant="minimal",
        size="lg",
        theme=theme,
    )
    table4.render(slide, left=5.5, top=4.5, width=4.0, height=2.5)

    # Label for minimal table
    label4 = TextBox(
        text="Minimal Variant", font_size=10, bold=True, alignment="center", theme=theme
    )
    label4.render(slide, left=5.5, top=4.3, width=4.0, height=0.3)


def create_text_showcase(prs, theme):
    """Showcase Text components (TextBox and BulletList)."""
    print("  • Creating Text Components showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Text Components"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # TextBox examples - Row 1
    text1 = TextBox(text="Simple Text Box", font_size=18, theme=theme)
    text1.render(slide, left=0.5, top=1.8, width=4, height=0.8)

    text2 = TextBox(
        text="Centered Bold",
        font_size=20,
        bold=True,
        alignment="center",
        color="primary.DEFAULT",
        theme=theme,
    )
    text2.render(slide, left=5, top=1.8, width=4, height=0.8)

    # TextBox with different styles - Row 2
    text3 = TextBox(text="Left Aligned", alignment="left", color="secondary.DEFAULT", theme=theme)
    text3.render(slide, left=0.5, top=2.8, width=2.8, height=0.6)

    text4 = TextBox(
        text="Right Aligned", alignment="right", italic=True, color="accent.DEFAULT", theme=theme
    )
    text4.render(slide, left=3.5, top=2.8, width=2.8, height=0.6)

    text5 = TextBox(text="Justified Text Example", alignment="justify", theme=theme)
    text5.render(slide, left=6.5, top=2.8, width=3, height=0.6)

    # BulletList examples
    bullets1 = BulletList(
        items=["Revenue Growth", "Cost Reduction", "Market Expansion"],
        font_size=16,
        color="primary.DEFAULT",
        theme=theme,
    )
    bullets1.render(slide, left=0.5, top=4, width=4, height=2.5)

    bullets2 = BulletList(
        items=["Design", "Develop", "Deploy", "Monitor"],
        bullet_char="→",
        font_size=16,
        color="secondary.DEFAULT",
        theme=theme,
    )
    bullets2.render(slide, left=5, top=4, width=4, height=2.5)


def create_images_showcase(prs, theme):
    """Showcase Image component with different layouts and effects."""
    import tempfile
    from PIL import Image as PILImage

    print("  • Creating Image Components showcase...")

    # Create temporary demo images
    temp_images = []
    try:
        # Create different colored sample images
        colors = [
            (255, 107, 107),  # Red
            (78, 205, 196),  # Cyan
            (255, 195, 113),  # Orange
            (99, 110, 250),  # Blue
        ]

        for i, color in enumerate(colors):
            # Create image
            img = PILImage.new("RGB", (400, 300), color=color)
            # Save to temp file
            temp_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            img.save(temp_file.name)
            temp_images.append(temp_file.name)
            temp_file.close()

        # Slide 1: Full-bleed / Hero Image
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        theme.apply_to_slide(slide1)

        # Full-screen background image
        img_full = Image(image_source=temp_images[0], theme=theme)
        img_full.render(slide1, left=0, top=0, width=10, height=7.5)

        # Overlay title
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
        text_frame = title_box.text_frame
        text_frame.text = "Full-Bleed Hero Image"
        p = text_frame.paragraphs[0]
        p.font.size = Inches(0.5)
        p.font.bold = True
        p.font.color.rgb = theme.get_color("background.DEFAULT")

        # Slide 2: Grid Layout
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])

        title_shape = slide2.shapes.title
        if title_shape:
            title_shape.text = "Image Grid Layouts"
            # Title color set by theme.apply_to_slide()

        # Apply theme after setting title text
        theme.apply_to_slide(slide2)

        # 2x2 Grid
        positions = [
            (0.5, 1.8, 4.5, 2.5),  # Top-left
            (5.2, 1.8, 4.5, 2.5),  # Top-right
            (0.5, 4.6, 4.5, 2.5),  # Bottom-left
            (5.2, 4.6, 4.5, 2.5),  # Bottom-right
        ]

        for idx, (left, top, width, height) in enumerate(positions):
            img = Image(
                image_source=temp_images[idx % len(temp_images)],
                shadow=idx % 2 == 0,  # Alternate shadow
                theme=theme,
            )
            img.render(slide2, left=left, top=top, width=width, height=height)

        # Slide 3: Different Sizes and Effects
        slide3 = prs.slides.add_slide(prs.slide_layouts[5])

        title_shape = slide3.shapes.title
        if title_shape:
            title_shape.text = "Image Sizes & Shadow Effects"
            # Title color set by theme.apply_to_slide()

        # Apply theme after setting title text
        theme.apply_to_slide(slide3)

        # Large image with shadow
        img_large = Image(image_source=temp_images[0], shadow=True, theme=theme)
        img_large.render(slide3, left=0.5, top=1.8, width=6, height=4)

        # Small images stacked vertically (no shadow)
        for i in range(3):
            img_small = Image(image_source=temp_images[i + 1], shadow=False, theme=theme)
            img_small.render(slide3, left=7, top=1.8 + i * 1.7, width=2.5, height=1.4)

        # Slide 4: Aspect Ratio Variations
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])

        title_shape = slide4.shapes.title
        if title_shape:
            title_shape.text = "Aspect Ratios & Sizing"
            # Title color set by theme.apply_to_slide()

        # Apply theme after setting title text
        theme.apply_to_slide(slide4)

        # Width only (maintains ratio)
        img_w = Image(image_source=temp_images[0], theme=theme)
        img_w.render(slide4, left=0.5, top=1.8, width=3)

        # Height only (maintains ratio)
        img_h = Image(image_source=temp_images[1], theme=theme)
        img_h.render(slide4, left=4, top=1.8, height=2.5)

        # Fixed width and height (may distort)
        img_fixed = Image(image_source=temp_images[2], shadow=True, theme=theme)
        img_fixed.render(slide4, left=6.5, top=1.8, width=3, height=5)

        # Slide 5: Image Filters
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])

        title_shape = slide5.shapes.title
        if title_shape:
            title_shape.text = "Image Filters & Effects"
            # Title color set by theme.apply_to_slide()

        # Apply theme after setting title text
        theme.apply_to_slide(slide5)

        # Create a sample photo for filter demos (using a gradient-like pattern)
        sample_img = PILImage.new("RGB", (400, 300), color=(255, 107, 107))
        # Add some visual interest with colored rectangles
        from PIL import ImageDraw

        draw = ImageDraw.Draw(sample_img)
        draw.rectangle([50, 50, 150, 150], fill=(78, 205, 196))
        draw.rectangle([200, 100, 350, 250], fill=(255, 195, 113))
        draw.ellipse([120, 150, 280, 280], fill=(99, 110, 250))

        sample_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
        sample_img.save(sample_file.name)
        sample_file.close()
        temp_images.append(sample_file.name)

        # Layout calculation for 3x4 grid within slide boundaries
        # Slide: 10" x 7.5", Title ends at ~1.2", margins 0.5" on sides
        # Image size: 2.0" x 1.5", Column spacing: 0.3", Row spacing: 0.2"
        col_positions = [0.5, 2.8, 5.1, 7.4]  # 4 columns
        row_positions = [1.5, 3.5, 5.5]  # 3 rows (bottom of row 3: 5.5+1.8=7.3" < 7.5" ✓)

        # Row 1: Basic filters
        filter_demos_row1 = [
            ("Original", {}, col_positions[0], row_positions[0]),
            ("Blur (r=10)", {"blur_radius": 10}, col_positions[1], row_positions[0]),
            ("Grayscale", {"grayscale": True}, col_positions[2], row_positions[0]),
            ("Sepia", {"sepia": True}, col_positions[3], row_positions[0]),
        ]

        # Row 2: Adjustments
        filter_demos_row2 = [
            ("Bright +50%", {"brightness": 1.5}, col_positions[0], row_positions[1]),
            ("Dark -30%", {"brightness": 0.7}, col_positions[1], row_positions[1]),
            ("High Contrast", {"contrast": 1.8}, col_positions[2], row_positions[1]),
            ("Saturated", {"saturation": 2.0}, col_positions[3], row_positions[1]),
        ]

        # Row 3: Special effects
        filter_demos_row3 = [
            ("Sharpen", {"sharpen": True}, col_positions[0], row_positions[2]),
            ("Invert", {"invert": True}, col_positions[1], row_positions[2]),
            ("Desaturate", {"saturation": 0.3}, col_positions[2], row_positions[2]),
            (
                "Combined",
                {"blur_radius": 3, "brightness": 1.2, "saturation": 1.5},
                col_positions[3],
                row_positions[2],
            ),
        ]

        filter_demos = filter_demos_row1 + filter_demos_row2 + filter_demos_row3

        for label, filters, left, top in filter_demos:
            # Add image with filter
            img = Image(image_source=sample_file.name, **filters, theme=theme)
            img.render(slide5, left=left, top=top + 0.3, width=2.0, height=1.5)

            # Add label
            label_box = slide5.shapes.add_textbox(
                Inches(left), Inches(top), Inches(2.0), Inches(0.25)
            )
            text_frame = label_box.text_frame
            text_frame.text = label
            p = text_frame.paragraphs[0]
            p.font.size = Inches(0.1)
            p.font.bold = True
            p.font.color.rgb = theme.get_color("foreground.DEFAULT")

    finally:
        # Cleanup temp images
        import os

        for temp_path in temp_images:
            try:
                os.unlink(temp_path)
            except OSError:
                pass


def create_text_with_grid_showcase(prs, theme):
    """Showcase Text components with Grid layout system."""
    print("  • Creating Text with Grid Layout showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Text Components + Grid Layout"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Use 12-column grid system
    grid = Grid(columns=12, gap="md")

    # Row 1: Header text spanning full width (12 cols)
    pos_header = grid.get_span(col_span=12, col_start=0, left=0.5, top=1.8, width=9.0, height=0.8)
    header = TextBox(
        text="12-Column Grid Layout for Text",
        font_size=24,
        bold=True,
        alignment="center",
        color="primary.DEFAULT",
        theme=theme,
    )
    header.render(slide, **pos_header)

    # Row 2: Two columns (6 + 6) - bullet lists side by side
    pos_left = grid.get_span(col_span=6, col_start=0, left=0.5, top=2.8, width=9.0, height=2.2)
    bullets_left = BulletList(
        items=["Revenue Growth: +15%", "Customer Acquisition: 2.5K", "Market Share: 23%"],
        color="primary.DEFAULT",
        bullet_char="→",
        theme=theme,
    )
    bullets_left.render(slide, **pos_left)

    pos_right = grid.get_span(col_span=6, col_start=6, left=0.5, top=2.8, width=9.0, height=2.2)
    bullets_right = BulletList(
        items=["Q1 Performance Strong", "New Product Launch", "Team Expansion Plan"],
        color="secondary.DEFAULT",
        bullet_char="•",
        theme=theme,
    )
    bullets_right.render(slide, **pos_right)

    # Row 3: Three equal columns (4 + 4 + 4) - text boxes
    text_items = [
        ("Left Column", "left", 0),
        ("Center Column", "center", 4),
        ("Right Column", "right", 8),
    ]

    for text, align, col_start in text_items:
        pos = grid.get_span(
            col_span=4, col_start=col_start, left=0.5, top=5.2, width=9.0, height=1.0
        )
        text_box = TextBox(
            text=text, alignment=align, color="primary.DEFAULT", font_size=16, theme=theme
        )
        text_box.render(slide, **pos)


def create_text_with_stack_showcase(prs, theme):
    """Showcase Text components with Stack layout system."""
    print("  • Creating Text with Stack Layout showcase...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Text Components + Stack Layout"
        # Title color set by theme.apply_to_slide()

    # Apply theme after setting title text
    theme.apply_to_slide(slide)

    # Left side: Vertical stack of text boxes
    text_boxes = []
    for i, text in enumerate(["Introduction", "Analysis", "Recommendations", "Conclusion"]):
        text_box = TextBox(text=text, font_size=18, bold=True, color="primary.DEFAULT", theme=theme)
        text_boxes.append(text_box)

    v_stack = Stack(direction="vertical", gap="md", align="start")
    v_stack.render_children(slide, text_boxes, left=0.5, top=2.0, item_width=4.0, item_height=0.8)

    # Right side: Horizontal stack of bullet lists
    # Create multiple small bullet lists
    bullet_lists = []
    data_sets = [["Q1", "Q2", "Q3"], ["Jan", "Feb", "Mar"], ["NY", "LA", "SF"]]

    for items in data_sets:
        bullets = BulletList(
            items=items, font_size=14, bullet_char="→", color="secondary.DEFAULT", theme=theme
        )
        bullet_lists.append(bullets)

    h_stack = Stack(direction="horizontal", gap="lg", align="start")
    h_stack.render_children(slide, bullet_lists, left=5.0, top=2.0, item_width=1.5, item_height=2.5)


def create_images_with_grid_showcase(prs, theme):
    """Showcase Image components with Grid layout system."""
    import tempfile
    from PIL import Image as PILImage

    print("  • Creating Images with Grid Layout showcase...")

    # Create temporary demo images
    temp_images = []
    try:
        # Create different colored sample images
        colors = [
            (255, 107, 107),  # Red
            (78, 205, 196),  # Cyan
            (255, 195, 113),  # Orange
            (99, 110, 250),  # Blue
            (186, 104, 200),  # Purple
            (72, 219, 251),  # Light Blue
        ]

        for color in colors:
            img = PILImage.new("RGB", (400, 300), color=color)
            temp_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            img.save(temp_file.name)
            temp_images.append(temp_file.name)
            temp_file.close()

        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = "Image Components + Grid Layout"
            # Title color set by theme.apply_to_slide()

        # Apply theme after setting title text
        theme.apply_to_slide(slide)

        # Use 12-column grid for image gallery
        grid = Grid(columns=12, gap="sm")

        # Row 1: Two large images (6 + 6 columns)
        for i in range(2):
            pos = grid.get_span(
                col_span=6, col_start=i * 6, left=0.5, top=1.8, width=9.0, height=2.0
            )
            img = Image(image_source=temp_images[i], shadow=True, theme=theme)
            img.render(slide, **pos)

        # Row 2: Four small images (3 + 3 + 3 + 3 columns)
        for i in range(4):
            pos = grid.get_span(
                col_span=3, col_start=i * 3, left=0.5, top=4.0, width=9.0, height=1.8
            )
            # Apply different filters to showcase variety
            filters = [
                {"grayscale": True},
                {"sepia": True},
                {"brightness": 1.3},
                {"blur_radius": 5},
            ][i]

            img = Image(image_source=temp_images[i + 2], **filters, theme=theme)
            img.render(slide, **pos)

    finally:
        # Cleanup temp images
        import os

        for temp_path in temp_images:
            try:
                os.unlink(temp_path)
            except OSError:
                pass


def create_images_with_stack_showcase(prs, theme):
    """Showcase Image components with Stack layout system."""
    import tempfile
    from PIL import Image as PILImage

    print("  • Creating Images with Stack Layout showcase...")

    # Create temporary demo images
    temp_images = []
    try:
        # Create different colored sample images
        colors = [
            (255, 159, 64),  # Orange
            (75, 192, 192),  # Teal
            (153, 102, 255),  # Purple
            (255, 99, 132),  # Pink
        ]

        for color in colors:
            img = PILImage.new("RGB", (400, 300), color=color)
            temp_file = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
            img.save(temp_file.name)
            temp_images.append(temp_file.name)
            temp_file.close()

        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Title
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.text = "Image Components + Stack Layout"
            # Title color set by theme.apply_to_slide()

        # Apply theme after setting title text
        theme.apply_to_slide(slide)

        # Vertical stack of images on left
        v_images = []
        for i in range(3):
            img = Image(image_source=temp_images[i], shadow=True, theme=theme)
            v_images.append(img)

        v_stack = Stack(direction="vertical", gap="sm", align="start")
        v_stack.render_children(slide, v_images, left=0.5, top=2.0, item_width=4.0, item_height=1.5)

        # Horizontal stack of images on right with filters
        h_images = []
        filters_list = [
            {"brightness": 1.4},
            {"contrast": 1.6},
            {"saturation": 1.8},
        ]

        for i, filters in enumerate(filters_list):
            img = Image(image_source=temp_images[i], **filters, theme=theme)
            h_images.append(img)

        h_stack = Stack(direction="horizontal", gap="sm", align="start")
        h_stack.render_children(slide, h_images, left=5.0, top=2.5, item_width=1.4, item_height=1.8)

    finally:
        # Cleanup temp images
        import os

        for temp_path in temp_images:
            try:
                os.unlink(temp_path)
            except OSError:
                pass


def main():
    """Generate comprehensive showcase presentation."""
    print("\n🎨 Creating Core Components Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    create_button_showcase(prs, theme)
    create_badge_showcase(prs, theme)
    create_alert_showcase(prs, theme)
    create_alert_composition_showcase(prs, theme)
    create_card_showcase(prs, theme)
    create_progress_icon_timeline_showcase(prs, theme)
    create_tile_avatar_showcase(prs, theme)
    create_shapes_showcase(prs, theme)
    create_connectors_showcase(prs, theme)
    create_smartart_showcase(prs, theme)
    create_table_showcase(prs, theme)
    create_text_showcase(prs, theme)
    create_text_with_grid_showcase(prs, theme)
    create_text_with_stack_showcase(prs, theme)
    create_images_showcase(prs, theme)
    create_images_with_grid_showcase(prs, theme)
    create_images_with_stack_showcase(prs, theme)
    create_combined_dashboard(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "core_components_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n🎨 Showcase Features:")
    print("  • Button Components (all variants, sizes, icons, groups)")
    print("  • Badge Components (all variants, dots, counts, tags)")
    print("  • Alert Components (all variants, composition patterns)")
    print("  • Card Components (variants, composition, metrics)")
    print("  • Progress, Icons & Timeline (PowerPoint-specific components)")
    print("  • Tiles & Avatars (dashboard elements)")
    print("  • Shape Components (25+ geometric shapes)")
    print("  • Connectors & Arrows (straight, elbow, curved)")
    print("  • SmartArt Diagrams (process, cycle, hierarchy)")
    print("  • Table Components (default, bordered, striped, minimal variants)")
    print("  • Text Components (text boxes, bullet lists, formatting)")
    print("  • Text + Grid Layout (12-column responsive text layouts)")
    print("  • Text + Stack Layout (vertical/horizontal text stacks)")
    print("  • Image Components (layouts, filters, effects, aspect ratios)")
    print("  • Images + Grid Layout (responsive image galleries)")
    print("  • Images + Stack Layout (stacked image arrangements)")
    print("  • Combined Dashboard (real-world usage)")
    print("\n💡 Demonstrates:")
    print("  • Component-based architecture")
    print("  • Theme-aware styling")
    print("  • Composition patterns (shadcn-style)")
    print("  • Variant system (cva-inspired)")
    print("  • Design tokens and semantic colors")
    print("  • PowerPoint-specific components (ProgressBar, Icon, Timeline, Tile, Avatar)")
    print("  • Layout System Integration (Grid + Stack with components)")


if __name__ == "__main__":
    main()
