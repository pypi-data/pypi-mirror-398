#!/usr/bin/env python3
"""
Chat Conversation Showcase - Demonstration of chat/messaging components.

Shows various chat UI scenarios:
- Customer support conversations
- Team messaging
- Chat bubbles with avatars
- System messages
- Different message variants
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.chat import ChatMessage, ChatConversation
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_basic_messages_slide(prs, theme):
    """Showcase basic message variants."""
    print("  • Creating basic messages slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Chat Message Variants"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Received message (left-aligned)
    received = ChatMessage(
        text="Hello! How can I help you today?",
        sender="Support Agent",
        timestamp="10:30 AM",
        variant="received",
        theme=theme.__dict__,
    )
    received.render(slide, left=0.5, top=2.0, width=7.0)

    # Sent message (right-aligned)
    sent = ChatMessage(
        text="I need help with my account settings",
        timestamp="10:31 AM",
        variant="sent",
        theme=theme.__dict__,
    )
    sent.render(slide, left=0.5, top=3.3, width=7.0)

    # System message (centered)
    system = ChatMessage(
        text="Support Agent joined the conversation", variant="system", theme=theme.__dict__
    )
    system.render(slide, left=0.5, top=4.5, width=7.0)


def create_avatar_messages_slide(prs, theme):
    """Showcase messages with avatars."""
    print("  • Creating messages with avatars...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Messages with Avatars"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Received with avatar
    msg1 = ChatMessage(
        text="Hi! I'm here to help with any questions you have.",
        sender="Sarah (Support)",
        avatar_text="SM",
        timestamp="2:15 PM",
        variant="received",
        show_avatar=True,
        theme=theme.__dict__,
    )
    msg1.render(slide, left=0.5, top=2.0, width=7.0)

    # User message
    msg2 = ChatMessage(
        text="Great! I have a question about billing.",
        timestamp="2:16 PM",
        variant="sent",
        theme=theme.__dict__,
    )
    msg2.render(slide, left=0.5, top=3.5, width=7.0)

    # Agent response with avatar
    msg3 = ChatMessage(
        text="I'd be happy to help with billing. Can you tell me more about your question?",
        sender="Sarah (Support)",
        avatar_text="SM",
        timestamp="2:17 PM",
        variant="received",
        show_avatar=True,
        theme=theme.__dict__,
    )
    msg3.render(slide, left=0.5, top=4.7, width=7.0)


def create_customer_support_conversation(prs, theme):
    """Full customer support conversation example."""
    print("  • Creating customer support conversation...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Customer Support Chat"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Create conversation
    messages = [
        {
            "text": "Welcome to TechSupport! How can we assist you?",
            "sender": "Bot",
            "avatar_icon": "user",
            "variant": "system",
        },
        {
            "text": "Hi, I'm having trouble logging into my account",
            "timestamp": "3:45 PM",
            "variant": "sent",
        },
        {
            "text": "I'm sorry to hear that. Let me help you with that. Can you tell me your email address?",
            "sender": "Alex (Support)",
            "avatar_text": "AS",
            "timestamp": "3:46 PM",
            "variant": "received",
            "show_avatar": True,
        },
        {"text": "Sure, it's john.doe@email.com", "timestamp": "3:47 PM", "variant": "sent"},
        {
            "text": "Thank you! I've found your account. I'm sending a password reset link to your email now.",
            "sender": "Alex (Support)",
            "avatar_text": "AS",
            "timestamp": "3:48 PM",
            "variant": "received",
            "show_avatar": True,
        },
        {"text": "Perfect, I received it. Thank you!", "timestamp": "3:50 PM", "variant": "sent"},
    ]

    # Reduce spacing and adjust positioning to fit within slide boundaries
    # Slide height is 7.5", title takes ~1.3", leaving ~6.2" for content
    conversation = ChatConversation(messages, spacing=0.15, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_team_messaging_conversation(prs, theme):
    """Team messaging example."""
    print("  • Creating team messaging conversation...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Team Chat - Project Discussion"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    messages = [
        {
            "text": "Hey team! How's the new feature coming along?",
            "sender": "Sarah (PM)",
            "avatar_text": "SM",
            "timestamp": "10:00 AM",
            "variant": "received",
            "show_avatar": True,
        },
        {
            "text": "Going well! Just finished the API integration.",
            "sender": "John (Dev)",
            "avatar_text": "JD",
            "timestamp": "10:02 AM",
            "variant": "received",
            "show_avatar": True,
        },
        {
            "text": "Great work! I'm starting on the UI components now.",
            "sender": "Alice (Dev)",
            "avatar_text": "AS",
            "timestamp": "10:03 AM",
            "variant": "received",
            "show_avatar": True,
        },
        {
            "text": "Awesome team! Let's sync up at 2pm to review.",
            "sender": "Sarah (PM)",
            "avatar_text": "SM",
            "timestamp": "10:05 AM",
            "variant": "received",
            "show_avatar": True,
        },
    ]

    conversation = ChatConversation(messages, spacing=0.2, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_product_demo_conversation(prs, theme):
    """Product demo conversation."""
    print("  • Creating product demo conversation...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "AI Assistant Demo"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    messages = [
        {
            "text": "Welcome! I'm your AI assistant. Ask me anything!",
            "sender": "AI Assistant",
            "avatar_icon": "lightbulb",
            "variant": "system",
        },
        {"text": "Can you help me analyze my sales data?", "timestamp": "Now", "variant": "sent"},
        {
            "text": "Of course! I can analyze your sales trends, identify top products, and provide insights. What would you like to know?",
            "sender": "AI Assistant",
            "avatar_icon": "lightbulb",
            "timestamp": "Now",
            "variant": "received",
            "show_avatar": True,
        },
        {"text": "Show me top 5 products from last month", "timestamp": "Now", "variant": "sent"},
        {
            "text": "Here are your top 5 products: 1. Product A ($45K), 2. Product B ($38K), 3. Product C ($32K), 4. Product D ($28K), 5. Product E ($25K)",
            "sender": "AI Assistant",
            "avatar_icon": "chart",
            "timestamp": "Now",
            "variant": "received",
            "show_avatar": True,
        },
    ]

    conversation = ChatConversation(messages, spacing=0.18, theme=theme.__dict__)
    conversation.render(slide, left=1.0, top=1.8, width=8.0)


def create_mixed_scenarios_slide(prs, theme):
    """Mixed chat scenarios on one slide."""
    print("  • Creating mixed scenarios slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Chat UI Patterns"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Short exchange - left column
    messages_left = [
        {"text": "Quick question!", "timestamp": "Now", "variant": "sent"},
        {
            "text": "Sure, what's up?",
            "sender": "Alex",
            "avatar_text": "AS",
            "variant": "received",
            "show_avatar": True,
        },
        {"text": "When is the deadline?", "variant": "sent"},
        {
            "text": "Friday at 5pm",
            "sender": "Alex",
            "avatar_text": "AS",
            "variant": "received",
            "show_avatar": True,
        },
    ]

    conv1 = ChatConversation(messages_left, spacing=0.18, theme=theme.__dict__)
    conv1.render(slide, left=0.5, top=1.8, width=4.5)

    # Short exchange - right column
    messages_right = [
        {"text": "Meeting started", "variant": "system"},
        {
            "text": "Thanks for joining!",
            "sender": "Sarah",
            "avatar_text": "SM",
            "variant": "received",
            "show_avatar": True,
        },
        {"text": "Happy to be here!", "variant": "sent"},
    ]

    conv2 = ChatConversation(messages_right, spacing=0.18, theme=theme.__dict__)
    conv2.render(slide, left=5.2, top=1.8, width=4.5)


def main():
    """Generate chat conversation showcase presentation."""
    print("\n💬 Creating Chat Conversation Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    create_basic_messages_slide(prs, theme)
    create_avatar_messages_slide(prs, theme)
    create_customer_support_conversation(prs, theme)
    create_team_messaging_conversation(prs, theme)
    create_product_demo_conversation(prs, theme)
    create_mixed_scenarios_slide(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "chat_conversation_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n💬 Showcase Features:")
    print("  • Basic message variants (sent, received, system)")
    print("  • Messages with avatars")
    print("  • Customer support conversation flow")
    print("  • Team messaging scenarios")
    print("  • AI assistant product demo")
    print("  • Mixed chat UI patterns")
    print("\n💡 Demonstrates:")
    print("  • ChatMessage component with all variants")
    print("  • ChatConversation for multi-message flows")
    print("  • Avatar integration in messages")
    print("  • Real-world chat scenarios")
    print("  • Customer support, team messaging, and AI demos")


if __name__ == "__main__":
    main()
