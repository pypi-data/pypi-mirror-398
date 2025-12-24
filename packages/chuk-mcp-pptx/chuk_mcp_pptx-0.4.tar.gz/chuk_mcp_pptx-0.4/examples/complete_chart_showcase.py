#!/usr/bin/env python3
"""
Complete Chart Showcase - All 18 Chart Types
===========================================

Demonstrates every chart component with realistic data and beautiful themes.
Each chart is shown with variants and best practices.

Chart Categories:
- Basic Charts: Column, Bar, Line, Area, Pie, Doughnut
- Statistical Charts: Scatter, Bubble, Matrix3D, Radar
- Specialized Charts: Combo, Sparkline, Waterfall, Gauge
- Business Charts: Funnel, Gantt, Heatmap

Output: complete_chart_showcase.pptx
"""

import asyncio
import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Import all chart components
from chuk_mcp_pptx.components.charts import (
    # Basic Charts
    ColumnChart,
    BarChart,
    LineChart,
    AreaChart,
    PieChart,
    DoughnutChart,
    # Statistical Charts
    ScatterChart,
    BubbleChart,
    Matrix3DChart,
    RadarChart,
    # Specialized Charts
    ComboChart,
    SparklineChart,
    WaterfallChart,
    GaugeChart,
    # Business Charts
    FunnelChart,
    GanttChart,
    HeatmapChart,
)
from chuk_mcp_pptx.themes import ThemeManager


def add_slide_header(slide, title: str, subtitle: str, theme):
    """Add a styled header to the slide."""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    para = title_frame.paragraphs[0]
    para.font.size = Pt(32)
    para.font.bold = True
    para.font.color.rgb = theme.get_color("primary.DEFAULT")

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.3))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        sub_para = sub_frame.paragraphs[0]
        sub_para.font.size = Pt(14)
        sub_para.font.color.rgb = theme.get_color("muted.foreground")


async def create_title_slide(prs, theme):
    """Create an attractive title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Complete Chart Showcase"
    para = title_frame.paragraphs[0]
    para.font.size = Pt(54)
    para.font.bold = True
    para.font.color.rgb = theme.get_color("primary.DEFAULT")
    para.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
    sub_frame = sub_box.text_frame
    sub_frame.text = "18 Chart Types • Beautiful Themes • Production Ready"
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = theme.get_color("muted.foreground")
    sub_para.alignment = PP_ALIGN.CENTER

    # Categories
    cat_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(6), Inches(1.5))
    cat_frame = cat_box.text_frame
    cat_frame.text = "Basic • Statistical • Specialized • Business"
    cat_para = cat_frame.paragraphs[0]
    cat_para.font.size = Pt(18)
    cat_para.font.color.rgb = theme.get_color("foreground.DEFAULT")
    cat_para.alignment = PP_ALIGN.CENTER


# ====================================================================================
# BASIC CHARTS
# ====================================================================================


async def demo_column_chart(prs, theme):
    """Column Chart - Vertical bar comparisons."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Column Chart", "Comparing categories with vertical bars", theme)

    # Sample data - Quarterly sales
    chart = ColumnChart(
        categories=["Q1 2024", "Q2 2024", "Q3 2024", "Q4 2024"],
        series={
            "Product A": [120, 145, 132, 158],
            "Product B": [98, 112, 125, 119],
            "Product C": [87, 95, 108, 121],
        },
        variant="clustered",
        title="Quarterly Product Sales (in thousands)",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_bar_chart(prs, theme):
    """Bar Chart - Horizontal comparisons."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Bar Chart", "Horizontal layout for easy label reading", theme)

    chart = BarChart(
        categories=[
            "Customer Satisfaction",
            "Product Quality",
            "Delivery Speed",
            "Value for Money",
            "Support",
        ],
        series={"Score": [87, 92, 78, 85, 90]},
        variant="clustered",
        title="Customer Satisfaction Metrics",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_line_chart(prs, theme):
    """Line Chart - Trends over time."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Line Chart", "Visualizing trends and patterns over time", theme)

    chart = LineChart(
        categories=[
            "Jan",
            "Feb",
            "Mar",
            "Apr",
            "May",
            "Jun",
            "Jul",
            "Aug",
            "Sep",
            "Oct",
            "Nov",
            "Dec",
        ],
        series={
            "2023": [45, 52, 48, 61, 58, 72, 78, 75, 68, 71, 65, 80],
            "2024": [52, 58, 63, 68, 75, 82, 88, 91, 85, 89, 93, 98],
        },
        variant="line",
        title="Monthly Revenue Trend (in thousands)",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_area_chart(prs, theme):
    """Area Chart - Cumulative trends."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Area Chart", "Showing cumulative values and proportions", theme)

    chart = AreaChart(
        categories=["Week 1", "Week 2", "Week 3", "Week 4", "Week 5", "Week 6"],
        series={
            "Email": [1200, 1350, 1280, 1450, 1520, 1680],
            "Social": [850, 920, 1050, 1180, 1240, 1350],
            "Direct": [420, 480, 510, 560, 590, 640],
        },
        variant="stacked",
        title="Traffic Sources (Weekly)",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_pie_chart(prs, theme):
    """Pie Chart - Part-to-whole relationships."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Pie Chart", "Showing proportions of a whole", theme)

    chart = PieChart(
        categories=["North America", "Europe", "Asia Pacific", "Latin America", "Middle East"],
        values=[42, 28, 18, 8, 4],
        title="Revenue by Region (%)",
        theme=theme,
    )
    chart.render(slide, left=2, top=1.5, width=6, height=4.5)


async def demo_doughnut_chart(prs, theme):
    """Doughnut Chart - Pie with a hole."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Doughnut Chart", "Modern alternative to pie charts", theme)

    chart = DoughnutChart(
        categories=["Product", "Engineering", "Marketing", "Sales", "Operations"],
        values=[35, 25, 15, 15, 10],
        hole_size=60,
        title="Budget Allocation (%)",
        theme=theme,
    )
    chart.render(slide, left=2, top=1.5, width=6, height=4.5)


# ====================================================================================
# STATISTICAL CHARTS
# ====================================================================================


async def demo_scatter_chart(prs, theme):
    """Scatter Chart - Correlation analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Scatter Chart", "Analyzing relationships between variables", theme)

    chart = ScatterChart(
        series_data=[
            {
                "name": "Dataset A",
                "x_values": [10, 20, 30, 40, 50, 60, 70],
                "y_values": [25, 35, 45, 50, 65, 75, 80],
            },
            {
                "name": "Dataset B",
                "x_values": [15, 25, 35, 45, 55, 65, 75],
                "y_values": [30, 28, 40, 45, 55, 60, 70],
            },
        ],
        title="Marketing Spend vs. Revenue Correlation",
        show_trendline=True,
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_bubble_chart(prs, theme):
    """Bubble Chart - Three-variable visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Bubble Chart", "Displaying three dimensions of data", theme)

    chart = BubbleChart(
        series_data=[
            {"name": "Product A", "points": [[30, 75, 50], [40, 80, 60], [50, 85, 70]]},
            {"name": "Product B", "points": [[35, 65, 45], [45, 70, 55], [55, 75, 65]]},
        ],
        title="Customer Satisfaction vs. Market Share (Size = Revenue)",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_matrix3d_chart(prs, theme):
    """Matrix3D Chart - Multi-dimensional analysis."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Matrix3D Chart", "Four-dimensional data visualization", theme)

    data_points = [
        {"price": 25, "quality": 80, "sales": 120, "category": "Premium"},
        {"price": 30, "quality": 85, "sales": 150, "category": "Premium"},
        {"price": 15, "quality": 65, "sales": 200, "category": "Standard"},
        {"price": 20, "quality": 70, "sales": 180, "category": "Standard"},
        {"price": 10, "quality": 50, "sales": 250, "category": "Budget"},
        {"price": 12, "quality": 55, "sales": 230, "category": "Budget"},
    ]

    chart = Matrix3DChart(
        data_points=data_points,
        x_field="price",
        y_field="quality",
        size_field="sales",
        color_field="category",
        title="Product Analysis: Price vs Quality (Size = Sales Volume)",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_radar_chart(prs, theme):
    """Radar Chart - Multi-criteria comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Radar Chart", "Comparing multiple criteria across items", theme)

    chart = RadarChart(
        categories=["Speed", "Security", "Usability", "Features", "Support", "Price"],
        series={
            "Our Product": [85, 90, 88, 92, 87, 75],
            "Competitor A": [90, 75, 82, 85, 80, 70],
            "Competitor B": [75, 85, 90, 78, 85, 80],
        },
        variant="filled",
        title="Competitive Product Analysis",
        theme=theme,
    )
    chart.render(slide, left=1.5, top=1.5, width=7, height=4.5)


# ====================================================================================
# SPECIALIZED CHARTS
# ====================================================================================


async def demo_combo_chart(prs, theme):
    """Combo Chart - Mixed column and line."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Combo Chart", "Combining columns and lines for dual metrics", theme)

    chart = ComboChart(
        categories=["Q1", "Q2", "Q3", "Q4"],
        column_series={"Revenue": [450, 520, 485, 610]},
        line_series={"Growth %": [12, 15, 8, 18]},
        secondary_axis=["Growth %"],
        title="Revenue & Growth Rate",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_sparkline_chart(prs, theme):
    """Sparkline Chart - Minimal inline trends."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Sparkline Chart", "Compact trend visualization for dashboards", theme)

    # Show multiple sparklines with integrated labels and values
    metrics = [
        ("Daily Active Users", [1200, 1250, 1180, 1350, 1420, 1380, 1520, 1480, 1650]),
        ("Page Views", [5200, 5450, 5180, 5850, 6120, 5980, 6520, 6280, 6950]),
        ("Conversion Rate", [2.5, 2.7, 2.4, 2.9, 3.1, 2.8, 3.3, 3.0, 3.5]),
    ]

    for i, (label, values) in enumerate(metrics):
        # Sparkline with integrated label and value
        chart = SparklineChart(values=values, label=label, show_value=True, theme=theme)
        chart.render(slide, left=1, top=2 + i * 1.2, width=8, height=0.8)


async def demo_waterfall_chart(prs, theme):
    """Waterfall Chart - Sequential value changes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(
        slide, "Waterfall Chart", "Showing cumulative effect of sequential values", theme
    )

    chart = WaterfallChart(
        categories=["Starting", "Revenue", "Cost of Sales", "Operating", "Marketing", "Ending"],
        values=[100, 85, -45, -15, -8, 117],
        title="Financial Statement Analysis (in millions)",
        theme=theme,
    )
    chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_gauge_chart(prs, theme):
    """Gauge Chart - KPI progress indicator."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Gauge Chart", "Visual KPI and progress indicators", theme)

    # Show multiple gauges with integrated labels and values
    kpis = [
        ("Customer Satisfaction", 87, 0, 100),
        ("Project Completion", 65, 0, 100),
        ("Revenue Target", 78, 0, 100),
    ]

    for i, (label, value, min_val, max_val) in enumerate(kpis):
        # Gauge with integrated title and value display
        chart = GaugeChart(
            value=value, min_value=min_val, max_value=max_val, title=label, theme=theme
        )
        chart.render(slide, left=1 + i * 3, top=2.5, width=2.5, height=2.5)


# ====================================================================================
# BUSINESS CHARTS
# ====================================================================================


async def demo_funnel_chart(prs, theme):
    """Funnel Chart - Conversion process visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Funnel Chart", "Tracking conversion through stages", theme)

    chart = FunnelChart(
        stages=["Website Visitors", "Product Views", "Add to Cart", "Checkout", "Purchase"],
        values=[10000, 4500, 1800, 950, 720],
        variant="standard",
        show_percentages=True,
        show_values=True,
        title="E-commerce Conversion Funnel",
        theme=theme,
    )
    await chart.render(slide, left=2, top=1.5, width=6, height=4.5)


async def demo_gantt_chart(prs, theme):
    """Gantt Chart - Project timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Gantt Chart", "Project timeline and task scheduling", theme)

    tasks = [
        {"name": "Planning", "start": "2024-01-01", "end": "2024-01-15", "progress": 1.0},
        {"name": "Design", "start": "2024-01-10", "end": "2024-02-05", "progress": 0.85},
        {"name": "Development", "start": "2024-01-25", "end": "2024-03-20", "progress": 0.60},
        {"name": "Testing", "start": "2024-03-01", "end": "2024-03-30", "progress": 0.30},
        {"name": "Deployment", "start": "2024-03-25", "end": "2024-04-10", "progress": 0.0},
    ]

    chart = GanttChart(
        tasks=tasks,
        start_date="2024-01-01",
        end_date="2024-04-30",
        title="Product Launch Timeline - Q1 2024",
        theme=theme,
    )
    await chart.render(slide, left=0.8, top=1.5, width=8.4, height=4.5)


async def demo_heatmap_chart(prs, theme):
    """Heatmap Chart - Matrix visualization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_header(slide, "Heatmap Chart", "Color-coded data matrix for pattern detection", theme)

    chart = HeatmapChart(
        x_labels=["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"],
        y_labels=["Week 1", "Week 2", "Week 3", "Week 4"],
        data=[
            [42, 58, 63, 71, 68, 45, 38],
            [48, 62, 68, 75, 72, 52, 41],
            [45, 59, 65, 73, 70, 48, 40],
            [51, 65, 72, 78, 75, 55, 45],
        ],
        color_scale="heat",
        show_values=True,
        title="Website Traffic Heatmap (visits per hour)",
        theme=theme,
    )
    await chart.render(slide, left=1.5, top=1.5, width=7, height=4.5)


# ====================================================================================
# MAIN EXECUTION
# ====================================================================================


async def main():
    """Create the complete chart showcase presentation."""
    print("🎨 Creating Complete Chart Showcase...")

    # Initialize
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("ocean")

    # Create all slides
    slides_created = 0

    print("📊 Creating slides...")

    await create_title_slide(prs, theme)
    slides_created += 1
    print("  ✓ Title slide")

    # Basic Charts
    await demo_column_chart(prs, theme)
    slides_created += 1
    print("  ✓ Column Chart")

    await demo_bar_chart(prs, theme)
    slides_created += 1
    print("  ✓ Bar Chart")

    await demo_line_chart(prs, theme)
    slides_created += 1
    print("  ✓ Line Chart")

    await demo_area_chart(prs, theme)
    slides_created += 1
    print("  ✓ Area Chart")

    await demo_pie_chart(prs, theme)
    slides_created += 1
    print("  ✓ Pie Chart")

    await demo_doughnut_chart(prs, theme)
    slides_created += 1
    print("  ✓ Doughnut Chart")

    # Statistical Charts
    await demo_scatter_chart(prs, theme)
    slides_created += 1
    print("  ✓ Scatter Chart")

    await demo_bubble_chart(prs, theme)
    slides_created += 1
    print("  ✓ Bubble Chart")

    await demo_matrix3d_chart(prs, theme)
    slides_created += 1
    print("  ✓ Matrix3D Chart")

    await demo_radar_chart(prs, theme)
    slides_created += 1
    print("  ✓ Radar Chart")

    # Specialized Charts
    await demo_combo_chart(prs, theme)
    slides_created += 1
    print("  ✓ Combo Chart")

    await demo_sparkline_chart(prs, theme)
    slides_created += 1
    print("  ✓ Sparkline Chart")

    await demo_waterfall_chart(prs, theme)
    slides_created += 1
    print("  ✓ Waterfall Chart")

    await demo_gauge_chart(prs, theme)
    slides_created += 1
    print("  ✓ Gauge Chart")

    # Business Charts
    await demo_funnel_chart(prs, theme)
    slides_created += 1
    print("  ✓ Funnel Chart")

    await demo_gantt_chart(prs, theme)
    slides_created += 1
    print("  ✓ Gantt Chart")

    await demo_heatmap_chart(prs, theme)
    slides_created += 1
    print("  ✓ Heatmap Chart")

    # Save
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_file = os.path.join(output_dir, "complete_chart_showcase.pptx")
    prs.save(output_file)

    print(f"\n✅ Complete! Created {slides_created} slides")
    print(f"📁 Saved to: {output_file}")
    print("📊 Charts showcased: 18 types")
    print("🎨 Theme: Ocean Light")


if __name__ == "__main__":
    asyncio.run(main())
