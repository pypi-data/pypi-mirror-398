#!/usr/bin/env python3
"""
Theme Showcase by Business Domain
Creates chart galleries for different business domains across all themes.
Demonstrates the shadcn-like design system for PowerPoint.
"""

import asyncio
import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

# Import chart components
from chuk_mcp_pptx.components.charts import (
    ColumnChart,
    BarChart,
    WaterfallChart,
    LineChart,
    AreaChart,
    PieChart,
    SunburstChart,
    BubbleChart,
    RadarChart,
    ComboChart,
    FunnelChart,
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


class DomainChartGallery:
    """Create domain-specific chart galleries with different themes."""

    def __init__(self):
        self.theme_manager = ThemeManager()
        self.available_themes = [
            "dark",
            "dark-blue",
            "dark-violet",
            "dark-green",
            "dark-purple",
            "light",
            "corporate",
            "light-warm",
        ]

    async def create_general_business_charts(self, prs: Presentation, theme_name: str):
        """Create general business/strategy charts."""
        theme = self.theme_manager.get_theme(theme_name)

        # Slide 1: Revenue Analysis
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide)

        title = slide.shapes.title
        title.text = "General Business Analytics"

        # Column Chart - Revenue by Region
        column_chart = ColumnChart(
            categories=["North", "South", "East", "West", "Central"],
            series={
                "Q1": [120, 95, 110, 125, 88],
                "Q2": [135, 102, 118, 140, 95],
                "Q3": [142, 108, 125, 155, 102],
                "Q4": [160, 115, 132, 175, 110],
            },
            variant="clustered",
            title="Revenue by Region (Millions)",
            theme=theme.__dict__,
        )
        await column_chart.render(slide, left=0.5, top=1.5, width=4.5, height=2.5)

        # Waterfall Chart - Profit Bridge
        waterfall = WaterfallChart(
            categories=["Start", "Sales", "COGS", "OpEx", "Tax", "Net"],
            values=[150, 85, -45, -30, -15, 145],
            title="Profit Bridge Analysis",
            theme=theme.__dict__,
        )
        await waterfall.render(slide, left=5.0, top=1.5, width=4.5, height=2.5)

        # Line Chart - Market Growth
        line_chart = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            series={
                "Market Share": [12, 14, 15, 17, 19, 22],
                "Competitor": [18, 17, 17, 16, 16, 15],
                "Industry Avg": [15, 15.5, 16, 16.5, 17, 17.5],
            },
            variant="smooth",
            title="Market Share Trends (%)",
            theme=theme.__dict__,
        )
        await line_chart.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

        # Sunburst - Product Portfolio
        sunburst = SunburstChart(
            data={
                "name": "Portfolio",
                "value": 1000,
                "children": [
                    {
                        "name": "Software",
                        "value": 450,
                        "children": [
                            {"name": "SaaS", "value": 280},
                            {"name": "Licenses", "value": 170},
                        ],
                    },
                    {
                        "name": "Services",
                        "value": 350,
                        "children": [
                            {"name": "Consulting", "value": 200},
                            {"name": "Support", "value": 150},
                        ],
                    },
                    {"name": "Hardware", "value": 200},
                ],
            },
            title="Product Portfolio Mix",
            theme=theme.__dict__,
        )
        await sunburst.render(slide, left=5.0, top=4.5, width=4.5, height=2.5)

        # Slide 2: Sales & Conversion Analysis
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide2)

        title2 = slide2.shapes.title
        title2.text = "Sales & Conversion Analytics"

        # Funnel Chart - Sales Pipeline
        funnel = FunnelChart(
            stages=["Leads", "Qualified", "Proposals", "Negotiation", "Closed"],
            values=[10000, 6000, 3000, 1500, 800],
            variant="standard",
            show_percentages=True,
            show_values=True,
            title="Sales Pipeline Conversion",
            theme=theme.__dict__,
        )
        await funnel.render(slide2, left=2.5, top=2.0, width=5.0, height=4.0)

    async def create_tech_team_charts(self, prs: Presentation, theme_name: str):
        """Create programming/tech team charts."""
        theme = self.theme_manager.get_theme(theme_name)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide)

        title = slide.shapes.title
        title.text = "Tech Team Analytics"

        # Burndown Chart (Line)
        burndown = LineChart(
            categories=[
                "Day 1",
                "Day 2",
                "Day 3",
                "Day 4",
                "Day 5",
                "Day 6",
                "Day 7",
                "Day 8",
                "Day 9",
                "Day 10",
            ],
            series={
                "Ideal": [100, 90, 80, 70, 60, 50, 40, 30, 20, 10],
                "Actual": [100, 95, 85, 78, 65, 58, 45, 38, 25, 12],
                "Projected": [100, 95, 85, 78, 65, 55, 42, 32, 20, 8],
            },
            variant="markers",
            title="Sprint Burndown Chart",
            theme=theme.__dict__,
        )
        await burndown.render(slide, left=0.5, top=1.5, width=4.5, height=2.5)

        # Velocity Chart (Column)
        velocity = ColumnChart(
            categories=["Sprint 1", "Sprint 2", "Sprint 3", "Sprint 4", "Sprint 5", "Sprint 6"],
            series={"Committed": [45, 48, 52, 50, 55, 58], "Completed": [42, 47, 48, 52, 53, 60]},
            variant="clustered",
            title="Team Velocity (Story Points)",
            theme=theme.__dict__,
        )
        await velocity.render(slide, left=5.0, top=1.5, width=4.5, height=2.5)

        # Stacked Bar - Features vs Backlog
        feature_chart = BarChart(
            categories=["Frontend", "Backend", "Mobile", "DevOps", "QA"],
            series={
                "Completed": [25, 30, 18, 22, 28],
                "In Progress": [8, 10, 5, 6, 7],
                "Backlog": [15, 12, 20, 10, 8],
            },
            variant="stacked",
            title="Feature Status by Team",
            theme=theme.__dict__,
        )
        await feature_chart.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

        # Area Chart - Cumulative Flow
        flow_chart = AreaChart(
            categories=["Week 1", "Week 2", "Week 3", "Week 4", "Week 5", "Week 6"],
            series={
                "Done": [10, 25, 45, 70, 95, 120],
                "Testing": [5, 10, 15, 18, 20, 22],
                "In Progress": [8, 12, 10, 8, 10, 12],
                "To Do": [50, 40, 30, 25, 20, 15],
            },
            variant="stacked",
            title="Cumulative Flow Diagram",
            theme=theme.__dict__,
        )
        await flow_chart.render(slide, left=5.0, top=4.5, width=4.5, height=2.5)

    async def create_finance_charts(self, prs: Presentation, theme_name: str):
        """Create finance/CFO charts."""
        theme = self.theme_manager.get_theme(theme_name)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide)

        title = slide.shapes.title
        title.text = "Financial Analytics"

        # Waterfall - EBITDA Bridge
        ebitda = WaterfallChart(
            categories=["Revenue", "COGS", "SG&A", "R&D", "Other", "EBITDA"],
            values=[500, -200, -120, -50, -10, 120],
            title="EBITDA Bridge (Millions)",
            theme=theme.__dict__,
        )
        await ebitda.render(slide, left=0.5, top=1.5, width=4.5, height=2.5)

        # Combo Chart - Revenue & Margin
        combo = ComboChart(
            categories=["Q1-22", "Q2-22", "Q3-22", "Q4-22", "Q1-23", "Q2-23"],
            column_series={"Revenue": [420, 445, 468, 495, 510, 535]},
            line_series={
                "Gross Margin %": [42, 43, 44, 45, 46, 47],
                "EBITDA Margin %": [18, 19, 20, 21, 22, 23],
            },
            title="Revenue & Margins",
            theme=theme.__dict__,
        )
        await combo.render(slide, left=5.0, top=1.5, width=4.5, height=2.5)

        # Stacked Column - Revenue by Business Unit
        revenue_chart = ColumnChart(
            categories=["2019", "2020", "2021", "2022", "2023"],
            series={
                "Americas": [180, 195, 220, 245, 280],
                "EMEA": [120, 130, 145, 165, 185],
                "APAC": [80, 95, 115, 135, 160],
            },
            variant="stacked",
            title="Revenue by Region (Millions)",
            theme=theme.__dict__,
        )
        await revenue_chart.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

        # Pie Chart - Expense Allocation
        expense_pie = PieChart(
            categories=["Personnel", "Technology", "Marketing", "Facilities", "Other"],
            values=[45, 20, 15, 12, 8],
            variant="exploded",
            title="Operating Expense Breakdown (%)",
            theme=theme.__dict__,
        )
        await expense_pie.render(slide, left=5.0, top=4.5, width=4.5, height=2.5)

    async def create_hr_charts(self, prs: Presentation, theme_name: str):
        """Create HR/People Analytics charts."""
        theme = self.theme_manager.get_theme(theme_name)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide)

        title = slide.shapes.title
        title.text = "People Analytics"

        # Bar Chart - Headcount by Department
        headcount = BarChart(
            categories=["Engineering", "Sales", "Marketing", "Operations", "HR", "Finance"],
            series={"2022": [120, 85, 45, 60, 25, 30], "2023": [145, 95, 52, 68, 28, 35]},
            variant="clustered",
            title="Headcount by Department",
            theme=theme.__dict__,
        )
        await headcount.render(slide, left=0.5, top=1.5, width=4.5, height=2.5)

        # Radar Chart - Engagement Survey
        radar = RadarChart(
            categories=["Work-Life", "Compensation", "Growth", "Leadership", "Culture", "Benefits"],
            series={
                "Current": [7.5, 7.2, 6.8, 7.8, 8.2, 7.5],
                "Previous": [7.2, 7.0, 6.5, 7.5, 7.8, 7.3],
                "Target": [8.0, 8.0, 8.0, 8.0, 8.0, 8.0],
            },
            variant="filled",
            max_value=10,
            title="Employee Engagement Score",
            theme=theme.__dict__,
        )
        await radar.render(slide, left=5.0, top=1.5, width=4.5, height=2.5)

        # Line Chart - Attrition Trend
        attrition = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep"],
            series={
                "Voluntary": [2.1, 2.3, 2.5, 2.2, 2.0, 1.8, 1.9, 2.1, 2.0],
                "Involuntary": [0.5, 0.4, 0.6, 0.5, 0.3, 0.4, 0.5, 0.4, 0.3],
                "Total": [2.6, 2.7, 3.1, 2.7, 2.3, 2.2, 2.4, 2.5, 2.3],
            },
            variant="markers",
            title="Monthly Attrition Rate (%)",
            theme=theme.__dict__,
        )
        await attrition.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

        # Stacked Bar - Diversity Metrics
        diversity = BarChart(
            categories=["Leadership", "Management", "Individual", "Entry Level"],
            series={
                "Male": [65, 58, 52, 48],
                "Female": [32, 38, 44, 48],
                "Non-Binary": [3, 4, 4, 4],
            },
            variant="stacked",
            title="Gender Distribution by Level (%)",
            theme=theme.__dict__,
        )
        await diversity.render(slide, left=5.0, top=4.5, width=4.5, height=2.5)

    async def create_project_mgmt_charts(self, prs: Presentation, theme_name: str):
        """Create Project Management charts."""
        theme = self.theme_manager.get_theme(theme_name)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide)

        title = slide.shapes.title
        title.text = "Project Management"

        # Milestone Timeline (using Column Chart creatively)
        milestone = ColumnChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
            series={
                "Planning": [100, 0, 0, 0, 0, 0],
                "Design": [0, 100, 100, 0, 0, 0],
                "Development": [0, 0, 0, 100, 100, 0],
                "Testing": [0, 0, 0, 0, 0, 100],
            },
            variant="stacked",
            title="Project Timeline",
            theme=theme.__dict__,
        )
        await milestone.render(slide, left=0.5, top=1.5, width=4.5, height=2.5)

        # Resource Histogram
        resource = ColumnChart(
            categories=["Week 1", "Week 2", "Week 3", "Week 4", "Week 5", "Week 6"],
            series={
                "Allocated": [25, 28, 32, 35, 30, 28],
                "Available": [30, 30, 35, 35, 35, 30],
                "Required": [25, 30, 38, 40, 32, 28],
            },
            variant="clustered",
            title="Resource Allocation (FTEs)",
            theme=theme.__dict__,
        )
        await resource.render(slide, left=5.0, top=1.5, width=4.5, height=2.5)

        # Additional Line Chart for Risk Trends
        risk_trend = LineChart(
            categories=["Week 1", "Week 2", "Week 3", "Week 4", "Week 5", "Week 6"],
            series={
                "High Risk Items": [2, 3, 4, 3, 2, 1],
                "Medium Risk Items": [5, 5, 4, 4, 3, 3],
                "Low Risk Items": [8, 7, 6, 7, 8, 9],
            },
            variant="markers",
            title="Risk Items Trend",
            theme=theme.__dict__,
        )
        await risk_trend.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

        # Progress Tracking
        progress = BarChart(
            categories=["Phase 1", "Phase 2", "Phase 3", "Phase 4"],
            series={"Completed": [100, 85, 45, 10], "Remaining": [0, 15, 55, 90]},
            variant="stacked",
            title="Project Phase Completion (%)",
            theme=theme.__dict__,
        )
        await progress.render(slide, left=5.0, top=4.5, width=4.5, height=2.5)

    async def create_stock_market_charts(self, prs: Presentation, theme_name: str):
        """Create Stock Market/Trading charts."""
        theme = self.theme_manager.get_theme(theme_name)

        slide = prs.slides.add_slide(prs.slide_layouts[5])
        theme.apply_to_slide(slide)

        title = slide.shapes.title
        title.text = "Stock Market Analytics"

        # Line Chart - Index Performance
        index_chart = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug"],
            series={
                "S&P 500": [4200, 4350, 4400, 4250, 4450, 4500, 4600, 4550],
                "NASDAQ": [13500, 14000, 14200, 13800, 14500, 14800, 15000, 14900],
                "DOW": [33500, 34000, 34200, 33800, 34500, 34800, 35200, 35000],
            },
            variant="smooth",
            title="Index Performance",
            theme=theme.__dict__,
        )
        await index_chart.render(slide, left=0.5, top=1.5, width=4.5, height=2.5)

        # Column Chart - Trading Volume
        volume = ColumnChart(
            categories=["Mon", "Tue", "Wed", "Thu", "Fri"],
            series={
                "Buy Volume": [125, 142, 138, 155, 168],
                "Sell Volume": [118, 135, 142, 148, 162],
            },
            variant="clustered",
            title="Weekly Trading Volume (Millions)",
            theme=theme.__dict__,
        )
        await volume.render(slide, left=5.0, top=1.5, width=4.5, height=2.5)

        # Scatter/Bubble - Risk vs Return
        risk_return = BubbleChart(
            series_data=[
                {
                    "name": "Growth Portfolio",
                    "points": [[2.5, 15, 25], [3.5, 18, 35], [3.0, 16, 30]],  # [risk, return, size]
                },
                {
                    "name": "Balanced Portfolio",
                    "points": [[1.8, 12, 40], [2.2, 10, 38], [2.0, 9, 35]],
                },
                {
                    "name": "Conservative Portfolio",
                    "points": [[1.5, 8, 30], [1.2, 6, 28], [1.0, 5, 25]],
                },
            ],
            size_scale=2.0,
            transparency=30,
            title="Portfolio Risk vs Return Analysis",
            theme=theme.__dict__,
        )
        await risk_return.render(slide, left=0.5, top=4.5, width=4.5, height=2.5)

        # Waterfall - P&L Breakdown
        pnl = WaterfallChart(
            categories=["Opening", "Equities", "Bonds", "Options", "FX", "Fees", "Closing"],
            values=[1000, 250, 150, -80, 120, -40, 1400],
            title="Daily P&L Breakdown ($000s)",
            theme=theme.__dict__,
        )
        await pnl.render(slide, left=5.0, top=4.5, width=4.5, height=2.5)

    async def generate_all_themes(self):
        """Generate galleries for all themes and domains."""

        domains = {
            "general_business": self.create_general_business_charts,
            "tech_teams": self.create_tech_team_charts,
            "finance": self.create_finance_charts,
            "hr": self.create_hr_charts,
            "project_mgmt": self.create_project_mgmt_charts,
            "stock_market": self.create_stock_market_charts,
        }

        output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs", "theme_galleries")
        os.makedirs(output_dir, exist_ok=True)

        print("\n🎨 Generating Theme Galleries by Business Domain")
        print("=" * 70)
        print(
            f"Creating {len(self.available_themes)} themes × {len(domains)} domains = "
            f"{len(self.available_themes) * len(domains)} presentations"
        )
        print()

        for theme_name in self.available_themes:
            print(f"\n📊 Theme: {theme_name}")
            print("-" * 40)

            for domain_name, domain_func in domains.items():
                prs = Presentation()
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)

                # Add title slide
                theme = self.theme_manager.get_theme(theme_name)
                title_slide = prs.slides.add_slide(prs.slide_layouts[0])
                theme.apply_to_slide(title_slide)

                title = title_slide.shapes.title
                subtitle = (
                    title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None
                )

                domain_display = domain_name.replace("_", " ").title()
                title.text = f"{domain_display} Charts"
                if subtitle:
                    subtitle.text = f"Theme: {theme_name}"

                # Add domain charts
                await domain_func(prs, theme_name)

                # Save presentation
                filename = f"{domain_name}_{theme_name.replace('-', '_')}.pptx"
                filepath = os.path.join(output_dir, filename)
                prs.save(filepath)
                print(f"  ✅ {domain_display}: {filename}")

        print("\n" + "=" * 70)
        print("✨ All theme galleries generated successfully!")
        print(f"📁 Output directory: {output_dir}")

        # Create index file
        index_path = os.path.join(output_dir, "README.md")
        with open(index_path, "w") as f:
            f.write("# Theme Galleries by Business Domain\n\n")
            f.write("## Available Themes\n")
            for theme in self.available_themes:
                f.write(f"- **{theme}**: {theme.replace('-', ' ').title()} theme\n")
            f.write("\n## Business Domains\n")
            for domain in domains.keys():
                f.write(f"- **{domain}**: {domain.replace('_', ' ').title()}\n")
            f.write("\n## Gallery Files\n")
            f.write("Each file shows the same business charts with different themes.\n")
            f.write("Compare them side-by-side to see the power of the design system!\n")

        print(f"📝 Created index: {index_path}")

        return output_dir


async def main():
    """Main function to generate all theme galleries."""
    gallery = DomainChartGallery()
    await gallery.generate_all_themes()

    print("\n🎯 Next Steps:")
    print("1. Open the presentations to see charts with different themes")
    print("2. Compare the same domain across themes")
    print("3. Use as reference for your own presentations")
    print("\n💡 This demonstrates our 'shadcn for PowerPoint' design system:")
    print("   - One component, multiple themes")
    print("   - Consistent behavior, different styles")
    print("   - Business-appropriate chart selections")


if __name__ == "__main__":
    asyncio.run(main())
