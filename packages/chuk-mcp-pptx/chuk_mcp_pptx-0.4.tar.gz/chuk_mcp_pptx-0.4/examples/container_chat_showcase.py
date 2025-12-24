#!/usr/bin/env python3
"""
Container Chat Showcase - Chat conversations inside realistic containers.

Demonstrates how chat conversations look inside various device and window containers:
- iMessage in iPhone
- Android Messages in Samsung
- WhatsApp in iPhone
- Slack in Browser
- Teams in Windows
- ChatGPT in Browser
- Generic chat in ChatContainer
"""

import sys
import os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.chat import (
    iMessageConversation,
    AndroidConversation,
    WhatsAppConversation,
    SlackConversation,
    TeamsConversation,
    ChatGPTConversation,
    FacebookMessengerConversation,
    AIMConversation,
    MSNConversation,
)
from chuk_mcp_pptx.components.containers import (
    iPhoneContainer,
    SamsungContainer,
    BrowserWindow,
    MacOSWindow,
    WindowsWindow,
    ChatContainer,
)
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


def create_iphone_imessage_slide(prs, theme):
    """iMessage conversation inside iPhone."""
    print("  • Creating iPhone with iMessage...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "iMessage in iPhone"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # iPhone container
    iphone = iPhoneContainer(show_notch=True, theme=theme.__dict__)
    content_area = iphone.render(slide, left=3.0, top=0.8, width=4.0, height=6.5)

    # iMessage conversation inside (3 messages to fit container)
    messages = [
        {"text": "Hey! Lunch today?", "variant": "received", "timestamp": "11:30 AM"},
        {"text": "Yes! New Italian place at 12pm?", "variant": "sent", "timestamp": "11:31 AM"},
        {"text": "Perfect! 👍", "variant": "received", "timestamp": "11:32 AM"},
    ]

    conversation = iMessageConversation(messages, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_samsung_android_slide(prs, theme):
    """Android Messages conversation inside Samsung phone."""
    print("  • Creating Samsung with Android Messages...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Android Messages in Samsung"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Samsung container
    samsung = SamsungContainer(variant="galaxy-s", theme=theme.__dict__)
    content_area = samsung.render(slide, left=3.0, top=0.8, width=4.0, height=6.5)

    # Android Messages conversation inside (3 messages to fit container)
    messages = [
        {
            "text": "Meeting in 10 min!",
            "sender": "Sarah",
            "variant": "received",
            "timestamp": "2:50 PM",
        },
        {"text": "On my way. Bring laptop?", "variant": "sent", "timestamp": "2:51 PM"},
        {"text": "Yes please", "sender": "Sarah", "variant": "received", "timestamp": "2:52 PM"},
    ]

    conversation = AndroidConversation(messages, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_browser_slack_slide(prs, theme):
    """Slack conversation inside browser window."""
    print("  • Creating Browser with Slack...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Slack in Browser"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Browser window
    browser = BrowserWindow(
        title="Slack - #general",
        url="slack.com/messages/general",
        browser_type="chrome",
        theme=theme.__dict__,
    )
    content_area = browser.render(slide, left=1.0, top=1.5, width=8.0, height=5.2)

    # Slack conversation inside (reduced to 3 messages to fit container)
    messages = [
        {
            "text": "Quick standup in 5!",
            "sender": "Sarah",
            "timestamp": "9:00 AM",
            "avatar_text": "SM",
        },
        {
            "text": "On my way",
            "sender": "John",
            "timestamp": "9:01 AM",
            "avatar_text": "JD",
            "reactions": ["👍 2"],
        },
        {
            "text": "Zoom link in calendar",
            "sender": "Sarah",
            "timestamp": "9:02 AM",
            "avatar_text": "SM",
            "reactions": ["✅ 3"],
        },
    ]

    conversation = SlackConversation(messages, spacing=0.15, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_windows_teams_slide(prs, theme):
    """Teams conversation inside Windows window."""
    print("  • Creating Windows with Teams...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Microsoft Teams in Windows"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Windows window
    windows = WindowsWindow(
        title="Microsoft Teams", app_icon="👥", show_menubar=False, theme=theme.__dict__
    )
    content_area = windows.render(slide, left=1.0, top=1.5, width=8.0, height=5.5)

    # Teams conversation inside
    messages = [
        {
            "text": "Notes shared",
            "sender": "Sarah Johnson",
            "timestamp": "Today at 10:30 AM",
            "avatar_text": "SJ",
        },
        {
            "text": "Great presentation!",
            "sender": "Mike Chen",
            "timestamp": "Today at 10:35 AM",
            "avatar_text": "MC",
            "reactions": ["👍 5"],
        },
        {
            "text": "Thanks! Sync tomorrow?",
            "sender": "Sarah Johnson",
            "timestamp": "Today at 10:40 AM",
            "avatar_text": "SJ",
            "reply_count": 3,
        },
    ]

    conversation = TeamsConversation(messages, spacing=0.15, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_macos_messages_slide(prs, theme):
    """iMessage conversation inside macOS Messages app window."""
    print("  • Creating macOS Messages app...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Messages in macOS"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # macOS window
    macos_window = MacOSWindow(
        title="Messages", app_icon="💬", show_toolbar=False, theme=theme.__dict__
    )
    content_area = macos_window.render(slide, left=1.5, top=1.5, width=7.0, height=5.5)

    # iMessage conversation inside (3 messages to fit container)
    messages = [
        {"text": "Working on the presentation", "variant": "received", "timestamp": "3:00 PM"},
        {"text": "How's it going?", "variant": "sent", "timestamp": "3:01 PM"},
        {"text": "Almost done! Sending it over now", "variant": "received", "timestamp": "3:02 PM"},
    ]

    conversation = iMessageConversation(messages, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_browser_chatgpt_slide(prs, theme):
    """ChatGPT conversation inside browser window."""
    print("  • Creating Browser with ChatGPT...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "ChatGPT in Browser"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Browser window
    browser = BrowserWindow(
        title="ChatGPT", url="chat.openai.com", browser_type="safari", theme=theme.__dict__
    )
    content_area = browser.render(slide, left=1.0, top=1.5, width=8.0, height=5.2)

    # ChatGPT conversation inside (2 messages to fit container - ChatGPT has large padding)
    messages = [
        {"text": "Explain recursion simply", "variant": "user"},
        {
            "text": "Recursion is when a function calls itself to solve smaller versions of the same problem. Each call works on a simpler case until reaching a base case that stops the recursion.",
            "variant": "assistant",
        },
    ]

    conversation = ChatGPTConversation(messages, spacing=0.2, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_generic_container_slide(prs, theme):
    """Generic chat container with WhatsApp conversation."""
    print("  • Creating generic ChatContainer...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "WhatsApp in Generic Container"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Generic container
    container = ChatContainer(
        title="Team Chat",
        show_header=True,
        show_border=True,
        variant="outlined",
        theme=theme.__dict__,
    )
    content_area = container.render(slide, left=2.0, top=1.5, width=6.0, height=5.5)

    # WhatsApp conversation inside
    messages = [
        {"text": "Project status?", "sender": "Alice", "variant": "received", "timestamp": "14:00"},
        {"text": "On track! 🎯", "variant": "sent", "timestamp": "14:01"},
        {"text": "Great work team!", "sender": "Bob", "variant": "received", "timestamp": "14:02"},
    ]

    conversation = WhatsAppConversation(messages, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_comparison_slide(prs, theme):
    """Side-by-side comparison of containers."""
    print("  • Creating container comparison...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Container Comparison"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # iPhone (left)
    iphone = iPhoneContainer(show_notch=True, theme=theme.__dict__)
    iphone_area = iphone.render(slide, left=0.5, top=1.8, width=3.0, height=5.0)

    iphone_msgs = [
        {"text": "iPhone", "variant": "sent"},
        {"text": "Container", "variant": "received"},
    ]
    iphone_conv = iMessageConversation(iphone_msgs, theme=theme.__dict__)
    iphone_conv.render(
        slide, left=iphone_area["left"], top=iphone_area["top"], width=iphone_area["width"]
    )

    # Samsung (middle-left)
    samsung = SamsungContainer(theme=theme.__dict__)
    samsung_area = samsung.render(slide, left=3.75, top=1.8, width=3.0, height=5.0)

    samsung_msgs = [
        {"text": "Samsung", "variant": "sent"},
        {"text": "Container", "variant": "received"},
    ]
    samsung_conv = AndroidConversation(samsung_msgs, theme=theme.__dict__)
    samsung_conv.render(
        slide, left=samsung_area["left"], top=samsung_area["top"], width=samsung_area["width"]
    )

    # Generic container (right)
    generic = ChatContainer(title="Generic", show_header=True, theme=theme.__dict__)
    generic_area = generic.render(slide, left=7.0, top=1.8, width=2.5, height=5.0)

    generic_msgs = [
        {"text": "Generic chat container", "variant": "received"},
    ]
    generic_conv = WhatsAppConversation(generic_msgs, theme=theme.__dict__)
    generic_conv.render(
        slide, left=generic_area["left"], top=generic_area["top"], width=generic_area["width"]
    )


def create_browser_facebook_slide(prs, theme):
    """Facebook Messenger conversation inside browser window."""
    print("  • Creating Browser with Facebook Messenger...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "Facebook Messenger in Browser"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Browser window
    browser = BrowserWindow(
        title="Messenger", url="messenger.com", browser_type="chrome", theme=theme.__dict__
    )
    content_area = browser.render(slide, left=1.0, top=1.5, width=8.0, height=5.2)

    # Facebook Messenger conversation inside (3 messages to fit container)
    messages = [
        {"text": "Hey! Want to grab coffee?", "variant": "received", "avatar_text": "JS"},
        {"text": "Sure! Tomorrow at 9am?", "variant": "sent"},
        {"text": "Perfect! See you then", "variant": "received", "avatar_text": "JS"},
    ]

    conversation = FacebookMessengerConversation(messages, spacing=0.12, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_generic_aol_slide(prs, theme):
    """AOL Instant Messenger conversation inside generic container."""
    print("  • Creating ChatContainer with AOL/AIM...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "AOL Instant Messenger"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Generic container
    container = ChatContainer(
        title="AIM - Buddy Chat",
        show_header=True,
        show_border=True,
        variant="outlined",
        theme=theme.__dict__,
    )
    content_area = container.render(slide, left=2.0, top=1.5, width=6.0, height=5.5)

    # AIM conversation inside (3 messages to fit container)
    messages = [
        {
            "text": "Hey! Want to hang out later?",
            "screen_name": "sk8rgrl2004",
            "variant": "received",
            "timestamp": "5:30 PM",
        },
        {
            "text": "Yeah! Let's go to the mall at 6",
            "screen_name": "xXCoolDude2003Xx",
            "variant": "sent",
            "timestamp": "5:31 PM",
        },
        {
            "text": "Cool! See you there!",
            "screen_name": "sk8rgrl2004",
            "variant": "received",
            "timestamp": "5:32 PM",
        },
    ]

    conversation = AIMConversation(messages, spacing=0.15, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def create_generic_msn_slide(prs, theme):
    """MSN Messenger conversation inside generic container."""
    print("  • Creating ChatContainer with MSN Messenger...")
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    theme.apply_to_slide(slide)

    # Title
    title_shape = slide.shapes.title
    if title_shape:
        title_shape.text = "MSN Messenger"
        title_shape.text_frame.paragraphs[0].font.color.rgb = theme.get_color("foreground.DEFAULT")

    # Generic container
    container = ChatContainer(
        title="MSN - Conversation",
        show_header=True,
        show_border=True,
        variant="outlined",
        theme=theme.__dict__,
    )
    content_area = container.render(slide, left=2.0, top=1.5, width=6.0, height=5.5)

    # MSN conversation inside (3 messages to fit container)
    messages = [
        {
            "text": "Hey! What are you up to?",
            "display_name": "AwesomeGirl",
            "variant": "received",
            "timestamp": "21:00",
            "emoticon": "😊",
        },
        {
            "text": "Just listening to some indie rock",
            "display_name": "CoolGuy",
            "variant": "sent",
            "timestamp": "21:01",
            "emoticon": "🎵",
        },
        {
            "text": "Nice! I love indie music too",
            "display_name": "AwesomeGirl",
            "variant": "received",
            "timestamp": "21:02",
        },
    ]

    conversation = MSNConversation(messages, spacing=0.12, theme=theme.__dict__)
    conversation.render(
        slide, left=content_area["left"], top=content_area["top"], width=content_area["width"]
    )


def main():
    """Generate container chat showcase presentation."""
    print("\n📱 Creating Container Chat Showcase")
    print("=" * 70)

    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get theme
    theme_manager = ThemeManager()
    theme = theme_manager.get_theme("dark-violet")

    # Create showcase slides
    create_iphone_imessage_slide(prs, theme)
    create_samsung_android_slide(prs, theme)
    create_browser_slack_slide(prs, theme)
    create_windows_teams_slide(prs, theme)
    create_macos_messages_slide(prs, theme)
    create_browser_chatgpt_slide(prs, theme)
    create_browser_facebook_slide(prs, theme)
    create_generic_aol_slide(prs, theme)
    create_generic_msn_slide(prs, theme)
    create_generic_container_slide(prs, theme)
    create_comparison_slide(prs, theme)

    # Save presentation
    output_dir = os.path.join(os.path.dirname(__file__), "..", "outputs")
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, "container_chat_showcase.pptx")
    prs.save(output_path)

    print(f"\n✅ Created {output_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Theme: {theme.name}")
    print("\n📱 Container Showcase Features:")
    print("  Mobile Devices:")
    print("     • iPhone with iMessage")
    print("     • Samsung with Android Messages")
    print("  Desktop Windows:")
    print("     • Browser with Slack")
    print("     • Browser with ChatGPT")
    print("     • Browser with Facebook Messenger")
    print("     • Windows with Teams")
    print("     • macOS with Messages")
    print("  Generic Containers:")
    print("     • ChatContainer with AOL/AIM")
    print("     • ChatContainer with MSN Messenger")
    print("     • ChatContainer with WhatsApp")
    print("     • Side-by-side comparison")
    print("\n💡 Demonstrates:")
    print("  • Conversations rendered inside authentic containers")
    print("  • Automatic content area calculation with padding")
    print("  • Theme-aware container styling")
    print("  • Multiple container types (mobile, desktop, generic)")
    print("  • Realistic device mockups with proper chrome")
    print("  • Modern & nostalgic chat platforms (iMessage, Slack, Teams, AOL, MSN)")


if __name__ == "__main__":
    main()
