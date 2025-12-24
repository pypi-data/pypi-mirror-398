"""
Tests for text_utils module.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from chuk_mcp_pptx.utilities.text_utils import (
    extract_slide_text,
    extract_presentation_text,
    format_text_frame,
    validate_text_fit,
    auto_fit_text,
)


@pytest.fixture
def presentation():
    """Create a test presentation."""
    return Presentation()


@pytest.fixture
def slide(presentation):
    """Create a test slide."""
    blank_layout = presentation.slide_layouts[6]  # Blank layout
    return presentation.slides.add_slide(blank_layout)


class TestExtractSlideText:
    """Tests for extract_slide_text function."""

    def test_extract_empty_slide(self, slide):
        """Test extracting text from empty slide."""
        result = extract_slide_text(slide)
        assert result["title"] is None
        assert result["subtitle"] is None
        assert result["body_text"] == []
        assert result["all_text"] == []
        assert result["combined_text"] == ""

    def test_extract_title_slide(self, presentation):
        """Test extracting text from title slide."""
        title_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_layout)

        # Add title
        slide.shapes.title.text = "Test Title"

        result = extract_slide_text(slide)
        assert result["title"] == "Test Title"
        assert "Test Title" in result["all_text"]
        assert result["combined_text"] == "Test Title"

    def test_extract_text_boxes(self, slide):
        """Test extracting text from text boxes."""
        # Add text boxes
        tb1 = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb1.text_frame.text = "Text Box 1"

        tb2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(4), Inches(1))
        tb2.text_frame.text = "Text Box 2"

        result = extract_slide_text(slide)
        assert len(result["text_boxes"]) == 2
        assert result["text_boxes"][0]["text"] == "Text Box 1"
        assert result["text_boxes"][1]["text"] == "Text Box 2"
        assert "Text Box 1" in result["all_text"]
        assert "Text Box 2" in result["all_text"]

    def test_extract_table_text(self, slide):
        """Test extracting text from tables."""
        # Add table
        rows, cols = 2, 3
        table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1), Inches(4), Inches(2))
        table = table_shape.table

        # Fill table with data
        table.cell(0, 0).text = "A1"
        table.cell(0, 1).text = "B1"
        table.cell(0, 2).text = "C1"
        table.cell(1, 0).text = "A2"
        table.cell(1, 1).text = "B2"
        table.cell(1, 2).text = "C2"

        result = extract_slide_text(slide)
        assert len(result["table_text"]) == 1
        assert result["table_text"][0] == [["A1", "B1", "C1"], ["A2", "B2", "C2"]]
        assert "A1" in result["all_text"]
        assert "C2" in result["all_text"]

    def test_extract_mixed_content(self, presentation):
        """Test extracting from slide with mixed content."""
        title_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_layout)

        # Add title
        slide.shapes.title.text = "Mixed Content"

        # Add text box
        tb = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        tb.text_frame.text = "Additional text"

        result = extract_slide_text(slide)
        assert result["title"] == "Mixed Content"
        assert len(result["text_boxes"]) == 1
        assert len(result["all_text"]) == 2
        assert "Mixed Content" in result["combined_text"]
        assert "Additional text" in result["combined_text"]


class TestExtractPresentationText:
    """Tests for extract_presentation_text function."""

    def test_extract_empty_presentation(self, presentation):
        """Test extracting from presentation with no slides."""
        result = extract_presentation_text(presentation)
        assert result["total_slides"] == 0
        assert result["slides"] == []
        assert result["all_titles"] == []
        assert result["all_text"] == []

    def test_extract_multiple_slides(self, presentation):
        """Test extracting from presentation with multiple slides."""
        # Add slides with titles
        for i in range(3):
            title_layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(title_layout)
            slide.shapes.title.text = f"Slide {i + 1}"

        result = extract_presentation_text(presentation)
        assert result["total_slides"] == 3
        assert len(result["slides"]) == 3
        assert len(result["all_titles"]) == 3
        assert result["all_titles"][0]["title"] == "Slide 1"
        assert result["all_titles"][1]["title"] == "Slide 2"
        assert result["all_titles"][2]["title"] == "Slide 3"

    def test_extract_with_combined_text(self, presentation):
        """Test combined text extraction."""
        # Add slides
        title_layout = presentation.slide_layouts[0]

        slide1 = presentation.slides.add_slide(title_layout)
        slide1.shapes.title.text = "First Slide"

        slide2 = presentation.slides.add_slide(title_layout)
        slide2.shapes.title.text = "Second Slide"

        result = extract_presentation_text(presentation)
        assert "=== Slide 1 ===" in result["combined_text"]
        assert "First Slide" in result["combined_text"]
        assert "=== Slide 2 ===" in result["combined_text"]
        assert "Second Slide" in result["combined_text"]


class TestFormatTextFrame:
    """Tests for format_text_frame function."""

    def test_format_font_name(self, slide):
        """Test setting font name."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, font_name="Arial")
        assert tb.text_frame.paragraphs[0].font.name == "Arial"

    def test_format_font_size(self, slide):
        """Test setting font size."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, font_size=24)
        assert tb.text_frame.paragraphs[0].font.size == Pt(24)

    def test_format_bold_italic(self, slide):
        """Test setting bold and italic."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, bold=True, italic=True)
        assert tb.text_frame.paragraphs[0].font.bold is True
        assert tb.text_frame.paragraphs[0].font.italic is True

    def test_format_color(self, slide):
        """Test setting color."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, color=(255, 0, 0))
        rgb = tb.text_frame.paragraphs[0].font.color.rgb
        assert rgb[0] == 255
        assert rgb[1] == 0
        assert rgb[2] == 0

    def test_format_alignment_left(self, slide):
        """Test left alignment."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, alignment="left")
        from pptx.enum.text import PP_ALIGN

        assert tb.text_frame.paragraphs[0].alignment == PP_ALIGN.LEFT

    def test_format_alignment_center(self, slide):
        """Test center alignment."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, alignment="center")
        from pptx.enum.text import PP_ALIGN

        assert tb.text_frame.paragraphs[0].alignment == PP_ALIGN.CENTER

    def test_format_alignment_right(self, slide):
        """Test right alignment."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, alignment="right")
        from pptx.enum.text import PP_ALIGN

        assert tb.text_frame.paragraphs[0].alignment == PP_ALIGN.RIGHT

    def test_format_alignment_justify(self, slide):
        """Test justify alignment."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        tb.text_frame.text = "Test"

        format_text_frame(tb.text_frame, alignment="justify")
        from pptx.enum.text import PP_ALIGN

        assert tb.text_frame.paragraphs[0].alignment == PP_ALIGN.JUSTIFY


class TestValidateTextFit:
    """Tests for validate_text_fit function."""

    def test_validate_shape_without_text_frame(self, slide):
        """Test validation on shape without text frame."""
        # Add a connector (which doesn't have a text frame)
        connector = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            Inches(1),
            Inches(1),
            Inches(3),
            Inches(3),
        )

        result = validate_text_fit(connector)
        assert result["fits"] is False
        assert "error" in result

    def test_validate_short_text_fits(self, slide):
        """Test that short text fits."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = "Short text"

        result = validate_text_fit(tb)
        assert result["fits"] is True
        assert "shape_width_inches" in result
        assert "shape_height_inches" in result

    def test_validate_returns_dimensions(self, slide):
        """Test that validation returns shape dimensions."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = "Test"

        result = validate_text_fit(tb)
        assert result["shape_width_inches"] == pytest.approx(4.0, rel=0.01)
        assert result["shape_height_inches"] == pytest.approx(2.0, rel=0.01)

    def test_validate_with_custom_font_size(self, slide):
        """Test validation with custom font size."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = "Test text"

        result = validate_text_fit(tb, font_size=36)
        assert result["current_font_size"] == 36

    def test_validate_suggests_adjustments_when_too_large(self, slide):
        """Test that suggestions are provided when text doesn't fit."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(0.5))
        tb.text_frame.text = (
            "This is a very long text that definitely will not fit in a small shape"
        )

        result = validate_text_fit(tb, font_size=24)
        if not result["fits"]:
            assert "suggestions" in result
            assert "reduce_font_size_to" in result["suggestions"]


class TestAutoFitText:
    """Tests for auto_fit_text function."""

    def test_auto_fit_on_shape_without_text_frame(self, slide):
        """Test auto-fit on shape without text frame."""
        line = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(2), Inches(0))
        # Should not raise an error
        auto_fit_text(line)

    def test_auto_fit_on_empty_text(self, slide):
        """Test auto-fit on empty text."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = ""
        # Should not raise an error
        auto_fit_text(tb)

    def test_auto_fit_sets_margins(self, slide):
        """Test that auto-fit sets margins."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
        tb.text_frame.text = "Test text for auto-fit"

        auto_fit_text(tb)

        # Check that margins were set
        assert tb.text_frame.margin_left == Inches(0.1)
        assert tb.text_frame.margin_right == Inches(0.1)
        assert tb.text_frame.margin_top == Inches(0.05)
        assert tb.text_frame.margin_bottom == Inches(0.05)

    def test_auto_fit_adjusts_font_size(self, slide):
        """Test that auto-fit adjusts font size."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
        tb.text_frame.text = "This is a long text that needs to be fitted"

        auto_fit_text(tb, min_font_size=8, max_font_size=20)

        # Font size should be set
        font_size = tb.text_frame.paragraphs[0].font.size
        assert font_size is not None
        assert font_size >= Pt(8)
        assert font_size <= Pt(20)

    def test_auto_fit_respects_min_font_size(self, slide):
        """Test that auto-fit respects minimum font size."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(0.3))
        tb.text_frame.text = "Very long text that cannot possibly fit in such a tiny space"

        auto_fit_text(tb, min_font_size=12, max_font_size=20)

        # Font size should be at least min_font_size
        font_size = tb.text_frame.paragraphs[0].font.size
        assert font_size >= Pt(12)

    def test_auto_fit_respects_max_font_size(self, slide):
        """Test that auto-fit respects maximum font size."""
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
        tb.text_frame.text = "Short"

        auto_fit_text(tb, min_font_size=10, max_font_size=24)

        # Font size should not exceed max_font_size
        font_size = tb.text_frame.paragraphs[0].font.size
        assert font_size <= Pt(24)
