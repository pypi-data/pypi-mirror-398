"""
Comprehensive tests for composition patterns.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.composition import (
    ComposableComponent,
    CardHeader,
    CardContent,
    CardFooter,
    CardTitle,
    CardDescription,
    Stack,
    Separator,
    Badge,
    compose,
    with_separator,
    CompositionBuilder,
)


class TestComposableComponent:
    """Test ComposableComponent base class."""

    def test_creation(self):
        """Test creating composable component."""
        comp = ComposableComponent()
        assert comp is not None
        assert len(comp._children) == 0

    def test_add_child(self):
        """Test adding child components."""
        parent = ComposableComponent()
        child = CardTitle("Test")

        parent.add_child(child)

        assert len(parent._children) == 1
        assert child.parent == parent

    def test_add_multiple_children(self):
        """Test adding multiple children."""
        parent = ComposableComponent()

        parent.add_child(CardTitle("Title"))
        parent.add_child(CardDescription("Description"))
        parent.add_child(CardContent("Content"))

        assert len(parent._children) == 3

    def test_get_children(self):
        """Test getting children."""
        parent = ComposableComponent()
        child1 = CardTitle("Title")
        child2 = CardContent("Content")

        parent.add_child(child1)
        parent.add_child(child2)

        children = parent.get_children()
        assert len(children) == 2
        assert children[0] == child1
        assert children[1] == child2

    def test_clear_children(self):
        """Test clearing children."""
        parent = ComposableComponent()
        parent.add_child(CardTitle("Title"))
        parent.add_child(CardContent("Content"))

        assert len(parent._children) == 2

        parent.clear_children()
        assert len(parent._children) == 0

    def test_theme_inheritance(self):
        """Test that children inherit parent theme."""
        theme = {"mode": "dark", "primary_hue": "blue"}
        parent = ComposableComponent(theme=theme)
        child = CardTitle("Title")

        parent.add_child(child)

        assert child.theme == theme


class TestSubComponents:
    """Test individual subcomponent classes."""

    def test_card_header(self):
        """Test CardHeader component."""
        header = CardHeader("Title", "Description")

        assert header.title == "Title"
        assert header.description == "Description"

    def test_card_header_without_description(self):
        """Test CardHeader without description."""
        header = CardHeader("Title")

        assert header.title == "Title"
        assert header.description is None

    def test_card_title(self):
        """Test CardTitle component."""
        title = CardTitle("My Title")

        assert title.text == "My Title"

    def test_card_description(self):
        """Test CardDescription component."""
        desc = CardDescription("My Description")

        assert desc.text == "My Description"

    def test_card_content(self):
        """Test CardContent component."""
        content = CardContent("Main content here")

        assert content.content == "Main content here"

    def test_card_footer(self):
        """Test CardFooter component."""
        footer = CardFooter("Footer text", align="center")

        assert footer.text == "Footer text"
        assert footer.align == "center"

    def test_card_footer_default_align(self):
        """Test CardFooter with default alignment."""
        footer = CardFooter("Footer")

        assert footer.align == "left"

    def test_separator(self):
        """Test Separator component."""
        sep = Separator()

        assert sep is not None

    def test_badge(self):
        """Test Badge component."""
        badge = Badge("New", variant="success")

        assert badge.text == "New"
        assert badge.variant == "success"

    def test_badge_default_variant(self):
        """Test Badge with default variant."""
        badge = Badge("Default")

        assert badge.variant == "default"

    def test_stack(self):
        """Test Stack component."""
        children = [CardTitle("Title"), CardContent("Content")]
        stack = Stack(children, spacing=0.2)

        assert len(stack.children) == 2
        assert stack.spacing == 0.2


class TestSubComponentRendering:
    """Test rendering subcomponents into text frames."""

    def setup_method(self):
        """Setup test presentation and slide."""
        self.prs = Presentation()
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.shape = self.slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
        self.text_frame = self.shape.text_frame

    def test_render_card_header(self):
        """Test rendering CardHeader."""
        header = CardHeader("Test Title", "Test Description")
        header.render_into(self.text_frame)

        # Should have created paragraphs
        assert len(self.text_frame.paragraphs) >= 2

    def test_render_card_title(self):
        """Test rendering CardTitle."""
        title = CardTitle("My Title")
        title.render_into(self.text_frame)

        # Should have text
        assert self.text_frame.paragraphs[0].text == "My Title"

    def test_render_card_content(self):
        """Test rendering CardContent."""
        content = CardContent("Content text")
        content.render_into(self.text_frame)

        # Should add a paragraph
        assert any("Content text" in p.text for p in self.text_frame.paragraphs)

    def test_render_separator(self):
        """Test rendering Separator."""
        sep = Separator()
        sep.render_into(self.text_frame)

        # Should add separator line
        assert len(self.text_frame.paragraphs) >= 1

    def test_render_badge(self):
        """Test rendering Badge."""
        badge = Badge("New", "success")
        badge.render_into(self.text_frame)

        # Should add text with badge
        assert any("[New]" in p.text for p in self.text_frame.paragraphs)


class TestCompositionHelpers:
    """Test composition helper functions."""

    def test_compose(self):
        """Test compose helper."""
        children = compose(
            CardTitle("Title"), CardDescription("Description"), CardContent("Content")
        )

        assert isinstance(children, list)
        assert len(children) == 3
        assert isinstance(children[0], CardTitle)
        assert isinstance(children[1], CardDescription)
        assert isinstance(children[2], CardContent)

    def test_compose_empty(self):
        """Test compose with no arguments."""
        children = compose()

        assert children == []

    def test_with_separator(self):
        """Test with_separator helper."""
        children = with_separator(
            CardTitle("Section 1"), CardTitle("Section 2"), CardTitle("Section 3")
        )

        # Should have original items plus separators
        assert len(children) == 5  # 3 titles + 2 separators

        # Every other item should be a separator (starting from index 1)
        assert isinstance(children[1], Separator)
        assert isinstance(children[3], Separator)

    def test_with_separator_single_item(self):
        """Test with_separator with single item."""
        children = with_separator(CardTitle("Only One"))

        # No separator needed for single item
        assert len(children) == 1
        assert isinstance(children[0], CardTitle)


class TestCompositionBuilder:
    """Test CompositionBuilder fluent API."""

    def test_builder_creation(self):
        """Test creating builder."""
        builder = CompositionBuilder()

        assert builder is not None
        assert len(builder._children) == 0

    def test_builder_with_theme(self):
        """Test builder with theme."""
        theme = {"mode": "dark"}
        builder = CompositionBuilder(theme)

        assert builder.theme == theme

    def test_builder_header(self):
        """Test adding header."""
        builder = CompositionBuilder()
        builder.header("Title", "Description")

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], CardHeader)

    def test_builder_title(self):
        """Test adding title."""
        builder = CompositionBuilder()
        builder.title("My Title")

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], CardTitle)

    def test_builder_description(self):
        """Test adding description."""
        builder = CompositionBuilder()
        builder.description("My Description")

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], CardDescription)

    def test_builder_content(self):
        """Test adding content."""
        builder = CompositionBuilder()
        builder.content("Content text")

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], CardContent)

    def test_builder_footer(self):
        """Test adding footer."""
        builder = CompositionBuilder()
        builder.footer("Footer text", align="center")

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], CardFooter)
        assert children[0].align == "center"

    def test_builder_badge(self):
        """Test adding badge."""
        builder = CompositionBuilder()
        builder.badge("New", "success")

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], Badge)

    def test_builder_separator(self):
        """Test adding separator."""
        builder = CompositionBuilder()
        builder.separator()

        children = builder.build()
        assert len(children) == 1
        assert isinstance(children[0], Separator)

    def test_builder_chaining(self):
        """Test method chaining."""
        builder = CompositionBuilder()
        children = (
            builder.title("Title")
            .description("Description")
            .separator()
            .content("Content")
            .badge("New", "success")
            .footer("Footer")
            .build()
        )

        assert len(children) == 6

    def test_builder_custom_component(self):
        """Test adding custom component."""
        builder = CompositionBuilder()
        custom = CardTitle("Custom")
        builder.custom(custom)

        children = builder.build()
        assert len(children) == 1
        assert children[0] == custom

    def test_builder_multiple_builds(self):
        """Test building multiple times."""
        builder = CompositionBuilder()
        builder.title("Title")

        children1 = builder.build()
        children2 = builder.build()

        # Should return same children
        assert len(children1) == len(children2)


class TestCompositionIntegration:
    """Test composition with actual components."""

    def test_full_composition_flow(self):
        """Test complete composition workflow."""
        # Create parent component
        parent = ComposableComponent()

        # Build composition
        builder = CompositionBuilder()
        children = (
            builder.header("Dashboard", "Real-time analytics")
            .separator()
            .content("Your metrics are trending upward")
            .badge("Live", "success")
            .footer("Updated 5 min ago")
            .build()
        )

        # Add to parent
        for child in children:
            parent.add_child(child)

        # Verify
        assert len(parent.get_children()) == 5

    def test_composition_with_direct_api(self):
        """Test using direct composition API."""
        parent = ComposableComponent()

        parent.add_child(CardHeader("Title", "Subtitle"))
        parent.add_child(CardContent("Content"))
        parent.add_child(CardFooter("Footer"))

        assert len(parent._children) == 3

    def test_composition_with_compose_helper(self):
        """Test using compose helper."""
        parent = ComposableComponent()

        children = compose(
            CardTitle("Welcome"), CardDescription("Get started"), CardContent("Instructions here")
        )

        for child in children:
            parent.add_child(child)

        assert len(parent._children) == 3


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
