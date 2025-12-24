# tests/slide_templates/test_slide_templates.py
"""
Tests for slide_templates module.

Tests the template registry, base class, and all slide templates.
"""

import pytest
from pptx import Presentation


# ============================================================================
# Test Registry
# ============================================================================


class TestTemplateRegistry:
    """Tests for template registry functionality."""

    def test_list_templates_returns_list(self):
        """Test that list_templates returns a list."""
        from chuk_mcp_pptx.slide_templates.registry import list_templates

        templates = list_templates()
        assert isinstance(templates, list)

    def test_list_templates_has_registered_templates(self):
        """Test that registered templates are listed."""
        from chuk_mcp_pptx.slide_templates.registry import list_templates

        # Import to trigger registration

        templates = list_templates()
        names = [t["name"] for t in templates]
        assert "MetricsDashboard" in names
        assert "ComparisonSlide" in names
        assert "TimelineSlide" in names
        assert "ContentGridSlide" in names

    def test_list_templates_filter_by_category(self):
        """Test filtering templates by category."""
        from chuk_mcp_pptx.slide_templates.registry import list_templates

        dashboard_templates = list_templates(category="dashboard")
        for t in dashboard_templates:
            assert t["category"] == "dashboard"

    def test_list_templates_filter_nonexistent_category(self):
        """Test filtering by non-existent category returns empty list."""
        from chuk_mcp_pptx.slide_templates.registry import list_templates

        templates = list_templates(category="nonexistent")
        assert templates == []

    def test_get_template_info_existing(self):
        """Test getting info for existing template."""
        from chuk_mcp_pptx.slide_templates.registry import get_template_info
        from chuk_mcp_pptx.slide_templates import MetricsDashboard  # noqa: F401

        info = get_template_info("MetricsDashboard")
        assert info is not None
        assert info["name"] == "MetricsDashboard"
        assert "description" in info
        assert "props" in info
        assert "category" in info

    def test_get_template_info_nonexistent(self):
        """Test getting info for non-existent template."""
        from chuk_mcp_pptx.slide_templates.registry import get_template_info

        info = get_template_info("NonExistentTemplate")
        assert info is None

    def test_get_template_existing(self):
        """Test getting template class."""
        from chuk_mcp_pptx.slide_templates.registry import get_template
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        template_class = get_template("MetricsDashboard")
        assert template_class is MetricsDashboard

    def test_get_template_nonexistent(self):
        """Test getting non-existent template class."""
        from chuk_mcp_pptx.slide_templates.registry import get_template

        template_class = get_template("NonExistentTemplate")
        assert template_class is None

    def test_get_all_categories(self):
        """Test getting all categories."""
        from chuk_mcp_pptx.slide_templates.registry import get_all_categories

        categories = get_all_categories()
        assert isinstance(categories, list)
        assert "dashboard" in categories
        assert "comparison" in categories
        assert "timeline" in categories


class TestTemplateCategory:
    """Tests for TemplateCategory enum."""

    def test_category_values(self):
        """Test that all expected categories exist."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateCategory

        assert TemplateCategory.OPENING.value == "opening"
        assert TemplateCategory.CONTENT.value == "content"
        assert TemplateCategory.DASHBOARD.value == "dashboard"
        assert TemplateCategory.COMPARISON.value == "comparison"
        assert TemplateCategory.TIMELINE.value == "timeline"
        assert TemplateCategory.CLOSING.value == "closing"
        assert TemplateCategory.LAYOUT.value == "layout"


class TestTemplateProp:
    """Tests for TemplateProp model."""

    def test_create_valid_prop(self):
        """Test creating a valid template prop."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateProp

        prop = TemplateProp(
            name="title",
            type="string",
            description="Slide title",
            required=True,
        )
        assert prop.name == "title"
        assert prop.type == "string"
        assert prop.required is True

    def test_prop_with_options(self):
        """Test creating prop with options."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateProp

        prop = TemplateProp(
            name="layout",
            type="string",
            description="Layout style",
            options=["grid", "row"],
            default="grid",
        )
        assert prop.options == ["grid", "row"]
        assert prop.default == "grid"

    def test_prop_invalid_type(self):
        """Test that invalid type raises error."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateProp
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            TemplateProp(
                name="test",
                type="invalid_type",
                description="Test",
            )

    def test_prop_empty_name(self):
        """Test that empty name raises error."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateProp
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            TemplateProp(
                name="",
                type="string",
                description="Test",
            )

    def test_prop_invalid_name_characters(self):
        """Test that invalid characters in name raise error."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateProp
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            TemplateProp(
                name="invalid name!",
                type="string",
                description="Test",
            )

    def test_prop_to_dict(self):
        """Test converting prop to dict."""
        from chuk_mcp_pptx.slide_templates.registry import TemplateProp

        prop = TemplateProp(
            name="test",
            type="string",
            description="Test prop",
        )
        d = prop.to_dict()
        assert d["name"] == "test"
        assert d["type"] == "string"
        assert d["description"] == "Test prop"


class TestTemplateMetadata:
    """Tests for TemplateMetadata model."""

    def test_create_valid_metadata(self):
        """Test creating valid metadata."""
        from chuk_mcp_pptx.slide_templates.registry import (
            TemplateMetadata,
            TemplateCategory,
            TemplateProp,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        metadata = TemplateMetadata(
            name="TestTemplate",
            category=TemplateCategory.DASHBOARD,
            description="A test template",
            props=[TemplateProp(name="title", type="string", description="Title")],
            examples=[],
            tags=["test"],
            class_ref=SlideTemplate,
        )
        assert metadata.name == "TestTemplate"
        assert metadata.category == TemplateCategory.DASHBOARD

    def test_metadata_empty_name(self):
        """Test that empty name raises error."""
        from chuk_mcp_pptx.slide_templates.registry import (
            TemplateMetadata,
            TemplateCategory,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            TemplateMetadata(
                name="",
                category=TemplateCategory.DASHBOARD,
                description="A test template",
                props=[],
                examples=[],
                tags=[],
                class_ref=SlideTemplate,
            )

    def test_metadata_empty_description(self):
        """Test that empty description raises error."""
        from chuk_mcp_pptx.slide_templates.registry import (
            TemplateMetadata,
            TemplateCategory,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            TemplateMetadata(
                name="Test",
                category=TemplateCategory.DASHBOARD,
                description="",
                props=[],
                examples=[],
                tags=[],
                class_ref=SlideTemplate,
            )

    def test_metadata_tags_validation(self):
        """Test that tags are lowercased and validated."""
        from chuk_mcp_pptx.slide_templates.registry import (
            TemplateMetadata,
            TemplateCategory,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        metadata = TemplateMetadata(
            name="Test",
            category=TemplateCategory.DASHBOARD,
            description="Test template",
            props=[],
            examples=[],
            tags=["TEST", "Tag2"],
            class_ref=SlideTemplate,
        )
        assert metadata.tags == ["test", "tag2"]

    def test_metadata_empty_tag(self):
        """Test that empty tag raises error."""
        from chuk_mcp_pptx.slide_templates.registry import (
            TemplateMetadata,
            TemplateCategory,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate
        from pydantic import ValidationError

        with pytest.raises(ValidationError):
            TemplateMetadata(
                name="Test",
                category=TemplateCategory.DASHBOARD,
                description="Test template",
                props=[],
                examples=[],
                tags=["valid", ""],
                class_ref=SlideTemplate,
            )

    def test_metadata_to_dict(self):
        """Test converting metadata to dict."""
        from chuk_mcp_pptx.slide_templates.registry import (
            TemplateMetadata,
            TemplateCategory,
            TemplateProp,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        metadata = TemplateMetadata(
            name="TestTemplate",
            category=TemplateCategory.DASHBOARD,
            description="A test template",
            props=[TemplateProp(name="title", type="string", description="Title")],
            examples=[{"example": "data"}],
            tags=["test"],
            class_ref=SlideTemplate,
        )
        d = metadata.to_dict()
        assert d["name"] == "TestTemplate"
        assert d["category"] == "dashboard"
        assert d["description"] == "A test template"
        assert len(d["props"]) == 1
        assert d["tags"] == ["test"]
        # class_ref should not be in dict
        assert "class_ref" not in d


class TestTemplateDecorator:
    """Tests for @template decorator."""

    def test_template_decorator_registers(self):
        """Test that decorator registers template."""
        from chuk_mcp_pptx.slide_templates.registry import (
            template,
            TemplateCategory,
            TemplateProp,
            get_template_info,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        @template(
            name="DecoratorTestTemplate",
            category=TemplateCategory.CONTENT,
            description="Test template created by decorator",
            props=[TemplateProp(name="text", type="string", description="Text content")],
            tags=["test", "decorator"],
        )
        class DecoratorTestTemplate(SlideTemplate):
            def render(self, prs):
                return 0

        info = get_template_info("DecoratorTestTemplate")
        assert info is not None
        assert info["name"] == "DecoratorTestTemplate"
        assert info["category"] == "content"
        assert info["tags"] == ["test", "decorator"]

    def test_template_decorator_adds_metadata_to_class(self):
        """Test that decorator adds _template_metadata to class."""
        from chuk_mcp_pptx.slide_templates.registry import (
            template,
            TemplateCategory,
        )
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        @template(
            name="MetadataTestTemplate",
            category=TemplateCategory.LAYOUT,
            description="Test metadata attachment",
            props=[],
            tags=[],
        )
        class MetadataTestTemplate(SlideTemplate):
            def render(self, prs):
                return 0

        assert hasattr(MetadataTestTemplate, "_template_metadata")
        assert MetadataTestTemplate._template_metadata.name == "MetadataTestTemplate"


# ============================================================================
# Test Base Class
# ============================================================================


class TestSlideTemplateBase:
    """Tests for SlideTemplate base class."""

    def test_init_with_theme(self):
        """Test initializing with theme."""
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        class ConcreteTemplate(SlideTemplate):
            def render(self, prs):
                return 0

        theme = {"colors": {"primary": "#FF0000"}}
        template = ConcreteTemplate(theme=theme)
        assert template.theme == theme

    def test_init_without_theme(self):
        """Test initializing without theme."""
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        class ConcreteTemplate(SlideTemplate):
            def render(self, prs):
                return 0

        template = ConcreteTemplate()
        assert template.theme is None

    def test_render_not_implemented(self):
        """Test that base render raises NotImplementedError."""
        from chuk_mcp_pptx.slide_templates.base import SlideTemplate

        # Can't instantiate abstract class directly
        # but we can test that it requires implementation
        class IncompleteTemplate(SlideTemplate):
            pass

        with pytest.raises(TypeError):
            IncompleteTemplate()


# ============================================================================
# Test MetricsDashboard
# ============================================================================


class TestMetricsDashboard:
    """Tests for MetricsDashboard template."""

    @pytest.fixture
    def presentation(self):
        """Create a blank presentation."""
        return Presentation()

    def test_init(self):
        """Test initialization."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        template = MetricsDashboard(
            title="Test Dashboard",
            metrics=[
                {"label": "Revenue", "value": "$1M"},
                {"label": "Users", "value": "10K"},
            ],
        )
        assert template.title == "Test Dashboard"
        assert len(template.metrics) == 2
        assert template.layout == "grid"

    def test_init_with_theme(self):
        """Test initialization with theme."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        theme = {"colors": {"primary": "#0066CC"}}
        template = MetricsDashboard(
            title="Themed Dashboard",
            metrics=[{"label": "Test", "value": "100"}],
            theme=theme,
        )
        assert template.theme == theme

    def test_init_with_row_layout(self):
        """Test initialization with row layout."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        template = MetricsDashboard(
            title="Row Dashboard",
            metrics=[{"label": "Test", "value": "100"}],
            layout="row",
        )
        assert template.layout == "row"

    def test_render_grid_layout(self, presentation):
        """Test rendering with grid layout."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        template = MetricsDashboard(
            title="Grid Dashboard",
            metrics=[
                {"label": "Revenue", "value": "$1M", "change": "+10%", "trend": "up"},
                {"label": "Users", "value": "10K", "change": "+5%", "trend": "up"},
                {"label": "NPS", "value": "72", "change": "-2pts", "trend": "down"},
            ],
            layout="grid",
        )
        slide_index = template.render(presentation)
        assert slide_index == 0
        assert len(presentation.slides) == 1

    def test_render_row_layout(self, presentation):
        """Test rendering with row layout."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        template = MetricsDashboard(
            title="Row Dashboard",
            metrics=[
                {"label": "Revenue", "value": "$1M"},
                {"label": "Users", "value": "10K"},
            ],
            layout="row",
        )
        slide_index = template.render(presentation)
        assert slide_index == 0
        assert len(presentation.slides) == 1

    def test_render_single_metric(self, presentation):
        """Test rendering with single metric."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        template = MetricsDashboard(
            title="Single Metric",
            metrics=[{"label": "Total", "value": "100"}],
        )
        slide_index = template.render(presentation)
        assert slide_index == 0

    def test_has_metadata(self):
        """Test that template has metadata."""
        from chuk_mcp_pptx.slide_templates import MetricsDashboard

        assert hasattr(MetricsDashboard, "_template_metadata")
        metadata = MetricsDashboard._template_metadata
        assert metadata.name == "MetricsDashboard"
        assert metadata.category.value == "dashboard"


# ============================================================================
# Test ComparisonSlide
# ============================================================================


class TestComparisonSlide:
    """Tests for ComparisonSlide template."""

    @pytest.fixture
    def presentation(self):
        """Create a blank presentation."""
        return Presentation()

    def test_init(self):
        """Test initialization."""
        from chuk_mcp_pptx.slide_templates import ComparisonSlide

        template = ComparisonSlide(
            title="Build vs Buy",
            left_title="Build",
            left_items=["Full control", "Higher cost"],
            right_title="Buy",
            right_items=["Quick deployment", "Lower cost"],
        )
        assert template.title == "Build vs Buy"
        assert template.left_title == "Build"
        assert template.right_title == "Buy"
        assert len(template.left_items) == 2
        assert len(template.right_items) == 2

    def test_init_with_theme(self):
        """Test initialization with theme."""
        from chuk_mcp_pptx.slide_templates import ComparisonSlide

        theme = {"colors": {"primary": "#0066CC"}}
        template = ComparisonSlide(
            title="Test",
            left_title="Left",
            left_items=["Item"],
            right_title="Right",
            right_items=["Item"],
            theme=theme,
        )
        assert template.theme == theme

    def test_render(self, presentation):
        """Test rendering."""
        from chuk_mcp_pptx.slide_templates import ComparisonSlide

        template = ComparisonSlide(
            title="Comparison",
            left_title="Option A",
            left_items=["Pro 1", "Pro 2", "Con 1"],
            right_title="Option B",
            right_items=["Pro 1", "Con 1", "Con 2"],
        )
        slide_index = template.render(presentation)
        assert slide_index == 0
        assert len(presentation.slides) == 1

    def test_has_metadata(self):
        """Test that template has metadata."""
        from chuk_mcp_pptx.slide_templates import ComparisonSlide

        assert hasattr(ComparisonSlide, "_template_metadata")
        metadata = ComparisonSlide._template_metadata
        assert metadata.name == "ComparisonSlide"
        assert metadata.category.value == "comparison"


# ============================================================================
# Test TimelineSlide
# ============================================================================


class TestTimelineSlide:
    """Tests for TimelineSlide template."""

    @pytest.fixture
    def presentation(self):
        """Create a blank presentation."""
        return Presentation()

    def test_init(self):
        """Test initialization."""
        from chuk_mcp_pptx.slide_templates import TimelineSlide

        template = TimelineSlide(
            title="Project Timeline",
            events=[
                {"date": "Q1", "title": "Phase 1", "description": "Start project"},
                {"date": "Q2", "title": "Phase 2", "description": "Development"},
            ],
        )
        assert template.title == "Project Timeline"
        assert len(template.events) == 2

    def test_render(self, presentation):
        """Test rendering."""
        from chuk_mcp_pptx.slide_templates import TimelineSlide

        template = TimelineSlide(
            title="Timeline",
            events=[
                {"date": "Jan", "title": "Start", "description": "Begin"},
                {"date": "Jun", "title": "Mid", "description": "Midpoint"},
                {"date": "Dec", "title": "End", "description": "Finish"},
            ],
        )
        slide_index = template.render(presentation)
        assert slide_index == 0
        assert len(presentation.slides) == 1

    def test_has_metadata(self):
        """Test that template has metadata."""
        from chuk_mcp_pptx.slide_templates import TimelineSlide

        assert hasattr(TimelineSlide, "_template_metadata")
        metadata = TimelineSlide._template_metadata
        assert metadata.name == "TimelineSlide"
        assert metadata.category.value == "timeline"


# ============================================================================
# Test ContentGridSlide
# ============================================================================


class TestContentGridSlide:
    """Tests for ContentGridSlide template."""

    @pytest.fixture
    def presentation(self):
        """Create a blank presentation."""
        return Presentation()

    def test_init(self):
        """Test initialization."""
        from chuk_mcp_pptx.slide_templates import ContentGridSlide

        template = ContentGridSlide(
            title="Features",
            items=[
                {"title": "Feature 1", "description": "Description 1"},
                {"title": "Feature 2", "description": "Description 2"},
            ],
        )
        assert template.title == "Features"
        assert len(template.items) == 2

    def test_init_with_columns(self):
        """Test initialization with custom columns."""
        from chuk_mcp_pptx.slide_templates import ContentGridSlide

        template = ContentGridSlide(
            title="Grid",
            items=[{"title": "Item", "description": "Desc"}],
            columns=4,
        )
        assert template.columns == 4

    def test_render(self, presentation):
        """Test rendering."""
        from chuk_mcp_pptx.slide_templates import ContentGridSlide

        template = ContentGridSlide(
            title="Content Grid",
            items=[
                {"title": "Item 1", "description": "Description 1"},
                {"title": "Item 2", "description": "Description 2"},
                {"title": "Item 3", "description": "Description 3"},
                {"title": "Item 4", "description": "Description 4"},
            ],
            columns=2,
        )
        slide_index = template.render(presentation)
        assert slide_index == 0
        assert len(presentation.slides) == 1

    def test_has_metadata(self):
        """Test that template has metadata."""
        from chuk_mcp_pptx.slide_templates import ContentGridSlide

        assert hasattr(ContentGridSlide, "_template_metadata")
        metadata = ContentGridSlide._template_metadata
        assert metadata.name == "ContentGridSlide"
        assert metadata.category.value == "content"


# ============================================================================
# Test Module Imports
# ============================================================================


class TestModuleImports:
    """Tests for module-level imports."""

    def test_import_all_templates(self):
        """Test importing all templates."""
        from chuk_mcp_pptx.slide_templates import (
            MetricsDashboard,
            ComparisonSlide,
            TimelineSlide,
            ContentGridSlide,
        )

        assert MetricsDashboard is not None
        assert ComparisonSlide is not None
        assert TimelineSlide is not None
        assert ContentGridSlide is not None

    def test_import_registry_functions(self):
        """Test importing registry functions."""
        from chuk_mcp_pptx.slide_templates import (
            template,
            TemplateCategory,
            TemplateProp,
            TemplateMetadata,
            list_templates,
            get_template_info,
            get_template,
        )

        assert callable(template)
        assert callable(list_templates)
        assert callable(get_template_info)
        assert callable(get_template)
        assert TemplateCategory is not None
        assert TemplateProp is not None
        assert TemplateMetadata is not None

    def test_import_base_class(self):
        """Test importing base class."""
        from chuk_mcp_pptx.slide_templates import SlideTemplate

        assert SlideTemplate is not None
