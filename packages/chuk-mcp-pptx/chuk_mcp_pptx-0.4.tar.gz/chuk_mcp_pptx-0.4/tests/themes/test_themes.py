"""
Comprehensive tests for the theme system.
"""

import pytest
import json
from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.themes.theme_manager import (
    ThemeManager,
    Theme,
    CyberpunkTheme,
    GradientTheme,
    MinimalTheme,
    CorporateTheme,
)
from chuk_mcp_pptx.tokens.colors import PALETTE, GRADIENTS


class TestTheme:
    """Test Theme base class."""

    def test_theme_creation(self):
        """Test creating a basic theme."""
        theme = Theme("test-theme", primary_hue="blue", mode="dark")

        assert theme.name == "test-theme"
        assert theme.primary_hue == "blue"
        assert theme.mode == "dark"
        assert theme.font_family == "Inter"

    def test_theme_with_custom_font(self):
        """Test theme with custom font family."""
        theme = Theme("custom", primary_hue="violet", mode="light", font_family="Helvetica")

        assert theme.font_family == "Helvetica"

    def test_theme_properties(self):
        """Test theme property accessors."""
        theme = Theme("test", primary_hue="emerald", mode="dark")

        # Should have token properties
        assert theme.background
        assert theme.foreground
        assert theme.primary
        assert theme.secondary
        assert theme.accent
        assert theme.chart

        # Check structure
        assert "DEFAULT" in theme.background
        assert "DEFAULT" in theme.primary

    def test_hex_to_rgb(self):
        """Test color conversion."""
        theme = Theme("test")

        rgb = theme.hex_to_rgb("#ff0000")
        assert rgb == (255, 0, 0)

        rgb2 = theme.hex_to_rgb("00ff00")  # Without #
        assert rgb2 == (0, 255, 0)

    def test_get_color(self):
        """Test getting colors from tokens."""
        theme = Theme("test", primary_hue="blue", mode="dark")

        # Should return RGBColor
        color = theme.get_color("primary.DEFAULT")
        assert color is not None

        # Background color
        bg = theme.get_color("background.DEFAULT")
        assert bg is not None

    def test_get_chart_colors(self):
        """Test getting chart colors."""
        theme = Theme("test")

        colors = theme.get_chart_colors()
        assert isinstance(colors, list)
        assert len(colors) > 0

    def test_to_dict(self):
        """Test theme to dictionary conversion."""
        theme = Theme("test", primary_hue="violet", mode="light", font_family="Arial")

        data = theme.to_dict()
        assert data["name"] == "test"
        assert data["primary_hue"] == "violet"
        assert data["mode"] == "light"
        assert data["font_family"] == "Arial"

    def test_export_json(self):
        """Test JSON export."""
        theme = Theme("test")

        json_str = theme.export_json()
        assert isinstance(json_str, str)

        # Should be valid JSON
        data = json.loads(json_str)
        assert data["name"] == "test"

    def test_from_dict(self):
        """Test creating theme from dictionary."""
        config = {
            "name": "custom",
            "primary_hue": "emerald",
            "mode": "light",
            "font_family": "Georgia",
        }

        theme = Theme.from_dict(config)
        assert theme.name == "custom"
        assert theme.primary_hue == "emerald"
        assert theme.mode == "light"
        assert theme.font_family == "Georgia"

    def test_apply_to_slide(self):
        """Test applying theme to a slide."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        theme = Theme("test", primary_hue="blue", mode="dark")
        theme.apply_to_slide(slide)

        # Slide background should be set
        assert slide.background.fill.type is not None

    def test_apply_to_shape(self):
        """Test applying theme to a shape."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        theme = Theme("test")
        theme.apply_to_shape(shape, "primary")

        # Shape should have styling applied
        assert shape.fill is not None


class TestThemeManager:
    """Test ThemeManager class."""

    def test_theme_manager_creation(self):
        """Test creating theme manager."""
        mgr = ThemeManager()
        assert mgr is not None
        assert len(mgr.themes) > 0

    def test_builtin_themes(self):
        """Test that built-in themes are registered."""
        mgr = ThemeManager()

        # Dark themes
        assert "dark" in mgr.themes
        assert "dark-blue" in mgr.themes
        assert "dark-violet" in mgr.themes

        # Light themes
        assert "light" in mgr.themes
        assert "light-blue" in mgr.themes

        # Special themes
        assert "cyberpunk" in mgr.themes
        assert "minimal" in mgr.themes
        assert "corporate" in mgr.themes

    def test_get_theme(self):
        """Test getting a theme."""
        mgr = ThemeManager()

        theme = mgr.get_theme("dark-violet")
        assert theme is not None
        assert theme.name == "dark-violet"
        assert theme.mode == "dark"
        assert theme.primary_hue == "violet"

    def test_get_nonexistent_theme(self):
        """Test getting a theme that doesn't exist."""
        mgr = ThemeManager()

        theme = mgr.get_theme("nonexistent")
        assert theme is None

    def test_register_custom_theme(self):
        """Test registering a custom theme."""
        mgr = ThemeManager()

        custom = Theme("my-theme", primary_hue="amber", mode="dark")
        mgr.register_theme(custom)

        assert "my-theme" in mgr.themes
        retrieved = mgr.get_theme("my-theme")
        assert retrieved.primary_hue == "amber"

    def test_set_current_theme(self):
        """Test setting current theme."""
        mgr = ThemeManager()

        mgr.set_current_theme("dark-green")
        assert mgr.current_theme is not None
        assert mgr.current_theme.name == "dark-green"

    def test_set_invalid_current_theme(self):
        """Test setting invalid current theme raises error."""
        mgr = ThemeManager()

        with pytest.raises(ValueError):
            mgr.set_current_theme("nonexistent")

    def test_list_themes(self):
        """Test listing all themes."""
        mgr = ThemeManager()

        themes = mgr.list_themes()
        assert isinstance(themes, list)
        assert len(themes) > 0
        assert "dark" in themes
        assert "light" in themes

    def test_list_themes_by_mode(self):
        """Test filtering themes by mode."""
        mgr = ThemeManager()

        dark_themes = mgr.list_themes_by_mode("dark")
        assert all(
            "dark" in theme or theme == "cyberpunk"
            for theme in dark_themes
            if theme != "sunset" and theme != "ocean" and theme != "aurora"
        )

        light_themes = mgr.list_themes_by_mode("light")
        assert "light" in light_themes
        assert "minimal" in light_themes
        assert "corporate" in light_themes

    def test_get_theme_info(self):
        """Test getting theme info as dict."""
        mgr = ThemeManager()

        info = mgr.get_theme_info("dark-blue")
        assert info is not None
        assert info["name"] == "dark-blue"
        assert info["mode"] == "dark"
        assert info["primary_hue"] == "blue"
        assert "colors" in info

    def test_export_theme(self):
        """Test exporting a theme."""
        mgr = ThemeManager()

        json_str = mgr.export_theme("dark-violet")
        assert json_str is not None

        data = json.loads(json_str)
        assert data["name"] == "dark-violet"

    def test_export_nonexistent_theme(self):
        """Test exporting nonexistent theme returns None."""
        mgr = ThemeManager()

        result = mgr.export_theme("nonexistent")
        assert result is None

    def test_export_all_themes(self):
        """Test exporting all themes."""
        mgr = ThemeManager()

        json_str = mgr.export_all_themes()
        assert json_str is not None

        data = json.loads(json_str)
        assert "dark" in data
        assert "light" in data

    def test_apply_to_slide(self):
        """Test applying theme to slide via manager."""
        mgr = ThemeManager()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        mgr.apply_to_slide(slide, "dark-blue")

        # Background should be set
        assert slide.background.fill.type is not None


class TestSpecialThemes:
    """Test special theme classes."""

    def test_cyberpunk_theme(self):
        """Test cyberpunk theme."""
        theme = CyberpunkTheme()

        assert theme.name == "cyberpunk"
        assert theme.mode == "dark"
        assert theme.primary_hue == "violet"

        # Should have cyberpunk colors (from PALETTE)
        assert theme.tokens["background"]["DEFAULT"] == "#020617"  # PALETTE["slate"][950]
        assert theme.tokens["foreground"]["DEFAULT"] == "#22d3ee"  # PALETTE["cyan"][400]

    def test_gradient_theme(self):
        """Test gradient theme."""
        gradient_colors = GRADIENTS["sunset"]
        theme = GradientTheme("sunset", gradient_colors)

        assert theme.name == "sunset"
        assert theme.gradient_colors == gradient_colors
        assert len(theme.gradient_colors) >= 3

    def test_minimal_theme(self):
        """Test minimal theme."""
        theme = MinimalTheme()

        assert theme.name == "minimal"
        assert theme.mode == "light"
        assert theme.tokens["background"]["DEFAULT"] == "#fafafa"  # PALETTE["zinc"][50]
        assert theme.tokens["foreground"]["DEFAULT"] == "#09090b"  # PALETTE["zinc"][950]

    def test_corporate_theme(self):
        """Test corporate theme."""
        theme = CorporateTheme()

        assert theme.name == "corporate"
        assert theme.mode == "light"
        assert theme.primary_hue == "blue"


class TestThemeVariations:
    """Test theme variations and edge cases."""

    def test_different_primary_hues(self):
        """Test themes with different primary hues."""
        hues = ["blue", "violet", "emerald", "orange", "red"]

        for hue in hues:
            theme = Theme(f"test-{hue}", primary_hue=hue, mode="dark")
            assert theme.primary_hue == hue

            # Primary color should be from that hue
            primary_color = theme.tokens["primary"]["DEFAULT"]
            assert primary_color in PALETTE[hue].values()

    def test_dark_vs_light_mode(self):
        """Test differences between dark and light modes."""
        dark = Theme("dark-test", primary_hue="blue", mode="dark")
        light = Theme("light-test", primary_hue="blue", mode="light")

        # Backgrounds should be different
        assert dark.tokens["background"]["DEFAULT"] != light.tokens["background"]["DEFAULT"]

        # Dark should have dark background
        dark_bg = dark.tokens["background"]["DEFAULT"]
        assert dark_bg in PALETTE["zinc"].values()

        # Light should have white or very light background
        light_bg = light.tokens["background"]["DEFAULT"]
        assert light_bg == "#ffffff" or light_bg in [PALETTE["zinc"][50]]

    def test_theme_immutability(self):
        """Test that theme tokens are properly isolated."""
        theme1 = Theme("theme1", primary_hue="blue")
        theme2 = Theme("theme2", primary_hue="violet")

        # Modifying one shouldn't affect the other
        theme1_primary = theme1.tokens["primary"]["DEFAULT"]
        theme2_primary = theme2.tokens["primary"]["DEFAULT"]

        assert theme1_primary != theme2_primary


class TestThemeIntegration:
    """Test theme integration with other systems."""

    def test_theme_with_components(self):
        """Test using theme with component system."""
        theme = Theme("test", primary_hue="emerald", mode="dark")

        # Should be able to convert to dict for components
        theme_dict = theme.to_dict()
        assert theme_dict["primary_hue"] == "emerald"

    def test_theme_manager_roundtrip(self):
        """Test exporting and reimporting theme."""
        mgr = ThemeManager()

        # Export a theme
        json_str = mgr.export_theme("dark-violet")
        data = json.loads(json_str)

        # Create new theme from exported data (without colors since those are derived)
        config = {
            "name": data["name"],
            "mode": data["mode"],
            "primary_hue": data["primary_hue"],
            "font_family": data["font_family"],
        }
        new_theme = Theme.from_dict(config)

        # Should match original
        assert new_theme.name == data["name"]
        assert new_theme.mode == data["mode"]
        assert new_theme.primary_hue == data["primary_hue"]


class TestThemeManagerBranches:
    """Test ThemeManager branch coverage."""

    def test_get_default_theme(self):
        """Test getting default theme."""
        mgr = ThemeManager()
        theme = mgr.get_default_theme()
        assert theme is not None
        assert theme.name == "dark"

    def test_get_theme_info_nonexistent(self):
        """Test get_theme_info for nonexistent theme."""
        mgr = ThemeManager()
        info = mgr.get_theme_info("nonexistent-theme")
        assert info is None

    def test_apply_to_slide_no_theme_name(self):
        """Test apply_to_slide without theme name uses current or default."""
        mgr = ThemeManager()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # No current theme set, should use default
        mgr.apply_to_slide(slide)
        assert slide.background.fill.type is not None

    def test_apply_to_slide_with_current_theme(self):
        """Test apply_to_slide uses current theme when set."""
        mgr = ThemeManager()
        mgr.set_current_theme("dark-violet")

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        mgr.apply_to_slide(slide)  # Uses current theme
        assert slide.background.fill.type is not None

    def test_apply_to_slide_no_theme_fallback(self):
        """Test apply_to_slide falls back to 'dark' theme."""
        mgr = ThemeManager()
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Clear current theme
        mgr.current_theme = None

        # Pass nonexistent theme name
        mgr.apply_to_slide(slide, "nonexistent")
        # Should fall back to dark theme
        assert slide.background.fill.type is not None


class TestThemeGetColorBranches:
    """Test Theme.get_color branch coverage."""

    def test_get_color_nested_path(self):
        """Test getting color with nested path."""
        theme = Theme("test", primary_hue="blue", mode="dark")

        # Get nested color
        color = theme.get_color("primary.DEFAULT")
        assert color is not None

    def test_get_color_invalid_path(self):
        """Test getting color with invalid path returns black."""
        theme = Theme("test")

        # Invalid path should return black
        color = theme.get_color("nonexistent.path.deep")
        assert color is not None
        # Should be black (0, 0, 0)
        assert color[0] == 0 and color[1] == 0 and color[2] == 0

    def test_get_color_value_not_dict(self):
        """Test get_color when intermediate value is not a dict."""
        theme = Theme("test")

        # Set a non-dict value at a path
        theme.tokens["simple"] = "#123456"

        color = theme.get_color("simple")
        assert color is not None

    def test_get_color_final_value_not_string(self):
        """Test get_color when final value is not a string."""
        theme = Theme("test")

        # Set a non-string value
        theme.tokens["weird"] = 123

        color = theme.get_color("weird")
        # Should return black
        assert color[0] == 0 and color[1] == 0 and color[2] == 0


class TestThemeApplyToSlideBranches:
    """Test Theme.apply_to_slide branch coverage."""

    def test_apply_to_slide_override_text_colors_true(self):
        """Test apply_to_slide with text color override."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add text shape
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Test"

        theme = Theme("test", mode="dark")
        theme.apply_to_slide(slide, override_text_colors=True)

        # Verify background set
        assert slide.background.fill.type is not None

    def test_apply_to_slide_override_text_colors_false(self):
        """Test apply_to_slide without text color override."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add text shape
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = "Test"

        theme = Theme("test", mode="dark")
        theme.apply_to_slide(slide, override_text_colors=False)

        # Background should be set but text colors preserved
        assert slide.background.fill.type is not None

    def test_apply_to_slide_with_multiple_runs(self):
        """Test apply_to_slide with multiple text runs."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add text shape with multiple runs
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = "First"
        run2 = p.add_run()
        run2.text = "Second"

        theme = Theme("test", mode="dark")
        theme.apply_to_slide(slide, override_text_colors=True)

        # Both runs should have color applied
        assert slide.background.fill.type is not None


class TestThemeApplyToShapeBranches:
    """Test Theme.apply_to_shape branch coverage."""

    def test_apply_to_shape_all_styles(self):
        """Test apply_to_shape with all style types."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        styles = ["card", "primary", "secondary", "accent", "muted", "unknown"]

        for style in styles:
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = "Test"

            theme = Theme("test")
            theme.apply_to_shape(shape, style)

            # Shape should have fill applied
            assert shape.fill.type is not None

    def test_apply_to_shape_without_fill(self):
        """Test apply_to_shape when shape has no fill attribute."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        from unittest.mock import MagicMock

        mock_shape = MagicMock(spec=["line", "text_frame"])  # No fill
        mock_shape.line = shape.line
        mock_shape.text_frame = shape.text_frame

        theme = Theme("test")
        theme.apply_to_shape(mock_shape, "primary")
        # Should not raise

    def test_apply_to_shape_without_line(self):
        """Test apply_to_shape when shape has no line attribute."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        from unittest.mock import MagicMock

        mock_shape = MagicMock(spec=["fill", "text_frame"])  # No line
        mock_shape.fill = shape.fill
        mock_shape.text_frame = shape.text_frame

        theme = Theme("test")
        theme.apply_to_shape(mock_shape, "primary")
        # Should not raise

    def test_apply_to_shape_without_text_frame(self):
        """Test apply_to_shape when shape has no text_frame."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        from unittest.mock import MagicMock

        mock_shape = MagicMock(spec=["fill", "line"])  # No text_frame
        mock_shape.fill = shape.fill
        mock_shape.line = shape.line

        theme = Theme("test")
        theme.apply_to_shape(mock_shape, "primary")
        # Should not raise

    def test_apply_to_shape_with_empty_text_frame(self):
        """Test apply_to_shape when text_frame is empty."""
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

        # Don't add any text
        theme = Theme("test")
        theme.apply_to_shape(shape, "primary")
        # Should not raise


class TestGradientThemeBranches:
    """Test GradientTheme branch coverage."""

    def test_gradient_theme_apply_to_slide(self):
        """Test GradientTheme.apply_to_slide uses first gradient color."""
        gradient_colors = GRADIENTS["sunset"]
        theme = GradientTheme("sunset", gradient_colors)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        theme.apply_to_slide(slide)

        # Background should be set to first gradient color
        assert slide.background.fill.type is not None


class TestThemeFromDictDefaults:
    """Test Theme.from_dict with missing values."""

    def test_from_dict_empty(self):
        """Test creating theme from empty dict uses defaults."""
        theme = Theme.from_dict({})

        assert theme.name == "custom"
        assert theme.primary_hue == "blue"
        assert theme.mode == "dark"
        assert theme.font_family == "Inter"

    def test_from_dict_partial(self):
        """Test creating theme from partial dict."""
        theme = Theme.from_dict({"name": "my-theme"})

        assert theme.name == "my-theme"
        assert theme.primary_hue == "blue"  # Default


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
