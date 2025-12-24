"""
Tests for themes/design_system.py

Comprehensive tests for design system resolution for >90% coverage.
"""

from unittest.mock import MagicMock
from pptx import Presentation
from pptx.util import Inches, Pt

from chuk_mcp_pptx.themes.design_system import (
    ResolvedDesignSystem,
    extract_template_design_system,
    extract_placeholder_styles,
    resolve_design_system,
    apply_design_system_to_shape,
    _hex_to_rgb,
)


class TestResolvedDesignSystem:
    """Tests for ResolvedDesignSystem dataclass."""

    def test_default_values(self):
        """Test default values for design system."""
        ds = ResolvedDesignSystem()

        # Colors
        assert ds.primary_color == "#4F46E5"
        assert ds.secondary_color == "#818CF8"
        assert ds.background_color == "#FFFFFF"
        assert ds.text_color == "#1F2937"
        assert ds.border_color == "#E5E7EB"

        # Typography
        assert ds.font_family == "Calibri"
        assert ds.font_size == 14
        assert ds.font_bold is False
        assert ds.font_italic is False

        # Spacing
        assert ds.padding == 0.2
        assert ds.margin == 0.1
        assert ds.gap == 0.15

        # Shape
        assert ds.border_radius == 0.1
        assert ds.border_width == 1.0

        # Source tracking
        assert ds.source == "default"
        assert ds.overrides == {}

    def test_custom_values(self):
        """Test creating design system with custom values."""
        ds = ResolvedDesignSystem(
            primary_color="#FF0000",
            secondary_color="#00FF00",
            background_color="#000000",
            text_color="#FFFFFF",
            border_color="#CCCCCC",
            font_family="Arial",
            font_size=18,
            font_bold=True,
            font_italic=True,
            padding=0.3,
            margin=0.2,
            gap=0.25,
            border_radius=0.2,
            border_width=2.0,
            source="custom",
            overrides={"font_size": "explicit"},
        )

        assert ds.primary_color == "#FF0000"
        assert ds.font_family == "Arial"
        assert ds.font_size == 18
        assert ds.font_bold is True
        assert ds.font_italic is True
        assert ds.source == "custom"
        assert ds.overrides == {"font_size": "explicit"}


class TestHexToRgb:
    """Tests for _hex_to_rgb function."""

    def test_hex_with_hash(self):
        """Test hex color with hash prefix."""
        assert _hex_to_rgb("#FF0000") == (255, 0, 0)
        assert _hex_to_rgb("#00FF00") == (0, 255, 0)
        assert _hex_to_rgb("#0000FF") == (0, 0, 255)

    def test_hex_without_hash(self):
        """Test hex color without hash prefix."""
        assert _hex_to_rgb("FF0000") == (255, 0, 0)
        assert _hex_to_rgb("00FF00") == (0, 255, 0)

    def test_lowercase_hex(self):
        """Test lowercase hex colors."""
        assert _hex_to_rgb("#ff0000") == (255, 0, 0)
        assert _hex_to_rgb("aabbcc") == (170, 187, 204)

    def test_black_and_white(self):
        """Test black and white colors."""
        assert _hex_to_rgb("#000000") == (0, 0, 0)
        assert _hex_to_rgb("#FFFFFF") == (255, 255, 255)


class TestExtractTemplateDesignSystem:
    """Tests for extract_template_design_system function."""

    def test_extract_from_basic_slide(self):
        """Test extracting design system from basic slide."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            result = extract_template_design_system(slide)
            # Should return something or None
            assert result is None or isinstance(result, ResolvedDesignSystem)

    def test_extract_handles_exceptions_gracefully(self):
        """Test extraction handles exceptions gracefully."""
        # Create a mock slide that will raise an exception
        mock_slide = MagicMock()
        mock_slide.part.package.presentation_part.presentation.slide_masters = []

        result = extract_template_design_system(mock_slide)
        # Should return None when slide_masters is empty
        assert result is None

    def test_extract_sets_source(self):
        """Test that extraction sets appropriate source."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            result = extract_template_design_system(slide)

            if result:
                assert "template" in result.source.lower()


class TestExtractPlaceholderStyles:
    """Tests for extract_placeholder_styles function."""

    def test_extract_from_none(self):
        """Test extracting styles from None placeholder."""
        styles = extract_placeholder_styles(None)
        assert styles == {}

    def test_extract_from_placeholder_with_text(self):
        """Test extracting styles from placeholder with text."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content

            if slide.placeholders:
                placeholder = list(slide.placeholders)[0]
                styles = extract_placeholder_styles(placeholder)
                # Should return dict (may be empty or have values)
                assert isinstance(styles, dict)

    def test_extract_handles_missing_text_frame(self):
        """Test extraction handles missing text_frame."""
        mock_placeholder = MagicMock()
        mock_placeholder.text_frame = None

        styles = extract_placeholder_styles(mock_placeholder)
        assert isinstance(styles, dict)

    def test_extract_handles_empty_paragraphs(self):
        """Test extraction handles empty paragraphs."""
        mock_placeholder = MagicMock()
        mock_placeholder.text_frame.paragraphs = []
        mock_placeholder.fill = None

        styles = extract_placeholder_styles(mock_placeholder)
        assert isinstance(styles, dict)

    def test_extract_handles_empty_runs(self):
        """Test extraction handles empty runs."""
        mock_placeholder = MagicMock()
        mock_para = MagicMock()
        mock_para.runs = []
        mock_placeholder.text_frame.paragraphs = [mock_para]
        mock_placeholder.fill = None

        styles = extract_placeholder_styles(mock_placeholder)
        assert isinstance(styles, dict)

    def test_extract_with_fill(self):
        """Test extracting fill color."""
        mock_placeholder = MagicMock()
        mock_placeholder.text_frame = None
        mock_placeholder.fill.fore_color.rgb = (255, 0, 0)

        # This should handle the fill extraction
        styles = extract_placeholder_styles(mock_placeholder)
        assert isinstance(styles, dict)


class TestResolveDesignSystem:
    """Tests for resolve_design_system function."""

    def test_resolve_with_no_params(self):
        """Test resolution with no parameters."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            result = resolve_design_system(slide)

            assert isinstance(result, ResolvedDesignSystem)

    def test_resolve_with_placeholder(self):
        """Test resolution with placeholder."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[1])

            if slide.placeholders:
                placeholder = list(slide.placeholders)[0]
                result = resolve_design_system(slide, placeholder=placeholder)

                assert isinstance(result, ResolvedDesignSystem)

    def test_resolve_with_param_overrides(self):
        """Test resolution with parameter overrides."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            params = {
                "bg_color": "#FF0000",
                "text_color": "#FFFFFF",
                "font_size": 20,
            }
            result = resolve_design_system(slide, params=params)

            assert result.background_color == "#FF0000"
            assert result.text_color == "#FFFFFF"
            assert result.font_size == 20
            assert "background_color" in result.overrides
            assert result.overrides["background_color"] == "explicit"

    def test_resolve_with_color_override(self):
        """Test resolution with primary color override."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            params = {"color": "#00FF00"}
            result = resolve_design_system(slide, params=params)

            assert result.primary_color == "#00FF00"

    def test_resolve_with_typography_overrides(self):
        """Test resolution with typography overrides."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            params = {
                "font_family": "Arial",
                "font_bold": True,
                "font_italic": True,
            }
            result = resolve_design_system(slide, params=params)

            assert result.font_family == "Arial"
            assert result.font_bold is True
            assert result.font_italic is True

    def test_resolve_with_spacing_overrides(self):
        """Test resolution with spacing overrides."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            params = {
                "padding": 0.5,
                "margin": 0.3,
                "gap": 0.4,
            }
            result = resolve_design_system(slide, params=params)

            assert result.padding == 0.5
            assert result.margin == 0.3
            assert result.gap == 0.4

    def test_resolve_with_border_overrides(self):
        """Test resolution with border overrides."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            params = {
                "border_color": "#CCCCCC",
                "border_radius": 0.2,
                "border_width": 2.0,
            }
            result = resolve_design_system(slide, params=params)

            assert result.border_color == "#CCCCCC"
            assert result.border_radius == 0.2
            assert result.border_width == 2.0

    def test_resolve_with_theme_string(self):
        """Test resolution with theme name string."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Use a theme name - should attempt to load theme
            result = resolve_design_system(slide, theme="dark-blue")

            assert isinstance(result, ResolvedDesignSystem)
            # Source should indicate theme
            assert "theme" in result.source.lower() or result.source == "default"

    def test_resolve_all_override_mappings(self):
        """Test all override mappings work."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Test all mappings
            params = {
                "bg_color": "#111111",
                "background_color": "#222222",  # Will override bg_color
                "color": "#333333",
                "primary_color": "#444444",  # Will override color
                "text_color": "#555555",
                "border_color": "#666666",
                "font_family": "Georgia",
                "font_size": 16,
                "font_bold": True,
                "font_italic": False,
                "padding": 0.1,
                "margin": 0.15,
                "gap": 0.2,
                "border_radius": 0.05,
                "border_width": 1.5,
            }
            result = resolve_design_system(slide, params=params)

            # Later mappings should override earlier ones
            assert result.background_color == "#222222"
            assert result.primary_color == "#444444"
            assert result.text_color == "#555555"


class TestApplyDesignSystemToShape:
    """Tests for apply_design_system_to_shape function."""

    def test_apply_to_shape(self):
        """Test applying design system to a shape."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

            ds = ResolvedDesignSystem(
                background_color="#FF0000",
                border_color="#00FF00",
                border_width=2.0,
                font_family="Arial",
                font_size=16,
                font_bold=True,
                font_italic=False,
                text_color="#0000FF",
            )

            apply_design_system_to_shape(shape, ds)

            # Verify fill was applied
            assert shape.fill.type is not None

    def test_apply_to_shape_with_text(self):
        """Test applying design system to shape with text."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = "Test Text"

            ds = ResolvedDesignSystem(
                font_family="Arial",
                font_size=18,
                font_bold=True,
                font_italic=True,
                text_color="#FF0000",
            )

            apply_design_system_to_shape(shape, ds)

            # Verify text styling was applied
            run = p.runs[0]
            assert run.font.name == "Arial"
            assert run.font.size == Pt(18)
            assert run.font.bold is True
            assert run.font.italic is True

    def test_apply_handles_shape_without_text_frame(self):
        """Test applying to shape without text_frame."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            from pptx.enum.shapes import MSO_SHAPE

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
            )

            ds = ResolvedDesignSystem()

            # Should not raise
            apply_design_system_to_shape(shape, ds)

    def test_apply_handles_errors_gracefully(self):
        """Test that apply handles errors gracefully."""
        mock_shape = MagicMock()
        mock_shape.fill.solid.side_effect = Exception("Test error")

        ds = ResolvedDesignSystem()

        # Should not raise
        apply_design_system_to_shape(mock_shape, ds)


class TestExtractTemplateDesignSystemAdvanced:
    """Advanced tests for extract_template_design_system."""

    def test_extract_with_theme_part(self):
        """Test extraction when theme_part exists."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        # Setup mock master with theme
        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # Mock theme_part with theme_element
        mock_theme_elem = MagicMock()
        mock_master.part.theme_part.theme_element = mock_theme_elem

        # Mock color scheme with accent colors
        mock_clr_scheme = MagicMock()
        mock_theme_elem.themeElements.clrScheme = mock_clr_scheme

        # Mock accent1 with srgbClr
        mock_clr_scheme.accent1.srgbClr.val = "FF0000"

        # Mock accent2 with srgbClr
        mock_clr_scheme.accent2.srgbClr.val = "00FF00"

        # Mock lt1 (background) with srgbClr
        mock_clr_scheme.lt1.srgbClr.val = "FFFFFF"

        # Mock dk1 (text) with sysClr
        del mock_clr_scheme.dk1.srgbClr  # Remove srgbClr to use sysClr
        mock_clr_scheme.dk1.sysClr.lastClr = "000000"

        # Mock shapes for typography extraction
        mock_shape = MagicMock()
        mock_run = MagicMock()
        mock_run.font.name = "Arial"
        mock_run.font.size.pt = 16
        mock_shape.text_frame.paragraphs = [MagicMock(runs=[mock_run])]
        mock_master.shapes = [mock_shape]

        result = extract_template_design_system(mock_slide)

        assert result is not None
        assert result.primary_color == "#FF0000"
        assert result.secondary_color == "#00FF00"
        assert result.background_color == "#FFFFFF"
        assert result.text_color == "#000000"
        assert result.font_family == "Arial"
        assert result.font_size == 16

    def test_extract_with_no_accent_colors(self):
        """Test extraction when accent colors don't exist."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # No accent colors in scheme
        mock_clr_scheme = MagicMock(spec=[])  # Empty spec means no attributes
        mock_master.part.theme_part.theme_element.themeElements.clrScheme = mock_clr_scheme

        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_with_no_theme_part(self):
        """Test extraction when theme_part doesn't exist."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # No theme_part
        del mock_master.part.theme_part

        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None
        assert result.source == "template"

    def test_extract_handles_theme_extraction_exception(self):
        """Test that exceptions in theme extraction are handled."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # Make theme_part access raise exception
        type(mock_master.part).theme_part = property(
            lambda self: (_ for _ in ()).throw(Exception("Test error"))
        )

        mock_master.shapes = []

        # Should handle exception gracefully
        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_with_no_shapes(self):
        """Test extraction when master has no shapes."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_with_shapes_no_text_frame(self):
        """Test extraction when shapes have no text_frame."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        mock_shape = MagicMock()
        mock_shape.text_frame = None
        mock_master.shapes = [mock_shape]

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_with_shapes_no_paragraphs(self):
        """Test extraction when text_frame has no paragraphs."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        mock_shape = MagicMock()
        mock_shape.text_frame.paragraphs = []
        mock_master.shapes = [mock_shape]

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_with_shapes_no_runs(self):
        """Test extraction when paragraphs have no runs."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        mock_para = MagicMock()
        mock_para.runs = []
        mock_shape = MagicMock()
        mock_shape.text_frame.paragraphs = [mock_para]
        mock_master.shapes = [mock_shape]

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_font_without_name(self):
        """Test extraction when font has no name."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        mock_run = MagicMock()
        mock_run.font.name = None
        mock_run.font.size = None
        mock_para = MagicMock(runs=[mock_run])
        mock_shape = MagicMock()
        mock_shape.text_frame.paragraphs = [mock_para]
        mock_master.shapes = [mock_shape]

        result = extract_template_design_system(mock_slide)
        assert result is not None


class TestExtractPlaceholderStylesAdvanced:
    """Advanced tests for extract_placeholder_styles."""

    def test_extract_with_font_styles(self):
        """Test extracting all font styles."""
        mock_placeholder = MagicMock()
        mock_run = MagicMock()
        mock_run.font.name = "Times New Roman"
        mock_run.font.size.pt = 20
        mock_run.font.bold = True
        mock_run.font.italic = True
        mock_run.font.color.rgb = (255, 128, 64)
        mock_para = MagicMock(runs=[mock_run])
        mock_placeholder.text_frame.paragraphs = [mock_para]
        mock_placeholder.fill = None

        styles = extract_placeholder_styles(mock_placeholder)

        assert styles["font_family"] == "Times New Roman"
        assert styles["font_size"] == 20
        assert styles["font_bold"] is True
        assert styles["font_italic"] is True
        assert styles["text_color"] == "#ff8040"

    def test_extract_with_fill_color(self):
        """Test extracting fill color."""
        mock_placeholder = MagicMock()
        # Need to mock text_frame such that accessing paragraphs doesn't raise
        mock_placeholder.text_frame.paragraphs = []
        mock_placeholder.fill.fore_color.rgb = (200, 100, 50)

        styles = extract_placeholder_styles(mock_placeholder)

        assert styles["background_color"] == "#c86432"

    def test_extract_handles_rgb_extraction_exception(self):
        """Test that RGB extraction exceptions are handled."""
        mock_placeholder = MagicMock()
        mock_run = MagicMock()
        mock_run.font.name = "Arial"
        mock_run.font.size = None
        mock_run.font.bold = None
        mock_run.font.italic = None
        # Make color.rgb raise exception
        type(mock_run.font.color).rgb = property(
            lambda self: (_ for _ in ()).throw(Exception("No RGB"))
        )
        mock_para = MagicMock(runs=[mock_run])
        mock_placeholder.text_frame.paragraphs = [mock_para]
        mock_placeholder.fill = None

        styles = extract_placeholder_styles(mock_placeholder)
        assert "font_family" in styles
        assert "text_color" not in styles

    def test_extract_handles_fill_exception(self):
        """Test that fill extraction exceptions are handled."""
        mock_placeholder = MagicMock()
        mock_placeholder.text_frame = None
        # Make fore_color.rgb raise exception
        type(mock_placeholder.fill.fore_color).rgb = property(
            lambda self: (_ for _ in ()).throw(Exception("No fill"))
        )

        styles = extract_placeholder_styles(mock_placeholder)
        assert "background_color" not in styles


class TestResolveDesignSystemAdvanced:
    """Advanced tests for resolve_design_system."""

    def test_resolve_with_theme_object(self):
        """Test resolution with theme object (not string)."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Create mock theme object
            mock_theme = MagicMock()
            mock_theme.colors = {
                "primary": "#123456",
                "secondary": "#789ABC",
                "background": "#FFFFFF",
                "text": "#000000",
            }
            mock_theme.typography = {
                "font_family": "Georgia",
                "font_size": 18,
            }

            result = resolve_design_system(slide, theme=mock_theme)

            assert result.primary_color == "#123456"
            assert result.secondary_color == "#789ABC"
            assert result.font_family == "Georgia"
            assert result.font_size == 18

    def test_resolve_with_none_theme_from_manager(self):
        """Test resolution when theme manager returns None."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            result = resolve_design_system(slide, theme="nonexistent-theme")

            # Should fall back to defaults
            assert isinstance(result, ResolvedDesignSystem)


class TestIntegration:
    """Integration tests for design system."""

    def test_full_workflow(self):
        """Test full workflow of resolving and applying design system."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.text = "Test"

            # Resolve design system with overrides
            ds = resolve_design_system(
                slide,
                params={
                    "bg_color": "#EEEEEE",
                    "text_color": "#333333",
                    "font_size": 14,
                },
            )

            # Apply to shape
            apply_design_system_to_shape(shape, ds)

            # Verify results
            assert ds.background_color == "#EEEEEE"
            assert ds.text_color == "#333333"
            assert ds.font_size == 14

    def test_design_system_priority_hierarchy(self):
        """Test that priority hierarchy is respected."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Explicit param should override defaults
            ds = resolve_design_system(
                slide,
                params={"font_size": 24},
            )

            assert ds.font_size == 24
            assert "font_size" in ds.overrides
            assert ds.overrides["font_size"] == "explicit"


class TestExtractColorFromScheme:
    """Tests for color scheme extraction branches."""

    def test_extract_color_returns_none_for_unknown_type(self):
        """Test that extract_color_from_scheme returns None for unknown color types."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # Create color scheme with accent1 that has neither srgbClr nor sysClr
        mock_clr_scheme = MagicMock()
        mock_accent = MagicMock(spec=[])  # Empty spec - no srgbClr or sysClr
        mock_clr_scheme.accent1 = mock_accent
        mock_clr_scheme.accent2 = mock_accent
        mock_clr_scheme.lt1 = mock_accent
        mock_clr_scheme.dk1 = mock_accent

        mock_master.part.theme_part.theme_element.themeElements.clrScheme = mock_clr_scheme
        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        # Should still return a design system with default colors
        assert result is not None

    def test_extract_no_theme_elements_attribute(self):
        """Test extraction when themeElements attribute doesn't exist."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # theme_element exists but has no themeElements
        mock_theme_elem = MagicMock(spec=[])  # No themeElements
        mock_master.part.theme_part.theme_element = mock_theme_elem
        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_no_clr_scheme_attribute(self):
        """Test extraction when clrScheme attribute doesn't exist."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # themeElements exists but has no clrScheme
        mock_theme_elements = MagicMock(spec=[])  # No clrScheme
        mock_master.part.theme_part.theme_element.themeElements = mock_theme_elements
        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_no_accent1_attribute(self):
        """Test extraction when accent1 doesn't exist in color scheme."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # clrScheme exists but has no accent1, accent2, lt1, dk1
        mock_clr_scheme = MagicMock(spec=[])
        mock_master.part.theme_part.theme_element.themeElements.clrScheme = mock_clr_scheme
        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None

    def test_extract_with_partial_accent_colors(self):
        """Test extraction with only some accent colors present."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock()
        mock_prs.slide_masters = [mock_master]

        # clrScheme with only accent1
        mock_clr_scheme = MagicMock(spec=["accent1"])  # Only has accent1
        mock_clr_scheme.accent1.srgbClr.val = "FF0000"
        mock_master.part.theme_part.theme_element.themeElements.clrScheme = mock_clr_scheme
        mock_master.shapes = []

        result = extract_template_design_system(mock_slide)
        assert result is not None
        assert result.primary_color == "#FF0000"

    def test_extract_without_slide_master_or_part(self):
        """Test extraction when master doesn't have slide_master or part attribute."""
        mock_slide = MagicMock()
        mock_prs = MagicMock()
        mock_slide.part.package.presentation_part.presentation = mock_prs

        mock_master = MagicMock(spec=[])  # No slide_master or part
        mock_prs.slide_masters = [mock_master]

        result = extract_template_design_system(mock_slide)
        # Should return template design system even without theme extraction
        assert result is not None

    def test_extract_main_exception_handler(self):
        """Test that main exception handler catches errors."""
        mock_slide = MagicMock()
        # Make the initial access raise an exception
        type(mock_slide).part = property(
            lambda self: (_ for _ in ()).throw(Exception("Access error"))
        )

        result = extract_template_design_system(mock_slide)
        assert result is None


class TestExtractPlaceholderStylesExceptions:
    """Tests for placeholder style extraction exception handling."""

    def test_extract_text_color_exception(self):
        """Test extraction handles text color RGB access exception."""
        mock_placeholder = MagicMock()
        mock_run = MagicMock()
        mock_run.font.name = "Arial"
        mock_run.font.size.pt = 14
        mock_run.font.bold = False
        mock_run.font.italic = False

        # Make color.rgb property raise exception
        mock_color = MagicMock()
        mock_color.rgb = None  # Will fail indexing
        mock_run.font.color = mock_color

        mock_para = MagicMock(runs=[mock_run])
        mock_placeholder.text_frame.paragraphs = [mock_para]
        mock_placeholder.fill = None

        styles = extract_placeholder_styles(mock_placeholder)
        assert "font_family" in styles
        # text_color should not be present due to exception

    def test_extract_fill_color_exception(self):
        """Test extraction handles fill color RGB access exception."""
        mock_placeholder = MagicMock()
        mock_placeholder.text_frame.paragraphs = []

        # Make fore_color.rgb access fail
        mock_fill = MagicMock()
        mock_fill.fore_color.rgb = None  # Will fail indexing
        mock_placeholder.fill = mock_fill

        extract_placeholder_styles(mock_placeholder)
        # background_color should not be present due to exception

    def test_extract_main_exception_handler(self):
        """Test that main exception handler catches errors."""
        mock_placeholder = MagicMock()
        # Make text_frame access raise exception
        type(mock_placeholder).text_frame = property(
            lambda self: (_ for _ in ()).throw(Exception("Access error"))
        )

        styles = extract_placeholder_styles(mock_placeholder)
        assert styles == {}


class TestResolveDesignSystemPlaceholder:
    """Tests for resolve_design_system with placeholder styles."""

    def test_resolve_applies_placeholder_styles(self):
        """Test that placeholder styles are applied."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Create mock placeholder with styles
            mock_placeholder = MagicMock()
            mock_run = MagicMock()
            mock_run.font.name = "Georgia"
            mock_run.font.size.pt = 20
            mock_run.font.bold = True
            mock_run.font.italic = False
            mock_run.font.color.rgb = (100, 100, 100)

            mock_para = MagicMock(runs=[mock_run])
            mock_placeholder.text_frame.paragraphs = [mock_para]
            mock_placeholder.fill = None

            result = resolve_design_system(slide, placeholder=mock_placeholder)

            assert result.font_family == "Georgia"
            assert result.font_size == 20
            assert result.font_bold is True
            assert "font_family" in result.overrides
            assert result.overrides["font_family"] == "placeholder"

    def test_resolve_placeholder_without_styles(self):
        """Test resolution with placeholder that has no extractable styles."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Empty placeholder
            mock_placeholder = MagicMock()
            mock_placeholder.text_frame.paragraphs = []
            mock_placeholder.fill = None

            result = resolve_design_system(slide, placeholder=mock_placeholder)
            # Should use defaults since no styles were extracted
            assert isinstance(result, ResolvedDesignSystem)


class TestResolveDesignSystemTheme:
    """Tests for resolve_design_system with theme."""

    def test_resolve_with_theme_missing_colors(self):
        """Test resolution with theme that has no colors attribute."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            mock_theme = MagicMock(spec=["typography"])  # No colors
            mock_theme.typography = {"font_family": "Arial", "font_size": 16}

            result = resolve_design_system(slide, theme=mock_theme)
            assert result.font_family == "Arial"

    def test_resolve_with_theme_missing_typography(self):
        """Test resolution with theme that has no typography attribute."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            mock_theme = MagicMock(spec=["colors"])  # No typography
            mock_theme.colors = {"primary": "#FF0000", "secondary": "#00FF00"}

            result = resolve_design_system(slide, theme=mock_theme)
            assert result.primary_color == "#FF0000"

    def test_resolve_with_theme_partial_colors(self):
        """Test resolution with theme that has partial color values."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            mock_theme = MagicMock()
            mock_theme.colors = {"primary": "#123456"}  # Only primary
            mock_theme.typography = {}

            result = resolve_design_system(slide, theme=mock_theme)
            assert result.primary_color == "#123456"
            # Other colors should remain default

    def test_resolve_with_theme_partial_typography(self):
        """Test resolution with theme that has partial typography values."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            mock_theme = MagicMock()
            mock_theme.colors = {}
            mock_theme.typography = {"font_family": "Georgia"}  # Only font_family

            result = resolve_design_system(slide, theme=mock_theme)
            assert result.font_family == "Georgia"


class TestApplyDesignSystemBranches:
    """Tests for apply_design_system_to_shape branches."""

    def test_apply_to_shape_without_fill(self):
        """Test applying to shape without fill attribute."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

            mock_shape = MagicMock(spec=["line", "text_frame"])  # No fill
            mock_shape.line = shape.line
            mock_shape.text_frame = None

            ds = ResolvedDesignSystem()
            apply_design_system_to_shape(mock_shape, ds)
            # Should not raise

    def test_apply_to_shape_without_line(self):
        """Test applying to shape without line attribute."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

            mock_shape = MagicMock(spec=["fill", "text_frame"])  # No line
            mock_shape.fill = shape.fill
            mock_shape.text_frame = None

            ds = ResolvedDesignSystem()
            apply_design_system_to_shape(mock_shape, ds)
            # Should not raise

    def test_apply_to_shape_with_empty_text_frame(self):
        """Test applying to shape with empty text_frame."""
        prs = Presentation()
        if prs.slide_layouts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))

            # Clear paragraphs
            shape.text_frame.paragraphs[0].clear()

            ds = ResolvedDesignSystem()
            apply_design_system_to_shape(shape, ds)
            # Should not raise
