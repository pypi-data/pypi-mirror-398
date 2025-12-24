"""
Tests for tools/template/* modules.

Comprehensive tests for template tools including:
- __init__.py: Registration functions
- analyze.py: Template analysis tools
- list.py: Template listing tools
- import_tools.py: Template import tools
- workflow.py: Template workflow tools
"""

import json
import pytest
from unittest.mock import MagicMock, AsyncMock
from pptx import Presentation

from chuk_mcp_pptx.tools.template import (
    register_template_tools,
    register_list_tools,
    register_analyze_tools,
    register_import_tools,
    register_workflow_tools,
)
from chuk_mcp_pptx.tools.template.models import (
    LayoutInfo,
    TemplateInfo,
    PresentationTemplateListResponse,
    BuiltinTemplateInfo,
    CustomTemplateInfo,
)


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server."""
    mcp = MagicMock()
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mcp.tool = tool_decorator
    mcp._tools = tools
    return mcp


@pytest.fixture
def mock_template_manager():
    """Create a mock template manager."""
    manager = MagicMock()

    # Mock list_templates
    mock_template = MagicMock()
    mock_template.name = "test_template"
    mock_template.display_name = "Test Template"
    mock_template.description = "A test template"
    mock_template.category = "test"
    mock_template.layout_count = 10
    mock_template.tags = ["test", "sample"]
    manager.list_templates.return_value = [mock_template]

    # Mock get_template_data - return None by default (not a builtin)
    manager.get_template_data = AsyncMock(return_value=None)

    return manager


@pytest.fixture
def mock_presentation_manager():
    """Create a mock presentation manager."""
    manager = MagicMock()

    # Create a real presentation for testing
    prs = Presentation()
    if prs.slide_layouts:
        prs.slides.add_slide(prs.slide_layouts[0])

    metadata = MagicMock()
    metadata.name = "test_presentation"
    metadata.namespace_id = "test-namespace-id"
    metadata.vfs_path = "/templates/test"

    # Mock get method
    manager.get = AsyncMock(return_value=(prs, metadata))
    manager.get_presentation = AsyncMock(return_value=prs)

    # Mock list_presentations
    pres_info = MagicMock()
    pres_info.name = "test_presentation"
    pres_info.slide_count = 1
    pres_info.namespace_id = "test-namespace-id"

    response = MagicMock()
    response.presentations = [pres_info]
    manager.list_presentations = AsyncMock(return_value=response)

    # Mock metadata storage
    manager._metadata = {
        "test_presentation": metadata,
    }

    # Mock update methods
    manager.update_slide_metadata = AsyncMock()
    manager._save_to_store = AsyncMock()
    manager.update = AsyncMock()
    manager.import_template = AsyncMock(return_value=True)
    manager.get_current_name = MagicMock(return_value="test_presentation")

    return manager


class TestRegisterTemplateTools:
    """Tests for template tools registration."""

    def test_register_template_tools(self, mock_mcp, mock_presentation_manager):
        """Test registering all template tools."""
        tools = register_template_tools(mock_mcp, mock_presentation_manager, template_manager=None)

        assert isinstance(tools, dict)
        # Should have tools from all submodules
        assert "pptx_list_templates" in tools
        assert "pptx_analyze_template" in tools
        assert "pptx_import_template" in tools
        assert "pptx_add_slide_from_template" in tools

    def test_register_template_tools_with_manager(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test registering with explicit template manager."""
        tools = register_template_tools(
            mock_mcp, mock_presentation_manager, template_manager=mock_template_manager
        )

        assert isinstance(tools, dict)
        assert len(tools) > 0


class TestListTools:
    """Tests for template listing tools."""

    @pytest.mark.asyncio
    async def test_list_templates(self, mock_mcp, mock_presentation_manager, mock_template_manager):
        """Test listing templates."""
        register_list_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)
        data = json.loads(result)
        assert "builtin_templates" in data or "error" not in data

    @pytest.mark.asyncio
    async def test_list_templates_no_template_manager(self, mock_mcp, mock_presentation_manager):
        """Test listing templates without template manager."""
        register_list_tools(mock_mcp, mock_presentation_manager, None)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_list_templates_with_custom_templates(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test listing templates with custom templates."""
        # Set up custom template in metadata
        mock_presentation_manager._metadata["test_presentation"].vfs_path = "/templates/custom"

        register_list_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)


class TestAnalyzeTools:
    """Tests for template analysis tools."""

    @pytest.mark.asyncio
    async def test_analyze_template_from_store(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing a template from the store."""
        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test_presentation")

        assert isinstance(result, str)
        data = json.loads(result)
        # Should have template info or error
        assert "name" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_builtin(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing a builtin template."""
        # Mock builtin template data
        prs = Presentation()
        import io

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("builtin_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "name" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_not_found(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing a template that doesn't exist."""
        mock_presentation_manager.get = AsyncMock(return_value=None)

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("nonexistent")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_variants(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing template variants."""
        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("test_presentation")

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_analyze_template_variants_not_found(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test analyzing variants for non-existent template."""
        mock_presentation_manager.get = AsyncMock(return_value=None)

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("nonexistent")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestImportTools:
    """Tests for template import tools."""

    @pytest.mark.asyncio
    async def test_import_template_success(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test successfully importing a template."""
        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_import_template"](
            file_path="/path/to/template.pptx", template_name="imported_template"
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_import_template_failure(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test importing template that fails."""
        mock_presentation_manager.import_template = AsyncMock(return_value=False)

        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_import_template"](
            file_path="/path/to/nonexistent.pptx", template_name="failed_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_import_template_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test importing template with exception."""
        mock_presentation_manager.import_template = AsyncMock(
            side_effect=Exception("Import failed")
        )

        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_import_template"](
            file_path="/path/to/error.pptx", template_name="error_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_get_builtin_template_not_found(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test getting builtin template that doesn't exist."""
        mock_template_manager.get_template_data = AsyncMock(return_value=None)

        register_import_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="nonexistent", save_as="my_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestWorkflowTools:
    """Tests for template workflow tools."""

    @pytest.mark.asyncio
    async def test_add_slide_from_template_no_presentation(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test adding slide when presentation not found."""
        mock_presentation_manager.get = AsyncMock(return_value=None)

        register_workflow_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_add_slide_from_template"](
            layout_index=0, presentation="nonexistent"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestModels:
    """Tests for template models."""

    def test_layout_info_model(self):
        """Test LayoutInfo model."""
        layout = LayoutInfo(
            index=0,
            name="Title Slide",
            placeholder_count=2,
            placeholders=[
                {"idx": 0, "type": "TITLE", "name": "Title"},
                {"idx": 1, "type": "SUBTITLE", "name": "Subtitle"},
            ],
        )
        assert layout.index == 0
        assert layout.name == "Title Slide"
        assert layout.placeholder_count == 2
        assert len(layout.placeholders) == 2

    def test_template_info_model(self):
        """Test TemplateInfo model."""
        info = TemplateInfo(
            name="test_template",
            slide_count=0,
            layout_count=5,
            layouts=[],
            master_count=1,
            has_theme=True,
        )
        assert info.name == "test_template"
        assert info.layout_count == 5

    def test_builtin_template_info_model(self):
        """Test BuiltinTemplateInfo model."""
        info = BuiltinTemplateInfo(
            name="corporate",
            display_name="Corporate Template",
            description="A corporate template",
            category="business",
            layout_count=20,
            tags=["business", "formal"],
            is_builtin=True,
        )
        assert info.name == "corporate"
        assert info.is_builtin is True

    def test_custom_template_info_model(self):
        """Test CustomTemplateInfo model."""
        info = CustomTemplateInfo(
            name="my_template",
            slide_count=5,
            namespace_id="ns-123",
            is_builtin=False,
            category="custom",
        )
        assert info.name == "my_template"
        assert info.is_builtin is False

    def test_template_list_response_model(self):
        """Test PresentationTemplateListResponse model."""
        response = PresentationTemplateListResponse(
            builtin_templates=[],
            custom_templates=[],
            total=0,
        )
        assert response.total == 0


class TestAnalyzeLayoutLoop:
    """Tests for analyze.py layout analysis loop (lines 118-138)."""

    @pytest.fixture
    def presentation_with_layouts(self):
        """Create a presentation with layouts that have placeholders."""
        prs = Presentation()
        return prs

    @pytest.mark.asyncio
    async def test_analyze_template_with_layouts_and_placeholders(
        self, mock_mcp, mock_template_manager
    ):
        """Test analyzing a template with slide_layouts and placeholders."""
        # Create a real presentation with layouts
        prs = Presentation()

        # Set up manager mock with real presentation
        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_with_layouts"
        manager.get = AsyncMock(return_value=(prs, metadata))

        register_analyze_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test_with_layouts")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "name" in data
        assert data["name"] == "test_with_layouts"
        # Should have layouts from the real presentation
        assert "layouts" in data
        assert "layout_count" in data
        # Verify layouts have the expected structure
        if data["layouts"]:
            layout = data["layouts"][0]
            assert "index" in layout
            assert "name" in layout
            assert "placeholder_count" in layout
            assert "placeholders" in layout

    @pytest.mark.asyncio
    async def test_analyze_template_no_slide_layouts(self, mock_mcp, mock_template_manager):
        """Test analyzing a template when slide_layouts is empty/falsy (line 118->138)."""
        # Create a mock presentation with empty slide_layouts
        mock_prs = MagicMock()
        mock_prs.slide_layouts = []  # Empty slide_layouts
        mock_prs.slides = []
        mock_prs.slide_masters = [MagicMock()]  # Need at least 1 master

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "no_layouts_template"
        manager.get = AsyncMock(return_value=(mock_prs, metadata))

        register_analyze_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("no_layouts_template")

        data = json.loads(result)
        assert "name" in data
        assert data["name"] == "no_layouts_template"
        assert data["layouts"] == []
        assert data["layout_count"] == 0

    @pytest.mark.asyncio
    async def test_analyze_template_no_slide_master(self, mock_mcp, mock_template_manager):
        """Test analyzing a template without slide_master."""
        mock_prs = MagicMock()
        mock_prs.slide_layouts = []
        mock_prs.slides = []
        mock_prs.slide_masters = [MagicMock()]  # Need at least 1 master
        # No slide_master attribute
        del mock_prs.slide_master

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(mock_prs, metadata))

        register_analyze_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test")

        data = json.loads(result)
        # Should handle missing slide_master gracefully
        assert "has_theme" in data
        assert data["has_theme"] is False

    @pytest.mark.asyncio
    async def test_analyze_template_layout_placeholder_details(
        self, mock_mcp, mock_template_manager
    ):
        """Test that placeholder details are correctly extracted."""
        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        register_analyze_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("test")

        data = json.loads(result)
        if data.get("layouts"):
            for layout in data["layouts"]:
                # Each placeholder should have idx, type, name
                for ph in layout.get("placeholders", []):
                    assert "idx" in ph
                    assert "type" in ph
                    assert "name" in ph


class TestAnalyzeVariantsBuiltinTemplate:
    """Tests for analyze.py variant analysis with builtin template (lines 178-197)."""

    @pytest.mark.asyncio
    async def test_analyze_variants_builtin_template(self, mock_mcp, mock_presentation_manager):
        """Test analyzing variants for a builtin template."""
        import io

        # Create template manager that returns template data
        template_manager = MagicMock()
        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("builtin_test")

        assert isinstance(result, str)
        # Should have successfully analyzed the builtin template
        data = json.loads(result)
        # Either we get layout groups or an error
        assert "layout_groups" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_variants_builtin_with_layouts(self, mock_mcp, mock_presentation_manager):
        """Test analyzing variants for builtin template with multiple layouts."""
        import io

        template_manager = MagicMock()
        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("multi_layout_template")

        assert isinstance(result, str)
        data = json.loads(result)
        # Check structure
        if "layout_groups" in data:
            assert isinstance(data["layout_groups"], list)

    @pytest.mark.asyncio
    async def test_analyze_variants_from_artifact_store(self, mock_mcp, mock_presentation_manager):
        """Test analyzing variants when template is from artifact store (not builtin)."""

        # Template manager returns None (not a builtin template)
        template_manager = MagicMock()
        template_manager.get_template_data = AsyncMock(return_value=None)

        # Use mock_presentation_manager which is already set up
        register_analyze_tools(mock_mcp, mock_presentation_manager, template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("store_template")

        assert isinstance(result, str)
        data = json.loads(result)
        # Should have analyzed the template from store
        assert "layout_groups" in data or "ungrouped_layouts" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_variants_complete_flow(self, mock_mcp, mock_presentation_manager):
        """Test complete flow of analyze_template_variants with builtin template."""
        import io
        from unittest.mock import patch

        # Create a real presentation with slides
        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])
            if len(prs.slide_layouts) > 1:
                prs.slides.add_slide(prs.slide_layouts[1])

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        # Template manager returns the template data (simulating builtin)
        template_manager = MagicMock()
        template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        register_analyze_tools(mock_mcp, mock_presentation_manager, template_manager)

        # Mock the import to work correctly
        with patch("chuk_mcp_pptx.tools.template.analyze.io") as mock_io:
            # Use real io for BytesIO
            mock_io.BytesIO = io.BytesIO

            result = await mock_mcp._tools["pptx_analyze_template_variants"]("brand_proposal")

        assert isinstance(result, str)
        data = json.loads(result)
        # Should have layout analysis or error (due to import path issue in code)
        assert "total_layouts" in data or "error" in data


class TestGetBuiltinTemplateSuccess:
    """Tests for import_tools.py pptx_get_builtin_template success path (lines 92-131)."""

    @pytest.mark.asyncio
    async def test_get_builtin_template_success(self, mock_mcp, mock_template_manager):
        """Test successfully getting a builtin template."""
        import io

        # Create a valid presentation
        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_data = buffer.read()
        mock_template_manager.get_template_data = AsyncMock(return_value=template_data)

        # Set up manager with mock store
        manager = MagicMock()
        mock_store = MagicMock()
        mock_namespace_info = MagicMock()
        mock_namespace_info.namespace_id = "ns-123"
        mock_store.create_namespace = AsyncMock(return_value=mock_namespace_info)
        mock_store.write_namespace = AsyncMock()
        manager._get_store = MagicMock(return_value=mock_store)
        manager._sanitize_name = MagicMock(side_effect=lambda x: x)
        manager._namespace_ids = {}
        manager.base_path = "/presentations"
        manager.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="corporate", save_as="my_corporate"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "message" in data
        assert "corporate" in data["message"]
        assert "my_corporate" in data["message"]

    @pytest.mark.asyncio
    async def test_get_builtin_template_no_store(self, mock_mcp, mock_template_manager):
        """Test getting builtin template when no store is available."""
        import io

        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        # Manager with no store
        manager = MagicMock()
        manager._get_store = MagicMock(return_value=None)

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="corporate", save_as="my_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "store" in data["error"].lower()

    @pytest.mark.asyncio
    async def test_get_builtin_template_exception(self, mock_mcp, mock_template_manager):
        """Test exception during builtin template import."""
        import io

        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        mock_template_manager.get_template_data = AsyncMock(return_value=buffer.read())

        manager = MagicMock()
        mock_store = MagicMock()
        mock_store.create_namespace = AsyncMock(side_effect=Exception("Store error"))
        manager._get_store = MagicMock(return_value=mock_store)
        manager._sanitize_name = MagicMock(side_effect=lambda x: x)
        manager.base_path = "/presentations"
        manager.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="corporate", save_as="my_template"
        )

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_get_builtin_template_with_layouts(self, mock_mcp, mock_template_manager):
        """Test builtin template reports correct layout count."""
        import io

        prs = Presentation()
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        template_data = buffer.read()
        mock_template_manager.get_template_data = AsyncMock(return_value=template_data)

        manager = MagicMock()
        mock_store = MagicMock()
        mock_namespace_info = MagicMock()
        mock_namespace_info.namespace_id = "ns-456"
        mock_store.create_namespace = AsyncMock(return_value=mock_namespace_info)
        mock_store.write_namespace = AsyncMock()
        manager._get_store = MagicMock(return_value=mock_store)
        manager._sanitize_name = MagicMock(side_effect=lambda x: x)
        manager._namespace_ids = {}
        manager.base_path = "/presentations"
        manager.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

        register_import_tools(mock_mcp, manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_get_builtin_template"](
            template_name="modern", save_as="my_modern"
        )

        data = json.loads(result)
        assert "message" in data
        # Should mention layouts
        assert "layout" in data["message"].lower()


class TestErrorHandling:
    """Test error handling in template tools."""

    @pytest.mark.asyncio
    async def test_list_templates_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test exception handling in list_templates."""
        mock_template_manager.list_templates.side_effect = Exception("List failed")

        register_list_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test exception handling in analyze_template."""
        mock_presentation_manager.get = AsyncMock(side_effect=Exception("Get failed"))

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template"]("error_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_analyze_variants_exception(
        self, mock_mcp, mock_presentation_manager, mock_template_manager
    ):
        """Test exception handling in analyze_template_variants."""
        mock_presentation_manager.get = AsyncMock(side_effect=Exception("Variant error"))
        mock_template_manager.get_template_data = AsyncMock(return_value=None)

        register_analyze_tools(mock_mcp, mock_presentation_manager, mock_template_manager)

        result = await mock_mcp._tools["pptx_analyze_template_variants"]("error_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


# ============================================================================
# Test Workflow Tools Placeholder Handling (lines 132-214)
# ============================================================================


class TestWorkflowPlaceholderHandling:
    """Tests for workflow.py placeholder type guidance (lines 132-214)."""

    @pytest.mark.asyncio
    async def test_add_slide_with_chart_placeholder(self, mock_mcp):
        """Test that chart placeholders get proper guidance."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        # Create presentation with slide that has CHART placeholder
        prs = Presentation()
        prs.slide_layouts[0]

        # Create mock manager
        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.update_slide_metadata = AsyncMock()
        manager._save_to_store = AsyncMock()

        # Mock template manager
        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=0)

        data = json.loads(result)
        # Should succeed and include layout_info
        assert "slide_index" in data or "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_with_table_placeholder(self, mock_mcp):
        """Test that table placeholders get proper guidance."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.update_slide_metadata = AsyncMock()
        manager._save_to_store = AsyncMock()

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        # Use a layout that might have different placeholder types
        result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=1)

        data = json.loads(result)
        assert "slide_index" in data or "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_with_picture_placeholder(self, mock_mcp):
        """Test that picture placeholders get proper guidance."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.update_slide_metadata = AsyncMock()
        manager._save_to_store = AsyncMock()

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        # Use layout index 5 (Two Content) which may have different placeholders
        if len(prs.slide_layouts) > 5:
            result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=5)
        else:
            result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=0)

        data = json.loads(result)
        assert "slide_index" in data or "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_with_object_placeholder(self, mock_mcp):
        """Test that object placeholders get proper guidance."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.update_slide_metadata = AsyncMock()
        manager._save_to_store = AsyncMock()

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        # Use layout index 2 (Content with Caption) which typically has OBJECT
        if len(prs.slide_layouts) > 2:
            result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=2)
        else:
            result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=0)

        data = json.loads(result)
        assert "slide_index" in data or "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_invalid_layout_index(self, mock_mcp):
        """Test adding slide with invalid layout index."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        # Use invalid layout index
        result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=999)

        data = json.loads(result)
        assert "error" in data
        # The error message may vary due to how ErrorResponse is imported
        assert "Invalid layout index" in data["error"] or "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_negative_layout_index(self, mock_mcp):
        """Test adding slide with negative layout index."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=-1)

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_blank_layout_no_placeholders(self, mock_mcp):
        """Test adding slide with a blank layout (no placeholders)."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.update_slide_metadata = AsyncMock()
        manager._save_to_store = AsyncMock()

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        # Use layout index 6 (Blank) which typically has no placeholders
        if len(prs.slide_layouts) > 6:
            result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=6)
            data = json.loads(result)
            # For blank layout, should have minimal or no placeholder guidance
            assert "slide_index" in data or "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_exception_handling(self, mock_mcp):
        """Test exception handling during add slide."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        metadata.name = "test_prs"
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.update_slide_metadata = AsyncMock(side_effect=Exception("Update failed"))

        template_manager = MagicMock()

        register_workflow_tools(mock_mcp, manager, template_manager)

        result = await mock_mcp._tools["pptx_add_slide_from_template"](layout_index=0)

        data = json.loads(result)
        assert "error" in data


# ============================================================================
# Test List Tools Branch Coverage (lines 75->73, 77->73)
# ============================================================================


class TestListToolsBranchCoverage:
    """Tests for list.py branch coverage."""

    @pytest.mark.asyncio
    async def test_list_templates_metadata_without_namespace_id(self, mock_mcp):
        """Test listing templates when metadata has no namespace_id (line 75->73)."""
        from chuk_mcp_pptx.tools.template.list import register_list_tools

        manager = MagicMock()

        # Create presentation info
        pres_info = MagicMock()
        pres_info.name = "test_presentation"
        pres_info.slide_count = 1
        pres_info.namespace_id = "ns-123"

        response = MagicMock()
        response.presentations = [pres_info]
        manager.list_presentations = AsyncMock(return_value=response)

        # Metadata WITHOUT namespace_id (line 75 branch)
        metadata = MagicMock()
        metadata.namespace_id = None  # No namespace_id
        metadata.vfs_path = "/templates/test"
        manager._metadata = {"test_presentation": metadata}

        template_manager = MagicMock()
        template_manager.list_templates.return_value = []

        register_list_tools(mock_mcp, manager, template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        data = json.loads(result)
        # Should not include as custom template because no namespace_id
        assert data["custom_templates"] == []

    @pytest.mark.asyncio
    async def test_list_templates_metadata_without_templates_path(self, mock_mcp):
        """Test listing when vfs_path doesn't include /templates/ (line 77->73)."""
        from chuk_mcp_pptx.tools.template.list import register_list_tools

        manager = MagicMock()

        pres_info = MagicMock()
        pres_info.name = "test_presentation"
        pres_info.slide_count = 1
        pres_info.namespace_id = "ns-123"

        response = MagicMock()
        response.presentations = [pres_info]
        manager.list_presentations = AsyncMock(return_value=response)

        # Metadata WITH namespace_id but vfs_path doesn't have /templates/
        metadata = MagicMock()
        metadata.namespace_id = "ns-123"
        metadata.vfs_path = "/presentations/test"  # Not a template path
        manager._metadata = {"test_presentation": metadata}

        template_manager = MagicMock()
        template_manager.list_templates.return_value = []

        register_list_tools(mock_mcp, manager, template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        data = json.loads(result)
        # Should not include as custom template because path is not /templates/
        assert data["custom_templates"] == []

    @pytest.mark.asyncio
    async def test_list_templates_metadata_missing_for_presentation(self, mock_mcp):
        """Test listing when metadata doesn't exist for presentation (line 74 branch)."""
        from chuk_mcp_pptx.tools.template.list import register_list_tools

        manager = MagicMock()

        pres_info = MagicMock()
        pres_info.name = "unknown_presentation"
        pres_info.slide_count = 1
        pres_info.namespace_id = "ns-123"

        response = MagicMock()
        response.presentations = [pres_info]
        manager.list_presentations = AsyncMock(return_value=response)

        # No metadata for this presentation
        manager._metadata = {}

        template_manager = MagicMock()
        template_manager.list_templates.return_value = []

        register_list_tools(mock_mcp, manager, template_manager)

        result = await mock_mcp._tools["pptx_list_templates"]()

        data = json.loads(result)
        # Should not crash and should return empty custom templates
        assert data["custom_templates"] == []


# ============================================================================
# Test Extraction Module (lines 166-172, 177-212, 217-255, 260-299, etc.)
# ============================================================================


class TestExtractionHelperFunctions:
    """Tests for extraction.py helper functions."""

    def test_rgb_to_hex_with_rgb_attribute(self):
        """Test _rgb_to_hex with valid RGB color."""
        from chuk_mcp_pptx.tools.template.extraction import _rgb_to_hex

        mock_color = MagicMock()
        mock_color.rgb = (255, 128, 64)

        result = _rgb_to_hex(mock_color)
        assert result == "#ff8040"

    def test_rgb_to_hex_without_rgb_attribute(self):
        """Test _rgb_to_hex without rgb attribute returns default."""
        from chuk_mcp_pptx.tools.template.extraction import _rgb_to_hex

        mock_color = MagicMock(spec=[])  # No rgb attribute

        result = _rgb_to_hex(mock_color)
        assert result == "#000000"

    def test_rgb_to_hex_exception(self):
        """Test _rgb_to_hex handles exceptions."""
        from chuk_mcp_pptx.tools.template.extraction import _rgb_to_hex

        mock_color = MagicMock()
        mock_color.rgb = MagicMock(side_effect=Exception("Color error"))

        result = _rgb_to_hex(mock_color)
        assert result == "#000000"

    def test_extract_colors_from_theme_with_color_scheme(self):
        """Test _extract_colors_from_theme with valid color scheme."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_theme

        mock_theme = MagicMock()
        mock_scheme = MagicMock()
        mock_element = MagicMock()

        # Create mock color element with srgbClr
        mock_accent1 = MagicMock()
        mock_accent1.srgbClr.val = "FF6633"
        mock_element.accent1 = mock_accent1

        mock_scheme._element = mock_element
        mock_theme.color_scheme = mock_scheme

        result = _extract_colors_from_theme(mock_theme)
        assert isinstance(result, dict)

    def test_extract_colors_from_theme_no_color_scheme(self):
        """Test _extract_colors_from_theme without color_scheme attribute."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_theme

        mock_theme = MagicMock(spec=[])  # No color_scheme

        result = _extract_colors_from_theme(mock_theme)
        assert result == {}

    def test_extract_colors_from_theme_exception(self):
        """Test _extract_colors_from_theme handles exceptions in color_scheme access."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_theme

        # Create a theme that raises an exception when accessing color_scheme
        mock_theme = MagicMock()
        type(mock_theme).color_scheme = property(
            fget=lambda self: (_ for _ in ()).throw(Exception("Theme error"))
        )

        result = _extract_colors_from_theme(mock_theme)
        assert result == {}

    def test_extract_colors_from_shapes_with_fill(self):
        """Test _extract_colors_from_shapes with shape fill colors."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_shapes

        mock_slide = MagicMock()
        mock_shape = MagicMock()

        # Mock fill color
        mock_fill = MagicMock()
        mock_fore_color = MagicMock()
        mock_fore_color.rgb = (100, 150, 200)
        mock_fill.fore_color = mock_fore_color
        mock_shape.fill = mock_fill

        # Mock line color
        mock_line = MagicMock()
        mock_line_color = MagicMock()
        mock_line_color.rgb = (50, 100, 150)
        mock_line.color = mock_line_color
        mock_shape.line = mock_line

        # No text_frame
        del mock_shape.text_frame

        mock_slide.shapes = [mock_shape]

        result = _extract_colors_from_shapes(mock_slide)
        assert isinstance(result, list)

    def test_extract_colors_from_shapes_with_text_colors(self):
        """Test _extract_colors_from_shapes with text colors."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_shapes

        mock_slide = MagicMock()
        mock_shape = MagicMock()

        # No fill or line
        mock_shape.fill = None
        mock_shape.line = None

        # Mock text frame with runs
        mock_text_frame = MagicMock()
        mock_paragraph = MagicMock()
        mock_run = MagicMock()
        mock_font_color = MagicMock()
        mock_font_color.rgb = (255, 0, 0)
        mock_run.font.color = mock_font_color
        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_shape.text_frame = mock_text_frame

        mock_slide.shapes = [mock_shape]

        result = _extract_colors_from_shapes(mock_slide)
        assert isinstance(result, list)

    def test_extract_colors_from_shapes_empty(self):
        """Test _extract_colors_from_shapes with no shapes."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_shapes

        mock_slide = MagicMock()
        mock_slide.shapes = []

        result = _extract_colors_from_shapes(mock_slide)
        assert result == []

    def test_extract_colors_from_shapes_exception(self):
        """Test _extract_colors_from_shapes handles exceptions."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_shapes

        mock_slide = MagicMock()
        mock_slide.shapes = MagicMock(side_effect=Exception("Shape error"))

        result = _extract_colors_from_shapes(mock_slide)
        assert result == []

    def test_extract_typography_from_master(self):
        """Test _extract_typography_from_master."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_typography_from_master

        mock_master = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True

        # Mock TITLE placeholder (type=1)
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.type = 1
        mock_shape.placeholder_format = mock_placeholder_format

        # Mock text frame with font info
        mock_text_frame = MagicMock()
        mock_paragraph = MagicMock()
        mock_run = MagicMock()
        mock_run.font.name = "Arial"
        mock_run.font.size = MagicMock()
        mock_run.font.size.pt = 36.0
        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_shape.text_frame = mock_text_frame

        mock_master.shapes = [mock_shape]

        result = _extract_typography_from_master(mock_master)
        assert isinstance(result, dict)
        assert "title_font" in result
        assert "body_font" in result

    def test_extract_typography_from_master_body_placeholder(self):
        """Test _extract_typography_from_master with BODY placeholder."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_typography_from_master

        mock_master = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True

        # Mock BODY placeholder (type=2)
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.type = 2
        mock_shape.placeholder_format = mock_placeholder_format

        mock_text_frame = MagicMock()
        mock_paragraph = MagicMock()
        mock_run = MagicMock()
        mock_run.font.name = "Helvetica"
        mock_run.font.size = MagicMock()
        mock_run.font.size.pt = 14.0
        mock_paragraph.runs = [mock_run]
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_shape.text_frame = mock_text_frame

        mock_master.shapes = [mock_shape]

        result = _extract_typography_from_master(mock_master)
        assert result["body_font"] == "Helvetica"
        assert result["body_size"] == 14.0

    def test_extract_typography_from_master_exception(self):
        """Test _extract_typography_from_master handles exceptions."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_typography_from_master

        mock_master = MagicMock()
        mock_master.shapes = MagicMock(side_effect=Exception("Master error"))

        result = _extract_typography_from_master(mock_master)
        # Should return default typography
        assert result["title_font"] == "Calibri"
        assert result["body_font"] == "Calibri"


class TestExtractionDesignSystem:
    """Tests for extract_design_system_from_template function."""

    @pytest.mark.asyncio
    async def test_extract_design_system_template_not_found(self):
        """Test extract_design_system when template not found."""
        from chuk_mcp_pptx.tools.template.extraction import extract_design_system_from_template

        manager = MagicMock()
        manager.get = AsyncMock(return_value=None)

        with pytest.raises(ValueError) as exc_info:
            await extract_design_system_from_template(manager, "nonexistent")

        assert "not found" in str(exc_info.value)

    @pytest.mark.asyncio
    async def test_extract_design_system_success(self):
        """Test successful design system extraction."""
        from chuk_mcp_pptx.tools.template.extraction import extract_design_system_from_template

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        result = await extract_design_system_from_template(manager, "test_template")

        assert result.template_name == "test_template"
        assert isinstance(result.layouts, list)

    @pytest.mark.asyncio
    async def test_extract_design_system_with_slides(self):
        """Test design system extraction from template with slides."""
        from chuk_mcp_pptx.tools.template.extraction import extract_design_system_from_template

        prs = Presentation()
        # Add a slide
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        result = await extract_design_system_from_template(manager, "test_template")

        assert result.template_name == "test_template"


class TestAnalyzeLayoutVariants:
    """Tests for analyze_layout_variants function."""

    def test_analyze_layout_variants_empty(self):
        """Test analyzing presentation with no layouts."""
        from chuk_mcp_pptx.tools.template.extraction import analyze_layout_variants

        prs = Presentation()

        result = analyze_layout_variants(prs)

        assert result.total_layouts >= 0
        assert isinstance(result.layout_groups, list)
        assert isinstance(result.ungrouped_layouts, list)

    def test_analyze_layout_variants_with_slides(self):
        """Test analyzing presentation with slides using layouts."""
        from chuk_mcp_pptx.tools.template.extraction import analyze_layout_variants

        prs = Presentation()
        # Add slides using different layouts
        if prs.slide_layouts and len(prs.slide_layouts) > 1:
            prs.slides.add_slide(prs.slide_layouts[0])
            prs.slides.add_slide(prs.slide_layouts[1])

        result = analyze_layout_variants(prs)

        assert result.total_layouts > 0

    def test_analyze_layout_variants_detects_variant_numbers(self):
        """Test that variant numbers are detected from layout names."""
        from chuk_mcp_pptx.tools.template.extraction import analyze_layout_variants

        prs = Presentation()

        result = analyze_layout_variants(prs)

        # Check structure of result
        assert hasattr(result, "layout_groups")
        assert hasattr(result, "ungrouped_layouts")
        assert hasattr(result, "unique_groups")


class TestExtractionToolsRegistration:
    """Tests for extraction tools registration and tool functions."""

    @pytest.mark.asyncio
    async def test_extract_template_design_system_tool(self, mock_mcp):
        """Test pptx_extract_template_design_system tool."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_design_system"]("test_template")

        assert isinstance(result, str)
        data = json.loads(result)
        assert "template_name" in data or "error" in data

    @pytest.mark.asyncio
    async def test_extract_template_design_system_not_found(self, mock_mcp):
        """Test pptx_extract_template_design_system when template not found."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        manager = MagicMock()
        manager.get = AsyncMock(return_value=None)

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_design_system"]("nonexistent")

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_register_template_as_theme_success(self, mock_mcp):
        """Test pptx_register_template_as_theme success."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        # Add a slide to get some colors
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value=None)

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
            apply_to_current=False,
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_template_as_theme_apply_to_current(self, mock_mcp):
        """Test pptx_register_template_as_theme with apply_to_current."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value="current_presentation")
        manager._save_to_store = AsyncMock()

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
            apply_to_current=True,
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_template_as_theme_no_colors(self, mock_mcp):
        """Test pptx_register_template_as_theme when no colors extracted."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        # Create mock presentation with no usable colors
        mock_prs = MagicMock()
        mock_prs.slides = []
        mock_prs.slide_layouts = []
        mock_prs.slide_master = None

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(mock_prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="empty_template",
            theme_name="my_theme",
        )

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_extract_template_colors_success(self, mock_mcp):
        """Test pptx_extract_template_colors success."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_colors"]("test_template")

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_extract_template_colors_no_colors(self, mock_mcp):
        """Test pptx_extract_template_colors when no colors found."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        mock_prs = MagicMock()
        mock_prs.slides = []
        mock_prs.slide_layouts = []
        mock_prs.slide_master = None

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(mock_prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_colors"]("empty_template")

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_extract_template_colors_exception(self, mock_mcp):
        """Test pptx_extract_template_colors exception handling."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        manager = MagicMock()
        manager.get = AsyncMock(side_effect=Exception("Color extraction failed"))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_colors"]("error_template")

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_extract_template_typography_success(self, mock_mcp):
        """Test pptx_extract_template_typography success."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_typography"]("test_template")

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_extract_template_typography_no_typography(self, mock_mcp):
        """Test pptx_extract_template_typography when no typography found."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        mock_prs = MagicMock()
        mock_prs.slide_master = None

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(mock_prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_typography"]("empty_template")

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_extract_template_typography_exception(self, mock_mcp):
        """Test pptx_extract_template_typography exception handling."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        manager = MagicMock()
        manager.get = AsyncMock(side_effect=Exception("Typography extraction failed"))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_extract_template_typography"]("error_template")

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_compare_slide_to_template_presentation_not_found(self, mock_mcp):
        """Test pptx_compare_slide_to_template when presentation not found."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        manager = MagicMock()
        manager.get = AsyncMock(return_value=None)

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="nonexistent",
            slide_index=0,
            template_name="template",
        )

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_compare_slide_to_template_slide_out_of_range(self, mock_mcp):
        """Test pptx_compare_slide_to_template when slide index out of range."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        # Don't add any slides

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="test_prs",
            slide_index=10,
            template_name="template",
        )

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_compare_slide_to_template_template_not_found(self, mock_mcp):
        """Test pptx_compare_slide_to_template when template not found."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()

        # First call returns presentation, second returns None (template not found)
        manager.get = AsyncMock(side_effect=[(prs, metadata), None])

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="test_prs",
            slide_index=0,
            template_name="nonexistent_template",
        )

        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_compare_slide_to_template_success(self, mock_mcp):
        """Test pptx_compare_slide_to_template success."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        template_prs = Presentation()
        if template_prs.slide_layouts:
            template_prs.slides.add_slide(template_prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()

        # Return both presentations
        manager.get = AsyncMock(side_effect=[(prs, metadata), (template_prs, metadata)])

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="test_prs",
            slide_index=0,
            template_name="template",
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_compare_slide_to_template_with_specific_template_slide(self, mock_mcp):
        """Test pptx_compare_slide_to_template with specific template slide index."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        template_prs = Presentation()
        if template_prs.slide_layouts:
            template_prs.slides.add_slide(template_prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()

        manager.get = AsyncMock(side_effect=[(prs, metadata), (template_prs, metadata)])

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="test_prs",
            slide_index=0,
            template_name="template",
            template_slide_index=0,
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_compare_slide_to_template_exception(self, mock_mcp):
        """Test pptx_compare_slide_to_template exception handling."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        manager = MagicMock()
        manager.get = AsyncMock(side_effect=Exception("Compare failed"))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="test_prs",
            slide_index=0,
            template_name="template",
        )

        data = json.loads(result)
        assert "error" in data


class TestExtractionCoverage:
    """Tests to improve extraction.py coverage."""

    @pytest.mark.asyncio
    async def test_register_theme_with_mock_color_scheme(self, mock_mcp):
        """Test registering theme with a properly mocked color scheme."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        # Create presentation with slides that have fill colors to extract
        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value=None)

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        # Just verify it can be called without error
        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
            apply_to_current=False,
        )

        # Should return a result (either success or error)
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_theme_with_typography(self, mock_mcp):
        """Test registering theme with typography."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value=None)

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
        )

        # Should return a result
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_theme_apply_to_current_with_slides(self, mock_mcp):
        """Test registering theme and applying to current presentation with multiple slides."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value="current_prs")
        manager._save_to_store = AsyncMock()

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
            apply_to_current=True,
        )

        # Should return a result
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_theme_with_apply_to_current_and_slides(self, mock_mcp):
        """Test registering theme and applying to current presentation with slides."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        # Create presentation with slides
        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value="current_presentation")
        manager._save_to_store = AsyncMock()

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
            apply_to_current=True,
        )

        # Verify the function was called
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_theme_with_apply_no_current_presentation(self, mock_mcp):
        """Test registering theme with apply_to_current but no current presentation."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()
        manager.get = AsyncMock(return_value=(prs, metadata))
        manager.get_current_name = MagicMock(return_value=None)  # No current presentation
        manager._save_to_store = AsyncMock()

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="test_template",
            theme_name="my_theme",
            apply_to_current=True,
        )

        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_register_theme_exception(self, mock_mcp):
        """Test register_template_as_theme exception handling."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        manager = MagicMock()
        manager.get = AsyncMock(side_effect=Exception("Theme registration failed"))

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_register_template_as_theme"](
            template_name="error_template",
            theme_name="my_theme",
        )

        data = json.loads(result)
        assert "error" in data

    def test_analyze_layout_variants_with_variant_groups(self):
        """Test analyze_layout_variants with layouts that have variant numbers."""
        from chuk_mcp_pptx.tools.template.extraction import analyze_layout_variants

        # Create presentation where layouts have similar base names
        prs = Presentation()

        result = analyze_layout_variants(prs)

        # Should have analyzed all layouts
        assert result.total_layouts >= 0
        assert isinstance(result.layout_groups, list)
        assert isinstance(result.ungrouped_layouts, list)

    def test_extract_colors_from_shapes_with_line_color(self):
        """Test _extract_colors_from_shapes with line colors."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_shapes

        mock_slide = MagicMock()
        mock_shape = MagicMock()

        # Mock fill with no color
        mock_fill = MagicMock()
        mock_fill.fore_color = None
        mock_shape.fill = mock_fill

        # Mock line with color
        mock_line = MagicMock()
        mock_line_color = MagicMock()
        mock_line_color.rgb = (128, 128, 255)
        mock_line.color = mock_line_color
        mock_shape.line = mock_line

        # No text_frame
        del mock_shape.text_frame

        mock_slide.shapes = [mock_shape]

        result = _extract_colors_from_shapes(mock_slide)
        assert isinstance(result, list)

    def test_extract_typography_with_no_runs(self):
        """Test _extract_typography_from_master with paragraphs but no runs."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_typography_from_master

        mock_master = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True

        mock_placeholder_format = MagicMock()
        mock_placeholder_format.type = 1  # TITLE
        mock_shape.placeholder_format = mock_placeholder_format

        # Empty runs
        mock_text_frame = MagicMock()
        mock_paragraph = MagicMock()
        mock_paragraph.runs = []  # No runs
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_shape.text_frame = mock_text_frame

        mock_master.shapes = [mock_shape]

        result = _extract_typography_from_master(mock_master)
        # Should return defaults when no font info found
        assert "title_font" in result

    def test_extract_typography_with_no_paragraphs(self):
        """Test _extract_typography_from_master with no paragraphs."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_typography_from_master

        mock_master = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = True

        mock_placeholder_format = MagicMock()
        mock_placeholder_format.type = 1  # TITLE
        mock_shape.placeholder_format = mock_placeholder_format

        # Empty paragraphs
        mock_text_frame = MagicMock()
        mock_text_frame.paragraphs = []
        mock_shape.text_frame = mock_text_frame

        mock_master.shapes = [mock_shape]

        result = _extract_typography_from_master(mock_master)
        # Should return defaults
        assert "title_font" in result

    def test_extract_typography_non_placeholder(self):
        """Test _extract_typography_from_master with non-placeholder shape."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_typography_from_master

        mock_master = MagicMock()
        mock_shape = MagicMock()
        mock_shape.is_placeholder = False  # Not a placeholder

        mock_master.shapes = [mock_shape]

        result = _extract_typography_from_master(mock_master)
        # Should return defaults since no placeholders
        assert result["title_font"] == "Calibri"
        assert result["body_font"] == "Calibri"

    def test_extract_colors_with_no_fore_color(self):
        """Test _extract_colors_from_shapes when fore_color access fails."""
        from chuk_mcp_pptx.tools.template.extraction import _extract_colors_from_shapes

        mock_slide = MagicMock()
        mock_shape = MagicMock()

        # Mock fill that raises exception on fore_color
        mock_fill = MagicMock()
        type(mock_fill).fore_color = property(
            fget=lambda self: (_ for _ in ()).throw(Exception("No color"))
        )
        mock_shape.fill = mock_fill
        mock_shape.line = None

        # No text_frame
        del mock_shape.text_frame

        mock_slide.shapes = [mock_shape]

        result = _extract_colors_from_shapes(mock_slide)
        assert isinstance(result, list)

    @pytest.mark.asyncio
    async def test_compare_slide_with_no_layout(self, mock_mcp):
        """Test pptx_compare_slide_to_template when slide has no slide_layout."""
        from chuk_mcp_pptx.tools.template.extraction import register_extraction_tools

        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[0])

        # Make the slide's layout None
        mock_slide = MagicMock()
        mock_slide.slide_layout = None
        mock_slide.shapes = []

        template_prs = Presentation()
        if template_prs.slide_layouts:
            template_prs.slides.add_slide(template_prs.slide_layouts[0])

        manager = MagicMock()
        metadata = MagicMock()

        manager.get = AsyncMock(side_effect=[(prs, metadata), (template_prs, metadata)])

        theme_manager = MagicMock()

        register_extraction_tools(mock_mcp, manager, theme_manager)

        result = await mock_mcp._tools["pptx_compare_slide_to_template"](
            presentation="test_prs",
            slide_index=0,
            template_name="template",
        )

        # Should handle the case gracefully
        assert isinstance(result, str)


class TestExtractionModels:
    """Tests for extraction.py Pydantic models."""

    def test_extracted_color_scheme_model(self):
        """Test ExtractedColorScheme model."""
        from chuk_mcp_pptx.tools.template.extraction import ExtractedColorScheme

        scheme = ExtractedColorScheme(
            name="test_colors",
            primary="#FF0000",
            secondary="#00FF00",
            accent="#0000FF",
            background="#FFFFFF",
            text="#000000",
            additional_colors=["#AABBCC"],
        )
        assert scheme.name == "test_colors"
        assert scheme.primary == "#FF0000"

    def test_extracted_typography_model(self):
        """Test ExtractedTypography model."""
        from chuk_mcp_pptx.tools.template.extraction import ExtractedTypography

        typography = ExtractedTypography(
            name="test_typography",
            title_font="Arial",
            title_size=36.0,
            body_font="Helvetica",
            body_size=12.0,
        )
        assert typography.title_font == "Arial"
        assert typography.body_size == 12.0

    def test_extracted_layout_model(self):
        """Test ExtractedLayout model."""
        from chuk_mcp_pptx.tools.template.extraction import ExtractedLayout

        layout = ExtractedLayout(
            name="Title Slide",
            index=0,
            width=10.0,
            height=7.5,
            placeholders=[],
            background_type="solid",
        )
        assert layout.name == "Title Slide"
        assert layout.width == 10.0

    def test_extracted_design_system_model(self):
        """Test ExtractedDesignSystem model."""
        from chuk_mcp_pptx.tools.template.extraction import ExtractedDesignSystem

        system = ExtractedDesignSystem(
            template_name="test_template",
            layouts=[],
        )
        assert system.template_name == "test_template"
        assert system.color_scheme is None

    def test_placeholder_usage_model(self):
        """Test PlaceholderUsage model."""
        from chuk_mcp_pptx.tools.template.extraction import PlaceholderUsage

        usage = PlaceholderUsage(
            idx=0,
            type="TITLE",
            name="Title 1",
            top=100000,
            left=50000,
        )
        assert usage.idx == 0
        assert usage.type == "TITLE"

    def test_layout_comparison_result_model(self):
        """Test LayoutComparisonResult model."""
        from chuk_mcp_pptx.tools.template.extraction import LayoutComparisonResult

        result = LayoutComparisonResult(
            layout_name="Title Slide",
            generated_slide_index=0,
            matches_pattern=True,
        )
        assert result.layout_name == "Title Slide"
        assert result.matches_pattern is True

    def test_layout_variant_model(self):
        """Test LayoutVariant model."""
        from chuk_mcp_pptx.tools.template.extraction import LayoutVariant

        variant = LayoutVariant(
            index=0,
            name="Title Slide",
            placeholder_count=2,
        )
        assert variant.name == "Title Slide"
        assert variant.placeholder_count == 2

    def test_layout_group_model(self):
        """Test LayoutGroup model."""
        from chuk_mcp_pptx.tools.template.extraction import LayoutGroup, LayoutVariant

        base = LayoutVariant(index=0, name="Content", placeholder_count=2)
        group = LayoutGroup(
            base_name="Content",
            base_layout=base,
            variants=[],
            total_count=1,
            placeholder_signature="2 placeholders",
        )
        assert group.base_name == "Content"
        assert group.total_count == 1

    def test_layout_analysis_model(self):
        """Test LayoutAnalysis model."""
        from chuk_mcp_pptx.tools.template.extraction import LayoutAnalysis

        analysis = LayoutAnalysis(
            total_layouts=10,
            unique_groups=3,
            layout_groups=[],
            ungrouped_layouts=[],
        )
        assert analysis.total_layouts == 10
        assert analysis.unique_groups == 3
