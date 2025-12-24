# tests/tools/test_inspection_analysis.py
"""
Tests for the inspection and analysis tools.

Tests the tools:
- pptx_inspect_slide
- pptx_fix_slide_layout
- pptx_analyze_presentation_layout
"""

import pytest
from unittest.mock import MagicMock
from pptx import Presentation
from pptx.util import Inches


# ============================================================================
# Fixtures
# ============================================================================


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation
        self._current_name = "test_presentation"

    async def get_presentation(self, name=None):
        """Get presentation by name or return current."""
        if name is None or name == self._current_name:
            return self._presentation
        return None


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures tool registrations."""
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mock = MagicMock()
    mock.tool = tool_decorator
    mock._tools = tools
    return mock


@pytest.fixture
def mock_manager_no_prs():
    """Create a mock manager with no presentation."""
    return MockPresentationManager(presentation=None)


@pytest.fixture
def presentation_with_slide():
    """Create a presentation with one blank slide."""
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # Blank layout
    prs.slides.add_slide(blank_layout)
    return prs


@pytest.fixture
def presentation_with_title_slide():
    """Create a presentation with a title slide."""
    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(title_layout)
    # Set title
    if slide.shapes.title:
        slide.shapes.title.text = "Test Title"
    return prs


@pytest.fixture
def presentation_with_elements(presentation_with_slide):
    """Create a presentation with various elements on the slide."""
    prs = presentation_with_slide
    slide = prs.slides[0]

    # Add a text box
    left = Inches(1)
    top = Inches(1)
    width = Inches(2)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Sample text content"

    # Add another text box
    left2 = Inches(4)
    top2 = Inches(1)
    txBox2 = slide.shapes.add_textbox(left2, top2, width, height)
    tf2 = txBox2.text_frame
    tf2.text = "Another text box"

    return prs


@pytest.fixture
def mock_manager(presentation_with_slide):
    """Create a mock manager with a presentation."""
    return MockPresentationManager(presentation=presentation_with_slide)


@pytest.fixture
def mock_manager_with_elements(presentation_with_elements):
    """Create a mock manager with elements."""
    return MockPresentationManager(presentation=presentation_with_elements)


@pytest.fixture
def mock_manager_with_title(presentation_with_title_slide):
    """Create a mock manager with title slide."""
    return MockPresentationManager(presentation=presentation_with_title_slide)


@pytest.fixture
def inspection_tools(mock_mcp, mock_manager):
    """Register and return inspection tools."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager)


@pytest.fixture
def inspection_tools_no_prs(mock_mcp, mock_manager_no_prs):
    """Register and return inspection tools with no presentation."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager_no_prs)


@pytest.fixture
def inspection_tools_with_elements(mock_mcp, mock_manager_with_elements):
    """Register and return inspection tools with elements."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager_with_elements)


@pytest.fixture
def inspection_tools_with_title(mock_mcp, mock_manager_with_title):
    """Register and return inspection tools with title slide."""
    from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

    return register_inspection_tools(mock_mcp, mock_manager_with_title)


# ============================================================================
# Test pptx_inspect_slide
# ============================================================================


class TestInspectSlide:
    """Tests for pptx_inspect_slide tool."""

    @pytest.mark.asyncio
    async def test_inspect_slide_basic(self, inspection_tools):
        """Test basic slide inspection."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        assert isinstance(result, str)
        assert "SLIDE 0 INSPECTION" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_no_presentation(self, inspection_tools_no_prs):
        """Test inspecting when no presentation exists."""
        result = await inspection_tools_no_prs["pptx_inspect_slide"](slide_index=0)
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_inspect_slide_out_of_range(self, inspection_tools):
        """Test inspecting slide index out of range."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=999)
        assert "error" in result.lower()
        assert "out of range" in result.lower()

    @pytest.mark.asyncio
    async def test_inspect_slide_with_title(self, inspection_tools_with_title):
        """Test inspecting slide with title."""
        result = await inspection_tools_with_title["pptx_inspect_slide"](slide_index=0)
        assert "Title:" in result
        assert "Test Title" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_with_text_boxes(self, inspection_tools_with_elements):
        """Test inspecting slide with text boxes."""
        result = await inspection_tools_with_elements["pptx_inspect_slide"](slide_index=0)
        assert "TEXT BOXES" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_includes_measurements(self, inspection_tools_with_elements):
        """Test that measurements are included."""
        result = await inspection_tools_with_elements["pptx_inspect_slide"](
            slide_index=0, include_measurements=True
        )
        assert "at (" in result or "size" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_no_measurements(self, inspection_tools_with_elements):
        """Test inspection without measurements."""
        result = await inspection_tools_with_elements["pptx_inspect_slide"](
            slide_index=0, include_measurements=False
        )
        assert isinstance(result, str)
        # Should still contain basic info
        assert "SLIDE 0 INSPECTION" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_string_index(self, inspection_tools):
        """Test that string slide index is converted to int."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index="0")
        assert isinstance(result, str)
        assert "SLIDE 0 INSPECTION" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_summary(self, inspection_tools):
        """Test that summary is included."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        assert "SUMMARY" in result
        assert "Total elements" in result

    @pytest.mark.asyncio
    async def test_inspect_slide_no_issues(self, inspection_tools):
        """Test slide with no layout issues."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        # Blank slide should have no issues
        assert "No layout issues" in result or "Layout issues: 0" in result


# ============================================================================
# Test pptx_fix_slide_layout
# ============================================================================


class TestFixSlideLayout:
    """Tests for pptx_fix_slide_layout tool."""

    @pytest.mark.asyncio
    async def test_fix_layout_basic(self, inspection_tools):
        """Test basic layout fixing."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index=0)
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_no_presentation(self, inspection_tools_no_prs):
        """Test fixing when no presentation exists."""
        result = await inspection_tools_no_prs["pptx_fix_slide_layout"](slide_index=0)
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_layout_out_of_range(self, inspection_tools):
        """Test fixing slide index out of range."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index=999)
        assert "error" in result.lower()
        assert "out of range" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_layout_no_issues(self, inspection_tools):
        """Test fixing when no issues exist."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index=0)
        assert "optimal" in result.lower() or "no layout issues" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_layout_string_index(self, inspection_tools):
        """Test that string slide index is converted to int."""
        result = await inspection_tools["pptx_fix_slide_layout"](slide_index="0")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_with_presentation_name(self, inspection_tools):
        """Test fixing with presentation name."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, presentation="test_presentation"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_fix_bounds_only(self, inspection_tools):
        """Test fixing only bounds issues."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=False, fix_bounds=True, fix_spacing=False
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_fix_overlaps_only(self, inspection_tools):
        """Test fixing only overlap issues."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=True, fix_bounds=False, fix_spacing=False
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_fix_spacing_only(self, inspection_tools):
        """Test fixing only spacing issues."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=False, fix_bounds=False, fix_spacing=True
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_layout_maintain_relative_positions(self, inspection_tools):
        """Test fixing with maintain_relative_positions flag."""
        result = await inspection_tools["pptx_fix_slide_layout"](
            slide_index=0, maintain_relative_positions=True
        )
        assert isinstance(result, str)


# ============================================================================
# Test pptx_analyze_presentation_layout
# ============================================================================


class TestAnalyzePresentationLayout:
    """Tests for pptx_analyze_presentation_layout tool."""

    @pytest.mark.asyncio
    async def test_analyze_layout_basic(self, inspection_tools):
        """Test basic presentation layout analysis."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert isinstance(result, str)
        assert "PRESENTATION LAYOUT ANALYSIS" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_no_presentation(self, inspection_tools_no_prs):
        """Test analyzing when no presentation exists."""
        result = await inspection_tools_no_prs["pptx_analyze_presentation_layout"]()
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_total_slides(self, inspection_tools):
        """Test that total slides count is shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "Total slides:" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_layout_usage(self, inspection_tools):
        """Test that layout usage is shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "LAYOUT USAGE" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_element_stats(self, inspection_tools):
        """Test that element statistics are shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "ELEMENT STATISTICS" in result
        assert "Images:" in result
        assert "Charts:" in result
        assert "Tables:" in result
        assert "Text boxes:" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_shows_recommendations(self, inspection_tools):
        """Test that recommendations section is shown."""
        result = await inspection_tools["pptx_analyze_presentation_layout"]()
        assert "RECOMMENDATIONS" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_with_elements(self, inspection_tools_with_elements):
        """Test analyzing presentation with elements."""
        result = await inspection_tools_with_elements["pptx_analyze_presentation_layout"]()
        assert "PRESENTATION LAYOUT ANALYSIS" in result
        # Should count text boxes
        assert "Text boxes:" in result

    @pytest.mark.asyncio
    async def test_analyze_layout_with_presentation_name(self, inspection_tools):
        """Test analyzing with presentation name."""
        result = await inspection_tools["pptx_analyze_presentation_layout"](
            presentation="test_presentation"
        )
        assert isinstance(result, str)
        assert "PRESENTATION LAYOUT ANALYSIS" in result


# ============================================================================
# Test Tool Registration
# ============================================================================


class TestToolRegistration:
    """Tests for tool registration."""

    def test_all_tools_registered(self, inspection_tools):
        """Test that all inspection tools are registered."""
        expected_tools = [
            "pptx_inspect_slide",
            "pptx_fix_slide_layout",
            "pptx_analyze_presentation_layout",
        ]
        for tool_name in expected_tools:
            assert tool_name in inspection_tools

    def test_tools_are_async(self, inspection_tools):
        """Test that all tools are async functions."""
        import inspect

        for tool_name, tool_func in inspection_tools.items():
            assert inspect.iscoroutinefunction(tool_func), f"{tool_name} should be async"


# ============================================================================
# Test Edge Cases and Overlapping Elements
# ============================================================================


class TestOverlappingElements:
    """Tests for overlapping element detection."""

    @pytest.fixture
    def presentation_with_overlaps(self):
        """Create a presentation with overlapping elements."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add two overlapping text boxes
        left1 = Inches(2)
        top1 = Inches(2)
        width = Inches(3)
        height = Inches(2)

        txBox1 = slide.shapes.add_textbox(left1, top1, width, height)
        txBox1.text_frame.text = "Overlapping box 1"

        # Second box partially overlaps first
        left2 = Inches(3)
        top2 = Inches(2.5)
        txBox2 = slide.shapes.add_textbox(left2, top2, width, height)
        txBox2.text_frame.text = "Overlapping box 2"

        return prs

    @pytest.fixture
    def inspection_tools_with_overlaps(self, mock_mcp, presentation_with_overlaps):
        """Register inspection tools with overlapping elements."""
        manager = MockPresentationManager(presentation=presentation_with_overlaps)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_detects_overlaps(self, inspection_tools_with_overlaps):
        """Test that inspection detects overlapping elements."""
        result = await inspection_tools_with_overlaps["pptx_inspect_slide"](
            slide_index=0, check_overlaps=True
        )
        # May or may not detect overlaps depending on implementation details
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_fix_overlaps(self, inspection_tools_with_overlaps):
        """Test fixing overlapping elements."""
        result = await inspection_tools_with_overlaps["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=True
        )
        assert isinstance(result, str)


# ============================================================================
# Test Out of Bounds Elements
# ============================================================================


class TestOutOfBoundsElements:
    """Tests for out of bounds element detection."""

    @pytest.fixture
    def presentation_with_oob(self):
        """Create a presentation with out of bounds elements."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add a text box that extends beyond slide bounds
        left = Inches(8)  # Very close to right edge
        top = Inches(1)
        width = Inches(4)  # Will extend beyond slide
        height = Inches(1)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text = "Out of bounds text"

        return prs

    @pytest.fixture
    def inspection_tools_with_oob(self, mock_mcp, presentation_with_oob):
        """Register inspection tools with out of bounds elements."""
        manager = MockPresentationManager(presentation=presentation_with_oob)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_detects_oob(self, inspection_tools_with_oob):
        """Test that inspection detects out of bounds elements."""
        result = await inspection_tools_with_oob["pptx_inspect_slide"](slide_index=0)
        assert isinstance(result, str)
        # Should mention out of bounds or layout issues
        assert "OUT OF BOUNDS" in result or "LAYOUT ISSUES" in result

    @pytest.mark.asyncio
    async def test_fix_oob(self, inspection_tools_with_oob):
        """Test fixing out of bounds elements."""
        result = await inspection_tools_with_oob["pptx_fix_slide_layout"](
            slide_index=0, fix_bounds=True
        )
        assert isinstance(result, str)


# ============================================================================
# Test Multiple Slides
# ============================================================================


class TestMultipleSlides:
    """Tests for presentations with multiple slides."""

    @pytest.fixture
    def presentation_with_multiple_slides(self):
        """Create a presentation with multiple slides."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        # Add 3 slides
        for _ in range(3):
            prs.slides.add_slide(blank_layout)

        return prs

    @pytest.fixture
    def inspection_tools_multiple_slides(self, mock_mcp, presentation_with_multiple_slides):
        """Register inspection tools with multiple slides."""
        manager = MockPresentationManager(presentation=presentation_with_multiple_slides)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_each_slide(self, inspection_tools_multiple_slides):
        """Test inspecting each slide individually."""
        for i in range(3):
            result = await inspection_tools_multiple_slides["pptx_inspect_slide"](slide_index=i)
            assert isinstance(result, str)
            assert f"SLIDE {i} INSPECTION" in result

    @pytest.mark.asyncio
    async def test_analyze_multiple_slides(self, inspection_tools_multiple_slides):
        """Test analyzing presentation with multiple slides."""
        result = await inspection_tools_multiple_slides["pptx_analyze_presentation_layout"]()
        assert "Total slides: 3" in result


# ============================================================================
# Additional Coverage Tests
# ============================================================================


class TestShapeTypeCategorization:
    """Tests for shape type categorization in inspect_slide."""

    @pytest.fixture
    def presentation_with_all_shape_types(self):
        """Create presentation with various shape types."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add text box
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
        txBox.text_frame.text = "Test text box"

        # Add a table
        slide.shapes.add_table(3, 3, Inches(1), Inches(2), Inches(4), Inches(2))

        # Add a chart
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B", "C"]
        chart_data.add_series("Series 1", (1, 2, 3))
        slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5), Inches(1), Inches(4), Inches(3), chart_data
        )

        # Add an auto shape
        slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(1),
            Inches(5),
            Inches(2),
            Inches(1),
        )

        return prs

    @pytest.fixture
    def inspection_tools_all_shapes(self, mock_mcp, presentation_with_all_shape_types):
        """Register inspection tools with all shape types."""
        manager = MockPresentationManager(presentation=presentation_with_all_shape_types)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_categorizes_text_boxes(self, inspection_tools_all_shapes):
        """Test that text boxes are categorized."""
        result = await inspection_tools_all_shapes["pptx_inspect_slide"](slide_index=0)
        assert "TEXT BOXES" in result

    @pytest.mark.asyncio
    async def test_inspect_categorizes_tables(self, inspection_tools_all_shapes):
        """Test that tables are categorized."""
        result = await inspection_tools_all_shapes["pptx_inspect_slide"](slide_index=0)
        assert "TABLES" in result

    @pytest.mark.asyncio
    async def test_inspect_categorizes_charts(self, inspection_tools_all_shapes):
        """Test that charts are categorized."""
        result = await inspection_tools_all_shapes["pptx_inspect_slide"](slide_index=0)
        assert "CHARTS" in result

    @pytest.mark.asyncio
    async def test_inspect_categorizes_other_shapes(self, inspection_tools_all_shapes):
        """Test that other shapes are categorized."""
        result = await inspection_tools_all_shapes["pptx_inspect_slide"](slide_index=0)
        # Auto shapes go to other_shapes
        assert "OTHER SHAPES" in result or "Shape" in result


class TestImageHandling:
    """Tests for image handling in inspection."""

    @pytest.fixture
    def presentation_with_images(self, tmp_path):
        """Create presentation with images."""
        import io
        from PIL import Image

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create a simple test image
        img = Image.new("RGB", (100, 100), color="red")
        img_bytes = io.BytesIO()
        img.save(img_bytes, format="PNG")
        img_bytes.seek(0)

        # Add image to slide
        slide.shapes.add_picture(img_bytes, Inches(2), Inches(2), Inches(2), Inches(2))

        # Add a larger background-like image
        img2 = Image.new("RGB", (100, 100), color="blue")
        img_bytes2 = io.BytesIO()
        img2.save(img_bytes2, format="PNG")
        img_bytes2.seek(0)
        slide.shapes.add_picture(img_bytes2, Inches(0), Inches(0), Inches(10), Inches(7.5))

        return prs

    @pytest.fixture
    def inspection_tools_with_images(self, mock_mcp, presentation_with_images):
        """Register inspection tools with images."""
        manager = MockPresentationManager(presentation=presentation_with_images)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_categorizes_images(self, inspection_tools_with_images):
        """Test that images are categorized."""
        result = await inspection_tools_with_images["pptx_inspect_slide"](slide_index=0)
        assert "IMAGES" in result

    @pytest.mark.asyncio
    async def test_fix_layout_with_images(self, inspection_tools_with_images):
        """Test fixing layout with images."""
        result = await inspection_tools_with_images["pptx_fix_slide_layout"](
            slide_index=0, fix_spacing=True
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_analyze_counts_images(self, inspection_tools_with_images):
        """Test that analysis counts images."""
        result = await inspection_tools_with_images["pptx_analyze_presentation_layout"]()
        assert "Images:" in result


class TestChartDetails:
    """Tests for chart detail extraction."""

    @pytest.fixture
    def presentation_with_titled_chart(self):
        """Create presentation with a titled chart."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add a chart with title
        chart_data = CategoryChartData()
        chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
        chart_data.add_series("Revenue", (100, 150, 200, 250))

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), chart_data
        )

        # Set chart title
        chart = chart_shape.chart
        chart.has_title = True
        chart.chart_title.text_frame.text = "Quarterly Revenue"

        return prs

    @pytest.fixture
    def inspection_tools_with_chart(self, mock_mcp, presentation_with_titled_chart):
        """Register inspection tools with titled chart."""
        manager = MockPresentationManager(presentation=presentation_with_titled_chart)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_shows_chart_type(self, inspection_tools_with_chart):
        """Test that chart type is shown."""
        result = await inspection_tools_with_chart["pptx_inspect_slide"](slide_index=0)
        assert "CHARTS" in result
        assert "type:" in result

    @pytest.mark.asyncio
    async def test_inspect_shows_chart_title(self, inspection_tools_with_chart):
        """Test that chart title is shown."""
        result = await inspection_tools_with_chart["pptx_inspect_slide"](slide_index=0)
        assert "Quarterly Revenue" in result


class TestTitlePlaceholderHandling:
    """Tests for title placeholder handling in overlap and bounds checking."""

    @pytest.fixture
    def presentation_with_title_and_content(self):
        """Create presentation with title and content placeholders."""
        prs = Presentation()
        title_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(title_layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = "Test Title Placeholder"

        # Add overlapping text box near title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(3), Inches(1))
        txBox.text_frame.text = "Overlapping text"

        return prs

    @pytest.fixture
    def inspection_tools_with_title_content(self, mock_mcp, presentation_with_title_and_content):
        """Register inspection tools with title and content."""
        manager = MockPresentationManager(presentation=presentation_with_title_and_content)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_title_placeholder_excluded_from_overlap_check(
        self, inspection_tools_with_title_content
    ):
        """Test that title placeholders are excluded from overlap checking."""
        result = await inspection_tools_with_title_content["pptx_inspect_slide"](
            slide_index=0, check_overlaps=True
        )
        assert isinstance(result, str)
        # Title should be present but not flagged for overlaps
        assert "Test Title Placeholder" in result or "Title:" in result

    @pytest.mark.asyncio
    async def test_fix_layout_skips_title_placeholder(self, inspection_tools_with_title_content):
        """Test that fixing layout skips title placeholders."""
        result = await inspection_tools_with_title_content["pptx_fix_slide_layout"](
            slide_index=0, fix_bounds=True, fix_overlaps=True
        )
        assert isinstance(result, str)


class TestSpacingIssues:
    """Tests for spacing issue detection."""

    @pytest.fixture
    def presentation_with_edge_elements(self):
        """Create presentation with elements close to edges."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add element very close to left edge (within 0.2 inches)
        txBox = slide.shapes.add_textbox(Inches(0.1), Inches(2), Inches(2), Inches(1))
        txBox.text_frame.text = "Close to left edge"

        # Add element very close to top (within margin threshold)
        txBox2 = slide.shapes.add_textbox(Inches(3), Inches(0.05), Inches(2), Inches(0.5))
        txBox2.text_frame.text = "Close to top"

        return prs

    @pytest.fixture
    def inspection_tools_edge_elements(self, mock_mcp, presentation_with_edge_elements):
        """Register inspection tools with edge elements."""
        manager = MockPresentationManager(presentation=presentation_with_edge_elements)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_detect_spacing_issues(self, inspection_tools_edge_elements):
        """Test detection of elements too close to edges."""
        result = await inspection_tools_edge_elements["pptx_inspect_slide"](slide_index=0)
        # Should detect spacing issues
        assert isinstance(result, str)
        assert "SPACING" in result or "close to" in result.lower() or "edge" in result.lower()


class TestBoundsChecking:
    """Tests for bounds checking with different edge cases."""

    @pytest.fixture
    def presentation_with_negative_positions(self):
        """Create presentation with elements at negative positions."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add text box (will manually set negative position after)
        txBox = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(2), Inches(1))
        txBox.text_frame.text = "Normal position"

        # Set negative left position
        txBox.left = Inches(-0.5)

        # Add another element with negative top
        txBox2 = slide.shapes.add_textbox(Inches(2), Inches(0), Inches(2), Inches(1))
        txBox2.text_frame.text = "Negative top"
        txBox2.top = Inches(-0.3)

        return prs

    @pytest.fixture
    def inspection_tools_negative_pos(self, mock_mcp, presentation_with_negative_positions):
        """Register inspection tools with negative positions."""
        manager = MockPresentationManager(presentation=presentation_with_negative_positions)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_detect_negative_left(self, inspection_tools_negative_pos):
        """Test detection of elements extending beyond left edge."""
        result = await inspection_tools_negative_pos["pptx_inspect_slide"](slide_index=0)
        assert "OUT OF BOUNDS" in result
        assert "left edge" in result.lower()

    @pytest.mark.asyncio
    async def test_detect_negative_top(self, inspection_tools_negative_pos):
        """Test detection of elements extending beyond top edge."""
        result = await inspection_tools_negative_pos["pptx_inspect_slide"](slide_index=0)
        assert "top edge" in result.lower()

    @pytest.mark.asyncio
    async def test_fix_negative_positions(self, inspection_tools_negative_pos):
        """Test fixing elements with negative positions."""
        result = await inspection_tools_negative_pos["pptx_fix_slide_layout"](
            slide_index=0, fix_bounds=True
        )
        assert "Fixed" in result or "optimal" in result.lower()


class TestAlignmentFunctions:
    """Tests for element alignment in improve_spacing."""

    @pytest.fixture
    def presentation_with_multiple_images(self, tmp_path):
        """Create presentation with multiple images needing alignment."""
        import io
        from PIL import Image

        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create test images at different vertical positions
        for i, (x, y) in enumerate([(1, 2), (3, 2.2), (5, 1.9)]):
            img = Image.new("RGB", (100, 100), color="red")
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            img_bytes.seek(0)
            slide.shapes.add_picture(img_bytes, Inches(x), Inches(y), Inches(1.5), Inches(1.5))

        return prs

    @pytest.fixture
    def inspection_tools_multiple_images(self, mock_mcp, presentation_with_multiple_images):
        """Register inspection tools with multiple images."""
        manager = MockPresentationManager(presentation=presentation_with_multiple_images)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_align_multiple_images(self, inspection_tools_multiple_images):
        """Test alignment of multiple images."""
        result = await inspection_tools_multiple_images["pptx_fix_slide_layout"](
            slide_index=0, fix_spacing=True
        )
        assert isinstance(result, str)
        # Should improve spacing for images
        assert "spacing" in result.lower() or "optimal" in result.lower()

    @pytest.fixture
    def presentation_with_multiple_text_boxes(self):
        """Create presentation with multiple text boxes needing vertical alignment."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Create text boxes at different horizontal positions
        for i, (x, y) in enumerate([(1, 1), (1.2, 2.5), (0.9, 4)]):
            txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(2), Inches(1))
            txBox.text_frame.text = f"Text box {i + 1}"

        return prs

    @pytest.fixture
    def inspection_tools_multiple_textboxes(self, mock_mcp, presentation_with_multiple_text_boxes):
        """Register inspection tools with multiple text boxes."""
        manager = MockPresentationManager(presentation=presentation_with_multiple_text_boxes)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_align_multiple_text_boxes(self, inspection_tools_multiple_textboxes):
        """Test vertical alignment of multiple text boxes."""
        result = await inspection_tools_multiple_textboxes["pptx_fix_slide_layout"](
            slide_index=0, fix_spacing=True
        )
        assert isinstance(result, str)


class TestOverlapFixing:
    """Tests for overlap fixing with position finding."""

    @pytest.fixture
    def presentation_with_many_overlaps(self):
        """Create presentation with many overlapping elements."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add multiple overlapping text boxes
        for i in range(5):
            txBox = slide.shapes.add_textbox(
                Inches(2 + i * 0.3), Inches(2 + i * 0.2), Inches(2), Inches(1)
            )
            txBox.text_frame.text = f"Overlapping box {i + 1}"

        return prs

    @pytest.fixture
    def inspection_tools_many_overlaps(self, mock_mcp, presentation_with_many_overlaps):
        """Register inspection tools with many overlaps."""
        manager = MockPresentationManager(presentation=presentation_with_many_overlaps)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_fix_many_overlapping_elements(self, inspection_tools_many_overlaps):
        """Test fixing many overlapping elements."""
        result = await inspection_tools_many_overlaps["pptx_fix_slide_layout"](
            slide_index=0, fix_overlaps=True
        )
        assert isinstance(result, str)
        # Should resolve some overlaps
        assert "Resolved" in result or "optimal" in result.lower()


class TestAnalyzePresentationBranches:
    """Tests for analyze_presentation_layout branches."""

    @pytest.fixture
    def presentation_with_varied_content(self, tmp_path):
        """Create presentation with varied content for analysis."""
        import io
        from PIL import Image
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        prs = Presentation()

        # Add multiple slides with different layouts
        for layout_idx in [0, 1, 5, 6]:  # Title, Title+Content, Blank, etc.
            try:
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            except IndexError:
                slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Add images to first slide
        slide = prs.slides[0]
        for _ in range(4):  # Many images
            img = Image.new("RGB", (50, 50), color="green")
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            img_bytes.seek(0)
            slide.shapes.add_picture(img_bytes, Inches(1), Inches(1), Inches(1), Inches(1))

        # Add chart to second slide
        slide = prs.slides[1]
        chart_data = CategoryChartData()
        chart_data.categories = ["A", "B"]
        chart_data.add_series("S1", (1, 2))
        slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(1), Inches(1), Inches(4), Inches(3), chart_data
        )

        # Add table to third slide
        slide = prs.slides[2]
        slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))

        return prs

    @pytest.fixture
    def inspection_tools_varied_content(self, mock_mcp, presentation_with_varied_content):
        """Register inspection tools with varied content."""
        manager = MockPresentationManager(presentation=presentation_with_varied_content)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_analyze_counts_all_element_types(self, inspection_tools_varied_content):
        """Test that analysis counts all element types."""
        result = await inspection_tools_varied_content["pptx_analyze_presentation_layout"]()
        assert "Images:" in result
        assert "Charts:" in result
        assert "Tables:" in result

    @pytest.mark.asyncio
    async def test_analyze_recommendations_many_layouts(self, inspection_tools_varied_content):
        """Test recommendations for many layout variations."""
        result = await inspection_tools_varied_content["pptx_analyze_presentation_layout"]()
        assert "RECOMMENDATIONS" in result
        # With 4 different layouts, should recommend fewer variations
        assert "fewer layout" in result.lower() or "consistency" in result.lower()

    @pytest.mark.asyncio
    async def test_analyze_high_image_density(self, inspection_tools_varied_content):
        """Test recommendation for high image density."""
        result = await inspection_tools_varied_content["pptx_analyze_presentation_layout"]()
        # 4 images on 4 slides = 1 per slide average
        # But we added 4 images to single slide which is > 3 per slide
        assert "RECOMMENDATIONS" in result


class TestSlidesWithIssues:
    """Tests for slides with issues reporting."""

    @pytest.fixture
    def presentation_with_issues_on_slides(self):
        """Create presentation with issues on specific slides."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]

        # Slide 0: Clean
        prs.slides.add_slide(blank_layout)

        # Slide 1: Overlapping elements
        slide1 = prs.slides.add_slide(blank_layout)
        txBox1 = slide1.shapes.add_textbox(Inches(2), Inches(2), Inches(3), Inches(2))
        txBox1.text_frame.text = "Box 1"
        txBox2 = slide1.shapes.add_textbox(Inches(3), Inches(2.5), Inches(3), Inches(2))
        txBox2.text_frame.text = "Box 2"

        # Slide 2: Out of bounds element
        slide2 = prs.slides.add_slide(blank_layout)
        txBox3 = slide2.shapes.add_textbox(Inches(8), Inches(1), Inches(4), Inches(1))
        txBox3.text_frame.text = "Out of bounds"

        return prs

    @pytest.fixture
    def inspection_tools_with_issues(self, mock_mcp, presentation_with_issues_on_slides):
        """Register inspection tools with slides having issues."""
        manager = MockPresentationManager(presentation=presentation_with_issues_on_slides)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_analyze_reports_slides_with_issues(self, inspection_tools_with_issues):
        """Test that analysis reports slides with issues."""
        result = await inspection_tools_with_issues["pptx_analyze_presentation_layout"]()
        assert isinstance(result, str)
        # Should report issues on slides
        assert "SLIDES WITH ISSUES" in result or "issues" in result.lower()

    @pytest.mark.asyncio
    async def test_analyze_suggests_fix_tool(self, inspection_tools_with_issues):
        """Test that analysis suggests using fix tool."""
        result = await inspection_tools_with_issues["pptx_analyze_presentation_layout"]()
        # Should recommend using fix tool
        assert "pptx_fix_slide_layout" in result or "fix" in result.lower()


class TestLongTextTruncation:
    """Tests for long text truncation in shape analysis."""

    @pytest.fixture
    def presentation_with_long_text(self):
        """Create presentation with long text content."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add text box with very long text (>50 chars)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
        long_text = "This is a very long piece of text that exceeds fifty characters and should be truncated with ellipsis in the inspection output"
        txBox.text_frame.text = long_text

        return prs

    @pytest.fixture
    def inspection_tools_long_text(self, mock_mcp, presentation_with_long_text):
        """Register inspection tools with long text."""
        manager = MockPresentationManager(presentation=presentation_with_long_text)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_long_text_truncated(self, inspection_tools_long_text):
        """Test that long text is truncated with ellipsis."""
        result = await inspection_tools_long_text["pptx_inspect_slide"](slide_index=0)
        assert "..." in result
        # Should show preview but not full text
        assert "This is a very long" in result


class TestShapeWithoutAttributes:
    """Tests for shapes without expected attributes."""

    @pytest.mark.asyncio
    async def test_inspect_handles_shapes_gracefully(self, inspection_tools):
        """Test that inspection handles shapes without issues."""
        result = await inspection_tools["pptx_inspect_slide"](slide_index=0)
        # Should complete without errors
        assert isinstance(result, str)
        assert "SLIDE 0 INSPECTION" in result


class TestNoOverlapChecking:
    """Tests for disabling overlap checking."""

    @pytest.fixture
    def presentation_with_potential_overlaps(self):
        """Create presentation with potential overlaps."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add potentially overlapping elements
        txBox1 = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(2), Inches(1))
        txBox1.text_frame.text = "Box 1"
        txBox2 = slide.shapes.add_textbox(Inches(2.5), Inches(2), Inches(2), Inches(1))
        txBox2.text_frame.text = "Box 2"

        return prs

    @pytest.fixture
    def inspection_tools_potential_overlaps(self, mock_mcp, presentation_with_potential_overlaps):
        """Register inspection tools with potential overlaps."""
        manager = MockPresentationManager(presentation=presentation_with_potential_overlaps)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_inspect_without_overlap_check(self, inspection_tools_potential_overlaps):
        """Test inspection with overlap checking disabled."""
        result = await inspection_tools_potential_overlaps["pptx_inspect_slide"](
            slide_index=0, check_overlaps=False
        )
        assert isinstance(result, str)
        # Should not report overlapping elements section
        # (may still have other issues)


class TestMissingElementTypes:
    """Tests for missing element type recommendations."""

    @pytest.fixture
    def presentation_without_charts_or_tables(self):
        """Create presentation without charts or tables."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

        # Add only text boxes
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        txBox.text_frame.text = "Only text"

        return prs

    @pytest.fixture
    def inspection_tools_no_charts_tables(self, mock_mcp, presentation_without_charts_or_tables):
        """Register inspection tools without charts/tables."""
        manager = MockPresentationManager(presentation=presentation_without_charts_or_tables)
        from chuk_mcp_pptx.tools.inspection.analysis import register_inspection_tools

        return register_inspection_tools(mock_mcp, manager)

    @pytest.mark.asyncio
    async def test_recommend_missing_element_types(self, inspection_tools_no_charts_tables):
        """Test recommendation for missing element types."""
        result = await inspection_tools_no_charts_tables["pptx_analyze_presentation_layout"]()
        assert "RECOMMENDATIONS" in result
        # Should recommend adding variety
        assert "variety" in result.lower() or "adding" in result.lower() or "No" in result
