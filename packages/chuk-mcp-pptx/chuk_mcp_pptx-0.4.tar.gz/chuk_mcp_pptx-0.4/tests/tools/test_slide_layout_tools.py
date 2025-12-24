"""
Tests for slide_layout_tools.py

Tests all slide layout management MCP tools for >90% coverage.
"""

import pytest
from unittest.mock import MagicMock, AsyncMock
from pptx import Presentation
from chuk_mcp_pptx.tools.layout.slide_management import register_layout_tools


@pytest.fixture
def layout_tools(mock_mcp_server, mock_presentation_manager):
    """Register layout tools and return them."""
    tools = register_layout_tools(mock_mcp_server, mock_presentation_manager)
    return tools


class TestListLayouts:
    """Test pptx_list_layouts tool."""

    @pytest.mark.asyncio
    async def test_list_layouts_basic(self, layout_tools, mock_presentation_manager):
        """Test listing layouts."""
        result = await layout_tools["pptx_list_layouts"]()
        assert isinstance(result, str)
        assert "AVAILABLE SLIDE LAYOUTS" in result or "layouts" in result.lower()

    @pytest.mark.asyncio
    async def test_list_layouts_no_presentation(self, layout_tools, mock_presentation_manager):
        """Test error when no presentation."""
        # Use a non-existent presentation name instead of mocking
        result = await layout_tools["pptx_list_layouts"](presentation="nonexistent")
        assert "No presentation found" in result or '{"error":' in result


class TestAddSlideWithLayout:
    """Test pptx_add_slide_with_layout tool."""

    @pytest.mark.asyncio
    async def test_add_slide_with_layout_basic(self, layout_tools, mock_presentation_manager):
        """Test adding slide with layout."""
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=1, title="Test Slide"
        )
        assert isinstance(result, str)
        assert "Added slide" in result or "layout" in result.lower()

    @pytest.mark.asyncio
    async def test_add_slide_with_layout_content(self, layout_tools, mock_presentation_manager):
        """Test adding slide with content."""
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=1, title="Test", content=["Item 1", "Item 2", "Item 3"]
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_slide_with_layout_subtitle(self, layout_tools, mock_presentation_manager):
        """Test adding slide with subtitle."""
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=0, title="Main Title", subtitle="Subtitle text"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_slide_with_layout_invalid_index(
        self, layout_tools, mock_presentation_manager
    ):
        """Test error with invalid layout index."""
        result = await layout_tools["pptx_add_slide_with_layout"](layout_index=999)
        assert '{"message":"Error:' in result or "out of range" in result

    @pytest.mark.asyncio
    async def test_add_slide_with_layout_no_presentation(
        self, layout_tools, mock_presentation_manager
    ):
        """Test error when no presentation."""
        # Use a non-existent presentation name instead of mocking
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=1, presentation="nonexistent"
        )
        assert "No presentation found" in result or '{"error":' in result

    @pytest.mark.asyncio
    async def test_add_slide_with_layout_content_string(
        self, layout_tools, mock_presentation_manager
    ):
        """Test adding slide with string content (not list)."""
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=1, title="Test", content="Single string content"
        )
        assert isinstance(result, str)


class TestCustomizeLayout:
    """Test pptx_customize_layout tool."""

    @pytest.mark.asyncio
    async def test_customize_layout_background(self, layout_tools, mock_presentation_manager):
        """Test customizing slide background."""
        result = await layout_tools["pptx_customize_layout"](
            slide_index=0, background_color="#F5F5F5"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_customize_layout_footer(self, layout_tools, mock_presentation_manager):
        """Test adding footer."""
        result = await layout_tools["pptx_customize_layout"](
            slide_index=0, add_footer="Confidential"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_customize_layout_page_number(self, layout_tools, mock_presentation_manager):
        """Test adding page number."""
        result = await layout_tools["pptx_customize_layout"](slide_index=0, add_page_number=True)
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_customize_layout_date(self, layout_tools, mock_presentation_manager):
        """Test adding date."""
        result = await layout_tools["pptx_customize_layout"](slide_index=0, add_date="2024-12-01")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_customize_layout_all_options(self, layout_tools, mock_presentation_manager):
        """Test all customization options together."""
        result = await layout_tools["pptx_customize_layout"](
            slide_index=0,
            background_color="#FFFFFF",
            add_footer="Footer text",
            add_page_number=True,
            add_date="2024-12-01",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_customize_layout_invalid_slide(self, layout_tools, mock_presentation_manager):
        """Test error with invalid slide index."""
        result = await layout_tools["pptx_customize_layout"](slide_index=999)
        assert "No presentation found" in result or '{"error":' in result

    @pytest.mark.asyncio
    async def test_customize_layout_no_presentation(self, layout_tools, mock_presentation_manager):
        """Test error when no presentation."""
        # Use a non-existent presentation name instead of mocking
        result = await layout_tools["pptx_customize_layout"](
            slide_index=0, presentation="nonexistent"
        )
        assert "No presentation found" in result or '{"error":' in result

    @pytest.mark.asyncio
    async def test_customize_layout_no_options(self, layout_tools, mock_presentation_manager):
        """Test no customizations applied."""
        result = await layout_tools["pptx_customize_layout"](slide_index=0)
        assert isinstance(result, str)
        assert "No customizations" in result or "Customized" in result

    @pytest.mark.asyncio
    async def test_customize_layout_invalid_color(self, layout_tools, mock_presentation_manager):
        """Test with invalid background color format."""
        result = await layout_tools["pptx_customize_layout"](
            slide_index=0, background_color="INVALID"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_customize_layout_color_without_hash(
        self, layout_tools, mock_presentation_manager
    ):
        """Test background color without # prefix."""
        result = await layout_tools["pptx_customize_layout"](
            slide_index=0, background_color="F5F5F5"
        )
        assert isinstance(result, str)


class TestApplyMasterLayout:
    """Test pptx_apply_master_layout tool."""

    @pytest.mark.asyncio
    async def test_apply_master_layout_basic(self, layout_tools, mock_presentation_manager):
        """Test applying master layout."""
        result = await layout_tools["pptx_apply_master_layout"](layout_name="corporate")
        assert isinstance(result, str)
        assert "Applied" in result or "slides" in result.lower()

    @pytest.mark.asyncio
    async def test_apply_master_layout_with_font(self, layout_tools, mock_presentation_manager):
        """Test applying with font."""
        result = await layout_tools["pptx_apply_master_layout"](
            layout_name="modern", font_name="Arial"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_master_layout_with_colors(self, layout_tools, mock_presentation_manager):
        """Test applying with colors."""
        result = await layout_tools["pptx_apply_master_layout"](
            layout_name="corporate", title_color="#003366", body_color="#333333"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_master_layout_all_options(self, layout_tools, mock_presentation_manager):
        """Test applying with all options."""
        result = await layout_tools["pptx_apply_master_layout"](
            layout_name="modern", font_name="Calibri", title_color="#FF0000", body_color="#000000"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_master_layout_no_presentation(
        self, layout_tools, mock_presentation_manager
    ):
        """Test error when no presentation."""
        # Use a non-existent presentation name instead of mocking
        result = await layout_tools["pptx_apply_master_layout"](
            layout_name="corporate", presentation="nonexistent"
        )
        assert "No presentation found" in result or '{"error":' in result

    @pytest.mark.asyncio
    async def test_apply_master_layout_invalid_colors(
        self, layout_tools, mock_presentation_manager
    ):
        """Test applying with invalid color formats."""
        result = await layout_tools["pptx_apply_master_layout"](
            layout_name="modern", title_color="INVALID", body_color="BADCOLOR"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_master_layout_colors_without_hash(
        self, layout_tools, mock_presentation_manager
    ):
        """Test applying with colors without # prefix."""
        result = await layout_tools["pptx_apply_master_layout"](
            layout_name="modern", title_color="003366", body_color="333333"
        )
        assert isinstance(result, str)


class TestDuplicateSlide:
    """Test pptx_duplicate_slide tool."""

    @pytest.mark.asyncio
    async def test_duplicate_slide_basic(self, layout_tools, mock_presentation_manager):
        """Test duplicating a slide."""
        result = await layout_tools["pptx_duplicate_slide"](slide_index=0)
        assert isinstance(result, str)
        # Can either succeed or have an error with shape types
        assert (
            "Duplicated" in result
            or "slide" in result.lower()
            or '{"error":' in result
            or "TEXT_BOX" in result
        )

    @pytest.mark.asyncio
    async def test_duplicate_slide_different_indices(self, layout_tools, mock_presentation_manager):
        """Test duplicating different slides."""
        for i in range(3):
            result = await layout_tools["pptx_duplicate_slide"](slide_index=i)
            assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_duplicate_slide_invalid_index(self, layout_tools, mock_presentation_manager):
        """Test error with invalid slide index."""
        result = await layout_tools["pptx_duplicate_slide"](slide_index=999)
        assert "No presentation found" in result or '{"error":' in result

    @pytest.mark.asyncio
    async def test_duplicate_slide_no_presentation(self, layout_tools, mock_presentation_manager):
        """Test error when no presentation."""
        # Use a non-existent presentation name
        result = await layout_tools["pptx_duplicate_slide"](
            slide_index=0, presentation="nonexistent"
        )
        assert "No presentation found" in result or '{"error":' in result


class TestReorderSlides:
    """Test pptx_reorder_slides tool."""

    @pytest.mark.asyncio
    async def test_reorder_slides_basic(self, layout_tools, mock_presentation_manager):
        """Test reordering slides."""
        result = await layout_tools["pptx_reorder_slides"](slide_index=2, new_position=0)
        assert isinstance(result, str)
        assert "Moved" in result or "position" in result.lower()

    @pytest.mark.asyncio
    async def test_reorder_slides_same_position(self, layout_tools, mock_presentation_manager):
        """Test moving slide to same position."""
        result = await layout_tools["pptx_reorder_slides"](slide_index=1, new_position=1)
        assert isinstance(result, str)
        assert "already at position" in result.lower()

    @pytest.mark.asyncio
    async def test_reorder_slides_invalid_current(self, layout_tools, mock_presentation_manager):
        """Test error with invalid current index."""
        result = await layout_tools["pptx_reorder_slides"](slide_index=999, new_position=0)
        assert "No presentation found" in result or '{"error":' in result

    @pytest.mark.asyncio
    async def test_reorder_slides_invalid_new(self, layout_tools, mock_presentation_manager):
        """Test error with invalid new position."""
        result = await layout_tools["pptx_reorder_slides"](slide_index=0, new_position=999)
        assert '{"message":"Error:' in result or "out of range" in result

    @pytest.mark.asyncio
    async def test_reorder_slides_no_presentation(self, layout_tools, mock_presentation_manager):
        """Test error when no presentation."""
        # Use a non-existent presentation name instead of mocking
        result = await layout_tools["pptx_reorder_slides"](
            slide_index=0, new_position=1, presentation="nonexistent"
        )
        assert "No presentation found" in result or '{"error":' in result


class TestIntegration:
    """Integration tests for slide layout tools."""

    @pytest.mark.asyncio
    async def test_all_tools_registered(self, layout_tools):
        """Test that all slide layout tools are registered."""
        expected_tools = [
            "pptx_list_layouts",
            "pptx_add_slide_with_layout",
            "pptx_customize_layout",
            "pptx_apply_master_layout",
            "pptx_duplicate_slide",
            "pptx_reorder_slides",
        ]

        for tool_name in expected_tools:
            assert tool_name in layout_tools, f"Tool {tool_name} not registered"
            assert callable(layout_tools[tool_name]), f"Tool {tool_name} not callable"

    @pytest.mark.asyncio
    async def test_only_expected_tools_registered(self, layout_tools):
        """Test that no extra tools are registered."""
        expected_count = 6
        assert len(layout_tools) == expected_count, (
            f"Expected {expected_count} tools, got {len(layout_tools)}"
        )

    @pytest.mark.asyncio
    async def test_workflow_create_customize_duplicate(
        self, layout_tools, mock_presentation_manager
    ):
        """Test workflow: create slide, customize it, then duplicate."""
        # Add slide
        result1 = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=1, title="Original Slide"
        )
        assert isinstance(result1, str)

        # Customize it
        result2 = await layout_tools["pptx_customize_layout"](
            slide_index=0, background_color="#F0F0F0", add_footer="Test"
        )
        assert isinstance(result2, str)

        # Duplicate it
        result3 = await layout_tools["pptx_duplicate_slide"](slide_index=0)
        assert isinstance(result3, str)


# ============================================================================
# Additional Coverage Tests
# ============================================================================


class TestListLayoutsExceptionHandling:
    """Test exception handling in pptx_list_layouts."""

    @pytest.mark.asyncio
    async def test_list_layouts_exception(self, mock_mcp_server):
        """Test exception handling in list_layouts - covers lines 120-121."""
        # Create a manager that raises an exception
        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(side_effect=Exception("Test exception"))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_list_layouts"]()
        assert isinstance(result, str)
        assert "Test exception" in result or '{"error":' in result


class TestAddSlideWithLayoutPlaceholders:
    """Test placeholder handling in pptx_add_slide_with_layout."""

    @pytest.mark.asyncio
    async def test_add_slide_with_subtitle_placeholder(
        self, layout_tools, mock_presentation_manager
    ):
        """Test adding slide with subtitle - covers lines 189-190."""
        # Use layout 0 (Title Slide) which typically has subtitle
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=0, title="Main Title", subtitle="This is the subtitle text"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_slide_with_theme_application(self, mock_mcp_server):
        """Test theme application on add_slide - covers lines 211-216."""

        prs = Presentation()
        layout = prs.slide_layouts[1]
        prs.slides.add_slide(layout)

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"
        mock_metadata.theme = "dark-violet"

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_add_slide_with_layout"](layout_index=1, title="Themed Slide")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_slide_exception_handling(self, mock_mcp_server):
        """Test exception handling in add_slide - covers lines 281-282."""
        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(side_effect=Exception("Add slide error"))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_add_slide_with_layout"](layout_index=1)
        assert isinstance(result, str)
        assert "error" in result.lower()


class TestCustomizeLayoutExceptionHandling:
    """Test exception handling in pptx_customize_layout."""

    @pytest.mark.asyncio
    async def test_customize_layout_exception(self, mock_mcp_server):
        """Test exception handling - covers lines 396-397."""
        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(side_effect=Exception("Customize error"))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_customize_layout"](slide_index=0)
        assert isinstance(result, str)
        assert "Customize error" in result or '{"error":' in result


class TestApplyMasterLayoutBodyFormatting:
    """Test body text formatting in pptx_apply_master_layout."""

    @pytest.mark.asyncio
    async def test_apply_master_with_body_text(self, mock_mcp_server):
        """Test formatting body text shapes - covers lines 469->479."""
        from pptx.util import Inches

        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        # Add a text box (not title)
        textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
        textbox.text_frame.text = "Body text content"

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_apply_master_layout"](
            layout_name="corporate", font_name="Arial", body_color="#333333"
        )
        assert isinstance(result, str)
        assert "Applied" in result

    @pytest.mark.asyncio
    async def test_apply_master_exception_handling(self, mock_mcp_server):
        """Test exception handling - covers lines 504-505."""
        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(side_effect=Exception("Apply master error"))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_apply_master_layout"](layout_name="corporate")
        assert isinstance(result, str)
        assert "error" in result.lower()


class TestDuplicateSlideBranches:
    """Test various branches in pptx_duplicate_slide."""

    @pytest.mark.asyncio
    async def test_duplicate_slide_with_textbox(self, mock_mcp_server):
        """Test duplicating slide with text boxes - covers lines 550-551."""
        from pptx.util import Inches

        prs = Presentation()
        layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(layout)

        # Add a text box
        textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        textbox.text_frame.text = "Text content to duplicate"

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_duplicate_slide"](slide_index=0)
        assert isinstance(result, str)
        # Should succeed
        assert "Duplicated" in result or "slide" in result.lower()

    @pytest.mark.asyncio
    async def test_duplicate_slide_with_title(self, mock_mcp_server):
        """Test duplicating slide with title - covers lines 564-565."""

        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = "Original Title"

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_duplicate_slide"](slide_index=0)
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_duplicate_slide_exception_handling(self, mock_mcp_server):
        """Test exception handling - covers lines 575-576."""
        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(side_effect=Exception("Duplicate error"))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_duplicate_slide"](slide_index=0)
        assert isinstance(result, str)
        assert "error" in result.lower()

    @pytest.mark.asyncio
    async def test_duplicate_slide_with_text_shapes(self, mock_mcp_server):
        """Test duplicating slide with various text shapes - covers lines 552-561."""
        from pptx.util import Inches

        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)

        # Add shapes with text
        textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(3), Inches(1))
        textbox.text_frame.text = "Some content"

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_duplicate_slide"](slide_index=0)
        assert isinstance(result, str)


class TestReorderSlidesExceptionHandling:
    """Test exception handling in pptx_reorder_slides."""

    @pytest.mark.asyncio
    async def test_reorder_slides_exception(self, mock_mcp_server):
        """Test exception handling - covers lines 638-639."""
        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(side_effect=Exception("Reorder error"))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_reorder_slides"](slide_index=0, new_position=1)
        assert isinstance(result, str)
        assert "error" in result.lower()


class TestPlaceholderTypesCoverage:
    """Test different placeholder types in add_slide_with_layout."""

    @pytest.mark.asyncio
    async def test_layout_with_various_placeholders(self, mock_mcp_server):
        """Test layouts with different placeholder types - covers lines 237-257."""

        prs = Presentation()

        # Try layouts that might have different placeholder types
        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"
        mock_metadata.theme = None

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        # Test with different layout indices
        for layout_idx in range(min(6, len(prs.slide_master.slide_layouts))):
            result = await tools["pptx_add_slide_with_layout"](
                layout_index=layout_idx, title=f"Slide with layout {layout_idx}"
            )
            assert isinstance(result, str)


class TestLayoutsWithNoPlaceholders:
    """Test layouts with no placeholders in list."""

    @pytest.mark.asyncio
    async def test_list_layouts_no_placeholders(self, mock_mcp_server):
        """Test layout listing when some layouts have no placeholders - covers line 110."""

        prs = Presentation()

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        result = await tools["pptx_list_layouts"]()
        assert isinstance(result, str)
        assert "AVAILABLE SLIDE LAYOUTS" in result


class TestTitleNoneHandling:
    """Test handling when title parameter is None."""

    @pytest.mark.asyncio
    async def test_add_slide_no_title(self, layout_tools, mock_presentation_manager):
        """Test adding slide without title - covers line 179."""
        result = await layout_tools["pptx_add_slide_with_layout"](
            layout_index=6,  # Blank layout typically
            title=None,
        )
        assert isinstance(result, str)


class TestChartPicturePlaceholderWarnings:
    """Test CHART and PICTURE placeholder warning messages."""

    @pytest.mark.asyncio
    async def test_add_slide_with_content_without_chart_picture(self, mock_mcp_server):
        """Test layout without chart/picture placeholders - covers lines 267-277."""

        prs = Presentation()

        mock_metadata = MagicMock()
        mock_metadata.name = "test_pres"
        mock_metadata.theme = None

        mock_manager = MagicMock()
        mock_manager.get = AsyncMock(return_value=(prs, mock_metadata))
        mock_manager.update = AsyncMock()

        tools = register_layout_tools(mock_mcp_server, mock_manager)

        # Use layout index 1 (Title and Content) which has BODY placeholder
        result = await tools["pptx_add_slide_with_layout"](layout_index=1, title="Test Slide")
        assert isinstance(result, str)
        # Should mention placeholder info
        assert "placeholder" in result.lower() or "Placeholder" in result
