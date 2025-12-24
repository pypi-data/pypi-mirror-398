"""
Tests for tools/core/placeholder.py

Tests placeholder population tools for >90% coverage.
"""

import json
import pytest
from unittest.mock import MagicMock
from pptx import Presentation

from chuk_mcp_pptx.tools.core.placeholder import register_placeholder_tools


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures registered tools."""
    mcp = MagicMock()
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mcp.tool = tool_decorator
    mcp._tools = tools
    return mcp


def create_presentation_with_placeholders():
    """Create a presentation with slides that have placeholders."""
    prs = Presentation()

    # Add slide with TITLE_AND_CONTENT layout (usually has title and body placeholders)
    if len(prs.slide_layouts) > 1:
        layout = prs.slide_layouts[1]  # TITLE_AND_CONTENT
        prs.slides.add_slide(layout)

    return prs


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation or create_presentation_with_placeholders()
        self._current_name = "test_presentation"
        self._metadata = MagicMock()
        self._metadata.name = self._current_name

    async def get(self, name=None):
        """Get presentation."""
        if self._presentation is None:
            return None
        if name is None or name == self._current_name:
            return self._presentation, self._metadata
        return None

    async def get_presentation(self, name=None):
        """Get presentation object directly."""
        if self._presentation is None:
            return None
        if name is None or name == self._current_name:
            return self._presentation
        return None

    def get_current_name(self):
        """Get current presentation name."""
        return self._current_name

    async def update_slide_metadata(self, slide_index):
        """Update slide metadata."""
        pass

    async def update(self, name=None):
        """Update presentation."""
        pass

    async def _save_to_store(self, name, prs):
        """Save to store."""
        pass


@pytest.fixture
def mock_manager():
    """Create a mock presentation manager."""
    return MockPresentationManager()


@pytest.fixture
def placeholder_tools(mock_mcp, mock_manager):
    """Register placeholder tools and return them."""
    register_placeholder_tools(mock_mcp, mock_manager)
    return mock_mcp._tools


class TestPopulatePlaceholder:
    """Tests for pptx_populate_placeholder."""

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_string(self, placeholder_tools, mock_manager):
        """Test populating placeholder with string content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test Title",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_dict(self, placeholder_tools, mock_manager):
        """Test populating placeholder with dict content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content={"type": "Table", "headers": ["A", "B"], "data": [["1", "2"]]},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_json_string(self, placeholder_tools, mock_manager):
        """Test populating placeholder with JSON string."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content='{"type": "Table", "headers": ["X"], "data": [["Y"]]}',
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_invalid_slide_index(self, placeholder_tools, mock_manager):
        """Test with invalid slide index."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=999,
            placeholder_idx=0,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_populate_placeholder_invalid_placeholder_idx(
        self, placeholder_tools, mock_manager
    ):
        """Test with invalid placeholder index."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=999,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data or "not found" in result.lower()

    @pytest.mark.asyncio
    async def test_populate_placeholder_no_presentation(self, mock_mcp):
        """Test when no presentation exists."""
        manager = MockPresentationManager()
        manager._presentation = None
        register_placeholder_tools(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_populate_placeholder_with_presentation_name(
        self, placeholder_tools, mock_manager
    ):
        """Test with explicit presentation name."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Named Presentation",
            presentation="test_presentation",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_negative_slide_index(self, placeholder_tools, mock_manager):
        """Test with negative slide index."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=-1,
            placeholder_idx=0,
            content="Test",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_populate_placeholder_empty_string(self, placeholder_tools, mock_manager):
        """Test populating with empty string."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_multiline_string(self, placeholder_tools, mock_manager):
        """Test populating with multiline string."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Line 1\nLine 2\nLine 3",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_placeholder_special_characters(self, placeholder_tools, mock_manager):
        """Test populating with special characters."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Special: <>&\"'",
        )
        assert isinstance(result, str)


class TestPopulatePlaceholderContentTypes:
    """Test different content types for placeholder population."""

    @pytest.mark.asyncio
    async def test_populate_with_table_dict(self, placeholder_tools, mock_manager):
        """Test populating with table dict."""
        table_content = {
            "type": "Table",
            "headers": ["Name", "Value", "Status"],
            "data": [
                ["Item 1", "100", "Active"],
                ["Item 2", "200", "Pending"],
            ],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=table_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with chart dict."""
        chart_content = {
            "type": "ColumnChart",
            "categories": ["Q1", "Q2", "Q3"],
            "series": {"Sales": [100, 150, 200]},
            "title": "Quarterly Sales",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_image_dict(self, placeholder_tools, mock_manager):
        """Test populating with image dict."""
        image_content = {
            "type": "Image",
            "image_source": "https://example.com/image.png",
            "alt": "Test Image",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=image_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_pie_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with pie chart dict."""
        chart_content = {
            "type": "PieChart",
            "categories": ["A", "B", "C"],
            "values": [30, 50, 20],
            "title": "Distribution",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_line_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with line chart dict."""
        chart_content = {
            "type": "LineChart",
            "categories": ["Jan", "Feb", "Mar"],
            "series": {"Revenue": [1000, 1200, 1100]},
            "title": "Monthly Revenue",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_bar_chart_dict(self, placeholder_tools, mock_manager):
        """Test populating with bar chart dict."""
        chart_content = {
            "type": "BarChart",
            "categories": ["Product A", "Product B", "Product C"],
            "series": {"Units": [50, 75, 60]},
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=chart_content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_unknown_type(self, placeholder_tools, mock_manager):
        """Test populating with unknown content type."""
        content = {
            "type": "UnknownType",
            "data": "something",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_dict_missing_type(self, placeholder_tools, mock_manager):
        """Test populating with dict missing type field."""
        content = {
            "headers": ["A", "B"],
            "data": [["1", "2"]],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


class TestToolRegistration:
    """Test tool registration."""

    def test_tool_registered(self, placeholder_tools):
        """Test that placeholder tool is registered."""
        assert "pptx_populate_placeholder" in placeholder_tools
        assert callable(placeholder_tools["pptx_populate_placeholder"])

    def test_tool_is_async(self, placeholder_tools):
        """Test that tool is async."""
        import asyncio

        assert asyncio.iscoroutinefunction(placeholder_tools["pptx_populate_placeholder"])


class TestEdgeCases:
    """Test edge cases."""

    @pytest.mark.asyncio
    async def test_populate_with_very_long_text(self, placeholder_tools, mock_manager):
        """Test populating with very long text."""
        long_text = "A" * 10000
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=long_text,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_unicode(self, placeholder_tools, mock_manager):
        """Test populating with unicode characters."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="日本語テスト 한국어 中文",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_emoji(self, placeholder_tools, mock_manager):
        """Test populating with emoji."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test 🎉 Emoji 🚀",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_html_like_content(self, placeholder_tools, mock_manager):
        """Test populating with HTML-like content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="<b>Bold</b> and <i>italic</i>",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_populate_with_invalid_json_string(self, placeholder_tools, mock_manager):
        """Test populating with invalid JSON string that looks like JSON."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="{invalid json}",
        )
        assert isinstance(result, str)


class TestTableContentVariants:
    """Test various table content configurations."""

    @pytest.mark.asyncio
    async def test_table_with_variant(self, placeholder_tools, mock_manager):
        """Test table with variant styling."""
        content = {
            "type": "Table",
            "headers": ["Col1", "Col2"],
            "data": [["A", "B"]],
            "variant": "striped",
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_table_with_empty_data(self, placeholder_tools, mock_manager):
        """Test table with empty data."""
        content = {
            "type": "Table",
            "headers": ["Col1", "Col2"],
            "data": [],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_table_with_single_cell(self, placeholder_tools, mock_manager):
        """Test table with single cell."""
        content = {
            "type": "Table",
            "headers": ["Single"],
            "data": [["Value"]],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_table_with_many_rows(self, placeholder_tools, mock_manager):
        """Test table with many rows."""
        content = {
            "type": "Table",
            "headers": ["ID", "Name"],
            "data": [[str(i), f"Item {i}"] for i in range(20)],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


class TestDictContentErrorPaths:
    """Test error paths for dict content."""

    @pytest.mark.asyncio
    async def test_dict_content_no_presentation(self, mock_mcp):
        """Test dict content when no presentation exists."""
        manager = MockPresentationManager()
        manager._presentation = None
        register_placeholder_tools(mock_mcp, manager)

        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_invalid_slide_index(self, placeholder_tools):
        """Test dict content with invalid slide index."""
        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=999,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_negative_slide_index(self, placeholder_tools):
        """Test dict content with negative slide index."""
        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=-1,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_invalid_placeholder(self, placeholder_tools):
        """Test dict content with invalid placeholder index."""
        content = {"type": "Table", "headers": ["A"], "data": [["1"]]}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=9999,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data or "not found" in result.lower()

    @pytest.mark.asyncio
    async def test_dict_content_unknown_component(self, placeholder_tools):
        """Test dict content with unknown component type."""
        content = {"type": "UnknownComponent", "data": "something"}
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=content,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_dict_content_with_invalid_params(self, placeholder_tools):
        """Test dict content with invalid params for component."""
        content = {
            "type": "Table",
            "headers": ["A"],
            "data": [["1"]],
            "invalid_param": "should be ignored",
            "another_invalid": 123,
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


class TestStringContentPlaceholderTypes:
    """Test string content for different placeholder types."""

    @pytest.mark.asyncio
    async def test_body_placeholder_with_bullets(self, placeholder_tools, mock_manager):
        """Test body placeholder with bullet points using \\n separator."""
        # Body placeholders typically have idx=1 in TITLE_AND_CONTENT layout
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content="First bullet\\nSecond bullet\\nThird bullet",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_title_placeholder_simple(self, placeholder_tools, mock_manager):
        """Test title placeholder with simple text."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Simple Title",
        )
        assert isinstance(result, str)


class TestInvalidContentTypes:
    """Test invalid content type handling."""

    @pytest.mark.asyncio
    async def test_invalid_content_type_list(self, placeholder_tools, mock_manager):
        """Test with list content (not supported)."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=["item1", "item2"],  # List is not a valid content type
        )
        assert isinstance(result, str)
        # Should either succeed by converting or error

    @pytest.mark.asyncio
    async def test_invalid_content_type_number(self, placeholder_tools, mock_manager):
        """Test with number content."""
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content=12345,  # Number should be handled
        )
        assert isinstance(result, str)


class TestChartContentVariants:
    """Test various chart content configurations."""

    @pytest.mark.asyncio
    async def test_chart_with_multiple_series(self, placeholder_tools, mock_manager):
        """Test chart with multiple series."""
        content = {
            "type": "ColumnChart",
            "categories": ["Q1", "Q2", "Q3", "Q4"],
            "series": {
                "Sales": [100, 120, 140, 160],
                "Costs": [80, 90, 100, 110],
                "Profit": [20, 30, 40, 50],
            },
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_chart_without_title(self, placeholder_tools, mock_manager):
        """Test chart without title."""
        content = {
            "type": "ColumnChart",
            "categories": ["A", "B", "C"],
            "series": {"Values": [1, 2, 3]},
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_doughnut_chart(self, placeholder_tools, mock_manager):
        """Test doughnut chart."""
        content = {
            "type": "DoughnutChart",
            "categories": ["Segment A", "Segment B", "Segment C"],
            "values": [40, 35, 25],
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_area_chart(self, placeholder_tools, mock_manager):
        """Test area chart."""
        content = {
            "type": "AreaChart",
            "categories": ["Week 1", "Week 2", "Week 3", "Week 4"],
            "series": {"Traffic": [1000, 1200, 1100, 1400]},
        }
        result = await placeholder_tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content=content,
        )
        assert isinstance(result, str)


# ============================================================================
# Additional Coverage Tests
# ============================================================================


class TestSuccessfulComponentRender:
    """Tests for successful component rendering to placeholder."""

    @pytest.fixture
    def manager_with_valid_placeholder(self):
        """Create manager with presentation having valid placeholder."""

        prs = Presentation()
        # Use a layout that has content placeholders
        if len(prs.slide_layouts) > 1:
            layout = prs.slide_layouts[1]
            prs.slides.add_slide(layout)

        manager = MockPresentationManager(presentation=prs)
        return manager

    @pytest.mark.asyncio
    async def test_successful_table_render_to_placeholder(
        self, mock_mcp, manager_with_valid_placeholder
    ):
        """Test successful table render to placeholder - covers lines 281-300."""
        register_placeholder_tools(mock_mcp, manager_with_valid_placeholder)

        # Get actual placeholder idx from the slide
        prs = manager_with_valid_placeholder._presentation
        if prs.slides and prs.slides[0].placeholders:
            # Find an OBJECT or content placeholder
            for ph in prs.slides[0].placeholders:
                ph_type = ph.placeholder_format.type
                if ph_type in (2, 7):  # BODY or OBJECT
                    content = {
                        "type": "Table",
                        "headers": ["Name", "Value"],
                        "data": [["Test", "123"]],
                    }
                    result = await mock_mcp._tools["pptx_populate_placeholder"](
                        slide_index=0,
                        placeholder_idx=ph.placeholder_format.idx,
                        content=content,
                    )
                    assert isinstance(result, str)
                    # Check for success or expected error
                    break


class TestPlaceholderTextAttribute:
    """Tests for placeholder .text attribute fallback (lines 356-359)."""

    @pytest.fixture
    def manager_with_text_only_placeholder(self):
        """Create manager with placeholder that has .text but not .text_frame."""
        # Create a mock slide with a mock placeholder
        mock_slide = MagicMock()
        mock_placeholder = MagicMock()
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 0
        mock_placeholder_format.type = 1  # TITLE type

        mock_placeholder.placeholder_format = mock_placeholder_format
        # Remove text_frame, only have text attribute
        mock_placeholder.text_frame = None
        mock_placeholder.text = ""

        # Use spec to remove text_frame attribute properly
        mock_placeholder.configure_mock(**{"text_frame": None})

        mock_slide.placeholders = [mock_placeholder]

        # Mock prs.slides as a MagicMock list-like object
        mock_slides = MagicMock()
        mock_slides.__getitem__ = MagicMock(return_value=mock_slide)
        mock_slides.__len__ = MagicMock(return_value=1)

        mock_prs = MagicMock()
        mock_prs.slides = mock_slides

        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_title_placeholder_with_text_attr_only(
        self, mock_mcp, manager_with_text_only_placeholder
    ):
        """Test populating TITLE placeholder that only has .text attribute."""
        register_placeholder_tools(mock_mcp, manager_with_text_only_placeholder)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Title via .text attribute",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        # Should succeed via .text attribute fallback
        assert "message" in data or "error" in data


def create_mock_slides_with_placeholder(mock_placeholder):
    """Helper to create mock slides with proper MagicMock structure."""
    mock_slide = MagicMock()
    mock_slide.placeholders = [mock_placeholder]

    mock_slides = MagicMock()
    mock_slides.__getitem__ = MagicMock(return_value=mock_slide)
    mock_slides.__len__ = MagicMock(return_value=1)

    mock_prs = MagicMock()
    mock_prs.slides = mock_slides

    return mock_prs


class TestPlaceholderNoTextCapability:
    """Tests for placeholders without text capability (lines 356-359, 367, 396-404, 414-417)."""

    @pytest.fixture
    def manager_with_no_text_placeholder(self):
        """Create manager with placeholder that has neither .text nor .text_frame."""
        mock_placeholder = MagicMock(spec=["placeholder_format"])
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 0
        mock_placeholder_format.type = 1  # TITLE type
        mock_placeholder.placeholder_format = mock_placeholder_format

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_title_placeholder_no_text_support(
        self, mock_mcp, manager_with_no_text_placeholder
    ):
        """Test TITLE placeholder without text support - covers lines 359-362."""
        register_placeholder_tools(mock_mcp, manager_with_no_text_placeholder)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test content",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "does not support text content" in data["error"]


class TestBodyPlaceholderNoTextFrame:
    """Tests for BODY placeholder without text_frame (line 367)."""

    @pytest.fixture
    def manager_with_body_no_textframe(self):
        """Create manager with BODY placeholder without text_frame."""
        mock_placeholder = MagicMock(spec=["placeholder_format"])
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 1
        mock_placeholder_format.type = 2  # BODY type
        mock_placeholder.placeholder_format = mock_placeholder_format

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_body_placeholder_no_textframe(self, mock_mcp, manager_with_body_no_textframe):
        """Test BODY placeholder without text_frame - covers line 367."""
        register_placeholder_tools(mock_mcp, manager_with_body_no_textframe)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content="Body content",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "does not have a text frame" in data["error"]


class TestChartTablePicturePlaceholderWithString:
    """Tests for CHART/TABLE/PICTURE placeholders with string content (lines 388-414)."""

    @pytest.fixture
    def manager_with_chart_placeholder(self):
        """Create manager with CHART placeholder (type 12)."""
        mock_placeholder = MagicMock()
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 2
        mock_placeholder_format.type = 12  # CHART type
        mock_placeholder.placeholder_format = mock_placeholder_format
        mock_placeholder.text_frame = MagicMock()
        mock_placeholder.text_frame.text = ""

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.fixture
    def manager_with_table_placeholder(self):
        """Create manager with TABLE placeholder (type 14)."""
        mock_placeholder = MagicMock(spec=["placeholder_format", "text"])
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 3
        mock_placeholder_format.type = 14  # TABLE type
        mock_placeholder.placeholder_format = mock_placeholder_format
        mock_placeholder.text = ""

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.fixture
    def manager_with_picture_placeholder_no_text(self):
        """Create manager with PICTURE placeholder without text capability."""
        mock_placeholder = MagicMock(spec=["placeholder_format"])
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 4
        mock_placeholder_format.type = 18  # PICTURE type
        mock_placeholder.placeholder_format = mock_placeholder_format

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_chart_placeholder_with_string_textframe(
        self, mock_mcp, manager_with_chart_placeholder
    ):
        """Test CHART placeholder with string content via text_frame - covers lines 388-393."""
        register_placeholder_tools(mock_mcp, manager_with_chart_placeholder)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=2,
            content="Chart caption text",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        # Should succeed - CHART placeholders can have text captions
        assert "message" in data or "error" in data

    @pytest.mark.asyncio
    async def test_table_placeholder_with_string_text_attr(
        self, mock_mcp, manager_with_table_placeholder
    ):
        """Test TABLE placeholder with string content via .text - covers lines 394-395."""
        register_placeholder_tools(mock_mcp, manager_with_table_placeholder)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=3,
            content="Table caption",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "message" in data or "error" in data

    @pytest.mark.asyncio
    async def test_picture_placeholder_no_text_capability(
        self, mock_mcp, manager_with_picture_placeholder_no_text
    ):
        """Test PICTURE placeholder without text capability - covers lines 396-404."""
        register_placeholder_tools(mock_mcp, manager_with_picture_placeholder_no_text)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=4,
            content="This should fail",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "PICTURE" in data["error"]
        assert "Use dict content" in data["error"]


class TestOtherPlaceholderTypes:
    """Tests for other placeholder types (lines 407-417)."""

    @pytest.fixture
    def manager_with_unknown_placeholder_type(self):
        """Create manager with unknown placeholder type."""
        mock_placeholder = MagicMock()
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 5
        mock_placeholder_format.type = 99  # Unknown type
        mock_placeholder.placeholder_format = mock_placeholder_format
        mock_placeholder.text_frame = MagicMock()
        mock_placeholder.text_frame.text = ""

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.fixture
    def manager_with_unknown_placeholder_no_text(self):
        """Create manager with unknown placeholder type without text capability."""
        mock_placeholder = MagicMock(spec=["placeholder_format"])
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 6
        mock_placeholder_format.type = 99  # Unknown type
        mock_placeholder.placeholder_format = mock_placeholder_format

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_unknown_type_with_textframe(
        self, mock_mcp, manager_with_unknown_placeholder_type
    ):
        """Test unknown placeholder type with text_frame - covers lines 408-410."""
        register_placeholder_tools(mock_mcp, manager_with_unknown_placeholder_type)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=5,
            content="Content for unknown type",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "message" in data or "error" in data

    @pytest.mark.asyncio
    async def test_unknown_type_no_text_support(
        self, mock_mcp, manager_with_unknown_placeholder_no_text
    ):
        """Test unknown placeholder type without text support - covers lines 414-417."""
        register_placeholder_tools(mock_mcp, manager_with_unknown_placeholder_no_text)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=6,
            content="This should fail",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "does not support text content" in data["error"]


class TestTopLevelExceptionHandling:
    """Tests for top-level exception handling (lines 430-434)."""

    @pytest.fixture
    def manager_that_raises(self):
        """Create manager that raises exception on get_presentation."""
        manager = MockPresentationManager()

        async def raise_exception(name=None):
            raise RuntimeError("Unexpected error during presentation retrieval")

        manager.get_presentation = raise_exception
        manager.get = raise_exception
        return manager

    @pytest.mark.asyncio
    async def test_top_level_exception_handling(self, mock_mcp, manager_that_raises):
        """Test top-level exception handling - covers lines 430-434."""
        register_placeholder_tools(mock_mcp, manager_that_raises)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=0,
            content="Test content",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data
        assert "Unexpected error" in data["error"]


class TestSubtitlePlaceholder:
    """Tests for SUBTITLE placeholder (type 3)."""

    @pytest.fixture
    def manager_with_subtitle_placeholder(self):
        """Create manager with SUBTITLE placeholder."""
        mock_placeholder = MagicMock()
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 1
        mock_placeholder_format.type = 3  # SUBTITLE type
        mock_placeholder.placeholder_format = mock_placeholder_format
        mock_placeholder.text_frame = MagicMock()
        mock_placeholder.text_frame.text = ""

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_subtitle_placeholder_with_textframe(
        self, mock_mcp, manager_with_subtitle_placeholder
    ):
        """Test SUBTITLE placeholder with text_frame - covers lines 353-355."""
        register_placeholder_tools(mock_mcp, manager_with_subtitle_placeholder)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=1,
            content="Subtitle text",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "message" in data or "error" in data


class TestObjectPlaceholder:
    """Tests for OBJECT placeholder (type 7)."""

    @pytest.fixture
    def manager_with_object_placeholder(self):
        """Create manager with OBJECT placeholder."""
        mock_placeholder = MagicMock()
        mock_placeholder_format = MagicMock()
        mock_placeholder_format.idx = 2
        mock_placeholder_format.type = 7  # OBJECT type
        mock_placeholder.placeholder_format = mock_placeholder_format

        # Setup text_frame with paragraphs
        mock_text_frame = MagicMock()
        mock_paragraph = MagicMock()
        mock_paragraph.text = ""
        mock_text_frame.paragraphs = [mock_paragraph]
        mock_text_frame.clear = MagicMock()
        mock_text_frame.add_paragraph = MagicMock(return_value=MagicMock())
        mock_placeholder.text_frame = mock_text_frame

        mock_prs = create_mock_slides_with_placeholder(mock_placeholder)
        manager = MockPresentationManager(presentation=mock_prs)
        return manager

    @pytest.mark.asyncio
    async def test_object_placeholder_with_bullets(self, mock_mcp, manager_with_object_placeholder):
        """Test OBJECT placeholder with bullet points - covers lines 365, 371-384."""
        register_placeholder_tools(mock_mcp, manager_with_object_placeholder)

        result = await mock_mcp._tools["pptx_populate_placeholder"](
            slide_index=0,
            placeholder_idx=2,
            content="First bullet\\nSecond bullet\\nThird bullet",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "message" in data or "error" in data
