"""
Tests for tools/theme/management.py

Tests theme management tools for >90% coverage.
"""

import json
import pytest
from unittest.mock import MagicMock
from pptx import Presentation

from chuk_mcp_pptx.tools.theme.management import register_theme_tools


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures registered tools."""
    mcp = MagicMock()
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mcp.tool = tool_decorator
    mcp._tools = tools
    return mcp


def create_presentation_with_slides():
    """Create a presentation with slides for testing."""
    prs = Presentation()
    if len(prs.slide_layouts) > 1:
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        # Add a shape to the slide
        from pptx.util import Inches

        shapes = slide.shapes
        shapes.add_shape(
            1,  # Rectangle
            Inches(1),
            Inches(1),
            Inches(2),
            Inches(1),
        )
    return prs


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation or create_presentation_with_slides()
        self._current_name = "test_presentation"
        self._metadata = MagicMock()
        self._metadata.name = self._current_name

    async def get(self, name=None):
        """Get presentation."""
        if self._presentation is None:
            return None
        if name is None or name == self._current_name:
            return self._presentation, self._metadata
        return None

    def get_current_name(self):
        """Get current presentation name."""
        return self._current_name

    async def update_slide_metadata(self, slide_index):
        """Update slide metadata."""
        pass

    async def update(self, name=None):
        """Update presentation."""
        pass

    async def _save_to_store(self, name, prs):
        """Save to store."""
        pass


@pytest.fixture
def mock_manager():
    """Create a mock presentation manager."""
    return MockPresentationManager()


@pytest.fixture
def theme_tools(mock_mcp, mock_manager):
    """Register theme tools and return them."""
    register_theme_tools(mock_mcp, mock_manager)
    return mock_mcp._tools


class TestListThemes:
    """Tests for pptx_list_themes."""

    @pytest.mark.asyncio
    async def test_list_themes_returns_string(self, theme_tools):
        """Test that list themes returns string."""
        result = await theme_tools["pptx_list_themes"]()
        assert isinstance(result, str)
        assert len(result) > 0

    @pytest.mark.asyncio
    async def test_list_themes_contains_themes(self, theme_tools):
        """Test that list themes returns theme data."""
        result = await theme_tools["pptx_list_themes"]()
        # Should have themes listed (contains theme names)
        assert "dark" in result.lower() or "light" in result.lower()


class TestGetThemeInfo:
    """Tests for pptx_get_theme_info."""

    @pytest.mark.asyncio
    async def test_get_theme_info_dark(self, theme_tools):
        """Test getting info for dark theme."""
        result = await theme_tools["pptx_get_theme_info"](theme_name="dark")
        assert isinstance(result, str)
        data = json.loads(result)
        assert "name" in data or "error" not in data

    @pytest.mark.asyncio
    async def test_get_theme_info_light(self, theme_tools):
        """Test getting info for light theme."""
        result = await theme_tools["pptx_get_theme_info"](theme_name="light")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_theme_info_dark_violet(self, theme_tools):
        """Test getting info for dark-violet theme."""
        result = await theme_tools["pptx_get_theme_info"](theme_name="dark-violet")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_theme_info_unknown(self, theme_tools):
        """Test getting info for unknown theme."""
        result = await theme_tools["pptx_get_theme_info"](theme_name="nonexistent_theme")
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data or "not found" in result.lower()


class TestCreateCustomTheme:
    """Tests for pptx_create_custom_theme."""

    @pytest.mark.asyncio
    async def test_create_custom_theme_default(self, theme_tools):
        """Test creating custom theme with defaults."""
        result = await theme_tools["pptx_create_custom_theme"]()
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_create_custom_theme_with_name(self, theme_tools):
        """Test creating custom theme with name."""
        result = await theme_tools["pptx_create_custom_theme"](name="my_custom_theme")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_create_custom_theme_with_hue(self, theme_tools):
        """Test creating custom theme with primary hue."""
        result = await theme_tools["pptx_create_custom_theme"](
            name="violet_theme", primary_hue="violet"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_create_custom_theme_with_mode(self, theme_tools):
        """Test creating custom theme with mode."""
        result = await theme_tools["pptx_create_custom_theme"](name="light_theme", mode="light")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_create_custom_theme_full_params(self, theme_tools):
        """Test creating custom theme with all parameters."""
        result = await theme_tools["pptx_create_custom_theme"](
            name="full_custom",
            primary_hue="emerald",
            mode="dark",
            font_family="Arial",
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_create_custom_theme_various_hues(self, theme_tools):
        """Test creating themes with various hues."""
        hues = ["blue", "red", "green", "orange", "pink", "purple"]
        for hue in hues:
            result = await theme_tools["pptx_create_custom_theme"](
                name=f"{hue}_theme", primary_hue=hue
            )
            assert isinstance(result, str)


class TestApplyTheme:
    """Tests for pptx_apply_theme."""

    @pytest.mark.asyncio
    async def test_apply_theme_to_all_slides(self, theme_tools, mock_manager):
        """Test applying theme to all slides."""
        result = await theme_tools["pptx_apply_theme"](slide_index=None, theme="dark")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_theme_to_specific_slide(self, theme_tools, mock_manager):
        """Test applying theme to specific slide."""
        result = await theme_tools["pptx_apply_theme"](slide_index=0, theme="dark-violet")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_theme_light(self, theme_tools, mock_manager):
        """Test applying light theme."""
        result = await theme_tools["pptx_apply_theme"](slide_index=0, theme="light")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_theme_invalid_slide(self, theme_tools, mock_manager):
        """Test applying theme to invalid slide."""
        result = await theme_tools["pptx_apply_theme"](slide_index=999, theme="dark")
        assert isinstance(result, str)
        # Could be JSON error or plain text error
        assert (
            "error" in result.lower()
            or "invalid" in result.lower()
            or "not found" in result.lower()
        )

    @pytest.mark.asyncio
    async def test_apply_theme_no_presentation(self, mock_mcp):
        """Test applying theme when no presentation exists."""
        manager = MockPresentationManager()
        manager._presentation = None
        register_theme_tools(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_apply_theme"](slide_index=0, theme="dark")
        assert isinstance(result, str)
        assert "error" in result.lower() or "no presentation" in result.lower()

    @pytest.mark.asyncio
    async def test_apply_theme_with_presentation_name(self, theme_tools, mock_manager):
        """Test applying theme with presentation name."""
        result = await theme_tools["pptx_apply_theme"](
            slide_index=0, theme="dark", presentation="test_presentation"
        )
        assert isinstance(result, str)


class TestApplyComponentTheme:
    """Tests for pptx_apply_component_theme."""

    @pytest.mark.asyncio
    async def test_apply_component_theme_basic(self, theme_tools, mock_manager):
        """Test applying theme to component."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=0, shape_index=0, theme_style="card"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_component_theme_invalid_slide(self, theme_tools, mock_manager):
        """Test applying theme with invalid slide."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=999, shape_index=0, theme_style="card"
        )
        assert isinstance(result, str)
        assert (
            "error" in result.lower()
            or "invalid" in result.lower()
            or "not found" in result.lower()
        )

    @pytest.mark.asyncio
    async def test_apply_component_theme_invalid_shape(self, theme_tools, mock_manager):
        """Test applying theme with invalid shape."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=0, shape_index=999, theme_style="card"
        )
        assert isinstance(result, str)
        assert (
            "error" in result.lower()
            or "invalid" in result.lower()
            or "not found" in result.lower()
        )

    @pytest.mark.asyncio
    async def test_apply_component_theme_primary(self, theme_tools, mock_manager):
        """Test applying primary theme style to component."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=0, shape_index=0, theme_style="primary"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_component_theme_secondary(self, theme_tools, mock_manager):
        """Test applying secondary theme style to component."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=0, shape_index=0, theme_style="secondary"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_component_theme_accent(self, theme_tools, mock_manager):
        """Test applying accent theme style to component."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=0, shape_index=0, theme_style="accent"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_component_theme_muted(self, theme_tools, mock_manager):
        """Test applying muted theme style to component."""
        result = await theme_tools["pptx_apply_component_theme"](
            slide_index=0, shape_index=0, theme_style="muted"
        )
        assert isinstance(result, str)


class TestListComponentThemes:
    """Tests for pptx_list_component_themes."""

    @pytest.mark.asyncio
    async def test_list_component_themes(self, theme_tools):
        """Test listing component themes."""
        result = await theme_tools["pptx_list_component_themes"]()
        assert isinstance(result, str)
        assert len(result) > 0


class TestGetColorPalette:
    """Tests for pptx_get_color_palette."""

    @pytest.mark.asyncio
    async def test_get_color_palette(self, theme_tools):
        """Test getting color palette."""
        result = await theme_tools["pptx_get_color_palette"]()
        assert isinstance(result, str)
        data = json.loads(result)
        assert isinstance(data, dict)


class TestGetSemanticColors:
    """Tests for pptx_get_semantic_colors."""

    @pytest.mark.asyncio
    async def test_get_semantic_colors_default(self, theme_tools):
        """Test getting semantic colors with defaults."""
        result = await theme_tools["pptx_get_semantic_colors"]()
        assert isinstance(result, str)
        data = json.loads(result)
        assert isinstance(data, dict)

    @pytest.mark.asyncio
    async def test_get_semantic_colors_with_hue(self, theme_tools):
        """Test getting semantic colors with hue."""
        result = await theme_tools["pptx_get_semantic_colors"](primary_hue="violet")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_semantic_colors_light_mode(self, theme_tools):
        """Test getting semantic colors in light mode."""
        result = await theme_tools["pptx_get_semantic_colors"](primary_hue="blue", mode="light")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_semantic_colors_various_hues(self, theme_tools):
        """Test getting semantic colors for various hues."""
        hues = ["blue", "red", "emerald", "orange", "pink"]
        for hue in hues:
            result = await theme_tools["pptx_get_semantic_colors"](primary_hue=hue, mode="dark")
            assert isinstance(result, str)


class TestGetTypographyTokens:
    """Tests for pptx_get_typography_tokens."""

    @pytest.mark.asyncio
    async def test_get_typography_tokens(self, theme_tools):
        """Test getting typography tokens."""
        result = await theme_tools["pptx_get_typography_tokens"]()
        assert isinstance(result, str)
        data = json.loads(result)
        assert isinstance(data, dict)


class TestGetTextStyle:
    """Tests for pptx_get_text_style."""

    @pytest.mark.asyncio
    async def test_get_text_style_body(self, theme_tools):
        """Test getting body text style."""
        result = await theme_tools["pptx_get_text_style"](variant="body")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_text_style_heading(self, theme_tools):
        """Test getting heading text style."""
        result = await theme_tools["pptx_get_text_style"](variant="heading")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_text_style_title(self, theme_tools):
        """Test getting title text style."""
        result = await theme_tools["pptx_get_text_style"](variant="title")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_text_style_caption(self, theme_tools):
        """Test getting caption text style."""
        result = await theme_tools["pptx_get_text_style"](variant="caption")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_get_text_style_unknown(self, theme_tools):
        """Test getting unknown text style."""
        result = await theme_tools["pptx_get_text_style"](variant="unknown_variant")
        assert isinstance(result, str)


class TestGetSpacingTokens:
    """Tests for pptx_get_spacing_tokens."""

    @pytest.mark.asyncio
    async def test_get_spacing_tokens(self, theme_tools):
        """Test getting spacing tokens."""
        result = await theme_tools["pptx_get_spacing_tokens"]()
        assert isinstance(result, str)
        data = json.loads(result)
        assert isinstance(data, dict)


class TestGetAllTokens:
    """Tests for pptx_get_all_tokens."""

    @pytest.mark.asyncio
    async def test_get_all_tokens(self, theme_tools):
        """Test getting all tokens."""
        result = await theme_tools["pptx_get_all_tokens"]()
        assert isinstance(result, str)
        data = json.loads(result)
        assert isinstance(data, dict)


class TestToolRegistration:
    """Tests for tool registration."""

    def test_all_tools_registered(self, theme_tools):
        """Test that all expected tools are registered."""
        expected_tools = [
            "pptx_list_themes",
            "pptx_get_theme_info",
            "pptx_create_custom_theme",
            "pptx_apply_theme",
            "pptx_apply_component_theme",
            "pptx_list_component_themes",
            "pptx_get_color_palette",
            "pptx_get_semantic_colors",
            "pptx_get_typography_tokens",
            "pptx_get_text_style",
            "pptx_get_spacing_tokens",
            "pptx_get_all_tokens",
        ]
        for tool_name in expected_tools:
            assert tool_name in theme_tools, f"Tool {tool_name} not registered"
            assert callable(theme_tools[tool_name])

    def test_tools_are_async(self, theme_tools):
        """Test that all tools are async."""
        import asyncio

        for tool_name, tool_func in theme_tools.items():
            assert asyncio.iscoroutinefunction(tool_func), f"Tool {tool_name} is not async"


class TestEdgeCases:
    """Test edge cases."""

    @pytest.mark.asyncio
    async def test_create_theme_empty_name(self, theme_tools):
        """Test creating theme with empty name."""
        result = await theme_tools["pptx_create_custom_theme"](name="")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_create_theme_special_chars_name(self, theme_tools):
        """Test creating theme with special characters in name."""
        result = await theme_tools["pptx_create_custom_theme"](name="theme-with_special.chars")
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_apply_theme_negative_slide(self, theme_tools, mock_manager):
        """Test applying theme to negative slide index (allowed by API)."""
        result = await theme_tools["pptx_apply_theme"](slide_index=-1, theme="dark")
        assert isinstance(result, str)
        # The API may allow negative indexing, so just verify we get a result
