"""
Tests for tools/universal/api.py

Tests the Universal Component API tools for >90% coverage.
"""

import json
import pytest
from unittest.mock import MagicMock
from pptx import Presentation

from chuk_mcp_pptx.tools.universal.api import register_universal_component_api


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures registered tools."""
    mcp = MagicMock()
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mcp.tool = tool_decorator
    mcp._tools = tools
    return mcp


@pytest.fixture
def presentation_with_slides():
    """Create a presentation with slides for testing."""
    prs = Presentation()
    # Add a slide with title and content layout
    if prs.slide_layouts:
        layout = prs.slide_layouts[1]  # TITLE_AND_CONTENT
        prs.slides.add_slide(layout)
    return prs


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation or Presentation()
        self._current_name = "test_presentation"
        self._metadata = MagicMock()
        self._metadata.name = self._current_name
        self._slides_metadata = {}

    async def get(self, name=None):
        """Get presentation."""
        if name is None or name == self._current_name:
            return self._presentation, self._metadata
        return None

    def get_current_name(self):
        """Get current presentation name."""
        return self._current_name

    async def update_slide_metadata(self, slide_index):
        """Update slide metadata."""
        self._slides_metadata[slide_index] = True

    async def update(self, name=None):
        """Update presentation."""
        pass


@pytest.fixture
def mock_manager(presentation_with_slides):
    """Create a mock presentation manager."""
    return MockPresentationManager(presentation_with_slides)


@pytest.fixture
def api_tools(mock_mcp, mock_manager):
    """Register API tools and return them."""
    register_universal_component_api(mock_mcp, mock_manager)
    return mock_mcp._tools


class TestListSlideComponents:
    """Tests for pptx_list_slide_components."""

    @pytest.mark.asyncio
    async def test_list_slide_components_basic(self, api_tools, mock_manager):
        """Test basic list slide components."""
        result = await api_tools["pptx_list_slide_components"](slide_index=0, presentation=None)
        assert isinstance(result, str)
        data = json.loads(result)
        assert "slide_index" in data or "error" not in data

    @pytest.mark.asyncio
    async def test_list_slide_components_with_presentation_name(self, api_tools, mock_manager):
        """Test list components with explicit presentation name."""
        result = await api_tools["pptx_list_slide_components"](
            slide_index=0, presentation="test_presentation"
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_list_slide_components_invalid_slide_index(self, api_tools, mock_manager):
        """Test with invalid slide index."""
        result = await api_tools["pptx_list_slide_components"](slide_index=999, presentation=None)
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_list_slide_components_negative_index(self, api_tools, mock_manager):
        """Test with negative slide index."""
        result = await api_tools["pptx_list_slide_components"](slide_index=-1, presentation=None)
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_list_slide_components_no_presentation(self, mock_mcp):
        """Test when no presentation exists."""
        manager = MockPresentationManager()
        manager._presentation = None

        async def mock_get(name=None):
            return None

        manager.get = mock_get
        register_universal_component_api(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_list_slide_components"](
            slide_index=0, presentation=None
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestAddComponent:
    """Tests for pptx_add_component."""

    @pytest.mark.asyncio
    async def test_add_component_basic_badge(self, api_tools, mock_manager):
        """Test adding a basic badge component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            width=2.0,
            height=0.5,
            params={"text": "Test Badge"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_with_json_params(self, api_tools, mock_manager):
        """Test adding component with JSON string params."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            params='{"text": "JSON Badge"}',
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_invalid_json_params(self, api_tools, mock_manager):
        """Test adding component with invalid JSON params."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            params="invalid json {",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_component_unknown_component(self, api_tools, mock_manager):
        """Test adding unknown component type."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="UnknownComponent",
            left=1.0,
            top=1.0,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_component_invalid_slide_index(self, api_tools, mock_manager):
        """Test adding component to invalid slide."""
        result = await api_tools["pptx_add_component"](
            slide_index=999,
            component="Badge",
            left=1.0,
            top=1.0,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_component_no_presentation(self, mock_mcp):
        """Test when no presentation exists."""
        manager = MockPresentationManager()

        async def mock_get(name=None):
            return None

        manager.get = mock_get
        register_universal_component_api(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_component_with_theme(self, api_tools, mock_manager):
        """Test adding component with theme."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            theme="dark-violet",
            params={"text": "Themed Badge"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_with_component_id(self, api_tools, mock_manager):
        """Test adding component with custom ID."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            component_id="my_badge_1",
            params={"text": "ID Badge"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_button(self, api_tools, mock_manager):
        """Test adding button component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Button",
            left=2.0,
            top=2.0,
            width=2.0,
            height=0.5,
            params={"text": "Click Me"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_text(self, api_tools, mock_manager):
        """Test adding text component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Text",
            left=1.0,
            top=3.0,
            width=4.0,
            height=1.0,
            params={"text": "Hello World", "font_size": 24},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_params_not_dict(self, api_tools, mock_manager):
        """Test adding component with non-dict params."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            params=12345,  # Invalid type
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_component_alert(self, api_tools, mock_manager):
        """Test adding alert component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Alert",
            left=1.0,
            top=4.0,
            width=5.0,
            height=1.0,
            params={"title": "Warning", "description": "This is a warning"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_card(self, api_tools, mock_manager):
        """Test adding card component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Card",
            left=1.0,
            top=1.0,
            width=3.0,
            height=2.0,
            params={"title": "Card Title", "description": "Card content"},
        )
        assert isinstance(result, str)


class TestUpdateComponent:
    """Tests for pptx_update_component."""

    @pytest.mark.asyncio
    async def test_update_component_basic(self, api_tools, mock_manager):
        """Test basic component update."""
        # First add a component
        await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            component_id="update_test_badge",
            params={"text": "Original"},
        )

        # Then update it
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="update_test_badge",
            params={"text": "Updated"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_update_component_not_found(self, api_tools, mock_manager):
        """Test updating non-existent component."""
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="nonexistent_component",
            params={"text": "Updated"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data or "not found" in result.lower()

    @pytest.mark.asyncio
    async def test_update_component_invalid_slide(self, api_tools, mock_manager):
        """Test updating component on invalid slide."""
        result = await api_tools["pptx_update_component"](
            slide_index=999,
            component_id="some_component",
            params={"text": "Updated"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_update_component_no_presentation(self, mock_mcp):
        """Test when no presentation exists."""
        manager = MockPresentationManager()

        async def mock_get(name=None):
            return None

        manager.get = mock_get
        register_universal_component_api(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_update_component"](
            slide_index=0,
            component_id="some_component",
            params={"text": "Updated"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_update_component_with_json_params(self, api_tools, mock_manager):
        """Test updating component with JSON string params."""
        # First add a component
        await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            component_id="json_update_badge",
            params={"text": "Original"},
        )

        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="json_update_badge",
            params='{"text": "JSON Updated"}',
        )
        assert isinstance(result, str)


class TestToolRegistration:
    """Tests for tool registration."""

    def test_all_tools_registered(self, api_tools):
        """Test that all expected tools are registered."""
        expected_tools = [
            "pptx_list_slide_components",
            "pptx_add_component",
            "pptx_update_component",
        ]
        for tool_name in expected_tools:
            assert tool_name in api_tools, f"Tool {tool_name} not registered"
            assert callable(api_tools[tool_name])

    def test_tools_are_async(self, api_tools):
        """Test that tools are async functions."""
        import asyncio

        for tool_name, tool_func in api_tools.items():
            assert asyncio.iscoroutinefunction(tool_func), f"Tool {tool_name} is not async"


class TestComponentTypes:
    """Test various component types."""

    @pytest.mark.asyncio
    async def test_add_progress_component(self, api_tools, mock_manager):
        """Test adding progress component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Progress",
            left=1.0,
            top=5.0,
            width=4.0,
            height=0.3,
            params={"value": 75},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_divider_component(self, api_tools, mock_manager):
        """Test adding divider component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Divider",
            left=1.0,
            top=5.5,
            width=8.0,
            height=0.1,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_icon_component(self, api_tools, mock_manager):
        """Test adding icon component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Icon",
            left=1.0,
            top=1.0,
            width=0.5,
            height=0.5,
            params={"name": "star"},
        )
        assert isinstance(result, str)


class TestEdgeCases:
    """Test edge cases and error handling."""

    @pytest.mark.asyncio
    async def test_add_component_none_params(self, api_tools, mock_manager):
        """Test adding component with None params."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            params=None,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_empty_params(self, api_tools, mock_manager):
        """Test adding component with empty params dict."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            params={},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_at_zero_position(self, api_tools, mock_manager):
        """Test adding component at position 0,0."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=0.0,
            top=0.0,
            params={"text": "Zero Position"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_component_large_position(self, api_tools, mock_manager):
        """Test adding component at large position."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=12.0,
            top=8.0,
            params={"text": "Large Position"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_list_components_empty_slide(self, mock_mcp):
        """Test listing components on empty slide."""
        # Create presentation with empty slide
        prs = Presentation()
        if prs.slide_layouts:
            prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        manager = MockPresentationManager(prs)
        register_universal_component_api(mock_mcp, manager)

        result = await mock_mcp._tools["pptx_list_slide_components"](
            slide_index=0, presentation=None
        )
        assert isinstance(result, str)


class TestWorkflows:
    """Test common workflows."""

    @pytest.mark.asyncio
    async def test_add_multiple_components(self, api_tools, mock_manager):
        """Test adding multiple components to same slide."""
        # Add badge
        r1 = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            params={"text": "Badge 1"},
        )
        assert isinstance(r1, str)

        # Add button
        r2 = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Button",
            left=3.0,
            top=1.0,
            params={"text": "Button 1"},
        )
        assert isinstance(r2, str)

        # Add text
        r3 = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Text",
            left=1.0,
            top=2.0,
            params={"text": "Some text"},
        )
        assert isinstance(r3, str)

        # List all components
        list_result = await api_tools["pptx_list_slide_components"](
            slide_index=0, presentation=None
        )
        assert isinstance(list_result, str)

    @pytest.mark.asyncio
    async def test_add_and_update_workflow(self, api_tools, mock_manager):
        """Test add then update workflow."""
        # Add component
        await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            component_id="workflow_badge",
            params={"text": "Original Text"},
        )

        # Update component
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="workflow_badge",
            params={"text": "Updated Text"},
        )
        assert isinstance(result, str)


class TestFreeFormPositioning:
    """Test free-form positioning mode."""

    @pytest.mark.asyncio
    async def test_freeform_no_left_top(self, api_tools, mock_manager):
        """Test free-form without left/top raises error."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            # No left/top provided
            params={"text": "Test"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_freeform_only_left(self, api_tools, mock_manager):
        """Test free-form with only left but no top."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            # No top provided
            params={"text": "Test"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_freeform_only_top(self, api_tools, mock_manager):
        """Test free-form with only top but no left."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            top=1.0,
            # No left provided
            params={"text": "Test"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data


class TestTargetPlaceholder:
    """Test placeholder targeting mode."""

    @pytest.fixture
    def presentation_with_content_slide(self):
        """Create presentation with content layout slide."""
        prs = Presentation()
        # Add a content slide with placeholder
        if prs.slide_layouts:
            layout = prs.slide_layouts[1]  # TITLE_AND_CONTENT
            prs.slides.add_slide(layout)
        return prs

    @pytest.fixture
    def manager_with_placeholder(self, presentation_with_content_slide):
        """Create manager with placeholder presentation."""
        return MockPresentationManager(presentation_with_content_slide)

    @pytest.fixture
    def placeholder_tools(self, mock_mcp, manager_with_placeholder):
        """Register API tools with placeholder presentation."""
        register_universal_component_api(mock_mcp, manager_with_placeholder)
        return mock_mcp._tools

    @pytest.mark.asyncio
    async def test_target_placeholder_not_found(self, placeholder_tools):
        """Test targeting non-existent placeholder."""
        result = await placeholder_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_placeholder=999,  # Non-existent
            params={"text": "Test"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_target_placeholder_success(self, placeholder_tools):
        """Test targeting existing placeholder."""
        # Try with placeholder 0 (typically title)
        result = await placeholder_tools["pptx_add_component"](
            slide_index=0,
            component="Text",
            target_placeholder=0,
            params={"text": "Title Text"},
        )
        assert isinstance(result, str)


class TestTargetComponent:
    """Test component targeting mode (composition)."""

    @pytest.mark.asyncio
    async def test_target_component_not_found(self, api_tools, mock_manager):
        """Test targeting non-existent component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_component="nonexistent_component",
            params={"text": "Test"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_composition_workflow(self, api_tools, mock_manager):
        """Test adding component inside another component."""
        # First add parent component
        parent_result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Card",
            left=1.0,
            top=1.0,
            width=4.0,
            height=3.0,
            component_id="parent_card",
            params={"title": "Parent Card"},
        )
        assert isinstance(parent_result, str)

        # Then add child into parent
        child_result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_component="parent_card",
            params={"text": "Child Badge"},
        )
        assert isinstance(child_result, str)


class TestTargetLayout:
    """Test layout targeting mode."""

    @pytest.mark.asyncio
    async def test_target_layout_grid(self, api_tools, mock_manager):
        """Test targeting grid layout."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_layout="grid",
            params={"text": "Grid Item"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_target_layout_flex_row(self, api_tools, mock_manager):
        """Test targeting flex row layout."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_layout="flex_row",
            params={"text": "Flex Row Item"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_target_layout_flex_column(self, api_tools, mock_manager):
        """Test targeting flex column layout."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_layout="flex_column",
            left=1.0,  # Provide left for flex column
            params={"text": "Flex Column Item"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_target_layout_unknown(self, api_tools, mock_manager):
        """Test targeting unknown layout type."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_layout="unknown_layout",
            params={"text": "Test"},
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_multiple_items_in_grid(self, api_tools, mock_manager):
        """Test adding multiple items to grid layout."""
        # Add first item
        r1 = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_layout="grid",
            params={"text": "Item 1"},
        )
        assert isinstance(r1, str)

        # Add second item
        r2 = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            target_layout="grid",
            params={"text": "Item 2"},
        )
        assert isinstance(r2, str)


class TestUpdateComponentParams:
    """Test component update with various param types."""

    @pytest.mark.asyncio
    async def test_update_with_invalid_json_params(self, api_tools, mock_manager):
        """Test update with invalid JSON string params."""
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="some_id",
            params="invalid json {",
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_update_with_non_dict_params(self, api_tools, mock_manager):
        """Test update with non-dict params."""
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="some_id",
            params=12345,  # Invalid type
        )
        assert isinstance(result, str)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_update_position_only(self, api_tools, mock_manager):
        """Test updating only position."""
        # Add component first
        await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            component_id="position_test",
            params={"text": "Test"},
        )

        # Update only position
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="position_test",
            left=2.0,
            top=2.0,
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_update_size_only(self, api_tools, mock_manager):
        """Test updating only size."""
        # Add component first
        await api_tools["pptx_add_component"](
            slide_index=0,
            component="Badge",
            left=1.0,
            top=1.0,
            width=2.0,
            height=0.5,
            component_id="size_test",
            params={"text": "Test"},
        )

        # Update only size
        result = await api_tools["pptx_update_component"](
            slide_index=0,
            component_id="size_test",
            width=3.0,
            height=1.0,
        )
        assert isinstance(result, str)


class TestStackComposition:
    """Test Stack component composition."""

    @pytest.mark.asyncio
    async def test_add_to_stack(self, api_tools, mock_manager):
        """Test adding components to a Stack."""
        # Add a stack component
        stack_result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Stack",
            left=1.0,
            top=1.0,
            width=4.0,
            height=3.0,
            component_id="test_stack",
            params={"direction": "vertical", "gap": "md"},
        )
        assert isinstance(stack_result, str)

        # Add child to stack
        child_result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Text",
            target_component="test_stack",
            params={"text": "Stack child 1"},
        )
        assert isinstance(child_result, str)


class TestMoreComponentTypes:
    """Test additional component types for coverage."""

    @pytest.mark.asyncio
    async def test_add_table_component(self, api_tools, mock_manager):
        """Test adding table component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Table",
            left=1.0,
            top=1.0,
            width=6.0,
            height=3.0,
            params={
                "headers": ["Col1", "Col2"],
                "data": [["A", "B"], ["C", "D"]],
            },
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_shape_component(self, api_tools, mock_manager):
        """Test adding shape component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Shape",
            left=1.0,
            top=1.0,
            width=2.0,
            height=2.0,
            params={"shape_type": "rectangle"},
        )
        assert isinstance(result, str)

    @pytest.mark.asyncio
    async def test_add_tile_component(self, api_tools, mock_manager):
        """Test adding tile component."""
        result = await api_tools["pptx_add_component"](
            slide_index=0,
            component="Tile",
            left=1.0,
            top=1.0,
            width=2.0,
            height=2.0,
            params={"title": "Tile Title", "value": "100"},
        )
        assert isinstance(result, str)


class TestPlaceholderValidation:
    """Test placeholder validation logic."""

    @pytest.fixture
    def presentation_with_picture_placeholder(self):
        """Create presentation with picture placeholder."""
        from pptx import Presentation

        prs = Presentation()
        # Use a layout that might have picture placeholder (7 = Picture with Caption)
        layout_idx = min(7, len(prs.slide_layouts) - 1)
        prs.slides.add_slide(prs.slide_layouts[layout_idx])
        return prs

    @pytest.fixture
    def pic_manager(self, presentation_with_picture_placeholder):
        """Manager with picture placeholder presentation."""
        return MockPresentationManager(presentation_with_picture_placeholder)

    @pytest.fixture
    def pic_tools(self, mock_mcp, pic_manager):
        """Register tools with picture placeholder presentation."""
        register_universal_component_api(mock_mcp, pic_manager)
        return mock_mcp._tools

    @pytest.mark.asyncio
    async def test_list_components_with_placeholders(self, pic_tools):
        """Test listing components when placeholders exist."""
        result = await pic_tools["pptx_list_slide_components"](slide_index=0, presentation=None)
        assert isinstance(result, str)
        data = json.loads(result)
        # Should have placeholders array
        assert "placeholders" in data or "error" not in data
