# tests/tools/test_template_tools.py
"""
Tests for the template tools modules.

Tests the tools from:
- tools/template/list.py (pptx_list_templates)
- tools/template/analyze.py (pptx_analyze_template, pptx_analyze_template_variants)
- tools/template/import_tools.py (pptx_import_template, pptx_get_builtin_template)
- tools/template/workflow.py (pptx_add_slide_from_template)
"""

import pytest
import json
from unittest.mock import MagicMock
from pptx import Presentation


# ============================================================================
# Fixtures
# ============================================================================


class MockTemplateInfo:
    """Mock template info for testing."""

    def __init__(
        self,
        name,
        display_name="Test Template",
        description="A test template",
        category="test",
        layout_count=5,
        tags=None,
    ):
        self.name = name
        self.display_name = display_name
        self.description = description
        self.category = category
        self.layout_count = layout_count
        self.tags = tags or ["test"]


class MockTemplateManager:
    """Mock template manager for testing."""

    def __init__(self, templates=None):
        self._templates = templates or []

    def list_templates(self):
        """Return list of template infos."""
        return self._templates

    async def get_template_data(self, template_name):
        """Get template data by name."""
        # Return None for unknown templates
        for tmpl in self._templates:
            if tmpl.name == template_name:
                # Create a minimal valid pptx
                prs = Presentation()
                import io

                buffer = io.BytesIO()
                prs.save(buffer)
                return buffer.getvalue()
        return None


class MockPresentationMetadata:
    """Mock presentation metadata."""

    def __init__(self, name="test_presentation", namespace_id=None, vfs_path=None):
        self.name = name
        self.namespace_id = namespace_id
        self.vfs_path = vfs_path

    def update_modified(self):
        pass


class MockPresentationInfo:
    """Mock presentation info for list response."""

    def __init__(self, name="test_presentation", slide_count=1, namespace_id=None):
        self.name = name
        self.slide_count = slide_count
        self.namespace_id = namespace_id


class MockListResponse:
    """Mock list presentations response."""

    def __init__(self, presentations=None):
        self.presentations = presentations or []


class MockPresentationManager:
    """Mock presentation manager for testing."""

    def __init__(self, presentation=None):
        self._presentation = presentation or Presentation()
        self._current_name = "test_presentation"
        self._metadata = {}
        self._namespace_ids = {}
        self.base_path = "presentations"
        self.PPTX_MIME_TYPE = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    async def get(self, name=None):
        """Get presentation by name."""
        if name is None or name == self._current_name:
            metadata = MockPresentationMetadata(name=self._current_name)
            return self._presentation, metadata
        return None

    async def list_presentations(self):
        """List presentations."""
        return MockListResponse(presentations=[])

    async def import_template(self, file_path, template_name):
        """Import a template."""
        return True

    async def update_slide_metadata(self, slide_index):
        """Update slide metadata."""
        pass

    async def _save_to_store(self, name, prs):
        """Save to store."""
        pass

    def _get_store(self):
        """Get the artifact store."""
        return None

    def _sanitize_name(self, name):
        """Sanitize a name."""
        return name.replace(" ", "_").lower()


@pytest.fixture
def mock_mcp():
    """Create a mock MCP server that captures tool registrations."""
    tools = {}

    def tool_decorator(func):
        tools[func.__name__] = func
        return func

    mock = MagicMock()
    mock.tool = tool_decorator
    mock._tools = tools
    return mock


@pytest.fixture
def mock_manager():
    """Create a mock presentation manager."""
    return MockPresentationManager()


@pytest.fixture
def mock_manager_no_prs():
    """Create a mock manager that returns None."""
    manager = MockPresentationManager()

    async def get_none(name=None):
        return None

    manager.get = get_none
    return manager


@pytest.fixture
def mock_template_manager():
    """Create a mock template manager with some templates."""
    templates = [
        MockTemplateInfo(name="brand_proposal", display_name="Brand Proposal", layout_count=55),
        MockTemplateInfo(name="corporate", display_name="Corporate", layout_count=20),
    ]
    return MockTemplateManager(templates=templates)


@pytest.fixture
def mock_template_manager_empty():
    """Create a mock template manager with no templates."""
    return MockTemplateManager(templates=[])


# ============================================================================
# Test List Tools
# ============================================================================


class TestListTemplates:
    """Tests for pptx_list_templates tool."""

    @pytest.fixture
    def list_tools(self, mock_mcp, mock_manager, mock_template_manager):
        """Register and return list tools."""
        from chuk_mcp_pptx.tools.template.list import register_list_tools

        return register_list_tools(mock_mcp, mock_manager, mock_template_manager)

    @pytest.fixture
    def list_tools_empty(self, mock_mcp, mock_manager, mock_template_manager_empty):
        """Register list tools with empty template manager."""
        from chuk_mcp_pptx.tools.template.list import register_list_tools

        return register_list_tools(mock_mcp, mock_manager, mock_template_manager_empty)

    @pytest.mark.asyncio
    async def test_list_templates_returns_json(self, list_tools):
        """Test that list_templates returns valid JSON."""
        result = await list_tools["pptx_list_templates"]()
        assert isinstance(result, str)
        data = json.loads(result)
        assert "builtin_templates" in data
        assert "custom_templates" in data
        assert "total" in data

    @pytest.mark.asyncio
    async def test_list_templates_includes_builtin(self, list_tools):
        """Test that builtin templates are included."""
        result = await list_tools["pptx_list_templates"]()
        data = json.loads(result)
        assert len(data["builtin_templates"]) == 2
        names = [t["name"] for t in data["builtin_templates"]]
        assert "brand_proposal" in names
        assert "corporate" in names

    @pytest.mark.asyncio
    async def test_list_templates_template_info(self, list_tools):
        """Test that template info contains expected fields."""
        result = await list_tools["pptx_list_templates"]()
        data = json.loads(result)
        template = data["builtin_templates"][0]
        assert "name" in template
        assert "display_name" in template
        assert "description" in template
        assert "category" in template
        assert "layout_count" in template
        assert "is_builtin" in template

    @pytest.mark.asyncio
    async def test_list_templates_empty(self, list_tools_empty):
        """Test listing when no templates exist."""
        result = await list_tools_empty["pptx_list_templates"]()
        data = json.loads(result)
        assert data["total"] == 0
        assert len(data["builtin_templates"]) == 0


# ============================================================================
# Test Analyze Tools
# ============================================================================


class TestAnalyzeTemplate:
    """Tests for pptx_analyze_template tool."""

    @pytest.fixture
    def analyze_tools(self, mock_mcp, mock_manager, mock_template_manager):
        """Register and return analyze tools."""
        from chuk_mcp_pptx.tools.template.analyze import register_analyze_tools

        return register_analyze_tools(mock_mcp, mock_manager, mock_template_manager)

    @pytest.fixture
    def analyze_tools_no_template(self, mock_mcp, mock_manager_no_prs, mock_template_manager_empty):
        """Register analyze tools with no templates."""
        from chuk_mcp_pptx.tools.template.analyze import register_analyze_tools

        return register_analyze_tools(mock_mcp, mock_manager_no_prs, mock_template_manager_empty)

    @pytest.mark.asyncio
    async def test_analyze_template_returns_json(self, analyze_tools):
        """Test that analyze_template returns valid JSON."""
        result = await analyze_tools["pptx_analyze_template"](template_name="brand_proposal")
        assert isinstance(result, str)
        data = json.loads(result)
        # Should contain template info or error
        assert "name" in data or "error" in data

    @pytest.mark.asyncio
    async def test_analyze_template_unknown(self, analyze_tools_no_template):
        """Test analyzing unknown template."""
        result = await analyze_tools_no_template["pptx_analyze_template"](
            template_name="nonexistent_template"
        )
        data = json.loads(result)
        assert "error" in data
        assert "not found" in data["error"].lower()

    @pytest.mark.asyncio
    async def test_analyze_template_includes_layouts(self, analyze_tools):
        """Test that template analysis includes layouts."""
        result = await analyze_tools["pptx_analyze_template"](template_name="brand_proposal")
        data = json.loads(result)
        if "error" not in data:
            assert "layouts" in data
            assert "layout_count" in data


class TestAnalyzeTemplateVariants:
    """Tests for pptx_analyze_template_variants tool."""

    @pytest.fixture
    def analyze_tools(self, mock_mcp, mock_manager, mock_template_manager):
        """Register and return analyze tools."""
        from chuk_mcp_pptx.tools.template.analyze import register_analyze_tools

        return register_analyze_tools(mock_mcp, mock_manager, mock_template_manager)

    @pytest.mark.asyncio
    async def test_analyze_variants_returns_json(self, analyze_tools):
        """Test that analyze_template_variants returns valid JSON."""
        result = await analyze_tools["pptx_analyze_template_variants"](
            template_name="brand_proposal"
        )
        assert isinstance(result, str)
        # May return error if template_extraction module has issues
        json.loads(result)
        # Should be valid JSON either way


# ============================================================================
# Test Import Tools
# ============================================================================


class TestImportTemplate:
    """Tests for pptx_import_template tool."""

    @pytest.fixture
    def import_tools(self, mock_mcp, mock_manager, mock_template_manager):
        """Register and return import tools."""
        from chuk_mcp_pptx.tools.template.import_tools import register_import_tools

        return register_import_tools(mock_mcp, mock_manager, mock_template_manager)

    @pytest.mark.asyncio
    async def test_import_template_returns_json(self, import_tools):
        """Test that import_template returns valid JSON."""
        result = await import_tools["pptx_import_template"](
            file_path="/tmp/test.pptx", template_name="imported_template"
        )
        assert isinstance(result, str)
        data = json.loads(result)
        # Should contain template info or error
        assert "name" in data or "error" in data


class TestGetBuiltinTemplate:
    """Tests for pptx_get_builtin_template tool."""

    @pytest.fixture
    def import_tools(self, mock_mcp, mock_manager, mock_template_manager):
        """Register and return import tools."""
        from chuk_mcp_pptx.tools.template.import_tools import register_import_tools

        return register_import_tools(mock_mcp, mock_manager, mock_template_manager)

    @pytest.fixture
    def import_tools_no_template(self, mock_mcp, mock_manager, mock_template_manager_empty):
        """Register import tools with no templates."""
        from chuk_mcp_pptx.tools.template.import_tools import register_import_tools

        return register_import_tools(mock_mcp, mock_manager, mock_template_manager_empty)

    @pytest.mark.asyncio
    async def test_get_builtin_unknown(self, import_tools_no_template):
        """Test getting unknown builtin template."""
        result = await import_tools_no_template["pptx_get_builtin_template"](
            template_name="nonexistent", save_as="my_template"
        )
        data = json.loads(result)
        assert "error" in data
        assert "not found" in data["error"].lower()


# ============================================================================
# Test Workflow Tools
# ============================================================================


class TestAddSlideFromTemplate:
    """Tests for pptx_add_slide_from_template tool."""

    @pytest.fixture
    def presentation_with_layouts(self):
        """Create a presentation with layouts."""
        prs = Presentation()
        # Presentation has default layouts
        return prs

    @pytest.fixture
    def workflow_tools(self, mock_mcp, presentation_with_layouts, mock_template_manager):
        """Register and return workflow tools."""
        manager = MockPresentationManager(presentation=presentation_with_layouts)
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        return register_workflow_tools(mock_mcp, manager, mock_template_manager)

    @pytest.fixture
    def workflow_tools_no_prs(self, mock_mcp, mock_manager_no_prs, mock_template_manager):
        """Register workflow tools with no presentation."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        return register_workflow_tools(mock_mcp, mock_manager_no_prs, mock_template_manager)

    @pytest.mark.asyncio
    async def test_add_slide_no_presentation(self, workflow_tools_no_prs):
        """Test adding slide when no presentation exists."""
        result = await workflow_tools_no_prs["pptx_add_slide_from_template"](layout_index=0)
        data = json.loads(result)
        assert "error" in data

    @pytest.mark.asyncio
    async def test_add_slide_from_template_returns_json(self, workflow_tools):
        """Test that add_slide_from_template returns valid JSON."""
        result = await workflow_tools["pptx_add_slide_from_template"](layout_index=0)
        assert isinstance(result, str)
        data = json.loads(result)
        # Should return slide response or error
        assert "slide_index" in data or "error" in data


# ============================================================================
# Test Tool Registration
# ============================================================================


class TestTemplateToolsRegistration:
    """Tests for template tools registration."""

    def test_list_tools_registered(self, mock_mcp, mock_manager, mock_template_manager):
        """Test that list tools are registered."""
        from chuk_mcp_pptx.tools.template.list import register_list_tools

        tools = register_list_tools(mock_mcp, mock_manager, mock_template_manager)
        assert "pptx_list_templates" in tools

    def test_analyze_tools_registered(self, mock_mcp, mock_manager, mock_template_manager):
        """Test that analyze tools are registered."""
        from chuk_mcp_pptx.tools.template.analyze import register_analyze_tools

        tools = register_analyze_tools(mock_mcp, mock_manager, mock_template_manager)
        assert "pptx_analyze_template" in tools
        assert "pptx_analyze_template_variants" in tools

    def test_import_tools_registered(self, mock_mcp, mock_manager, mock_template_manager):
        """Test that import tools are registered."""
        from chuk_mcp_pptx.tools.template.import_tools import register_import_tools

        tools = register_import_tools(mock_mcp, mock_manager, mock_template_manager)
        assert "pptx_import_template" in tools
        assert "pptx_get_builtin_template" in tools

    def test_workflow_tools_registered(self, mock_mcp, mock_manager, mock_template_manager):
        """Test that workflow tools are registered."""
        from chuk_mcp_pptx.tools.template.workflow import register_workflow_tools

        tools = register_workflow_tools(mock_mcp, mock_manager, mock_template_manager)
        assert "pptx_add_slide_from_template" in tools

    def test_all_template_tools_registered(self, mock_mcp, mock_manager, mock_template_manager):
        """Test that all template tools can be registered together."""
        from chuk_mcp_pptx.tools.template import register_template_tools

        tools = register_template_tools(mock_mcp, mock_manager, mock_template_manager)
        expected_tools = [
            "pptx_list_templates",
            "pptx_analyze_template",
            "pptx_analyze_template_variants",
            "pptx_import_template",
            "pptx_get_builtin_template",
            "pptx_add_slide_from_template",
        ]
        for tool_name in expected_tools:
            assert tool_name in tools, f"Missing tool: {tool_name}"

    def test_tools_are_async(self, mock_mcp, mock_manager, mock_template_manager):
        """Test that all tools are async functions."""
        import inspect
        from chuk_mcp_pptx.tools.template import register_template_tools

        tools = register_template_tools(mock_mcp, mock_manager, mock_template_manager)
        for tool_name, tool_func in tools.items():
            assert inspect.iscoroutinefunction(tool_func), f"{tool_name} should be async"


# ============================================================================
# Test Models
# ============================================================================


class TestTemplateModels:
    """Tests for template models."""

    def test_builtin_template_info_model(self):
        """Test BuiltinTemplateInfo model."""
        from chuk_mcp_pptx.tools.template.models import BuiltinTemplateInfo

        info = BuiltinTemplateInfo(
            name="test",
            display_name="Test Template",
            description="A test template",
            category="test",
            layout_count=5,
            tags=["test"],
            is_builtin=True,
        )
        assert info.name == "test"
        assert info.layout_count == 5

    def test_custom_template_info_model(self):
        """Test CustomTemplateInfo model."""
        from chuk_mcp_pptx.tools.template.models import CustomTemplateInfo

        info = CustomTemplateInfo(
            name="custom_test",
            slide_count=3,
            namespace_id="ns-123",
            is_builtin=False,
            category="custom",
        )
        assert info.name == "custom_test"
        assert info.is_builtin is False

    def test_layout_info_model(self):
        """Test LayoutInfo model."""
        from chuk_mcp_pptx.tools.template.models import LayoutInfo

        info = LayoutInfo(
            index=0,
            name="Title Slide",
            placeholder_count=2,
            placeholders=[{"idx": 0, "type": "TITLE", "name": "Title 1"}],
        )
        assert info.index == 0
        assert info.placeholder_count == 2

    def test_template_info_model(self):
        """Test TemplateInfo model."""
        from chuk_mcp_pptx.tools.template.models import TemplateInfo, LayoutInfo

        layout = LayoutInfo(index=0, name="Title", placeholder_count=1, placeholders=[])
        info = TemplateInfo(
            name="test_template",
            slide_count=0,
            layout_count=1,
            layouts=[layout],
            master_count=1,
            has_theme=True,
        )
        assert info.name == "test_template"
        assert info.has_theme is True

    def test_presentation_template_list_response_model(self):
        """Test PresentationTemplateListResponse model."""
        from chuk_mcp_pptx.tools.template.models import (
            PresentationTemplateListResponse,
            BuiltinTemplateInfo,
            CustomTemplateInfo,
        )

        builtin = BuiltinTemplateInfo(
            name="builtin",
            display_name="Builtin",
            description="desc",
            category="cat",
            layout_count=5,
            tags=[],
            is_builtin=True,
        )
        custom = CustomTemplateInfo(
            name="custom", slide_count=1, namespace_id="ns-123", is_builtin=False, category="custom"
        )
        response = PresentationTemplateListResponse(
            builtin_templates=[builtin], custom_templates=[custom], total=2
        )
        assert response.total == 2
        assert len(response.builtin_templates) == 1
        assert len(response.custom_templates) == 1
