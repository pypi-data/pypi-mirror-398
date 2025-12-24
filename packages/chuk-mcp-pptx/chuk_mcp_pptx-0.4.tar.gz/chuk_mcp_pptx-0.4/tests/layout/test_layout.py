"""
Comprehensive tests for Layout components.
Tests Grid, Container, Stack, Spacer, and Divider.
"""

import pytest
from pptx import Presentation

from chuk_mcp_pptx.components.core import (
    Grid,
    Container,
    Stack,
    Spacer,
    Divider,
)
from chuk_mcp_pptx.layout import SLIDE_WIDTH
from chuk_mcp_pptx.themes import ThemeManager


class TestGrid:
    """Test Grid component."""

    def test_grid_creation(self):
        """Test creating a grid."""
        grid = Grid(columns=12, gap="md")
        assert grid.columns == 12
        assert grid.gap_inches > 0

    def test_grid_custom_columns(self):
        """Test grid with custom column count."""
        grid = Grid(columns=3, gap="lg")
        assert grid.columns == 3

    def test_grid_gap_sizes(self):
        """Test different gap sizes."""
        gaps = ["none", "xs", "sm", "md", "lg", "xl"]
        for gap in gaps:
            grid = Grid(gap=gap)
            assert grid.gap_inches >= 0

    def test_grid_cell_positions(self, mock_slide):
        """Test getting cell positions."""
        grid = Grid(columns=3, gap="md")
        cells = grid.get_cell_positions(mock_slide, left=0.5, top=2.0, width=9.0, height=3.0)

        assert len(cells) == 3
        for cell in cells:
            assert "left" in cell
            assert "top" in cell
            assert "width" in cell
            assert "height" in cell

    def test_grid_span_single_column(self):
        """Test spanning a single column."""
        grid = Grid(columns=12, gap="md")
        pos = grid.get_span(col_span=1, col_start=0, left=0.5, top=2.0, width=9.0)

        assert pos["left"] == 0.5
        assert pos["width"] < 1.0

    def test_grid_span_full_width(self):
        """Test spanning all columns."""
        grid = Grid(columns=12, gap="md")
        pos = grid.get_span(col_span=12, col_start=0, left=0.5, top=2.0, width=9.0, height=2.0)

        # Should span nearly full width minus gaps
        assert pos["width"] > 8.0

    def test_grid_span_half_width(self):
        """Test spanning half width."""
        grid = Grid(columns=12, gap="md")
        pos1 = grid.get_span(col_span=6, col_start=0, left=0.5, width=9.0)
        pos2 = grid.get_span(col_span=6, col_start=6, left=0.5, width=9.0)

        # Both should have similar widths
        assert abs(pos1["width"] - pos2["width"]) < 0.1

        # Second should start after first
        assert pos2["left"] > pos1["left"]

    def test_grid_multi_row(self):
        """Test grid with multiple rows."""
        grid = Grid(columns=3, rows=2, gap="sm")
        cells = grid.get_cell_positions(None, left=0.5, top=2.0, width=9.0, height=4.0)

        assert len(cells) == 6  # 3 cols × 2 rows
        assert cells[0]["row"] == 0
        assert cells[3]["row"] == 1


class TestContainer:
    """Test Container component."""

    def test_container_creation(self):
        """Test creating a container."""
        container = Container(size="lg", padding="md")
        assert container.size == "lg"

    def test_container_sizes(self):
        """Test different container sizes."""
        sizes = ["sm", "md", "lg", "xl", "2xl", "full"]
        for size in sizes:
            container = Container(size=size)
            assert container.size == size

    def test_container_render(self, mock_slide):
        """Test rendering container."""
        container = Container(size="lg", padding="md", center=True)
        bounds = container.render(mock_slide, top=2.0, height=3.0)

        assert "left" in bounds
        assert "top" in bounds
        assert "width" in bounds
        assert "height" in bounds

    def test_container_centering(self, mock_slide):
        """Test container centering."""
        # Use no padding to see clear difference
        container_centered = Container(size="md", center=True, padding="none")
        container_centered.render(mock_slide, top=2.0)

        container_not_centered = Container(size="md", center=False, padding="none")
        container_not_centered.render(mock_slide, top=2.0)

        # Centered should have different left position
        # Centered: (10 - 9) / 2 = 0.5
        # Not centered: 0.5
        # Actually they're the same for md size! Let's use a smaller size

        container_centered_sm = Container(size="sm", center=True, padding="none")
        bounds_centered_sm = container_centered_sm.render(mock_slide, top=2.0)

        container_not_centered_sm = Container(size="sm", center=False, padding="none")
        bounds_not_centered_sm = container_not_centered_sm.render(mock_slide, top=2.0)

        # Centered small should be more to the right than non-centered
        assert bounds_centered_sm["left"] > bounds_not_centered_sm["left"]

    def test_container_padding(self, mock_slide):
        """Test container padding."""
        container_no_pad = Container(size="lg", padding="none")
        bounds_no_pad = container_no_pad.render(mock_slide, top=2.0, height=3.0)

        container_with_pad = Container(size="lg", padding="lg")
        bounds_with_pad = container_with_pad.render(mock_slide, top=2.0, height=3.0)

        # With padding should have smaller dimensions
        assert bounds_with_pad["width"] < bounds_no_pad["width"]

    def test_container_small_size(self, mock_slide):
        """Test small container."""
        container = Container(size="sm")
        bounds = container.render(mock_slide, top=2.0)

        # Small container should be narrower than slide
        assert bounds["width"] < SLIDE_WIDTH

    def test_container_full_size(self, mock_slide):
        """Test full width container."""
        container = Container(size="full", padding="none")
        bounds = container.render(mock_slide, top=2.0)

        # Full size should be close to slide width
        assert bounds["width"] >= SLIDE_WIDTH - 1.0


class TestStack:
    """Test Stack component."""

    def test_stack_creation(self):
        """Test creating a stack."""
        stack = Stack(direction="vertical", gap="md")
        assert stack.direction == "vertical"

    def test_stack_directions(self):
        """Test different stack directions."""
        v_stack = Stack(direction="vertical")
        h_stack = Stack(direction="horizontal")

        assert v_stack.direction == "vertical"
        assert h_stack.direction == "horizontal"

    def test_stack_vertical_distribution(self):
        """Test vertical stack distribution."""
        stack = Stack(direction="vertical", gap="md")
        positions = stack.distribute(num_items=3, item_height=1.0, top=2.0)

        assert len(positions) == 3

        # Items should be stacked vertically
        assert positions[0]["top"] < positions[1]["top"]
        assert positions[1]["top"] < positions[2]["top"]

    def test_stack_horizontal_distribution(self):
        """Test horizontal stack distribution."""
        stack = Stack(direction="horizontal", gap="lg")
        positions = stack.distribute(num_items=4, item_width=2.0, left=0.5, top=3.0)

        assert len(positions) == 4

        # Items should be arranged horizontally
        assert positions[0]["left"] < positions[1]["left"]
        assert positions[1]["left"] < positions[2]["left"]

    def test_stack_alignment_start(self):
        """Test stack with start alignment."""
        stack = Stack(direction="vertical", gap="sm", align="start")
        positions = stack.distribute(num_items=2, item_height=1.0, top=2.0)

        # Should align to left/top
        assert positions[0]["left"] >= 0

    def test_stack_alignment_center(self):
        """Test stack with center alignment."""
        stack = Stack(direction="vertical", gap="sm", align="center")
        positions = stack.distribute(num_items=2, item_width=4.0, item_height=1.0, top=2.0)

        # Should be centered
        assert positions[0]["left"] > 2.0

    def test_stack_alignment_end_vertical(self):
        """Test vertical stack with end alignment."""
        stack = Stack(direction="vertical", gap="sm", align="end")
        positions = stack.distribute(
            num_items=2, item_width=4.0, item_height=1.0, top=2.0, left=1.0
        )

        # Should align to right edge (end)
        # Items should be near the right edge of slide
        assert positions[0]["left"] > 5.0

    def test_stack_alignment_end_horizontal(self):
        """Test horizontal stack with end alignment."""
        stack = Stack(direction="horizontal", gap="sm", align="end")
        positions = stack.distribute(
            num_items=2, item_width=2.0, item_height=1.0, top=2.0, left=1.0
        )

        # Should align to bottom edge (end)
        # Items should be near the bottom edge of slide (SLIDE_HEIGHT - height - margin)
        assert positions[0]["top"] > 4.0

    def test_stack_gap_consistency(self):
        """Test gap consistency between items."""
        stack = Stack(direction="vertical", gap="lg")
        positions = stack.distribute(num_items=3, item_height=1.0, top=2.0)

        # Calculate gaps
        gap1 = positions[1]["top"] - (positions[0]["top"] + positions[0]["height"])
        gap2 = positions[2]["top"] - (positions[1]["top"] + positions[1]["height"])

        # Gaps should be equal
        assert abs(gap1 - gap2) < 0.01

    def test_stack_empty(self):
        """Test stack with no items."""
        stack = Stack(direction="vertical", gap="md")
        positions = stack.distribute(num_items=0, top=2.0)

        assert len(positions) == 0

    def test_stack_render_children_vertical(self, mock_slide):
        """Test render_children with vertical stack."""

        # Create mock components with render method
        class MockComponent:
            def __init__(self, name):
                self.name = name
                self.rendered = False
                self.render_args = None

            def render(self, slide, **kwargs):
                self.rendered = True
                self.render_args = kwargs
                return f"shape_{self.name}"

        components = [MockComponent(f"item_{i}") for i in range(3)]

        stack = Stack(direction="vertical", gap="md")
        shapes = stack.render_children(
            mock_slide, components, left=0.5, top=2.0, item_width=4.0, item_height=1.0
        )

        # All components should be rendered
        assert len(shapes) == 3
        for component in components:
            assert component.rendered

        # Verify vertical stacking (increasing top positions)
        assert components[0].render_args["top"] < components[1].render_args["top"]
        assert components[1].render_args["top"] < components[2].render_args["top"]

        # Verify consistent dimensions
        for component in components:
            assert component.render_args["width"] == 4.0
            assert component.render_args["height"] == 1.0

    def test_stack_render_children_horizontal(self, mock_slide):
        """Test render_children with horizontal stack."""

        class MockComponent:
            def __init__(self, name):
                self.name = name
                self.rendered = False
                self.render_args = None

            def render(self, slide, **kwargs):
                self.rendered = True
                self.render_args = kwargs
                return f"shape_{self.name}"

        components = [MockComponent(f"item_{i}") for i in range(4)]

        stack = Stack(direction="horizontal", gap="lg")
        shapes = stack.render_children(
            mock_slide, components, left=0.5, top=2.0, item_width=2.0, item_height=1.5
        )

        # All components should be rendered
        assert len(shapes) == 4
        for component in components:
            assert component.rendered

        # Verify horizontal stacking (increasing left positions)
        assert components[0].render_args["left"] < components[1].render_args["left"]
        assert components[1].render_args["left"] < components[2].render_args["left"]
        assert components[2].render_args["left"] < components[3].render_args["left"]

        # Verify consistent dimensions
        for component in components:
            assert component.render_args["width"] == 2.0
            assert component.render_args["height"] == 1.5

    def test_stack_render_children_empty(self, mock_slide):
        """Test render_children with empty list."""
        stack = Stack(direction="vertical", gap="md")
        shapes = stack.render_children(mock_slide, [], left=0.5, top=2.0)

        assert len(shapes) == 0

    def test_stack_render_children_matches_distribute(self, mock_slide):
        """Test that render_children positioning matches distribute."""

        class MockComponent:
            def __init__(self):
                self.render_args = None

            def render(self, slide, **kwargs):
                self.render_args = kwargs
                return "shape"

        components = [MockComponent() for _ in range(3)]

        stack = Stack(direction="vertical", gap="sm")

        # Get positions from distribute
        positions = stack.distribute(
            num_items=3, item_width=3.0, item_height=1.0, left=1.0, top=2.5
        )

        # Render using render_children
        stack.render_children(
            mock_slide, components, left=1.0, top=2.5, item_width=3.0, item_height=1.0
        )

        # Verify positions match
        for i, (component, expected_pos) in enumerate(zip(components, positions)):
            assert component.render_args["left"] == expected_pos["left"]
            assert component.render_args["top"] == expected_pos["top"]
            assert component.render_args["width"] == expected_pos["width"]
            assert component.render_args["height"] == expected_pos["height"]

    def test_stack_render_children_without_render_method(self, mock_slide):
        """Test render_children skips objects without render method."""

        class MockComponent:
            def render(self, slide, **kwargs):
                return "shape"

        class InvalidComponent:
            pass  # No render method

        components = [MockComponent(), InvalidComponent(), MockComponent()]

        stack = Stack(direction="vertical", gap="md")
        shapes = stack.render_children(
            mock_slide, components, left=0.5, top=2.0, item_width=3.0, item_height=1.0
        )

        # Should only render components with render method
        assert len(shapes) == 2


class TestSpacer:
    """Test Spacer component."""

    def test_spacer_creation(self):
        """Test creating a spacer."""
        spacer = Spacer(size="md", direction="vertical")
        assert spacer.size == "md"
        assert spacer.direction == "vertical"

    def test_spacer_sizes(self):
        """Test different spacer sizes."""
        sizes = ["xs", "sm", "md", "lg", "xl", "2xl"]
        prev_size = 0

        for size in sizes:
            spacer = Spacer(size=size)
            current_size = spacer.get_size()

            # Sizes should increase
            assert current_size > prev_size
            prev_size = current_size

    def test_spacer_vertical(self):
        """Test vertical spacer."""
        spacer = Spacer(size="lg", direction="vertical")
        size = spacer.get_size()

        assert size > 0

    def test_spacer_horizontal(self):
        """Test horizontal spacer."""
        spacer = Spacer(size="xl", direction="horizontal")
        size = spacer.get_size()

        assert size > 0

    def test_spacer_render(self, mock_slide):
        """Test spacer render returns dimensions."""
        spacer = Spacer(size="md", direction="vertical")
        result = spacer.render(mock_slide, left=0, top=0)

        assert "height" in result or "width" in result

    def test_spacer_render_vertical(self, mock_slide):
        """Test vertical spacer returns height."""
        spacer = Spacer(size="lg", direction="vertical")
        result = spacer.render(mock_slide, left=0, top=0)

        assert "height" in result
        assert result["height"] > 0
        assert result["width"] == 0

    def test_spacer_render_horizontal(self, mock_slide):
        """Test horizontal spacer returns width."""
        spacer = Spacer(size="xl", direction="horizontal")
        result = spacer.render(mock_slide, left=0, top=0)

        assert "width" in result
        assert result["width"] > 0
        assert result["height"] == 0


class TestDivider:
    """Test Divider component."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_divider_creation(self):
        """Test creating a divider."""
        divider = Divider(orientation="horizontal", thickness=1)
        assert divider.orientation == "horizontal"
        assert divider.thickness == 1

    def test_divider_orientations(self):
        """Test different divider orientations."""
        h_divider = Divider(orientation="horizontal")
        v_divider = Divider(orientation="vertical")

        assert h_divider.orientation == "horizontal"
        assert v_divider.orientation == "vertical"

    def test_horizontal_divider_render(self, slide):
        """Test rendering horizontal divider."""
        divider = Divider(orientation="horizontal", thickness=1)
        shape = divider.render(slide, left=0.5, top=3.0, width=9.0)

        assert shape is not None
        assert shape in slide.shapes

    def test_vertical_divider_render(self, slide):
        """Test rendering vertical divider."""
        divider = Divider(orientation="vertical", thickness=2)
        shape = divider.render(slide, left=5.0, top=2.0, height=4.0)

        assert shape is not None
        assert shape in slide.shapes

    def test_divider_thickness(self, slide):
        """Test divider with different thickness."""
        thin_divider = Divider(orientation="horizontal", thickness=0.5)
        thick_divider = Divider(orientation="horizontal", thickness=3)

        thin_shape = thin_divider.render(slide, left=0.5, top=2.0, width=9.0)
        thick_shape = thick_divider.render(slide, left=0.5, top=3.0, width=9.0)

        assert thin_shape is not None
        assert thick_shape is not None

    def test_divider_with_theme(self, slide):
        """Test divider with theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("dark-violet")

        divider = Divider(orientation="horizontal", theme=theme.__dict__)
        shape = divider.render(slide, left=0.5, top=3.0, width=9.0)

        assert shape is not None


class TestLayoutIntegration:
    """Test layout components working together."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_grid_with_dividers(self, slide):
        """Test grid layout with dividers."""
        grid = Grid(columns=2, gap="md")
        cells = grid.get_cell_positions(slide, left=0.5, top=2.0, width=9.0, height=3.0)

        # Add divider between columns
        divider = Divider(orientation="vertical")
        shape = divider.render(slide, left=5.0, top=2.0, height=3.0)

        assert len(cells) == 2
        assert shape is not None

    def test_container_with_stack(self, slide):
        """Test container with stacked content."""
        container = Container(size="lg", padding="md", center=True)
        bounds = container.render(slide, top=2.0, height=4.0)

        stack = Stack(direction="vertical", gap="sm")
        positions = stack.distribute(
            num_items=3,
            item_width=bounds["width"],
            item_height=1.0,
            left=bounds["left"],
            top=bounds["top"],
        )

        assert len(positions) == 3

    def test_responsive_layout_pattern(self, slide):
        """Test responsive dashboard pattern."""
        # Main grid
        grid = Grid(columns=12, gap="md")

        # Main content (8 cols)
        main = grid.get_span(col_span=8, col_start=0, left=0.5, top=2.0, width=9.0)

        # Sidebar (4 cols)
        sidebar = grid.get_span(col_span=4, col_start=8, left=0.5, top=2.0, width=9.0)

        # Main should be wider
        assert main["width"] > sidebar["width"]

        # Sidebar stack
        stack = Stack(direction="vertical", gap="sm")
        items = stack.distribute(
            num_items=3,
            item_width=sidebar["width"],
            item_height=0.8,
            left=sidebar["left"],
            top=sidebar["top"],
        )

        assert len(items) == 3


class TestLayoutEdgeCases:
    """Test layout edge cases."""

    def test_grid_zero_gap(self):
        """Test grid with no gap."""
        grid = Grid(columns=3, gap="none")
        cells = grid.get_cell_positions(None, left=0.5, top=2.0, width=9.0, height=3.0)

        # Cells should be adjacent
        assert abs(cells[0]["left"] + cells[0]["width"] - cells[1]["left"]) < 0.01

    def test_stack_single_item(self):
        """Test stack with single item."""
        stack = Stack(direction="vertical", gap="lg")
        positions = stack.distribute(num_items=1, item_height=1.0, top=2.0)

        assert len(positions) == 1

    def test_container_extreme_sizes(self, mock_slide):
        """Test container with extreme sizes."""
        # Very small
        small = Container(size="sm", padding="none")
        small_bounds = small.render(mock_slide, top=2.0)

        # Very large
        large = Container(size="full", padding="none")
        large_bounds = large.render(mock_slide, top=2.0)

        assert small_bounds["width"] < large_bounds["width"]

    def test_spacer_minimum_size(self):
        """Test smallest spacer."""
        spacer = Spacer(size="xs")
        size = spacer.get_size()

        assert size > 0

    def test_grid_column_overflow(self):
        """Test grid span exceeding columns."""
        grid = Grid(columns=12, gap="md")

        # This should still work, just might extend beyond
        pos = grid.get_span(
            col_span=15,  # More than available
            col_start=0,
            left=0.5,
            width=9.0,
        )

        assert pos is not None

    def test_grid_with_bounds(self):
        """Test grid initialized with bounds."""
        bounds = {"left": 1.0, "top": 2.0, "width": 8.0, "height": 5.0}
        grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)

        # Get cell without specifying bounds
        pos = grid.get_cell(col_span=6, col_start=0, row_start=0)

        # Should use bounds from initialization
        assert pos["left"] >= bounds["left"]
        assert pos["top"] >= bounds["top"]
        assert pos["width"] <= bounds["width"]
        assert "height" not in pos  # auto_height=True by default

    def test_grid_get_cell_auto_height(self):
        """Test get_cell with auto_height=True."""
        grid = Grid(columns=12, gap="md")

        # Auto-height (default)
        pos = grid.get_cell(col_span=4, col_start=0, left=0.5, top=2.0, width=9.0, height=4.0)

        assert "left" in pos
        assert "top" in pos
        assert "width" in pos
        assert "height" not in pos  # Should be omitted

    def test_grid_get_cell_fixed_height(self):
        """Test get_cell with auto_height=False."""
        grid = Grid(columns=12, gap="md")

        # Fixed height
        pos = grid.get_cell(
            col_span=4, col_start=0, left=0.5, top=2.0, width=9.0, height=4.0, auto_height=False
        )

        assert "left" in pos
        assert "top" in pos
        assert "width" in pos
        assert "height" in pos  # Should be included

    def test_grid_get_cell_multiple_rows(self):
        """Test get_cell with multiple rows."""
        bounds = {"left": 0.5, "top": 1.8, "width": 9.0, "height": 4.5}
        grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)

        # Row 0
        row0_pos = grid.get_cell(col_span=4, col_start=0, row_start=0)

        # Row 1
        row1_pos = grid.get_cell(col_span=4, col_start=0, row_start=1)

        # Row 1 should be below Row 0
        assert row1_pos["top"] > row0_pos["top"]

    def test_grid_bounds_override(self):
        """Test that explicit parameters override bounds."""
        bounds = {"left": 1.0, "top": 2.0, "width": 8.0, "height": 5.0}
        grid = Grid(columns=12, gap="md", bounds=bounds)

        # Override with explicit left
        pos = grid.get_cell(col_span=6, col_start=0, left=0.5)

        # Should use overridden left, but bounds for others
        assert pos["left"] == 0.5 or pos["left"] > 0.5  # Might be offset by column
        assert pos["top"] >= bounds["top"]
