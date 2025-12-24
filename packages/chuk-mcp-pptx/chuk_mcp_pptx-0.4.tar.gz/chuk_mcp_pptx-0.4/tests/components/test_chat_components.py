"""
Comprehensive tests for chat components.

Tests for all chat bubble and conversation components across platforms.
Coverage target: 90%+
"""

from __future__ import annotations

import pytest
from pptx import Presentation


class TestChatComponentImports:
    """Test that all chat components can be imported."""

    def test_import_generic(self) -> None:
        """Test generic chat imports."""
        from chuk_mcp_pptx.components.chat import ChatMessage, ChatConversation

        assert ChatMessage is not None
        assert ChatConversation is not None

    def test_import_ios(self) -> None:
        """Test iOS chat imports."""
        from chuk_mcp_pptx.components.chat import iMessageBubble, iMessageConversation

        assert iMessageBubble is not None
        assert iMessageConversation is not None

    def test_import_android(self) -> None:
        """Test Android chat imports."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble, AndroidConversation

        assert AndroidMessageBubble is not None
        assert AndroidConversation is not None

    def test_import_whatsapp(self) -> None:
        """Test WhatsApp chat imports."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble, WhatsAppConversation

        assert WhatsAppBubble is not None
        assert WhatsAppConversation is not None

    def test_import_chatgpt(self) -> None:
        """Test ChatGPT chat imports."""
        from chuk_mcp_pptx.components.chat import ChatGPTMessage, ChatGPTConversation

        assert ChatGPTMessage is not None
        assert ChatGPTConversation is not None

    def test_import_slack(self) -> None:
        """Test Slack chat imports."""
        from chuk_mcp_pptx.components.chat import SlackMessage, SlackConversation

        assert SlackMessage is not None
        assert SlackConversation is not None

    def test_import_teams(self) -> None:
        """Test Teams chat imports."""
        from chuk_mcp_pptx.components.chat import TeamsMessage, TeamsConversation

        assert TeamsMessage is not None
        assert TeamsConversation is not None

    def test_import_facebook(self) -> None:
        """Test Facebook Messenger chat imports."""
        from chuk_mcp_pptx.components.chat import (
            FacebookMessengerBubble,
            FacebookMessengerConversation,
        )

        assert FacebookMessengerBubble is not None
        assert FacebookMessengerConversation is not None

    def test_import_aol(self) -> None:
        """Test AIM (AOL) chat imports."""
        from chuk_mcp_pptx.components.chat import AIMBubble, AIMConversation

        assert AIMBubble is not None
        assert AIMConversation is not None

    def test_import_msn(self) -> None:
        """Test MSN chat imports."""
        from chuk_mcp_pptx.components.chat import MSNBubble, MSNConversation

        assert MSNBubble is not None
        assert MSNConversation is not None


class TestiMessageBubble:
    """Tests for iMessage bubble component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_sent(self) -> None:
        """Test initializing a sent message."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Hello!", variant="sent")
        assert msg.text == "Hello!"
        assert msg.variant == "sent"

    def test_init_received(self) -> None:
        """Test initializing a received message."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Hi there!", variant="received")
        assert msg.text == "Hi there!"
        assert msg.variant == "received"

    def test_init_with_timestamp(self) -> None:
        """Test initializing with timestamp."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Hello!", timestamp="10:30 AM")
        assert msg.timestamp == "10:30 AM"

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Hello!", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Hi there!", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_timestamp(self, slide) -> None:
        """Test rendering with timestamp."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Hello!", timestamp="10:30 AM")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_get_bubble_color(self) -> None:
        """Test getting bubble color."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Test", variant="sent")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_text_color_sent(self) -> None:
        """Test getting text color for sent message."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Test", variant="sent")
        color = msg._get_text_color()
        assert color is not None

    def test_get_text_color_received(self) -> None:
        """Test getting text color for received message."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Test", variant="received")
        color = msg._get_text_color()
        assert color is not None

    def test_calculate_bubble_height(self) -> None:
        """Test calculating bubble height."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        msg = iMessageBubble(text="Short text")
        height = msg._calculate_bubble_height(6.0)
        assert height > 0

    def test_calculate_bubble_height_long_text(self) -> None:
        """Test calculating bubble height for long text."""
        from chuk_mcp_pptx.components.chat import iMessageBubble

        long_text = "This is a very long message that should span multiple lines when rendered in the chat bubble component."
        msg = iMessageBubble(text=long_text)
        height = msg._calculate_bubble_height(6.0)
        assert height > 0


class TestiMessageConversation:
    """Tests for iMessage conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import iMessageConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi there!", "variant": "received"},
        ]
        conv = iMessageConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import iMessageConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi!", "variant": "received"},
        ]
        conv = iMessageConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestAndroidMessageBubble:
    """Tests for Android message bubble component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_sent(self) -> None:
        """Test initializing a sent message."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Hello!", variant="sent")
        assert msg.text == "Hello!"
        assert msg.variant == "sent"

    def test_init_received(self) -> None:
        """Test initializing a received message."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Hi!", variant="received")
        assert msg.variant == "received"

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Hello!", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Hi!", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_timestamp(self, slide) -> None:
        """Test rendering with timestamp."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Hello!", timestamp="10:30 AM")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_get_bubble_color(self) -> None:
        """Test getting bubble color."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Test", variant="sent")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_text_color(self) -> None:
        """Test getting text color."""
        from chuk_mcp_pptx.components.chat import AndroidMessageBubble

        msg = AndroidMessageBubble(text="Test", variant="sent")
        color = msg._get_text_color()
        assert color is not None


class TestAndroidConversation:
    """Tests for Android conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import AndroidConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi there!", "variant": "received"},
        ]
        conv = AndroidConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import AndroidConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi!", "variant": "received"},
        ]
        conv = AndroidConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestWhatsAppBubble:
    """Tests for WhatsApp bubble component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_sent(self) -> None:
        """Test initializing a sent message."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Hello!", variant="sent")
        assert msg.text == "Hello!"
        assert msg.variant == "sent"

    def test_init_received(self) -> None:
        """Test initializing a received message."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Hi!", variant="received")
        assert msg.variant == "received"

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Hello!", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Hi!", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_timestamp(self, slide) -> None:
        """Test rendering with timestamp."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Hello!", timestamp="10:30 AM")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_get_bubble_color(self) -> None:
        """Test getting bubble color."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Test", variant="sent")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_text_color(self) -> None:
        """Test getting text color."""
        from chuk_mcp_pptx.components.chat import WhatsAppBubble

        msg = WhatsAppBubble(text="Test", variant="sent")
        color = msg._get_text_color()
        assert color is not None


class TestWhatsAppConversation:
    """Tests for WhatsApp conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import WhatsAppConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi there!", "variant": "received"},
        ]
        conv = WhatsAppConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import WhatsAppConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi!", "variant": "received"},
        ]
        conv = WhatsAppConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestChatGPTMessage:
    """Tests for ChatGPT message component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_user(self) -> None:
        """Test initializing a user message."""
        from chuk_mcp_pptx.components.chat import ChatGPTMessage

        msg = ChatGPTMessage(text="Hello!", variant="user")
        assert msg.text == "Hello!"
        assert msg.variant == "user"

    def test_init_assistant(self) -> None:
        """Test initializing an assistant message."""
        from chuk_mcp_pptx.components.chat import ChatGPTMessage

        msg = ChatGPTMessage(text="Hi!", variant="assistant")
        assert msg.variant == "assistant"

    def test_render_user(self, slide) -> None:
        """Test rendering a user message."""
        from chuk_mcp_pptx.components.chat import ChatGPTMessage

        msg = ChatGPTMessage(text="Hello!", variant="user")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_assistant(self, slide) -> None:
        """Test rendering an assistant message."""
        from chuk_mcp_pptx.components.chat import ChatGPTMessage

        msg = ChatGPTMessage(text="Hi!", variant="assistant")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)


class TestChatGPTConversation:
    """Tests for ChatGPT conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import ChatGPTConversation

        messages = [
            {"text": "What is Python?", "variant": "user"},
            {"text": "Python is a programming language.", "variant": "assistant"},
        ]
        conv = ChatGPTConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import ChatGPTConversation

        messages = [
            {"text": "Hello!", "variant": "user"},
            {"text": "Hi!", "variant": "assistant"},
        ]
        conv = ChatGPTConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestSlackMessage:
    """Tests for Slack message component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a Slack message."""
        from chuk_mcp_pptx.components.chat import SlackMessage

        msg = SlackMessage(text="Hello!", sender="john_doe", timestamp="10:30 AM")
        assert msg.text == "Hello!"
        assert msg.sender == "john_doe"
        assert msg.timestamp == "10:30 AM"

    def test_render(self, slide) -> None:
        """Test rendering a Slack message."""
        from chuk_mcp_pptx.components.chat import SlackMessage

        msg = SlackMessage(text="Hello team!", sender="john_doe", timestamp="10:30 AM")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_reactions(self, slide) -> None:
        """Test rendering with reactions."""
        from chuk_mcp_pptx.components.chat import SlackMessage

        msg = SlackMessage(
            text="Great work!", sender="john_doe", timestamp="10:30 AM", reactions=["👍 3", "❤️ 2"]
        )
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)


class TestSlackConversation:
    """Tests for Slack conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a Slack conversation."""
        from chuk_mcp_pptx.components.chat import SlackConversation

        messages = [
            {"text": "Hello!", "sender": "alice", "timestamp": "10:00 AM"},
            {"text": "Hi!", "sender": "bob", "timestamp": "10:01 AM"},
        ]
        conv = SlackConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a Slack conversation."""
        from chuk_mcp_pptx.components.chat import SlackConversation

        messages = [
            {"text": "Hello!", "sender": "alice", "timestamp": "10:00 AM"},
            {"text": "Hi!", "sender": "bob", "timestamp": "10:01 AM"},
        ]
        conv = SlackConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestTeamsMessage:
    """Tests for Teams message component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a Teams message."""
        from chuk_mcp_pptx.components.chat import TeamsMessage

        msg = TeamsMessage(text="Hello!", sender="John Doe", timestamp="10:30 AM")
        assert msg.text == "Hello!"
        assert msg.sender == "John Doe"
        assert msg.timestamp == "10:30 AM"

    def test_render(self, slide) -> None:
        """Test rendering a Teams message."""
        from chuk_mcp_pptx.components.chat import TeamsMessage

        msg = TeamsMessage(text="Hello team!", sender="John Doe", timestamp="10:30 AM")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_reactions(self, slide) -> None:
        """Test rendering with reactions."""
        from chuk_mcp_pptx.components.chat import TeamsMessage

        msg = TeamsMessage(
            text="Hello!", sender="John Doe", timestamp="10:30 AM", reactions=["👍 5"]
        )
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)


class TestTeamsConversation:
    """Tests for Teams conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a Teams conversation."""
        from chuk_mcp_pptx.components.chat import TeamsConversation

        messages = [
            {"text": "Hello!", "sender": "Alice", "timestamp": "10:00 AM"},
            {"text": "Hi!", "sender": "Bob", "timestamp": "10:01 AM"},
        ]
        conv = TeamsConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a Teams conversation."""
        from chuk_mcp_pptx.components.chat import TeamsConversation

        messages = [
            {"text": "Hello!", "sender": "Alice", "timestamp": "10:00 AM"},
            {"text": "Hi!", "sender": "Bob", "timestamp": "10:01 AM"},
        ]
        conv = TeamsConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestFacebookMessengerBubble:
    """Tests for Facebook Messenger bubble component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_sent(self) -> None:
        """Test initializing a sent message."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Hello!", variant="sent")
        assert msg.text == "Hello!"
        assert msg.variant == "sent"

    def test_init_received(self) -> None:
        """Test initializing a received message."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Hi!", variant="received")
        assert msg.variant == "received"

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Hello!", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Hi!", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_get_bubble_color(self) -> None:
        """Test getting bubble color."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Test", variant="sent")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_text_color(self) -> None:
        """Test getting text color."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Test", variant="sent")
        color = msg._get_text_color()
        assert color is not None

    def test_render_received_with_avatar(self, slide) -> None:
        """Test rendering a received message with avatar."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(
            text="Hey! How are you?", variant="received", avatar_text="JS"
        )
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)
        # Should include avatar shape
        assert len(shapes) >= 2

    def test_render_received_without_avatar(self, slide) -> None:
        """Test rendering a received message without avatar."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Just the message", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_calculate_bubble_height(self) -> None:
        """Test calculating bubble height."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        msg = FacebookMessengerBubble(text="Short text", variant="sent")
        height = msg._calculate_bubble_height(5.0)
        assert height > 0

    def test_calculate_bubble_height_long_text(self) -> None:
        """Test calculating bubble height for long text."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerBubble

        long_text = "This is a very long message that should span multiple lines in the bubble."
        msg = FacebookMessengerBubble(text=long_text, variant="sent")
        height = msg._calculate_bubble_height(5.0)
        assert height > 0


class TestFacebookMessengerConversation:
    """Tests for Facebook Messenger conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi!", "variant": "received"},
        ]
        conv = FacebookMessengerConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi!", "variant": "received"},
        ]
        conv = FacebookMessengerConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_avatars(self, slide) -> None:
        """Test rendering a conversation with avatars."""
        from chuk_mcp_pptx.components.chat import FacebookMessengerConversation

        messages = [
            {"text": "Hey!", "variant": "received", "avatar_text": "JS"},
            {"text": "Hi there!", "variant": "sent"},
            {"text": "What's up?", "variant": "received", "avatar_text": "JS"},
        ]
        conv = FacebookMessengerConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestAIMBubble:
    """Tests for AIM (AOL Instant Messenger) bubble component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_sent(self) -> None:
        """Test initializing a sent message."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Hello!", variant="sent", screen_name="cooluser123")
        assert msg.text == "Hello!"
        assert msg.variant == "sent"

    def test_init_received(self) -> None:
        """Test initializing a received message."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Hi!", variant="received", screen_name="buddy456")
        assert msg.variant == "received"

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Hello!", variant="sent", screen_name="user")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Hi!", variant="received", screen_name="buddy")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_timestamp(self, slide) -> None:
        """Test rendering a message with timestamp."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(
            text="Hey what's up?",
            variant="sent",
            screen_name="xXCoolDude2003Xx",
            timestamp="5:30 PM",
        )
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)
        # Should have more shapes when timestamp is present
        assert len(shapes) >= 2

    def test_calculate_message_height_with_timestamp(self) -> None:
        """Test calculating message height with timestamp."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg_no_ts = AIMBubble(text="Hello!", screen_name="user", variant="sent")
        msg_with_ts = AIMBubble(
            text="Hello!", screen_name="user", variant="sent", timestamp="10:30 AM"
        )

        height_no_ts = msg_no_ts._calculate_message_height(6.0)
        height_with_ts = msg_with_ts._calculate_message_height(6.0)

        # Height should be greater with timestamp
        assert height_with_ts > height_no_ts

    def test_get_screen_name_color_sent(self) -> None:
        """Test getting screen name color for sent message."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Test", screen_name="user", variant="sent")
        color = msg._get_screen_name_color()
        assert color is not None

    def test_get_screen_name_color_received(self) -> None:
        """Test getting screen name color for received message."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Test", screen_name="buddy", variant="received")
        color = msg._get_screen_name_color()
        assert color is not None

    def test_get_text_color(self) -> None:
        """Test getting text color."""
        from chuk_mcp_pptx.components.chat import AIMBubble

        msg = AIMBubble(text="Test", screen_name="user", variant="sent")
        color = msg._get_text_color()
        assert color is not None


class TestAIMConversation:
    """Tests for AIM conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import AIMConversation

        messages = [
            {"text": "ASL?", "variant": "sent", "screen_name": "cooluser123"},
            {"text": "14/f/cali", "variant": "received", "screen_name": "buddy456"},
        ]
        conv = AIMConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import AIMConversation

        messages = [
            {"text": "Hello!", "variant": "sent", "screen_name": "user1"},
            {"text": "Hi!", "variant": "received", "screen_name": "user2"},
        ]
        conv = AIMConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)

    def test_render_with_timestamps(self, slide) -> None:
        """Test rendering a conversation with timestamps."""
        from chuk_mcp_pptx.components.chat import AIMConversation

        messages = [
            {
                "text": "Hey! Want to hang out later?",
                "screen_name": "sk8rgrl2004",
                "variant": "received",
                "timestamp": "5:30 PM",
            },
            {
                "text": "Yeah! Let's go to the mall",
                "screen_name": "xXCoolDude2003Xx",
                "variant": "sent",
                "timestamp": "5:31 PM",
            },
        ]
        conv = AIMConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestMSNBubble:
    """Tests for MSN Messenger bubble component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_sent(self) -> None:
        """Test initializing a sent message."""
        from chuk_mcp_pptx.components.chat import MSNBubble

        msg = MSNBubble(text="Hello!", variant="sent", display_name="John")
        assert msg.text == "Hello!"
        assert msg.variant == "sent"

    def test_init_received(self) -> None:
        """Test initializing a received message."""
        from chuk_mcp_pptx.components.chat import MSNBubble

        msg = MSNBubble(text="Hi!", variant="received", display_name="Jane")
        assert msg.variant == "received"

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat import MSNBubble

        msg = MSNBubble(text="Hello!", variant="sent", display_name="John")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat import MSNBubble

        msg = MSNBubble(text="Hi!", variant="received", display_name="Jane")
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)


class TestMSNConversation:
    """Tests for MSN conversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat import MSNConversation

        messages = [
            {"text": "Hello!", "variant": "sent", "display_name": "John"},
            {"text": "Hi!", "variant": "received", "display_name": "Jane"},
        ]
        conv = MSNConversation(messages=messages)
        assert len(conv.messages) == 2

    def test_render(self, slide) -> None:
        """Test rendering a conversation."""
        from chuk_mcp_pptx.components.chat import MSNConversation

        messages = [
            {"text": "Hello!", "variant": "sent", "display_name": "User1"},
            {"text": "Hi!", "variant": "received", "display_name": "User2"},
        ]
        conv = MSNConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)


class TestGenericChatMessage:
    """Tests for the generic ChatMessage component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!")
        assert msg.text == "Hello!"
        assert msg.sender is None
        assert msg.timestamp is None
        assert msg.variant == "received"
        assert msg.show_avatar is False

    def test_init_sent_variant(self) -> None:
        """Test initializing sent message variant."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="sent")
        assert msg.variant == "sent"

    def test_init_system_variant(self) -> None:
        """Test initializing system message variant."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="System message", variant="system")
        assert msg.variant == "system"

    def test_init_with_sender(self) -> None:
        """Test initializing with sender."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", sender="John Doe")
        assert msg.sender == "John Doe"

    def test_init_with_timestamp(self) -> None:
        """Test initializing with timestamp."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", timestamp="10:30 AM")
        assert msg.timestamp == "10:30 AM"

    def test_init_with_avatar(self) -> None:
        """Test initializing with avatar."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", avatar_text="JD", avatar_icon="user", show_avatar=True)
        assert msg.avatar_text == "JD"
        assert msg.avatar_icon == "user"
        assert msg.show_avatar is True

    def test_get_bubble_color_sent(self) -> None:
        """Test getting bubble color for sent message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="sent")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_bubble_color_received(self) -> None:
        """Test getting bubble color for received message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="received")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_bubble_color_system(self) -> None:
        """Test getting bubble color for system message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="System message", variant="system")
        color = msg._get_bubble_color()
        assert color is not None

    def test_get_text_color_sent(self) -> None:
        """Test getting text color for sent message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="sent")
        color = msg._get_text_color()
        assert color is not None

    def test_get_text_color_received(self) -> None:
        """Test getting text color for received message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="received")
        color = msg._get_text_color()
        assert color is not None

    def test_get_meta_color_sent(self) -> None:
        """Test getting meta color for sent message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="sent")
        color = msg._get_meta_color()
        assert color is not None

    def test_get_meta_color_received(self) -> None:
        """Test getting meta color for received message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="received")
        color = msg._get_meta_color()
        assert color is not None

    def test_calculate_bubble_height_short_text(self) -> None:
        """Test calculating bubble height for short text."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hi")
        height = msg._calculate_bubble_height(4.0)
        assert height > 0

    def test_calculate_bubble_height_long_text(self) -> None:
        """Test calculating bubble height for long text."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(
            text="This is a very long message that should span multiple lines when rendered in the chat bubble."
        )
        height = msg._calculate_bubble_height(4.0)
        assert height > 0

    def test_calculate_bubble_height_with_metadata(self) -> None:
        """Test calculating bubble height with sender and timestamp."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg_no_meta = ChatMessage(text="Hello!")
        msg_with_meta = ChatMessage(text="Hello!", sender="John", timestamp="10:30 AM")

        height_no_meta = msg_no_meta._calculate_bubble_height(4.0)
        height_with_meta = msg_with_meta._calculate_bubble_height(4.0)

        # Height should be greater with metadata
        assert height_with_meta > height_no_meta

    def test_render_received(self, slide) -> None:
        """Test rendering a received message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)
        assert len(shapes) >= 1

    def test_render_sent(self, slide) -> None:
        """Test rendering a sent message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)
        assert len(shapes) >= 1

    def test_render_system(self, slide) -> None:
        """Test rendering a system message."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Welcome to the chat!", variant="system")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)
        assert len(shapes) >= 1

    def test_render_with_sender(self, slide) -> None:
        """Test rendering with sender name."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", sender="John Doe", variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)

    def test_render_with_timestamp(self, slide) -> None:
        """Test rendering with timestamp."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", timestamp="10:30 AM", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)

    def test_render_sent_with_timestamp(self, slide) -> None:
        """Test rendering sent message with timestamp."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", timestamp="10:30 AM", variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)

    def test_render_system_with_sender(self, slide) -> None:
        """Test rendering system message with sender."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="System notification", sender="Bot", variant="system")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)

    def test_render_system_with_timestamp(self, slide) -> None:
        """Test rendering system message with timestamp."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="System notification", timestamp="10:30 AM", variant="system")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)

    def test_render_with_avatar_received(self, slide) -> None:
        """Test rendering received message with avatar."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", avatar_text="JD", show_avatar=True, variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)
        # Should include avatar shape
        assert len(shapes) >= 2

    def test_render_with_avatar_sent(self, slide) -> None:
        """Test rendering sent message with avatar."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", avatar_text="ME", show_avatar=True, variant="sent")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)
        # Should include avatar shape
        assert len(shapes) >= 2

    def test_render_with_avatar_icon(self, slide) -> None:
        """Test rendering with avatar icon."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(text="Hello!", avatar_icon="user", show_avatar=True, variant="received")
        shapes = msg.render(slide, left=1.0, top=2.0, width=6.0)
        assert isinstance(shapes, list)

    def test_render_full_message(self, slide) -> None:
        """Test rendering a full message with all options."""
        from chuk_mcp_pptx.components.chat.generic import ChatMessage

        msg = ChatMessage(
            text="This is a complete test message!",
            sender="John Doe",
            timestamp="10:30 AM",
            avatar_text="JD",
            show_avatar=True,
            variant="received",
        )
        shapes = msg.render(slide, left=1.0, top=2.0, width=7.0)
        assert isinstance(shapes, list)


class TestGenericChatConversation:
    """Tests for the generic ChatConversation component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing a conversation."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi there!", "variant": "received"},
        ]
        conv = ChatConversation(messages=messages)
        assert len(conv.messages) == 2
        assert conv.spacing == 0.25

    def test_init_with_spacing(self) -> None:
        """Test initializing with custom spacing."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [{"text": "Hello!", "variant": "sent"}]
        conv = ChatConversation(messages=messages, spacing=0.5)
        assert conv.spacing == 0.5

    def test_render_simple(self, slide) -> None:
        """Test rendering a simple conversation."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi there!", "variant": "received"},
        ]
        conv = ChatConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=6.0)
        assert isinstance(shapes, list)
        # Should have shapes from both messages
        assert len(shapes) >= 2

    def test_render_with_metadata(self, slide) -> None:
        """Test rendering conversation with metadata."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Hello!", "sender": "John", "timestamp": "10:30 AM", "variant": "sent"},
            {"text": "Hi!", "sender": "Jane", "timestamp": "10:31 AM", "variant": "received"},
        ]
        conv = ChatConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=6.0)
        assert isinstance(shapes, list)

    def test_render_with_avatars(self, slide) -> None:
        """Test rendering conversation with avatars."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Hello!", "avatar_text": "JD", "show_avatar": True, "variant": "sent"},
            {"text": "Hi!", "avatar_text": "JN", "show_avatar": True, "variant": "received"},
        ]
        conv = ChatConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=7.0)
        assert isinstance(shapes, list)

    def test_render_mixed_variants(self, slide) -> None:
        """Test rendering conversation with mixed variants."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Welcome!", "variant": "system"},
            {"text": "Hello!", "variant": "sent"},
            {"text": "Hi!", "variant": "received"},
        ]
        conv = ChatConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=6.0)
        assert isinstance(shapes, list)

    def test_render_complex_conversation(self, slide) -> None:
        """Test rendering a complex conversation."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Welcome to support!", "variant": "system"},
            {
                "text": "I need help with my order",
                "sender": "John",
                "avatar_text": "JD",
                "timestamp": "10:30 AM",
                "variant": "sent",
                "show_avatar": True,
            },
            {
                "text": "I'd be happy to help! What's your order number?",
                "sender": "Support",
                "avatar_text": "SA",
                "timestamp": "10:31 AM",
                "variant": "received",
                "show_avatar": True,
            },
            {"text": "#12345", "timestamp": "10:32 AM", "variant": "sent"},
        ]
        conv = ChatConversation(messages=messages, spacing=0.3)
        shapes = conv.render(slide, left=1.0, top=1.0, width=7.0)
        assert isinstance(shapes, list)

    def test_render_empty_conversation(self, slide) -> None:
        """Test rendering an empty conversation."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        conv = ChatConversation(messages=[])
        shapes = conv.render(slide, left=1.0, top=1.5, width=6.0)
        assert isinstance(shapes, list)
        assert len(shapes) == 0

    def test_render_single_message(self, slide) -> None:
        """Test rendering a single message conversation."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [{"text": "Single message"}]
        conv = ChatConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=6.0)
        assert isinstance(shapes, list)
        assert len(shapes) >= 1

    def test_render_with_avatar_icon(self, slide) -> None:
        """Test rendering conversation with avatar icons."""
        from chuk_mcp_pptx.components.chat.generic import ChatConversation

        messages = [
            {"text": "Hello!", "avatar_icon": "user", "show_avatar": True, "variant": "received"},
        ]
        conv = ChatConversation(messages=messages)
        shapes = conv.render(slide, left=1.0, top=1.5, width=6.0)
        assert isinstance(shapes, list)
