"""
Tests for container components.

Tests for device mockups (iPhone, Samsung) and window containers
(browser, macOS, Windows) along with generic chat containers.
"""

import pytest
from pptx import Presentation


class TestContainerImports:
    """Tests for container component imports."""

    def test_import_iphone(self) -> None:
        """Test importing iPhoneContainer."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        assert iPhoneContainer is not None

    def test_import_samsung(self) -> None:
        """Test importing SamsungContainer."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        assert SamsungContainer is not None

    def test_import_browser(self) -> None:
        """Test importing BrowserWindow."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        assert BrowserWindow is not None

    def test_import_macos(self) -> None:
        """Test importing MacOSWindow."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        assert MacOSWindow is not None

    def test_import_windows(self) -> None:
        """Test importing WindowsWindow."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        assert WindowsWindow is not None

    def test_import_generic(self) -> None:
        """Test importing ChatContainer."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        assert ChatContainer is not None


class TestiPhoneContainer:
    """Tests for iPhone container component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        assert container.title == "iPhone"
        assert container.show_notch is True
        assert container.variant == "pro"

    def test_init_with_title(self) -> None:
        """Test initializing with title."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer(title="iMessage")
        assert container.title == "iMessage"

    def test_init_without_notch(self) -> None:
        """Test initializing without notch."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer(show_notch=False)
        assert container.show_notch is False

    def test_init_variants(self) -> None:
        """Test initializing with different variants."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        for variant in ["pro", "pro-max", "standard"]:
            container = iPhoneContainer(variant=variant)
            assert container.variant == variant

    def test_get_device_color_default(self) -> None:
        """Test getting device color without theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        color = container._get_device_color()
        assert color is not None

    def test_get_device_color_light_theme(self) -> None:
        """Test getting device color with light theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        container = iPhoneContainer(theme=theme)
        color = container._get_device_color()
        assert color is not None

    def test_get_device_color_dark_theme(self) -> None:
        """Test getting device color with dark theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = iPhoneContainer(theme=theme)
        color = container._get_device_color()
        assert color is not None

    def test_get_screen_bg_color_default(self) -> None:
        """Test getting screen background color without theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        color = container._get_screen_bg_color()
        assert color is not None

    def test_get_screen_bg_color_with_theme(self) -> None:
        """Test getting screen background color with theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [240, 240, 240]}}}
        container = iPhoneContainer(theme=theme)
        color = container._get_screen_bg_color()
        assert color is not None

    def test_get_text_color_default(self) -> None:
        """Test getting text color without theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        color = container._get_text_color()
        assert color is not None

    def test_get_text_color_with_theme(self) -> None:
        """Test getting text color with theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"foreground": {"DEFAULT": [50, 50, 50]}}}
        container = iPhoneContainer(theme=theme)
        color = container._get_text_color()
        assert color is not None

    def test_is_dark_mode_false(self) -> None:
        """Test dark mode detection for light theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        container = iPhoneContainer(theme=theme)
        assert container._is_dark_mode() is False

    def test_is_dark_mode_true(self) -> None:
        """Test dark mode detection for dark theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = iPhoneContainer(theme=theme)
        assert container._is_dark_mode() is True

    def test_is_dark_mode_no_theme(self) -> None:
        """Test dark mode detection without theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        assert container._is_dark_mode() is False

    def test_render(self, slide) -> None:
        """Test rendering an iPhone container."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        assert "left" in content_area
        assert "top" in content_area
        assert "width" in content_area
        assert "height" in content_area

    def test_render_with_notch(self, slide) -> None:
        """Test rendering with notch."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer(show_notch=True)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_without_notch(self, slide) -> None:
        """Test rendering without notch."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer(show_notch=False)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        container = iPhoneContainer()
        content_area = container.render(slide, left=2.0, top=2.0, width=4.0, height=8.0)
        assert isinstance(content_area, dict)

    def test_render_dark_theme(self, slide) -> None:
        """Test rendering with dark theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = iPhoneContainer(theme=theme)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_light_theme(self, slide) -> None:
        """Test rendering with light theme."""
        from chuk_mcp_pptx.components.containers import iPhoneContainer

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        container = iPhoneContainer(theme=theme)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)


class TestSamsungContainer:
    """Tests for Samsung container component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer()
        assert container.title == "Samsung"

    def test_init_with_title(self) -> None:
        """Test initializing with title."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer(title="Messages")
        assert container.title == "Messages"

    def test_get_device_color_default(self) -> None:
        """Test getting device color without theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer()
        color = container._get_device_color()
        assert color is not None

    def test_get_device_color_light_theme(self) -> None:
        """Test getting device color with light theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        container = SamsungContainer(theme=theme)
        color = container._get_device_color()
        assert color is not None

    def test_get_device_color_dark_theme(self) -> None:
        """Test getting device color with dark theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = SamsungContainer(theme=theme)
        color = container._get_device_color()
        assert color is not None

    def test_get_screen_bg_color_default(self) -> None:
        """Test getting screen background color without theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer()
        color = container._get_screen_bg_color()
        assert color is not None

    def test_get_screen_bg_color_with_theme(self) -> None:
        """Test getting screen background color with theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        theme = {"colors": {"background": {"DEFAULT": [240, 240, 240]}}}
        container = SamsungContainer(theme=theme)
        color = container._get_screen_bg_color()
        assert color is not None

    def test_get_text_color_default(self) -> None:
        """Test getting text color without theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer()
        color = container._get_text_color()
        assert color is not None

    def test_get_text_color_with_theme(self) -> None:
        """Test getting text color with theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        theme = {"colors": {"foreground": {"DEFAULT": [50, 50, 50]}}}
        container = SamsungContainer(theme=theme)
        color = container._get_text_color()
        assert color is not None

    def test_is_dark_mode(self) -> None:
        """Test dark mode detection."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = SamsungContainer(theme=theme)
        assert container._is_dark_mode() is True

    def test_render(self, slide) -> None:
        """Test rendering a Samsung container."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer()
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        assert "left" in content_area
        assert "top" in content_area
        assert "width" in content_area
        assert "height" in content_area

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        container = SamsungContainer()
        content_area = container.render(slide, left=2.0, top=2.0, width=4.0, height=8.0)
        assert isinstance(content_area, dict)

    def test_render_dark_theme(self, slide) -> None:
        """Test rendering with dark theme."""
        from chuk_mcp_pptx.components.containers import SamsungContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = SamsungContainer(theme=theme)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)


class TestBrowserWindow:
    """Tests for browser window component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        assert window.title == "Browser"
        assert window.url == "example.com"
        assert window.browser_type == "chrome"
        assert window.show_tabs is False

    def test_init_with_url(self) -> None:
        """Test initializing with URL."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(url="https://google.com")
        assert window.url == "https://google.com"

    def test_init_with_title(self) -> None:
        """Test initializing with title."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(title="My Browser")
        assert window.title == "My Browser"

    def test_get_chrome_color_default(self) -> None:
        """Test getting chrome color without theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        color = window._get_chrome_color()
        assert color is not None

    def test_get_chrome_color_light_theme(self) -> None:
        """Test getting chrome color with light theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        window = BrowserWindow(theme=theme)
        color = window._get_chrome_color()
        assert color is not None

    def test_get_chrome_color_dark_theme(self) -> None:
        """Test getting chrome color with dark theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = BrowserWindow(theme=theme)
        color = window._get_chrome_color()
        assert color is not None

    def test_get_content_bg_color_default(self) -> None:
        """Test getting content background color without theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        color = window._get_content_bg_color()
        assert color is not None

    def test_get_content_bg_color_with_theme(self) -> None:
        """Test getting content background color with theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [240, 240, 240]}}}
        window = BrowserWindow(theme=theme)
        color = window._get_content_bg_color()
        assert color is not None

    def test_get_text_color_default(self) -> None:
        """Test getting text color without theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        color = window._get_text_color()
        assert color is not None

    def test_get_text_color_with_theme(self) -> None:
        """Test getting text color with theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"foreground": {"DEFAULT": [50, 50, 50]}}}
        window = BrowserWindow(theme=theme)
        color = window._get_text_color()
        assert color is not None

    def test_is_dark_mode(self) -> None:
        """Test dark mode detection."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = BrowserWindow(theme=theme)
        assert window._is_dark_mode() is True

    def test_render(self, slide) -> None:
        """Test rendering a browser window."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        assert "left" in content_area
        assert "top" in content_area
        assert "width" in content_area
        assert "height" in content_area

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        content_area = window.render(slide, left=0.5, top=0.5, width=8.0, height=5.0)
        assert isinstance(content_area, dict)

    def test_render_dark_theme(self, slide) -> None:
        """Test rendering with dark theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = BrowserWindow(theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_init_safari(self) -> None:
        """Test initializing Safari browser."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(browser_type="safari")
        assert window.browser_type == "safari"

    def test_init_firefox(self) -> None:
        """Test initializing Firefox browser."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(browser_type="firefox")
        assert window.browser_type == "firefox"

    def test_init_with_tabs(self) -> None:
        """Test initializing with tabs enabled."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(show_tabs=True)
        assert window.show_tabs is True

    def test_render_safari(self, slide) -> None:
        """Test rendering Safari browser with macOS controls."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(browser_type="safari", title="Apple")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_safari_with_tabs(self, slide) -> None:
        """Test rendering Safari browser with tabs."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(browser_type="safari", show_tabs=True, title="Apple Tab")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_with_tabs(self, slide) -> None:
        """Test rendering browser with tab bar."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(show_tabs=True, title="My Tab")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        # Content area should be smaller due to tab bar
        assert content_area["top"] > 1.0

    def test_render_chrome_with_tabs(self, slide) -> None:
        """Test rendering Chrome browser with tabs."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(browser_type="chrome", show_tabs=True, title="Google Chrome")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_firefox_with_tabs(self, slide) -> None:
        """Test rendering Firefox browser with tabs."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow(browser_type="firefox", show_tabs=True, title="Firefox Tab")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_get_address_bar_color(self) -> None:
        """Test getting address bar color."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        window = BrowserWindow()
        color = window._get_address_bar_color()
        assert color is not None

    def test_get_address_bar_color_dark(self) -> None:
        """Test getting address bar color in dark mode."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = BrowserWindow(theme=theme)
        color = window._get_address_bar_color()
        assert color is not None

    def test_render_safari_dark_theme(self, slide) -> None:
        """Test rendering Safari with dark theme."""
        from chuk_mcp_pptx.components.containers import BrowserWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = BrowserWindow(browser_type="safari", theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)


class TestMacOSWindow:
    """Tests for macOS window component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow()
        assert window.title == "Application"
        assert window.app_icon is None
        assert window.show_toolbar is False

    def test_init_with_title(self) -> None:
        """Test initializing with title."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow(title="My App")
        assert window.title == "My App"

    def test_get_titlebar_color_default(self) -> None:
        """Test getting titlebar color without theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow()
        color = window._get_titlebar_color()
        assert color is not None

    def test_get_titlebar_color_light_theme(self) -> None:
        """Test getting titlebar color with light theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        window = MacOSWindow(theme=theme)
        color = window._get_titlebar_color()
        assert color is not None

    def test_get_titlebar_color_dark_theme(self) -> None:
        """Test getting titlebar color with dark theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = MacOSWindow(theme=theme)
        color = window._get_titlebar_color()
        assert color is not None

    def test_get_content_bg_color_default(self) -> None:
        """Test getting content background color without theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow()
        color = window._get_content_bg_color()
        assert color is not None

    def test_get_content_bg_color_with_theme(self) -> None:
        """Test getting content background color with theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"background": {"DEFAULT": [240, 240, 240]}}}
        window = MacOSWindow(theme=theme)
        color = window._get_content_bg_color()
        assert color is not None

    def test_get_text_color_default(self) -> None:
        """Test getting text color without theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow()
        color = window._get_text_color()
        assert color is not None

    def test_get_text_color_with_theme(self) -> None:
        """Test getting text color with theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"foreground": {"DEFAULT": [50, 50, 50]}}}
        window = MacOSWindow(theme=theme)
        color = window._get_text_color()
        assert color is not None

    def test_is_dark_mode(self) -> None:
        """Test dark mode detection."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = MacOSWindow(theme=theme)
        assert window._is_dark_mode() is True

    def test_render(self, slide) -> None:
        """Test rendering a macOS window."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow()
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        assert "left" in content_area
        assert "top" in content_area
        assert "width" in content_area
        assert "height" in content_area

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow()
        content_area = window.render(slide, left=0.5, top=0.5, width=8.0, height=5.0)
        assert isinstance(content_area, dict)

    def test_render_dark_theme(self, slide) -> None:
        """Test rendering with dark theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = MacOSWindow(theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_init_with_toolbar(self) -> None:
        """Test initializing with toolbar."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow(show_toolbar=True)
        assert window.show_toolbar is True

    def test_init_with_app_icon(self) -> None:
        """Test initializing with app icon."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow(title="Messages", app_icon="💬")
        assert window.app_icon == "💬"

    def test_render_with_toolbar(self, slide) -> None:
        """Test rendering with toolbar enabled."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow(show_toolbar=True, title="Finder")
        content_area = window.render(slide, left=1.0, top=1.0, width=7.0, height=5.0)
        assert isinstance(content_area, dict)
        # Content top should be lower due to toolbar
        assert content_area["top"] > 1.0 + 0.45  # titlebar_height + some

    def test_render_with_toolbar_dark_theme(self, slide) -> None:
        """Test rendering with toolbar in dark theme."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = MacOSWindow(show_toolbar=True, title="Finder", theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0, width=7.0, height=5.0)
        assert isinstance(content_area, dict)

    def test_render_with_app_icon(self, slide) -> None:
        """Test rendering with app icon."""
        from chuk_mcp_pptx.components.containers import MacOSWindow

        window = MacOSWindow(title="Messages", app_icon="💬")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)


class TestWindowsWindow:
    """Tests for Windows window component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow()
        assert window.title == "Application"
        assert window.app_icon is None
        assert window.show_menubar is False

    def test_init_with_title(self) -> None:
        """Test initializing with title."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow(title="My App")
        assert window.title == "My App"

    def test_get_titlebar_color_default(self) -> None:
        """Test getting titlebar color without theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow()
        color = window._get_titlebar_color()
        assert color is not None

    def test_get_titlebar_color_light_theme(self) -> None:
        """Test getting titlebar color with light theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        window = WindowsWindow(theme=theme)
        color = window._get_titlebar_color()
        assert color is not None

    def test_get_titlebar_color_dark_theme(self) -> None:
        """Test getting titlebar color with dark theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = WindowsWindow(theme=theme)
        color = window._get_titlebar_color()
        assert color is not None

    def test_get_content_bg_color_default(self) -> None:
        """Test getting content background color without theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow()
        color = window._get_content_bg_color()
        assert color is not None

    def test_get_content_bg_color_with_theme(self) -> None:
        """Test getting content background color with theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [240, 240, 240]}}}
        window = WindowsWindow(theme=theme)
        color = window._get_content_bg_color()
        assert color is not None

    def test_get_text_color_default(self) -> None:
        """Test getting text color without theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow()
        color = window._get_text_color()
        assert color is not None

    def test_get_text_color_with_theme(self) -> None:
        """Test getting text color with theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"foreground": {"DEFAULT": [50, 50, 50]}}}
        window = WindowsWindow(theme=theme)
        color = window._get_text_color()
        assert color is not None

    def test_is_dark_mode(self) -> None:
        """Test dark mode detection."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = WindowsWindow(theme=theme)
        assert window._is_dark_mode() is True

    def test_render(self, slide) -> None:
        """Test rendering a Windows window."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow()
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        assert "left" in content_area
        assert "top" in content_area
        assert "width" in content_area
        assert "height" in content_area

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow()
        content_area = window.render(slide, left=0.5, top=0.5, width=8.0, height=5.0)
        assert isinstance(content_area, dict)

    def test_render_dark_theme(self, slide) -> None:
        """Test rendering with dark theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = WindowsWindow(theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_init_with_menubar(self) -> None:
        """Test initializing with menubar."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow(show_menubar=True)
        assert window.show_menubar is True

    def test_init_with_app_icon(self) -> None:
        """Test initializing with app icon."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow(title="Microsoft Teams", app_icon="👥")
        assert window.app_icon == "👥"

    def test_render_with_menubar(self, slide) -> None:
        """Test rendering with menubar enabled."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow(show_menubar=True, title="Notepad")
        content_area = window.render(slide, left=1.0, top=1.0, width=7.0, height=5.0)
        assert isinstance(content_area, dict)
        # Content top should be lower due to menubar
        assert content_area["top"] > 1.0 + 0.4  # titlebar_height + menubar

    def test_render_with_menubar_light_theme(self, slide) -> None:
        """Test rendering with menubar in light theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        window = WindowsWindow(show_menubar=True, title="Explorer", theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0, width=7.0, height=5.0)
        assert isinstance(content_area, dict)

    def test_render_with_menubar_dark_theme(self, slide) -> None:
        """Test rendering with menubar in dark theme."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        window = WindowsWindow(show_menubar=True, title="Explorer", theme=theme)
        content_area = window.render(slide, left=1.0, top=1.0, width=7.0, height=5.0)
        assert isinstance(content_area, dict)

    def test_render_with_app_icon(self, slide) -> None:
        """Test rendering with app icon."""
        from chuk_mcp_pptx.components.containers import WindowsWindow

        window = WindowsWindow(title="Microsoft Teams", app_icon="👥")
        content_area = window.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)


class TestChatContainer:
    """Tests for generic chat container component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        assert container.title is None
        assert container.show_header is False
        assert container.show_border is True
        assert container.variant == "default"

    def test_init_with_title(self) -> None:
        """Test initializing with title."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer(title="Conversation")
        assert container.title == "Conversation"

    def test_init_with_header(self) -> None:
        """Test initializing with header enabled."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer(show_header=True, title="Chat")
        assert container.show_header is True

    def test_init_without_border(self) -> None:
        """Test initializing without border."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer(show_border=False)
        assert container.show_border is False

    def test_init_variants(self) -> None:
        """Test initializing with different variants."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        for variant in ["default", "outlined", "filled"]:
            container = ChatContainer(variant=variant)
            assert container.variant == variant

    def test_is_dark_mode_false(self) -> None:
        """Test dark mode detection for light theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        container = ChatContainer(theme=theme)
        assert container._is_dark_mode() is False

    def test_is_dark_mode_true(self) -> None:
        """Test dark mode detection for dark theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        container = ChatContainer(theme=theme)
        assert container._is_dark_mode() is True

    def test_is_dark_mode_no_theme(self) -> None:
        """Test dark mode detection without theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        assert container._is_dark_mode() is False

    def test_get_container_bg_color_default(self) -> None:
        """Test getting container background color without theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        color = container._get_container_bg_color()
        assert color is not None

    def test_get_container_bg_color_filled(self) -> None:
        """Test getting container background color with filled variant."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"card": {"DEFAULT": [250, 250, 250]}}}
        container = ChatContainer(variant="filled", theme=theme)
        color = container._get_container_bg_color()
        assert color is not None

    def test_get_container_bg_color_with_theme(self) -> None:
        """Test getting container background color with theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"background": {"DEFAULT": [240, 240, 240]}}}
        container = ChatContainer(theme=theme)
        color = container._get_container_bg_color()
        assert color is not None

    def test_get_border_color_default(self) -> None:
        """Test getting border color without theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        color = container._get_border_color()
        assert color is not None

    def test_get_border_color_with_theme(self) -> None:
        """Test getting border color with theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"border": {"DEFAULT": [200, 200, 200]}}}
        container = ChatContainer(theme=theme)
        color = container._get_border_color()
        assert color is not None

    def test_get_header_bg_color_default(self) -> None:
        """Test getting header background color without theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        color = container._get_header_bg_color()
        assert color is not None

    def test_get_header_bg_color_with_theme(self) -> None:
        """Test getting header background color with theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"muted": {"DEFAULT": [200, 200, 200]}}}
        container = ChatContainer(theme=theme)
        color = container._get_header_bg_color()
        assert color is not None

    def test_get_text_color_default(self) -> None:
        """Test getting text color without theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        color = container._get_text_color()
        assert color is not None

    def test_get_text_color_with_theme(self) -> None:
        """Test getting text color with theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"foreground": {"DEFAULT": [50, 50, 50]}}}
        container = ChatContainer(theme=theme)
        color = container._get_text_color()
        assert color is not None

    def test_render_simple(self, slide) -> None:
        """Test rendering a simple container."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)
        assert "left" in content_area
        assert "top" in content_area
        assert "width" in content_area
        assert "height" in content_area

    def test_render_with_header(self, slide) -> None:
        """Test rendering with header."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer(show_header=True, title="Chat")
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_without_border(self, slide) -> None:
        """Test rendering without border."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer(show_border=False)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        content_area = container.render(slide, left=0.5, top=0.5, width=8.0, height=6.0)
        assert isinstance(content_area, dict)

    def test_render_with_theme(self, slide) -> None:
        """Test rendering with theme."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        theme = {"colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        container = ChatContainer(theme=theme)
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_render_filled_variant(self, slide) -> None:
        """Test rendering with filled variant."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer(variant="filled")
        content_area = container.render(slide, left=1.0, top=1.0)
        assert isinstance(content_area, dict)

    def test_content_area_values(self, slide) -> None:
        """Test that content area has reasonable values."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container = ChatContainer()
        content_area = container.render(slide, left=1.0, top=1.0, width=6.0, height=5.0)

        assert content_area["left"] > 1.0  # Has padding
        assert content_area["top"] > 1.0  # Has padding
        assert content_area["width"] < 6.0  # Has padding
        assert content_area["height"] < 5.0  # Has padding

    def test_content_area_with_header(self, slide) -> None:
        """Test that header reduces content area."""
        from chuk_mcp_pptx.components.containers import ChatContainer

        container_no_header = ChatContainer(show_header=False)
        container_with_header = ChatContainer(show_header=True, title="Chat")

        area_no_header = container_no_header.render(slide, left=1.0, top=1.0, width=6.0, height=5.0)
        area_with_header = container_with_header.render(
            slide, left=1.0, top=1.0, width=6.0, height=5.0
        )

        # Header should reduce content height
        assert area_with_header["height"] < area_no_header["height"]
        # Header should push content down
        assert area_with_header["top"] > area_no_header["top"]
