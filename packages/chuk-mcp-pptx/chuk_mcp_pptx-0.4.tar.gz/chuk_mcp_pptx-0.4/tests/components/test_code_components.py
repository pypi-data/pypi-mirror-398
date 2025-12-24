"""
Tests for code block components.

Tests for CodeBlock, InlineCode, and Terminal components
used to display code snippets in presentations.
"""

import pytest
from pptx import Presentation


class TestCodeBlockImports:
    """Tests for code component imports."""

    def test_import_code_block(self) -> None:
        """Test importing CodeBlock."""
        from chuk_mcp_pptx.components.code import CodeBlock

        assert CodeBlock is not None

    def test_import_inline_code(self) -> None:
        """Test importing InlineCode."""
        from chuk_mcp_pptx.components.code import InlineCode

        assert InlineCode is not None

    def test_import_terminal(self) -> None:
        """Test importing Terminal."""
        from chuk_mcp_pptx.components.code import Terminal

        assert Terminal is not None


class TestCodeBlock:
    """Tests for CodeBlock component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="print('hello')")
        assert code.code == "print('hello')"
        assert code.language == "text"
        assert code.show_line_numbers is False

    def test_init_with_language(self) -> None:
        """Test initializing with language."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="def foo(): pass", language="python")
        assert code.language == "python"

    def test_init_with_line_numbers(self) -> None:
        """Test initializing with line numbers."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="print('hello')", show_line_numbers=True)
        assert code.show_line_numbers is True

    def test_init_language_lowercased(self) -> None:
        """Test that language is lowercased."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="code", language="PYTHON")
        assert code.language == "python"

    def test_format_code_no_line_numbers(self) -> None:
        """Test formatting code without line numbers."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="print('hello')\nprint('world')")
        formatted = code._format_code()
        assert formatted == "print('hello')\nprint('world')"

    def test_format_code_with_line_numbers(self) -> None:
        """Test formatting code with line numbers."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="print('hello')\nprint('world')", show_line_numbers=True)
        formatted = code._format_code()
        assert "1 │" in formatted
        assert "2 │" in formatted
        assert "print('hello')" in formatted
        assert "print('world')" in formatted

    def test_format_code_with_many_lines(self) -> None:
        """Test formatting code with many lines for proper padding."""
        from chuk_mcp_pptx.components.code import CodeBlock

        lines = "\n".join([f"line {i}" for i in range(15)])
        code = CodeBlock(code=lines, show_line_numbers=True)
        formatted = code._format_code()
        # With 15 lines, line numbers should be right-justified to 2 chars
        assert " 1 │" in formatted
        assert "15 │" in formatted

    def test_get_language_color(self) -> None:
        """Test getting language color."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="code", language="python")
        color = code.get_language_color()
        assert color is not None
        assert isinstance(color, tuple)
        assert len(color) == 3

    def test_get_language_color_unknown(self) -> None:
        """Test getting language color for unknown language."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="code", language="unknown_lang_xyz")
        color = code.get_language_color()
        assert color is not None

    def test_render(self, slide) -> None:
        """Test rendering a code block."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="print('hello')", language="python")
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_with_line_numbers(self, slide) -> None:
        """Test rendering with line numbers."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(
            code="def foo():\n    return 42", language="python", show_line_numbers=True
        )
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(code="code")
        shape = code.render(slide, left=0.5, top=0.5, width=8.0, height=4.0)
        assert shape is not None

    def test_render_multiline(self, slide) -> None:
        """Test rendering multiline code."""
        from chuk_mcp_pptx.components.code import CodeBlock

        code = CodeBlock(
            code="def hello():\n    print('hello')\n    return True", language="python"
        )
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_with_theme_light(self, slide) -> None:
        """Test rendering with light theme."""
        from chuk_mcp_pptx.components.code import CodeBlock

        theme = {"mode": "light", "colors": {"background": {"DEFAULT": [255, 255, 255]}}}
        code = CodeBlock(code="code", theme=theme)
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_with_theme_dark(self, slide) -> None:
        """Test rendering with dark theme."""
        from chuk_mcp_pptx.components.code import CodeBlock

        theme = {"mode": "dark", "colors": {"background": {"DEFAULT": [30, 30, 30]}}}
        code = CodeBlock(code="code", theme=theme)
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_various_languages(self, slide) -> None:
        """Test rendering with various languages."""
        from chuk_mcp_pptx.components.code import CodeBlock

        languages = ["python", "javascript", "typescript", "java", "go", "rust"]
        for lang in languages:
            code = CodeBlock(code="code", language=lang)
            shape = code.render(slide, left=1.0, top=1.0)
            assert shape is not None


class TestInlineCode:
    """Tests for InlineCode component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init(self) -> None:
        """Test initializing inline code."""
        from chuk_mcp_pptx.components.code import InlineCode

        code = InlineCode(code="variable")
        assert code.code == "variable"

    def test_init_with_theme(self) -> None:
        """Test initializing with theme."""
        from chuk_mcp_pptx.components.code import InlineCode

        theme = {"colors": {"muted": {"DEFAULT": [200, 200, 200]}}}
        code = InlineCode(code="code", theme=theme)
        assert code.theme == theme

    def test_render(self, slide) -> None:
        """Test rendering inline code."""
        from chuk_mcp_pptx.components.code import InlineCode

        code = InlineCode(code="variable_name")
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_auto_width(self, slide) -> None:
        """Test rendering with auto-calculated width."""
        from chuk_mcp_pptx.components.code import InlineCode

        code = InlineCode(code="short")
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_auto_width_long(self, slide) -> None:
        """Test rendering with auto-calculated width for long code."""
        from chuk_mcp_pptx.components.code import InlineCode

        code = InlineCode(code="this_is_a_very_long_variable_name_that_should_be_wide")
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_explicit_width(self, slide) -> None:
        """Test rendering with explicit width."""
        from chuk_mcp_pptx.components.code import InlineCode

        code = InlineCode(code="code")
        shape = code.render(slide, left=1.0, top=1.0, width=3.0)
        assert shape is not None

    def test_render_custom_height(self, slide) -> None:
        """Test rendering with custom height."""
        from chuk_mcp_pptx.components.code import InlineCode

        code = InlineCode(code="code")
        shape = code.render(slide, left=1.0, top=1.0, height=0.5)
        assert shape is not None

    def test_render_with_theme(self, slide) -> None:
        """Test rendering with theme."""
        from chuk_mcp_pptx.components.code import InlineCode

        theme = {"colors": {"muted": {"DEFAULT": [200, 200, 200]}}}
        code = InlineCode(code="code", theme=theme)
        shape = code.render(slide, left=1.0, top=1.0)
        assert shape is not None


class TestTerminal:
    """Tests for Terminal component."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_init_defaults(self) -> None:
        """Test initializing with default values."""
        from chuk_mcp_pptx.components.code import Terminal

        terminal = Terminal(output="ls -la")
        assert terminal.code == "ls -la"  # Stored in parent's code attribute
        assert terminal.prompt == "$"
        assert terminal.language == "shell"

    def test_init_with_prompt(self) -> None:
        """Test initializing with custom prompt."""
        from chuk_mcp_pptx.components.code import Terminal

        terminal = Terminal(output="command", prompt="#")
        assert terminal.prompt == "#"

    def test_init_with_theme(self) -> None:
        """Test initializing with theme."""
        from chuk_mcp_pptx.components.code import Terminal

        theme = {"colors": {}}
        terminal = Terminal(output="cmd", theme=theme)
        assert terminal.theme == theme

    def test_render(self, slide) -> None:
        """Test rendering terminal output."""
        from chuk_mcp_pptx.components.code import Terminal

        terminal = Terminal(output="ls -la\ndrwxr-xr-x  5 user  staff  160 Dec 15 10:00 .")
        shape = terminal.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_multiline(self, slide) -> None:
        """Test rendering multiline terminal output."""
        from chuk_mcp_pptx.components.code import Terminal

        output = """git status
On branch main
nothing to commit, working tree clean"""
        terminal = Terminal(output=output)
        shape = terminal.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_with_indented_output(self, slide) -> None:
        """Test rendering with indented output lines."""
        from chuk_mcp_pptx.components.code import Terminal

        output = """ls -la
  file1.txt
  file2.txt
cd /home"""
        terminal = Terminal(output=output)
        shape = terminal.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_custom_size(self, slide) -> None:
        """Test rendering with custom size."""
        from chuk_mcp_pptx.components.code import Terminal

        terminal = Terminal(output="echo hello")
        shape = terminal.render(slide, left=0.5, top=0.5, width=7.0, height=4.0)
        assert shape is not None

    def test_render_custom_prompt(self, slide) -> None:
        """Test rendering with custom prompt."""
        from chuk_mcp_pptx.components.code import Terminal

        terminal = Terminal(output="whoami\nroot", prompt="#")
        shape = terminal.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_empty_line(self, slide) -> None:
        """Test rendering with empty lines."""
        from chuk_mcp_pptx.components.code import Terminal

        output = """echo hello

echo world"""
        terminal = Terminal(output=output)
        shape = terminal.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_render_with_theme(self, slide) -> None:
        """Test rendering with theme."""
        from chuk_mcp_pptx.components.code import Terminal

        theme = {"colors": {}}
        terminal = Terminal(output="cmd", theme=theme)
        shape = terminal.render(slide, left=1.0, top=1.0)
        assert shape is not None
