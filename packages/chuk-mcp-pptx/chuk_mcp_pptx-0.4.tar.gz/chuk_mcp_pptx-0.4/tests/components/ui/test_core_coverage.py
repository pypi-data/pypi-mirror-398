"""
Additional tests for core components to improve coverage.
Focuses on edge cases and less-used code paths.
"""

import pytest
from pptx import Presentation

from chuk_mcp_pptx.components.core.progress import ProgressBar
from chuk_mcp_pptx.components.core.tile import Tile, IconTile, ValueTile
from chuk_mcp_pptx.components.core.timeline import Timeline
from chuk_mcp_pptx.components.core.avatar import Avatar, AvatarWithLabel, AvatarGroup
from chuk_mcp_pptx.components.core.button import IconButton, ButtonGroup
from chuk_mcp_pptx.components.core.badge import DotBadge, CountBadge
from chuk_mcp_pptx.components.core.card import Card, MetricCard


@pytest.fixture
def presentation():
    """Create a presentation for rendering tests."""
    return Presentation()


@pytest.fixture
def slide(presentation):
    """Create a slide for rendering tests."""
    blank_layout = presentation.slide_layouts[6]
    return presentation.slides.add_slide(blank_layout)


class TestProgressBarCoverage:
    """Additional tests for ProgressBar to improve coverage."""

    def test_render_with_label(self, slide, dark_theme):
        """Test rendering with label."""
        progress = ProgressBar(value=75, label="Loading", theme=dark_theme)
        shapes = progress.render(slide, left=1, top=2, width=4)
        # Should have background, fill, and label
        assert len(shapes) >= 2

    def test_render_with_percentage(self, slide, dark_theme):
        """Test rendering with percentage display."""
        progress = ProgressBar(value=75, show_percentage=True, theme=dark_theme)
        shapes = progress.render(slide, left=1, top=2, width=4)
        # Should have background, fill, and percentage text
        assert len(shapes) >= 2

    def test_render_with_label_and_percentage(self, slide, dark_theme):
        """Test rendering with both label and percentage."""
        progress = ProgressBar(value=75, label="Progress", show_percentage=True, theme=dark_theme)
        shapes = progress.render(slide, left=1, top=2, width=4)
        # Should have all elements
        assert len(shapes) >= 3

    def test_render_segmented_style(self, slide, dark_theme):
        """Test rendering segmented progress bar."""
        progress = ProgressBar(value=60, style="segmented", segments=10, theme=dark_theme)
        shapes = progress.render(slide, left=1, top=2, width=4)
        # Should have 10 segments
        assert len(shapes) == 10

    def test_render_segmented_with_label(self, slide, dark_theme):
        """Test segmented with label."""
        progress = ProgressBar(
            value=50, style="segmented", segments=5, label="Milestones", theme=dark_theme
        )
        shapes = progress.render(slide, left=1, top=2, width=4)
        # Label + 5 segments
        assert len(shapes) >= 5

    def test_render_zero_value(self, slide, dark_theme):
        """Test rendering with zero value."""
        progress = ProgressBar(value=0, theme=dark_theme)
        shapes = progress.render(slide, left=1, top=2, width=4)
        # Should still render background
        assert len(shapes) >= 1

    def test_render_full_value(self, slide, dark_theme):
        """Test rendering with 100% value."""
        progress = ProgressBar(value=100, theme=dark_theme)
        shapes = progress.render(slide, left=1, top=2, width=4)
        assert len(shapes) >= 2


class TestTileCoverage:
    """Additional tests for Tile components to improve coverage."""

    def test_tile_with_icon_name(self, slide, dark_theme):
        """Test tile with icon name."""
        tile = Tile(label="Task", icon="check", theme=dark_theme)
        shape = tile.render(slide, left=1, top=2)
        assert shape is not None

    def test_icon_tile_all_variants(self, slide, dark_theme):
        """Test IconTile with all variants."""
        for variant in ["default", "filled", "outlined", "ghost"]:
            tile = IconTile("star", label="Star", variant=variant, theme=dark_theme)
            shape = tile.render(slide, left=1, top=2)
            assert shape is not None

    def test_icon_tile_all_color_variants(self, slide, dark_theme):
        """Test IconTile with all color variants."""
        for color in ["default", "primary", "success", "warning", "destructive"]:
            tile = IconTile("check", label="OK", color_variant=color, theme=dark_theme)
            shape = tile.render(slide, left=1, top=2)
            assert shape is not None

    def test_icon_tile_all_sizes(self, slide, dark_theme):
        """Test IconTile with all sizes."""
        for size in ["sm", "md", "lg"]:
            tile = IconTile("star", label="Star", size=size, theme=dark_theme)
            shape = tile.render(slide, left=1, top=2)
            assert shape is not None

    def test_value_tile_variants(self, slide, dark_theme):
        """Test ValueTile with different variants."""
        for variant in ["default", "outlined", "filled"]:
            tile = ValueTile("42", label="Score", variant=variant, theme=dark_theme)
            shape = tile.render(slide, left=1, top=2)
            assert shape is not None

    def test_value_tile_without_label(self, slide, dark_theme):
        """Test ValueTile without label."""
        tile = ValueTile("99", theme=dark_theme)
        shape = tile.render(slide, left=1, top=2)
        assert shape is not None


class TestTimelineCoverage:
    """Additional tests for Timeline to improve coverage."""

    def test_timeline_with_highlights(self, slide, dark_theme):
        """Test timeline with highlighted events."""
        events = [
            {"date": "Q1", "title": "Start"},
            {"date": "Q2", "title": "Milestone", "highlight": True},
            {"date": "Q3", "title": "Launch", "highlight": True},
            {"date": "Q4", "title": "Complete"},
        ]
        timeline = Timeline(events, theme=dark_theme)
        shapes = timeline.render(slide, left=1, top=2, width=6)
        assert len(shapes) > 0

    def test_timeline_different_styles(self, slide, dark_theme):
        """Test timeline with different styles."""
        events = [
            {"date": "2020", "title": "Event 1"},
            {"date": "2021", "title": "Event 2"},
            {"date": "2022", "title": "Event 3"},
        ]

        for style in ["line", "arrow", "dots"]:
            timeline = Timeline(events, style=style, theme=dark_theme)
            shapes = timeline.render(slide, left=1, top=2, width=6)
            assert len(shapes) > 0

    def test_timeline_single_event(self, slide, dark_theme):
        """Test timeline with single event."""
        events = [{"date": "2024", "title": "Only Event"}]
        timeline = Timeline(events, theme=dark_theme)
        shapes = timeline.render(slide, left=1, top=2, width=6)
        assert len(shapes) > 0

    def test_timeline_many_events(self, slide, dark_theme):
        """Test timeline with many events."""
        events = [{"date": f"Q{i}", "title": f"Event {i}"} for i in range(1, 9)]
        timeline = Timeline(events, theme=dark_theme)
        shapes = timeline.render(slide, left=1, top=2, width=8)
        assert len(shapes) > 0


class TestAvatarCoverage:
    """Additional tests for Avatar components to improve coverage."""

    def test_avatar_with_long_text(self, slide, dark_theme):
        """Test avatar with long text (truncation)."""
        avatar = Avatar(text="ABCDEFGH", theme=dark_theme)
        shape = avatar.render(slide, left=1, top=2)
        assert shape is not None

    def test_avatar_all_sizes(self, slide, dark_theme):
        """Test avatar with all sizes."""
        for size in ["xs", "sm", "md", "lg", "xl"]:
            avatar = Avatar(text="AB", size=size, theme=dark_theme)
            shape = avatar.render(slide, left=1, top=2)
            assert shape is not None

    def test_avatar_all_variants(self, slide, dark_theme):
        """Test avatar with all variants."""
        for variant in ["filled", "outlined", "default"]:
            avatar = Avatar(text="XY", variant=variant, theme=dark_theme)
            shape = avatar.render(slide, left=1, top=2)
            assert shape is not None

    def test_avatar_with_label_vertical(self, slide, dark_theme):
        """Test AvatarWithLabel vertical orientation."""
        avatar_label = AvatarWithLabel(
            text="JD",
            label="John Doe",
            sublabel="Developer",
            orientation="vertical",
            theme=dark_theme,
        )
        shapes = avatar_label.render(slide, left=1, top=2, width=2)
        assert len(shapes) >= 1

    def test_avatar_with_label_horizontal(self, slide, dark_theme):
        """Test AvatarWithLabel horizontal orientation."""
        avatar_label = AvatarWithLabel(
            text="JD",
            label="John Doe",
            sublabel="Manager",
            orientation="horizontal",
            theme=dark_theme,
        )
        shapes = avatar_label.render(slide, left=1, top=2, width=3)
        assert len(shapes) >= 1

    def test_avatar_with_label_no_sublabel(self, slide, dark_theme):
        """Test AvatarWithLabel without sublabel."""
        avatar_label = AvatarWithLabel(
            text="AB", label="Only Label", orientation="horizontal", theme=dark_theme
        )
        shapes = avatar_label.render(slide, left=1, top=2, width=3)
        assert len(shapes) >= 1

    def test_avatar_group_no_overlap(self, slide, dark_theme):
        """Test AvatarGroup without overlap."""
        members = [{"text": "A"}, {"text": "B"}, {"text": "C"}]
        group = AvatarGroup(members, overlap=False, theme=dark_theme)
        shapes = group.render(slide, left=1, top=2)
        assert len(shapes) == 3

    def test_avatar_group_with_overlap(self, slide, dark_theme):
        """Test AvatarGroup with overlap."""
        members = [{"text": "A"}, {"text": "B"}, {"text": "C"}]
        group = AvatarGroup(members, overlap=True, theme=dark_theme)
        shapes = group.render(slide, left=1, top=2)
        assert len(shapes) == 3

    def test_avatar_group_max_display_exceeded(self, slide, dark_theme):
        """Test AvatarGroup with max_display limit."""
        members = [{"text": "A"}, {"text": "B"}, {"text": "C"}, {"text": "D"}, {"text": "E"}]
        group = AvatarGroup(members, max_display=3, theme=dark_theme)
        shapes = group.render(slide, left=1, top=2)
        # Should show 3 avatars + "+2" indicator
        assert len(shapes) == 4


class TestButtonCoverage:
    """Additional tests for Button to improve coverage."""

    def test_icon_button_sizes(self, slide, dark_theme):
        """Test IconButton with all sizes."""
        for size in ["sm", "md", "lg"]:
            btn = IconButton(icon="star", size=size, theme=dark_theme)
            shape = btn.render(slide, left=1, top=2)
            assert shape is not None

    def test_button_group_vertical(self, slide, dark_theme):
        """Test ButtonGroup with vertical orientation."""
        buttons = [
            {"text": "Button 1", "variant": "default", "size": "md"},
            {"text": "Button 2", "variant": "secondary", "size": "md"},
        ]
        group = ButtonGroup(buttons, orientation="vertical", theme=dark_theme)
        shapes = group.render(slide, left=1, top=2)
        assert len(shapes) == 2


class TestBadgeCoverage:
    """Additional tests for Badge to improve coverage."""

    def test_count_badge_over_max(self, slide, dark_theme):
        """Test CountBadge with count over 99."""
        badge = CountBadge(count=150, theme=dark_theme)
        shape = badge.render(slide, left=1, top=2)
        assert shape is not None

    def test_dot_badge_all_variants(self, slide, dark_theme):
        """Test DotBadge with all variants."""
        for variant in ["default", "primary", "success", "warning", "destructive"]:
            badge = DotBadge(variant=variant, theme=dark_theme)
            shape = badge.render(slide, left=1, top=2)
            assert shape is not None


class TestCardCoverage:
    """Additional tests for Card to improve coverage."""

    def test_card_with_nested_children(self, slide, dark_theme):
        """Test Card with multiple children."""
        card = Card(variant="elevated", theme=dark_theme)
        card.add_child(Card.Title("Title"))
        card.add_child(Card.Description("Description"))
        card.add_child(Card.Content("Content"))
        card.add_child(Card.Footer("Footer"))
        shape = card.render(slide, left=1, top=2, width=4, height=3)
        assert shape is not None

    def test_metric_card_all_trends(self, slide, dark_theme):
        """Test MetricCard with all trend types."""
        for trend in ["up", "down", "neutral"]:
            metric = MetricCard(
                label="Metric", value="100", change="+5%", trend=trend, theme=dark_theme
            )
            shape = metric.render(slide, left=1, top=2, width=3, height=2)
            assert shape is not None

    def test_card_auto_height(self, slide, dark_theme):
        """Test Card without specifying height."""
        card = Card(variant="default", theme=dark_theme)
        card.add_child(Card.Title("Auto Height"))
        card.add_child(Card.Description("This card calculates its own height"))
        shape = card.render(slide, left=1, top=2, width=4)
        assert shape is not None
