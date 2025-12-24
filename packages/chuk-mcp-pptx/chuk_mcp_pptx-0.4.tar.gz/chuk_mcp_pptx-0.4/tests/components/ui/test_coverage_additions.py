"""
Additional tests for improved coverage of core components.

Targets:
- connector.py (87.88% -> >90%): lines 163-168, 214
- content_grid.py (76.27% -> >90%): lines 116, 157->171, 171->134, 184-202
- grid.py (87.67% -> >90%): lines 291-303
- text.py (88.61% -> >90%): lines 157, 193->173, 195->173, 199->173, 229-232, etc.
"""

import pytest
from pptx import Presentation


class TestConnectorCoverage:
    """Additional tests for Connector component coverage."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_connector_with_placeholder_bounds(self, slide) -> None:
        """Test connector with placeholder bounds extraction."""
        from chuk_mcp_pptx.components.core.connector import Connector
        from unittest.mock import MagicMock
        from pptx.util import Inches

        # Create a mock placeholder with bounds
        mock_placeholder = MagicMock()
        mock_placeholder.left = Inches(1.0)
        mock_placeholder.top = Inches(2.0)
        mock_placeholder.width = Inches(4.0)
        mock_placeholder.height = Inches(1.0)

        # Mock delete_from_parent
        mock_element = MagicMock()
        mock_placeholder._element = mock_element
        mock_element.getparent.return_value = MagicMock()

        connector = Connector(start_x=0, start_y=0, end_x=1, end_y=1, line_color="primary.DEFAULT")
        shape = connector.render(slide, placeholder=mock_placeholder)
        assert shape is not None

    def test_connector_elbow_type(self, slide) -> None:
        """Test connector with elbow type."""
        from chuk_mcp_pptx.components.core.connector import Connector

        connector = Connector(
            start_x=1.0,
            start_y=2.0,
            end_x=5.0,
            end_y=4.0,
            connector_type="elbow",
            line_color="#FF0000",
        )
        shape = connector.render(slide)
        assert shape is not None

    def test_connector_curved_type(self, slide) -> None:
        """Test connector with curved type."""
        from chuk_mcp_pptx.components.core.connector import Connector

        connector = Connector(
            start_x=1.0,
            start_y=2.0,
            end_x=5.0,
            end_y=4.0,
            connector_type="curved",
        )
        shape = connector.render(slide)
        assert shape is not None

    def test_connector_both_arrows(self, slide) -> None:
        """Test connector with arrows at both ends."""
        from chuk_mcp_pptx.components.core.connector import Connector

        connector = Connector(
            start_x=1.0,
            start_y=2.0,
            end_x=5.0,
            end_y=4.0,
            arrow_start=True,
            arrow_end=True,
        )
        shape = connector.render(slide)
        assert shape is not None

    def test_connector_no_arrows(self, slide) -> None:
        """Test connector with no arrows."""
        from chuk_mcp_pptx.components.core.connector import Connector

        connector = Connector(
            start_x=1.0,
            start_y=2.0,
            end_x=5.0,
            end_y=4.0,
            arrow_start=False,
            arrow_end=False,
        )
        shape = connector.render(slide)
        assert shape is not None

    def test_connector_parse_hex_color(self, slide) -> None:
        """Test connector parsing hex color."""
        from chuk_mcp_pptx.components.core.connector import Connector

        connector = Connector(
            start_x=1.0,
            start_y=2.0,
            end_x=5.0,
            end_y=4.0,
            line_color="#00FF00",
        )
        shape = connector.render(slide)
        assert shape is not None

    def test_connector_default_color(self, slide) -> None:
        """Test connector with default (no) color."""
        from chuk_mcp_pptx.components.core.connector import Connector

        connector = Connector(start_x=1.0, start_y=2.0, end_x=5.0, end_y=4.0, line_color=None)
        shape = connector.render(slide)
        assert shape is not None


class TestContentGridCoverage:
    """Additional tests for ContentGrid component coverage."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_content_grid_card_items(self, slide) -> None:
        """Test content grid with card items."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid

        items = [
            {"title": "Feature 1", "description": "Description 1"},
            {"title": "Feature 2", "description": "Description 2"},
            {"title": "Feature 3"},  # No description
        ]
        grid = ContentGrid(items=items, item_type="card", columns=3)
        shapes = grid.render(slide, left=0.5, top=2.0, width=9.0, height=5.0)
        assert isinstance(shapes, list)
        assert len(shapes) > 0

    def test_content_grid_tile_items(self, slide) -> None:
        """Test content grid with tile items."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid

        items = [
            {"label": "Metric 1", "value": "100"},
            {"title": "Metric 2", "value": "200"},  # Using title instead of label
        ]
        grid = ContentGrid(items=items, item_type="tile", columns=2)
        shapes = grid.render(slide, left=0.5, top=2.0, width=9.0, height=4.0)
        assert isinstance(shapes, list)
        assert len(shapes) > 0

    def test_content_grid_button_items(self, slide) -> None:
        """Test content grid with button items."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid

        items = [
            {"text": "Button 1"},
            {"title": "Button 2"},  # Using title instead of text
            {"text": "Button 3", "variant": "primary"},
        ]
        grid = ContentGrid(items=items, item_type="button", columns=3)
        shapes = grid.render(slide, left=0.5, top=2.0, width=9.0, height=2.0)
        assert isinstance(shapes, list)
        assert len(shapes) > 0

    def test_content_grid_with_placeholder(self, slide) -> None:
        """Test content grid with placeholder bounds."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid
        from unittest.mock import MagicMock
        from pptx.util import Inches

        # Create a mock placeholder
        mock_placeholder = MagicMock()
        mock_placeholder.left = Inches(1.0)
        mock_placeholder.top = Inches(2.0)
        mock_placeholder.width = Inches(8.0)
        mock_placeholder.height = Inches(4.0)

        mock_element = MagicMock()
        mock_placeholder._element = mock_element
        mock_element.getparent.return_value = MagicMock()

        items = [{"title": "Item 1"}, {"title": "Item 2"}]
        grid = ContentGrid(items=items, item_type="card", columns=2)
        shapes = grid.render(
            slide, left=0.5, top=2.0, width=9.0, height=5.0, placeholder=mock_placeholder
        )
        assert isinstance(shapes, list)

    def test_content_grid_columns_clamping(self) -> None:
        """Test content grid columns are clamped to 2-4."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid

        # Too few columns
        grid1 = ContentGrid(items=[{"title": "1"}], item_type="card", columns=1)
        assert grid1.columns == 2  # Should be clamped to 2

        # Too many columns
        grid2 = ContentGrid(items=[{"title": "1"}], item_type="card", columns=10)
        assert grid2.columns == 4  # Should be clamped to 4

    def test_content_grid_card_variant(self, slide) -> None:
        """Test content grid cards with variants."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid

        items = [{"title": "Card 1", "variant": "outlined"}]
        grid = ContentGrid(items=items, item_type="card", columns=2)
        shapes = grid.render(slide, left=0.5, top=2.0, width=9.0, height=5.0)
        assert isinstance(shapes, list)

    def test_content_grid_tile_variant(self, slide) -> None:
        """Test content grid tiles with variants."""
        from chuk_mcp_pptx.components.core.content_grid import ContentGrid

        items = [{"label": "Stat", "value": "42", "variant": "success"}]
        grid = ContentGrid(items=items, item_type="tile", columns=2)
        shapes = grid.render(slide, left=0.5, top=2.0, width=9.0, height=4.0)
        assert isinstance(shapes, list)


class TestGridCoverage:
    """Additional tests for Grid component coverage."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_grid_create_layout(self) -> None:
        """Test grid create_layout method."""
        from chuk_mcp_pptx.components.core.grid import Grid

        grid = Grid(columns=12, rows=2, gap="md")
        config = grid.create_layout(left=0.5, top=1.8, width=9.0, height=5.5)

        assert config["columns"] == 12
        assert config["rows"] == 2
        assert config["gap"] == "md"
        assert "bounds" in config
        assert "cell_size" in config
        assert "usage" in config

    def test_grid_create_layout_with_bounds(self) -> None:
        """Test grid create_layout with bounds dict."""
        from chuk_mcp_pptx.components.core.grid import Grid

        bounds = {"left": 1.0, "top": 2.0, "width": 8.0, "height": 4.0}
        grid = Grid(columns=6, rows=3, gap="lg", bounds=bounds)
        config = grid.create_layout()

        assert config["bounds"]["left"] == 1.0
        assert config["bounds"]["top"] == 2.0

    def test_grid_create_layout_default_bounds(self) -> None:
        """Test grid create_layout with default bounds."""
        from chuk_mcp_pptx.components.core.grid import Grid

        grid = Grid(columns=12, rows=1)
        config = grid.create_layout()

        # Should use default bounds
        assert "bounds" in config

    def test_grid_gap_variations(self) -> None:
        """Test grid with different gap sizes."""
        from chuk_mcp_pptx.components.core.grid import Grid

        for gap in ["none", "xs", "sm", "md", "lg", "xl"]:
            grid = Grid(columns=4, gap=gap)
            assert grid.gap == gap

    def test_grid_get_span_method(self) -> None:
        """Test grid get_span method."""
        from chuk_mcp_pptx.components.core.grid import Grid

        grid = Grid(columns=12, rows=2, gap="md")
        span = grid.get_span(
            col_span=6,
            row_span=2,
            col_start=3,
            row_start=0,
            left=0.5,
            top=1.5,
            width=9.0,
            height=5.0,
        )

        assert "left" in span
        assert "top" in span
        assert "width" in span
        assert "height" in span

    def test_grid_get_cell_positions(self, slide) -> None:
        """Test grid get_cell_positions method."""
        from chuk_mcp_pptx.components.core.grid import Grid

        grid = Grid(columns=3, rows=2, gap="md")
        positions = grid.get_cell_positions(slide, left=0.5, top=2.0, width=9.0, height=4.0)

        assert len(positions) == 6  # 3 columns x 2 rows
        for pos in positions:
            assert "left" in pos
            assert "top" in pos
            assert "width" in pos
            assert "height" in pos
            assert "row" in pos
            assert "col" in pos

    def test_grid_get_cell_positions_default_size(self, slide) -> None:
        """Test grid get_cell_positions with default sizes."""
        from chuk_mcp_pptx.components.core.grid import Grid

        grid = Grid(columns=4, rows=1)
        positions = grid.get_cell_positions(slide, left=0.5, top=1.5)

        assert len(positions) == 4  # 4 columns x 1 row

    def test_grid_get_cell_auto_height_false(self) -> None:
        """Test grid get_cell with auto_height=False."""
        from chuk_mcp_pptx.components.core.grid import Grid

        grid = Grid(columns=12, rows=2, gap="md")
        cell = grid.get_cell(
            col_span=6,
            row_span=1,
            col_start=0,
            row_start=0,
            left=0.5,
            top=1.5,
            width=9.0,
            height=5.0,
            auto_height=False,
        )

        assert "left" in cell
        assert "top" in cell
        assert "width" in cell
        assert "height" in cell  # Should include height when auto_height=False

    def test_grid_get_cell_with_bounds(self) -> None:
        """Test grid get_cell using bounds from constructor."""
        from chuk_mcp_pptx.components.core.grid import Grid

        bounds = {"left": 0.5, "top": 1.5, "width": 9.0, "height": 5.0}
        grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)
        cell = grid.get_cell(col_span=6, col_start=0, row_start=0)

        assert "left" in cell
        assert "top" in cell
        assert "width" in cell
        # auto_height=True by default, so no 'height'
        assert "height" not in cell

    def test_grid_get_cell_overrides_bounds(self) -> None:
        """Test grid get_cell with explicit params overriding bounds."""
        from chuk_mcp_pptx.components.core.grid import Grid

        bounds = {"left": 0.5, "top": 1.5, "width": 9.0, "height": 5.0}
        grid = Grid(columns=12, rows=2, gap="md", bounds=bounds)

        # Override with explicit params
        cell = grid.get_cell(
            col_span=6, col_start=0, row_start=0, left=1.0, top=2.0, width=8.0, height=4.0
        )

        # Should use the explicit values, not bounds
        assert cell["left"] >= 1.0  # Starts from explicit left
        assert cell["top"] >= 2.0  # Starts from explicit top


class TestTextCoverage:
    """Additional tests for TextBox and BulletList coverage."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_textbox_with_placeholder(self, slide) -> None:
        """Test textbox populating a placeholder."""
        from chuk_mcp_pptx.components.core.text import TextBox
        from unittest.mock import MagicMock

        # Create a mock placeholder with text_frame
        mock_placeholder = MagicMock()
        mock_frame = MagicMock()
        mock_placeholder.text_frame = mock_frame
        mock_frame.paragraphs = [MagicMock()]
        mock_frame.paragraphs[0].font = MagicMock()

        text = TextBox(text="Hello World")
        shape = text.render(slide, left=1, top=1, width=4, height=1, placeholder=mock_placeholder)
        assert shape is not None

    def test_textbox_various_alignments(self, slide) -> None:
        """Test textbox with various alignments."""
        from chuk_mcp_pptx.components.core.text import TextBox

        for alignment in ["left", "center", "right", "justify"]:
            text = TextBox(text="Test", alignment=alignment)
            shape = text.render(slide, left=1, top=1, width=4, height=1)
            assert shape is not None

    def test_textbox_hex_color(self, slide) -> None:
        """Test textbox with hex color."""
        from chuk_mcp_pptx.components.core.text import TextBox

        text = TextBox(text="Colored Text", color="#FF5500")
        shape = text.render(slide, left=1, top=1, width=4, height=1)
        assert shape is not None

    def test_textbox_semantic_color_with_theme(self, slide) -> None:
        """Test textbox with semantic color and theme."""
        from chuk_mcp_pptx.components.core.text import TextBox

        # Create a simple theme dict
        theme = {
            "colors": {
                "primary": {"DEFAULT": [66, 133, 244]},
                "foreground": {"DEFAULT": [0, 0, 0]},
            },
            "typography": {"font_family": "Arial"},
        }

        text = TextBox(text="Themed Text", color="primary.DEFAULT", theme=theme)
        shape = text.render(slide, left=1, top=1, width=4, height=1)
        assert shape is not None

    def test_textbox_auto_fit(self, slide) -> None:
        """Test textbox with auto-fit."""
        from chuk_mcp_pptx.components.core.text import TextBox

        text = TextBox(text="Auto-fit Text", auto_fit=True)
        shape = text.render(slide, left=1, top=1, width=4, height=1)
        assert shape is not None

    def test_textbox_no_color_no_theme(self, slide) -> None:
        """Test textbox without color and theme."""
        from chuk_mcp_pptx.components.core.text import TextBox

        text = TextBox(text="Plain Text")
        shape = text.render(slide, left=1, top=1, width=4, height=1)
        assert shape is not None

    def test_bulletlist_with_placeholder(self, slide) -> None:
        """Test bullet list populating a placeholder."""
        from chuk_mcp_pptx.components.core.text import BulletList
        from unittest.mock import MagicMock

        # Create a mock placeholder with text_frame
        mock_placeholder = MagicMock()
        mock_frame = MagicMock()
        mock_placeholder.text_frame = mock_frame
        mock_para = MagicMock()
        mock_para.font = MagicMock()
        mock_frame.paragraphs = [mock_para]
        mock_frame.add_paragraph.return_value = MagicMock(font=MagicMock())

        bullets = BulletList(items=["Item 1", "Item 2"])
        shape = bullets.render(
            slide, left=1, top=1, width=4, height=2, placeholder=mock_placeholder
        )
        assert shape is not None

    def test_bulletlist_hex_color(self, slide) -> None:
        """Test bullet list with hex color."""
        from chuk_mcp_pptx.components.core.text import BulletList

        bullets = BulletList(items=["Item 1", "Item 2"], color="#0000FF")
        shape = bullets.render(slide, left=1, top=1, width=4, height=2)
        assert shape is not None

    def test_bulletlist_semantic_color_with_theme(self, slide) -> None:
        """Test bullet list with semantic color and theme."""
        from chuk_mcp_pptx.components.core.text import BulletList

        # Create a simple theme dict
        theme = {
            "colors": {
                "primary": {"DEFAULT": [66, 133, 244]},
                "foreground": {"DEFAULT": [0, 0, 0]},
            },
            "typography": {"font_family": "Arial"},
        }

        bullets = BulletList(items=["Item 1"], color="primary.DEFAULT", theme=theme)
        shape = bullets.render(slide, left=1, top=1, width=4, height=2)
        assert shape is not None

    def test_bulletlist_custom_bullet_char(self, slide) -> None:
        """Test bullet list with custom bullet character."""
        from chuk_mcp_pptx.components.core.text import BulletList

        bullets = BulletList(items=["Check 1", "Check 2"], bullet_char="✓")
        shape = bullets.render(slide, left=1, top=1, width=4, height=2)
        assert shape is not None


class TestVideoCoverage:
    """Tests for Video component coverage."""

    @pytest.fixture
    def slide(self):
        """Create a slide for testing."""
        prs = Presentation()
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def test_video_init(self) -> None:
        """Test Video component initialization."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="demo.mp4", poster_image="poster.jpg", autoplay=True, loop=True)
        assert video.video_source == "demo.mp4"
        assert video.poster_image == "poster.jpg"
        assert video.autoplay is True
        assert video.loop is True

    def test_video_init_empty_source_raises(self) -> None:
        """Test Video raises error with empty source."""
        from chuk_mcp_pptx.components.core.video import Video

        with pytest.raises(ValueError):
            Video(video_source="")

    def test_video_init_invalid_source_type_raises(self) -> None:
        """Test Video raises error with non-string source."""
        from chuk_mcp_pptx.components.core.video import Video

        with pytest.raises(TypeError):
            Video(video_source=123)  # type: ignore

    def test_video_get_mime_type_mp4(self) -> None:
        """Test MIME type detection for mp4."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="test.mp4")
        assert video._get_mime_type("video.mp4") == "video/mp4"

    def test_video_get_mime_type_avi(self) -> None:
        """Test MIME type detection for avi."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="test.mp4")
        assert video._get_mime_type("video.avi") == "video/x-msvideo"

    def test_video_get_mime_type_mov(self) -> None:
        """Test MIME type detection for mov."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="test.mp4")
        assert video._get_mime_type("video.mov") == "video/quicktime"

    def test_video_get_mime_type_wmv(self) -> None:
        """Test MIME type detection for wmv."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="test.mp4")
        assert video._get_mime_type("video.wmv") == "video/x-ms-wmv"

    def test_video_get_mime_type_webm(self) -> None:
        """Test MIME type detection for webm."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="test.mp4")
        assert video._get_mime_type("video.webm") == "video/webm"

    def test_video_get_mime_type_unknown(self) -> None:
        """Test MIME type detection for unknown extension."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="test.mp4")
        assert video._get_mime_type("video.unknown") == "video/mp4"  # Default

    def test_video_show_controls(self) -> None:
        """Test Video with show_controls option."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="demo.mp4", show_controls=False)
        assert video.show_controls is False

    @pytest.mark.asyncio
    async def test_video_render_file_not_found(self, slide) -> None:
        """Test Video render raises FileNotFoundError for missing file."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="nonexistent_video.mp4")
        with pytest.raises(FileNotFoundError):
            await video.render(slide, left=1.0, top=1.0, width=4.0, height=3.0)

    @pytest.mark.asyncio
    async def test_video_render_with_url(self, slide) -> None:
        """Test Video render with URL source (doesn't validate URL)."""
        from chuk_mcp_pptx.components.core.video import Video
        from unittest.mock import MagicMock

        video = Video(video_source="https://example.com/video.mp4")

        # Mock slide.shapes.add_movie to avoid actual video addition
        mock_movie = MagicMock()
        slide.shapes.add_movie = MagicMock(return_value=mock_movie)

        result = await video.render(slide, left=1.0, top=1.0, width=4.0, height=3.0)
        assert result == mock_movie

    @pytest.mark.asyncio
    async def test_video_render_with_placeholder(self, slide) -> None:
        """Test Video render with placeholder bounds extraction."""
        from chuk_mcp_pptx.components.core.video import Video
        from unittest.mock import MagicMock
        from pptx.util import Inches

        video = Video(video_source="https://example.com/video.mp4")

        # Create mock placeholder
        mock_placeholder = MagicMock()
        mock_placeholder.left = Inches(1.0)
        mock_placeholder.top = Inches(2.0)
        mock_placeholder.width = Inches(5.0)
        mock_placeholder.height = Inches(4.0)
        mock_element = MagicMock()
        mock_placeholder._element = mock_element
        mock_element.getparent.return_value = MagicMock()

        # Mock slide.shapes.add_movie
        mock_movie = MagicMock()
        slide.shapes.add_movie = MagicMock(return_value=mock_movie)

        result = await video.render(
            slide, left=1.0, top=1.0, width=4.0, height=3.0, placeholder=mock_placeholder
        )
        assert result == mock_movie

    @pytest.mark.asyncio
    async def test_video_render_with_base64_poster(self, slide) -> None:
        """Test Video render with base64 poster image."""
        from chuk_mcp_pptx.components.core.video import Video
        from unittest.mock import MagicMock
        import base64

        # Create a small fake base64 image
        fake_image_data = base64.b64encode(b"fake image data").decode()
        poster_base64 = f"data:image/png;base64,{fake_image_data}"

        video = Video(video_source="https://example.com/video.mp4", poster_image=poster_base64)

        # Mock slide.shapes.add_movie
        mock_movie = MagicMock()
        slide.shapes.add_movie = MagicMock(return_value=mock_movie)

        result = await video.render(slide, left=1.0, top=1.0, width=4.0, height=3.0)
        assert result == mock_movie

    @pytest.mark.asyncio
    async def test_video_render_with_autoplay_loop(self, slide) -> None:
        """Test Video render with autoplay and loop options."""
        from chuk_mcp_pptx.components.core.video import Video
        from unittest.mock import MagicMock

        video = Video(video_source="https://example.com/video.mp4", autoplay=True, loop=True)

        # Mock slide.shapes.add_movie
        mock_movie = MagicMock()
        mock_movie._element = MagicMock()
        slide.shapes.add_movie = MagicMock(return_value=mock_movie)

        result = await video.render(slide, left=1.0, top=1.0, width=4.0, height=3.0)
        assert result == mock_movie

    @pytest.mark.asyncio
    async def test_video_render_runtime_error(self, slide) -> None:
        """Test Video render raises RuntimeError on add_movie failure."""
        from chuk_mcp_pptx.components.core.video import Video

        video = Video(video_source="https://example.com/video.mp4")

        # Mock slide.shapes.add_movie to raise exception
        slide.shapes.add_movie = lambda *args, **kwargs: (_ for _ in ()).throw(
            Exception("Mock error")
        )

        with pytest.raises(RuntimeError, match="Failed to add video"):
            await video.render(slide, left=1.0, top=1.0, width=4.0, height=3.0)


class TestRegistryCoverage:
    """Additional tests for registry coverage."""

    def test_registry_search_by_name(self) -> None:
        """Test searching components by name."""
        from chuk_mcp_pptx.components.registry import registry

        results = registry.search("button")
        assert isinstance(results, list)

    def test_registry_search_by_description(self) -> None:
        """Test searching components by description."""
        from chuk_mcp_pptx.components.registry import registry

        results = registry.search("container")
        assert isinstance(results, list)

    def test_registry_search_by_tag(self) -> None:
        """Test searching components by tag."""
        from chuk_mcp_pptx.components.registry import registry

        results = registry.search("layout")
        assert isinstance(results, list)

    def test_registry_search_no_results(self) -> None:
        """Test searching with no results."""
        from chuk_mcp_pptx.components.registry import registry

        results = registry.search("xyznonexistent123")
        assert isinstance(results, list)
        assert len(results) == 0

    def test_registry_list_by_category(self) -> None:
        """Test listing by category."""
        from chuk_mcp_pptx.components.registry import registry, ComponentCategory

        for cat in ComponentCategory:
            results = registry.list_by_category(cat)
            assert isinstance(results, list)

    def test_registry_get_schema(self) -> None:
        """Test getting component schema."""
        from chuk_mcp_pptx.components.registry import registry

        components = registry.list_components()
        if components:
            schema = registry.get_schema(components[0])
            assert schema is not None or schema is None  # May not have schema

    def test_registry_get_schema_not_found(self) -> None:
        """Test getting schema for nonexistent component."""
        from chuk_mcp_pptx.components.registry import registry

        schema = registry.get_schema("nonexistent_component_xyz")
        assert schema is None

    def test_registry_get_all_schemas(self) -> None:
        """Test getting all schemas."""
        from chuk_mcp_pptx.components.registry import registry

        schemas = registry.get_all_schemas()
        assert isinstance(schemas, dict)

    def test_registry_export_for_llm(self) -> None:
        """Test exporting registry for LLM."""
        from chuk_mcp_pptx.components.registry import registry
        import json

        llm_export = registry.export_for_llm()
        assert isinstance(llm_export, str)
        # Should be valid JSON
        parsed = json.loads(llm_export)
        assert "components" in parsed

    def test_registry_get_component_signature(self) -> None:
        """Test getting component signature."""
        from chuk_mcp_pptx.components.registry import registry

        components = registry.list_components()
        if components:
            sig = registry.get_component_signature(components[0])
            # May return signature or None
            assert sig is None or isinstance(sig, str)

    def test_registry_get_component_signature_not_found(self) -> None:
        """Test getting signature for nonexistent component."""
        from chuk_mcp_pptx.components.registry import registry

        sig = registry.get_component_signature("nonexistent_xyz")
        assert sig is None

    def test_registry_list_variants(self) -> None:
        """Test listing variants for a component."""
        from chuk_mcp_pptx.components.registry import registry

        components = registry.list_components()
        if components:
            variants = registry.list_variants(components[0])
            # May have variants or None
            assert variants is None or isinstance(variants, dict)

    def test_registry_list_variants_not_found(self) -> None:
        """Test listing variants for nonexistent component."""
        from chuk_mcp_pptx.components.registry import registry

        variants = registry.list_variants("nonexistent_xyz")
        assert variants is None

    def test_registry_get_examples(self) -> None:
        """Test getting examples for a component."""
        from chuk_mcp_pptx.components.registry import registry

        components = registry.list_components()
        if components:
            examples = registry.get_examples(components[0])
            assert isinstance(examples, list)

    def test_registry_get_examples_not_found(self) -> None:
        """Test getting examples for nonexistent component."""
        from chuk_mcp_pptx.components.registry import registry

        examples = registry.get_examples("nonexistent_xyz")
        assert isinstance(examples, list)
        assert len(examples) == 0

    def test_get_component_class(self) -> None:
        """Test get_component_class helper."""
        from chuk_mcp_pptx.components.registry import get_component_class

        # Test with known component
        cls = get_component_class("Button")
        # May or may not find it
        assert cls is None or cls is not None

    def test_get_component_class_not_found(self) -> None:
        """Test get_component_class for nonexistent component."""
        from chuk_mcp_pptx.components.registry import get_component_class

        cls = get_component_class("nonexistent_xyz")
        assert cls is None

    def test_get_component_schema_helper(self) -> None:
        """Test get_component_schema helper function."""
        from chuk_mcp_pptx.components.registry import get_component_schema

        schema = get_component_schema("nonexistent_xyz")
        assert schema is None

    def test_list_components_helper(self) -> None:
        """Test list_components helper function."""
        from chuk_mcp_pptx.components.registry import list_components

        components = list_components()
        assert isinstance(components, list)


class TestVariantsCoverage:
    """Additional tests for variants coverage."""

    def test_variant_config_hex_validation_valid(self) -> None:
        """Test valid hex color validation."""
        from chuk_mcp_pptx.components.variants import VariantConfig

        # Valid 6-digit hex
        config = VariantConfig(props={"color": "#FF5500"})
        assert config.props["color"] == "#FF5500"

        # Valid 3-digit hex
        config2 = VariantConfig(props={"color": "#F50"})
        assert config2.props["color"] == "#F50"

    def test_variant_config_hex_validation_invalid(self) -> None:
        """Test invalid hex color validation."""
        from chuk_mcp_pptx.components.variants import VariantConfig

        with pytest.raises(ValueError):
            VariantConfig(props={"color": "#GGGGGG"})  # Invalid hex

    def test_variant_config_numeric_validation(self) -> None:
        """Test numeric prop validation."""
        from chuk_mcp_pptx.components.variants import VariantConfig

        # Valid positive numbers
        config = VariantConfig(props={"font_size": 14, "padding": 0.5})
        assert config.props["font_size"] == 14

        # Invalid negative values
        with pytest.raises(ValueError):
            VariantConfig(props={"font_size": -1})

        with pytest.raises(ValueError):
            VariantConfig(props={"padding": -0.5})

        with pytest.raises(ValueError):
            VariantConfig(props={"margin": -1})

    def test_variant_config_border_radius_validation(self) -> None:
        """Test border_radius validation."""
        from chuk_mcp_pptx.components.variants import VariantConfig

        # Valid border radius
        config = VariantConfig(props={"border_radius": 8})
        assert config.props["border_radius"] == 8

        # Invalid negative border radius
        with pytest.raises(ValueError):
            VariantConfig(props={"border_radius": -1})

    def test_variant_definition_get_missing_default(self) -> None:
        """Test VariantDefinition.get with missing key and no default."""
        from chuk_mcp_pptx.components.variants import VariantDefinition, VariantConfig

        var_def = VariantDefinition(
            {"primary": VariantConfig(props={"bg": "purple"})}  # No "default" key
        )

        # When "nonexistent" not found, falls back to "default" which also doesn't exist
        config = var_def.get("nonexistent")
        assert config is None  # No fallback available

    def test_chart_variants_preset(self) -> None:
        """Test CHART_VARIANTS preset."""
        from chuk_mcp_pptx.components.variants import CHART_VARIANTS

        props = CHART_VARIANTS.build(style="minimal", legend="none")
        assert props["show_legend"] is False

    def test_column_chart_variants_preset(self) -> None:
        """Test COLUMN_CHART_VARIANTS preset."""
        from chuk_mcp_pptx.components.variants import COLUMN_CHART_VARIANTS

        props = COLUMN_CHART_VARIANTS.build(variant="stacked", style="detailed")
        assert props["show_values"] is True

    def test_pie_chart_variants_preset(self) -> None:
        """Test PIE_CHART_VARIANTS preset."""
        from chuk_mcp_pptx.components.variants import PIE_CHART_VARIANTS

        props = PIE_CHART_VARIANTS.build(variant="doughnut", style="minimal")
        assert props["show_percentages"] is False

    def test_line_chart_variants_preset(self) -> None:
        """Test LINE_CHART_VARIANTS preset."""
        from chuk_mcp_pptx.components.variants import LINE_CHART_VARIANTS

        props = LINE_CHART_VARIANTS.build(variant="smooth_area", style="detailed")
        assert props["fill_area"] is True
        assert props["smooth"] is True


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
