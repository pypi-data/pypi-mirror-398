"""
Comprehensive tests for Button component.
Tests variants, sizes, composition, and rendering.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.button import Button, IconButton, ButtonGroup
from chuk_mcp_pptx.themes import ThemeManager
from chuk_mcp_pptx.components.registry import get_component_schema, list_components


class TestButtonCreation:
    """Test Button component creation and initialization."""

    def test_button_creation_default(self):
        """Test creating a button with default settings."""
        button = Button(text="Click me")
        assert button.text == "Click me"
        assert button.variant == "default"
        assert button.size == "md"

    def test_button_creation_with_variant(self):
        """Test creating buttons with different variants."""
        variants = ["default", "secondary", "outline", "ghost", "destructive"]
        for variant in variants:
            button = Button(text="Test", variant=variant)
            assert button.variant == variant
            assert button.variant_props is not None

    def test_button_creation_with_size(self):
        """Test creating buttons with different sizes."""
        sizes = ["sm", "md", "lg"]
        for size in sizes:
            button = Button(text="Test", size=size)
            assert button.size == size

    def test_button_with_custom_theme(self):
        """Test button with custom theme."""
        from chuk_mcp_pptx.themes import ThemeManager

        theme_manager = ThemeManager()
        custom_theme = theme_manager.get_theme("default-light")
        button = Button(text="Test", theme=custom_theme)
        assert button.theme is not None

    def test_button_variant_props_loaded(self):
        """Test that variant props are properly loaded."""
        button = Button(text="Test", variant="default", size="md")
        assert "bg_color" in button.variant_props
        assert "fg_color" in button.variant_props
        assert "height" in button.variant_props


class TestButtonVariants:
    """Test all button variants and their properties."""

    def test_default_variant(self):
        """Test default variant has correct properties."""
        button = Button(text="Test", variant="default")
        props = button.variant_props
        assert props.get("bg_color") == "primary.DEFAULT"
        assert props.get("fg_color") == "primary.foreground"

    def test_secondary_variant(self):
        """Test secondary variant has correct properties."""
        button = Button(text="Test", variant="secondary")
        props = button.variant_props
        assert props.get("bg_color") == "secondary.DEFAULT"

    def test_outline_variant(self):
        """Test outline variant has border and transparent background."""
        button = Button(text="Test", variant="outline")
        props = button.variant_props
        assert props.get("border_width", 0) > 0
        assert props.get("bg_color") == "transparent"

    def test_ghost_variant(self):
        """Test ghost variant has transparent background."""
        button = Button(text="Test", variant="ghost")
        props = button.variant_props
        assert props.get("bg_color") == "transparent"

    def test_destructive_variant(self):
        """Test destructive variant has destructive colors."""
        button = Button(text="Test", variant="destructive")
        props = button.variant_props
        assert props.get("bg_color") == "destructive.DEFAULT"
        assert props.get("fg_color") == "destructive.foreground"


class TestButtonSizes:
    """Test button size variations."""

    def test_small_size(self):
        """Test small button size properties."""
        button = Button(text="Test", size="sm")
        props = button.variant_props
        # Small size should have smaller height than lg
        assert props.get("height", 0.5) <= 0.8  # Small height

    def test_medium_size(self):
        """Test medium button size properties."""
        button = Button(text="Test", size="md")
        props = button.variant_props
        height = props.get("height", 0.5)
        # Medium size has some defined height
        assert height > 0

    def test_large_size(self):
        """Test large button size properties."""
        button = Button(text="Test", size="lg")
        props = button.variant_props
        # Large size has defined height
        assert props.get("height", 0) > 0  # Large height

    def test_size_affects_font_size(self):
        """Test that size affects font size."""
        small = Button(text="Test", size="sm")
        medium = Button(text="Test", size="md")
        large = Button(text="Test", size="lg")

        assert small.variant_props.get("font_size") < medium.variant_props.get("font_size")
        assert medium.variant_props.get("font_size") < large.variant_props.get("font_size")


class TestButtonRendering:
    """Test button rendering to PowerPoint slides."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_button_renders_to_slide(self, slide):
        """Test that button renders successfully to slide."""
        button = Button(text="Click me", variant="default")
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape is not None
        assert shape in slide.shapes

    def test_button_position(self, slide):
        """Test button is positioned correctly."""
        button = Button(text="Test")
        shape = button.render(slide, left=2.0, top=3.0)
        assert shape.left == Inches(2.0)
        assert shape.top == Inches(3.0)

    def test_button_custom_dimensions(self, slide):
        """Test button with custom width and height."""
        button = Button(text="Test")
        shape = button.render(slide, left=1.0, top=1.0, width=3.0, height=0.8)
        assert shape.width == Inches(3.0)
        assert shape.height == Inches(0.8)

    def test_button_has_text(self, slide):
        """Test button contains the correct text."""
        button = Button(text="Click me")
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape.text_frame.text == "Click me"

    def test_button_text_is_centered(self, slide):
        """Test button text is centered."""
        button = Button(text="Test")
        shape = button.render(slide, left=1.0, top=1.0)
        paragraph = shape.text_frame.paragraphs[0]
        from pptx.enum.text import PP_ALIGN

        assert paragraph.alignment == PP_ALIGN.CENTER

    def test_button_is_rounded_rectangle(self, slide):
        """Test button shape is autoshape (rounded rectangle variant)."""
        button = Button(text="Test")
        shape = button.render(slide, left=1.0, top=1.0)
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        assert shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE

    def test_button_with_different_themes(self, slide):
        """Test button renders with different themes."""
        theme_manager = ThemeManager()
        themes = ["default-light", "default-dark", "ocean-light", "forest-light"]

        for theme_name in themes:
            theme = theme_manager.get_theme(theme_name)
            button = Button(text="Test", theme=theme)
            shape = button.render(slide, left=1.0, top=1.0)
            assert shape is not None

    def test_all_variant_combinations_render(self, slide):
        """Test all variant and size combinations render."""
        variants = ["default", "secondary", "outline", "ghost", "destructive"]
        sizes = ["sm", "md", "lg"]

        top = 0.5
        for variant in variants:
            for size in sizes:
                button = Button(text="Test", variant=variant, size=size)
                shape = button.render(slide, left=1.0, top=top)
                assert shape is not None
                top += 0.6


class TestIconButton:
    """Test IconButton component."""

    def test_icon_button_creation(self):
        """Test creating an icon button."""
        button = IconButton(icon="play")
        assert button.text == "▶"  # Unicode play icon

    def test_icon_button_with_custom_icon(self):
        """Test icon button with custom Unicode."""
        button = IconButton(icon="♠")
        assert button.text == "♠"

    def test_icon_button_mappings(self):
        """Test all icon mappings work."""
        icons = [
            "play",
            "pause",
            "stop",
            "next",
            "previous",
            "plus",
            "minus",
            "close",
            "check",
            "star",
            "heart",
            "settings",
            "menu",
            "search",
        ]

        for icon_name in icons:
            button = IconButton(icon=icon_name)
            assert button.text in IconButton.ICONS.values()

    def test_icon_button_is_square(self, presentation):
        """Test icon button renders as square."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        button = IconButton(icon="play")
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape.width == shape.height

    def test_icon_button_sizes(self, presentation):
        """Test icon button different sizes are square."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        sizes = ["sm", "md", "lg"]

        for size in sizes:
            button = IconButton(icon="play", size=size)
            shape = button.render(slide, left=1.0, top=1.0)
            assert shape.width == shape.height

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()


class TestButtonGroup:
    """Test ButtonGroup component."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_button_group_creation(self):
        """Test creating a button group."""
        buttons = [{"text": "Save", "variant": "default"}, {"text": "Cancel", "variant": "ghost"}]
        group = ButtonGroup(buttons=buttons)
        assert len(group.buttons) == 2
        assert group.orientation == "horizontal"

    def test_button_group_horizontal_layout(self, slide):
        """Test horizontal button group layout."""
        buttons = [
            {"text": "One", "variant": "default"},
            {"text": "Two", "variant": "secondary"},
            {"text": "Three", "variant": "outline"},
        ]
        group = ButtonGroup(buttons=buttons, orientation="horizontal")
        shapes = group.render(slide, left=1.0, top=1.0)

        assert len(shapes) == 3
        # Check they're laid out horizontally (increasing left position)
        assert shapes[1].left > shapes[0].left
        assert shapes[2].left > shapes[1].left
        # Check same top position
        assert shapes[0].top == shapes[1].top == shapes[2].top

    def test_button_group_vertical_layout(self, slide):
        """Test vertical button group layout."""
        buttons = [
            {"text": "First", "variant": "default"},
            {"text": "Second", "variant": "secondary"},
        ]
        group = ButtonGroup(buttons=buttons, orientation="vertical")
        shapes = group.render(slide, left=1.0, top=1.0)

        assert len(shapes) == 2
        # Check they're laid out vertically (increasing top position)
        assert shapes[1].top > shapes[0].top
        # Check same left position
        assert shapes[0].left == shapes[1].left

    def test_button_group_spacing(self, slide):
        """Test button group spacing between buttons."""
        buttons = [{"text": "A", "variant": "default"}, {"text": "B", "variant": "default"}]
        spacing = 0.2
        group = ButtonGroup(buttons=buttons, spacing=spacing)
        shapes = group.render(slide, left=1.0, top=1.0)

        # Calculate actual spacing
        gap = shapes[1].left.inches - (shapes[0].left.inches + shapes[0].width.inches)
        assert abs(gap - spacing) < 0.01  # Allow small float precision error

    def test_button_group_empty(self, slide):
        """Test button group with no buttons."""
        group = ButtonGroup(buttons=[])
        shapes = group.render(slide, left=1.0, top=1.0)
        assert len(shapes) == 0


class TestButtonThemeIntegration:
    """Test button integration with theme system."""

    def test_button_uses_theme_colors(self):
        """Test button uses theme colors correctly."""
        from chuk_mcp_pptx.themes import ThemeManager

        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("default-light")
        button = Button(text="Test", variant="default", theme=theme)
        # Just verify it can get a color without errors
        color = button.get_color("primary.DEFAULT")
        assert color is not None

    def test_button_with_all_built_in_themes(self, presentation):
        """Test button renders with all built-in themes."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        theme_manager = ThemeManager()
        themes = theme_manager.list_themes()

        for theme_name in themes[:5]:  # Test first 5 themes
            theme = theme_manager.get_theme(theme_name)
            button = Button(text="Test", theme=theme)
            shape = button.render(slide, left=1.0, top=1.0)
            assert shape is not None

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()


class TestButtonRegistry:
    """Test button component registry integration."""

    def test_button_is_registered(self):
        """Test Button is in component registry."""
        components = list_components()
        assert "Button" in components

    def test_button_schema_exists(self):
        """Test Button has a valid schema."""
        schema = get_component_schema("Button")
        assert schema is not None
        assert schema["name"] == "Button"
        assert "schema" in schema
        assert "variants" in schema

    def test_button_schema_has_props(self):
        """Test Button schema has properties in JSON schema."""
        schema = get_component_schema("Button")
        # Schema uses JSON Schema format, props are in schema.schema.properties
        assert "schema" in schema
        assert "properties" in schema["schema"]
        prop_names = list(schema["schema"]["properties"].keys())
        assert "text" in prop_names
        assert "variant" in prop_names
        assert "size" in prop_names
        assert "left" in prop_names
        assert "top" in prop_names

    def test_button_schema_has_variants(self):
        """Test Button schema has variant definitions."""
        schema = get_component_schema("Button")
        assert "variant" in schema["variants"]
        assert "size" in schema["variants"]
        assert "default" in schema["variants"]["variant"]

    def test_button_schema_has_examples(self):
        """Test Button schema has usage examples."""
        schema = get_component_schema("Button")
        assert "examples" in schema
        assert len(schema["examples"]) > 0

    def test_icon_button_is_registered(self):
        """Test IconButton is in component registry."""
        components = list_components()
        assert "IconButton" in components

    def test_button_group_is_registered(self):
        """Test ButtonGroup is in component registry."""
        components = list_components()
        assert "ButtonGroup" in components


class TestButtonEdgeCases:
    """Test button edge cases and error handling."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_button_empty_text(self, slide):
        """Test button with empty text."""
        button = Button(text="")
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_button_very_long_text(self, slide):
        """Test button with very long text."""
        long_text = "This is a very long button text that might cause layout issues"
        button = Button(text=long_text)
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape is not None
        assert long_text in shape.text_frame.text

    def test_button_special_characters(self, slide):
        """Test button with special characters."""
        special = "Test™ © ® € £ ¥"
        button = Button(text=special)
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_button_unicode_emoji(self, slide):
        """Test button with emoji."""
        emoji = "Click me 👍 🎉"
        button = Button(text=emoji)
        shape = button.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_button_at_slide_boundaries(self, slide):
        """Test button positioned at slide boundaries."""
        button = Button(text="Test")
        # Near right edge
        shape = button.render(slide, left=9.0, top=1.0)
        assert shape is not None
        # Near bottom edge
        shape = button.render(slide, left=1.0, top=6.5)
        assert shape is not None

    def test_button_very_small_dimensions(self, slide):
        """Test button with very small custom dimensions."""
        button = Button(text="T")
        shape = button.render(slide, left=1.0, top=1.0, width=0.3, height=0.2)
        assert shape.width == Inches(0.3)
        assert shape.height == Inches(0.2)

    def test_button_very_large_dimensions(self, slide):
        """Test button with very large custom dimensions."""
        button = Button(text="Large Button")
        shape = button.render(slide, left=1.0, top=1.0, width=8.0, height=2.0)
        assert shape.width == Inches(8.0)
        assert shape.height == Inches(2.0)
