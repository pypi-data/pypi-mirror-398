"""
Comprehensive tests for Badge component.
Tests variants, rendering, and all badge types.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.badge import Badge, DotBadge, CountBadge
from chuk_mcp_pptx.themes import ThemeManager
from chuk_mcp_pptx.components.registry import get_component_schema, list_components


class TestBadgeCreation:
    """Test Badge component creation and initialization."""

    def test_badge_creation_default(self):
        """Test creating a badge with default settings."""
        badge = Badge(text="New")
        assert badge.text == "New"
        assert badge.variant == "default"

    def test_badge_creation_with_variant(self):
        """Test creating badges with different variants."""
        variants = ["default", "secondary", "success", "warning", "destructive", "outline"]
        for variant in variants:
            badge = Badge(text="Test", variant=variant)
            assert badge.variant == variant
            assert badge.variant_props is not None

    def test_badge_with_custom_theme(self):
        """Test badge with custom theme."""
        theme_manager = ThemeManager()
        custom_theme = theme_manager.get_theme("ocean-dark")
        badge = Badge(text="Custom", theme=custom_theme)
        assert badge.theme is not None

    def test_badge_variant_props_loaded(self):
        """Test that variant props are properly loaded."""
        badge = Badge(text="Test", variant="success")
        assert "bg_color" in badge.variant_props
        assert "fg_color" in badge.variant_props


class TestBadgeVariants:
    """Test all badge variants and their properties."""

    def test_default_variant(self):
        """Test default variant has correct properties."""
        badge = Badge(text="Default", variant="default")
        props = badge.variant_props
        assert props.get("bg_color") == "primary.DEFAULT"

    def test_secondary_variant(self):
        """Test secondary variant has correct properties."""
        badge = Badge(text="Secondary", variant="secondary")
        props = badge.variant_props
        assert props.get("bg_color") == "secondary.DEFAULT"

    def test_success_variant(self):
        """Test success variant has correct properties."""
        badge = Badge(text="Success", variant="success")
        props = badge.variant_props
        assert props.get("bg_color") == "success.DEFAULT"

    def test_warning_variant(self):
        """Test warning variant has correct properties."""
        badge = Badge(text="Warning", variant="warning")
        props = badge.variant_props
        assert props.get("bg_color") == "warning.DEFAULT"

    def test_destructive_variant(self):
        """Test destructive variant has correct properties."""
        badge = Badge(text="Error", variant="destructive")
        props = badge.variant_props
        assert props.get("bg_color") == "destructive.DEFAULT"

    def test_outline_variant(self):
        """Test outline variant has border."""
        badge = Badge(text="Outline", variant="outline")
        props = badge.variant_props
        assert props.get("border_width", 0) > 0


class TestBadgeRendering:
    """Test badge rendering to PowerPoint slides."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_badge_renders_to_slide(self, slide):
        """Test that badge renders successfully to slide."""
        badge = Badge(text="Active", variant="success")
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None
        assert shape in slide.shapes

    def test_badge_position(self, slide):
        """Test badge is positioned correctly."""
        badge = Badge(text="Test")
        shape = badge.render(slide, left=2.0, top=3.0)
        assert shape.left == Inches(2.0)
        assert shape.top == Inches(3.0)

    def test_badge_custom_dimensions(self, slide):
        """Test badge with custom width and height."""
        badge = Badge(text="Custom")
        shape = badge.render(slide, left=1.0, top=1.0, width=2.0, height=0.4)
        assert shape.width == Inches(2.0)
        assert shape.height == Inches(0.4)

    def test_badge_auto_width(self, slide):
        """Test badge auto-calculates width based on text."""
        short_badge = Badge(text="Hi")
        long_badge = Badge(text="This is a long badge")

        short_shape = short_badge.render(slide, left=1.0, top=1.0)
        long_shape = long_badge.render(slide, left=1.0, top=2.0)

        assert long_shape.width > short_shape.width

    def test_badge_has_text(self, slide):
        """Test badge contains the correct text."""
        badge = Badge(text="Active")
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape.text_frame.text == "Active"

    def test_badge_text_is_centered(self, slide):
        """Test badge text is centered."""
        badge = Badge(text="Test")
        shape = badge.render(slide, left=1.0, top=1.0)
        paragraph = shape.text_frame.paragraphs[0]
        from pptx.enum.text import PP_ALIGN

        assert paragraph.alignment == PP_ALIGN.CENTER

    def test_badge_is_rounded_rectangle(self, slide):
        """Test badge shape is autoshape (rounded rectangle)."""
        badge = Badge(text="Test")
        shape = badge.render(slide, left=1.0, top=1.0)
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        assert shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE

    def test_all_variant_combinations_render(self, slide):
        """Test all variants render successfully."""
        variants = ["default", "secondary", "success", "warning", "destructive", "outline"]

        top = 0.5
        for variant in variants:
            badge = Badge(text=variant.title(), variant=variant)
            shape = badge.render(slide, left=1.0, top=top)
            assert shape is not None
            top += 0.4


class TestDotBadge:
    """Test DotBadge component."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_dot_badge_creation(self):
        """Test creating a dot badge."""
        dot = DotBadge(variant="success")
        assert dot.variant == "success"
        assert dot.size == 0.15

    def test_dot_badge_custom_size(self):
        """Test dot badge with custom size."""
        dot = DotBadge(variant="default", size=0.3)
        assert dot.size == 0.3

    def test_dot_badge_renders(self, slide):
        """Test dot badge renders to slide."""
        dot = DotBadge(variant="success")
        shape = dot.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_dot_badge_is_circular(self, slide):
        """Test dot badge is circular (width == height)."""
        dot = DotBadge(variant="warning")
        shape = dot.render(slide, left=1.0, top=1.0)
        assert shape.width == shape.height

    def test_dot_badge_variants(self, slide):
        """Test all dot badge variants."""
        variants = ["default", "success", "warning", "destructive"]

        for variant in variants:
            dot = DotBadge(variant=variant)
            shape = dot.render(slide, left=1.0, top=1.0)
            assert shape is not None


class TestCountBadge:
    """Test CountBadge component."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_count_badge_creation(self):
        """Test creating a count badge."""
        badge = CountBadge(count=5)
        assert badge.count == 5
        assert badge.text == "5"

    def test_count_badge_with_max(self):
        """Test count badge respects max count."""
        badge = CountBadge(count=150, max_count=99)
        assert badge.count == 150
        assert badge.text == "99+"

    def test_count_badge_under_max(self):
        """Test count badge under max shows actual count."""
        badge = CountBadge(count=42, max_count=99)
        assert badge.text == "42"

    def test_count_badge_renders(self, slide):
        """Test count badge renders to slide."""
        badge = CountBadge(count=5)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None
        assert "5" in shape.text_frame.text

    def test_count_badge_single_digit(self, slide):
        """Test single digit count badge is compact."""
        badge = CountBadge(count=7)
        shape = badge.render(slide, left=1.0, top=1.0)
        # Single digit should be roughly circular
        assert abs(shape.width.inches - 0.3) < 0.1

    def test_count_badge_double_digit(self, slide):
        """Test double digit count badge is wider."""
        badge = CountBadge(count=42)
        shape = badge.render(slide, left=1.0, top=1.0)
        # Double digit should be wider
        assert shape.width.inches >= 0.4

    def test_count_badge_max_display(self, slide):
        """Test count badge showing '99+'."""
        badge = CountBadge(count=200, max_count=99)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert "99+" in shape.text_frame.text


class TestBadgeThemeIntegration:
    """Test badge integration with theme system."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_badge_uses_theme_colors(self):
        """Test badge uses theme colors correctly."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("default-light")
        badge = Badge(text="Themed", variant="success", theme=theme)
        color = badge.get_color("success.DEFAULT")
        assert color is not None

    def test_badge_with_all_built_in_themes(self, slide):
        """Test badge renders with all built-in themes."""
        theme_manager = ThemeManager()
        themes = theme_manager.list_themes()

        for theme_name in themes[:5]:  # Test first 5 themes
            theme = theme_manager.get_theme(theme_name)
            badge = Badge(text="Test", variant="success", theme=theme)
            shape = badge.render(slide, left=1.0, top=1.0)
            assert shape is not None


class TestBadgeRegistry:
    """Test badge component registry integration."""

    def test_badge_is_registered(self):
        """Test Badge is in component registry."""
        components = list_components()
        assert "Badge" in components

    def test_badge_schema_exists(self):
        """Test Badge has a valid schema."""
        schema = get_component_schema("Badge")
        assert schema is not None
        assert schema["name"] == "Badge"

    def test_badge_schema_has_props(self):
        """Test Badge schema has properties."""
        schema = get_component_schema("Badge")
        assert "schema" in schema
        prop_names = list(schema["schema"]["properties"].keys())
        assert "text" in prop_names
        assert "variant" in prop_names

    def test_badge_schema_has_variants(self):
        """Test Badge schema has variant definitions."""
        schema = get_component_schema("Badge")
        assert "variants" in schema
        assert "variant" in schema["variants"]

    def test_badge_schema_has_examples(self):
        """Test Badge schema has usage examples."""
        schema = get_component_schema("Badge")
        assert "examples" in schema
        assert len(schema["examples"]) > 0

    def test_dot_badge_is_registered(self):
        """Test DotBadge is in component registry."""
        components = list_components()
        assert "DotBadge" in components

    def test_count_badge_is_registered(self):
        """Test CountBadge is in component registry."""
        components = list_components()
        assert "CountBadge" in components


class TestBadgeEdgeCases:
    """Test badge edge cases and error handling."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_badge_empty_text(self, slide):
        """Test badge with empty text."""
        badge = Badge(text="")
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_badge_very_long_text(self, slide):
        """Test badge with very long text."""
        long_text = "This is an extremely long badge text that might cause issues"
        badge = Badge(text=long_text)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None
        assert long_text in shape.text_frame.text

    def test_badge_special_characters(self, slide):
        """Test badge with special characters."""
        special = "v1.0.0-beta+build.123"
        badge = Badge(text=special)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_badge_unicode_emoji(self, slide):
        """Test badge with emoji."""
        emoji = "🔥 Hot"
        badge = Badge(text=emoji)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_badge_numbers_only(self, slide):
        """Test badge with only numbers."""
        badge = Badge(text="12345")
        shape = badge.render(slide, left=1.0, top=1.0)
        assert "12345" in shape.text_frame.text

    def test_badge_single_character(self, slide):
        """Test badge with single character."""
        badge = Badge(text="A")
        shape = badge.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_badge_very_small_dimensions(self, slide):
        """Test badge with very small custom dimensions."""
        badge = Badge(text="X")
        shape = badge.render(slide, left=1.0, top=1.0, width=0.2, height=0.15)
        assert shape.width == Inches(0.2)

    def test_count_badge_zero(self, slide):
        """Test count badge with zero count."""
        badge = CountBadge(count=0)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert "0" in shape.text_frame.text

    def test_count_badge_large_number(self, slide):
        """Test count badge with very large number."""
        badge = CountBadge(count=999999, max_count=999)
        shape = badge.render(slide, left=1.0, top=1.0)
        assert "999+" in shape.text_frame.text

    def test_dot_badge_very_small(self, slide):
        """Test dot badge with very small size."""
        dot = DotBadge(variant="success", size=0.05)
        shape = dot.render(slide, left=1.0, top=1.0)
        assert shape.width == Inches(0.05)

    def test_dot_badge_very_large(self, slide):
        """Test dot badge with large size."""
        dot = DotBadge(variant="warning", size=0.5)
        shape = dot.render(slide, left=1.0, top=1.0)
        assert shape.width == Inches(0.5)
