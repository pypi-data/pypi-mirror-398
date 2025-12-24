"""
Tests for shape components (Shape, Connector, Image, SmartArt).
"""

import pytest
import tempfile
import os
from pptx import Presentation
from PIL import Image as PILImage

from chuk_mcp_pptx.components.core import (
    Shape,
    Connector,
    Image,
    ProcessFlow,
    CycleDiagram,
    HierarchyDiagram,
)


@pytest.fixture
def presentation():
    """Create a presentation for rendering tests."""
    return Presentation()


@pytest.fixture
def slide(presentation):
    """Create a slide for rendering tests."""
    blank_layout = presentation.slide_layouts[6]
    return presentation.slides.add_slide(blank_layout)


class TestShape:
    """Tests for Shape component."""

    def test_init(self):
        """Test Shape initialization."""
        shape = Shape(shape_type="rectangle", text="Test")
        assert shape.shape_type == "rectangle"
        assert shape.text == "Test"

    def test_shape_types(self, slide, dark_theme):
        """Test various shape types."""
        shape_types = [
            "rectangle",
            "rounded_rectangle",
            "oval",
            "circle",
            "diamond",
            "triangle",
            "star",
            "hexagon",
            "heart",
        ]

        for shape_type in shape_types:
            shape = Shape(shape_type=shape_type, theme=dark_theme)
            rendered = shape.render(slide, left=1, top=1, width=2, height=1.5)
            assert rendered is not None

    def test_shape_with_text(self, slide, dark_theme):
        """Test shape with text content."""
        shape = Shape(shape_type="rounded_rectangle", text="Hello World", theme=dark_theme)
        rendered = shape.render(slide, left=1, top=1, width=3, height=2)
        assert rendered is not None
        assert rendered.has_text_frame

    def test_shape_with_hex_color(self, slide, dark_theme):
        """Test shape with hex color."""
        shape = Shape(
            shape_type="rectangle", fill_color="#FF5733", line_color="#000000", theme=dark_theme
        )
        rendered = shape.render(slide, left=1, top=1, width=2, height=1.5)
        assert rendered is not None

    def test_shape_with_semantic_color(self, slide, dark_theme):
        """Test shape with semantic color."""
        shape = Shape(
            shape_type="oval",
            fill_color="primary.DEFAULT",
            line_color="border.DEFAULT",
            theme=dark_theme,
        )
        rendered = shape.render(slide, left=1, top=1, width=2, height=2)
        assert rendered is not None

    def test_shape_custom_line_width(self, slide, dark_theme):
        """Test shape with custom line width."""
        shape = Shape(shape_type="rectangle", line_width=5.0, theme=dark_theme)
        rendered = shape.render(slide, left=1, top=1, width=2, height=1.5)
        assert rendered is not None


class TestConnector:
    """Tests for Connector component."""

    def test_init(self):
        """Test Connector initialization."""
        connector = Connector(start_x=1.0, start_y=2.0, end_x=5.0, end_y=3.0)
        assert connector.start_x == 1.0
        assert connector.start_y == 2.0
        assert connector.end_x == 5.0
        assert connector.end_y == 3.0

    def test_connector_types(self, slide, dark_theme):
        """Test various connector types."""
        connector_types = ["straight", "elbow", "curved"]

        for connector_type in connector_types:
            connector = Connector(
                start_x=1.0,
                start_y=2.0,
                end_x=5.0,
                end_y=3.0,
                connector_type=connector_type,
                theme=dark_theme,
            )
            rendered = connector.render(slide)
            assert rendered is not None

    def test_connector_with_arrows(self, slide, dark_theme):
        """Test connector with arrow configurations."""
        # Arrow at end
        connector1 = Connector(
            start_x=1.0,
            start_y=2.0,
            end_x=5.0,
            end_y=2.0,
            arrow_end=True,
            arrow_start=False,
            theme=dark_theme,
        )
        rendered1 = connector1.render(slide)
        assert rendered1 is not None

        # Arrow at start
        connector2 = Connector(
            start_x=1.0,
            start_y=3.0,
            end_x=5.0,
            end_y=3.0,
            arrow_end=False,
            arrow_start=True,
            theme=dark_theme,
        )
        rendered2 = connector2.render(slide)
        assert rendered2 is not None

        # Arrows at both ends
        connector3 = Connector(
            start_x=1.0,
            start_y=4.0,
            end_x=5.0,
            end_y=4.0,
            arrow_end=True,
            arrow_start=True,
            theme=dark_theme,
        )
        rendered3 = connector3.render(slide)
        assert rendered3 is not None

    def test_connector_with_colors(self, slide, dark_theme):
        """Test connector with custom colors."""
        # Hex color
        connector1 = Connector(
            start_x=1.0, start_y=2.0, end_x=5.0, end_y=3.0, line_color="#FF0000", theme=dark_theme
        )
        rendered1 = connector1.render(slide)
        assert rendered1 is not None

        # Semantic color
        connector2 = Connector(
            start_x=1.0,
            start_y=3.0,
            end_x=5.0,
            end_y=4.0,
            line_color="primary.DEFAULT",
            theme=dark_theme,
        )
        rendered2 = connector2.render(slide)
        assert rendered2 is not None

    def test_connector_line_width(self, slide, dark_theme):
        """Test connector with custom line width."""
        connector = Connector(
            start_x=1.0, start_y=2.0, end_x=5.0, end_y=3.0, line_width=5.0, theme=dark_theme
        )
        rendered = connector.render(slide)
        assert rendered is not None


class TestImage:
    """Tests for Image component."""

    @pytest.fixture
    def temp_image(self):
        """Create a temporary test image."""
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as f:
            # Create a simple test image
            img = PILImage.new("RGB", (100, 100), color="red")
            img.save(f.name)
            yield f.name
            # Cleanup
            os.unlink(f.name)

    def test_init(self, temp_image):
        """Test Image initialization."""
        image = Image(image_source=temp_image)
        assert image.image_source == temp_image
        assert not image.shadow

    @pytest.mark.asyncio
    async def test_render_from_file(self, slide, temp_image, dark_theme):
        """Test rendering image from file path."""
        image = Image(image_source=temp_image, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_render_with_width_and_height(self, slide, temp_image, dark_theme):
        """Test rendering with both width and height."""
        image = Image(image_source=temp_image, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=4, height=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_render_with_width_only(self, slide, temp_image, dark_theme):
        """Test rendering with width only (maintains ratio)."""
        image = Image(image_source=temp_image, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=4)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_render_with_height_only(self, slide, temp_image, dark_theme):
        """Test rendering with height only (maintains ratio)."""
        image = Image(image_source=temp_image, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, height=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_render_with_shadow(self, slide, temp_image, dark_theme):
        """Test rendering with shadow effect."""
        image = Image(image_source=temp_image, shadow=True, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_file_not_found(self, slide, dark_theme):
        """Test error handling for non-existent file."""
        image = Image(image_source="/nonexistent/path.png", theme=dark_theme)
        with pytest.raises(FileNotFoundError):
            await image.render(slide, left=1, top=1)

    @pytest.mark.asyncio
    async def test_blur_filter(self, slide, temp_image, dark_theme):
        """Test blur filter."""
        image = Image(image_source=temp_image, blur_radius=10, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_grayscale_filter(self, slide, temp_image, dark_theme):
        """Test grayscale filter."""
        image = Image(image_source=temp_image, grayscale=True, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_sepia_filter(self, slide, temp_image, dark_theme):
        """Test sepia filter."""
        image = Image(image_source=temp_image, sepia=True, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_brightness_filter(self, slide, temp_image, dark_theme):
        """Test brightness adjustment."""
        image = Image(image_source=temp_image, brightness=1.5, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_contrast_filter(self, slide, temp_image, dark_theme):
        """Test contrast adjustment."""
        image = Image(image_source=temp_image, contrast=1.8, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_saturation_filter(self, slide, temp_image, dark_theme):
        """Test saturation adjustment."""
        image = Image(image_source=temp_image, saturation=2.0, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_sharpen_filter(self, slide, temp_image, dark_theme):
        """Test sharpen filter."""
        image = Image(image_source=temp_image, sharpen=True, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_invert_filter(self, slide, temp_image, dark_theme):
        """Test invert filter."""
        image = Image(image_source=temp_image, invert=True, theme=dark_theme)
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None

    @pytest.mark.asyncio
    async def test_combined_filters(self, slide, temp_image, dark_theme):
        """Test multiple filters combined."""
        image = Image(
            image_source=temp_image,
            blur_radius=5,
            brightness=1.2,
            contrast=1.3,
            saturation=1.5,
            shadow=True,
            theme=dark_theme,
        )
        rendered = await image.render(slide, left=1, top=1, width=3)
        assert rendered is not None


class TestProcessFlow:
    """Tests for ProcessFlow diagram."""

    def test_init(self):
        """Test ProcessFlow initialization."""
        process = ProcessFlow(items=["Step 1", "Step 2", "Step 3"])
        assert len(process.items) == 3

    def test_render(self, slide, dark_theme):
        """Test ProcessFlow rendering."""
        process = ProcessFlow(items=["Research", "Design", "Develop", "Test"], theme=dark_theme)
        shapes = process.render(slide, left=1, top=2, width=8, height=2)
        assert len(shapes) == 4

    def test_empty_items(self, slide, dark_theme):
        """Test ProcessFlow with empty items."""
        process = ProcessFlow(items=[], theme=dark_theme)
        shapes = process.render(slide, left=1, top=2, width=8, height=2)
        assert len(shapes) == 0

    def test_single_item(self, slide, dark_theme):
        """Test ProcessFlow with single item."""
        process = ProcessFlow(items=["Only One"], theme=dark_theme)
        shapes = process.render(slide, left=1, top=2, width=8, height=2)
        assert len(shapes) == 1

    def test_many_items(self, slide, dark_theme):
        """Test ProcessFlow with many items."""
        items = [f"Step {i}" for i in range(1, 8)]
        process = ProcessFlow(items=items, theme=dark_theme)
        shapes = process.render(slide, left=0.5, top=2, width=9, height=2)
        assert len(shapes) == 7


class TestCycleDiagram:
    """Tests for CycleDiagram."""

    def test_init(self):
        """Test CycleDiagram initialization."""
        cycle = CycleDiagram(items=["Plan", "Do", "Check", "Act"])
        assert len(cycle.items) == 4

    def test_render(self, slide, dark_theme):
        """Test CycleDiagram rendering."""
        cycle = CycleDiagram(items=["Plan", "Do", "Check", "Act"], theme=dark_theme)
        shapes = cycle.render(slide, left=1, top=1, width=6, height=5)
        # 4 shapes + 4 connectors = 8
        assert len(shapes) == 8

    def test_different_sizes(self, slide, dark_theme):
        """Test CycleDiagram with different numbers of items."""
        # Small cycle (4 items)
        cycle_small = CycleDiagram(items=["A", "B", "C", "D"], theme=dark_theme)
        shapes_small = cycle_small.render(slide, left=0.5, top=0.5, width=4, height=3)
        assert len(shapes_small) == 8  # 4 shapes + 4 connectors

        # Medium cycle (6 items)
        cycle_medium = CycleDiagram(items=["1", "2", "3", "4", "5", "6"], theme=dark_theme)
        shapes_medium = cycle_medium.render(slide, left=5, top=0.5, width=4, height=3)
        assert len(shapes_medium) == 12  # 6 shapes + 6 connectors

    def test_empty_items(self, slide, dark_theme):
        """Test CycleDiagram with empty items."""
        cycle = CycleDiagram(items=[], theme=dark_theme)
        shapes = cycle.render(slide, left=1, top=1, width=6, height=5)
        assert len(shapes) == 0


class TestHierarchyDiagram:
    """Tests for HierarchyDiagram."""

    def test_init(self):
        """Test HierarchyDiagram initialization."""
        hierarchy = HierarchyDiagram(items=["CEO", "CTO", "CFO"])
        assert len(hierarchy.items) == 3

    def test_render(self, slide, dark_theme):
        """Test HierarchyDiagram rendering."""
        hierarchy = HierarchyDiagram(items=["CEO", "CTO", "CFO", "COO"], theme=dark_theme)
        shapes = hierarchy.render(slide, left=1, top=1, width=8, height=3)
        # 1 root + 3 children + 3 connectors = 7
        assert len(shapes) == 7

    def test_single_item(self, slide, dark_theme):
        """Test HierarchyDiagram with only root."""
        hierarchy = HierarchyDiagram(items=["CEO"], theme=dark_theme)
        shapes = hierarchy.render(slide, left=1, top=1, width=8, height=3)
        assert len(shapes) == 1  # Only root

    def test_empty_items(self, slide, dark_theme):
        """Test HierarchyDiagram with empty items."""
        hierarchy = HierarchyDiagram(items=[], theme=dark_theme)
        shapes = hierarchy.render(slide, left=1, top=1, width=8, height=3)
        assert len(shapes) == 0

    def test_many_children(self, slide, dark_theme):
        """Test HierarchyDiagram with many children."""
        items = ["CEO"] + [f"Exec {i}" for i in range(1, 6)]
        hierarchy = HierarchyDiagram(items=items, theme=dark_theme)
        shapes = hierarchy.render(slide, left=0.5, top=1, width=9, height=3)
        # 1 root + 5 children + 5 connectors = 11
        assert len(shapes) == 11
