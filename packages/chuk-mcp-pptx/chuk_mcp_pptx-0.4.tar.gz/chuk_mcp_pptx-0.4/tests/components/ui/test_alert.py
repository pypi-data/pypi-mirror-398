"""
Comprehensive tests for Alert component.
Tests variants, composition, and rendering.
"""

import pytest
from pptx import Presentation
from pptx.util import Inches

from chuk_mcp_pptx.components.core.alert import Alert
from chuk_mcp_pptx.themes import ThemeManager
from chuk_mcp_pptx.components.registry import get_component_schema, list_components


class TestAlertCreation:
    """Test Alert component creation and initialization."""

    def test_alert_creation_default(self):
        """Test creating an alert with default settings."""
        alert = Alert()
        assert alert.variant == "info"
        assert alert.show_icon

    def test_alert_creation_with_variant(self):
        """Test creating alerts with different variants."""
        variants = ["default", "info", "success", "warning", "error"]
        for variant in variants:
            alert = Alert(variant=variant)
            assert alert.variant == variant

    def test_alert_with_title_and_description(self):
        """Test alert with title and description."""
        alert = Alert(
            variant="success", title="Success!", description="Your changes have been saved."
        )
        assert alert.title_text == "Success!"
        assert alert.description_text == "Your changes have been saved."
        assert len(alert._children) == 2

    def test_alert_without_icon(self):
        """Test alert without icon."""
        alert = Alert(variant="info", show_icon=False)
        assert not alert.show_icon

    def test_alert_with_custom_theme(self):
        """Test alert with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("forest-dark")
        alert = Alert(variant="warning", theme=theme)
        assert alert.theme is not None


class TestAlertVariants:
    """Test all alert variants and their properties."""

    def test_default_variant(self):
        """Test default variant has correct properties."""
        alert = Alert(variant="default")
        props = alert.variant_props
        assert props.get("bg_color") == "background.secondary"
        assert props.get("icon") == "ℹ"

    def test_info_variant(self):
        """Test info variant has correct properties."""
        alert = Alert(variant="info")
        props = alert.variant_props
        assert props.get("bg_color") == "info.DEFAULT"
        assert props.get("icon") == "ℹ"

    def test_success_variant(self):
        """Test success variant has correct properties."""
        alert = Alert(variant="success")
        props = alert.variant_props
        assert props.get("bg_color") == "success.DEFAULT"
        assert props.get("icon") == "✓"

    def test_warning_variant(self):
        """Test warning variant has correct properties."""
        alert = Alert(variant="warning")
        props = alert.variant_props
        assert props.get("bg_color") == "warning.DEFAULT"
        assert props.get("icon") == "⚠"

    def test_error_variant(self):
        """Test error variant has correct properties."""
        alert = Alert(variant="error")
        props = alert.variant_props
        assert props.get("bg_color") == "destructive.DEFAULT"
        assert props.get("icon") == "✖"


class TestAlertComposition:
    """Test Alert composition patterns."""

    def test_alert_with_title_child(self):
        """Test adding title as child component."""
        alert = Alert(variant="info")
        alert.add_child(Alert.Title("Information"))
        assert len(alert._children) == 1

    def test_alert_with_description_child(self):
        """Test adding description as child component."""
        alert = Alert(variant="warning")
        alert.add_child(Alert.Description("Please be careful."))
        assert len(alert._children) == 1

    def test_alert_with_both_children(self):
        """Test adding both title and description."""
        alert = Alert(variant="success")
        alert.add_child(Alert.Title("Complete!"))
        alert.add_child(Alert.Description("The process finished successfully."))
        assert len(alert._children) == 2

    def test_alert_no_children(self):
        """Test alert without any children."""
        alert = Alert(variant="default")
        assert len(alert._children) == 0

    def test_alert_multiple_children(self):
        """Test alert with multiple children."""
        alert = Alert(variant="info")
        alert.add_child(Alert.Title("Title 1"))
        alert.add_child(Alert.Description("Description 1"))
        alert.add_child(Alert.Description("Description 2"))
        assert len(alert._children) == 3


class TestAlertRendering:
    """Test alert rendering to PowerPoint slides."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_alert_renders_to_slide(self, slide):
        """Test that alert renders successfully."""
        alert = Alert(variant="success", title="Success", description="Done!")
        shape = alert.render(slide, left=1.0, top=1.0)
        assert shape is not None
        assert shape in slide.shapes

    def test_alert_position(self, slide):
        """Test alert is positioned correctly."""
        alert = Alert(variant="info")
        shape = alert.render(slide, left=2.0, top=3.0)
        assert shape.left == Inches(2.0)
        assert shape.top == Inches(3.0)

    def test_alert_custom_dimensions(self, slide):
        """Test alert with custom width and height."""
        alert = Alert(variant="warning")
        shape = alert.render(slide, left=1.0, top=1.0, width=6.0, height=2.0)
        assert shape.width == Inches(6.0)
        assert shape.height == Inches(2.0)

    def test_alert_default_dimensions(self, slide):
        """Test alert with default dimensions."""
        alert = Alert(variant="error")
        shape = alert.render(slide, left=1.0, top=1.0)
        # Default width and height from function signature
        assert shape.width == Inches(5.0)
        assert shape.height == Inches(1.5)

    def test_alert_with_icon_renders(self, slide):
        """Test alert with icon renders."""
        alert = Alert(variant="success", title="Done", show_icon=True)
        shape = alert.render(slide, left=1.0, top=1.0)
        # Icon should be in the text (✓ for success)
        assert "✓" in shape.text_frame.text

    def test_alert_without_icon_renders(self, slide):
        """Test alert without icon renders."""
        alert = Alert(variant="success", title="Done", show_icon=False)
        shape = alert.render(slide, left=1.0, top=1.0)
        # Icon should NOT be in the text
        assert "✓" not in shape.text_frame.text

    def test_all_variant_combinations_render(self, slide):
        """Test all variants render successfully."""
        variants = ["default", "info", "success", "warning", "error"]

        top = 0.5
        for variant in variants:
            alert = Alert(variant=variant, title=variant.title())
            shape = alert.render(slide, left=1.0, top=top)
            assert shape is not None
            top += 1.7


class TestAlertThemeIntegration:
    """Test alert integration with theme system."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_alert_uses_theme_colors(self):
        """Test alert uses theme colors correctly."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("default-dark")
        alert = Alert(variant="success", theme=theme)
        color = alert.get_color("success.DEFAULT")
        assert color is not None

    def test_alert_with_different_themes(self, slide):
        """Test alert renders with different themes."""
        theme_manager = ThemeManager()
        themes = ["default-light", "default-dark", "ocean-light", "forest-dark"]

        for theme_name in themes:
            theme = theme_manager.get_theme(theme_name)
            alert = Alert(variant="info", title="Test", theme=theme)
            shape = alert.render(slide, left=1.0, top=1.0)
            assert shape is not None


class TestAlertRegistry:
    """Test alert component registry integration."""

    def test_alert_is_registered(self):
        """Test Alert is in component registry."""
        components = list_components()
        assert "Alert" in components

    def test_alert_schema_exists(self):
        """Test Alert has a valid schema."""
        schema = get_component_schema("Alert")
        assert schema is not None
        assert schema["name"] == "Alert"

    def test_alert_schema_has_props(self):
        """Test Alert schema has properties."""
        schema = get_component_schema("Alert")
        assert "schema" in schema
        prop_names = list(schema["schema"]["properties"].keys())
        assert "variant" in prop_names
        assert "title" in prop_names
        assert "description" in prop_names
        assert "show_icon" in prop_names

    def test_alert_schema_has_variants(self):
        """Test Alert schema has variant definitions."""
        schema = get_component_schema("Alert")
        assert "variants" in schema
        assert "variant" in schema["variants"]
        assert "default" in schema["variants"]["variant"]
        assert "success" in schema["variants"]["variant"]

    def test_alert_schema_has_composition(self):
        """Test Alert schema has composition info."""
        schema = get_component_schema("Alert")
        assert "composition" in schema
        # Alert supports composition
        assert schema["composition"] is not None

    def test_alert_schema_has_examples(self):
        """Test Alert schema has usage examples."""
        schema = get_component_schema("Alert")
        assert "examples" in schema
        assert len(schema["examples"]) >= 2


class TestAlertEdgeCases:
    """Test alert edge cases and error handling."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_alert_empty_title(self, slide):
        """Test alert with empty title."""
        alert = Alert(variant="info", title="")
        shape = alert.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_alert_empty_description(self, slide):
        """Test alert with empty description."""
        alert = Alert(variant="warning", description="")
        shape = alert.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_alert_only_title(self, slide):
        """Test alert with only title."""
        alert = Alert(variant="success", title="Success!")
        shape = alert.render(slide, left=1.0, top=1.0)
        assert "Success!" in shape.text_frame.text

    def test_alert_only_description(self, slide):
        """Test alert with only description."""
        alert = Alert(variant="info", description="Please note this.")
        shape = alert.render(slide, left=1.0, top=1.0)
        assert "Please note this." in shape.text_frame.text

    def test_alert_very_long_title(self, slide):
        """Test alert with very long title."""
        long_title = "This is an extremely long title that might cause layout issues in the alert"
        alert = Alert(variant="warning", title=long_title)
        shape = alert.render(slide, left=1.0, top=1.0, width=8.0, height=2.0)
        assert shape is not None

    def test_alert_very_long_description(self, slide):
        """Test alert with very long description."""
        long_desc = "This is a very long description with lots of details that might need to wrap across multiple lines in the alert box to display properly."
        alert = Alert(variant="error", description=long_desc)
        shape = alert.render(slide, left=1.0, top=1.0, width=7.0, height=3.0)
        assert shape is not None

    def test_alert_special_characters(self, slide):
        """Test alert with special characters."""
        alert = Alert(
            variant="info",
            title="Error #404",
            description="Path: /api/v1/users?filter=active&sort=desc",
        )
        shape = alert.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_alert_unicode_emoji(self, slide):
        """Test alert with emoji."""
        alert = Alert(
            variant="success",
            title="🎉 Congratulations!",
            description="You've completed the tutorial 🚀",
        )
        shape = alert.render(slide, left=1.0, top=1.0)
        assert shape is not None

    def test_alert_multiline_description(self, slide):
        """Test alert with multiline description."""
        alert = Alert(variant="warning", title="Important", description="Line 1\nLine 2\nLine 3")
        shape = alert.render(slide, left=1.0, top=1.0, height=2.5)
        assert shape is not None

    def test_alert_very_small_dimensions(self, slide):
        """Test alert with very small dimensions."""
        alert = Alert(variant="info", title="Tiny")
        shape = alert.render(slide, left=1.0, top=1.0, width=1.0, height=0.5)
        assert shape.width == Inches(1.0)
        assert shape.height == Inches(0.5)

    def test_alert_very_large_dimensions(self, slide):
        """Test alert with very large dimensions."""
        alert = Alert(
            variant="success",
            title="Large Alert",
            description="This alert takes up a lot of space.",
        )
        shape = alert.render(slide, left=0.5, top=0.5, width=9.0, height=6.0)
        assert shape.width == Inches(9.0)
        assert shape.height == Inches(6.0)

    def test_alert_at_slide_boundaries(self, slide):
        """Test alert positioned at slide boundaries."""
        alert = Alert(variant="error", title="Edge")
        # Near right edge
        shape = alert.render(slide, left=5.0, top=1.0, width=4.0)
        assert shape is not None
        # Near bottom edge
        shape = alert.render(slide, left=1.0, top=6.0, height=1.0)
        assert shape is not None
