"""
Tests for text components (TextBox, BulletList).
"""

import pytest
from pptx import Presentation

from chuk_mcp_pptx.components.core import TextBox, BulletList


@pytest.fixture
def presentation():
    """Create a presentation for rendering tests."""
    return Presentation()


@pytest.fixture
def slide(presentation):
    """Create a slide for rendering tests."""
    blank_layout = presentation.slide_layouts[6]
    return presentation.slides.add_slide(blank_layout)


class TestTextBox:
    """Tests for TextBox component."""

    def test_init(self):
        """Test TextBox initialization."""
        textbox = TextBox(text="Hello World")
        assert textbox.text == "Hello World"
        assert textbox.font_size == 18
        assert textbox.alignment == "left"

    def test_render_simple(self, slide, dark_theme):
        """Test simple textbox rendering."""
        textbox = TextBox(text="Test Text", theme=dark_theme)
        rendered = textbox.render(slide, left=1, top=1, width=4, height=1)
        assert rendered is not None
        assert rendered.has_text_frame
        assert rendered.text_frame.text == "Test Text"

    def test_render_with_formatting(self, slide, dark_theme):
        """Test textbox with formatting."""
        textbox = TextBox(
            text="Bold Title", font_size=24, bold=True, alignment="center", theme=dark_theme
        )
        rendered = textbox.render(slide, left=1, top=1, width=6, height=1.5)
        assert rendered is not None
        assert rendered.text_frame.text == "Bold Title"

    def test_render_with_hex_color(self, slide, dark_theme):
        """Test textbox with hex color."""
        textbox = TextBox(text="Colored Text", color="#FF5733", theme=dark_theme)
        rendered = textbox.render(slide, left=1, top=1, width=4, height=1)
        assert rendered is not None

    def test_render_with_semantic_color(self, slide, dark_theme):
        """Test textbox with semantic color."""
        textbox = TextBox(text="Theme Text", color="primary.DEFAULT", theme=dark_theme)
        rendered = textbox.render(slide, left=1, top=1, width=4, height=1)
        assert rendered is not None

    def test_render_with_alignment(self, slide, dark_theme):
        """Test textbox with different alignments."""
        alignments = ["left", "center", "right", "justify"]

        for alignment in alignments:
            textbox = TextBox(
                text=f"{alignment.title()} aligned", alignment=alignment, theme=dark_theme
            )
            rendered = textbox.render(slide, left=1, top=1, width=4, height=1)
            assert rendered is not None

    def test_render_with_auto_fit(self, slide, dark_theme):
        """Test textbox with auto-fit."""
        textbox = TextBox(
            text="This is a longer text that should auto-fit", auto_fit=True, theme=dark_theme
        )
        rendered = textbox.render(slide, left=1, top=1, width=4, height=1)
        assert rendered is not None


class TestBulletList:
    """Tests for BulletList component."""

    def test_init(self):
        """Test BulletList initialization."""
        bullets = BulletList(items=["Item 1", "Item 2", "Item 3"])
        assert len(bullets.items) == 3
        assert bullets.bullet_char == "•"

    def test_render_simple(self, slide, dark_theme):
        """Test simple bullet list rendering."""
        bullets = BulletList(items=["First", "Second", "Third"], theme=dark_theme)
        rendered = bullets.render(slide, left=1, top=2, width=8, height=4)
        assert rendered is not None
        assert rendered.has_text_frame

    def test_render_with_custom_bullet(self, slide, dark_theme):
        """Test bullet list with custom bullet character."""
        bullets = BulletList(
            items=["Step 1", "Step 2", "Step 3"], bullet_char="→", theme=dark_theme
        )
        rendered = bullets.render(slide, left=1, top=2, width=8, height=4)
        assert rendered is not None

    def test_render_with_formatting(self, slide, dark_theme):
        """Test bullet list with formatting."""
        bullets = BulletList(
            items=["Important", "Critical", "Essential"],
            font_size=20,
            color="primary.DEFAULT",
            theme=dark_theme,
        )
        rendered = bullets.render(slide, left=1, top=2, width=8, height=4)
        assert rendered is not None

    def test_render_empty_list(self, slide, dark_theme):
        """Test bullet list with empty items."""
        bullets = BulletList(items=[], theme=dark_theme)
        rendered = bullets.render(slide, left=1, top=2, width=8, height=4)
        assert rendered is not None

    def test_render_single_item(self, slide, dark_theme):
        """Test bullet list with single item."""
        bullets = BulletList(items=["Only one"], theme=dark_theme)
        rendered = bullets.render(slide, left=1, top=2, width=8, height=4)
        assert rendered is not None

    def test_render_many_items(self, slide, dark_theme):
        """Test bullet list with many items."""
        items = [f"Item {i}" for i in range(1, 11)]
        bullets = BulletList(items=items, theme=dark_theme)
        rendered = bullets.render(slide, left=1, top=1, width=8, height=6)
        assert rendered is not None

    def test_render_with_hex_color(self, slide, dark_theme):
        """Test bullet list with hex color."""
        bullets = BulletList(items=["Red", "Green", "Blue"], color="#FF0000", theme=dark_theme)
        rendered = bullets.render(slide, left=1, top=2, width=8, height=4)
        assert rendered is not None
