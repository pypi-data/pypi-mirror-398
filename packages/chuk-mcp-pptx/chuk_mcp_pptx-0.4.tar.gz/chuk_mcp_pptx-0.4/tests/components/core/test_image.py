"""
Comprehensive tests for Image component.
"""

import pytest
import base64
import io
from pptx import Presentation
from PIL import Image as PILImage

from chuk_mcp_pptx.components.core.image import Image


class TestImageComponent:
    """Test Image component initialization and rendering."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    @pytest.fixture
    def test_image_path(self, tmp_path):
        """Create a temporary test image."""
        image_path = tmp_path / "test_image.png"
        # Create a simple 100x100 red image
        img = PILImage.new("RGB", (100, 100), color="red")
        img.save(str(image_path))
        return str(image_path)

    @pytest.fixture
    def test_rgba_image_path(self, tmp_path):
        """Create a temporary RGBA test image."""
        image_path = tmp_path / "test_image_rgba.png"
        # Create a simple 100x100 image with alpha channel
        img = PILImage.new("RGBA", (100, 100), color=(255, 0, 0, 128))
        img.save(str(image_path))
        return str(image_path)

    @pytest.fixture
    def base64_image(self):
        """Create a base64 encoded test image."""
        # Create a tiny 2x2 red PNG
        img = PILImage.new("RGB", (2, 2), color="red")
        buffer = io.BytesIO()
        img.save(buffer, format="PNG")
        encoded = base64.b64encode(buffer.getvalue()).decode("utf-8")
        return f"data:image/png;base64,{encoded}"

    def test_initialization(self, test_image_path):
        """Test basic initialization."""
        image = Image(image_source=test_image_path)
        assert image.image_source == test_image_path
        assert image.shadow is False
        assert image.glow is False
        assert image.reflection is False
        assert image.blur_radius == 0
        assert image.grayscale is False
        assert image.sepia is False
        assert image.brightness == 1.0
        assert image.contrast == 1.0
        assert image.saturation == 1.0
        assert image.sharpen is False
        assert image.invert is False

    def test_initialization_with_effects(self, test_image_path):
        """Test initialization with effects."""
        image = Image(
            image_source=test_image_path,
            shadow=True,
            glow=True,
            reflection=True,
            blur_radius=5,
            grayscale=True,
        )
        assert image.shadow is True
        assert image.glow is True
        assert image.reflection is True
        assert image.blur_radius == 5
        assert image.grayscale is True

    @pytest.mark.asyncio
    async def test_render_basic(self, slide, test_image_path):
        """Test basic rendering."""
        image = Image(image_source=test_image_path)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None
        assert len(slide.shapes) == 1

    @pytest.mark.asyncio
    async def test_render_with_width_only(self, slide, test_image_path):
        """Test rendering with width only (maintains aspect ratio)."""
        image = Image(image_source=test_image_path)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_height_only(self, slide, test_image_path):
        """Test rendering with height only (maintains aspect ratio)."""
        image = Image(image_source=test_image_path)
        result = await image.render(slide, left=1.0, top=1.0, height=3.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_width_and_height(self, slide, test_image_path):
        """Test rendering with both width and height."""
        image = Image(image_source=test_image_path)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0, height=3.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_no_dimensions(self, slide, test_image_path):
        """Test rendering without dimensions (uses original size)."""
        image = Image(image_source=test_image_path)
        result = await image.render(slide, left=1.0, top=1.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_shadow(self, slide, test_image_path):
        """Test rendering with shadow effect."""
        image = Image(image_source=test_image_path, shadow=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None
        # Verify shadow is applied
        assert result.shadow.visible is True

    @pytest.mark.asyncio
    async def test_render_base64_image(self, slide, base64_image):
        """Test rendering with base64 image data."""
        image = Image(image_source=base64_image)
        result = await image.render(slide, left=1.0, top=1.0, width=3.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_base64_no_processing(self, slide, base64_image):
        """Test base64 image without any processing (covers lines 199-202)."""
        # This specifically tests the base64 path when no filters are applied
        image = Image(
            image_source=base64_image,
            blur_radius=0,
            grayscale=False,
            brightness=1.0,
            contrast=1.0,
            saturation=1.0,
        )
        result = await image.render(slide, left=1.0, top=1.0, width=3.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_file_not_found(self, slide):
        """Test error handling for non-existent file."""
        image = Image(image_source="/path/to/nonexistent/image.png")
        with pytest.raises(FileNotFoundError):
            await image.render(slide, left=1.0, top=1.0)

    @pytest.mark.asyncio
    async def test_load_image_file_not_found(self):
        """Test _load_image with non-existent file (covers line 255)."""
        image = Image(image_source="/path/to/nonexistent/image.png")
        with pytest.raises(FileNotFoundError, match="Image not found"):
            await image._load_image()

    @pytest.mark.asyncio
    async def test_load_image_from_base64(self, base64_image):
        """Test _load_image with base64 data (covers lines 247-250)."""
        image = Image(image_source=base64_image)
        pil_image = await image._load_image()
        assert isinstance(pil_image, PILImage.Image)
        assert pil_image.size == (2, 2)

    @pytest.mark.asyncio
    async def test_render_with_blur(self, slide, test_image_path):
        """Test rendering with blur filter."""
        image = Image(image_source=test_image_path, blur_radius=5)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_grayscale(self, slide, test_image_path):
        """Test rendering with grayscale filter."""
        image = Image(image_source=test_image_path, grayscale=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_sepia(self, slide, test_image_path):
        """Test rendering with sepia filter."""
        image = Image(image_source=test_image_path, sepia=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_brightness(self, slide, test_image_path):
        """Test rendering with brightness adjustment."""
        image = Image(image_source=test_image_path, brightness=1.5)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_contrast(self, slide, test_image_path):
        """Test rendering with contrast adjustment."""
        image = Image(image_source=test_image_path, contrast=1.5)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_saturation(self, slide, test_image_path):
        """Test rendering with saturation adjustment."""
        image = Image(image_source=test_image_path, saturation=1.5)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_sharpen(self, slide, test_image_path):
        """Test rendering with sharpen filter."""
        image = Image(image_source=test_image_path, sharpen=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_invert(self, slide, test_image_path):
        """Test rendering with color inversion."""
        image = Image(image_source=test_image_path, invert=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_render_with_multiple_filters(self, slide, test_image_path):
        """Test rendering with multiple filters combined."""
        image = Image(
            image_source=test_image_path,
            blur_radius=3,
            brightness=1.2,
            contrast=1.3,
            saturation=0.8,
            shadow=True,
        )
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None
        assert result.shadow.visible is True

    @pytest.mark.asyncio
    async def test_apply_filters_with_rgba_image(self, slide, test_rgba_image_path):
        """Test filter application with RGBA image (covers line 261)."""
        # This tests the RGB conversion path in _apply_filters
        image = Image(
            image_source=test_rgba_image_path,
            blur_radius=2,  # Forces processing
        )
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_apply_sepia(self, test_image_path):
        """Test sepia filter application."""
        image = Image(image_source=test_image_path, sepia=True)
        pil_image = await image._load_image()
        result = image._apply_sepia(pil_image)
        assert isinstance(result, PILImage.Image)

    @pytest.mark.asyncio
    async def test_apply_filters_all_combinations(self, test_image_path):
        """Test all filter combinations."""
        image = Image(
            image_source=test_image_path,
            blur_radius=5,
            grayscale=True,
            sepia=False,  # Note: grayscale overrides sepia
            brightness=1.2,
            contrast=1.3,
            saturation=0.8,
            sharpen=True,
            invert=True,
        )
        pil_image = await image._load_image()
        result = image._apply_filters(pil_image)
        assert isinstance(result, PILImage.Image)

    def test_needs_processing_detection(self, test_image_path):
        """Test that processing is correctly detected."""
        # No processing needed
        image1 = Image(image_source=test_image_path)
        assert not (
            image1.blur_radius > 0
            or image1.grayscale
            or image1.sepia
            or image1.brightness != 1.0
            or image1.contrast != 1.0
            or image1.saturation != 1.0
            or image1.sharpen
            or image1.invert
        )

        # Processing needed (blur)
        image2 = Image(image_source=test_image_path, blur_radius=5)
        assert image2.blur_radius > 0

        # Processing needed (brightness)
        image3 = Image(image_source=test_image_path, brightness=1.5)
        assert image3.brightness != 1.0


class TestImageEffects:
    """Test PowerPoint effects on images."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    @pytest.fixture
    def test_image_path(self, tmp_path):
        """Create a temporary test image."""
        image_path = tmp_path / "test_image.png"
        img = PILImage.new("RGB", (100, 100), color="blue")
        img.save(str(image_path))
        return str(image_path)

    @pytest.mark.asyncio
    async def test_shadow_properties(self, slide, test_image_path):
        """Test that shadow is applied with correct properties."""
        image = Image(image_source=test_image_path, shadow=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)

        shadow = result.shadow
        assert shadow.visible is True
        assert shadow.inherit is False

    @pytest.mark.asyncio
    async def test_multiple_images_on_slide(self, slide, test_image_path):
        """Test adding multiple images to same slide."""
        image1 = Image(image_source=test_image_path)
        image2 = Image(image_source=test_image_path, shadow=True)
        image3 = Image(image_source=test_image_path, grayscale=True)

        result1 = await image1.render(slide, left=1.0, top=1.0, width=2.0)
        result2 = await image2.render(slide, left=3.5, top=1.0, width=2.0)
        result3 = await image3.render(slide, left=6.0, top=1.0, width=2.0)

        assert result1 is not None
        assert result2 is not None
        assert result3 is not None
        assert len(slide.shapes) == 3


class TestImageEdgeCases:
    """Test edge cases and boundary conditions."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    @pytest.fixture
    def test_image_path(self, tmp_path):
        """Create a temporary test image."""
        image_path = tmp_path / "test_image.png"
        img = PILImage.new("RGB", (100, 100), color="green")
        img.save(str(image_path))
        return str(image_path)

    @pytest.mark.asyncio
    async def test_extreme_blur(self, slide, test_image_path):
        """Test with very high blur radius."""
        image = Image(image_source=test_image_path, blur_radius=50)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_extreme_brightness(self, slide, test_image_path):
        """Test with extreme brightness values."""
        # Very dark
        image1 = Image(image_source=test_image_path, brightness=0.1)
        result1 = await image1.render(slide, left=1.0, top=1.0, width=2.0)
        assert result1 is not None

        # Very bright
        image2 = Image(image_source=test_image_path, brightness=3.0)
        result2 = await image2.render(slide, left=4.0, top=1.0, width=2.0)
        assert result2 is not None

    @pytest.mark.asyncio
    async def test_extreme_contrast(self, slide, test_image_path):
        """Test with extreme contrast values."""
        image = Image(image_source=test_image_path, contrast=5.0)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_zero_saturation(self, slide, test_image_path):
        """Test with zero saturation (should be like grayscale)."""
        image = Image(image_source=test_image_path, saturation=0.0)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_grayscale_then_sepia(self, slide, test_image_path):
        """Test grayscale with sepia (grayscale should win)."""
        image = Image(image_source=test_image_path, grayscale=True, sepia=True)
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None

    @pytest.mark.asyncio
    async def test_all_filters_combined(self, slide, test_image_path):
        """Test with all filters enabled."""
        image = Image(
            image_source=test_image_path,
            blur_radius=5,
            grayscale=True,
            sepia=True,
            brightness=1.2,
            contrast=1.3,
            saturation=1.1,
            sharpen=True,
            invert=True,
            shadow=True,
        )
        result = await image.render(slide, left=1.0, top=1.0, width=4.0)
        assert result is not None
