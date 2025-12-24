"""
Comprehensive tests for Timeline component.
"""

import pytest
from pptx import Presentation

from chuk_mcp_pptx.components.core.timeline import Timeline
from chuk_mcp_pptx.themes import ThemeManager


class TestTimelineComponent:
    """Test Timeline component initialization and rendering."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    @pytest.fixture
    def simple_events(self):
        """Create simple event data."""
        return [
            {"date": "Jan 2024", "title": "Project Start"},
            {"date": "Mar 2024", "title": "Beta Release"},
            {"date": "Jun 2024", "title": "Launch"},
        ]

    @pytest.fixture
    def detailed_events(self):
        """Create detailed event data with descriptions."""
        return [
            {
                "date": "Q1 2024",
                "title": "Planning",
                "description": "Initial planning and research",
            },
            {"date": "Q2 2024", "title": "Development", "description": "Build core features"},
            {
                "date": "Q3 2024",
                "title": "Launch",
                "description": "Public release",
                "highlight": True,
            },
        ]

    @pytest.fixture
    def highlighted_events(self):
        """Create events with highlights."""
        return [
            {"date": "2023", "title": "Founded", "highlight": False},
            {"date": "2024", "title": "Series A", "highlight": True},
            {"date": "2025", "title": "Expansion", "highlight": False},
        ]

    def test_initialization(self, simple_events):
        """Test basic initialization."""
        timeline = Timeline(simple_events)
        assert timeline.events == simple_events
        assert timeline.variant == "default"
        assert timeline.style == "line"
        assert timeline.show_descriptions is False

    def test_initialization_with_variant(self, simple_events):
        """Test initialization with variant."""
        timeline = Timeline(simple_events, variant="highlighted")
        assert timeline.variant == "highlighted"

    def test_initialization_with_style(self, simple_events):
        """Test initialization with style."""
        timeline = Timeline(simple_events, style="arrow")
        assert timeline.style == "arrow"

    def test_initialization_with_descriptions(self, detailed_events):
        """Test initialization with show_descriptions enabled."""
        timeline = Timeline(detailed_events, show_descriptions=True)
        assert timeline.show_descriptions is True

    def test_render_basic_timeline(self, slide, simple_events):
        """Test rendering a basic timeline."""
        timeline = Timeline(simple_events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0
        # Should have 1 line + (marker + date + title) * 3 events = 10 shapes
        assert len(shapes) == 10

    def test_render_empty_events(self, slide):
        """Test rendering with empty events list."""
        timeline = Timeline([])
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) == 0

    def test_render_single_event(self, slide):
        """Test rendering with single event."""
        events = [{"date": "2024", "title": "Event"}]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0
        # Should have 1 line + marker + date + title = 4 shapes
        assert len(shapes) == 4

    def test_render_with_descriptions(self, slide, detailed_events):
        """Test rendering with descriptions enabled (covers lines 173-179, 291-309)."""
        timeline = Timeline(detailed_events, show_descriptions=True)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0
        # Should have 1 line + (marker + date + title + description) * 3 events = 13 shapes
        assert len(shapes) == 13

    def test_render_with_descriptions_missing_key(self, slide):
        """Test rendering with descriptions enabled but event missing description key."""
        events = [
            {"date": "Q1", "title": "Event 1", "description": "Has description"},
            {"date": "Q2", "title": "Event 2"},  # No description
        ]
        timeline = Timeline(events, show_descriptions=True)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        # Should have 1 line + (marker + date + title + desc) + (marker + date + title) = 8 shapes
        assert len(shapes) == 8

    def test_render_arrow_style(self, slide, simple_events):
        """Test rendering with arrow style."""
        timeline = Timeline(simple_events, style="arrow")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_render_segmented_style(self, slide, simple_events):
        """Test rendering with segmented style."""
        timeline = Timeline(simple_events, style="segmented")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_render_highlighted_variant(self, slide, highlighted_events):
        """Test rendering with highlighted variant."""
        timeline = Timeline(highlighted_events, variant="highlighted")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_render_minimal_variant(self, slide, simple_events):
        """Test rendering with minimal variant."""
        timeline = Timeline(simple_events, variant="minimal")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_render_with_theme(self, slide, simple_events):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean")

        timeline = Timeline(simple_events, theme=theme)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_render_custom_width(self, slide, simple_events):
        """Test rendering with custom width."""
        timeline = Timeline(simple_events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=6.0)
        assert len(shapes) > 0

    def test_render_custom_position(self, slide, simple_events):
        """Test rendering at custom position."""
        timeline = Timeline(simple_events)
        shapes = timeline.render(slide, left=2.0, top=4.0, width=8.0)
        assert len(shapes) > 0

    def test_get_line_color(self, simple_events):
        """Test line color retrieval."""
        timeline = Timeline(simple_events)
        color = timeline._get_line_color()
        assert color is not None

    def test_get_marker_color_normal(self, simple_events):
        """Test marker color for normal events."""
        timeline = Timeline(simple_events)
        color = timeline._get_marker_color(is_highlighted=False)
        assert color is not None

    def test_get_marker_color_highlighted(self, simple_events):
        """Test marker color for highlighted events."""
        timeline = Timeline(simple_events)
        color = timeline._get_marker_color(is_highlighted=True)
        assert color is not None

    def test_get_marker_color_highlighted_variant(self, simple_events):
        """Test marker color with highlighted variant."""
        timeline = Timeline(simple_events, variant="highlighted")
        color = timeline._get_marker_color(is_highlighted=False)
        assert color is not None

    def test_get_highlight_color(self, simple_events):
        """Test highlight color retrieval."""
        timeline = Timeline(simple_events)
        color = timeline._get_highlight_color()
        assert color is not None

    def test_render_marker_normal(self, slide, simple_events):
        """Test marker rendering for normal events."""
        timeline = Timeline(simple_events)
        marker = timeline._render_marker(slide, 3.0, 3.0, is_highlighted=False)
        assert marker is not None

    def test_render_marker_highlighted(self, slide, simple_events):
        """Test marker rendering for highlighted events."""
        timeline = Timeline(simple_events)
        marker = timeline._render_marker(slide, 3.0, 3.0, is_highlighted=True)
        assert marker is not None

    def test_render_date(self, slide, simple_events):
        """Test date rendering."""
        timeline = Timeline(simple_events)
        date_shape = timeline._render_date(slide, 3.0, 2.5, "Jan 2024")
        assert date_shape is not None

    def test_render_title_normal(self, slide, simple_events):
        """Test title rendering for normal events."""
        timeline = Timeline(simple_events)
        title_shape = timeline._render_title(slide, 3.0, 3.0, "Event Title", is_highlighted=False)
        assert title_shape is not None

    def test_render_title_highlighted(self, slide, simple_events):
        """Test title rendering for highlighted events."""
        timeline = Timeline(simple_events)
        title_shape = timeline._render_title(slide, 3.0, 3.0, "Event Title", is_highlighted=True)
        assert title_shape is not None

    def test_render_description(self, slide, simple_events):
        """Test description rendering (covers lines 291-309)."""
        timeline = Timeline(simple_events, show_descriptions=True)
        desc_shape = timeline._render_description(slide, 3.0, 3.5, "Event description text")
        assert desc_shape is not None

    def test_render_line_default_style(self, slide, simple_events):
        """Test line rendering with default style."""
        timeline = Timeline(simple_events, style="line")
        shapes = timeline._render_line(slide, 1.0, 3.0, 8.0)
        assert len(shapes) == 1

    def test_render_line_arrow_style(self, slide, simple_events):
        """Test line rendering with arrow style."""
        timeline = Timeline(simple_events, style="arrow")
        shapes = timeline._render_line(slide, 1.0, 3.0, 8.0)
        assert len(shapes) == 1


class TestTimelineEdgeCases:
    """Test edge cases and boundary conditions."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_many_events(self, slide):
        """Test timeline with many events."""
        events = [{"date": f"Day {i}", "title": f"Event {i}"} for i in range(10)]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_long_text_values(self, slide):
        """Test timeline with very long text."""
        events = [
            {
                "date": "Q1 2024",
                "title": "Very long event title that spans multiple words and characters",
                "description": "This is an extremely long description that contains a lot of text and information about the event. It should wrap properly within the text box.",
            }
        ]
        timeline = Timeline(events, show_descriptions=True)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_empty_text_values(self, slide):
        """Test timeline with empty text values."""
        events = [{"date": "", "title": ""}, {"date": "Q1", "title": ""}]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_all_events_highlighted(self, slide):
        """Test timeline with all events highlighted."""
        events = [
            {"date": "Q1", "title": "Event 1", "highlight": True},
            {"date": "Q2", "title": "Event 2", "highlight": True},
            {"date": "Q3", "title": "Event 3", "highlight": True},
        ]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_minimal_width(self, slide):
        """Test timeline with minimal width."""
        events = [{"date": "2024", "title": "Event"}]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=1.0)
        assert len(shapes) > 0

    def test_maximum_width(self, slide):
        """Test timeline with maximum width."""
        events = [{"date": "2020", "title": "Start"}, {"date": "2025", "title": "End"}]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=0.5, top=3.0, width=9.0)
        assert len(shapes) > 0

    def test_descriptions_with_special_characters(self, slide):
        """Test descriptions with special characters."""
        events = [
            {
                "date": "2024",
                "title": "Event!",
                "description": "Description with special chars: @#$%^&*()",
            }
        ]
        timeline = Timeline(events, show_descriptions=True)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_mixed_highlight_states(self, slide):
        """Test timeline with mixed highlight states."""
        events = [
            {"date": "Q1", "title": "Event 1", "highlight": True},
            {"date": "Q2", "title": "Event 2", "highlight": False},
            {"date": "Q3", "title": "Event 3"},  # No highlight key
            {"date": "Q4", "title": "Event 4", "highlight": True},
        ]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_two_events(self, slide):
        """Test timeline with exactly two events."""
        events = [{"date": "Start", "title": "Beginning"}, {"date": "End", "title": "Finish"}]
        timeline = Timeline(events)
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0


class TestTimelineVariantCombinations:
    """Test different variant and style combinations."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    @pytest.fixture
    def test_events(self):
        """Create test events."""
        return [
            {"date": "Q1", "title": "Phase 1", "description": "First phase"},
            {"date": "Q2", "title": "Phase 2", "description": "Second phase", "highlight": True},
            {"date": "Q3", "title": "Phase 3", "description": "Third phase"},
        ]

    def test_default_variant_line_style(self, slide, test_events):
        """Test default variant with line style."""
        timeline = Timeline(test_events, variant="default", style="line")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_default_variant_arrow_style(self, slide, test_events):
        """Test default variant with arrow style."""
        timeline = Timeline(test_events, variant="default", style="arrow")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_highlighted_variant_line_style(self, slide, test_events):
        """Test highlighted variant with line style."""
        timeline = Timeline(test_events, variant="highlighted", style="line")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_highlighted_variant_arrow_style(self, slide, test_events):
        """Test highlighted variant with arrow style."""
        timeline = Timeline(test_events, variant="highlighted", style="arrow")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_minimal_variant_line_style(self, slide, test_events):
        """Test minimal variant with line style."""
        timeline = Timeline(test_events, variant="minimal", style="line")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_minimal_variant_arrow_style(self, slide, test_events):
        """Test minimal variant with arrow style."""
        timeline = Timeline(test_events, variant="minimal", style="arrow")
        shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
        assert len(shapes) > 0

    def test_all_combinations_with_descriptions(self, slide, test_events):
        """Test all variant/style combinations with descriptions enabled."""
        variants = ["default", "minimal", "highlighted"]
        styles = ["line", "arrow", "segmented"]

        for variant in variants:
            for style in styles:
                timeline = Timeline(
                    test_events, variant=variant, style=style, show_descriptions=True
                )
                shapes = timeline.render(slide, left=1.0, top=3.0, width=8.0)
                assert len(shapes) > 0
