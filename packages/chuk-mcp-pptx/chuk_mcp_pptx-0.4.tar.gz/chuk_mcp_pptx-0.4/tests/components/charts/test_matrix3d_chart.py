"""
Tests for Matrix3DChart component.
"""

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from chuk_mcp_pptx.components.charts.scatter_bubble import Matrix3DChart


class TestMatrix3DChart:
    """Test Matrix3DChart component."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_initialization_without_color_field(self):
        """Test matrix 3D chart initialization without color grouping."""
        data_points = [
            {"x": 1, "y": 2, "size": 10},
            {"x": 2, "y": 3, "size": 20},
            {"x": 3, "y": 4, "size": 30},
        ]
        chart = Matrix3DChart(data_points=data_points, x_field="x", y_field="y", size_field="size")
        assert chart.x_field == "x"
        assert chart.y_field == "y"
        assert chart.size_field == "size"
        assert chart.color_field is None

    def test_initialization_with_color_field(self):
        """Test matrix 3D chart initialization with color grouping."""
        data_points = [
            {"x": 1, "y": 2, "size": 10, "category": "A"},
            {"x": 2, "y": 3, "size": 20, "category": "B"},
        ]
        chart = Matrix3DChart(
            data_points=data_points,
            x_field="x",
            y_field="y",
            size_field="size",
            color_field="category",
        )
        assert chart.color_field == "category"

    def test_convert_to_series_single_series(self):
        """Test converting data to single series format."""
        data_points = [{"x": 1, "y": 2, "size": 10}, {"x": 2, "y": 3, "size": 20}]
        chart = Matrix3DChart(data_points=data_points, x_field="x", y_field="y", size_field="size")
        series_data = chart._convert_to_series()
        assert len(series_data) == 1
        assert series_data[0]["name"] == "Data"
        assert len(series_data[0]["points"]) == 2
        # Check format is [x, y, size]
        assert series_data[0]["points"][0] == [1, 2, 10]

    def test_convert_to_series_grouped(self):
        """Test converting data with color grouping."""
        data_points = [
            {"x": 1, "y": 2, "size": 10, "category": "A"},
            {"x": 2, "y": 3, "size": 20, "category": "A"},
            {"x": 3, "y": 4, "size": 30, "category": "B"},
        ]
        chart = Matrix3DChart(
            data_points=data_points,
            x_field="x",
            y_field="y",
            size_field="size",
            color_field="category",
        )
        series_data = chart._convert_to_series()
        assert len(series_data) == 2
        # Find series by name
        series_names = {s["name"] for s in series_data}
        assert "A" in series_names
        assert "B" in series_names

        # Check A has 2 points, B has 1
        for series in series_data:
            if series["name"] == "A":
                assert len(series["points"]) == 2
            elif series["name"] == "B":
                assert len(series["points"]) == 1

    def test_convert_to_series_handles_missing_values(self):
        """Test data conversion with missing field values."""
        data_points = [
            {"x": 1, "y": 2},  # Missing size
            {"x": 2, "size": 20},  # Missing y
        ]
        chart = Matrix3DChart(data_points=data_points, x_field="x", y_field="y", size_field="size")
        series_data = chart._convert_to_series()
        assert len(series_data) == 1
        # Missing values should default to 0 or 1
        points = series_data[0]["points"]
        assert len(points) == 2
        # First point: x=1, y=2, size=1 (default)
        assert points[0] == [1, 2, 1]
        # Second point: x=2, y=0 (default), size=20
        assert points[1] == [2, 0, 20]

    def test_chart_type_is_bubble(self):
        """Test that Matrix3D uses bubble chart type."""
        data_points = [{"x": 1, "y": 2, "size": 10}]
        chart = Matrix3DChart(data_points=data_points, x_field="x", y_field="y", size_field="size")
        assert chart.chart_type == XL_CHART_TYPE.BUBBLE

    def test_render_basic_matrix3d(self, slide):
        """Test rendering a basic matrix 3D chart."""
        data_points = [
            {"x": 1, "y": 2, "size": 10},
            {"x": 2, "y": 3, "size": 20},
            {"x": 3, "y": 4, "size": 30},
        ]
        chart = Matrix3DChart(
            data_points=data_points, x_field="x", y_field="y", size_field="size", title="3D Matrix"
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.chart_type == XL_CHART_TYPE.BUBBLE

    def test_render_with_color_grouping(self, slide):
        """Test rendering with color-based grouping."""
        data_points = [
            {"x": 1, "y": 2, "size": 10, "type": "Alpha"},
            {"x": 2, "y": 3, "size": 20, "type": "Beta"},
            {"x": 3, "y": 4, "size": 30, "type": "Alpha"},
        ]
        chart = Matrix3DChart(
            data_points=data_points, x_field="x", y_field="y", size_field="size", color_field="type"
        )
        result = chart.render(slide)
        assert result is not None
        # Should have 2 series (Alpha and Beta)
        assert len(result.chart.series) == 2

    def test_render_many_data_points(self, slide):
        """Test rendering with many data points."""
        data_points = [{"x": i, "y": i**2, "size": i * 5 + 1} for i in range(1, 20)]
        chart = Matrix3DChart(data_points=data_points, x_field="x", y_field="y", size_field="size")
        result = chart.render(slide)
        assert result is not None

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        data_points = [{"x": 1, "y": 2, "size": 10}]
        chart = Matrix3DChart(data_points=data_points, x_field="x", y_field="y", size_field="size")
        chart.render(slide, left=2.0, top=3.0)
        assert len(slide.shapes) > 0

    def test_multiple_groups(self, slide):
        """Test with multiple distinct groups."""
        data_points = [
            {"x": 1, "y": 1, "size": 10, "group": "G1"},
            {"x": 2, "y": 2, "size": 20, "group": "G2"},
            {"x": 3, "y": 3, "size": 30, "group": "G3"},
        ]
        chart = Matrix3DChart(
            data_points=data_points,
            x_field="x",
            y_field="y",
            size_field="size",
            color_field="group",
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 3

    def test_convert_unknown_color_field(self):
        """Test grouping with Unknown for missing color field."""
        data_points = [
            {"x": 1, "y": 1, "size": 10, "type": "A"},
            {"x": 2, "y": 2, "size": 20},  # Missing type field
        ]
        chart = Matrix3DChart(
            data_points=data_points, x_field="x", y_field="y", size_field="size", color_field="type"
        )
        series_data = chart._convert_to_series()
        series_names = {s["name"] for s in series_data}
        assert "A" in series_names
        assert "Unknown" in series_names
