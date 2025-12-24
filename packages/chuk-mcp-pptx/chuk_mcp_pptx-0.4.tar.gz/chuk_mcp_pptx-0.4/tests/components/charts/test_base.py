"""
Comprehensive tests for base chart component functionality.
"""

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

from chuk_mcp_pptx.components.charts.base import ChartComponent
from chuk_mcp_pptx.themes import ThemeManager
from chuk_mcp_pptx.layout.boundaries import validate_boundaries, adjust_to_boundaries


class _TestChartBase(ChartComponent):
    """Test implementation of ChartComponent (not a pytest test class)."""

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.categories = kwargs.get("categories", ["A", "B", "C"])
        self.values = kwargs.get("values", [1, 2, 3])

    def _prepare_chart_data(self):
        """Prepare simple test data."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories
        chart_data.add_series("Test", self.values)
        return chart_data


class TestChartComponentInitialization:
    """Test ChartComponent initialization."""

    def test_initialization_defaults(self):
        """Test chart component with default settings."""
        chart = _TestChartBase()
        assert chart.title is None
        assert chart.data is None
        assert chart.style == "default"
        assert chart.legend == "right"
        assert chart.variant_props is not None

    def test_initialization_with_title(self):
        """Test chart with title."""
        chart = _TestChartBase(title="Test Chart")
        assert chart.title == "Test Chart"

    def test_initialization_with_style(self):
        """Test chart with different styles."""
        styles = ["default", "minimal", "detailed"]
        for style in styles:
            chart = _TestChartBase(style=style)
            assert chart.style == style
            assert chart.variant_props is not None

    def test_initialization_with_legend_position(self):
        """Test chart with different legend positions."""
        positions = ["right", "bottom", "top", "none"]
        for position in positions:
            chart = _TestChartBase(legend=position)
            assert chart.legend == position

    def test_initialization_with_theme(self):
        """Test chart with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("default-light")
        chart = _TestChartBase(theme=theme)
        assert chart.theme is not None


class TestChartComponentVariants:
    """Test chart variant system."""

    def test_default_style_variant(self):
        """Test default style variant props."""
        chart = _TestChartBase(style="default")
        props = chart.variant_props
        assert props.get("show_legend") is True
        assert props.get("show_grid") is True

    def test_minimal_style_variant(self):
        """Test minimal style variant."""
        chart = _TestChartBase(style="minimal")
        props = chart.variant_props
        assert props.get("show_grid") is False

    def test_detailed_style_variant(self):
        """Test detailed style variant."""
        chart = _TestChartBase(style="detailed")
        props = chart.variant_props
        assert props.get("show_values") is True
        assert props.get("show_grid") is True

    def test_legend_variant_right(self):
        """Test legend positioned right."""
        chart = _TestChartBase(legend="right")
        props = chart.variant_props
        assert props.get("legend_position") == "right"
        assert props.get("show_legend") is True

    def test_legend_variant_none(self):
        """Test legend hidden."""
        chart = _TestChartBase(legend="none")
        props = chart.variant_props
        assert props.get("show_legend") is False


class TestChartComponentValidation:
    """Test chart data validation."""

    def test_validate_returns_true_for_base(self):
        """Test base validation returns True."""
        chart = _TestChartBase()
        is_valid, error = chart.validate()
        assert is_valid is True
        assert error is None

    def test_validate_data_returns_true_for_base(self):
        """Test base validate_data returns True."""
        chart = _TestChartBase()
        is_valid, error = chart.validate_data()
        assert is_valid is True
        assert error is None


class TestChartComponentRendering:
    """Test chart rendering."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_chart(self, slide):
        """Test rendering a basic chart."""
        chart = _TestChartBase(title="Test Chart")
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_title is True

    def test_render_with_default_position(self, slide):
        """Test chart renders with default position."""
        chart = _TestChartBase()
        result = chart.render(slide)
        assert result is not None
        # Chart should be added to slide
        assert len(slide.shapes) > 0

    def test_render_with_custom_position(self, slide):
        """Test chart with custom position."""
        chart = _TestChartBase()
        result = chart.render(slide, left=2.0, top=3.0)
        assert result is not None

    def test_render_with_custom_size(self, slide):
        """Test chart with custom width and height."""
        chart = _TestChartBase()
        result = chart.render(slide, width=6.0, height=3.0)
        assert result is not None

    def test_render_applies_title(self, slide):
        """Test chart title is applied."""
        chart = _TestChartBase(title="Revenue Chart")
        result = chart.render(slide)
        assert result.chart.has_title is True
        assert result.chart.chart_title.text_frame.text == "Revenue Chart"

    def test_render_with_different_styles(self, slide):
        """Test rendering with different style variants."""
        styles = ["default", "minimal", "detailed"]
        for style in styles:
            chart = _TestChartBase(style=style)
            result = chart.render(slide)
            assert result is not None

    def test_render_with_different_legend_positions(self, slide):
        """Test rendering with different legend positions."""
        positions = ["right", "bottom", "top"]
        for position in positions:
            chart = _TestChartBase(legend=position)
            result = chart.render(slide)
            assert result is not None
            assert result.chart.has_legend is True

    def test_render_without_legend(self, slide):
        """Test rendering without legend."""
        chart = _TestChartBase(legend="none")
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_legend is False


class TestChartComponentThemeIntegration:
    """Test chart integration with theme system."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_chart_uses_theme_colors(self, slide):
        """Test chart applies theme colors."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("default-light")
        chart = _TestChartBase(theme=theme, title="Test")
        result = chart.render(slide)
        assert result is not None

    def test_chart_with_different_themes(self, slide):
        """Test chart renders with different themes."""
        theme_manager = ThemeManager()
        themes = ["default-light", "default-dark", "ocean-light"]

        for theme_name in themes:
            theme = theme_manager.get_theme(theme_name)
            chart = _TestChartBase(theme=theme)
            result = chart.render(slide)
            assert result is not None

    def test_chart_gets_font_family_from_theme(self):
        """Test chart extracts font family from theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("default-light")
        chart = _TestChartBase(theme=theme)
        font = chart._get_font_family()
        assert font is not None
        assert isinstance(font, str)

    def test_chart_applies_theme_colors(self, slide):
        """Test apply_theme_colors method."""
        chart = _TestChartBase()
        result = chart.render(slide)
        # Should not raise errors
        chart.apply_theme_colors(result.chart)


class TestChartComponentHelpers:
    """Test chart helper methods."""

    def test_get_font_family_default(self):
        """Test font family defaults to Calibri when no theme."""
        chart = _TestChartBase()
        font = chart._get_font_family()
        assert font == "Calibri"  # Default when no theme is set

    def test_get_font_family_from_theme(self):
        """Test font family from theme."""
        theme = {"font_family": "Arial"}
        chart = _TestChartBase(theme=theme)
        font = chart._get_font_family()
        assert font == "Arial"

    def test_hex_to_rgb_conversion(self):
        """Test hex to RGB conversion."""
        chart = _TestChartBase()
        rgb = chart.hex_to_rgb("#FF0000")
        assert rgb == (255, 0, 0)

        rgb = chart.hex_to_rgb("#00FF00")
        assert rgb == (0, 255, 0)

        rgb = chart.hex_to_rgb("#0000FF")
        assert rgb == (0, 0, 255)

    def test_get_color_returns_rgb(self):
        """Test get_color returns RGBColor."""
        chart = _TestChartBase()
        color = chart.get_color("primary.DEFAULT")
        assert color is not None


class TestChartComponentBoundaries:
    """Test chart boundary validation."""

    def test_validate_boundaries_valid(self):
        """Test valid boundaries."""
        result = validate_boundaries(2.0, 2.0, 5.0, 3.0)
        assert result == (True, None)

    def test_validate_boundaries_exceeds_width(self):
        """Test exceeding width boundary."""
        result = validate_boundaries(1.0, 1.0, 15.0, 3.0)
        is_valid, error = result
        assert is_valid is False
        assert "exceeds" in error.lower()

    def test_adjust_to_boundaries_no_adjustment(self):
        """Test no adjustment needed."""
        adjusted = adjust_to_boundaries(2.0, 2.0, 5.0, 3.0)
        assert adjusted == (2.0, 2.0, 5.0, 3.0)

    def test_adjust_to_boundaries_with_adjustment(self):
        """Test adjustment when needed."""
        adjusted = adjust_to_boundaries(8.0, 2.0, 5.0, 3.0)
        left, top, width, height = adjusted
        assert left < 8.0  # Should be adjusted


class TestChartComponentDefaults:
    """Test chart default values."""

    def test_default_width(self):
        """Test default chart width."""
        chart = _TestChartBase()
        assert chart.DEFAULT_WIDTH == 8.0

    def test_default_height(self):
        """Test default chart height."""
        chart = _TestChartBase()
        assert chart.DEFAULT_HEIGHT == 3.0

    def test_default_left_position(self):
        """Test default left position."""
        chart = _TestChartBase()
        assert chart.DEFAULT_LEFT == 1.0

    def test_default_top_position(self):
        """Test default top position."""
        chart = _TestChartBase()
        assert chart.DEFAULT_TOP == 2.0

    def test_default_chart_type(self):
        """Test default chart type."""
        chart = _TestChartBase()
        assert chart.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED


class TestChartComponentComposition:
    """Test chart composition support."""

    def test_chart_extends_composable_component(self):
        """Test ChartComponent extends ComposableComponent."""
        from chuk_mcp_pptx.components.composition import ComposableComponent

        chart = _TestChartBase()
        assert isinstance(chart, ComposableComponent)

    def test_chart_has_children_list(self):
        """Test chart has _children attribute."""
        chart = _TestChartBase()
        assert hasattr(chart, "_children")
        assert isinstance(chart._children, list)

    def test_chart_can_add_children(self):
        """Test chart can add child components."""
        from chuk_mcp_pptx.components.composition import CardTitle

        chart = _TestChartBase()
        child = CardTitle("Test")
        chart.add_child(child)
        assert len(chart._children) == 1

    def test_chart_get_children(self):
        """Test getting chart children."""
        from chuk_mcp_pptx.components.composition import CardTitle

        chart = _TestChartBase()
        child = CardTitle("Test")
        chart.add_child(child)
        children = chart.get_children()
        assert len(children) == 1
        assert children[0] == child

    def test_chart_clear_children(self):
        """Test clearing chart children."""
        from chuk_mcp_pptx.components.composition import CardTitle

        chart = _TestChartBase()
        chart.add_child(CardTitle("Test"))
        chart.clear_children()
        assert len(chart._children) == 0


class TestChartComponentEdgeCases:
    """Test edge cases and error handling."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a test slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_chart_with_none_title(self, slide):
        """Test chart with None title."""
        chart = _TestChartBase(title=None)
        result = chart.render(slide)
        assert result is not None

    def test_chart_with_empty_title(self, slide):
        """Test chart with empty string title."""
        chart = _TestChartBase(title="")
        result = chart.render(slide)
        assert result is not None

    def test_chart_with_very_long_title(self, slide):
        """Test chart with very long title."""
        long_title = "This is a very long chart title that might cause layout issues"
        chart = _TestChartBase(title=long_title)
        result = chart.render(slide)
        assert result is not None
        assert result.chart.chart_title.text_frame.text == long_title

    def test_chart_with_special_characters_in_title(self, slide):
        """Test chart with special characters in title."""
        special = "Test™ © ® € £ ¥"
        chart = _TestChartBase(title=special)
        result = chart.render(slide)
        assert result is not None

    def test_chart_at_slide_boundaries(self, slide):
        """Test chart positioned at slide edges."""
        chart = _TestChartBase()
        # Near right edge
        result = chart.render(slide, left=8.0, top=1.0, width=2.0)
        assert result is not None

    def test_chart_with_minimal_dimensions(self, slide):
        """Test chart with very small dimensions."""
        chart = _TestChartBase()
        result = chart.render(slide, width=1.0, height=1.0)
        assert result is not None

    def test_chart_with_zero_position(self, slide):
        """Test chart at position (0, 0)."""
        chart = _TestChartBase()
        result = chart.render(slide, left=0.0, top=0.0)
        assert result is not None


class TestChartComponentTokens:
    """Test chart token access."""

    def test_chart_has_tokens(self):
        """Test chart has tokens attribute."""
        chart = _TestChartBase()
        assert hasattr(chart, "tokens")
        assert chart.tokens is not None

    def test_chart_tokens_have_colors(self):
        """Test chart tokens include colors."""
        chart = _TestChartBase()
        assert "chart" in chart.tokens or len(chart.tokens) > 0

    def test_chart_can_get_semantic_color(self):
        """Test chart can access semantic colors."""
        chart = _TestChartBase()
        color = chart.get_color("primary.DEFAULT")
        assert color is not None
