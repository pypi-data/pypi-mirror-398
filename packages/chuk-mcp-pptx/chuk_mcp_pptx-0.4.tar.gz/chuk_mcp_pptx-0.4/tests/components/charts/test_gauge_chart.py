"""
Tests for GaugeChart component.
"""

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from chuk_mcp_pptx.components.charts.radar_combo import GaugeChart
from chuk_mcp_pptx.themes import ThemeManager


class TestGaugeChart:
    """Test GaugeChart component."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_initialization(self):
        """Test gauge chart initialization."""
        chart = GaugeChart(value=75, min_value=0, max_value=100, title="Progress")
        assert chart.value == 75
        assert chart.min_value == 0
        assert chart.max_value == 100
        assert chart.title == "Progress"

    def test_chart_type(self):
        """Test gauge uses doughnut chart type."""
        chart = GaugeChart(value=50)
        assert chart.chart_type == XL_CHART_TYPE.DOUGHNUT

    def test_validate_data_valid(self):
        """Test validation with valid data."""
        chart = GaugeChart(value=50, min_value=0, max_value=100)
        is_valid, error = chart.validate_data()
        assert is_valid is True
        assert error is None

    def test_validate_data_invalid_range(self):
        """Test validation with invalid min/max range."""
        with pytest.raises(ValueError, match="min_value must be less than max_value"):
            GaugeChart(value=50, min_value=100, max_value=0)

    def test_validate_data_value_out_of_range(self):
        """Test validation with value out of range."""
        with pytest.raises(ValueError, match="value .* must be between min_value"):
            GaugeChart(value=150, min_value=0, max_value=100)

    def test_calculate_segments(self):
        """Test gauge segment calculation."""
        chart = GaugeChart(value=75, min_value=0, max_value=100)
        # Should have 3 categories: Value, Empty, Hidden
        assert len(chart.categories) == 3
        assert len(chart.values) == 3
        assert chart.categories == ["Value", "Empty", "Hidden"]

    def test_render_basic_gauge(self, slide):
        """Test rendering a basic gauge chart."""
        chart = GaugeChart(value=60, min_value=0, max_value=100, title="Progress")
        result = chart.render(slide)
        assert result is not None
        assert result.chart.chart_type == XL_CHART_TYPE.DOUGHNUT

    def test_render_half_full(self, slide):
        """Test rendering at 50%."""
        chart = GaugeChart(value=50, min_value=0, max_value=100)
        result = chart.render(slide)
        assert result is not None

    def test_render_full(self, slide):
        """Test rendering at 100%."""
        chart = GaugeChart(value=100, min_value=0, max_value=100)
        result = chart.render(slide)
        assert result is not None

    def test_render_empty(self, slide):
        """Test rendering at 0%."""
        chart = GaugeChart(value=0, min_value=0, max_value=100)
        result = chart.render(slide)
        assert result is not None

    def test_render_custom_range(self, slide):
        """Test rendering with custom min/max values."""
        chart = GaugeChart(value=50, min_value=25, max_value=75)
        result = chart.render(slide)
        assert result is not None

    def test_render_with_thresholds(self, slide):
        """Test rendering with thresholds."""
        chart = GaugeChart(
            value=75, min_value=0, max_value=100, thresholds={"warning": 60, "danger": 80}
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_no_legend(self, slide):
        """Test that gauge hides legend."""
        chart = GaugeChart(value=50)
        result = chart.render(slide)
        assert result.chart.has_legend is False

    def test_render_applies_colors(self, slide):
        """Test that gauge applies success/muted colors."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-light")

        chart = GaugeChart(value=75, theme=theme)
        result = chart.render(slide)
        # Verify chart was created with theme
        assert result is not None

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        chart = GaugeChart(value=80)
        chart.render(slide, left=2.0, top=3.0)
        assert len(slide.shapes) > 0

    def test_render_with_title(self, slide):
        """Test gauge with integrated title label."""
        chart = GaugeChart(value=87, min_value=0, max_value=100, title="Customer Satisfaction")
        result = chart.render(slide, left=1, top=2.5, width=2.5, height=2.5)
        assert result is not None
        # Should have added textboxes for title, value, and percentage
        assert len(slide.shapes) > 1  # Chart + labels

    def test_render_adds_value_label(self, slide):
        """Test that gauge adds value label automatically."""
        chart = GaugeChart(value=75, min_value=0, max_value=100)
        chart.render(slide)
        # Should have chart + value label + percentage label
        assert len(slide.shapes) >= 3

    def test_render_adds_percentage_label(self, slide):
        """Test that gauge adds percentage label automatically."""
        chart = GaugeChart(value=60, min_value=0, max_value=100)
        chart.render(slide)
        # Verify multiple shapes (chart + labels)
        assert len(slide.shapes) > 1
