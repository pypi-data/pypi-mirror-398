"""
Comprehensive rendering tests for Scatter and Bubble charts.
Tests actual PowerPoint rendering to boost coverage.
"""

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from chuk_mcp_pptx.components.charts.scatter_bubble import ScatterChart, BubbleChart
from chuk_mcp_pptx.themes import ThemeManager


class TestScatterChartRendering:
    """Test ScatterChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_scatter(self, slide):
        """Test rendering a basic scatter chart."""
        chart = ScatterChart(
            series_data=[
                {"name": "Series 1", "x_values": [1, 2, 3, 4, 5], "y_values": [2, 4, 3, 5, 4]}
            ],
            title="Scatter Plot",
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_title

    def test_render_default_variant(self, slide):
        """Test default variant rendering."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], variant="default"
        )
        result = chart.render(slide)
        # Default variant is XY_SCATTER
        assert result.chart.chart_type in [XL_CHART_TYPE.XY_SCATTER, XL_CHART_TYPE.XY_SCATTER_LINES]

    def test_render_smooth_variant(self, slide):
        """Test smooth variant rendering."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 4, 2]}], variant="smooth"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.XY_SCATTER_SMOOTH

    def test_render_smooth_markers_variant(self, slide):
        """Test smooth markers variant rendering."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], variant="smooth_markers"
        )
        result = chart.render(slide)
        # Smooth markers variant
        assert result.chart.chart_type in [
            XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS,
            XL_CHART_TYPE.XY_SCATTER_SMOOTH,
        ]

    def test_render_lines_variant(self, slide):
        """Test lines variant rendering."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], variant="lines"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.XY_SCATTER_LINES

    def test_render_lines_markers_variant(self, slide):
        """Test lines markers variant rendering."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], variant="lines_markers"
        )
        result = chart.render(slide)
        # Lines markers variant
        assert result.chart.chart_type in [
            XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS,
            XL_CHART_TYPE.XY_SCATTER_LINES,
        ]

    def test_render_multiple_series(self, slide):
        """Test rendering with multiple data series."""
        chart = ScatterChart(
            series_data=[
                {"name": "Group A", "x_values": [1, 2, 3], "y_values": [2, 4, 3]},
                {"name": "Group B", "x_values": [1, 2, 3], "y_values": [3, 1, 4]},
            ]
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_with_trendline(self, slide):
        """Test rendering with trendline."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3, 4, 5], "y_values": [2, 4, 3, 5, 6]}],
            show_trendline=True,
        )
        chart.render(slide)
        assert chart.show_trendline is True

    def test_render_without_trendline(self, slide):
        """Test rendering without trendline."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], show_trendline=False
        )
        chart.render(slide)
        assert chart.show_trendline is False

    def test_render_custom_marker_size(self, slide):
        """Test rendering with custom marker size."""
        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], marker_size=12
        )
        chart.render(slide)
        assert chart.marker_size == 12

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        chart = ScatterChart(series_data=[{"x_values": [1, 2], "y_values": [1, 2]}])
        chart.render(slide, left=2.0, top=3.0)
        assert len(slide.shapes) > 0

    def test_render_custom_size(self, slide):
        """Test rendering with custom size."""
        chart = ScatterChart(series_data=[{"x_values": [1, 2], "y_values": [1, 2]}])
        result = chart.render(slide, width=6.0, height=6.0)
        assert result is not None

    def test_render_with_theme(self, slide):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-light")

        chart = ScatterChart(
            series_data=[{"x_values": [1, 2, 3], "y_values": [1, 2, 3]}], theme=theme
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_without_series_names(self, slide):
        """Test rendering without series names."""
        chart = ScatterChart(series_data=[{"x_values": [1, 2, 3], "y_values": [2, 3, 4]}])
        result = chart.render(slide)
        assert result is not None

    def test_render_with_negative_values(self, slide):
        """Test scatter with negative values."""
        chart = ScatterChart(series_data=[{"x_values": [-1, 0, 1], "y_values": [-2, 0, 2]}])
        result = chart.render(slide)
        assert result is not None


class TestBubbleChartRendering:
    """Test BubbleChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_bubble(self, slide):
        """Test rendering a basic bubble chart."""
        chart = BubbleChart(
            series_data=[{"name": "Series 1", "points": [[1, 2, 10], [2, 3, 20], [3, 4, 30]]}],
            title="Bubble Chart",
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_title

    def test_render_multiple_bubble_series(self, slide):
        """Test rendering with multiple bubble series."""
        chart = BubbleChart(
            series_data=[
                {"name": "Group A", "points": [[1, 1, 10], [2, 2, 20]]},
                {"name": "Group B", "points": [[2, 3, 15], [3, 4, 25]]},
            ]
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_3d_bubble(self, slide):
        """Test 3D bubble - uses standard bubble chart."""
        chart = BubbleChart(series_data=[{"points": [[1, 1, 10], [2, 2, 20], [3, 3, 30]]}])
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.BUBBLE

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        chart = BubbleChart(series_data=[{"points": [[1, 1, 10], [2, 2, 20]]}])
        chart.render(slide, left=1.5, top=2.5)
        assert len(slide.shapes) > 0

    def test_render_with_theme(self, slide):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-dark")

        chart = BubbleChart(
            series_data=[{"points": [[1, 1, 10], [2, 2, 20], [3, 3, 30]]}], theme=theme
        )
        result = chart.render(slide)
        assert result is not None


class TestScatterChartValidation:
    """Test ScatterChart validation."""

    def test_validate_empty_series_data(self):
        """Test validation with empty series data."""
        with pytest.raises(ValueError, match="No series data provided"):
            ScatterChart(series_data=[])

    def test_validate_missing_x_values(self):
        """Test validation with missing x_values."""
        with pytest.raises(ValueError, match="Series 0 missing x_values"):
            ScatterChart(series_data=[{"y_values": [1, 2, 3]}])

    def test_validate_missing_y_values(self):
        """Test validation with missing y_values."""
        with pytest.raises(ValueError, match="Series 0 missing y_values"):
            ScatterChart(series_data=[{"x_values": [1, 2, 3]}])

    def test_validate_mismatched_lengths(self):
        """Test validation with mismatched x/y lengths."""
        with pytest.raises(ValueError, match="x_values.*and y_values.*must have same length"):
            ScatterChart(series_data=[{"x_values": [1, 2], "y_values": [1, 2, 3]}])

    def test_validate_empty_data_points(self):
        """Test validation with empty data points."""
        with pytest.raises(ValueError, match="Series 0 has no data points"):
            ScatterChart(series_data=[{"x_values": [], "y_values": []}])


class TestBubbleChartValidation:
    """Test BubbleChart validation."""

    def test_validate_empty_series_data(self):
        """Test validation with empty series data."""
        with pytest.raises(ValueError, match="No series data provided"):
            BubbleChart(series_data=[])

    def test_validate_missing_points(self):
        """Test validation with missing points."""
        with pytest.raises(ValueError, match="Series 0 missing points data"):
            BubbleChart(series_data=[{"name": "Test"}])

    def test_validate_empty_points(self):
        """Test validation with empty points list."""
        with pytest.raises(ValueError, match="Series 0 has no data points"):
            BubbleChart(series_data=[{"points": []}])

    def test_validate_invalid_point_format(self):
        """Test validation with invalid point format."""
        with pytest.raises(ValueError, match="must be \\[x, y, size\\] format"):
            BubbleChart(series_data=[{"points": [[1, 2]]}])  # Only 2 values

    def test_validate_non_numeric_values(self):
        """Test validation with non-numeric values."""
        with pytest.raises(ValueError, match="x, y, size must be numeric"):
            BubbleChart(series_data=[{"points": [["a", 2, 10]]}])

    def test_validate_negative_size(self):
        """Test validation with negative size."""
        with pytest.raises(ValueError, match="size must be positive"):
            BubbleChart(series_data=[{"points": [[1, 2, -10]]}])

    def test_validate_zero_size(self):
        """Test validation with zero size."""
        with pytest.raises(ValueError, match="size must be positive"):
            BubbleChart(series_data=[{"points": [[1, 2, 0]]}])


class TestScatterBubbleEdgeCases:
    """Test edge cases for scatter and bubble charts."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_scatter_single_point(self, slide):
        """Test scatter with single data point."""
        chart = ScatterChart(series_data=[{"x_values": [1], "y_values": [1]}])
        result = chart.render(slide)
        assert result is not None

    def test_scatter_many_points(self, slide):
        """Test scatter with many data points."""
        x_values = list(range(100))
        y_values = [i**2 for i in range(100)]

        chart = ScatterChart(series_data=[{"x_values": x_values, "y_values": y_values}])
        result = chart.render(slide)
        assert result is not None

    def test_bubble_small_size(self, slide):
        """Test bubble with small size."""
        chart = BubbleChart(series_data=[{"points": [[1, 1, 1], [2, 2, 10]]}])
        result = chart.render(slide)
        assert result is not None

    def test_scatter_decimal_values(self, slide):
        """Test scatter with decimal values."""
        chart = ScatterChart(
            series_data=[{"x_values": [1.5, 2.7, 3.2], "y_values": [1.1, 2.3, 3.9]}]
        )
        result = chart.render(slide)
        assert result is not None
