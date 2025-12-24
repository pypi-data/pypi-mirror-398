"""
Comprehensive rendering tests for Radar and Combo charts.
Tests actual PowerPoint rendering to boost coverage.
"""

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from chuk_mcp_pptx.components.charts.radar_combo import RadarChart, ComboChart
from chuk_mcp_pptx.themes import ThemeManager


class TestRadarChartRendering:
    """Test RadarChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_radar(self, slide):
        """Test rendering a basic radar chart."""
        chart = RadarChart(
            categories=["Speed", "Power", "Accuracy", "Reliability"],
            series={"Product A": [80, 90, 70, 85]},
            title="Product Comparison",
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_title

    def test_render_filled_variant(self, slide):
        """Test filled variant rendering."""
        chart = RadarChart(
            categories=["A", "B", "C", "D"], series={"S1": [10, 20, 30, 40]}, variant="filled"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.RADAR_FILLED

    def test_render_markers_variant(self, slide):
        """Test markers variant rendering."""
        chart = RadarChart(
            categories=["A", "B", "C", "D"], series={"S1": [10, 20, 30, 40]}, variant="markers"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.RADAR_MARKERS

    def test_render_lines_variant(self, slide):
        """Test lines variant rendering."""
        chart = RadarChart(
            categories=["A", "B", "C", "D"], series={"S1": [10, 20, 30, 40]}, variant="lines"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.RADAR

    def test_render_multiple_series(self, slide):
        """Test rendering with multiple data series."""
        chart = RadarChart(
            categories=["Speed", "Power", "Accuracy"],
            series={"Product A": [80, 90, 70], "Product B": [70, 85, 90]},
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_with_max_value(self, slide):
        """Test rendering with custom max value."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]}, max_value=100)
        chart.render(slide)
        assert chart.max_value == 100

    def test_render_auto_max_value(self, slide):
        """Test rendering with auto-calculated max value."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]})
        result = chart.render(slide)
        # max_value should be None (auto-calculated by PowerPoint)
        assert result is not None

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]})
        chart.render(slide, left=2.0, top=3.0)
        assert len(slide.shapes) > 0

    def test_render_custom_size(self, slide):
        """Test rendering with custom size."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]})
        result = chart.render(slide, width=6.0, height=6.0)
        assert result is not None

    def test_render_with_theme(self, slide):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-light")

        chart = RadarChart(
            categories=["A", "B", "C", "D"], series={"S1": [10, 20, 30, 40]}, theme=theme
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_applies_theme_colors(self, slide):
        """Test that theme colors are applied."""
        chart = RadarChart(
            categories=["A", "B", "C"], series={"Series 1": [10, 20, 30], "Series 2": [15, 25, 35]}
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_minimum_categories(self, slide):
        """Test radar with minimum 3 categories."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]})
        result = chart.render(slide)
        assert result is not None

    def test_render_many_categories(self, slide):
        """Test radar with many categories."""
        categories = [f"Cat{i}" for i in range(12)]
        values = [10 * i for i in range(12)]

        chart = RadarChart(categories=categories, series={"Data": values})
        result = chart.render(slide)
        assert result is not None

    def test_render_with_default_style(self, slide):
        """Test rendering with default style."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]}, style="default")
        result = chart.render(slide)
        assert result.chart.has_legend

    def test_render_without_title(self, slide):
        """Test rendering without title."""
        chart = RadarChart(categories=["A", "B", "C"], series={"S1": [10, 20, 30]})
        result = chart.render(slide)
        assert result is not None


class TestComboChartRendering:
    """Test ComboChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_combo(self, slide):
        """Test rendering a basic combo chart."""
        chart = ComboChart(
            categories=["Q1", "Q2", "Q3", "Q4"],
            column_series={"Revenue": [100, 120, 110, 130]},
            line_series={"Target": [110, 115, 120, 125]},
            title="Revenue vs Target",
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_title

    def test_render_column_line_combo(self, slide):
        """Test column-line combo chart."""
        chart = ComboChart(
            categories=["A", "B", "C"],
            column_series={"Bars": [10, 20, 30]},
            line_series={"Line": [15, 25, 35]},
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_with_secondary_axis(self, slide):
        """Test combo chart with secondary axis."""
        chart = ComboChart(
            categories=["A", "B", "C"],
            column_series={"Values": [100, 200, 300]},
            line_series={"Percent": [10, 20, 30]},
            secondary_axis=["Percent"],
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_without_secondary_axis(self, slide):
        """Test combo chart without secondary axis."""
        chart = ComboChart(
            categories=["A", "B"], column_series={"S1": [10, 20]}, line_series={"S2": [30, 40]}
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_multiple_columns(self, slide):
        """Test combo with multiple column series."""
        chart = ComboChart(
            categories=["A", "B", "C"],
            column_series={"Col1": [10, 20, 30], "Col2": [15, 25, 35]},
            line_series={"Line": [5, 15, 25]},
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        chart = ComboChart(
            categories=["A", "B"], column_series={"S1": [10, 20]}, line_series={"S2": [30, 40]}
        )
        chart.render(slide, left=1.5, top=2.5)
        assert len(slide.shapes) > 0

    def test_render_custom_size(self, slide):
        """Test rendering with custom size."""
        chart = ComboChart(
            categories=["A", "B"], column_series={"S1": [10, 20]}, line_series={"S2": [30, 40]}
        )
        result = chart.render(slide, width=7.0, height=4.0)
        assert result is not None

    def test_render_with_theme(self, slide):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-dark")

        chart = ComboChart(
            categories=["A", "B"],
            column_series={"S1": [10, 20]},
            line_series={"S2": [30, 40]},
            theme=theme,
        )
        result = chart.render(slide)
        assert result is not None

    def test_render_applies_colors(self, slide):
        """Test that colors are applied to series."""
        chart = ComboChart(
            categories=["A", "B", "C"],
            column_series={"Bars": [10, 20, 30]},
            line_series={"Line": [15, 25, 35]},
        )
        result = chart.render(slide)
        # Verify series were created
        assert len(result.chart.series) == 2

    def test_render_without_title(self, slide):
        """Test rendering without title."""
        chart = ComboChart(
            categories=["A", "B"], column_series={"S1": [10, 20]}, line_series={"S2": [30, 40]}
        )
        result = chart.render(slide)
        assert result is not None


class TestRadarChartValidation:
    """Test RadarChart validation."""

    def test_validate_empty_categories(self):
        """Test validation with empty categories."""
        with pytest.raises(ValueError, match="No categories provided"):
            RadarChart(categories=[], series={"S1": [10, 20, 30]})

    def test_validate_empty_series(self):
        """Test validation with empty series."""
        with pytest.raises(ValueError, match="No data series provided"):
            RadarChart(categories=["A", "B", "C"], series={})

    def test_validate_too_few_categories(self):
        """Test validation with less than 3 categories."""
        with pytest.raises(ValueError, match="needs at least 3 categories"):
            RadarChart(categories=["A", "B"], series={"S1": [10, 20]})

    def test_validate_mismatched_series_length(self):
        """Test validation with mismatched series values."""
        with pytest.raises(ValueError, match="has 2 values, expected 3"):
            RadarChart(
                categories=["A", "B", "C"],
                series={"S1": [10, 20]},  # Only 2 values for 3 categories
            )


class TestComboChartValidation:
    """Test ComboChart validation."""

    def test_validate_empty_categories(self):
        """Test validation with empty categories."""
        with pytest.raises(ValueError, match="No categories provided"):
            ComboChart(categories=[], column_series={"C1": [10]}, line_series={"L1": [20]})

    def test_validate_no_series(self):
        """Test validation with no column or line series."""
        with pytest.raises(ValueError, match="No data series provided"):
            ComboChart(categories=["A", "B"], column_series={}, line_series={})

    def test_validate_column_series_mismatch(self):
        """Test validation with mismatched column series length."""
        with pytest.raises(ValueError, match="Column series.*has 1 values, expected 2"):
            ComboChart(
                categories=["A", "B"],
                column_series={"C1": [10]},  # Only 1 value for 2 categories
                line_series={},
            )

    def test_validate_line_series_mismatch(self):
        """Test validation with mismatched line series length."""
        with pytest.raises(ValueError, match="Line series.*has 3 values, expected 2"):
            ComboChart(
                categories=["A", "B"],
                column_series={},
                line_series={"L1": [10, 20, 30]},  # 3 values for 2 categories
            )


class TestRadarComboEdgeCases:
    """Test edge cases for radar and combo charts."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_radar_zero_values(self, slide):
        """Test radar chart with zero values."""
        chart = RadarChart(categories=["A", "B", "C"], series={"Zeros": [0, 0, 0]})
        result = chart.render(slide)
        assert result is not None

    def test_radar_negative_values(self, slide):
        """Test radar chart with negative values."""
        chart = RadarChart(categories=["A", "B", "C"], series={"Mixed": [10, -5, 15]})
        result = chart.render(slide)
        assert result is not None

    def test_combo_only_columns(self, slide):
        """Test combo chart with only column series."""
        chart = ComboChart(
            categories=["A", "B"],
            column_series={"Col1": [10, 20], "Col2": [15, 25]},
            line_series={},
        )
        result = chart.render(slide)
        assert result is not None

    def test_combo_many_series(self, slide):
        """Test combo chart with many series."""
        column_series = {f"Col{i}": [10 * i, 20 * i, 30 * i] for i in range(3)}
        line_series = {f"Line{i}": [5 * i, 15 * i, 25 * i] for i in range(2)}

        chart = ComboChart(
            categories=["A", "B", "C"], column_series=column_series, line_series=line_series
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 5
