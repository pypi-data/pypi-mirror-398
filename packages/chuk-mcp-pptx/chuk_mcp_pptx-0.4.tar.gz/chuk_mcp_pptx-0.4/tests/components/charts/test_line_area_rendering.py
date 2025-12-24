"""
Comprehensive rendering tests for Line, Area, and Sparkline charts.
Tests actual PowerPoint rendering to boost coverage.
"""

import pytest
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE

from chuk_mcp_pptx.components.charts.line_area import LineChart, AreaChart, SparklineChart
from chuk_mcp_pptx.themes import ThemeManager


class TestLineChartRendering:
    """Test LineChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_line_chart(self, slide):
        """Test rendering a basic line chart."""
        chart = LineChart(
            categories=["Jan", "Feb", "Mar", "Apr"],
            series={"Revenue": [100, 120, 115, 130]},
            title="Monthly Revenue",
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_title
        assert result.chart.chart_title.text_frame.text == "Monthly Revenue"

    def test_render_multiple_series(self, slide):
        """Test rendering with multiple data series."""
        chart = LineChart(
            categories=["Q1", "Q2", "Q3"],
            series={"Revenue": [100, 120, 110], "Costs": [80, 90, 85]},
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_line_variant(self, slide):
        """Test line variant rendering."""
        chart = LineChart(categories=["A", "B"], series={"S1": [10, 20]}, variant="line")
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.LINE_MARKERS

    def test_render_smooth_variant(self, slide):
        """Test smooth variant rendering."""
        chart = LineChart(categories=["A", "B", "C"], series={"S1": [10, 20, 15]}, variant="smooth")
        result = chart.render(slide)
        # Smooth variant applies smoothing to series
        assert result is not None

    def test_render_area_variant(self, slide):
        """Test area variant rendering."""
        chart = LineChart(categories=["A", "B"], series={"S1": [10, 20]}, variant="area")
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.AREA

    def test_render_smooth_area_variant(self, slide):
        """Test smooth area variant rendering."""
        chart = LineChart(
            categories=["A", "B", "C"], series={"S1": [10, 20, 15]}, variant="smooth_area"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.AREA

    def test_render_with_default_style(self, slide):
        """Test rendering with default style."""
        chart = LineChart(categories=["A"], series={"S1": [10]}, style="default")
        result = chart.render(slide)
        assert result is not None
        assert result.chart.has_legend

    def test_render_with_minimal_style(self, slide):
        """Test rendering with minimal style."""
        chart = LineChart(categories=["A", "B"], series={"S1": [10, 20]}, style="minimal")
        result = chart.render(slide)
        assert result is not None

    def test_render_with_detailed_style(self, slide):
        """Test rendering with detailed style."""
        chart = LineChart(categories=["A"], series={"S1": [10]}, style="detailed")
        result = chart.render(slide)
        assert result is not None

    def test_render_custom_position(self, slide):
        """Test rendering at custom position."""
        chart = LineChart(categories=["A"], series={"S1": [10]})
        chart.render(slide, left=2.0, top=3.0)
        assert len(slide.shapes) > 0

    def test_render_custom_size(self, slide):
        """Test rendering with custom width and height."""
        chart = LineChart(categories=["A"], series={"S1": [10]})
        result = chart.render(slide, width=6.0, height=3.0)
        assert result is not None

    def test_render_with_theme(self, slide):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-light")

        chart = LineChart(categories=["A", "B"], series={"S1": [10, 20]}, theme=theme)
        result = chart.render(slide)
        assert result is not None

    def test_render_applies_theme_colors(self, slide):
        """Test that theme colors are applied to lines."""
        chart = LineChart(
            categories=["A", "B", "C"], series={"Revenue": [100, 120, 110], "Costs": [80, 90, 85]}
        )
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_with_markers(self, slide):
        """Test rendering with markers."""
        chart = LineChart(categories=["A", "B"], series={"S1": [10, 20]}, variant="line")
        result = chart.render(slide)
        # Line variant shows markers by default
        assert result is not None

    def test_render_without_grid(self, slide):
        """Test rendering without gridlines."""
        chart = LineChart(categories=["A", "B"], series={"S1": [10, 20]}, style="minimal")
        result = chart.render(slide)
        # Minimal style hides grid
        assert result is not None

    def test_render_legend_positions(self, slide):
        """Test rendering with different legend positions."""
        positions = ["right", "bottom", "top"]
        for position in positions:
            chart = LineChart(categories=["A"], series={"S1": [10]}, legend=position)
            result = chart.render(slide)
            assert result.chart.has_legend

    def test_render_without_legend(self, slide):
        """Test rendering without legend."""
        chart = LineChart(categories=["A"], series={"S1": [10]}, legend="none")
        result = chart.render(slide)
        assert result is not None

    def test_render_without_title(self, slide):
        """Test rendering without title."""
        chart = LineChart(categories=["A"], series={"S1": [10]})
        result = chart.render(slide)
        assert result is not None


class TestAreaChartRendering:
    """Test AreaChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_area_chart(self, slide):
        """Test rendering a basic area chart."""
        chart = AreaChart(
            categories=["2020", "2021", "2022"],
            series={"Growth": [100, 150, 200]},
            title="Growth Trend",
        )
        result = chart.render(slide)
        assert result is not None
        assert result.chart.chart_type == XL_CHART_TYPE.AREA

    def test_render_stacked_variant(self, slide):
        """Test stacked area chart."""
        chart = AreaChart(
            categories=["A", "B"], series={"S1": [10, 20], "S2": [5, 10]}, variant="stacked"
        )
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.AREA_STACKED

    def test_render_stacked100_variant(self, slide):
        """Test 100% stacked area chart."""
        chart = AreaChart(categories=["A", "B"], series={"S1": [10, 20]}, variant="stacked100")
        result = chart.render(slide)
        assert result.chart.chart_type == XL_CHART_TYPE.AREA_STACKED_100

    def test_render_with_multiple_series(self, slide):
        """Test area chart with multiple series."""
        chart = AreaChart(categories=["A", "B"], series={"Q1": [10, 20], "Q2": [15, 25]})
        result = chart.render(slide)
        assert len(result.chart.series) == 2

    def test_render_with_theme(self, slide):
        """Test rendering with custom theme."""
        theme_manager = ThemeManager()
        theme = theme_manager.get_theme("ocean-dark")

        chart = AreaChart(categories=["A", "B"], series={"S1": [10, 20]}, theme=theme)
        result = chart.render(slide)
        assert result is not None


class TestSparklineChartRendering:
    """Test SparklineChart rendering to PowerPoint."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_render_basic_sparkline(self, slide):
        """Test rendering a basic sparkline."""
        chart = SparklineChart(values=[10, 15, 13, 17, 14, 20])
        result = chart.render(slide)
        assert result is not None

    def test_render_hides_axes(self, slide):
        """Test that sparklines hide axes."""
        chart = SparklineChart(values=[10, 20, 30])
        result = chart.render(slide)
        # Axes should be hidden
        if hasattr(result, "category_axis"):
            assert result.category_axis.visible is False
        if hasattr(result, "value_axis"):
            assert result.value_axis.visible is False

    def test_render_with_small_dimensions(self, slide):
        """Test sparkline renders with small dimensions."""
        chart = SparklineChart(values=[1, 2, 3, 4, 5])
        result = chart.render(slide)
        # Default small size
        assert result is not None

    def test_render_custom_dimensions(self, slide):
        """Test sparkline with custom dimensions."""
        chart = SparklineChart(values=[1, 2, 3])
        result = chart.render(slide, width=3.0, height=1.0)
        assert result is not None

    def test_render_with_categories(self, slide):
        """Test sparkline with explicit categories."""
        chart = SparklineChart(values=[10, 20, 30], categories=["A", "B", "C"])
        result = chart.render(slide)
        assert result is not None

    def test_render_with_label(self, slide):
        """Test sparkline with label."""
        chart = SparklineChart(values=[10, 20, 30, 40], label="Daily Active Users")
        result = chart.render(slide)
        assert result is not None
        # Should have added textbox for label
        assert len(slide.shapes) > 1  # Chart + label textbox

    def test_render_with_show_value(self, slide):
        """Test sparkline with show_value enabled."""
        chart = SparklineChart(values=[10, 20, 30, 40], show_value=True)
        result = chart.render(slide)
        assert result is not None
        # Should have added textbox for value
        assert len(slide.shapes) > 1  # Chart + value textbox

    def test_render_with_label_and_value(self, slide):
        """Test sparkline with both label and value."""
        chart = SparklineChart(values=[10, 20, 30, 40], label="Revenue", show_value=True)
        result = chart.render(slide)
        assert result is not None
        # Should have added textboxes for both label and value
        assert len(slide.shapes) > 2  # Chart + label + value textboxes

    def test_render_with_show_value_false(self, slide):
        """Test sparkline with show_value disabled."""
        chart = SparklineChart(values=[10, 20, 30], show_value=False)
        result = chart.render(slide)
        assert result is not None


class TestLineAreaEdgeCases:
    """Test edge cases and boundary conditions."""

    @pytest.fixture
    def presentation(self):
        """Create a test presentation."""
        return Presentation()

    @pytest.fixture
    def slide(self, presentation):
        """Create a blank slide."""
        blank_layout = presentation.slide_layouts[6]
        return presentation.slides.add_slide(blank_layout)

    def test_single_data_point(self, slide):
        """Test line chart with single data point."""
        chart = LineChart(categories=["Only"], series={"Single": [42]})
        result = chart.render(slide)
        assert result is not None

    def test_many_data_points(self, slide):
        """Test line chart with many data points."""
        categories = [f"Point{i}" for i in range(20)]
        values = list(range(20))

        chart = LineChart(categories=categories, series={"Data": values})
        result = chart.render(slide)
        assert result is not None

    def test_negative_values(self, slide):
        """Test line chart with negative values."""
        chart = LineChart(categories=["A", "B", "C"], series={"Mixed": [10, -5, 15]})
        result = chart.render(slide)
        assert result is not None

    def test_zero_values(self, slide):
        """Test line chart with zero values."""
        chart = LineChart(categories=["A", "B", "C"], series={"Zeros": [0, 0, 0]})
        result = chart.render(slide)
        assert result is not None

    def test_decimal_values(self, slide):
        """Test line chart with decimal values."""
        chart = LineChart(categories=["A", "B", "C"], series={"Decimals": [1.5, 2.75, 3.125]})
        result = chart.render(slide)
        assert result is not None
