"""
Tests for Table component.
"""

import pytest
from pptx import Presentation
from pptx.util import Pt

from chuk_mcp_pptx.components.core import Table
from chuk_mcp_pptx.themes.theme_manager import ThemeManager


@pytest.fixture
def presentation():
    """Create a test presentation."""
    return Presentation()


@pytest.fixture
def slide(presentation):
    """Create a test slide."""
    blank_layout = presentation.slide_layouts[6]  # Blank layout
    return presentation.slides.add_slide(blank_layout)


@pytest.fixture
def theme():
    """Create a test theme."""
    theme_manager = ThemeManager()
    return theme_manager.get_theme("dark-violet")


@pytest.fixture
def sample_data():
    """Sample table data."""
    return {
        "headers": ["Product", "Q1", "Q2", "Q3"],
        "data": [
            ["Laptops", "$100K", "$120K", "$110K"],
            ["Phones", "$80K", "$90K", "$95K"],
            ["Tablets", "$60K", "$65K", "$70K"],
        ],
    }


class TestTableInit:
    """Tests for Table initialization."""

    def test_init_basic(self, sample_data):
        """Test basic table initialization."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"])
        assert table.headers == sample_data["headers"]
        assert table.data == sample_data["data"]
        assert table.variant == "default"
        assert table.size == "md"

    def test_init_with_variant(self, sample_data):
        """Test initialization with variant."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], variant="bordered")
        assert table.variant == "bordered"

    def test_init_with_size(self, sample_data):
        """Test initialization with size."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], size="lg")
        assert table.size == "lg"

    def test_init_with_theme(self, sample_data, theme):
        """Test initialization with theme."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        assert table.theme is not None


class TestTableVariants:
    """Tests for Table variants."""

    def test_default_variant(self, slide, sample_data, theme):
        """Test default variant."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], variant="default", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None

    def test_bordered_variant(self, slide, sample_data, theme):
        """Test bordered variant."""
        table = Table(
            headers=sample_data["headers"],
            data=sample_data["data"],
            variant="bordered",
            theme=theme,
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None

    def test_striped_variant(self, slide, sample_data, theme):
        """Test striped variant."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], variant="striped", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None

    def test_minimal_variant(self, slide, sample_data, theme):
        """Test minimal variant."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], variant="minimal", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None


class TestTableSizes:
    """Tests for Table sizes."""

    def test_small_size(self, slide, sample_data, theme):
        """Test small size."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], size="sm", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None

        # Check that font size is smaller
        pptx_table = table_shape.table
        cell = pptx_table.cell(1, 0)  # First data cell
        assert cell.text_frame.paragraphs[0].font.size == Pt(10)

    def test_medium_size(self, slide, sample_data, theme):
        """Test medium size."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], size="md", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None

        # Check that font size is medium
        pptx_table = table_shape.table
        cell = pptx_table.cell(1, 0)  # First data cell
        assert cell.text_frame.paragraphs[0].font.size == Pt(12)

    def test_large_size(self, slide, sample_data, theme):
        """Test large size."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], size="lg", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)
        assert table_shape is not None

        # Check that font size is larger
        pptx_table = table_shape.table
        cell = pptx_table.cell(1, 0)  # First data cell
        assert cell.text_frame.paragraphs[0].font.size == Pt(14)


class TestTableRender:
    """Tests for Table rendering."""

    def test_render_creates_table(self, slide, sample_data, theme):
        """Test that render creates a table shape."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        assert table_shape is not None
        assert hasattr(table_shape, "table")
        pptx_table = table_shape.table
        assert pptx_table is not None

    def test_render_correct_dimensions(self, slide, sample_data, theme):
        """Test that table has correct dimensions."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        pptx_table = table_shape.table
        # Should have 4 rows (1 header + 3 data)
        assert len(pptx_table.rows) == 4
        # Should have 4 columns
        assert len(pptx_table.columns) == 4

    def test_render_headers(self, slide, sample_data, theme):
        """Test that headers are rendered correctly."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        pptx_table = table_shape.table
        # Check header text
        for col_idx, header in enumerate(sample_data["headers"]):
            cell = pptx_table.cell(0, col_idx)
            assert cell.text == header

    def test_render_data(self, slide, sample_data, theme):
        """Test that data is rendered correctly."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        pptx_table = table_shape.table
        # Check data text
        for row_idx, row_data in enumerate(sample_data["data"]):
            for col_idx, value in enumerate(row_data):
                cell = pptx_table.cell(row_idx + 1, col_idx)  # +1 for header
                assert cell.text == value

    def test_render_header_bold(self, slide, sample_data, theme):
        """Test that headers are bold."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        pptx_table = table_shape.table
        # Check that header is bold
        cell = pptx_table.cell(0, 0)
        assert cell.text_frame.paragraphs[0].font.bold is True

    def test_render_position(self, slide, sample_data, theme):
        """Test that table is positioned correctly."""
        left, top, width, height = 1, 2, 6, 3
        table = Table(headers=sample_data["headers"], data=sample_data["data"], theme=theme)
        table_shape = table.render(slide, left=left, top=top, width=width, height=height)

        # Check position (within tolerance for rounding)
        assert abs(table_shape.left.inches - left) < 0.01
        assert abs(table_shape.top.inches - top) < 0.01
        assert abs(table_shape.width.inches - width) < 0.01
        assert abs(table_shape.height.inches - height) < 0.01


class TestTableWithoutTheme:
    """Tests for Table without theme."""

    def test_table_without_theme(self, slide, sample_data):
        """Test that table works without explicit theme."""
        table = Table(headers=sample_data["headers"], data=sample_data["data"])
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        assert table_shape is not None
        pptx_table = table_shape.table
        assert len(pptx_table.rows) == 4
        assert len(pptx_table.columns) == 4


class TestTableEdgeCases:
    """Tests for Table edge cases."""

    def test_single_row(self, slide, theme):
        """Test table with single data row."""
        table = Table(headers=["Name", "Value"], data=[["Item 1", "100"]], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=4, height=2)

        pptx_table = table_shape.table
        assert len(pptx_table.rows) == 2  # Header + 1 data row

    def test_many_columns(self, slide, theme):
        """Test table with many columns."""
        headers = [f"Col{i}" for i in range(10)]
        data = [[f"Val{i}" for i in range(10)]]

        table = Table(headers=headers, data=data, theme=theme)
        table_shape = table.render(slide, left=0.5, top=2, width=9, height=2)

        pptx_table = table_shape.table
        assert len(pptx_table.columns) == 10

    def test_many_rows(self, slide, theme):
        """Test table with many rows."""
        headers = ["Name", "Value"]
        data = [[f"Item {i}", str(i * 100)] for i in range(20)]

        table = Table(headers=headers, data=data, theme=theme)
        table_shape = table.render(slide, left=1, top=1, width=4, height=6)

        pptx_table = table_shape.table
        assert len(pptx_table.rows) == 21  # Header + 20 data rows

    def test_empty_cells(self, slide, theme):
        """Test table with empty cells."""
        table = Table(headers=["Name", "Value"], data=[["Item 1", ""], ["", "200"]], theme=theme)
        table_shape = table.render(slide, left=1, top=2, width=4, height=2)

        pptx_table = table_shape.table
        assert pptx_table.cell(1, 1).text == ""
        assert pptx_table.cell(2, 0).text == ""


class TestTableStripedVariant:
    """Tests specific to striped variant."""

    def test_striped_alternating_rows(self, slide, sample_data, theme):
        """Test that striped variant alternates row colors."""
        table = Table(
            headers=sample_data["headers"], data=sample_data["data"], variant="striped", theme=theme
        )
        table_shape = table.render(slide, left=1, top=2, width=6, height=3)

        pptx_table = table_shape.table
        # Row 1 (index 1) should be normal
        # Row 2 (index 2) should be alternating color
        # Row 3 (index 3) should be normal
        # This is tested by checking the fill is solid for even rows
        cell_even = pptx_table.cell(2, 0)
        # The striped variant should have alternating background
        assert cell_even.fill.type is not None
