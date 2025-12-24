"""
Chart utilities for reusable chart configuration.

Extracted from chart components to avoid duplication and
provide consistent chart styling across all chart types.

Also contains legacy API functions for MCP tool compatibility.
"""

from typing import List, Optional, Dict, Any
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Import design system typography tokens
from ..tokens.typography import FONT_SIZES, FONT_FAMILIES


def configure_legend(
    chart,
    position: str = "right",
    show: bool = True,
    font_family: Optional[str] = None,
    font_size: Optional[int] = None,
):
    """
    Configure chart legend with consistent styling using design tokens.

    Args:
        chart: PowerPoint chart object
        position: Legend position (right, left, top, bottom, corner)
        show: Whether to show legend
        font_family: Font family for legend text (uses design token default if None)
        font_size: Font size in points (uses design token default if None)
    """
    # Use design system defaults
    if font_family is None:
        font_family = FONT_FAMILIES["sans"][0]
    if font_size is None:
        font_size = FONT_SIZES["xs"]  # 10pt for legends

    chart.has_legend = show

    if show:
        position_map = {
            "right": XL_LEGEND_POSITION.RIGHT,
            "left": XL_LEGEND_POSITION.LEFT,
            "top": XL_LEGEND_POSITION.TOP,
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "corner": XL_LEGEND_POSITION.CORNER,
        }

        chart.legend.position = position_map.get(position.lower(), XL_LEGEND_POSITION.RIGHT)
        chart.legend.include_in_layout = False

        # Apply font styling
        if hasattr(chart.legend, "font"):
            chart.legend.font.name = font_family
            chart.legend.font.size = Pt(font_size)


def configure_axes(
    chart,
    show_gridlines: bool = True,
    gridline_color: Optional[RGBColor] = None,
    label_font_family: Optional[str] = None,
    label_font_size: Optional[int] = None,
    label_color: Optional[RGBColor] = None,
):
    """
    Configure chart axes with consistent styling using design tokens.

    Args:
        chart: PowerPoint chart object
        show_gridlines: Whether to show gridlines
        gridline_color: Color for gridlines
        label_font_family: Font family for axis labels (uses design token default if None)
        label_font_size: Font size for axis labels (uses design token default if None)
        label_color: Color for axis labels
    """
    # Use design system defaults
    if label_font_family is None:
        label_font_family = FONT_FAMILIES["sans"][0]
    if label_font_size is None:
        label_font_size = FONT_SIZES["xs"]  # 10pt for axis labels

    try:
        # Configure value axis
        if hasattr(chart, "value_axis"):
            value_axis = chart.value_axis
            value_axis.has_major_gridlines = show_gridlines

            if show_gridlines and gridline_color:
                value_axis.major_gridlines.format.line.color.rgb = gridline_color
                value_axis.major_gridlines.format.line.width = Pt(0.5)

            # Format axis labels
            if hasattr(value_axis, "tick_labels"):
                value_axis.tick_labels.font.name = label_font_family
                value_axis.tick_labels.font.size = Pt(label_font_size)
                if label_color:
                    value_axis.tick_labels.font.color.rgb = label_color

        # Configure category axis
        if hasattr(chart, "category_axis"):
            cat_axis = chart.category_axis
            if hasattr(cat_axis, "tick_labels"):
                cat_axis.tick_labels.font.name = label_font_family
                cat_axis.tick_labels.font.size = Pt(label_font_size)
                if label_color:
                    cat_axis.tick_labels.font.color.rgb = label_color
    except ValueError:
        # Chart type doesn't have axes (e.g., pie charts)
        pass


def set_chart_title(
    chart,
    title: str,
    font_family: Optional[str] = None,
    font_size: Optional[int] = None,
    font_color: Optional[RGBColor] = None,
    bold: bool = True,
):
    """
    Set and style chart title consistently using design tokens.

    Args:
        chart: PowerPoint chart object
        title: Chart title text
        font_family: Font family for title (uses design token default if None)
        font_size: Font size in points (auto-reduced for long titles, uses design tokens)
        font_color: Title color
        bold: Whether to make title bold
    """
    if title:
        chart.has_title = True
        chart_title = chart.chart_title
        chart_title.text_frame.text = title

        # Use design system defaults if not specified
        if font_family is None:
            font_family = FONT_FAMILIES["sans"][0]

        # Auto-scale font size based on title length using design tokens
        # Approximate character thresholds based on typical chart widths
        if font_size is None:
            if len(title) > 60:
                font_size = FONT_SIZES["sm"]  # 12pt - Very long title
            elif len(title) > 45:
                font_size = FONT_SIZES["base"]  # 14pt - Long title
            else:
                font_size = FONT_SIZES["lg"]  # 16pt - Default chart title size

        # Style the title
        para = chart_title.text_frame.paragraphs[0]
        para.font.name = font_family
        para.font.size = Pt(font_size)
        para.font.bold = bold
        if font_color:
            para.font.color.rgb = font_color

        # Disable text wrapping to prevent multi-line titles that overlap charts
        chart_title.text_frame.word_wrap = False


def apply_chart_colors(chart, colors: List[RGBColor]):
    """
    Apply color palette to chart series.

    Args:
        chart: PowerPoint chart object
        colors: List of RGBColor objects for series
    """
    for idx, series in enumerate(chart.series):
        if idx < len(colors):
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = colors[idx]


# =============================================================================
# Legacy API Functions (for MCP tool compatibility)
# =============================================================================

# Chart type mapping
CHART_TYPES = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "bubble": XL_CHART_TYPE.BUBBLE,
}


def add_chart(
    slide,
    chart_type: str,
    left: float,
    top: float,
    width: float,
    height: float,
    categories: List[str],
    series_data: Dict[str, List[float]],
    title: Optional[str] = None,
    has_legend: bool = True,
    legend_position: str = "right",
) -> Any:
    """
    Add a chart to a slide (legacy API).

    Args:
        slide: Slide to add chart to
        chart_type: Type of chart (column, bar, line, pie, etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: Category labels
        series_data: Dictionary of series names to data values
        title: Chart title
        has_legend: Whether to show legend
        legend_position: Legend position (right, left, top, bottom)

    Returns:
        The created chart shape
    """
    xl_chart_type = CHART_TYPES.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = categories

    for series_name, values in series_data.items():
        chart_data.add_series(series_name, values)

    chart_shape = slide.shapes.add_chart(
        xl_chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )

    chart = chart_shape.chart

    if title:
        # Use design tokens (xl = 18pt for legacy compatibility)
        set_chart_title(chart, title, font_size=FONT_SIZES["xl"])

    configure_legend(
        chart,
        position=legend_position,
        show=has_legend,
        # Uses design token defaults: Inter font, 10pt
    )

    return chart_shape


def add_scatter_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    series_data: List[Dict[str, Any]],
    title: Optional[str] = None,
    has_legend: bool = True,
) -> Any:
    """
    Add a scatter plot to a slide (legacy API).

    Args:
        slide: Slide to add chart to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        series_data: List of series, each with 'name', 'x_values', 'y_values'
        title: Chart title
        has_legend: Whether to show legend

    Returns:
        The created chart shape
    """
    chart_data = XyChartData()

    for series in series_data:
        series_obj = chart_data.add_series(series.get("name", "Series"))

        x_values = series.get("x_values", [])
        y_values = series.get("y_values", [])

        for x, y in zip(x_values, y_values):
            series_obj.add_data_point(x, y)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
        chart_data,
    )

    chart = chart_shape.chart

    if title:
        # Use design tokens (xl = 18pt for legacy compatibility)
        set_chart_title(chart, title, font_size=FONT_SIZES["xl"])

    configure_legend(chart, position="right", show=has_legend)
    # Uses design token defaults

    return chart_shape


def add_pie_chart(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    categories: List[str],
    values: List[float],
    title: Optional[str] = None,
    show_percentages: bool = True,
    explode_slice: Optional[int] = None,
) -> Any:
    """
    Add a pie chart to a slide (legacy API).

    Args:
        slide: Slide to add chart to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: Category labels
        values: Data values
        title: Chart title
        show_percentages: Whether to show percentage labels
        explode_slice: Index of slice to explode (separate from pie)

    Returns:
        The created chart shape
    """
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("", values)

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )

    chart = chart_shape.chart

    if title:
        set_chart_title(chart, title, font_family="Calibri", font_size=18)

    configure_legend(chart, position="right", show=True, font_family="Calibri")

    if show_percentages:
        chart.plots[0].has_data_labels = True
        data_labels = chart.plots[0].data_labels
        data_labels.show_percentage = True
        data_labels.show_value = False
        data_labels.font.size = Pt(10)

    if explode_slice is not None and 0 <= explode_slice < len(categories):
        chart.series[0].points[explode_slice].explosion = 20

    return chart_shape


def add_data_table(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    headers: List[str],
    data: List[List[Any]],
    style: str = "medium",
) -> Any:
    """
    Add a formatted data table to a slide (legacy API).

    Args:
        slide: Slide to add table to
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        headers: Column headers
        data: Table data (list of rows)
        style: Table style (light, medium, dark)

    Returns:
        The created table shape
    """
    rows = len(data) + 1  # +1 for header
    cols = len(headers)

    table_shape = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    ).table

    # Set headers
    for col_idx, header in enumerate(headers):
        cell = table_shape.cell(0, col_idx)
        cell.text = str(header)

        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = Pt(12)

        if style == "dark":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(68, 68, 68)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
        elif style == "medium":
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(217, 217, 217)

    # Add data
    for row_idx, row_data in enumerate(data):
        for col_idx, value in enumerate(row_data):
            cell = table_shape.cell(row_idx + 1, col_idx)
            cell.text = str(value)

            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(10)

            if style == "medium" and row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

    return table_shape
