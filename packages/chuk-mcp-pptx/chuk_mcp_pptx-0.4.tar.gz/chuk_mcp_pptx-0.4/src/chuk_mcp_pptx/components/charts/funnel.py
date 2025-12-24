"""
Funnel chart component for sales pipeline and conversion visualization.
"""

from typing import List, Dict, Any, Optional, Tuple, TYPE_CHECKING

if TYPE_CHECKING:
    pass  # For type hints only
from pptx.slide import Slide
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from .base import ChartComponent


class FunnelChart(ChartComponent):
    """
    Funnel chart for visualizing sales pipeline, conversion rates, or process stages.
    """

    def __init__(
        self,
        stages: List[str],
        values: List[float],
        variant: str = "standard",
        show_percentages: bool = True,
        show_values: bool = True,
        title: Optional[str] = None,
        **kwargs,
    ):
        """
        Initialize funnel chart.

        Args:
            stages: List of stage names
            values: List of values for each stage
            variant: Chart variant ('standard', 'cylinder', 'inverted')
            show_percentages: Show conversion percentages
            show_values: Show absolute values
            title: Chart title
            **kwargs: Additional ChartComponent parameters
        """
        super().__init__(**kwargs)
        self.stages = stages
        self.values = values
        self.variant = variant
        self.show_percentages = show_percentages
        self.show_values = show_values
        self.title = title

        # Validate data during initialization
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

    def _calculate_dimensions(self, width: float, height: float) -> List[Tuple[float, float]]:
        """Calculate dimensions for each funnel segment."""
        dimensions = []
        max_value = max(self.values)

        # Reserve space for title if present
        chart_top = 0.5 if self.title else 0.2
        chart_height = height - chart_top - 0.2

        # Calculate segment height
        segment_height = chart_height / len(self.stages)

        for i, value in enumerate(self.values):
            # Calculate width based on value
            segment_width = (value / max_value) * width * 0.8  # 80% of available width

            # Center the segment
            left_offset = (width - segment_width) / 2

            # Calculate vertical position
            top_offset = chart_top + (i * segment_height)

            dimensions.append((left_offset, top_offset, segment_width, segment_height * 0.9))

        return dimensions

    async def render(
        self,
        slide: Slide,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
        placeholder: Optional[Any] = None,
    ):
        """
        Render funnel chart to a slide (shape-based, not native chart).

        Args:
            slide: Slide to add chart to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            None (no chart object since this uses shapes)
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            import logging

            logger = logging.getLogger(__name__)
            logger.info(
                f"FunnelChart targeting placeholder - using bounds: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
            )

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Use defaults if not specified
        left = left if left is not None else self.DEFAULT_LEFT
        top = top if top is not None else self.DEFAULT_TOP
        width = width if width is not None else self.DEFAULT_WIDTH
        height = height if height is not None else self.DEFAULT_HEIGHT

        # Validate data
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Chart data validation failed: {error}")

        # Call the shape-based rendering
        self._render_sync(slide, left=left, top=top, width=width, height=height)
        return None

    def _render_sync(self, slide: Slide, **kwargs) -> None:
        """
        Render funnel chart synchronously using shapes.

        Since PowerPoint doesn't have native funnel charts, we'll create it using shapes.
        """
        left = kwargs.get("left", 1.0)
        top = kwargs.get("top", 1.5)
        width = kwargs.get("width", 8.0)
        height = kwargs.get("height", 5.0)

        # Convert to EMU
        left_emu = Inches(left)
        top_emu = Inches(top)

        # Add title if present
        if self.title:
            title_box = slide.shapes.add_textbox(left_emu, top_emu, Inches(width), Inches(0.5))
            title_frame = title_box.text_frame
            title_frame.text = self.title
            title_frame.paragraphs[0].font.size = Pt(18)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Apply theme color
            if self.tokens:
                fg_color = self.tokens.get("foreground", {}).get("DEFAULT", "#000000")
                rgb = self.hex_to_rgb(fg_color)
                title_frame.paragraphs[0].font.color.rgb = RGBColor(*rgb)

        # Calculate dimensions for each segment
        dimensions = self._calculate_dimensions(width, height)

        # Get theme colors
        colors = self._get_chart_colors()

        # Draw funnel segments
        for i, (stage, value, (seg_left, seg_top, seg_width, seg_height)) in enumerate(
            zip(self.stages, self.values, dimensions)
        ):
            # Create trapezoid shape for funnel segment
            if self.variant == "standard":
                # For standard funnel, use rectangles with gradient effect
                shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(left + seg_left),
                    Inches(top + seg_top),
                    Inches(seg_width),
                    Inches(seg_height),
                )
            else:
                # Use rectangle for other variants
                shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(left + seg_left),
                    Inches(top + seg_top),
                    Inches(seg_width),
                    Inches(seg_height),
                )

            # Apply color
            fill = shape.fill
            fill.solid()
            color_idx = i % len(colors)
            color_hex = colors[color_idx]
            rgb = self.hex_to_rgb(color_hex)
            fill.fore_color.rgb = RGBColor(*rgb)

            # Add border
            line = shape.line
            line.color.rgb = RGBColor(255, 255, 255)
            line.width = Pt(1)

            # Add text to segment
            text_frame = shape.text_frame
            text_frame.clear()  # Clear any default text

            # Enable word wrap for narrow segments, disable for wide ones
            text_frame.word_wrap = seg_width < 1.5

            text_frame.margin_left = Inches(0.02)
            text_frame.margin_right = Inches(0.02)
            text_frame.margin_top = Inches(0.02)
            text_frame.margin_bottom = Inches(0.02)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            # Aggressive font size based on segment width
            if seg_width > 2.5:
                stage_font_size = 11
                value_font_size = 9
                pct_font_size = 8
            elif seg_width > 1.5:
                stage_font_size = 9
                value_font_size = 8
                pct_font_size = 7
            elif seg_width > 0.8:
                stage_font_size = 7
                value_font_size = 6
                pct_font_size = 6
            else:
                # Very narrow - use smallest fonts
                stage_font_size = 6
                value_font_size = 5
                pct_font_size = 5

            # Add stage name
            p = text_frame.paragraphs[0]
            p.text = stage
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(stage_font_size)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)

            # Add value if requested
            if self.show_values:
                p = text_frame.add_paragraph()
                p.text = f"{value:,.0f}"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(value_font_size)
                p.font.color.rgb = RGBColor(255, 255, 255)

            # Add percentage if requested (conversion rate from previous stage)
            if self.show_percentages and i > 0:
                conversion_rate = (value / self.values[i - 1]) * 100
                p = text_frame.add_paragraph()
                p.text = f"{conversion_rate:.1f}%"
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(pct_font_size)
                p.font.color.rgb = RGBColor(255, 255, 255)

    def _get_chart_colors(self) -> List[str]:
        """Get colors for funnel segments."""
        if self.tokens and "chart" in self.tokens:
            colors = self.tokens["chart"]
            if isinstance(colors, list) and colors:
                return colors

        # Default funnel colors (gradient from dark to light)
        return ["#1e40af", "#2563eb", "#3b82f6", "#60a5fa", "#93bbfc", "#c7dbfe"]

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate funnel chart data."""
        if not self.stages or not self.values:
            return False, "Funnel chart requires stages and values"

        if len(self.stages) != len(self.values):
            return False, "Number of stages must match number of values"

        if any(v < 0 for v in self.values):
            return False, "Funnel values must be non-negative"

        # Check for logical funnel progression (values should generally decrease)
        if self.variant != "inverted":
            # Allow some flexibility but warn if values increase significantly
            for i in range(1, len(self.values)):
                if self.values[i] > self.values[i - 1] * 1.1:  # More than 10% increase
                    pass  # Just a warning, not an error

        return True, None

    def validate(self) -> Tuple[bool, Optional[str]]:
        """Validate funnel chart configuration."""
        return self.validate_data()


class GanttChart(ChartComponent):
    """
    Gantt chart for project timeline visualization.
    """

    def __init__(
        self,
        tasks: List[Dict[str, Any]],
        start_date: str,
        end_date: str,
        show_dependencies: bool = True,
        show_milestones: bool = True,
        title: Optional[str] = None,
        **kwargs,
    ):
        """
        Initialize Gantt chart.

        Args:
            tasks: List of task dictionaries with 'name', 'start', 'end', 'progress' keys
            start_date: Project start date
            end_date: Project end date
            show_dependencies: Show task dependencies
            show_milestones: Highlight milestones
            title: Chart title
            **kwargs: Additional ChartComponent parameters
        """
        super().__init__(**kwargs)
        self.tasks = tasks
        self.start_date = start_date
        self.end_date = end_date
        self.show_dependencies = show_dependencies
        self.show_milestones = show_milestones
        self.title = title

        # Validate data during initialization
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

    async def render(
        self,
        slide: Slide,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
        placeholder: Optional[Any] = None,
    ):
        """
        Render Gantt chart to a slide (shape-based, not native chart).

        Args:
            slide: Slide to add chart to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            None (no chart object since this uses shapes)
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            import logging

            logger = logging.getLogger(__name__)
            logger.info(
                f"GanttChart targeting placeholder - using bounds: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
            )

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Use defaults if not specified
        left = left if left is not None else self.DEFAULT_LEFT
        top = top if top is not None else self.DEFAULT_TOP
        width = width if width is not None else self.DEFAULT_WIDTH
        height = height if height is not None else self.DEFAULT_HEIGHT

        # Validate data
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Chart data validation failed: {error}")

        # Call the shape-based rendering
        self._render_sync(slide, left=left, top=top, width=width, height=height)
        return None

    def _render_sync(self, slide: Slide, **kwargs) -> None:
        """
        Render Gantt chart using shapes and lines.

        Since PowerPoint doesn't have native Gantt charts, we create it with shapes.
        """
        # For now, create a simplified Gantt using horizontal bars
        # This would be expanded in a full implementation
        left = kwargs.get("left", 1.0)
        top = kwargs.get("top", 1.5)
        width = kwargs.get("width", 8.0)
        kwargs.get("height", 5.0)

        # Add title
        if self.title:
            title_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(0.5)
            )
            title_frame = title_box.text_frame
            title_frame.text = self.title
            title_frame.paragraphs[0].font.size = Pt(16)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Simplified implementation - would need full timeline logic
        # For demonstration, just create task bars
        colors = self._get_chart_colors()
        task_height = 0.3

        for i, task in enumerate(self.tasks):
            # Create task bar
            bar_top = top + 1.0 + (i * (task_height + 0.1))
            bar_left = left + 2.0  # Leave space for task names
            bar_width = 3.0  # Simplified - would calculate based on dates

            # Task name
            name_box = slide.shapes.add_textbox(
                Inches(left), Inches(bar_top), Inches(1.8), Inches(task_height)
            )
            name_box.text = task.get("name", f"Task {i + 1}")

            # Task bar
            bar = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(bar_left),
                Inches(bar_top),
                Inches(bar_width * task.get("progress", 1.0)),
                Inches(task_height),
            )

            # Apply color
            fill = bar.fill
            fill.solid()
            color_hex = colors[i % len(colors)]
            rgb = self.hex_to_rgb(color_hex)
            fill.fore_color.rgb = RGBColor(*rgb)

    def _get_chart_colors(self) -> List[str]:
        """Get colors for Gantt bars."""
        if self.tokens and "chart" in self.tokens:
            colors = self.tokens["chart"]
            if isinstance(colors, list) and colors:
                return colors
        return ["#3b82f6", "#10b981", "#f59e0b", "#ef4444", "#8b5cf6", "#ec4899"]

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate Gantt chart data."""
        if not self.tasks:
            return False, "Gantt chart requires tasks"

        for task in self.tasks:
            if "name" not in task:
                return False, "Each task must have a name"

        return True, None

    def validate(self) -> Tuple[bool, Optional[str]]:
        """Validate Gantt chart configuration."""
        return self.validate_data()


class HeatmapChart(ChartComponent):
    """
    Heatmap chart for visualizing data density or correlations.
    """

    def __init__(
        self,
        x_labels: List[str],
        y_labels: List[str],
        data: List[List[float]],
        color_scale: str = "heat",
        show_values: bool = True,
        title: Optional[str] = None,
        **kwargs,
    ):
        """
        Initialize heatmap chart.

        Args:
            x_labels: Labels for x-axis
            y_labels: Labels for y-axis
            data: 2D array of values
            color_scale: Color scale ('heat', 'cool', 'diverging')
            show_values: Show values in cells
            title: Chart title
            **kwargs: Additional ChartComponent parameters
        """
        super().__init__(**kwargs)
        self.x_labels = x_labels
        self.y_labels = y_labels
        self.data = data
        self.color_scale = color_scale
        self.show_values = show_values
        self.title = title

        # Validate data during initialization
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

    async def render(
        self,
        slide: Slide,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
        placeholder: Optional[Any] = None,
    ):
        """
        Render heatmap chart to a slide (shape-based, not native chart).

        Args:
            slide: Slide to add chart to
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            None (no chart object since this uses shapes)
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            import logging

            logger = logging.getLogger(__name__)
            logger.info(
                f"HeatmapChart targeting placeholder - using bounds: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
            )

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Use defaults if not specified
        left = left if left is not None else self.DEFAULT_LEFT
        top = top if top is not None else self.DEFAULT_TOP
        width = width if width is not None else self.DEFAULT_WIDTH
        height = height if height is not None else self.DEFAULT_HEIGHT

        # Validate data
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Chart data validation failed: {error}")

        # Call the shape-based rendering
        self._render_sync(slide, left=left, top=top, width=width, height=height)
        return None

    def _render_sync(self, slide: Slide, **kwargs) -> None:
        """Render heatmap using a table with colored cells."""
        left = kwargs.get("left", 1.5)
        top = kwargs.get("top", 1.5)
        width = kwargs.get("width", 7.0)
        height = kwargs.get("height", 4.5)

        # Calculate cell dimensions
        rows = len(self.y_labels) + 1  # +1 for header
        cols = len(self.x_labels) + 1  # +1 for row labels

        # Add table
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        table = table_shape.table

        # Set column headers (x labels)
        for i, label in enumerate(self.x_labels):
            cell = table.cell(0, i + 1)
            cell.text = str(label)
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.bold = True

        # Set row headers (y labels) and data
        for i, y_label in enumerate(self.y_labels):
            # Row header
            cell = table.cell(i + 1, 0)
            cell.text = str(y_label)
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            cell.text_frame.paragraphs[0].font.bold = True

            # Data cells
            for j, value in enumerate(self.data[i]):
                cell = table.cell(i + 1, j + 1)
                if self.show_values:
                    cell.text = str(int(value))
                    cell.text_frame.paragraphs[0].font.size = Pt(8)

                # Color based on value (simple heat scale)
                if self.data:
                    all_values = [val for row in self.data for val in row]
                    min_val = min(all_values)
                    max_val = max(all_values)
                    if max_val > min_val:
                        normalized = (value - min_val) / (max_val - min_val)
                        # Heat color scale: blue (low) to red (high)
                        if normalized < 0.5:
                            # Blue to yellow
                            r = int(255 * (normalized * 2))
                            g = int(255 * (normalized * 2))
                            b = 255
                        else:
                            # Yellow to red
                            r = 255
                            g = int(255 * (1 - (normalized - 0.5) * 2))
                            b = 0

                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(r, g, b)

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate heatmap chart data."""
        if not self.data or not self.x_labels or not self.y_labels:
            return False, "Heatmap requires data, x_labels, and y_labels"

        if len(self.data) != len(self.y_labels):
            return False, "Data rows must match y_labels length"

        for row in self.data:
            if len(row) != len(self.x_labels):
                return False, "Data columns must match x_labels length"

        return True, None

    def validate(self) -> Tuple[bool, Optional[str]]:
        """Validate heatmap configuration."""
        return self.validate_data()
