"""
Browser window container component for PowerPoint presentations.

Provides authentic browser window mockups (Chrome, Safari, Firefox) for
displaying web-based chat interfaces and applications.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES
from ...tokens.platform_colors import (
    get_browser_color,
    MACOS_CONTROLS,
    get_container_ui_color,
)
from ...constants import Platform, ColorKey, Theme


class BrowserWindow(Component):
    """
    Browser window container component.

    Creates a realistic browser window with:
    - Browser chrome (top bar with controls)
    - Address bar
    - Tab bar (optional)
    - Window controls (close, minimize, maximize)
    - Different browser styles (Chrome, Safari, Firefox)
    - Dark or light theme support

    Examples:
        # Chrome browser window
        browser = BrowserWindow(
            title="Slack - Team Chat",
            url="slack.com/messages",
            browser_type="chrome",
            theme=theme
        )

        content_area = browser.render(slide, left=1, top=1, width=8, height=6)

        # Render content inside browser
        conversation.render(
            slide,
            left=content_area['left'],
            top=content_area['top'],
            width=content_area['width']
        )
    """

    def __init__(
        self,
        title: str = "Browser",
        url: str = "example.com",
        browser_type: str = "chrome",  # chrome, safari, firefox
        show_tabs: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize browser window.

        Args:
            title: Page title
            url: URL shown in address bar
            browser_type: Browser type (chrome, safari, firefox)
            show_tabs: Whether to show tab bar
            theme: Optional theme (respects dark/light mode)
        """
        super().__init__(theme)
        self.title = title
        self.url = url
        self.browser_type = browser_type
        self.show_tabs = show_tabs

    def _is_dark_mode(self) -> bool:
        """Check if theme is dark mode."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return sum(bg) < 384
        return False

    def _get_chrome_color(self) -> RGBColor:
        """Get browser chrome color based on theme and browser type."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_browser_color(self.browser_type, ColorKey.BORDER, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color based on theme."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.CHROME, ColorKey.TEXT, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_address_bar_color(self) -> RGBColor:
        """Get address bar color."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.CHROME, ColorKey.ADDRESSBAR, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_content_bg_color(self) -> RGBColor:
        """Get content background color from theme."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return RGBColor(bg[0], bg[1], bg[2])
        hex_color = get_container_ui_color(Platform.CHROME, ColorKey.PLACEHOLDER, Theme.LIGHT)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def render(
        self, slide, left: float, top: float, width: float = 8.0, height: float = 6.0
    ) -> Dict[str, float]:
        """
        Render browser window to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Window width in inches
            height: Window height in inches

        Returns:
            Dict with content area bounds
        """
        shapes = []

        # Window frame
        window_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        window_frame.fill.solid()
        window_frame.fill.fore_color.rgb = self._get_chrome_color()
        hex_color = get_container_ui_color(Platform.CHROME, ColorKey.BORDER, Theme.LIGHT)
        window_frame.line.color.rgb = RGBColor(*self.hex_to_rgb(hex_color))
        window_frame.line.width = Pt(0.5)
        shapes.append(window_frame)

        # Chrome/Title bar
        chrome_height = 0.35
        current_y = top

        # Window controls (macOS style for Safari, Windows style for others)
        control_size = 0.12
        control_y = top + 0.12

        if self.browser_type == "safari":
            # macOS style (left side)
            control_spacing = 0.15
            control_x = left + 0.15

            # Red (close)
            close_btn = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(control_x),
                Inches(control_y),
                Inches(control_size),
                Inches(control_size),
            )
            close_btn.fill.solid()
            close_btn.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(MACOS_CONTROLS["close"]))
            close_btn.line.fill.background()
            shapes.append(close_btn)

            # Yellow (minimize)
            min_btn = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(control_x + control_spacing),
                Inches(control_y),
                Inches(control_size),
                Inches(control_size),
            )
            min_btn.fill.solid()
            min_btn.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(MACOS_CONTROLS["minimize"]))
            min_btn.line.fill.background()
            shapes.append(min_btn)

            # Green (maximize)
            max_btn = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(control_x + control_spacing * 2),
                Inches(control_y),
                Inches(control_size),
                Inches(control_size),
            )
            max_btn.fill.solid()
            max_btn.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(MACOS_CONTROLS["maximize"]))
            max_btn.line.fill.background()
            shapes.append(max_btn)

        current_y += chrome_height

        # Tab bar (if enabled)
        if self.show_tabs:
            tab_height = 0.3
            tab_width = 2.0

            tab = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left + 0.1),
                Inches(current_y - 0.05),
                Inches(tab_width),
                Inches(tab_height),
            )
            tab.fill.solid()
            tab.fill.fore_color.rgb = self._get_address_bar_color()
            tab.line.fill.background()

            # Tab text
            tab_text = tab.text_frame
            tab_text.text = self.title
            tab_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            tab_p = tab_text.paragraphs[0]
            tab_p.alignment = PP_ALIGN.CENTER
            tab_p.font.size = Pt(FONT_SIZES["xs"])
            tab_p.font.color.rgb = self._get_text_color()
            shapes.append(tab)

            current_y += tab_height

        # Address bar
        address_bar_height = 0.35
        address_bar_margin = 0.15

        address_bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + address_bar_margin),
            Inches(current_y),
            Inches(width - address_bar_margin * 2),
            Inches(address_bar_height),
        )
        address_bar.fill.solid()
        address_bar.fill.fore_color.rgb = self._get_address_bar_color()
        address_bar.line.color.rgb = self.get_color("border.DEFAULT")
        address_bar.line.width = Pt(0.5)

        # Address bar text
        address_text = address_bar.text_frame
        address_text.text = f"  🔒 {self.url}"
        address_text.vertical_anchor = MSO_ANCHOR.MIDDLE
        address_p = address_text.paragraphs[0]
        address_p.alignment = PP_ALIGN.LEFT
        address_p.font.size = Pt(FONT_SIZES["xs"])
        address_p.font.color.rgb = (
            self._get_text_color()
            if not self._is_dark_mode()
            else self.get_color("muted.foreground")
        )
        shapes.append(address_bar)

        current_y += address_bar_height + 0.1

        # Content area background
        content_height = height - (current_y - top) - 0.1
        content = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left + 0.05),
            Inches(current_y),
            Inches(width - 0.1),
            Inches(content_height),
        )
        content.fill.solid()
        content.fill.fore_color.rgb = self._get_content_bg_color()
        content.line.fill.background()
        shapes.append(content)

        # Return content area bounds with padding
        padding_h = 0.4  # Horizontal padding (increased)
        padding_v = 0.4  # Vertical padding (increased)

        return {
            "left": left + padding_h,
            "top": current_y + padding_v,
            "width": width - (padding_h * 2),
            "height": content_height - (padding_v * 2),
        }


# TODO: Register component when registry is implemented
