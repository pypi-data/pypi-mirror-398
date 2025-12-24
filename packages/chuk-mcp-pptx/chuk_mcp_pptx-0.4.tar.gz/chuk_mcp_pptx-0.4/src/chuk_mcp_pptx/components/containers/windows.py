"""
Windows application window container for PowerPoint presentations.

Provides authentic Windows window mockups for displaying desktop applications.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES
from ...tokens.platform_colors import (
    WINDOWS_CONTROLS,
    get_container_ui_color,
)
from ...constants import Theme, Platform, ColorKey


class WindowsWindow(Component):
    """
    Windows application window container.

    Creates a realistic Windows window with:
    - Windows 11 style title bar
    - Window controls (minimize, maximize, close)
    - Optional menu/ribbon bar
    - Rounded corners (Windows 11)
    - Dark or light mode support

    Examples:
        # Windows Teams app window
        windows_window = WindowsWindow(
            title="Microsoft Teams",
            app_icon="👥",
            show_menubar=True,
            theme=theme
        )

        content_area = windows_window.render(slide, left=1, top=1, width=7, height=5)
    """

    def __init__(
        self,
        title: str = "Application",
        app_icon: Optional[str] = None,
        show_menubar: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Windows window.

        Args:
            title: Window title
            app_icon: Optional app icon (emoji or text)
            show_menubar: Whether to show menu bar below title bar
            theme: Optional theme (respects dark/light mode)
        """
        super().__init__(theme)
        self.title = title
        self.app_icon = app_icon
        self.show_menubar = show_menubar

    def _is_dark_mode(self) -> bool:
        """Check if theme is dark mode."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return sum(bg) < 384
        return False

    def _get_titlebar_color(self) -> RGBColor:
        """Get title bar color based on theme."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.WINDOWS, ColorKey.TITLEBAR, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color based on theme."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.WINDOWS, ColorKey.TEXT, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_content_bg_color(self) -> RGBColor:
        """Get content background color."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return RGBColor(bg[0], bg[1], bg[2])

        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.WINDOWS, ColorKey.MENUBAR, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def render(
        self, slide, left: float, top: float, width: float = 7.0, height: float = 5.0
    ) -> Dict[str, float]:
        """
        Render Windows window to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Window width in inches
            height: Window height in inches

        Returns:
            Dict with content area bounds
        """
        shapes = []

        # Window frame (Windows 11 rounded corners)
        window_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        window_frame.fill.solid()
        window_frame.fill.fore_color.rgb = self._get_content_bg_color()
        window_frame.line.color.rgb = RGBColor(150, 150, 150)
        window_frame.line.width = Pt(0.5)

        # Windows 11 subtle shadow
        window_frame.shadow.visible = True
        window_frame.shadow.blur_radius = Pt(6)
        window_frame.shadow.distance = Pt(2)
        window_frame.shadow.angle = 90
        window_frame.shadow.transparency = 0.4

        shapes.append(window_frame)

        # Title bar
        titlebar_height = 0.4
        titlebar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(titlebar_height),
        )
        titlebar.fill.solid()
        titlebar.fill.fore_color.rgb = self._get_titlebar_color()
        titlebar.line.fill.background()
        shapes.append(titlebar)

        # App icon and title (left side)
        title_text = f"{self.app_icon} {self.title}" if self.app_icon else self.title

        title_box = slide.shapes.add_textbox(
            Inches(left + 0.15),
            Inches(top + 0.05),
            Inches(width - 2.0),
            Inches(titlebar_height - 0.1),
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.LEFT
        title_p.font.size = Pt(FONT_SIZES["sm"])
        title_p.font.color.rgb = self._get_text_color()
        shapes.append(title_box)

        # Window controls (right side - Windows 11 style)
        control_width = 0.45
        control_height = titlebar_height
        control_x = left + width - (control_width * 3)

        # Minimize button
        min_btn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(control_x),
            Inches(top),
            Inches(control_width),
            Inches(control_height),
        )
        min_btn.fill.solid()
        min_btn.fill.fore_color.rgb = self._get_titlebar_color()
        min_btn.line.fill.background()

        # Minimize icon (horizontal line)
        min_icon = slide.shapes.add_textbox(
            Inches(control_x), Inches(top), Inches(control_width), Inches(control_height)
        )
        min_icon.text_frame.text = "─"
        min_p = min_icon.text_frame.paragraphs[0]
        min_p.alignment = PP_ALIGN.CENTER
        min_icon.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        min_p.font.size = Pt(14)
        min_p.font.color.rgb = self._get_text_color()
        shapes.extend([min_btn, min_icon])

        # Maximize button
        max_btn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(control_x + control_width),
            Inches(top),
            Inches(control_width),
            Inches(control_height),
        )
        max_btn.fill.solid()
        max_btn.fill.fore_color.rgb = self._get_titlebar_color()
        max_btn.line.fill.background()

        # Maximize icon (square)
        max_icon = slide.shapes.add_textbox(
            Inches(control_x + control_width),
            Inches(top),
            Inches(control_width),
            Inches(control_height),
        )
        max_icon.text_frame.text = "□"
        max_p = max_icon.text_frame.paragraphs[0]
        max_p.alignment = PP_ALIGN.CENTER
        max_icon.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        max_p.font.size = Pt(14)
        max_p.font.color.rgb = self._get_text_color()
        shapes.extend([max_btn, max_icon])

        # Close button (red on hover in real Windows 11)
        close_btn = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(control_x + control_width * 2),
            Inches(top),
            Inches(control_width),
            Inches(control_height),
        )
        close_btn.fill.solid()
        close_btn.fill.fore_color.rgb = RGBColor(
            *self.hex_to_rgb(WINDOWS_CONTROLS["close"])
        )  # Windows red
        close_btn.line.fill.background()

        # Close icon (X)
        close_icon = slide.shapes.add_textbox(
            Inches(control_x + control_width * 2),
            Inches(top),
            Inches(control_width),
            Inches(control_height),
        )
        close_icon.text_frame.text = "✕"
        close_p = close_icon.text_frame.paragraphs[0]
        close_p.alignment = PP_ALIGN.CENTER
        close_icon.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        close_p.font.size = Pt(14)
        close_p.font.color.rgb = RGBColor(255, 255, 255)
        shapes.extend([close_btn, close_icon])

        current_y = top + titlebar_height

        # Menu bar (if enabled)
        menubar_height = 0.0
        if self.show_menubar:
            menubar_height = 0.3

            menubar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left),
                Inches(current_y),
                Inches(width),
                Inches(menubar_height),
            )
            menubar.fill.solid()

            if self._is_dark_mode():
                menubar.fill.fore_color.rgb = RGBColor(45, 45, 45)
            else:
                menubar.fill.fore_color.rgb = RGBColor(250, 250, 250)

            menubar.line.fill.background()
            shapes.append(menubar)

            # Menu items
            menu_text = slide.shapes.add_textbox(
                Inches(left + 0.15), Inches(current_y), Inches(width - 0.3), Inches(menubar_height)
            )
            menu_frame = menu_text.text_frame
            menu_frame.text = "File   Edit   View   Help"
            menu_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            menu_p = menu_frame.paragraphs[0]
            menu_p.alignment = PP_ALIGN.LEFT
            menu_p.font.size = Pt(FONT_SIZES["xs"])
            menu_p.font.color.rgb = self._get_text_color()
            shapes.append(menu_text)

            current_y += menubar_height

        # Content area with padding
        content_height = height - (current_y - top)
        padding_h = 0.3  # Horizontal padding
        padding_v = 0.3  # Vertical padding

        return {
            "left": left + padding_h,
            "top": current_y + padding_v,
            "width": width - (padding_h * 2),
            "height": content_height - (padding_v * 2),
        }


# TODO: Register component when registry is implemented
