"""
macOS application window container for PowerPoint presentations.

Provides authentic macOS window mockups for displaying desktop applications.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES
from ...tokens.platform_colors import (
    MACOS_CONTROLS,
    get_container_ui_color,
)
from ...constants import Theme, Platform, ColorKey


class MacOSWindow(Component):
    """
    macOS application window container.

    Creates a realistic macOS window with:
    - Traffic light controls (red, yellow, green)
    - Title bar with app name
    - Optional toolbar
    - macOS Big Sur/Monterey/Ventura style
    - Dark or light mode support

    Examples:
        # macOS Messages app window
        macos_window = MacOSWindow(
            title="Messages",
            app_icon="💬",
            show_toolbar=True,
            theme=theme
        )

        content_area = macos_window.render(slide, left=1, top=1, width=7, height=5)
    """

    def __init__(
        self,
        title: str = "Application",
        app_icon: Optional[str] = None,
        show_toolbar: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize macOS window.

        Args:
            title: Window title
            app_icon: Optional app icon (emoji or text)
            show_toolbar: Whether to show toolbar below title bar
            theme: Optional theme (respects dark/light mode)
        """
        super().__init__(theme)
        self.title = title
        self.app_icon = app_icon
        self.show_toolbar = show_toolbar

    def _is_dark_mode(self) -> bool:
        """Check if theme is dark mode."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return sum(bg) < 384
        return False

    def _get_titlebar_color(self) -> RGBColor:
        """Get title bar color based on theme."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.MACOS, ColorKey.TITLEBAR, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color based on theme."""
        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.MACOS, ColorKey.TEXT, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_content_bg_color(self) -> RGBColor:
        """Get content background color."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return RGBColor(bg[0], bg[1], bg[2])
        hex_color = get_container_ui_color(Platform.MACOS, ColorKey.CONTENT_BG, Theme.LIGHT)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def render(
        self, slide, left: float, top: float, width: float = 7.0, height: float = 5.0
    ) -> Dict[str, float]:
        """
        Render macOS window to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Window width in inches
            height: Window height in inches

        Returns:
            Dict with content area bounds
        """
        shapes = []

        # Window frame with shadow
        window_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        window_frame.fill.solid()
        window_frame.fill.fore_color.rgb = self._get_content_bg_color()
        hex_color = get_container_ui_color(Platform.MACOS, ColorKey.BORDER, Theme.LIGHT)
        window_frame.line.color.rgb = RGBColor(*self.hex_to_rgb(hex_color))
        window_frame.line.width = Pt(0.5)

        # macOS-style shadow
        window_frame.shadow.visible = True
        window_frame.shadow.blur_radius = Pt(8)
        window_frame.shadow.distance = Pt(3)
        window_frame.shadow.angle = 90
        window_frame.shadow.transparency = 0.3

        shapes.append(window_frame)

        # Title bar
        titlebar_height = 0.45
        titlebar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(titlebar_height),
        )
        titlebar.fill.solid()
        titlebar.fill.fore_color.rgb = self._get_titlebar_color()
        titlebar.line.fill.background()
        shapes.append(titlebar)

        # Traffic light controls (left side)
        control_size = 0.12
        control_y = top + 0.15
        control_x = left + 0.15
        control_spacing = 0.18

        # Red (close)
        close_btn = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(control_x),
            Inches(control_y),
            Inches(control_size),
            Inches(control_size),
        )
        close_btn.fill.solid()
        close_btn.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(MACOS_CONTROLS["close"]))
        # Border is slightly darker than fill
        close_btn.line.color.rgb = RGBColor(220, 85, 76)
        close_btn.line.width = Pt(0.5)
        shapes.append(close_btn)

        # Yellow (minimize)
        min_btn = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(control_x + control_spacing),
            Inches(control_y),
            Inches(control_size),
            Inches(control_size),
        )
        min_btn.fill.solid()
        min_btn.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(MACOS_CONTROLS["minimize"]))
        # Border is slightly darker than fill
        min_btn.line.color.rgb = RGBColor(225, 169, 41)
        min_btn.line.width = Pt(0.5)
        shapes.append(min_btn)

        # Green (maximize)
        max_btn = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(control_x + control_spacing * 2),
            Inches(control_y),
            Inches(control_size),
            Inches(control_size),
        )
        max_btn.fill.solid()
        max_btn.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(MACOS_CONTROLS["maximize"]))
        # Border is slightly darker than fill
        max_btn.line.color.rgb = RGBColor(35, 181, 57)
        max_btn.line.width = Pt(0.5)
        shapes.append(max_btn)

        # Window title (centered)
        title_text = f"{self.app_icon} {self.title}" if self.app_icon else self.title

        title_box = slide.shapes.add_textbox(
            Inches(left + 1.0),
            Inches(top + 0.05),
            Inches(width - 2.0),
            Inches(titlebar_height - 0.1),
        )
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_p = title_frame.paragraphs[0]
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.size = Pt(FONT_SIZES["sm"])
        title_p.font.bold = False
        title_p.font.color.rgb = self._get_text_color()
        shapes.append(title_box)

        current_y = top + titlebar_height

        # Toolbar (if enabled)
        toolbar_height = 0.0
        if self.show_toolbar:
            toolbar_height = 0.35

            toolbar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(left),
                Inches(current_y),
                Inches(width),
                Inches(toolbar_height),
            )
            toolbar.fill.solid()

            # Toolbar is slightly lighter/darker than title bar
            theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
            hex_color = get_container_ui_color(Platform.MACOS, ColorKey.TOOLBAR, theme_mode)
            toolbar.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(hex_color))

            toolbar.line.fill.background()
            shapes.append(toolbar)

            current_y += toolbar_height

        # Content area with padding
        content_height = height - (current_y - top)
        padding_h = 0.3  # Horizontal padding
        padding_v = 0.3  # Vertical padding

        # Return content area bounds (no need to draw, already part of window frame)
        return {
            "left": left + padding_h,
            "top": current_y + padding_v,
            "width": width - (padding_h * 2),
            "height": content_height - (padding_v * 2),
        }


# TODO: Register component when registry is implemented
