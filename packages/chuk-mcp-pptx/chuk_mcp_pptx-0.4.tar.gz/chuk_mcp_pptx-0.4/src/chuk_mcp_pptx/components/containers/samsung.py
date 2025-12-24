"""
Samsung device container component for PowerPoint presentations.

Provides an authentic Samsung/Android phone mockup frame for displaying
chat conversations and mobile content.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES


class SamsungContainer(Component):
    """
    Samsung/Android device container component.

    Creates a realistic Samsung phone frame with:
    - Rounded corners matching Samsung design
    - Punch-hole camera (modern Samsung style)
    - Status bar with Android icons
    - Navigation bar at bottom
    - Dark or light theme support
    - Content area for rendering conversations

    Examples:
        # Samsung with Android Messages conversation
        samsung = SamsungContainer(
            title="Messages",
            variant="galaxy-s",
            theme=theme
        )

        # Render the container and get content area bounds
        content_area = samsung.render(slide, left=1, top=1)

        # Render conversation inside
        conversation.render(
            slide,
            left=content_area['left'],
            top=content_area['top'],
            width=content_area['width']
        )
    """

    def __init__(
        self,
        title: str = "Samsung",
        variant: str = "galaxy-s",  # galaxy-s, galaxy-note
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Samsung container.

        Args:
            title: Title shown in the status bar area
            variant: Samsung variant (galaxy-s, galaxy-note)
            theme: Optional theme (respects dark/light mode)
        """
        super().__init__(theme)
        self.title = title
        self.variant = variant

    def _get_device_color(self) -> RGBColor:
        """Get device frame color based on theme."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                if sum(bg) < 384:
                    return RGBColor(30, 30, 30)  # Phantom black
                else:
                    return RGBColor(255, 255, 255)  # Phantom white

        return RGBColor(30, 30, 30)

    def _get_screen_bg_color(self) -> RGBColor:
        """Get screen background color based on theme."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return RGBColor(bg[0], bg[1], bg[2])

        return RGBColor(255, 255, 255)

    def _get_text_color(self) -> RGBColor:
        """Get text color based on theme."""
        if self.theme and isinstance(self.theme, dict):
            fg = self.theme.get("colors", {}).get("foreground", {}).get("DEFAULT")
            if fg and isinstance(fg, (list, tuple)) and len(fg) >= 3:
                return RGBColor(fg[0], fg[1], fg[2])

        return RGBColor(0, 0, 0)

    def _is_dark_mode(self) -> bool:
        """Check if theme is dark mode."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return sum(bg) < 384
        return False

    def render(
        self, slide, left: float, top: float, width: float = 3.5, height: float = 6.5
    ) -> Dict[str, float]:
        """
        Render Samsung container to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Device width in inches
            height: Device height in inches

        Returns:
            Dict with content area bounds
        """
        shapes = []

        # Device frame
        device_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        device_frame.fill.solid()
        device_frame.fill.fore_color.rgb = self._get_device_color()
        device_frame.line.color.rgb = self._get_device_color()
        device_frame.line.width = Pt(2)
        shapes.append(device_frame)

        # Screen area
        screen_margin = 0.06
        screen_left = left + screen_margin
        screen_top = top + screen_margin
        screen_width = width - (screen_margin * 2)
        screen_height = height - (screen_margin * 2)

        screen = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(screen_left),
            Inches(screen_top),
            Inches(screen_width),
            Inches(screen_height),
        )
        screen.fill.solid()
        screen.fill.fore_color.rgb = self._get_screen_bg_color()
        screen.line.fill.background()
        shapes.append(screen)

        # Punch-hole camera (centered at top)
        camera_size = 0.12
        camera_left = screen_left + (screen_width - camera_size) / 2

        camera = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(camera_left),
            Inches(screen_top + 0.05),
            Inches(camera_size),
            Inches(camera_size),
        )
        camera.fill.solid()
        camera.fill.fore_color.rgb = self._get_device_color()
        camera.line.fill.background()
        shapes.append(camera)

        # Status bar
        status_bar_height = 0.25
        status_bar = slide.shapes.add_textbox(
            Inches(screen_left + 0.1),
            Inches(screen_top + 0.05),
            Inches(screen_width - 0.2),
            Inches(status_bar_height),
        )
        status_frame = status_bar.text_frame
        status_frame.text = "10:30"
        status_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        status_p = status_frame.paragraphs[0]
        status_p.alignment = PP_ALIGN.CENTER
        status_p.font.size = Pt(FONT_SIZES["xs"])
        status_p.font.bold = True
        status_p.font.color.rgb = self._get_text_color()
        shapes.append(status_bar)

        # Navigation bar (bottom - Android style)
        nav_bar_height = 0.25
        nav_bar_top = screen_top + screen_height - nav_bar_height

        # Draw three navigation buttons
        button_y = nav_bar_top + 0.08
        button_height = 0.08
        button_color = RGBColor(90, 90, 90) if self._is_dark_mode() else RGBColor(60, 60, 60)

        # Back button (triangle pointing left)
        back_left = screen_left + screen_width * 0.2
        back_button = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(back_left),
            Inches(button_y),
            Inches(button_height),
            Inches(button_height),
        )
        back_button.fill.solid()
        back_button.fill.fore_color.rgb = button_color
        back_button.line.fill.background()
        back_button.rotation = 270  # Point left
        shapes.append(back_button)

        # Home button (circle)
        home_left = screen_left + (screen_width - button_height) / 2
        home_button = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(home_left),
            Inches(button_y),
            Inches(button_height),
            Inches(button_height),
        )
        home_button.fill.solid()
        home_button.fill.fore_color.rgb = button_color
        home_button.line.fill.background()
        shapes.append(home_button)

        # Recent apps button (square)
        recent_left = screen_left + screen_width * 0.7
        recent_button = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(recent_left),
            Inches(button_y),
            Inches(button_height),
            Inches(button_height),
        )
        recent_button.fill.solid()
        recent_button.fill.fore_color.rgb = button_color
        recent_button.line.fill.background()
        shapes.append(recent_button)

        # Calculate content area
        content_top = screen_top + status_bar_height + 0.1
        content_bottom = nav_bar_top - 0.05
        content_height = content_bottom - content_top

        # Add padding for realistic spacing
        padding_h = 0.15  # Samsung has tighter horizontal padding
        padding_v = 0.15

        content_left = screen_left + padding_h
        content_width = screen_width - (padding_h * 2)

        return {
            "left": content_left,
            "top": content_top + padding_v,
            "width": content_width,
            "height": content_height - (padding_v * 2),
        }


# TODO: Register component when registry is implemented
