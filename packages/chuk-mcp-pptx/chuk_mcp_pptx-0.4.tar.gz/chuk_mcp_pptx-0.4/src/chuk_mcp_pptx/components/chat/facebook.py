"""
Facebook Messenger chat components for PowerPoint presentations.

Provides realistic Facebook Messenger-style chat interfaces.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, FONT_FAMILIES
from ...tokens.platform_colors import get_chat_color, CHAT_COLORS
from ...constants import Theme, Platform, ColorKey


class FacebookMessengerBubble(Component):
    """
    Facebook Messenger-style chat bubble component.

    Replicates Facebook Messenger's look with:
    - Blue bubbles for sent messages (right-aligned)
    - Gray bubbles for received messages (left-aligned)
    - Circular avatars
    - Facebook Messenger blue (#0084FF)
    - Compact bubble design

    Examples:
        # Sent message
        msg = FacebookMessengerBubble(
            text="Hey! Want to grab coffee?",
            variant="sent",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Received message with avatar
        msg = FacebookMessengerBubble(
            text="Sure! When works for you?",
            variant="received",
            avatar_text="JS",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(
        self,
        text: str,
        variant: str = "received",
        avatar_text: Optional[str] = None,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Facebook Messenger bubble.

        Args:
            text: Message text
            variant: Message variant (sent, received)
            avatar_text: Avatar initials (for received messages)
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.variant = variant
        self.avatar_text = avatar_text

    def _get_bubble_color(self) -> RGBColor:
        """Get Facebook Messenger bubble color."""
        hex_color = get_chat_color(Platform.FACEBOOK, self.variant, Theme.LIGHT)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color."""
        hex_color = CHAT_COLORS[Platform.FACEBOOK][ColorKey.TEXT]
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _calculate_bubble_height(self, width: float) -> float:
        """Estimate bubble height."""
        chars_per_line = int(width * 10)  # More generous
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        line_height = 0.2  # Increased from 0.17
        padding = 0.28  # More padding
        return (lines * line_height) + padding

    def render(self, slide, left: float, top: float, width: float = 6.0) -> list:
        """Render Facebook Messenger bubble."""
        shapes = []

        # Avatar settings
        avatar_size = 0.4
        avatar_gap = 0.1

        # Calculate dimensions
        bubble_width = min(width * 0.6, width - 0.8)
        bubble_height = self._calculate_bubble_height(bubble_width)

        # Position based on variant
        if self.variant == "sent":
            bubble_left = left + width - bubble_width
            show_avatar = False
        else:
            # Avatar for received messages
            if self.avatar_text:
                bubble_left = left + avatar_size + avatar_gap
                bubble_width -= avatar_size + avatar_gap
                show_avatar = True
            else:
                bubble_left = left
                show_avatar = False

        # Render avatar (for received)
        if show_avatar:
            avatar = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left),
                Inches(top + bubble_height - avatar_size),  # Bottom align
                Inches(avatar_size),
                Inches(avatar_size),
            )
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = RGBColor(
                *self.hex_to_rgb(CHAT_COLORS[Platform.FACEBOOK][ColorKey.SENT])
            )
            avatar.line.fill.background()

            # Avatar text
            av_frame = avatar.text_frame
            av_frame.text = self.avatar_text
            av_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            av_p = av_frame.paragraphs[0]
            av_p.alignment = PP_ALIGN.CENTER
            av_p.font.size = Pt(FONT_SIZES["sm"])
            av_p.font.bold = True
            av_p.font.color.rgb = self.get_color("primary.foreground")
            shapes.append(avatar)

        # Create bubble
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bubble_left),
            Inches(top),
            Inches(bubble_width),
            Inches(bubble_height),
        )

        # Style bubble - Facebook style is very rounded
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = self._get_bubble_color()
        bubble.line.fill.background()

        # Text frame
        text_frame = bubble.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        padding = 0.13
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        # Message text
        p = text_frame.paragraphs[0]
        p.text = self.text
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(FONT_SIZES["lg"])
        p.font.name = FONT_FAMILIES["sans"][0]
        p.font.color.rgb = self._get_text_color()
        p.line_spacing = 1.3

        shapes.append(bubble)

        return shapes


class FacebookMessengerConversation(Component):
    """
    Facebook Messenger-style conversation component.

    Displays a complete Messenger conversation with automatic spacing.

    Examples:
        messages = [
            {"text": "Hey!", "variant": "received", "avatar_text": "JS"},
            {"text": "Hi there!", "variant": "sent"}
        ]
        conversation = FacebookMessengerConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=8)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.1,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Facebook Messenger conversation.

        Args:
            messages: List of message dicts with FacebookMessengerBubble params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render Facebook Messenger conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = FacebookMessengerBubble(
                text=msg_data.get("text", ""),
                variant=msg_data.get("variant", "received"),
                avatar_text=msg_data.get("avatar_text"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            bubble_width = min(width * 0.6, width - 0.8)
            msg_height = message._calculate_bubble_height(bubble_width)
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# FacebookMessengerBubble - Variants: sent, received
# FacebookMessengerConversation - Full conversation flow
