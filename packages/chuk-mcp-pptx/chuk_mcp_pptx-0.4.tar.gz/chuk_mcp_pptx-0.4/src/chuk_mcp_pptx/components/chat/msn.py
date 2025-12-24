"""
MSN Messenger chat components for PowerPoint presentations.

Provides nostalgic MSN Messenger (Windows Live Messenger)-style chat interfaces.
Recreates the classic mid-2000s Microsoft messaging experience.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, FONT_FAMILIES


class MSNBubble(Component):
    """
    MSN Messenger-style chat bubble component.

    Recreates the classic MSN interface with:
    - Colored display names
    - "says:" text between name and message
    - Emoticon support indicators
    - No chat bubbles - simple text layout
    - Classic MSN green/orange theme
    - Nostalgic 2000s design

    Examples:
        # Sent message
        msg = MSNBubble(
            text="Hey! Want to play some games?",
            display_name="CoolGuy",
            variant="sent",
            timestamp="21:30",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Received message
        msg = MSNBubble(
            text="Sure! What do you want to play?",
            display_name="AwesomeGirl",
            variant="received",
            timestamp="21:31",
            emoticon="😊",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(
        self,
        text: str,
        display_name: str,
        variant: str = "received",
        timestamp: Optional[str] = None,
        emoticon: Optional[str] = None,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize MSN bubble.

        Args:
            text: Message text
            display_name: MSN display name
            variant: Message variant (sent, received)
            timestamp: Optional timestamp (e.g., "21:30")
            emoticon: Optional emoticon to display
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.display_name = display_name
        self.variant = variant
        self.timestamp = timestamp
        self.emoticon = emoticon

    def _get_display_name_color(self) -> RGBColor:
        """Get display name color."""
        if self.variant == "sent":
            # MSN green for sent
            return self.get_color("success.DEFAULT")
        else:
            # MSN orange for received
            return self.get_color("warning.DEFAULT")

    def _get_text_color(self) -> RGBColor:
        """Get message text color."""
        # Black text
        return self.get_color("foreground.DEFAULT")

    def _calculate_message_height(self, width: float) -> float:
        """Estimate message height."""
        # Display name + "says:" line
        header_height = 0.22  # Increased from 0.17

        # Calculate text height
        chars_per_line = int(width * 11)  # More generous
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        text_height = lines * 0.19  # Increased from 0.16

        # Emoticon
        emoticon_height = 0.18 if self.emoticon else 0

        return header_height + text_height + emoticon_height + 0.08  # More padding

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render MSN message."""
        shapes = []
        current_top = top

        # Display name + "says:" (e.g., "CoolGuy says:")
        header_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(current_top),
            Inches(width),
            Inches(0.22),  # Increased from 0.17
        )
        header_frame = header_box.text_frame
        header_frame.word_wrap = False
        header_p = header_frame.paragraphs[0]
        header_p.alignment = PP_ALIGN.LEFT

        # Add display name in color
        if self.timestamp:
            header_frame.text = f"{self.display_name} says: ({self.timestamp})"
        else:
            header_frame.text = f"{self.display_name} says:"

        header_p.font.size = Pt(FONT_SIZES["sm"])  # Larger
        header_p.font.bold = True
        header_p.font.name = FONT_FAMILIES["sans"][0]  # Classic MSN used Tahoma
        header_p.font.color.rgb = self._get_display_name_color()
        shapes.append(header_box)
        current_top += 0.20  # More spacing

        # Message text
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(current_top), Inches(width), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        # Add emoticon at start if present
        if self.emoticon:
            text_frame.text = f"{self.emoticon} {self.text}"
        else:
            text_frame.text = self.text

        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.LEFT
        text_p.font.size = Pt(FONT_SIZES["sm"])  # Larger
        text_p.font.name = FONT_FAMILIES["sans"][0]
        text_p.font.color.rgb = self._get_text_color()
        text_p.line_spacing = 1.3  # Better line spacing
        shapes.append(text_box)

        return shapes


class MSNConversation(Component):
    """
    MSN Messenger-style conversation component.

    Displays a complete MSN Messenger conversation with classic styling.

    Examples:
        messages = [
            {
                "text": "Hey! What are you up to?",
                "display_name": "AwesomeGirl",
                "variant": "received",
                "timestamp": "21:00",
                "emoticon": "😊"
            },
            {
                "text": "Not much, just listening to music",
                "display_name": "CoolGuy",
                "variant": "sent",
                "timestamp": "21:01",
                "emoticon": "🎵"
            }
        ]
        conversation = MSNConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=7)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.12,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize MSN conversation.

        Args:
            messages: List of message dicts with MSNBubble params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render MSN Messenger conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = MSNBubble(
                text=msg_data.get("text", ""),
                display_name=msg_data.get("display_name", "User"),
                variant=msg_data.get("variant", "received"),
                timestamp=msg_data.get("timestamp"),
                emoticon=msg_data.get("emoticon"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            msg_height = message._calculate_message_height(width)
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# MSNBubble - Classic MSN Messenger style with "says:" and emoticons
# MSNConversation - Full conversation flow
