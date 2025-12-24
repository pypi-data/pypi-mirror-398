"""
iMessage (iOS) chat components for PowerPoint presentations.

Provides realistic iMessage-style chat interfaces.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, FONT_FAMILIES
from ...tokens.platform_colors import get_chat_color, CHAT_COLORS
from ...constants import Theme, Platform, ColorKey


class iMessageBubble(Component):
    """
    iMessage-style chat bubble component.

    Replicates the iOS Messages app appearance with:
    - Blue bubbles for sent messages (right-aligned)
    - Gray bubbles for received messages (left-aligned)
    - Rounded corners with tail
    - White text on blue, dark text on gray
    - Timestamps below bubbles

    Examples:
        # Sent message (blue bubble)
        msg = iMessageBubble(
            text="Hey! How are you?",
            variant="sent",
            timestamp="10:30 AM",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Received message (gray bubble)
        msg = iMessageBubble(
            text="I'm doing great, thanks!",
            variant="received",
            timestamp="10:31 AM",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(
        self,
        text: str,
        variant: str = "received",
        timestamp: Optional[str] = None,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize iMessage bubble.

        Args:
            text: Message text
            variant: Message variant (sent, received)
            timestamp: Optional timestamp
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.variant = variant
        self.timestamp = timestamp

    def _get_bubble_color(self) -> RGBColor:
        """Get iMessage-specific bubble color."""
        hex_color = get_chat_color(Platform.IOS, self.variant, Theme.LIGHT)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color."""
        if self.variant == ColorKey.SENT:
            hex_color = CHAT_COLORS[Platform.IOS][ColorKey.TEXT_SENT]
        else:
            hex_color = CHAT_COLORS[Platform.IOS][ColorKey.TEXT_RECEIVED]
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _calculate_bubble_height(self, width: float) -> float:
        """Estimate bubble height based on text length."""
        chars_per_line = int(width * 11)  # More conservative
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        line_height = 0.2  # Increased from 0.18
        padding = 0.3  # Increased from 0.25
        return (lines * line_height) + padding

    def render(self, slide, left: float, top: float, width: float = 6.0) -> list:
        """Render iMessage bubble."""
        shapes = []

        # Calculate dimensions
        bubble_width = min(width * 0.65, width - 1.0)
        bubble_height = self._calculate_bubble_height(bubble_width)

        # Position based on variant
        if self.variant == "sent":
            bubble_left = left + width - bubble_width
        else:
            bubble_left = left

        # Create bubble
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bubble_left),
            Inches(top),
            Inches(bubble_width),
            Inches(bubble_height),
        )

        # Style bubble - iMessage has more rounded corners
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = self._get_bubble_color()
        bubble.line.fill.background()  # No border

        # Text frame
        text_frame = bubble.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Top align for consistency

        padding = 0.12
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        # Message text
        p = text_frame.paragraphs[0]
        p.text = self.text
        p.alignment = PP_ALIGN.LEFT  # All text left-aligned within bubble
        p.font.size = Pt(FONT_SIZES["lg"])  # 16pt (was 15pt, using closest design token)
        p.font.name = FONT_FAMILIES["sans"][0]  # Use design system sans font
        p.font.color.rgb = self._get_text_color()

        shapes.append(bubble)

        # Timestamp (below bubble, subtle)
        if self.timestamp:
            ts_box = slide.shapes.add_textbox(
                Inches(bubble_left),
                Inches(top + bubble_height + 0.05),
                Inches(bubble_width),
                Inches(0.2),
            )
            ts_frame = ts_box.text_frame
            ts_frame.text = self.timestamp
            ts_p = ts_frame.paragraphs[0]
            ts_p.alignment = PP_ALIGN.CENTER
            ts_p.font.size = Pt(FONT_SIZES["xs"])
            # Use muted text color from design system
            ts_p.font.color.rgb = self.get_color("muted.foreground")
            shapes.append(ts_box)

        return shapes


class iMessageConversation(Component):
    """
    iMessage-style conversation component.

    Displays a complete iMessage conversation with automatic spacing.

    Examples:
        messages = [
            {"text": "Hey!", "variant": "received", "timestamp": "10:30 AM"},
            {"text": "Hi there!", "variant": "sent", "timestamp": "10:31 AM"}
        ]
        conversation = iMessageConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=8)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.15,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize iMessage conversation.

        Args:
            messages: List of message dicts with iMessageBubble params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render iMessage conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = iMessageBubble(
                text=msg_data.get("text", ""),
                variant=msg_data.get("variant", "received"),
                timestamp=msg_data.get("timestamp"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            bubble_width = min(width * 0.65, width - 1.0)
            msg_height = message._calculate_bubble_height(bubble_width)
            # Add extra space for timestamp if present
            timestamp_space = 0.25 if msg_data.get("timestamp") else 0
            current_top += msg_height + timestamp_space + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# iMessageBubble - Variants: sent, received
# iMessageConversation - Full conversation flow
