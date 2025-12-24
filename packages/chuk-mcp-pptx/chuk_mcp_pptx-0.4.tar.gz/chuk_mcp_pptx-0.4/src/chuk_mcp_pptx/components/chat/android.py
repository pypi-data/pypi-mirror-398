"""
Android Messages chat components for PowerPoint presentations.

Provides realistic Android Messages-style chat interfaces with Material Design.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, FONT_FAMILIES
from ...tokens.platform_colors import get_chat_color, CHAT_COLORS
from ...constants import ComponentSizing, Theme, Platform, ColorKey


class AndroidMessageBubble(Component):
    """
    Android Messages-style chat bubble component.

    Replicates Google's Material Design messaging with:
    - Purple/blue bubbles for sent messages (right-aligned)
    - Light gray bubbles for received messages (left-aligned)
    - Subtle shadows
    - Material Design colors and typography

    Examples:
        # Sent message
        msg = AndroidMessageBubble(
            text="On my way!",
            variant="sent",
            timestamp="2:15 PM",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Received message with sender
        msg = AndroidMessageBubble(
            text="Great! See you soon.",
            sender="John",
            variant="received",
            timestamp="2:16 PM",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(
        self,
        text: str,
        sender: Optional[str] = None,
        variant: str = "received",
        timestamp: Optional[str] = None,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Android Messages bubble.

        Args:
            text: Message text
            sender: Sender name (for received messages)
            variant: Message variant (sent, received)
            timestamp: Optional timestamp
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.sender = sender
        self.variant = variant
        self.timestamp = timestamp

    def _get_bubble_color(self) -> RGBColor:
        """Get Android Messages bubble color."""
        hex_color = get_chat_color(Platform.ANDROID, self.variant, Theme.LIGHT)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color."""
        if self.variant == ColorKey.SENT:
            hex_color = CHAT_COLORS[Platform.ANDROID][ColorKey.TEXT_SENT]
        else:
            hex_color = CHAT_COLORS[Platform.ANDROID][ColorKey.TEXT_RECEIVED]
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _calculate_bubble_height(self, width: float) -> float:
        """Estimate bubble height."""
        chars_per_line = int(width * 11)  # More conservative
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        line_height = 0.2  # Increased from 0.18
        padding = 0.3  # Increased from 0.25
        sender_height = 0.22 if self.sender else 0  # Increased from 0.2
        return (lines * line_height) + padding + sender_height

    def render(self, slide, left: float, top: float, width: float = 6.0) -> list:
        """Render Android Messages bubble."""
        shapes = []

        # Calculate dimensions
        bubble_width = min(width * 0.7, width - 1.0)
        bubble_height = self._calculate_bubble_height(bubble_width)

        # Position based on variant
        if self.variant == "sent":
            bubble_left = left + width - bubble_width
        else:
            bubble_left = left

        # Create bubble
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bubble_left),
            Inches(top),
            Inches(bubble_width),
            Inches(bubble_height),
        )

        # Style bubble - Material Design
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = self._get_bubble_color()
        bubble.line.fill.background()

        # Subtle shadow for Material Design
        bubble.shadow.visible = True
        bubble.shadow.blur_radius = Pt(ComponentSizing.SHADOW_BLUR_MD)
        bubble.shadow.distance = Pt(ComponentSizing.SHADOW_DISTANCE_MD)
        bubble.shadow.angle = 90
        bubble.shadow.transparency = 0.4

        # Text frame
        text_frame = bubble.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        padding = 0.12
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        current_p = text_frame.paragraphs[0]

        # Sender name (for received messages)
        if self.sender and self.variant == "received":
            current_p.text = self.sender
            current_p.alignment = PP_ALIGN.LEFT
            current_p.font.size = Pt(FONT_SIZES["sm"])
            current_p.font.bold = True
            current_p.font.color.rgb = RGBColor(
                *self.hex_to_rgb(CHAT_COLORS[Platform.ANDROID][ColorKey.TIMESTAMP])
            )
            current_p = text_frame.add_paragraph()
            current_p.space_before = Pt(ComponentSizing.SPACE_SM)

        # Message text
        current_p.text = self.text
        current_p.alignment = PP_ALIGN.LEFT
        current_p.font.size = Pt(FONT_SIZES["base"])
        current_p.font.name = FONT_FAMILIES["sans"][0]
        current_p.font.color.rgb = self._get_text_color()

        shapes.append(bubble)

        # Timestamp (in corner, small)
        if self.timestamp:
            ts_width = 0.8
            if self.variant == "sent":
                ts_left = bubble_left + bubble_width - ts_width - 0.05
            else:
                ts_left = bubble_left + 0.05

            ts_box = slide.shapes.add_textbox(
                Inches(ts_left), Inches(top + bubble_height + 0.05), Inches(ts_width), Inches(0.2)
            )
            ts_frame = ts_box.text_frame
            ts_frame.text = self.timestamp
            ts_p = ts_frame.paragraphs[0]
            ts_p.alignment = PP_ALIGN.LEFT if self.variant == ColorKey.RECEIVED else PP_ALIGN.RIGHT
            ts_p.font.size = Pt(FONT_SIZES["xs"])
            ts_p.font.color.rgb = RGBColor(
                *self.hex_to_rgb(CHAT_COLORS[Platform.ANDROID][ColorKey.TIMESTAMP])
            )
            shapes.append(ts_box)

        return shapes


class AndroidConversation(Component):
    """
    Android Messages-style conversation component.

    Displays a complete Android Messages conversation with automatic spacing.

    Examples:
        messages = [
            {"text": "Meeting soon?", "sender": "Alice", "variant": "received", "timestamp": "2:15 PM"},
            {"text": "Yes, on my way!", "variant": "sent", "timestamp": "2:16 PM"}
        ]
        conversation = AndroidConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=8)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.15,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Android conversation.

        Args:
            messages: List of message dicts with AndroidMessageBubble params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render Android Messages conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = AndroidMessageBubble(
                text=msg_data.get("text", ""),
                sender=msg_data.get("sender"),
                variant=msg_data.get("variant", "received"),
                timestamp=msg_data.get("timestamp"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            bubble_width = min(width * 0.7, width - 1.0)
            msg_height = message._calculate_bubble_height(bubble_width)
            # Add extra space for timestamp if present
            timestamp_space = 0.25 if msg_data.get("timestamp") else 0
            current_top += msg_height + timestamp_space + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# AndroidMessageBubble - Variants: sent, received
# AndroidConversation - Full conversation flow
