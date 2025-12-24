"""
Code block component for PowerPoint presentations.
Provides syntax-highlighted code display.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from .base import Component
from ..tokens.typography import FONT_FAMILIES, FONT_SIZES
from ..tokens.platform_colors import get_language_color, TERMINAL_COLORS
from ..constants import ComponentSizing


class CodeBlock(Component):
    """
    Code block component with syntax highlighting appearance.

    Features:
    - Language label
    - Monospace font
    - Theme-aware colors
    - Line numbers (optional)
    """

    def __init__(
        self,
        code: str,
        language: str = "text",
        show_line_numbers: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize code block component.

        Args:
            code: Code content
            language: Programming language
            show_line_numbers: Whether to show line numbers
            theme: Theme configuration
        """
        super().__init__(theme)
        self.code = code
        self.language = language.lower()
        self.show_line_numbers = show_line_numbers

    def _format_code(self) -> str:
        """Format code with optional line numbers."""
        lines = self.code.split("\n")

        if self.show_line_numbers:
            max_line_num = len(lines)
            width = len(str(max_line_num))
            formatted_lines = []

            for i, line in enumerate(lines, 1):
                line_num = str(i).rjust(width)
                formatted_lines.append(f"{line_num} │ {line}")

            return "\n".join(formatted_lines)

        return self.code

    def get_language_color(self) -> Any:
        """Get color for language badge."""
        hex_color = get_language_color(self.language)
        return self.hex_to_rgb(hex_color)

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float = 6.0,
        height: float = 3.0,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render code block to slide or replace a placeholder.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            Shape object containing the code
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Create container
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Apply dark background for code
        container.fill.solid()
        if self.get_theme_attr("mode") == "light":
            # Use a dark background even in light mode
            container.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb("#1e1e1e"))
        else:
            container.fill.fore_color.rgb = self.get_color("card.DEFAULT")

        # Add subtle border
        container.line.color.rgb = self.get_color("border.DEFAULT")
        container.line.width = Pt(ComponentSizing.BORDER_WIDTH_THIN)

        # Add language badge (as part of text)
        text_frame = container.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_bottom = Inches(0.15)

        # Add language label
        p = text_frame.paragraphs[0]
        p.text = f"// {self.language}"
        p.font.name = FONT_FAMILIES["mono"][0]  # Use design system mono font
        p.font.size = Pt(FONT_SIZES["xs"])
        p.font.color.rgb = self.get_color("muted.foreground")

        # Add formatted code
        formatted_code = self._format_code()
        for line in formatted_code.split("\n"):
            p = text_frame.add_paragraph()
            p.text = line
            p.font.name = FONT_FAMILIES["mono"][0]  # Use design system mono font
            p.font.size = Pt(FONT_SIZES["sm"])

            # Use light text on dark background
            if self.get_theme_attr("mode") == "light":
                # Use design system foreground instead of hardcoded hex
                p.font.color.rgb = self.get_color("foreground.DEFAULT")
            else:
                p.font.color.rgb = self.get_color("foreground.DEFAULT")

            p.level = 0
            p.space_before = Pt(0)
            p.space_after = Pt(0)

        return container


class InlineCode(Component):
    """
    Inline code component for small code snippets.
    """

    def __init__(self, code: str, theme: Optional[Dict[str, Any]] = None):
        """
        Initialize inline code component.

        Args:
            code: Code snippet
            theme: Theme configuration
        """
        super().__init__(theme)
        self.code = code

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: float = 0.4,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render inline code to slide or replace a placeholder.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Width (auto-calculated if None)
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            Shape object containing the code
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Calculate width based on text length if not provided
        if width is None:
            # Rough estimation: 0.1 inch per character
            width = max(1.0, min(6.0, len(self.code) * 0.08))

        # Create container
        code_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Apply muted background
        code_shape.fill.solid()
        code_shape.fill.fore_color.rgb = self.get_color("muted.DEFAULT")

        # Remove border
        code_shape.line.fill.background()

        # Add code text
        text_frame = code_shape.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)

        p = text_frame.paragraphs[0]
        p.text = self.code
        p.font.name = FONT_FAMILIES["mono"][0]  # Use design system mono font
        p.font.size = Pt(FONT_SIZES["sm"])
        p.font.color.rgb = self.get_color("muted.foreground")
        p.alignment = PP_ALIGN.CENTER

        return code_shape


class Terminal(CodeBlock):
    """
    Terminal/console output component.
    """

    def __init__(self, output: str, prompt: str = "$", theme: Optional[Dict[str, Any]] = None):
        """
        Initialize terminal component.

        Args:
            output: Terminal output text
            prompt: Command prompt character
            theme: Theme configuration
        """
        super().__init__(output, "shell", False, theme)
        self.prompt = prompt

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float = 6.0,
        height: float = 3.0,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """Render terminal output or replace a placeholder."""
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Create container with black background
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Black background for terminal
        container.fill.solid()
        container.fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(TERMINAL_COLORS["background"]))

        # Green border for classic terminal look
        container.line.color.rgb = RGBColor(*self.hex_to_rgb(TERMINAL_COLORS["border"]))
        container.line.width = Pt(ComponentSizing.BORDER_WIDTH_MEDIUM)

        # Add terminal content
        text_frame = container.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.2)
        text_frame.margin_right = Inches(0.2)
        text_frame.margin_top = Inches(0.15)
        text_frame.margin_bottom = Inches(0.15)

        # Add terminal header
        p = text_frame.paragraphs[0]
        p.text = "Terminal"
        p.font.name = FONT_FAMILIES["mono"][0]  # Use design system mono font
        p.font.size = Pt(FONT_SIZES["xs"])
        p.font.color.rgb = RGBColor(*self.hex_to_rgb(TERMINAL_COLORS["text"]))
        p.font.bold = True

        # Add output lines
        for line in self.code.split("\n"):
            p = text_frame.add_paragraph()
            # Add prompt for command lines
            if line and not line.startswith(" "):
                p.text = f"{self.prompt} {line}"
            else:
                p.text = line

            p.font.name = FONT_FAMILIES["mono"][0]  # Use design system mono font
            p.font.size = Pt(FONT_SIZES["sm"])
            p.font.color.rgb = RGBColor(*self.hex_to_rgb(TERMINAL_COLORS["text"]))
            p.level = 0
            p.space_before = Pt(0)
            p.space_after = Pt(0)

        return container
