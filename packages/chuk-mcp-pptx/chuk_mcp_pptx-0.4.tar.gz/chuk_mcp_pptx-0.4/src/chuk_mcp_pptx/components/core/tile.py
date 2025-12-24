# src/chuk_mcp_pptx/components/core/tile.py
"""
Tile components for PowerPoint presentations.

Simple, compact tiles for dashboard layouts. Perfect for icons, status indicators,
single values, or minimal content displays in grid-based layouts.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, PARAGRAPH_SPACING


class Tile(Component):
    """
    Tile component - compact display element for dashboards.

    A simpler alternative to Card, perfect for grid layouts with icons,
    status indicators, or single values. Designed for dashboard UIs.

    Variants:
        - default: Standard tile with subtle background
        - outlined: Border with transparent background
        - filled: Solid background color
        - ghost: Minimal styling

    Sizes:
        - sm: 1.5" x 1.5"
        - md: 2.0" x 2.0" (default)
        - lg: 2.5" x 2.5"
        - xl: 3.0" x 3.0"

    Examples:
        # Simple icon tile
        tile = Tile(icon="rocket", variant="filled", theme=theme)
        tile.render(slide, left=1, top=2)

        # Value tile
        tile = Tile(text="42", label="Tasks", variant="outlined", theme=theme)
        tile.render(slide, left=3, top=2)

        # Status tile
        tile = Tile(text="Active", icon="check", variant="default",
                   color_variant="success", theme=theme)
        tile.render(slide, left=5, top=2)
    """

    # Size mapping in inches (width, height)
    SIZE_MAP = {
        "sm": (1.5, 1.5),
        "md": (2.0, 2.0),
        "lg": (2.5, 2.5),
        "xl": (3.0, 3.0),
    }

    def __init__(
        self,
        text: Optional[str] = None,
        label: Optional[str] = None,
        icon: Optional[str] = None,
        variant: str = "default",
        size: str = "md",
        color_variant: str = "default",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize tile.

        Args:
            text: Main text content (optional)
            label: Secondary label text (optional)
            icon: Icon name (optional)
            variant: Visual variant (default, outlined, filled, ghost)
            size: Tile size (sm, md, lg, xl)
            color_variant: Color variant (default, primary, success, warning, destructive)
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.label = label
        self.icon = icon
        self.variant = variant
        self.size = size
        self.color_variant = color_variant

    def _get_bg_color(self) -> Optional[RGBColor]:
        """Get background color based on variant."""
        if self.variant == "ghost":
            return None
        elif self.variant == "filled":
            if self.color_variant == "primary":
                return self.get_color("primary.DEFAULT")
            elif self.color_variant == "success":
                return self.get_color("success.DEFAULT")
            elif self.color_variant == "warning":
                return self.get_color("warning.DEFAULT")
            elif self.color_variant == "destructive":
                return self.get_color("destructive.DEFAULT")
            else:
                return self.get_color("card.DEFAULT")
        else:  # default, outlined
            return self.get_color("card.DEFAULT")

    def _get_border_color(self) -> RGBColor:
        """Get border color based on variant."""
        if self.color_variant == "primary":
            return self.get_color("primary.DEFAULT")
        elif self.color_variant == "success":
            return self.get_color("success.DEFAULT")
        elif self.color_variant == "warning":
            return self.get_color("warning.DEFAULT")
        elif self.color_variant == "destructive":
            return self.get_color("destructive.DEFAULT")
        else:
            return self.get_color("border.DEFAULT")

    def _get_text_color(self) -> RGBColor:
        """Get text color based on variant."""
        if self.variant == "filled":
            # Use white/light text on filled backgrounds
            if self.color_variant in ["primary", "success", "destructive"]:
                return self.get_color("primary.foreground")
            else:
                return self.get_color("card.foreground")
        else:
            # Use theme foreground for other variants
            if self.color_variant == "primary":
                return self.get_color("primary.DEFAULT")
            elif self.color_variant == "success":
                return self.get_color("success.DEFAULT")
            elif self.color_variant == "warning":
                return self.get_color("warning.DEFAULT")
            elif self.color_variant == "destructive":
                return self.get_color("destructive.DEFAULT")
            else:
                return self.get_color("card.foreground")

    def _get_label_color(self) -> RGBColor:
        """Get label color."""
        if self.variant == "filled" and self.color_variant in ["primary", "success", "destructive"]:
            return self.get_color("primary.foreground")
        else:
            return self.get_color("muted.foreground")

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render tile to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Tile width in inches (optional, uses size if not provided)
            height: Tile height in inches (optional, uses size if not provided)
            placeholder: Optional placeholder to replace

        Returns:
            Shape object representing the tile
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Get dimensions from size if not provided
        default_width, default_height = self.SIZE_MAP.get(self.size, (2.0, 2.0))
        tile_width = width if width is not None else default_width
        tile_height = height if height is not None else default_height

        # Create tile shape
        tile = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(tile_width),
            Inches(tile_height),
        )

        # Apply background
        bg_color = self._get_bg_color()
        if bg_color:
            tile.fill.solid()
            tile.fill.fore_color.rgb = bg_color
        else:
            tile.fill.background()

        # Apply border
        if self.variant in ["outlined", "default"]:
            border_width = 1.5 if self.variant == "outlined" else 0.5
            tile.line.color.rgb = self._get_border_color()
            tile.line.width = Pt(border_width)
        else:
            tile.line.fill.background()

        # Setup text frame
        text_frame = tile.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Minimal margins for compact look
        padding = 0.15
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        # Render content based on what's provided
        font_family = self._get_font_family()
        current_p = text_frame.paragraphs[0]
        has_content = False

        # Icon (if provided)
        if self.icon:
            from .icon import ICON_SYMBOLS

            symbol = ICON_SYMBOLS.get(self.icon, self.icon)
            current_p.text = symbol
            current_p.alignment = PP_ALIGN.CENTER
            current_p.font.name = font_family
            current_p.font.size = Pt(24 if self.size in ["lg", "xl"] else 20)
            current_p.font.color.rgb = self._get_text_color()
            has_content = True

        # Main text (if provided)
        if self.text:
            if has_content:
                current_p = text_frame.add_paragraph()
                current_p.space_before = Pt(PARAGRAPH_SPACING["xs"])

            current_p.text = self.text
            current_p.alignment = PP_ALIGN.CENTER

            # Size based on tile size
            if self.size == "sm":
                font_size = 14
            elif self.size in ["lg", "xl"]:
                font_size = 20
            else:
                font_size = 16

            current_p.font.name = font_family
            current_p.font.size = Pt(font_size)
            current_p.font.bold = True
            current_p.font.color.rgb = self._get_text_color()
            has_content = True

        # Label (if provided)
        if self.label:
            if has_content:
                current_p = text_frame.add_paragraph()
                current_p.space_before = Pt(PARAGRAPH_SPACING["xs"])

            current_p.text = self.label
            current_p.alignment = PP_ALIGN.CENTER
            current_p.font.name = font_family
            current_p.font.size = Pt(FONT_SIZES["xs"])
            current_p.font.color.rgb = self._get_label_color()

        return tile


class IconTile(Tile):
    """
    Icon tile - specialized tile focused on icon display.

    Perfect for feature indicators, navigation elements, or status displays.

    Examples:
        # Simple icon tile
        tile = IconTile("rocket", label="Fast", theme=theme)
        tile.render(slide, left=1, top=2)

        # Colored icon tile
        tile = IconTile("check", label="Complete",
                       variant="filled", color_variant="success", theme=theme)
        tile.render(slide, left=3, top=2)
    """

    def __init__(
        self,
        icon: str,
        label: Optional[str] = None,
        variant: str = "default",
        size: str = "md",
        color_variant: str = "default",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize icon tile.

        Args:
            icon: Icon name
            label: Optional label text
            variant: Visual variant
            size: Tile size
            color_variant: Color variant
            theme: Optional theme
        """
        super().__init__(
            text=None,
            label=label,
            icon=icon,
            variant=variant,
            size=size,
            color_variant=color_variant,
            theme=theme,
        )


class ValueTile(Tile):
    """
    Value tile - specialized tile for displaying single values/metrics.

    Simpler than MetricCard, perfect for compact dashboard displays.

    Examples:
        # Simple value
        tile = ValueTile("42", label="Tasks", theme=theme)
        tile.render(slide, left=1, top=2)

        # Highlighted value
        tile = ValueTile("98%", label="Uptime",
                        variant="filled", color_variant="success", theme=theme)
        tile.render(slide, left=3, top=2)
    """

    def __init__(
        self,
        value: str,
        label: Optional[str] = None,
        variant: str = "default",
        size: str = "md",
        color_variant: str = "default",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize value tile.

        Args:
            value: Value to display
            label: Optional label text
            variant: Visual variant
            size: Tile size
            color_variant: Color variant
            theme: Optional theme
        """
        super().__init__(
            text=value,
            label=label,
            icon=None,
            variant=variant,
            size=size,
            color_variant=color_variant,
            theme=theme,
        )


# TODO: Register component when registry is implemented
# Component metadata for documentation:
# Tile - Variants: default, outlined, filled, ghost
# Tile - Sizes: sm, md, lg, xl
# Tile - Color variants: default, primary, success, warning, destructive
# IconTile - Specialized tile for icons
# ValueTile - Specialized tile for values/metrics
