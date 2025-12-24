# src/chuk_mcp_pptx/components/core/avatar.py
"""
Avatar components for PowerPoint presentations.

Provides avatar displays for user profiles, team members, and contact information.
Uses initials, icons, or colored circles for visual representation.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import get_text_style


class Avatar(Component):
    """
    Avatar component - visual representation of users/profiles.

    Perfect for team slides, contact information, author attribution,
    and any user-related content in presentations.

    Variants:
        - default: Standard avatar with border
        - filled: Solid background color
        - outlined: Border only
        - minimal: No border or background

    Sizes:
        - xs: 0.5" diameter
        - sm: 0.75" diameter
        - md: 1.0" diameter (default)
        - lg: 1.5" diameter
        - xl: 2.0" diameter

    Display modes:
        - initials: Shows user initials (e.g., "JD" for John Doe)
        - icon: Shows an icon (e.g., user icon)
        - empty: Shows empty circle (for placeholders)

    Examples:
        # Avatar with initials
        avatar = Avatar(text="JD", theme=theme)
        avatar.render(slide, left=1, top=2)

        # Large filled avatar
        avatar = Avatar(text="AS", variant="filled", size="lg",
                       color_variant="primary", theme=theme)
        avatar.render(slide, left=3, top=2)

        # Icon avatar
        avatar = Avatar(icon="user", variant="outlined", theme=theme)
        avatar.render(slide, left=5, top=2)

        # Avatar with label
        avatar = AvatarWithLabel(text="JD", label="John Doe",
                                sublabel="Designer", theme=theme)
        avatar.render(slide, left=1, top=2)
    """

    # Size mapping in inches (diameter)
    SIZE_MAP = {
        "xs": 0.5,
        "sm": 0.75,
        "md": 1.0,
        "lg": 1.5,
        "xl": 2.0,
    }

    # Font size mapping for initials
    FONT_SIZE_MAP = {
        "xs": 8,
        "sm": 10,
        "md": 14,
        "lg": 20,
        "xl": 28,
    }

    def __init__(
        self,
        text: Optional[str] = None,
        icon: Optional[str] = None,
        variant: str = "default",
        size: str = "md",
        color_variant: str = "default",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize avatar.

        Args:
            text: Initials or text to display (optional)
            icon: Icon name (optional, overrides text)
            variant: Visual variant (default, filled, outlined, minimal)
            size: Avatar size (xs, sm, md, lg, xl)
            color_variant: Color variant (default, primary, success, warning, destructive)
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.icon = icon
        self.variant = variant
        self.size = size
        self.color_variant = color_variant

    def _get_bg_color(self) -> Optional[RGBColor]:
        """Get background color based on variant."""
        if self.variant == "minimal":
            return None
        elif self.variant == "filled":
            if self.color_variant == "primary":
                return self.get_color("primary.DEFAULT")
            elif self.color_variant == "success":
                return self.get_color("success.DEFAULT")
            elif self.color_variant == "warning":
                return self.get_color("warning.DEFAULT")
            elif self.color_variant == "destructive":
                return self.get_color("destructive.DEFAULT")
            else:
                return self.get_color("muted.DEFAULT")
        else:  # default, outlined
            return self.get_color("card.DEFAULT")

    def _get_border_color(self) -> RGBColor:
        """Get border color based on variant."""
        if self.color_variant == "primary":
            return self.get_color("primary.DEFAULT")
        elif self.color_variant == "success":
            return self.get_color("success.DEFAULT")
        elif self.color_variant == "warning":
            return self.get_color("warning.DEFAULT")
        elif self.color_variant == "destructive":
            return self.get_color("destructive.DEFAULT")
        else:
            return self.get_color("border.DEFAULT")

    def _get_text_color(self) -> RGBColor:
        """Get text color based on variant."""
        if self.variant == "filled":
            # Use white/light text on filled backgrounds
            if self.color_variant in ["primary", "success", "destructive"]:
                return self.get_color("primary.foreground")
            else:
                return self.get_color("foreground.DEFAULT")
        else:
            # Use theme colors for other variants
            if self.color_variant == "primary":
                return self.get_color("primary.DEFAULT")
            elif self.color_variant == "success":
                return self.get_color("success.DEFAULT")
            elif self.color_variant == "warning":
                return self.get_color("warning.DEFAULT")
            elif self.color_variant == "destructive":
                return self.get_color("destructive.DEFAULT")
            else:
                return self.get_color("foreground.DEFAULT")

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self,
        slide,
        left: float,
        top: float,
        diameter: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render avatar to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            diameter: Avatar diameter in inches (optional, uses size if not provided)
            placeholder: Optional placeholder to replace

        Returns:
            Shape object representing the avatar
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            diameter = min(width, height)  # Use smaller dimension for circular avatar

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Get diameter from size if not provided
        avatar_diameter = diameter if diameter is not None else self.SIZE_MAP.get(self.size, 1.0)

        # Create circular avatar shape
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left),
            Inches(top),
            Inches(avatar_diameter),
            Inches(avatar_diameter),
        )

        # Apply background
        bg_color = self._get_bg_color()
        if bg_color:
            avatar.fill.solid()
            avatar.fill.fore_color.rgb = bg_color
        else:
            avatar.fill.background()

        # Apply border
        if self.variant in ["outlined", "default"]:
            border_width = 2 if self.variant == "outlined" else 1
            avatar.line.color.rgb = self._get_border_color()
            avatar.line.width = Pt(border_width)
        else:
            avatar.line.fill.background()

        # Setup text frame for content
        text_frame = avatar.text_frame
        text_frame.clear()
        text_frame.word_wrap = False
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # No margins - center content
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0

        font_family = self._get_font_family()
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER

        # Render icon or text
        if self.icon:
            from .icon import ICON_SYMBOLS

            symbol = ICON_SYMBOLS.get(self.icon, self.icon)
            p.text = symbol
            p.font.name = font_family
            p.font.size = Pt(self.FONT_SIZE_MAP.get(self.size, 14))
        elif self.text:
            p.text = self.text.upper()[:2]  # Max 2 characters for initials
            p.font.name = font_family
            p.font.size = Pt(self.FONT_SIZE_MAP.get(self.size, 14))
            p.font.bold = True

        p.font.color.rgb = self._get_text_color()

        return avatar


class AvatarWithLabel(Component):
    """
    Avatar with label component - avatar + name/title display.

    Combines an avatar with text labels for complete profile display.
    Perfect for team slides, speaker introductions, and contact cards.

    Examples:
        # Simple avatar with name
        avatar = AvatarWithLabel(text="JD", label="John Doe", theme=theme)
        avatar.render(slide, left=1, top=2)

        # Avatar with name and title
        avatar = AvatarWithLabel(
            text="AS",
            label="Alice Smith",
            sublabel="Product Designer",
            variant="filled",
            color_variant="primary",
            theme=theme
        )
        avatar.render(slide, left=1, top=2)
    """

    def __init__(
        self,
        text: Optional[str] = None,
        icon: Optional[str] = None,
        label: Optional[str] = None,
        sublabel: Optional[str] = None,
        variant: str = "default",
        size: str = "md",
        color_variant: str = "default",
        orientation: str = "horizontal",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize avatar with label.

        Args:
            text: Initials or text to display (optional)
            icon: Icon name (optional)
            label: Primary label text (e.g., name)
            sublabel: Secondary label text (e.g., title/role)
            variant: Visual variant
            size: Avatar size
            color_variant: Color variant
            orientation: Layout orientation (horizontal, vertical)
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.icon = icon
        self.label = label
        self.sublabel = sublabel
        self.variant = variant
        self.size = size
        self.color_variant = color_variant
        self.orientation = orientation

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self, slide, left: float, top: float, width: float = 3.0, placeholder: Optional[Any] = None
    ) -> list:
        """
        Render avatar with label to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Total width for the component in inches
            placeholder: Optional placeholder to replace

        Returns:
            List of rendered shapes [avatar, label_box]
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        shapes = []

        # Create avatar
        avatar_component = Avatar(
            text=self.text,
            icon=self.icon,
            variant=self.variant,
            size=self.size,
            color_variant=self.color_variant,
            theme=self.theme,
        )
        avatar_shape = avatar_component.render(slide, left, top)
        shapes.append(avatar_shape)

        # Get avatar diameter for positioning
        avatar_diameter = Avatar.SIZE_MAP.get(self.size, 1.0)

        if self.orientation == "vertical":
            # Vertical layout: avatar on top, labels below
            label_left = left
            label_top = top + avatar_diameter + 0.1
            label_width = width
            text_align = PP_ALIGN.CENTER
        else:
            # Horizontal layout: avatar on left, labels on right
            label_left = left + avatar_diameter + 0.15
            label_top = top
            label_width = width - avatar_diameter - 0.15
            text_align = PP_ALIGN.LEFT

        # Create label text box
        if self.label or self.sublabel:
            font_family = self._get_font_family()
            label_height = 0.5 if self.sublabel else 0.3

            label_box = slide.shapes.add_textbox(
                Inches(label_left), Inches(label_top), Inches(label_width), Inches(label_height)
            )

            text_frame = label_box.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            text_frame.vertical_anchor = (
                MSO_ANCHOR.TOP if self.orientation == "vertical" else MSO_ANCHOR.MIDDLE
            )
            text_frame.margin_left = 0
            text_frame.margin_top = 0

            # Primary label
            if self.label:
                p = text_frame.paragraphs[0]
                p.text = self.label
                p.alignment = text_align
                style = get_text_style("body")
                p.font.name = font_family
                p.font.size = Pt(style["font_size"])
                p.font.bold = True
                p.font.color.rgb = self.get_color("foreground.DEFAULT")

            # Sublabel
            if self.sublabel:
                p = text_frame.add_paragraph()
                p.text = self.sublabel
                p.alignment = text_align
                p.space_before = Pt(2)
                style = get_text_style("small")
                p.font.name = font_family
                p.font.size = Pt(style["font_size"])
                p.font.color.rgb = self.get_color("muted.foreground")

            shapes.append(label_box)

        return shapes


class AvatarGroup(Component):
    """
    Avatar group component - display multiple avatars in a row.

    Shows multiple avatars with optional overlap for compact team display.
    Perfect for showing contributors, team members, or participants.

    Examples:
        # Simple avatar group
        members = [
            {"text": "JD"},
            {"text": "AS"},
            {"text": "BM"}
        ]
        group = AvatarGroup(members, theme=theme)
        group.render(slide, left=1, top=2)

        # Overlapping avatars with max display
        members = [
            {"text": "JD", "color_variant": "primary"},
            {"text": "AS", "color_variant": "success"},
            {"text": "BM", "color_variant": "warning"},
            {"text": "KL", "color_variant": "destructive"},
            {"text": "MN", "color_variant": "default"}
        ]
        group = AvatarGroup(members, max_display=3, overlap=True, theme=theme)
        group.render(slide, left=1, top=2)
    """

    def __init__(
        self,
        members: list,
        size: str = "md",
        variant: str = "default",
        max_display: Optional[int] = None,
        overlap: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize avatar group.

        Args:
            members: List of member dicts with keys: text/icon, color_variant (optional)
            size: Avatar size
            variant: Visual variant
            max_display: Maximum number to display (shows "+N" for extras)
            overlap: Whether to overlap avatars
            theme: Optional theme
        """
        super().__init__(theme)
        self.members = members
        self.size = size
        self.variant = variant
        self.max_display = max_display
        self.overlap = overlap

    def render(self, slide, left: float, top: float, placeholder: Optional[Any] = None) -> list:
        """
        Render avatar group to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            placeholder: Optional placeholder to replace

        Returns:
            List of rendered avatar shapes
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        shapes = []
        avatar_diameter = Avatar.SIZE_MAP.get(self.size, 1.0)

        # Calculate spacing
        if self.overlap:
            spacing = avatar_diameter * 0.7  # 30% overlap
        else:
            spacing = avatar_diameter + 0.1  # Small gap

        # Determine how many to show
        total_members = len(self.members)
        if self.max_display and total_members > self.max_display:
            display_count = self.max_display
            remaining = total_members - self.max_display
        else:
            display_count = total_members
            remaining = 0

        # Render avatars
        for i in range(display_count):
            member = self.members[i]
            avatar_left = left + (i * spacing)

            avatar = Avatar(
                text=member.get("text"),
                icon=member.get("icon"),
                variant=self.variant,
                size=self.size,
                color_variant=member.get("color_variant", "default"),
                theme=self.theme,
            )
            avatar_shape = avatar.render(slide, avatar_left, top)
            shapes.append(avatar_shape)

        # Show "+N" for remaining
        if remaining > 0:
            remaining_left = left + (display_count * spacing)
            remaining_avatar = Avatar(
                text=f"+{remaining}",
                variant=self.variant,
                size=self.size,
                color_variant="default",
                theme=self.theme,
            )
            remaining_shape = remaining_avatar.render(slide, remaining_left, top)
            shapes.append(remaining_shape)

        return shapes


# TODO: Register component when registry is implemented
# Component metadata for documentation:
# Avatar - Variants: default, filled, outlined, minimal
# Avatar - Sizes: xs, sm, md, lg, xl
# Avatar - Color variants: default, primary, success, warning, destructive
# AvatarWithLabel - Orientations: horizontal, vertical
# AvatarGroup - Props: members, max_display, overlap
