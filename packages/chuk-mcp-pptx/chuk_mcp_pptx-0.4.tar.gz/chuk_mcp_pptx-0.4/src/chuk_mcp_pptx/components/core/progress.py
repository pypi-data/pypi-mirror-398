# src/chuk_mcp_pptx/components/core/progress.py
"""
Progress bar components for PowerPoint presentations.

Provides visual progress indicators for showing completion status,
KPI achievement, and task progress.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import get_text_style


class ProgressBar(Component):
    """
    Progress bar component - visual progress indicator.

    Perfect for showing completion status, KPI achievement, task progress.
    Supports linear and segmented styles with optional labels.

    Variants:
        - default: Standard progress bar with primary color
        - success: Green success color
        - warning: Yellow/orange warning color
        - error: Red error color

    Styles:
        - linear: Smooth linear progress bar (default)
        - segmented: Divided into segments

    Examples:
        # Simple progress bar
        progress = ProgressBar(value=65, theme=theme)
        progress.render(slide, left=1, top=2, width=6)

        # With label and percentage
        progress = ProgressBar(
            value=85,
            label="Project Completion",
            show_percentage=True,
            variant="success",
            theme=theme
        )
        progress.render(slide, left=1, top=2, width=6)

        # Segmented progress
        progress = ProgressBar(
            value=60,
            segments=10,
            style="segmented",
            theme=theme
        )
        progress.render(slide, left=1, top=2, width=6)
    """

    def __init__(
        self,
        value: float,
        label: Optional[str] = None,
        show_percentage: bool = False,
        variant: str = "default",
        style: str = "linear",
        segments: int = 1,
        height: float = 0.3,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize progress bar.

        Args:
            value: Progress value (0-100)
            label: Optional label text
            show_percentage: Show percentage value
            variant: Color variant (default, success, warning, error)
            style: Bar style (linear, segmented)
            segments: Number of segments for segmented style
            height: Bar height in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.value = max(0, min(100, value))  # Clamp between 0-100
        self.label = label
        self.show_percentage = show_percentage
        self.variant = variant
        self.style = style
        self.segments = max(1, segments)
        self.height = height

    def _get_progress_color(self) -> RGBColor:
        """Get color based on variant."""
        color_map = {
            "default": "primary.DEFAULT",
            "success": "success.DEFAULT",
            "warning": "warning.DEFAULT",
            "error": "destructive.DEFAULT",
        }
        color_path = color_map.get(self.variant, "primary.DEFAULT")
        return self.get_color(color_path)

    def _get_background_color(self) -> RGBColor:
        """Get background color for unfilled portion."""
        return self.get_color("muted.DEFAULT")

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self, slide, left: float, top: float, width: float = 6.0, placeholder: Optional[Any] = None
    ) -> Any:
        """
        Render progress bar to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Total width in inches
            placeholder: Optional placeholder to replace

        Returns:
            List of rendered shapes
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        from pptx.util import Inches

        shapes = []
        current_top = top

        # Render label if provided
        font_family = self._get_font_family()
        if self.label:
            label_box = slide.shapes.add_textbox(
                Inches(left), Inches(current_top), Inches(width), Inches(0.3)
            )
            label_frame = label_box.text_frame
            label_frame.text = self.label
            label_frame.margin_top = 0
            label_frame.margin_bottom = Inches(0.05)

            p = label_frame.paragraphs[0]
            style = get_text_style("small")
            p.font.name = font_family
            p.font.size = Pt(style["font_size"])
            p.font.color.rgb = self.get_color("foreground.DEFAULT")
            p.font.bold = style.get("font_weight", 400) >= 600

            shapes.append(label_box)
            current_top += 0.35

        # Render progress bar
        if self.style == "segmented":
            shapes.extend(self._render_segmented(slide, left, current_top, width))
        else:
            shapes.extend(self._render_linear(slide, left, current_top, width))

        # Render percentage if requested
        if self.show_percentage:
            pct_top = current_top + self.height + 0.05
            pct_box = slide.shapes.add_textbox(
                Inches(left + width - 0.8), Inches(pct_top), Inches(0.8), Inches(0.25)
            )
            pct_frame = pct_box.text_frame
            pct_frame.text = f"{self.value:.0f}%"
            pct_frame.margin_top = 0

            p = pct_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            style = get_text_style("small")
            p.font.name = font_family
            p.font.size = Pt(style["font_size"])
            p.font.color.rgb = self.get_color("muted.foreground")
            p.font.bold = True

            shapes.append(pct_box)

        return shapes

    def _render_linear(self, slide, left: float, top: float, width: float):
        """Render linear progress bar."""
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE

        shapes = []

        # Background bar (unfilled)
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(self.height),
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = self._get_background_color()
        bg_shape.line.color.rgb = self._get_background_color()
        shapes.append(bg_shape)

        # Filled portion
        if self.value > 0:
            filled_width = width * (self.value / 100)
            fill_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(filled_width),
                Inches(self.height),
            )
            fill_shape.fill.solid()
            fill_shape.fill.fore_color.rgb = self._get_progress_color()
            fill_shape.line.color.rgb = self._get_progress_color()
            shapes.append(fill_shape)

        return shapes

    def _render_segmented(self, slide, left: float, top: float, width: float):
        """Render segmented progress bar."""
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE

        shapes = []
        gap = 0.05
        segment_width = (width - (gap * (self.segments - 1))) / self.segments
        filled_segments = int((self.value / 100) * self.segments)

        for i in range(self.segments):
            seg_left = left + (i * (segment_width + gap))
            is_filled = i < filled_segments

            seg_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(seg_left),
                Inches(top),
                Inches(segment_width),
                Inches(self.height),
            )
            seg_shape.fill.solid()

            if is_filled:
                seg_shape.fill.fore_color.rgb = self._get_progress_color()
                seg_shape.line.color.rgb = self._get_progress_color()
            else:
                seg_shape.fill.fore_color.rgb = self._get_background_color()
                seg_shape.line.color.rgb = self._get_background_color()

            shapes.append(seg_shape)

        return shapes


# TODO: Register component when registry is implemented
# Component metadata for documentation:
# - Variants: default, success, warning, error
# - Styles: linear, segmented
# - Props: value, label, show_percentage, variant, style, segments, height
