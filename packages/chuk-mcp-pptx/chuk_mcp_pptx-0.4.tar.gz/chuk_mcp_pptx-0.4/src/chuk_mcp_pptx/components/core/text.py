# src/chuk_mcp_pptx/components/core/text.py
"""
Text components for PowerPoint presentations.

Provides text box and bullet list components with formatting.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

from ..base import Component
from ..registry import component, ComponentCategory, prop, example


@component(
    name="TextBox",
    category=ComponentCategory.UI,
    description="Text box component with formatting and styling",
    props=[
        prop("text", "string", "Text content", required=True),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("width", "number", "Width in inches", required=True),
        prop("height", "number", "Height in inches", required=True),
        prop("font_name", "string", "Font family name", default="Calibri"),
        prop("font_size", "number", "Font size in points", default=18),
        prop("bold", "boolean", "Bold text", default=False),
        prop("italic", "boolean", "Italic text", default=False),
        prop("color", "string", "Text color (semantic or hex)", default=None),
        prop(
            "alignment", "string", "Text alignment (left, center, right, justify)", default="left"
        ),
        prop("auto_fit", "boolean", "Auto-fit text to shape", default=False),
    ],
    examples=[
        example(
            "Simple text box",
            """
text = TextBox(text="Hello World")
text.render(slide, left=2, top=2, width=4, height=1)
            """,
            text="Hello World",
        ),
        example(
            "Formatted text box",
            """
text = TextBox(
    text="Important Message",
    font_size=24,
    bold=True,
    color="primary.DEFAULT",
    alignment="center"
)
text.render(slide, left=1, top=3, width=8, height=1.5)
            """,
            text="Important Message",
            bold=True,
            alignment="center",
        ),
    ],
    tags=["text", "textbox", "label", "heading"],
)
class TextBox(Component):
    """
    Text box component for adding formatted text to slides.

    Features:
    - Custom font family, size, and styling
    - Text alignment (left, center, right, justify)
    - Semantic or hex color support
    - Auto-fit text option
    - Word wrapping

    Usage:
        # Simple text
        text = TextBox(text="Hello World")
        text.render(slide, left=2, top=2, width=4, height=1)

        # Styled text
        text = TextBox(
            text="Title",
            font_size=32,
            bold=True,
            color="primary.DEFAULT",
            alignment="center"
        )
        text.render(slide, left=1, top=1, width=8, height=1.5)
    """

    def __init__(
        self,
        text: str,
        font_name: str = "Calibri",
        font_size: int = 18,
        bold: bool = False,
        italic: bool = False,
        color: Optional[str] = None,
        alignment: str = "left",
        auto_fit: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize text box component.

        Args:
            text: Text content
            font_name: Font family name
            font_size: Font size in points
            bold: Bold text
            italic: Italic text
            color: Text color (semantic color like "primary.DEFAULT" or hex like "#FF0000")
            alignment: Text alignment (left, center, right, justify)
            auto_fit: Auto-fit text to shape
            theme: Optional theme override
        """
        super().__init__(theme)
        self.text = text
        self.font_name = font_name
        self.font_size = font_size
        self.bold = bold
        self.italic = italic
        self.color = color
        self.alignment = alignment
        self.auto_fit = auto_fit

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render text box to slide or populate a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to populate

        Returns:
            Text box shape object
        """
        # Use placeholder if provided, otherwise create new textbox
        if placeholder is not None and hasattr(placeholder, "text_frame"):
            text_box = placeholder
        else:
            # Delete placeholder if it exists but doesn't have text_frame
            self._delete_placeholder_if_needed(placeholder)
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )

        text_frame = text_box.text_frame
        text_frame.text = self.text
        text_frame.word_wrap = True

        # Get font family from theme if not explicitly set
        font_family = self._get_font_family() if self.font_name == "Calibri" else self.font_name

        # Format text
        for paragraph in text_frame.paragraphs:
            # Set alignment
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY,
            }
            paragraph.alignment = alignment_map.get(self.alignment.lower(), PP_ALIGN.LEFT)

            # Format font
            font = paragraph.font
            font.name = font_family
            font.size = Pt(self.font_size)
            font.bold = self.bold
            font.italic = self.italic

            # Apply color
            if self.color:
                rgb = self._parse_color(self.color)
                if rgb:
                    font.color.rgb = rgb
            elif self.theme:
                # Default to theme foreground color if no color specified
                try:
                    rgb = self.theme.get_color("foreground.DEFAULT")
                    if rgb:
                        font.color.rgb = rgb
                except (AttributeError, KeyError, TypeError, ValueError):
                    pass

        # Apply auto-fit if requested
        if self.auto_fit:
            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.05)
            text_frame.margin_bottom = Inches(0.05)

        return text_box

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        """Parse color string (semantic or hex) to RGBColor."""
        # Handle hex colors
        if color_str.startswith("#"):
            hex_color = color_str.lstrip("#")
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)

        # Handle semantic colors
        if self.theme and "." in color_str:
            try:
                color_rgb = self.theme.get_color(color_str)
                return color_rgb
            except (AttributeError, KeyError, TypeError, ValueError):
                pass

        return None


@component(
    name="BulletList",
    category=ComponentCategory.UI,
    description="Bullet list component for presenting items",
    props=[
        prop("items", "array", "List items", required=True),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("width", "number", "Width in inches", required=True),
        prop("height", "number", "Height in inches", required=True),
        prop("font_size", "number", "Font size in points", default=16),
        prop("color", "string", "Text color (semantic or hex)", default=None),
        prop("bullet_char", "string", "Bullet character", default="•"),
        prop("spacing", "number", "Space after each item in points", default=6),
    ],
    examples=[
        example(
            "Simple bullet list",
            """
bullets = BulletList(items=["First item", "Second item", "Third item"])
bullets.render(slide, left=1, top=2, width=8, height=4)
            """,
            items=["First", "Second", "Third"],
        ),
        example(
            "Styled bullet list",
            """
bullets = BulletList(
    items=["Increase revenue", "Reduce costs", "Improve quality"],
    font_size=18,
    color="primary.DEFAULT",
    bullet_char="→"
)
bullets.render(slide, left=1, top=2, width=8, height=3)
            """,
            items=["Revenue", "Costs", "Quality"],
            bullet_char="→",
        ),
    ],
    tags=["bullet", "list", "items", "enumeration"],
)
class BulletList(Component):
    """
    Bullet list component for presenting items.

    Features:
    - Custom bullet characters
    - Font size and color control
    - Adjustable item spacing
    - Semantic color support
    - Word wrapping

    Usage:
        # Simple list
        bullets = BulletList(items=["Item 1", "Item 2", "Item 3"])
        bullets.render(slide, left=1, top=2, width=8, height=4)

        # Styled list
        bullets = BulletList(
            items=["Revenue", "Costs", "Quality"],
            font_size=18,
            color="primary.DEFAULT",
            bullet_char="✓"
        )
        bullets.render(slide, left=1, top=2, width=8, height=3)
    """

    def __init__(
        self,
        items: List[str],
        font_size: int = 16,
        color: Optional[str] = None,
        bullet_char: str = "•",
        spacing: int = 6,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize bullet list component.

        Args:
            items: List of items to display
            font_size: Font size in points
            color: Text color (semantic or hex)
            bullet_char: Character to use for bullets
            spacing: Space after each item in points
            theme: Optional theme override
        """
        super().__init__(theme)
        self.items = items
        self.font_size = font_size
        self.color = color
        self.bullet_char = bullet_char
        self.spacing = spacing

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render bullet list to slide or populate a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to populate

        Returns:
            Text box shape object
        """
        # Use placeholder if provided, otherwise create new textbox
        if placeholder is not None and hasattr(placeholder, "text_frame"):
            text_box = placeholder
        else:
            # Delete placeholder if it exists but doesn't have text_frame
            self._delete_placeholder_if_needed(placeholder)
            text_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )

        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        font_family = self._get_font_family()

        # Add items
        for idx, item in enumerate(self.items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = f"{self.bullet_char} {item}"
            p.font.name = font_family
            p.font.size = Pt(self.font_size)
            p.space_after = Pt(self.spacing)

            # Apply color
            if self.color:
                rgb = self._parse_color(self.color)
                if rgb:
                    p.font.color.rgb = rgb
            elif self.theme:
                # Default to theme foreground color if no color specified
                try:
                    rgb = self.theme.get_color("foreground.DEFAULT")
                    if rgb:
                        p.font.color.rgb = rgb
                except (AttributeError, KeyError, TypeError, ValueError):
                    pass

        return text_box

    def _parse_color(self, color_str: str) -> Optional[RGBColor]:
        """Parse color string (semantic or hex) to RGBColor."""
        # Handle hex colors
        if color_str.startswith("#"):
            hex_color = color_str.lstrip("#")
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            return RGBColor(r, g, b)

        # Handle semantic colors
        if self.theme and "." in color_str:
            try:
                color_rgb = self.theme.get_color(color_str)
                return color_rgb
            except (AttributeError, KeyError, TypeError, ValueError):
                pass

        return None
