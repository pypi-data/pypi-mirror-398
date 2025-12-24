# src/chuk_mcp_pptx/components/core/table.py
"""
Table component with variants and theme support.
"""

from typing import Optional, Any, List
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from ..composition import ComposableComponent
from ..variants import create_variants
from ..registry import component, ComponentCategory, prop, example


# Table variants configuration
TABLE_VARIANTS = create_variants(
    base={
        "header_bold": True,
    },
    variants={
        "variant": {
            "default": {
                "header_bg": "card.DEFAULT",
                "header_fg": "card.foreground",
                "header_bold": True,
                "cell_bg": "background.DEFAULT",
                "cell_fg": "foreground.DEFAULT",
                "border_color": "border.DEFAULT",
                "border_width": 1,
            },
            "bordered": {
                "header_bg": "primary.DEFAULT",
                "header_fg": "primary.foreground",
                "header_bold": True,
                "cell_bg": "background.DEFAULT",
                "cell_fg": "foreground.DEFAULT",
                "border_color": "border.DEFAULT",
                "border_width": 2,
            },
            "striped": {
                "header_bg": "muted.DEFAULT",
                "header_fg": "muted.foreground",
                "header_bold": True,
                "cell_bg": "background.DEFAULT",
                "cell_fg": "foreground.DEFAULT",
                "alt_bg": "muted.DEFAULT",  # Alternating row color
                "border_color": "border.secondary",
                "border_width": 0.5,
            },
            "minimal": {
                "header_bg": "transparent",
                "header_fg": "foreground.DEFAULT",
                "header_bold": True,
                "cell_bg": "transparent",
                "cell_fg": "foreground.DEFAULT",
                "border_color": "border.secondary",
                "border_width": 0.5,
            },
        },
        "size": {
            "sm": {
                "font_size": 10,
                "header_font_size": 11,
                "padding": 0.05,
            },
            "md": {
                "font_size": 12,
                "header_font_size": 13,
                "padding": 0.08,
            },
            "lg": {
                "font_size": 14,
                "header_font_size": 16,
                "padding": 0.12,
            },
        },
    },
    default_variants={"variant": "default", "size": "md"},
)


@component(
    name="Table",
    category=ComponentCategory.DATA,
    description="Data table with headers and rows, supporting multiple variants and theme-aware styling",
    props=[
        prop("headers", "list", "Column header labels", required=True, example=["Name", "Value"]),
        prop(
            "data",
            "list",
            "Table data as list of rows",
            required=True,
            example=[["Item 1", "100"], ["Item 2", "200"]],
        ),
        prop(
            "variant",
            "string",
            "Visual variant",
            options=["default", "bordered", "striped", "minimal"],
            default="default",
            example="default",
        ),
        prop(
            "size", "string", "Table size", options=["sm", "md", "lg"], default="md", example="md"
        ),
        prop("left", "number", "Left position in inches", required=True, example=1.0),
        prop("top", "number", "Top position in inches", required=True, example=2.0),
        prop("width", "number", "Width in inches", required=True, example=8.0),
        prop("height", "number", "Height in inches", required=True, example=4.0),
    ],
    variants={"variant": ["default", "bordered", "striped", "minimal"], "size": ["sm", "md", "lg"]},
    examples=[
        example(
            "Basic data table",
            """
table = Table(
    headers=["Product", "Q1", "Q2", "Q3"],
    data=[
        ["Laptops", "$100K", "$120K", "$110K"],
        ["Phones", "$80K", "$90K", "$95K"]
    ],
    variant="default",
    size="md"
)
table.render(slide, left=1, top=2, width=8, height=3)
            """,
            headers=["Product", "Q1", "Q2", "Q3"],
            data=[["Laptops", "$100K", "$120K", "$110K"], ["Phones", "$80K", "$90K", "$95K"]],
            variant="default",
            size="md",
        ),
        example(
            "Striped comparison table",
            """
table = Table(
    headers=["Feature", "Basic", "Pro", "Enterprise"],
    data=[
        ["Storage", "10GB", "100GB", "Unlimited"],
        ["Users", "1", "10", "Unlimited"],
        ["Support", "Email", "Priority", "24/7"]
    ],
    variant="striped",
    size="lg"
)
table.render(slide, left=1, top=2, width=8, height=4)
            """,
            headers=["Feature", "Basic", "Pro", "Enterprise"],
            data=[["Storage", "10GB", "100GB", "Unlimited"]],
            variant="striped",
            size="lg",
        ),
    ],
    tags=["table", "data", "grid", "display"],
)
class Table(ComposableComponent):
    """
    Data table component with theme-aware styling.

    Features:
    - Multiple variants (default, bordered, striped, minimal)
    - Three sizes (sm, md, lg)
    - Theme-aware colors
    - Automatic cell sizing
    - Header formatting

    Usage:
        # Simple table
        table = Table(
            headers=["Name", "Value"],
            data=[["Item 1", "100"], ["Item 2", "200"]],
            variant="default",
            size="md"
        )
        table.render(slide, left=1, top=2, width=6, height=3)

        # Striped table with theme
        table = Table(
            headers=["Product", "Price", "Stock"],
            data=[["Widget", "$10", "50"], ["Gadget", "$25", "30"]],
            variant="striped",
            size="lg",
            theme=theme
        )
        table.render(slide, left=1, top=1, width=8, height=4)
    """

    def __init__(
        self,
        headers: List[str],
        data: List[List[str]],
        variant: str = "default",
        size: str = "md",
        theme: Optional[Any] = None,
    ):
        """
        Initialize table component.

        Args:
            headers: List of column header strings
            data: List of rows, where each row is a list of cell values
            variant: Visual variant (default, bordered, striped, minimal)
            size: Table size (sm, md, lg)
            theme: Optional theme override
        """
        super().__init__(theme)

        # Validate required parameters
        if not headers:
            raise ValueError("Table requires 'headers' (list of column names)")
        if not isinstance(headers, list):
            raise TypeError(f"Table 'headers' must be a list, got {type(headers).__name__}")
        if data is None:
            raise ValueError("Table requires 'data' (list of rows)")
        if not isinstance(data, list):
            raise TypeError(f"Table 'data' must be a list of rows, got {type(data).__name__}")

        # Validate each row has same number of columns as headers
        for i, row in enumerate(data):
            if not isinstance(row, list):
                raise TypeError(f"Table row {i} must be a list, got {type(row).__name__}")
            if len(row) != len(headers):
                raise ValueError(
                    f"Table row {i} has {len(row)} columns but headers has {len(headers)} columns. "
                    f"Each row must have the same number of values as headers."
                )

        self.headers = headers
        self.data = data
        self.variant = variant
        self.size = size

        # Get variant props
        self.variant_props = TABLE_VARIANTS.build(variant=variant, size=size)

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render table to slide or replace a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            Table shape object
        """
        # If placeholder provided, use its bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            import logging

            logger = logging.getLogger(__name__)
            logger.info(
                f"Table targeting placeholder - using bounds: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
            )

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        rows = len(self.data) + 1  # +1 for header
        cols = len(self.headers)

        # Create table
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        table = table_shape.table

        # Apply header styling
        self._apply_header_styles(table)

        # Apply data cell styling
        self._apply_data_styles(table)

        return table_shape

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def _apply_header_styles(self, table):
        """Apply styling to header row."""
        props = self.variant_props
        font_family = self._get_font_family()

        for col_idx, header_text in enumerate(self.headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header_text)

            # Text formatting
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = font_family
            paragraph.font.size = Pt(props.get("header_font_size", 13))

            if props.get("header_bold", True):
                paragraph.font.bold = True

            # Text color
            header_fg = props.get("header_fg", "card.foreground")
            paragraph.font.color.rgb = self.get_color(header_fg)

            # Cell background
            header_bg = props.get("header_bg", "card.DEFAULT")
            if header_bg != "transparent":
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.get_color(header_bg)

            # Cell margins
            padding = props.get("padding", 0.08)
            cell.text_frame.margin_left = Inches(padding)
            cell.text_frame.margin_right = Inches(padding)
            cell.text_frame.margin_top = Inches(padding / 2)
            cell.text_frame.margin_bottom = Inches(padding / 2)

    def _apply_data_styles(self, table):
        """Apply styling to data cells."""
        props = self.variant_props
        font_family = self._get_font_family()

        for row_idx, row_data in enumerate(self.data):
            actual_row_idx = row_idx + 1  # Skip header row

            for col_idx, value in enumerate(row_data):
                cell = table.cell(actual_row_idx, col_idx)
                cell.text = str(value)

                # Text formatting
                paragraph = cell.text_frame.paragraphs[0]
                paragraph.font.name = font_family
                paragraph.font.size = Pt(props.get("font_size", 12))

                # Text color
                cell_fg = props.get("cell_fg", "foreground.DEFAULT")
                paragraph.font.color.rgb = self.get_color(cell_fg)

                # Cell background (with alternating rows for striped variant)
                if self.variant == "striped" and actual_row_idx % 2 == 0:
                    # Alternate row
                    alt_bg = props.get("alt_bg", "muted.DEFAULT")
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.get_color(alt_bg)
                else:
                    # Normal row
                    cell_bg = props.get("cell_bg", "background.DEFAULT")
                    if cell_bg != "transparent":
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = self.get_color(cell_bg)

                # Cell margins
                padding = props.get("padding", 0.08)
                cell.text_frame.margin_left = Inches(padding)
                cell.text_frame.margin_right = Inches(padding)
                cell.text_frame.margin_top = Inches(padding / 2)
                cell.text_frame.margin_bottom = Inches(padding / 2)

                # Text alignment - left for first column, center for others
                if col_idx == 0:
                    paragraph.alignment = PP_ALIGN.LEFT
                else:
                    paragraph.alignment = PP_ALIGN.CENTER
