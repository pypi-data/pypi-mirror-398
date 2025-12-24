# src/chuk_mcp_pptx/components/core/alert.py
"""
Alert component for PowerPoint presentations.
Display important messages and notifications.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..composition import ComposableComponent, SubComponent, CardTitle, CardDescription
from ..variants import create_variants
from ..registry import component, ComponentCategory, prop, example
from ...tokens.typography import FONT_SIZES, PARAGRAPH_SPACING


# Alert-specific variants
ALERT_VARIANTS = create_variants(
    base={
        "border_radius": 8,
        "padding": 0.4,
        "border_width": 1,
    },
    variants={
        "variant": {
            "default": {
                "bg_color": "background.secondary",
                "fg_color": "foreground.DEFAULT",
                "border_color": "border.DEFAULT",
                "icon": "ℹ",
            },
            "info": {
                "bg_color": "info.DEFAULT",
                "fg_color": "info.foreground",
                "border_color": "info.DEFAULT",
                "icon": "ℹ",
            },
            "success": {
                "bg_color": "success.DEFAULT",
                "fg_color": "success.foreground",
                "border_color": "success.DEFAULT",
                "icon": "✓",
            },
            "warning": {
                "bg_color": "warning.DEFAULT",
                "fg_color": "warning.foreground",
                "border_color": "warning.DEFAULT",
                "icon": "⚠",
            },
            "error": {
                "bg_color": "destructive.DEFAULT",
                "fg_color": "destructive.foreground",
                "border_color": "destructive.DEFAULT",
                "icon": "✖",
            },
        }
    },
    default_variants={"variant": "default"},
)


@component(
    name="Alert",
    category=ComponentCategory.UI,
    description="Alert/notification component for displaying important messages",
    props=[
        prop(
            "variant",
            "string",
            "Alert type",
            options=["default", "info", "success", "warning", "error"],
            default="info",
            example="success",
        ),
        prop("title", "string", "Alert title", example="Success!"),
        prop("description", "string", "Alert description", example="Your changes have been saved."),
        prop("show_icon", "boolean", "Show icon indicator", default=True),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("width", "number", "Width in inches", default=5.0),
        prop("height", "number", "Height in inches", default=1.5),
    ],
    variants={"variant": ["default", "info", "success", "warning", "error"]},
    composition={
        "supports": ["AlertTitle", "AlertDescription"],
        "pattern": "Alert can be composed with title and description subcomponents",
    },
    examples=[
        example(
            description="Success alert",
            code="""
alert = Alert(
    variant="success",
    title="Success!",
    description="Your changes have been saved."
)
alert.render(slide, left=1, top=1)
            """,
            variant="success",
            title="Success!",
        ),
        example(
            description="Warning alert",
            code="""
alert = Alert(
    variant="warning",
    title="Warning",
    description="This action cannot be undone."
)
alert.render(slide, left=1, top=3)
            """,
            variant="warning",
            title="Warning",
        ),
        example(
            description="Composed alert",
            code="""
alert = Alert(variant="info")
alert.add_child(Alert.Title("Notice"))
alert.add_child(Alert.Description("Please review the changes."))
alert.render(slide, left=1, top=5)
            """,
            variant="info",
        ),
    ],
    tags=["alert", "notification", "message", "ui", "feedback"],
)
class Alert(ComposableComponent):
    """
    Alert component for important messages and notifications.

    Features:
    - Multiple variants (default, info, success, warning, error)
    - Optional icon indicators
    - Composition support (Alert.Title, Alert.Description)
    - Theme-aware colors

    Usage:
        # Simple alert
        alert = Alert(variant="success", title="Done!", description="Task completed")
        alert.render(slide, left=1, top=1)

        # Composed alert
        alert = Alert(variant="warning")
        alert.add_child(Alert.Title("Warning"))
        alert.add_child(Alert.Description("This is important"))
        alert.render(slide, left=1, top=3, width=6, height=2)
    """

    # Expose composition components as class attributes
    class Title(CardTitle):
        """Alert title subcomponent."""

        pass

    class Description(CardDescription):
        """Alert description subcomponent."""

        pass

    def __init__(
        self,
        variant: str = "info",
        title: Optional[str] = None,
        description: Optional[str] = None,
        show_icon: bool = True,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize alert component.

        Args:
            variant: Alert type (default, info, success, warning, error)
            title: Optional alert title
            description: Optional alert description
            show_icon: Whether to show icon indicator
            theme: Optional theme override
        """
        super().__init__(theme)
        self.variant = variant
        self.title_text = title
        self.description_text = description
        self.show_icon = show_icon

        # Get variant props
        self.variant_props = ALERT_VARIANTS.build(variant=variant)

        # Add default children if title/description provided
        if title:
            self.add_child(self.Title(title, theme))
        if description:
            self.add_child(self.Description(description, theme))

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float = 5.0,
        height: float = 1.5,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render alert to slide.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Alert width in inches
            height: Alert height in inches
            placeholder: Optional placeholder to replace

        Returns:
            Shape object representing the alert
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Create alert container
        alert = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Apply variant styling
        self._apply_variant_styles(alert)

        # Setup text frame
        text_frame = alert.text_frame
        text_frame.clear()
        text_frame.word_wrap = True

        padding = self.variant_props.get("padding", 0.4)
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        # Render children (title, description, etc.)
        if self._children:
            # Add icon inline with title if enabled
            if self.show_icon:
                icon = self.variant_props.get("icon", "ℹ")
                # Prepend icon to first child
                p = text_frame.paragraphs[0]
                p.text = f"{icon}  "
                p.alignment = PP_ALIGN.LEFT  # Explicitly left-align icon
                p.font.size = Pt(FONT_SIZES["base"])
                p.font.color.rgb = self.get_color(self.variant_props.get("fg_color"))

            for child in self._children:
                child.render_into(text_frame, self.theme)

        return alert

    def _apply_variant_styles(self, shape):
        """Apply variant-based styling to alert."""
        props = self.variant_props

        # Background
        bg_color = props.get("bg_color")
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.get_color(bg_color)

        # Border
        border_color = props.get("border_color", "border.DEFAULT")
        border_width = props.get("border_width", 1)
        shape.line.color.rgb = self.get_color(border_color)
        shape.line.width = Pt(border_width)


class AlertTitle(SubComponent):
    """Alert title subcomponent."""

    def __init__(self, text: str, theme: Optional[Dict[str, Any]] = None):
        super().__init__(theme)
        self.text = text

    def render_into(self, text_frame, theme: Optional[Dict[str, Any]] = None) -> Any:
        """Render title into alert text frame."""
        theme = theme or self.theme

        p = (
            text_frame.add_paragraph()
            if text_frame.paragraphs[0].text
            else text_frame.paragraphs[0]
        )
        p.text = self.text
        p.alignment = PP_ALIGN.LEFT  # Ensure left alignment
        style = self.get_text_style("h5")
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])
        p.font.bold = True

        return p


class AlertDescription(SubComponent):
    """Alert description subcomponent."""

    def __init__(self, text: str, theme: Optional[Dict[str, Any]] = None):
        super().__init__(theme)
        self.text = text

    def render_into(self, text_frame, theme: Optional[Dict[str, Any]] = None) -> Any:
        """Render description into alert text frame."""
        theme = theme or self.theme

        p = text_frame.add_paragraph()
        p.text = self.text
        p.alignment = PP_ALIGN.LEFT  # Ensure left alignment
        p.space_before = Pt(PARAGRAPH_SPACING["sm"])
        style = self.get_text_style("body-sm")
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])

        return p


# Add to Alert class
Alert.Title = AlertTitle  # type: ignore[misc]
Alert.Description = AlertDescription  # type: ignore[misc]
