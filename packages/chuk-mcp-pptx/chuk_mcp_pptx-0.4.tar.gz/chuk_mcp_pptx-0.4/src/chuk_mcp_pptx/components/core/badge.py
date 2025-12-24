# src/chuk_mcp_pptx/components/core/badge.py
"""
Badge component for PowerPoint presentations.
Small status indicators and labels.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

from ..base import Component
from ...constants import ComponentSizing
from ..variants import BADGE_VARIANTS
from ..registry import component, ComponentCategory, prop, example


@component(
    name="Badge",
    category=ComponentCategory.UI,
    description="Small badge/label component for status indicators and tags",
    props=[
        prop("text", "string", "Badge text", required=True, example="New"),
        prop(
            "variant",
            "string",
            "Visual variant",
            options=["default", "secondary", "success", "warning", "destructive", "outline"],
            default="default",
            example="success",
        ),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
    ],
    variants={"variant": ["default", "secondary", "success", "warning", "destructive", "outline"]},
    examples=[
        example(
            "Success badge",
            """
badge = Badge(text="Active", variant="success")
badge.render(slide, left=1, top=1)
            """,
            text="Active",
            variant="success",
        ),
        example(
            "Warning badge",
            """
badge = Badge(text="Beta", variant="warning")
badge.render(slide, left=2, top=1)
            """,
            text="Beta",
            variant="warning",
        ),
        example(
            "Destructive badge",
            """
badge = Badge(text="Deprecated", variant="destructive")
badge.render(slide, left=3, top=1)
            """,
            text="Deprecated",
            variant="destructive",
        ),
    ],
    tags=["badge", "label", "status", "tag", "ui"],
)
class Badge(Component):
    """
    Badge component for status indicators and labels.

    Features:
    - Multiple variants (default, secondary, success, warning, destructive, outline)
    - Auto-sized based on text
    - Theme-aware colors
    - Small, compact design

    Usage:
        # Simple badge
        badge = Badge(text="New", variant="default")
        badge.render(slide, left=1, top=1)

        # Status badge
        badge = Badge(text="Active", variant="success", theme=theme_dict)
        badge.render(slide, left=2, top=2)
    """

    def __init__(self, text: str, variant: str = "default", theme: Optional[Dict[str, Any]] = None):
        """
        Initialize badge component.

        Args:
            text: Badge text
            variant: Visual variant (default, secondary, success, warning, destructive, outline)
            theme: Optional theme override
        """
        super().__init__(theme)
        self.text = text
        self.variant = variant

        # Get variant props
        self.variant_props = BADGE_VARIANTS.build(variant=variant)

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render badge to slide or replace a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Optional width override (auto-calculated if None)
            height: Optional height override (default 0.3)
            placeholder: Optional placeholder shape to replace

        Returns:
            Shape object representing the badge
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Calculate dimensions
        badge_width = width or self._calculate_width()
        badge_height = height or 0.3

        # Create badge shape (using rounded rectangle)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(badge_width),
            Inches(badge_height),
        )

        # Apply variant styling
        self._apply_variant_styles(badge)

        # Add text
        self._add_text(badge)

        return badge

    def _calculate_width(self) -> float:
        """Calculate width based on text length."""
        # More generous estimation to prevent text cutoff
        char_width = ComponentSizing.BADGE_CHAR_WIDTH  # inches per character (increased from 0.05)
        min_width = 0.6
        padding = ComponentSizing.BADGE_PADDING  # More padding (increased from 0.3)

        calculated = len(self.text) * char_width + padding
        return max(min_width, calculated)

    def _apply_variant_styles(self, shape):
        """Apply variant-based styling to badge shape."""
        props = self.variant_props

        # Background color
        bg_color = props.get("bg_color")
        if bg_color and bg_color != "transparent":
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.get_color(bg_color)
        else:
            shape.fill.background()  # Transparent

        # Border
        border_width = props.get("border_width", 0)
        if border_width > 0:
            border_color = props.get("border_color", "border.DEFAULT")
            shape.line.color.rgb = self.get_color(border_color)
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()

    def _add_text(self, shape):
        """Add text to badge with proper styling."""
        props = self.variant_props

        text_frame = shape.text_frame
        text_frame.clear()

        # Disable word wrap and enable auto-size for single-line text
        text_frame.word_wrap = False
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        # Set tight padding for compact appearance
        padding = props.get("padding", 0.15)
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)

        # Add text
        paragraph = text_frame.paragraphs[0]
        paragraph.text = self.text
        paragraph.alignment = PP_ALIGN.CENTER

        # Font styling from variants
        paragraph.font.size = Pt(props.get("font_size", 10))
        # Handle both Theme objects and dict themes
        if hasattr(self.theme, "typography"):
            font_family = self.theme.typography.get("font_family", "Inter")
        else:
            font_family = (
                self.theme.get("font_family", "Inter") if isinstance(self.theme, dict) else "Inter"
            )
        paragraph.font.name = font_family

        # Text color
        fg_color = props.get("fg_color", "foreground.DEFAULT")
        paragraph.font.color.rgb = self.get_color(fg_color)

        # Font weight
        if props.get("font_weight", 400) >= 600:
            paragraph.font.bold = True


@component(
    name="DotBadge",
    category=ComponentCategory.UI,
    description="Small dot indicator for status (like notification dots)",
    props=[
        prop(
            "variant",
            "string",
            "Color variant",
            options=["default", "success", "warning", "destructive"],
            default="default",
        ),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("size", "number", "Dot size in inches", default=0.15),
    ],
    examples=[
        example(
            "Success dot",
            """
dot = DotBadge(variant="success")
dot.render(slide, left=1, top=1)
            """,
            variant="success",
        )
    ],
    tags=["badge", "dot", "indicator", "status"],
)
class DotBadge(Component):
    """
    Dot badge - small colored dot for status indicators.
    Perfect for notification dots, online status, etc.
    """

    def __init__(
        self, variant: str = "default", size: float = 0.15, theme: Optional[Dict[str, Any]] = None
    ):
        """
        Initialize dot badge.

        Args:
            variant: Color variant (default, success, warning, destructive)
            size: Dot size in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.variant = variant
        self.size = size

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render dot badge or replace a placeholder.

        Args:
            slide: PowerPoint slide
            left: Left position
            top: Top position
            width: Optional width (uses size if not provided)
            height: Optional height (uses size if not provided)
            placeholder: Optional placeholder shape to replace

        Returns:
            Oval shape representing the dot
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Use size if width/height not provided
        dot_size = width or height or self.size

        # Create circular dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(dot_size), Inches(dot_size)
        )

        # Get color based on variant
        color_map = {
            "default": "primary.DEFAULT",
            "success": "success.DEFAULT",
            "warning": "warning.DEFAULT",
            "destructive": "destructive.DEFAULT",
        }

        color = self.get_color(color_map.get(self.variant, "primary.DEFAULT"))

        # Apply color
        dot.fill.solid()
        dot.fill.fore_color.rgb = color

        # No border
        dot.line.fill.background()

        return dot


@component(
    name="CountBadge",
    category=ComponentCategory.UI,
    description="Badge showing a count/number (like notification counts)",
    props=[
        prop("count", "number", "Count to display", required=True, example=5),
        prop(
            "variant",
            "string",
            "Visual variant",
            options=["default", "destructive"],
            default="destructive",
        ),
        prop("max", "number", "Maximum count before showing '+'", default=99),
    ],
    examples=[
        example(
            "Notification count",
            """
badge = CountBadge(count=5, variant="destructive")
badge.render(slide, left=1, top=1)
            """,
            count=5,
            variant="destructive",
        ),
        example(
            "High count with max",
            """
badge = CountBadge(count=150, max=99)  # Shows "99+"
badge.render(slide, left=2, top=1)
            """,
            count=150,
            max=99,
        ),
    ],
    tags=["badge", "count", "notification", "number"],
)
class CountBadge(Badge):
    """
    Count badge - displays a number (e.g., notification count).
    Shows "+ " when count exceeds max.
    """

    def __init__(
        self,
        count: int,
        variant: str = "destructive",
        max_count: int = 99,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize count badge.

        Args:
            count: Number to display
            variant: Visual variant
            max_count: Maximum before showing "+"
            theme: Optional theme
        """
        # Format text
        if count > max_count:
            text = f"{max_count}+"
        else:
            text = str(count)

        super().__init__(text, variant, theme)
        self.count = count
        self.max_count = max_count

    def _calculate_width(self) -> float:
        """Count badges are more compact and circular."""
        # Make roughly circular for small counts
        if self.count < 10:
            return 0.3
        elif self.count < 100:
            return 0.4
        else:
            return 0.5
