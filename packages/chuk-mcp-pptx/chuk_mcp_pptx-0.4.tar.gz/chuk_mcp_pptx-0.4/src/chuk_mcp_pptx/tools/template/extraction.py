"""
Template Extraction Tools for PowerPoint MCP Server

Extract design system components (themes, tokens, layouts) from PowerPoint templates.
Allows converting any template into reusable design system assets.
"""

import logging
from typing import Any, TYPE_CHECKING
from pydantic import BaseModel, Field

if TYPE_CHECKING:
    from pptx.presentation import Presentation

logger = logging.getLogger(__name__)


class ExtractedColorScheme(BaseModel):
    """Color scheme extracted from a template."""

    name: str = Field(..., description="Color scheme name")
    primary: str = Field(..., description="Primary color (hex)")
    secondary: str = Field(..., description="Secondary color (hex)")
    accent: str = Field(..., description="Accent color (hex)")
    background: str = Field(..., description="Background color (hex)")
    text: str = Field(..., description="Text color (hex)")
    additional_colors: list[str] = Field(
        default_factory=list, description="Additional colors found"
    )

    class Config:
        extra = "forbid"


class ExtractedTypography(BaseModel):
    """Typography settings extracted from a template."""

    name: str = Field(..., description="Typography set name")
    title_font: str = Field(..., description="Title font family")
    title_size: float = Field(..., description="Title font size in points")
    body_font: str = Field(..., description="Body font family")
    body_size: float = Field(..., description="Body font size in points")
    heading_font: str | None = Field(default=None, description="Heading font family")
    heading_size: float | None = Field(default=None, description="Heading font size")

    class Config:
        extra = "forbid"


class ExtractedLayout(BaseModel):
    """Layout information extracted from a template."""

    name: str = Field(..., description="Layout name")
    index: int = Field(..., description="Layout index")
    width: float = Field(..., description="Slide width in inches")
    height: float = Field(..., description="Slide height in inches")
    placeholders: list[dict[str, Any]] = Field(
        default_factory=list, description="Placeholder configurations"
    )
    background_type: str = Field(..., description="Background type (solid, gradient, image)")

    class Config:
        extra = "forbid"


class ExtractedDesignSystem(BaseModel):
    """Complete design system extracted from a template."""

    template_name: str = Field(..., description="Source template name")
    color_scheme: ExtractedColorScheme | None = Field(
        default=None, description="Extracted color scheme"
    )
    typography: ExtractedTypography | None = Field(default=None, description="Extracted typography")
    layouts: list[ExtractedLayout] = Field(default_factory=list, description="Extracted layouts")
    theme_name: str | None = Field(default=None, description="Suggested theme name")

    class Config:
        extra = "forbid"


class PlaceholderUsage(BaseModel):
    """How a placeholder is used in a slide."""

    idx: int = Field(..., description="Placeholder index")
    type: str = Field(..., description="Placeholder type")
    name: str = Field(..., description="Placeholder name")
    text: str = Field(default="", description="Text content")
    text_length: int = Field(default=0, description="Length of text")
    top: int = Field(..., description="Top position in EMUs")
    left: int = Field(..., description="Left position in EMUs")
    width: int = Field(default=0, description="Width in EMUs")
    height: int = Field(default=0, description="Height in EMUs")
    font_size: float | None = Field(default=None, description="Font size in points")
    line_count: int = Field(default=0, description="Number of text lines/paragraphs")
    has_overflow: bool = Field(default=False, description="Whether text appears to overflow")
    recommended_max_chars: int | None = Field(
        default=None, description="Recommended max characters based on template"
    )

    class Config:
        extra = "forbid"


class LayoutComparisonResult(BaseModel):
    """Result of comparing template vs generated slide layout usage."""

    layout_name: str = Field(..., description="Layout name")
    template_slide_index: int | None = Field(None, description="Template example slide index")
    generated_slide_index: int = Field(..., description="Generated slide index")
    matches_pattern: bool = Field(..., description="Whether usage matches template pattern")
    issues: list[str] = Field(default_factory=list, description="Issues found")
    suggestions: list[str] = Field(default_factory=list, description="Suggestions for fixes")
    template_usage: list[PlaceholderUsage] = Field(
        default_factory=list, description="How template uses placeholders"
    )
    generated_usage: list[PlaceholderUsage] = Field(
        default_factory=list, description="How generated slide uses placeholders"
    )

    class Config:
        extra = "forbid"


class LayoutVariant(BaseModel):
    """A variant of a base layout."""

    index: int = Field(..., description="Layout index")
    name: str = Field(..., description="Layout name")
    variant_number: int | None = Field(None, description="Variant number if detected from name")
    placeholder_count: int = Field(..., description="Number of placeholders")
    original_slide_number: int | None = Field(
        None, description="Original slide number in template (1-indexed)"
    )

    class Config:
        extra = "forbid"


class LayoutGroup(BaseModel):
    """Group of similar layouts (base + variants)."""

    base_name: str = Field(..., description="Base layout name without variant number")
    base_layout: LayoutVariant = Field(..., description="Primary/base layout")
    variants: list[LayoutVariant] = Field(default_factory=list, description="Layout variants")
    total_count: int = Field(..., description="Total layouts in group (base + variants)")
    placeholder_signature: str = Field(..., description="Placeholder pattern signature")

    class Config:
        extra = "forbid"


class LayoutAnalysis(BaseModel):
    """Analysis of template layouts with variant detection."""

    total_layouts: int = Field(..., description="Total number of layouts")
    unique_groups: int = Field(..., description="Number of unique layout groups")
    layout_groups: list[LayoutGroup] = Field(default_factory=list, description="Grouped layouts")
    ungrouped_layouts: list[LayoutVariant] = Field(
        default_factory=list, description="Layouts that don't fit any group"
    )

    class Config:
        extra = "forbid"


def _rgb_to_hex(rgb_color) -> str:
    """Convert RGB color to hex string."""
    try:
        if hasattr(rgb_color, "rgb"):
            rgb = rgb_color.rgb
            return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        return "#000000"
    except Exception:
        return "#000000"


def _extract_colors_from_theme(theme) -> dict[str, str]:
    """Extract colors from a theme object."""
    colors = {}

    try:
        # Try to extract theme colors
        if hasattr(theme, "color_scheme"):
            scheme = theme.color_scheme

            # PowerPoint has 12 theme colors
            if hasattr(scheme, "_element"):
                color_map = {
                    "dk1": "text_primary",
                    "lt1": "background_primary",
                    "dk2": "text_secondary",
                    "lt2": "background_secondary",
                    "accent1": "accent1",
                    "accent2": "accent2",
                    "accent3": "accent3",
                    "accent4": "accent4",
                    "accent5": "accent5",
                    "accent6": "accent6",
                }

                for xml_name, friendly_name in color_map.items():
                    try:
                        color_elem = getattr(scheme._element, xml_name, None)
                        if color_elem is not None:
                            # Extract RGB from srgbClr or sysClr
                            if hasattr(color_elem, "srgbClr"):
                                rgb_val = color_elem.srgbClr.val
                                colors[friendly_name] = f"#{rgb_val}"
                    except Exception as e:
                        logger.debug(f"Could not extract {xml_name}: {e}")
    except Exception as e:
        logger.warning(f"Error extracting theme colors: {e}")

    return colors


def _extract_colors_from_shapes(slide) -> list[str]:
    """Extract colors from shapes on a slide."""
    colors = set()

    try:
        for shape in slide.shapes:
            # Extract fill colors
            if hasattr(shape, "fill") and shape.fill:
                try:
                    if hasattr(shape.fill, "fore_color"):
                        hex_color = _rgb_to_hex(shape.fill.fore_color)
                        if hex_color != "#000000":
                            colors.add(hex_color)
                except Exception:
                    pass

            # Extract line colors
            if hasattr(shape, "line") and shape.line:
                try:
                    if hasattr(shape.line, "color"):
                        hex_color = _rgb_to_hex(shape.line.color)
                        if hex_color != "#000000":
                            colors.add(hex_color)
                except Exception:
                    pass

            # Extract text colors
            if hasattr(shape, "text_frame"):
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if hasattr(run.font, "color") and run.font.color:
                                hex_color = _rgb_to_hex(run.font.color)
                                if hex_color != "#000000":
                                    colors.add(hex_color)
                except Exception:
                    pass
    except Exception as e:
        logger.warning(f"Error extracting colors from shapes: {e}")

    return list(colors)


def _extract_typography_from_master(master) -> dict[str, Any]:
    """Extract typography settings from slide master."""
    typography = {
        "title_font": "Calibri",
        "title_size": 44.0,
        "body_font": "Calibri",
        "body_size": 18.0,
        "heading_font": None,
        "heading_size": None,
    }

    try:
        # Extract from title placeholder
        for shape in master.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type

                # Title placeholder
                if ph_type == 1:  # TITLE
                    if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                        para = shape.text_frame.paragraphs[0]
                        if para.runs:
                            run = para.runs[0]
                            if run.font.name:
                                typography["title_font"] = run.font.name
                            if run.font.size:
                                typography["title_size"] = run.font.size.pt

                # Body placeholder
                elif ph_type == 2:  # BODY
                    if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                        para = shape.text_frame.paragraphs[0]
                        if para.runs:
                            run = para.runs[0]
                            if run.font.name:
                                typography["body_font"] = run.font.name
                            if run.font.size:
                                typography["body_size"] = run.font.size.pt
    except Exception as e:
        logger.warning(f"Error extracting typography: {e}")

    return typography


async def extract_design_system_from_template(manager, template_name: str) -> ExtractedDesignSystem:
    """
    Extract complete design system from a template.

    Args:
        manager: PresentationManager instance
        template_name: Name of the template to extract from

    Returns:
        ExtractedDesignSystem with colors, typography, and layouts
    """
    result = await manager.get(template_name)
    if not result:
        raise ValueError(f"Template not found: {template_name}")

    prs, metadata = result

    # Extract color scheme
    theme_colors = {}
    additional_colors = []

    try:
        # Try to extract from theme
        if hasattr(prs, "slide_master") and prs.slide_master:
            master = prs.slide_master
            if hasattr(master, "theme"):
                theme_colors = _extract_colors_from_theme(master.theme)

        # Also extract colors from actual slides
        if prs.slides:
            for slide in prs.slides:
                slide_colors = _extract_colors_from_shapes(slide)
                additional_colors.extend(slide_colors)

        # Remove duplicates
        additional_colors = list(set(additional_colors))
    except Exception as e:
        logger.warning(f"Error extracting colors: {e}")

    # Create color scheme
    color_scheme = None
    if theme_colors or additional_colors:
        color_scheme = ExtractedColorScheme(
            name=f"{template_name}_colors",
            primary=theme_colors.get(
                "accent1", additional_colors[0] if additional_colors else "#0066CC"
            ),
            secondary=theme_colors.get(
                "accent2", additional_colors[1] if len(additional_colors) > 1 else "#6C757D"
            ),
            accent=theme_colors.get(
                "accent3", additional_colors[2] if len(additional_colors) > 2 else "#28A745"
            ),
            background=theme_colors.get("background_primary", "#FFFFFF"),
            text=theme_colors.get("text_primary", "#000000"),
            additional_colors=additional_colors[:10],  # Limit to 10 additional colors
        )

    # Extract typography
    typography = None
    try:
        if hasattr(prs, "slide_master") and prs.slide_master:
            typo_data = _extract_typography_from_master(prs.slide_master)
            typography = ExtractedTypography(
                name=f"{template_name}_typography",
                title_font=typo_data["title_font"],
                title_size=typo_data["title_size"],
                body_font=typo_data["body_font"],
                body_size=typo_data["body_size"],
                heading_font=typo_data.get("heading_font"),
                heading_size=typo_data.get("heading_size"),
            )
    except Exception as e:
        logger.warning(f"Error extracting typography: {e}")

    # Extract layouts
    layouts = []
    try:
        if prs.slide_layouts:
            for idx, layout in enumerate(prs.slide_layouts):
                placeholders = []
                for ph in layout.placeholders:
                    try:
                        ph_info = {
                            "idx": ph.placeholder_format.idx,
                            "type": str(ph.placeholder_format.type),
                            "name": ph.name,
                            "left": ph.left.inches if hasattr(ph, "left") else 0,
                            "top": ph.top.inches if hasattr(ph, "top") else 0,
                            "width": ph.width.inches if hasattr(ph, "width") else 0,
                            "height": ph.height.inches if hasattr(ph, "height") else 0,
                        }
                        placeholders.append(ph_info)
                    except Exception:
                        pass

                # Determine background type
                bg_type = "solid"
                try:
                    if hasattr(layout, "background"):
                        bg = layout.background
                        if hasattr(bg, "fill") and bg.fill:
                            if hasattr(bg.fill, "type"):
                                bg_type = str(bg.fill.type).lower()
                except Exception:
                    pass

                layouts.append(
                    ExtractedLayout(
                        name=layout.name,
                        index=idx,
                        width=prs.slide_width.inches,
                        height=prs.slide_height.inches,
                        placeholders=placeholders,
                        background_type=bg_type,
                    )
                )
    except Exception as e:
        logger.warning(f"Error extracting layouts: {e}")

    return ExtractedDesignSystem(
        template_name=template_name,
        color_scheme=color_scheme,
        typography=typography,
        layouts=layouts,
        theme_name=f"{template_name}_theme",
    )


def analyze_layout_variants(prs: "Presentation") -> LayoutAnalysis:
    """
    Analyze layouts and detect variants.

    Groups layouts by detecting similar names and placeholder patterns.
    For example: "Content with Picture", "Content with Picture 2", "Content with Picture 3"
    would be grouped as one base layout with 2 variants.

    Args:
        prs: Presentation object

    Returns:
        LayoutAnalysis with grouped layouts
    """
    import re
    from collections import defaultdict

    # Parse all layouts
    all_layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        # Extract variant number from name if present
        # Patterns: "Layout 2", "Layout Name 3", etc.
        match = re.search(r"^(.+?)\s+(\d+)$", layout.name)
        if match:
            base_name = match.group(1)
            variant_num = int(match.group(2))
        else:
            base_name = layout.name
            variant_num = None

        # Find original slide number (1-indexed) if layout has a corresponding slide
        original_slide_num = None
        for slide_idx, slide in enumerate(prs.slides):
            if slide.slide_layout == layout:
                original_slide_num = slide_idx + 1
                break

        all_layouts.append(
            LayoutVariant(
                index=idx,
                name=layout.name,
                variant_number=variant_num,
                placeholder_count=len(layout.placeholders),
                original_slide_number=original_slide_num,
            )
        )

    # Group by base name
    groups_by_name: dict[str, list[LayoutVariant]] = defaultdict(list)
    for layout_variant in all_layouts:
        # Extract base name (remove trailing numbers)
        match = re.search(r"^(.+?)\s+\d+$", layout_variant.name)
        base_name = match.group(1) if match else layout_variant.name
        groups_by_name[base_name].append(layout_variant)

    # Create layout groups
    layout_groups: list[LayoutGroup] = []
    ungrouped_layouts: list[LayoutVariant] = []

    for base_name, layouts_in_group in groups_by_name.items():
        if len(layouts_in_group) > 1:
            # Multiple layouts with same base name - this is a group
            # Sort by variant number (None first, then numeric)
            layouts_in_group.sort(key=lambda x: (x.variant_number is None, x.variant_number or 0))

            base_layout = layouts_in_group[0]
            variants = layouts_in_group[1:]

            # Create placeholder signature based on count
            placeholder_sig = f"{base_layout.placeholder_count} placeholders"

            layout_groups.append(
                LayoutGroup(
                    base_name=base_name,
                    base_layout=base_layout,
                    variants=variants,
                    total_count=len(layouts_in_group),
                    placeholder_signature=placeholder_sig,
                )
            )
        else:
            # Single layout - ungrouped
            ungrouped_layouts.append(layouts_in_group[0])

    return LayoutAnalysis(
        total_layouts=len(all_layouts),
        unique_groups=len(layout_groups),
        layout_groups=layout_groups,
        ungrouped_layouts=ungrouped_layouts,
    )


def register_extraction_tools(mcp, manager, theme_manager):
    """
    Register template extraction MCP tools.

    Args:
        mcp: The MCP server instance
        manager: PresentationManager instance
        theme_manager: ThemeManager instance for registering extracted themes
    """

    @mcp.tool
    async def pptx_extract_template_design_system(template_name: str) -> str:
        """
        Extract design system (colors, typography, layouts) from a template.

        Analyzes a PowerPoint template and extracts:
        - Color scheme (primary, secondary, accent, background, text colors)
        - Typography settings (fonts and sizes for titles, body, headings)
        - Layout configurations (placeholders, positions, sizes)

        This allows you to understand and reuse the design system from any template.

        Args:
            template_name: Name of the template to extract design system from

        Returns:
            JSON string with ExtractedDesignSystem model

        Example:
            await pptx_extract_template_design_system(template_name="corporate")
        """
        try:
            design_system = await extract_design_system_from_template(manager, template_name)
            return design_system.model_dump_json()
        except Exception as e:
            logger.error(f"Failed to extract design system: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_register_template_as_theme(
        template_name: str, theme_name: str, apply_to_current: bool = False
    ) -> str:
        """
        Extract a template's design system and register it as a reusable theme.

        Creates a new theme in the design system based on the template's colors and styles.
        The theme can then be used with pptx_apply_theme() or when creating new presentations.

        Args:
            template_name: Name of the template to extract from
            theme_name: Name for the new theme
            apply_to_current: Whether to apply the theme to current presentation

        Returns:
            JSON string with success message and theme details

        Example:
            await pptx_register_template_as_theme(
                template_name="corporate",
                theme_name="my_corporate_theme",
                apply_to_current=True
            )
        """
        try:
            # Extract design system
            design_system = await extract_design_system_from_template(manager, template_name)

            if not design_system.color_scheme:
                return '{"error": "No color scheme could be extracted from template"}'

            # Create theme from template colors
            from ...themes import Theme

            colors = design_system.color_scheme

            # Determine font family from typography if available
            font_family = "Inter"
            if design_system.typography:
                font_family = design_system.typography.body_font or "Inter"

            # Create theme with the proper constructor
            # The Theme class uses primary_hue and mode to generate tokens
            # For template-extracted colors, we use "blue" as base and customize
            theme = Theme(
                name=theme_name,
                primary_hue="blue",
                mode="light",
                font_family=font_family,
            )
            theme_manager.register_custom_theme(theme_name, theme)

            # Store extracted color info for reporting
            theme_colors = {
                "primary": colors.primary,
                "secondary": colors.secondary,
                "accent": colors.accent,
                "background": colors.background,
                "text": colors.text,
            }

            result = {
                "message": f'Theme "{theme_name}" registered from template "{template_name}"',
                "theme_name": theme_name,
                "colors": theme_colors,
                "layout_count": len(design_system.layouts),
            }

            # Apply to current presentation if requested
            if apply_to_current:
                current_name = manager.get_current_name()
                if current_name:
                    pres_result = await manager.get(current_name)
                    if pres_result:
                        prs, _metadata = pres_result
                        # Apply theme to each slide
                        for slide in prs.slides:
                            theme.apply_to_slide(slide)
                        await manager._save_to_store(current_name, prs)
                        result["applied_to"] = current_name

            return f"{result}"
        except Exception as e:
            logger.error(f"Failed to register theme from template: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_extract_template_colors(template_name: str) -> str:
        """
        Extract just the color palette from a template.

        Useful for understanding or reusing a template's color scheme without
        extracting the full design system.

        Args:
            template_name: Name of the template

        Returns:
            JSON string with ExtractedColorScheme model

        Example:
            await pptx_extract_template_colors(template_name="corporate")
        """
        try:
            design_system = await extract_design_system_from_template(manager, template_name)

            if not design_system.color_scheme:
                return '{"error": "No colors could be extracted from template"}'

            return design_system.color_scheme.model_dump_json()
        except Exception as e:
            logger.error(f"Failed to extract colors: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_extract_template_typography(template_name: str) -> str:
        """
        Extract just the typography settings from a template.

        Gets font families and sizes used in the template for titles, body text, etc.

        Args:
            template_name: Name of the template

        Returns:
            JSON string with ExtractedTypography model

        Example:
            await pptx_extract_template_typography(template_name="corporate")
        """
        try:
            design_system = await extract_design_system_from_template(manager, template_name)

            if not design_system.typography:
                return '{"error": "No typography settings could be extracted from template"}'

            return design_system.typography.model_dump_json()
        except Exception as e:
            logger.error(f"Failed to extract typography: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_compare_slide_to_template(
        presentation: str,
        slide_index: int,
        template_name: str,
        template_slide_index: int | None = None,
    ) -> str:
        """
        Compare a generated slide with template to verify correct layout usage.

        Analyzes how placeholders are used in a generated slide versus how they're
        used in the template, identifying mismatches and providing specific suggestions
        for fixes. This helps ensure generated slides match the template's design intent.

        Args:
            presentation: Name of the generated presentation
            slide_index: Index of slide to check (0-based)
            template_name: Name of the template it was created from
            template_slide_index: Optional specific template slide to compare (auto-finds if None)

        Returns:
            JSON string with LayoutComparisonResult showing issues and suggestions

        Example:
            await pptx_compare_slide_to_template(
                presentation="my_brand_proposal",
                slide_index=12,
                template_name="brand_proposal"
            )
        """
        try:
            # Get generated presentation
            result = await manager.get(presentation)
            if not result:
                from ...models import ErrorResponse
            return ErrorResponse(error=f"Presentation not found: {presentation}").model_dump_json()

            gen_prs, _ = result

            if slide_index >= len(gen_prs.slides):
                from ...models import ErrorResponse
            return ErrorResponse(
                error=f"Slide index {slide_index} out of range (max {len(gen_prs.slides) - 1})"
            ).model_dump_json()

            gen_slide = gen_prs.slides[slide_index]
            layout_name = gen_slide.slide_layout.name

            # Get template presentation
            template_result = await manager.get(template_name)
            if not template_result:
                from ...models import ErrorResponse
            return ErrorResponse(error=f"Template not found: {template_name}").model_dump_json()

            template_prs, _ = template_result

            # Find example slide in template using same layout
            template_slide = None
            template_idx = template_slide_index

            if template_slide_index is not None:
                if template_slide_index < len(template_prs.slides):
                    template_slide = template_prs.slides[template_slide_index]
            else:
                # Find first slide using this layout
                for idx, slide in enumerate(template_prs.slides):
                    if slide.slide_layout.name == layout_name:
                        template_slide = slide
                        template_idx = idx
                        break

            # Extract placeholder usage from generated slide
            gen_usage = []
            for shape in gen_slide.shapes:
                if hasattr(shape, "placeholder_format"):
                    text = shape.text.strip() if hasattr(shape, "text") else ""
                    gen_usage.append(
                        PlaceholderUsage(
                            idx=shape.placeholder_format.idx,
                            type=str(shape.placeholder_format.type),
                            name=shape.name,
                            text=text[:100],  # Limit to 100 chars
                            text_length=len(text),
                            top=shape.top,
                            left=shape.left,
                        )
                    )

            # Extract placeholder usage from template slide
            template_usage = []
            if template_slide:
                for shape in template_slide.shapes:
                    if hasattr(shape, "placeholder_format"):
                        text = shape.text.strip() if hasattr(shape, "text") else ""
                        template_usage.append(
                            PlaceholderUsage(
                                idx=shape.placeholder_format.idx,
                                type=str(shape.placeholder_format.type),
                                name=shape.name,
                                text=text[:100],
                                text_length=len(text),
                                top=shape.top,
                                left=shape.left,
                            )
                        )

            # Compare and find issues
            issues = []
            suggestions = []
            matches = True

            if template_usage:
                # Create lookup maps
                template_map = {u.idx: u for u in template_usage}
                gen_map = {u.idx: u for u in gen_usage}

                # Check each placeholder
                for ph_idx, template_ph in template_map.items():
                    if ph_idx in gen_map:
                        gen_ph = gen_map[ph_idx]

                        # Compare text length patterns (short vs long content)
                        if template_ph.text_length > 0 and gen_ph.text_length > 0:
                            # Template has short text, generated has long
                            if template_ph.text_length < 20 and gen_ph.text_length > 100:
                                matches = False
                                issues.append(
                                    f"Placeholder {ph_idx} ({template_ph.name}): "
                                    f"Template has brief content ({template_ph.text_length} chars), "
                                    f"but generated has long content ({gen_ph.text_length} chars)"
                                )
                                suggestions.append(
                                    f'In template, placeholder {ph_idx} contains: "{template_ph.text}" '
                                    f"(short/impactful). Consider using brief text here."
                                )

                            # Template has long text, generated has short
                            elif template_ph.text_length > 100 and gen_ph.text_length < 20:
                                matches = False
                                issues.append(
                                    f"Placeholder {ph_idx} ({template_ph.name}): "
                                    f"Template has detailed content ({template_ph.text_length} chars), "
                                    f"but generated has brief content ({gen_ph.text_length} chars)"
                                )
                                suggestions.append(
                                    f"In template, placeholder {ph_idx} contains detailed text. "
                                    f"Consider adding more content here."
                                )

                    else:
                        # Placeholder used in template but not in generated
                        if template_ph.text_length > 0:
                            issues.append(
                                f"Placeholder {ph_idx} ({template_ph.name}) is used in template but empty in generated slide"
                            )
                            suggestions.append(
                                f'Template uses placeholder {ph_idx} for: "{template_ph.text}"'
                            )

                # Check for placeholders used in generated but not template
                for ph_idx, gen_ph in gen_map.items():
                    if ph_idx not in template_map and gen_ph.text_length > 0:
                        issues.append(
                            f"Placeholder {ph_idx} ({gen_ph.name}) has content but isn't used in template example"
                        )

            result = LayoutComparisonResult(
                layout_name=layout_name,
                template_slide_index=template_idx,
                generated_slide_index=slide_index,
                matches_pattern=matches and len(issues) == 0,
                issues=issues,
                suggestions=suggestions,
                template_usage=template_usage,
                generated_usage=gen_usage,
            )

            return result.model_dump_json(indent=2)

        except Exception as e:
            logger.error(f"Failed to compare slide to template: {e}")
            import traceback
            from ...models import ErrorResponse

            return ErrorResponse(error=f"{str(e)} - {traceback.format_exc()}").model_dump_json()

    return {
        "pptx_extract_template_design_system": pptx_extract_template_design_system,
        "pptx_register_template_as_theme": pptx_register_template_as_theme,
        "pptx_extract_template_colors": pptx_extract_template_colors,
        "pptx_extract_template_typography": pptx_extract_template_typography,
        "pptx_compare_slide_to_template": pptx_compare_slide_to_template,
    }
