"""
Template analysis tools.

Provides tools to analyze template layouts, placeholders, and structure.
All responses use Pydantic models for type safety.
"""

import io
import logging
from pptx import Presentation
from .models import LayoutInfo, LayoutPlaceholderInfo, TemplateInfo

logger = logging.getLogger(__name__)


def register_analyze_tools(mcp, manager, template_manager):
    """Register template analysis tools."""

    @mcp.tool
    async def pptx_analyze_template(template_name: str) -> str:
        """
        Analyze a template's layouts and structure.

        CRITICAL FIRST STEP when working with templates. This tool provides detailed
        information about ALL available slide layouts, their placeholders, and structure.
        You MUST call this after creating a presentation from a template to understand
        what layouts are available.

        Args:
            template_name: Name of the template to analyze

        Returns:
            JSON string with TemplateInfo model containing:
            - name: Template name
            - layout_count: Number of available layouts
            - layouts: Array of layout objects, each with:
                - index: Layout index (use this in pptx_add_slide_from_template)
                - name: Human-readable layout name (e.g., "Title Slide", "Content with Picture")
                - placeholder_count: Number of placeholders in this layout
                - placeholders: Array of placeholder objects with idx, type, and name

        TEMPLATE WORKFLOW:
            1. Create presentation from template:
               await pptx_create(name="my_deck", template_name="brand_proposal")

            2. IMMEDIATELY analyze template to see all layouts:
               layouts = await pptx_analyze_template("brand_proposal")
               # Shows ALL 55 layouts with names and placeholders

            3. Choose appropriate layouts for each slide you want to create:
               # Example response shows:
               # Layout 0: "Title with Picture" - placeholders: idx=0 (TITLE), idx=1 (SUBTITLE), idx=2 (PICTURE)
               # Layout 1: "Content" - placeholders: idx=0 (TITLE), idx=1 (OBJECT)
               # Layout 10: "Agenda" - placeholders: idx=0 (TITLE), idx=13 (OBJECT)
               # ... and 52 more layouts

            4. Add slides using specific layout indices:
               await pptx_add_slide_from_template(layout_index=0)  # Title slide
               await pptx_add_slide_from_template(layout_index=1)  # Content slide
               await pptx_add_slide_from_template(layout_index=10) # Agenda slide

        IMPORTANT: Different layouts have different placeholders. Always review
        the placeholder information for each layout to know which placeholder_idx
        values to use when calling pptx_populate_placeholder.

        Example:
            # Analyze brand_proposal template
            result = await pptx_analyze_template("brand_proposal")
            # Returns:
            # {
            #   "name": "brand_proposal",
            #   "layout_count": 55,
            #   "layouts": [
            #     {
            #       "index": 0,
            #       "name": "Title with Picture",
            #       "placeholder_count": 3,
            #       "placeholders": [
            #         {"idx": 0, "type": "TITLE", "name": "Title 1"},
            #         {"idx": 1, "type": "SUBTITLE", "name": "Subtitle 2"},
            #         {"idx": 2, "type": "PICTURE", "name": "Picture Placeholder 3"}
            #       ]
            #     },
            #     {
            #       "index": 1,
            #       "name": "Content",
            #       "placeholder_count": 2,
            #       "placeholders": [
            #         {"idx": 0, "type": "TITLE", "name": "Title 1"},
            #         {"idx": 1, "type": "OBJECT", "name": "Content Placeholder 2"}
            #       ]
            #     },
            #     ... 53 more layouts
            #   ]
            # }
        """
        try:
            # First check if it's a builtin template
            template_data = await template_manager.get_template_data(template_name)
            if template_data:
                # Load builtin template
                buffer = io.BytesIO(template_data)
                prs = Presentation(buffer)
                metadata = None
            else:
                # Get from artifact store
                result = await manager.get(template_name)
                if not result:
                    from ...models import ErrorResponse

                    return ErrorResponse(
                        error=f"Template not found: {template_name}"
                    ).model_dump_json()
                prs, metadata = result

            # Analyze layouts
            layouts = []
            if prs.slide_layouts:
                for idx, layout in enumerate(prs.slide_layouts):
                    placeholders: list[LayoutPlaceholderInfo] = []
                    for placeholder in layout.placeholders:
                        ph_info = LayoutPlaceholderInfo(
                            idx=placeholder.placeholder_format.idx,
                            type=str(placeholder.placeholder_format.type),
                            name=placeholder.name,
                        )
                        placeholders.append(ph_info)

                    layout_info = LayoutInfo(
                        index=idx,
                        name=layout.name,
                        placeholder_count=len(placeholders),
                        placeholders=placeholders,
                    )
                    layouts.append(layout_info)

            # Check for theme
            has_theme = hasattr(prs, "slide_master") and prs.slide_master is not None

            template_info = TemplateInfo(
                name=template_name,
                slide_count=len(prs.slides),
                layout_count=len(prs.slide_layouts) if prs.slide_layouts else 0,
                layouts=layouts,
                master_count=len(prs.slide_masters) if hasattr(prs, "slide_masters") else 1,
                has_theme=has_theme,
            )

            return template_info.model_dump_json()
        except Exception as e:
            logger.error(f"Failed to analyze template: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_analyze_template_variants(template_name: str) -> str:
        """
        Analyze a template's layouts and detect variants.

        Groups similar layouts together (e.g., "Content with Picture", "Content with Picture 2")
        and provides a structured view of unique layouts vs. variants. This is useful for
        understanding which layouts are variations of the same base design.

        Args:
            template_name: Name of the template to analyze

        Returns:
            JSON string with LayoutAnalysis model containing grouped layouts

        Example:
            await pptx_analyze_template_variants(template_name="brand_proposal")
        """
        try:
            from .extraction import analyze_layout_variants

            # First check if it's a builtin template
            template_data = await template_manager.get_template_data(template_name)
            if template_data:
                # Load builtin template
                buffer = io.BytesIO(template_data)
                prs = Presentation(buffer)
            else:
                # Get from artifact store
                result = await manager.get(template_name)
                if not result:
                    from ...models import ErrorResponse

                    return ErrorResponse(
                        error=f"Template not found: {template_name}"
                    ).model_dump_json()
                prs, metadata = result

            # Analyze layout variants
            analysis = analyze_layout_variants(prs)

            return analysis.model_dump_json()
        except Exception as e:
            logger.error(f"Failed to analyze template variants: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    return {
        "pptx_analyze_template": pptx_analyze_template,
        "pptx_analyze_template_variants": pptx_analyze_template_variants,
    }
