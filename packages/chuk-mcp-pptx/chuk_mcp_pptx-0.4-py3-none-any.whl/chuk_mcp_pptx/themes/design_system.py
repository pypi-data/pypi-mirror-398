"""
Design System Resolution for Universal Component API

Resolves design system properties with proper priority hierarchy:
1. Template design system (extracted from template)
2. Placeholder styles (inherited from placeholder)
3. Presentation theme (fallback)
4. Explicit theme override (parameter)
5. Individual property overrides (params)
"""

import logging
from typing import Any
from dataclasses import dataclass, field

logger = logging.getLogger(__name__)


@dataclass
class ResolvedDesignSystem:
    """
    Resolved design system with all properties.

    This represents the final computed design system after applying
    all priority levels.
    """

    # Colors
    primary_color: str = "#4F46E5"
    secondary_color: str = "#818CF8"
    background_color: str = "#FFFFFF"
    text_color: str = "#1F2937"
    border_color: str = "#E5E7EB"

    # Typography
    font_family: str = "Calibri"
    font_size: int = 14
    font_bold: bool = False
    font_italic: bool = False

    # Spacing
    padding: float = 0.2  # inches
    margin: float = 0.1
    gap: float = 0.15

    # Shape
    border_radius: float = 0.1
    border_width: float = 1.0

    # Source tracking (for debugging)
    source: str = "default"
    overrides: dict = field(default_factory=dict)


def extract_template_design_system(slide) -> ResolvedDesignSystem | None:
    """
    Extract design system from template/master slide.

    Analyzes the slide master and layouts to extract colors, fonts,
    and spacing that define the template's design system.
    """
    try:
        prs = slide.part.package.presentation_part.presentation

        if not prs.slide_masters:
            return None

        master = prs.slide_masters[0]

        # Extract theme colors from the Office theme
        design_system = ResolvedDesignSystem(source="template")

        # Try to extract theme colors from the XML theme
        if hasattr(master, "slide_master") or hasattr(master, "part"):
            try:
                # Access the theme from master's part
                theme_part = master.part.theme_part if hasattr(master.part, "theme_part") else None

                if theme_part and hasattr(theme_part, "theme_element"):
                    theme_elem = theme_part.theme_element

                    # Extract color scheme
                    if hasattr(theme_elem, "themeElements"):
                        theme_elements = theme_elem.themeElements
                        if hasattr(theme_elements, "clrScheme"):
                            clr_scheme = theme_elements.clrScheme

                            # Extract accent colors (commonly used in templates)
                            def extract_color_from_scheme(color_elem):
                                """Extract RGB from color scheme element."""
                                if hasattr(color_elem, "srgbClr"):
                                    rgb_hex = color_elem.srgbClr.val
                                    return f"#{rgb_hex}"
                                elif hasattr(color_elem, "sysClr"):
                                    # System color - use last color value
                                    return f"#{color_elem.sysClr.lastClr}"
                                return None

                            # Try to get accent colors
                            if hasattr(clr_scheme, "accent1"):
                                color = extract_color_from_scheme(clr_scheme.accent1)
                                if color:
                                    design_system.primary_color = color

                            if hasattr(clr_scheme, "accent2"):
                                color = extract_color_from_scheme(clr_scheme.accent2)
                                if color:
                                    design_system.secondary_color = color

                            # Background and text colors
                            if hasattr(clr_scheme, "lt1"):  # Light 1 = background
                                color = extract_color_from_scheme(clr_scheme.lt1)
                                if color:
                                    design_system.background_color = color

                            if hasattr(clr_scheme, "dk1"):  # Dark 1 = text
                                color = extract_color_from_scheme(clr_scheme.dk1)
                                if color:
                                    design_system.text_color = color

                            design_system.source = "template_theme_colors"
                            logger.debug(
                                f"Extracted theme colors: primary={design_system.primary_color}, bg={design_system.background_color}"
                            )

            except Exception as e:
                logger.debug(f"Could not extract theme colors from XML: {e}")

        # Extract typography from master shapes
        if hasattr(master, "shapes"):
            for shape in master.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    # Extract font from first paragraph
                    if shape.text_frame.paragraphs:
                        para = shape.text_frame.paragraphs[0]
                        if para.runs:
                            run = para.runs[0]
                            if hasattr(run.font, "name") and run.font.name:
                                design_system.font_family = run.font.name
                            if hasattr(run.font, "size") and run.font.size:
                                design_system.font_size = int(run.font.size.pt)
                            break

        return design_system

    except Exception as e:
        logger.warning(f"Could not extract template design system: {e}")
        return None


def extract_placeholder_styles(placeholder) -> dict[str, Any]:
    """
    Extract styles from a placeholder.

    Returns dictionary of style properties that can override
    the design system.
    """
    styles = {}

    try:
        # Extract text styles if placeholder has text
        if hasattr(placeholder, "text_frame"):
            text_frame = placeholder.text_frame
            if text_frame.paragraphs:
                para = text_frame.paragraphs[0]
                if para.runs:
                    run = para.runs[0]
                    font = run.font

                    if hasattr(font, "name") and font.name:
                        styles["font_family"] = font.name
                    if hasattr(font, "size") and font.size:
                        styles["font_size"] = int(font.size.pt)
                    if hasattr(font, "bold"):
                        styles["font_bold"] = font.bold
                    if hasattr(font, "italic"):
                        styles["font_italic"] = font.italic

                    # Extract color if available
                    if hasattr(font, "color") and hasattr(font.color, "rgb"):
                        try:
                            rgb = font.color.rgb
                            styles["text_color"] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                        except Exception:
                            pass

        # Extract fill color
        if hasattr(placeholder, "fill"):
            fill = placeholder.fill
            if hasattr(fill, "fore_color") and hasattr(fill.fore_color, "rgb"):
                try:
                    rgb = fill.fore_color.rgb
                    styles["background_color"] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                except Exception:
                    pass

    except Exception as e:
        logger.warning(f"Could not extract placeholder styles: {e}")

    return styles


def resolve_design_system(
    slide, placeholder=None, theme=None, params: dict | None = None
) -> ResolvedDesignSystem:
    """
    Resolve design system with proper priority hierarchy.

    Priority (lowest to highest):
    1. Template design system (extracted from template)
    2. Placeholder styles (inherited from placeholder)
    3. Presentation theme (if no template)
    4. Explicit theme parameter (override)
    5. Individual property overrides in params

    Args:
        slide: The slide to render on
        placeholder: Optional placeholder to inherit styles from
        theme: Optional theme object or name to override with
        params: Optional dict of explicit property overrides

    Returns:
        ResolvedDesignSystem with all properties computed
    """
    params = params or {}

    # Start with default design system
    design_system = ResolvedDesignSystem(source="default")

    # Priority 1: Extract from template (if exists)
    template_ds = extract_template_design_system(slide)
    if template_ds:
        design_system = template_ds
        logger.debug("Applied template design system")

    # Priority 2: Inherit from placeholder (if targeting one)
    if placeholder:
        placeholder_styles = extract_placeholder_styles(placeholder)
        for key, value in placeholder_styles.items():
            if hasattr(design_system, key):
                setattr(design_system, key, value)
                design_system.overrides[key] = "placeholder"
        if placeholder_styles:
            logger.debug(f"Applied placeholder styles: {list(placeholder_styles.keys())}")

    # Priority 3: Presentation theme (if no template and theme provided)
    if theme:
        from .theme_manager import ThemeManager

        theme_manager = ThemeManager()

        # Get theme object
        if isinstance(theme, str):
            theme_obj = theme_manager.get_theme(theme)
        else:
            theme_obj = theme

        if theme_obj:
            # Apply theme properties
            if hasattr(theme_obj, "colors"):
                design_system.primary_color = theme_obj.colors.get(
                    "primary", design_system.primary_color
                )
                design_system.secondary_color = theme_obj.colors.get(
                    "secondary", design_system.secondary_color
                )
                design_system.background_color = theme_obj.colors.get(
                    "background", design_system.background_color
                )
                design_system.text_color = theme_obj.colors.get("text", design_system.text_color)

            if hasattr(theme_obj, "typography"):
                design_system.font_family = theme_obj.typography.get(
                    "font_family", design_system.font_family
                )
                design_system.font_size = theme_obj.typography.get(
                    "font_size", design_system.font_size
                )

            design_system.source = f"theme:{theme if isinstance(theme, str) else 'custom'}"
            logger.debug(f"Applied theme: {design_system.source}")

    # Priority 4: Individual property overrides from params
    override_mappings = {
        "bg_color": "background_color",
        "background_color": "background_color",
        "color": "primary_color",
        "primary_color": "primary_color",
        "text_color": "text_color",
        "border_color": "border_color",
        "font_family": "font_family",
        "font_size": "font_size",
        "font_bold": "font_bold",
        "font_italic": "font_italic",
        "padding": "padding",
        "margin": "margin",
        "gap": "gap",
        "border_radius": "border_radius",
        "border_width": "border_width",
    }

    for param_key, ds_key in override_mappings.items():
        if param_key in params:
            setattr(design_system, ds_key, params[param_key])
            design_system.overrides[ds_key] = "explicit"

    if design_system.overrides:
        logger.debug(f"Applied explicit overrides: {list(design_system.overrides.keys())}")

    return design_system


def apply_design_system_to_shape(shape, design_system: ResolvedDesignSystem):
    """
    Apply resolved design system to a shape.

    Args:
        shape: python-pptx shape to style
        design_system: Resolved design system to apply
    """
    try:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        # Apply fill color
        if hasattr(shape, "fill"):
            shape.fill.solid()
            rgb = _hex_to_rgb(design_system.background_color)
            shape.fill.fore_color.rgb = RGBColor(*rgb)

        # Apply line/border
        if hasattr(shape, "line"):
            shape.line.color.rgb = RGBColor(*_hex_to_rgb(design_system.border_color))
            shape.line.width = Pt(design_system.border_width)

        # Apply text styles
        if hasattr(shape, "text_frame"):
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = design_system.font_family
                    run.font.size = Pt(design_system.font_size)
                    run.font.bold = design_system.font_bold
                    run.font.italic = design_system.font_italic
                    run.font.color.rgb = RGBColor(*_hex_to_rgb(design_system.text_color))

        logger.debug(f"Applied design system to shape (source: {design_system.source})")

    except Exception as e:
        logger.warning(f"Could not fully apply design system to shape: {e}")


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip("#")
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return (r, g, b)
