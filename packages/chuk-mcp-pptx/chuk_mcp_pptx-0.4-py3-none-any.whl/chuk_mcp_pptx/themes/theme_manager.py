# src/chuk_mcp_pptx/themes/theme_manager.py
"""
Theme manager for PowerPoint presentations.
Central system for managing and applying themes.

The theme system provides:
- Built-in themes (dark/light variants, special themes)
- Custom theme creation
- Theme registration and discovery
- Application to slides and components
- Export/import for sharing

Usage:
    from chuk_mcp_pptx.themes import ThemeManager

    # Get theme manager
    mgr = ThemeManager()

    # List available themes
    themes = mgr.list_themes()

    # Get and apply a theme
    theme = mgr.get_theme("dark-violet")
    theme.apply_to_slide(slide)

    # Create custom theme
    custom = Theme("my-theme", primary_hue="emerald", mode="dark")
    mgr.register_theme(custom)
"""

from typing import Dict, Any, Optional, List
import json
from pptx.util import Pt
from pptx.dml.color import RGBColor

from ..tokens.colors import get_semantic_tokens, GRADIENTS, PALETTE


class ThemeManager:
    """
    Manages themes for PowerPoint presentations.
    Provides theme registration, selection, and application.
    """

    def __init__(self):
        """Initialize theme manager with built-in themes."""
        self.themes = {}
        self._register_builtin_themes()
        self.current_theme = None

    def _register_builtin_themes(self):
        """Register all built-in themes."""
        # Dark themes
        self.register_theme(Theme("dark", primary_hue="blue", mode="dark"))
        self.register_theme(Theme("dark-blue", primary_hue="blue", mode="dark"))
        self.register_theme(Theme("dark-violet", primary_hue="violet", mode="dark"))
        self.register_theme(Theme("dark-green", primary_hue="emerald", mode="dark"))
        self.register_theme(Theme("dark-orange", primary_hue="orange", mode="dark"))
        self.register_theme(Theme("dark-red", primary_hue="red", mode="dark"))
        self.register_theme(Theme("dark-pink", primary_hue="pink", mode="dark"))
        self.register_theme(Theme("dark-purple", primary_hue="purple", mode="dark"))

        # Light themes
        self.register_theme(Theme("light", primary_hue="blue", mode="light"))
        self.register_theme(Theme("light-blue", primary_hue="blue", mode="light"))
        self.register_theme(Theme("light-violet", primary_hue="violet", mode="light"))
        self.register_theme(Theme("light-green", primary_hue="emerald", mode="light"))
        self.register_theme(Theme("light-orange", primary_hue="orange", mode="light"))
        self.register_theme(Theme("light-warm", primary_hue="amber", mode="light"))

        # Special themes
        self.register_theme(CyberpunkTheme())
        self.register_theme(GradientTheme("sunset", GRADIENTS["sunset"]))
        self.register_theme(GradientTheme("ocean", GRADIENTS["ocean"]))
        self.register_theme(GradientTheme("aurora", GRADIENTS["aurora"]))
        self.register_theme(MinimalTheme())
        self.register_theme(CorporateTheme())

    def register_theme(self, theme: "Theme"):
        """Register a theme."""
        self.themes[theme.name] = theme

    def get_theme(self, name: str) -> Optional["Theme"]:
        """Get theme by name."""
        return self.themes.get(name)

    def get_default_theme(self) -> "Theme":
        """Get the default theme."""
        theme = self.get_theme("dark")
        if theme is None:
            raise ValueError("Default theme 'dark' not found")
        return theme

    def set_current_theme(self, name: str):
        """Set the current active theme."""
        theme = self.get_theme(name)
        if theme:
            self.current_theme = theme
            return theme
        raise ValueError(f"Theme '{name}' not found")

    def list_themes(self) -> List[str]:
        """List all available theme names."""
        return list(self.themes.keys())

    def list_themes_by_mode(self, mode: str) -> List[str]:
        """List themes filtered by mode (dark/light)."""
        return [name for name, theme in self.themes.items() if theme.mode == mode]

    def get_theme_info(self, name: str) -> Optional[Dict[str, Any]]:
        """
        Get theme information as a dictionary.

        Args:
            name: Theme name

        Returns:
            Dictionary with theme info or None
        """
        theme = self.get_theme(name)
        if not theme:
            return None

        return {
            "name": theme.name,
            "mode": theme.mode,
            "primary_hue": theme.primary_hue,
            "font_family": theme.font_family,
            "colors": {
                "background": theme.background,
                "foreground": theme.foreground,
                "primary": theme.primary,
                "secondary": theme.secondary,
                "accent": theme.accent,
                "chart": theme.chart,
            },
        }

    def export_theme(self, name: str) -> Optional[str]:
        """
        Export theme as JSON string.

        Args:
            name: Theme name

        Returns:
            JSON string or None
        """
        info = self.get_theme_info(name)
        if info:
            return json.dumps(info, indent=2)
        return None

    def export_all_themes(self) -> str:
        """Export all themes as JSON."""
        all_themes = {name: self.get_theme_info(name) for name in self.list_themes()}
        return json.dumps(all_themes, indent=2)

    def apply_to_slide(self, slide, theme_name: Optional[str] = None):
        """
        Apply theme to a slide.

        Args:
            slide: PowerPoint slide object
            theme_name: Theme name or None for current theme
        """
        theme = self.get_theme(theme_name) if theme_name else self.current_theme
        if not theme:
            theme = self.get_theme("dark")  # Default

        if theme:
            theme.apply_to_slide(slide)


class Theme:
    """
    Base theme class.
    """

    def __init__(
        self, name: str, primary_hue: str = "blue", mode: str = "dark", font_family: str = "Inter"
    ):
        """
        Initialize theme.

        Args:
            name: Theme name
            primary_hue: Primary color hue
            mode: Color mode (dark/light)
            font_family: Primary font family
        """
        self.name = name
        self.primary_hue = primary_hue
        self.mode = mode
        self.font_family = font_family
        self.tokens = get_semantic_tokens(primary_hue, mode)

    # Properties to expose tokens as direct attributes for compatibility
    @property
    def background(self):
        return self.tokens.get("background", {})

    @property
    def foreground(self):
        return self.tokens.get("foreground", {})

    @property
    def primary(self):
        return self.tokens.get("primary", {})

    @property
    def secondary(self):
        return self.tokens.get("secondary", {})

    @property
    def accent(self):
        return self.tokens.get("accent", {})

    @property
    def chart(self):
        return self.tokens.get("chart", [])

    def hex_to_rgb(self, hex_color: str) -> tuple:
        """Convert hex to RGB."""
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

    def get_color(self, path: str) -> RGBColor:
        """Get color from tokens."""
        parts = path.split(".")
        value = self.tokens

        for part in parts:
            if isinstance(value, dict):
                value = value.get(part, "#000000")
            else:
                break

        if isinstance(value, str):
            return RGBColor(*self.hex_to_rgb(value))
        return RGBColor(0, 0, 0)

    def apply_to_slide(self, slide, override_text_colors: bool = True):
        """
        Apply theme to slide.

        Args:
            slide: PowerPoint slide object
            override_text_colors: If True, override all text colors with theme foreground.
                                 If False, only set background (useful for slides with pre-styled components).
        """
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.get_color("background.DEFAULT")

        # Optionally set default text color for all existing text shapes
        if override_text_colors:
            foreground_color = self.get_color("foreground.DEFAULT")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    # Set color for existing text
                    for paragraph in shape.text_frame.paragraphs:
                        # Set default font for paragraph
                        paragraph.font.color.rgb = foreground_color
                        # Also set for any existing runs
                        for run in paragraph.runs:
                            run.font.color.rgb = foreground_color

    def apply_to_shape(self, shape, style: str = "card"):
        """Apply theme to shape."""
        style_map = {
            "card": ("card.DEFAULT", "card.foreground", "border.DEFAULT"),
            "primary": ("primary.DEFAULT", "primary.foreground", "primary.DEFAULT"),
            "secondary": ("secondary.DEFAULT", "secondary.foreground", "border.DEFAULT"),
            "accent": ("accent.DEFAULT", "accent.foreground", "accent.DEFAULT"),
            "muted": ("muted.DEFAULT", "muted.foreground", "border.secondary"),
        }

        bg_path, fg_path, border_path = style_map.get(style, style_map["card"])

        # Apply fill
        if hasattr(shape, "fill"):
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.get_color(bg_path)

        # Apply border
        if hasattr(shape, "line"):
            shape.line.color.rgb = self.get_color(border_path)
            shape.line.width = Pt(1)

        # Apply text color
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = self.get_color(fg_path)
                    run.font.name = self.font_family

    def get_chart_colors(self) -> List[RGBColor]:
        """Get chart colors for data visualization."""
        chart_colors = self.tokens.get("chart", [])
        return [RGBColor(*self.hex_to_rgb(color)) for color in chart_colors]

    def to_dict(self) -> Dict[str, Any]:
        """
        Convert theme to dictionary representation.

        Returns:
            Dictionary with theme configuration
        """
        return {
            "name": self.name,
            "mode": self.mode,
            "primary_hue": self.primary_hue,
            "font_family": self.font_family,
        }

    def export_json(self) -> str:
        """
        Export theme as JSON string.

        Returns:
            JSON representation of theme
        """
        return json.dumps(self.to_dict(), indent=2)

    @classmethod
    def from_dict(cls, config: Dict[str, Any]) -> "Theme":
        """
        Create theme from dictionary.

        Args:
            config: Theme configuration dictionary

        Returns:
            Theme instance
        """
        return cls(
            name=config.get("name", "custom"),
            primary_hue=config.get("primary_hue", "blue"),
            mode=config.get("mode", "dark"),
            font_family=config.get("font_family", "Inter"),
        )


class CyberpunkTheme(Theme):
    """Cyberpunk theme with neon colors."""

    def __init__(self):
        super().__init__("cyberpunk", primary_hue="violet", mode="dark")

        # Override with cyberpunk colors using PALETTE
        self.tokens.update(
            {
                "background": {"DEFAULT": PALETTE["slate"][950]},  # Very dark
                "foreground": {"DEFAULT": PALETTE["cyan"][400]},  # Bright cyan
                "primary": {
                    "DEFAULT": PALETTE["fuchsia"][500],  # Magenta
                    "foreground": PALETTE["slate"][950],
                },
                "accent": {
                    "DEFAULT": PALETTE["yellow"][400],  # Electric yellow
                    "foreground": PALETTE["slate"][950],
                },
                "border": {"DEFAULT": PALETTE["fuchsia"][500]},
                "chart": [
                    PALETTE["fuchsia"][500],  # Magenta
                    PALETTE["cyan"][400],  # Cyan
                    PALETTE["yellow"][400],  # Yellow
                    PALETTE["pink"][500],  # Pink
                    PALETTE["lime"][400],  # Green
                ],
            }
        )
        self.font_family = "Orbitron"


class GradientTheme(Theme):
    """Theme with gradient backgrounds."""

    def __init__(self, name: str, gradient_colors: List[str]):
        super().__init__(name, mode="dark")
        self.gradient_colors = gradient_colors

    def apply_to_slide(self, slide):
        """Apply gradient background (using first color as fallback)."""
        # PowerPoint doesn't easily support gradients via python-pptx
        # Use first color as solid background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*self.hex_to_rgb(self.gradient_colors[0]))


class MinimalTheme(Theme):
    """Minimal black and white theme."""

    def __init__(self):
        super().__init__("minimal", primary_hue="zinc", mode="light")
        self.tokens = {
            "background": {"DEFAULT": PALETTE["zinc"][50]},  # Pure white
            "foreground": {"DEFAULT": PALETTE["zinc"][950]},  # Pure black
            "card": {
                "DEFAULT": PALETTE["zinc"][100],
                "foreground": PALETTE["zinc"][950],
            },
            "primary": {
                "DEFAULT": PALETTE["zinc"][950],
                "foreground": PALETTE["zinc"][50],
            },
            "secondary": {
                "DEFAULT": PALETTE["zinc"][200],
                "foreground": PALETTE["zinc"][900],
            },
            "accent": {
                "DEFAULT": PALETTE["zinc"][950],
                "foreground": PALETTE["zinc"][50],
            },
            "muted": {
                "DEFAULT": PALETTE["zinc"][100],
                "foreground": PALETTE["zinc"][600],
            },
            "border": {
                "DEFAULT": PALETTE["zinc"][300],
                "secondary": PALETTE["zinc"][400],
            },
            "chart": [
                PALETTE["zinc"][950],
                PALETTE["zinc"][800],
                PALETTE["zinc"][600],
                PALETTE["zinc"][400],
                PALETTE["zinc"][300],
            ],
        }
        self.font_family = "Helvetica Neue"


class CorporateTheme(Theme):
    """Professional corporate theme."""

    def __init__(self):
        super().__init__("corporate", primary_hue="blue", mode="light")
        self.tokens.update(
            {
                "background": {"DEFAULT": PALETTE["slate"][50]},
                "foreground": {"DEFAULT": PALETTE["slate"][900]},
                "primary": {
                    "DEFAULT": PALETTE["blue"][600],
                    "foreground": PALETTE["slate"][50],
                },
                "secondary": {
                    "DEFAULT": PALETTE["slate"][200],
                    "foreground": PALETTE["slate"][800],
                },
                "accent": {
                    "DEFAULT": PALETTE["green"][500],
                    "foreground": PALETTE["slate"][50],
                },
                "chart": [
                    PALETTE["blue"][600],
                    PALETTE["green"][500],
                    PALETTE["orange"][500],
                    PALETTE["violet"][500],
                    PALETTE["teal"][500],
                    PALETTE["red"][500],
                ],
            }
        )
        self.font_family = "Segoe UI"
