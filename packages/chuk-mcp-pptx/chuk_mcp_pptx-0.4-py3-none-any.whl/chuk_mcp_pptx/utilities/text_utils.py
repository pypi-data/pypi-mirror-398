"""
Text Utilities for PowerPoint MCP Server

Provides text extraction, formatting, and validation utilities.

Note: Text creation functions have been moved to components.core.text (TextBox, BulletList).
This module contains only extraction and utility functions.
"""

from typing import Dict, Any, Tuple, Optional
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor


def extract_slide_text(slide) -> Dict[str, Any]:
    """
    Extract all text content from a slide.

    Args:
        slide: PowerPoint slide object

    Returns:
        Dictionary containing extracted text organized by type
    """
    result: Dict[str, Any] = {
        "title": None,
        "subtitle": None,
        "body_text": [],
        "placeholders": {},
        "text_boxes": [],
        "table_text": [],
        "all_text": [],
    }

    # Extract from shapes
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text_frame.text.strip()
        if not text:
            continue

        result["all_text"].append(text)

        # Check if it's title
        if shape == slide.shapes.title:
            result["title"] = text
        # Check if it's a placeholder
        elif shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            result["placeholders"][str(placeholder_type)] = text

            # Check for subtitle
            if placeholder_type == 2:  # SUBTITLE
                result["subtitle"] = text
            elif placeholder_type == 7:  # BODY
                result["body_text"].append(text)
        else:
            # Regular text box
            result["text_boxes"].append(
                {
                    "text": text,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }
            )

    # Extract from tables
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            table_text = []
            for row_idx, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    row_text.append(cell_text)
                    result["all_text"].append(cell_text)
                table_text.append(row_text)
            result["table_text"].append(table_text)

    # Combine all text
    result["combined_text"] = "\n".join(result["all_text"])

    return result


def extract_presentation_text(prs) -> Dict[str, Any]:
    """
    Extract all text from a presentation.

    Args:
        prs: Presentation object

    Returns:
        Dictionary containing all text organized by slide
    """
    result: Dict[str, Any] = {
        "total_slides": len(prs.slides),
        "slides": [],
        "all_titles": [],
        "all_text": [],
    }

    for idx, slide in enumerate(prs.slides):
        slide_text = extract_slide_text(slide)
        slide_info = {"slide_index": idx, "slide_number": idx + 1, "text": slide_text}
        result["slides"].append(slide_info)

        if slide_text["title"]:
            result["all_titles"].append({"slide": idx + 1, "title": slide_text["title"]})

        result["all_text"].extend(slide_text["all_text"])

    # Create combined text
    result["combined_text"] = "\n\n".join(
        [
            f"=== Slide {s['slide_number']} ===\n{s['text']['combined_text']}"
            for s in result["slides"]
            if s["text"]["combined_text"]
        ]
    )

    return result


def format_text_frame(
    text_frame,
    font_name: Optional[str] = None,
    font_size: Optional[int] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    color: Optional[Tuple[int, int, int]] = None,
    alignment: Optional[str] = None,
):
    """
    Format a text frame with specified styling.

    Args:
        text_frame: Text frame to format
        font_name: Font family name
        font_size: Font size in points
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple
        alignment: Text alignment (left, center, right, justify)
    """
    for paragraph in text_frame.paragraphs:
        # Set alignment
        if alignment:
            if alignment.lower() == "left":
                paragraph.alignment = PP_ALIGN.LEFT
            elif alignment.lower() == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif alignment.lower() == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            elif alignment.lower() == "justify":
                paragraph.alignment = PP_ALIGN.JUSTIFY

        # Format font
        font = paragraph.font
        if font_name:
            font.name = font_name
        if font_size:
            font.size = Pt(font_size)
        if bold is not None:
            font.bold = bold
        if italic is not None:
            font.italic = italic
        if color:
            r, g, b = color
            font.color.rgb = RGBColor(r, g, b)


def validate_text_fit(
    shape, text: Optional[str] = None, font_size: Optional[int] = None
) -> Dict[str, Any]:
    """
    Validate if text fits within a shape and suggest adjustments.

    Args:
        shape: Shape containing the text
        text: Optional text to validate (uses existing text if not provided)
        font_size: Optional font size to test

    Returns:
        Dictionary with validation results and suggestions
    """
    if not shape.has_text_frame:
        return {"fits": False, "error": "Shape does not have text frame"}

    text_frame = shape.text_frame

    # Get or set text
    if text is not None:
        text_frame.text = text
    current_text = text_frame.text

    # Get shape dimensions
    width_inches = shape.width / 914400  # EMUs to inches
    height_inches = shape.height / 914400

    # Estimate text dimensions
    if font_size:
        test_font_size = font_size
    else:
        # Get current font size
        try:
            test_font_size = text_frame.paragraphs[0].font.size.pt
        except (AttributeError, IndexError):
            test_font_size = 18  # Default

    # Simple estimation (rough approximation)
    chars_per_line = int(width_inches * 72 / (test_font_size * 0.6))
    lines_needed = len(current_text) // chars_per_line + 1
    height_needed = lines_needed * test_font_size * 1.2 / 72  # in inches

    fits = height_needed <= height_inches

    result = {
        "fits": fits,
        "shape_width_inches": width_inches,
        "shape_height_inches": height_inches,
        "estimated_lines": lines_needed,
        "estimated_height_needed": height_needed,
        "current_font_size": test_font_size,
        "text_length": len(current_text),
    }

    if not fits:
        # Suggest adjustments
        suggested_font_size = int(test_font_size * height_inches / height_needed)
        result["suggestions"] = {
            "reduce_font_size_to": max(suggested_font_size, 8),
            "reduce_text_by": int((1 - height_inches / height_needed) * len(current_text)),
            "increase_shape_height_to": height_needed,
        }

    return result


def auto_fit_text(shape, min_font_size: int = 10, max_font_size: int = 44):
    """
    Automatically adjust font size to fit text within shape.

    Args:
        shape: Shape containing the text
        min_font_size: Minimum allowed font size
        max_font_size: Maximum allowed font size
    """
    if not shape.has_text_frame:
        return

    text_frame = shape.text_frame

    # Enable auto-fit
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    # Set margins
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)

    # Try to validate and adjust font size
    current_text = text_frame.text
    if not current_text:
        return

    # Binary search for optimal font size
    left, right = min_font_size, max_font_size
    optimal_size = max_font_size

    while left <= right:
        mid = (left + right) // 2
        validation = validate_text_fit(shape, font_size=mid)

        if validation["fits"]:
            optimal_size = mid
            left = mid + 1
        else:
            right = mid - 1

    # Apply optimal font size
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(optimal_size)
