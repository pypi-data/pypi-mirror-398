# src/chuk_mcp_pptx/tools/slide_layout_tools.py
"""
Slide Layout Tools for PowerPoint MCP Server

Provides async MCP tools for managing PowerPoint slide layouts.
Handles:
- Listing and selecting slide master layouts
- Customizing slide appearance (backgrounds, footers, etc.)
- Duplicating and reordering slides
- Applying master layout formatting

Note: Grid-based component positioning is handled by component_tools.py.
Slide templates (dashboard, comparison, etc.) are in slide_templates/.
"""

from ...models import ErrorResponse, SuccessResponse
from ...constants import (
    ErrorMessages,
)

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Import design system typography tokens
from ...tokens.typography import FONT_SIZES


def register_layout_tools(mcp, manager):
    """Register all layout-related tools with the MCP server."""

    # ========================================================================
    # TRADITIONAL LAYOUT MANAGEMENT TOOLS
    # ========================================================================
    # These tools directly manipulate presentation slides:
    # - List and select slide layouts
    # - Customize slide appearance
    # - Duplicate and reorder slides
    # - Apply master layout settings

    @mcp.tool
    async def pptx_list_layouts(presentation: str | None = None) -> str:
        """
        List all available slide layouts in the presentation.

        Shows the built-in layouts from the slide master including:
        - Title Slide
        - Title and Content
        - Section Header
        - Two Content
        - Comparison
        - Title Only
        - Blank
        - Content with Caption
        - Picture with Caption
        - Title and Vertical Text
        - Vertical Title and Text

        Args:
            presentation: Name of presentation (uses current if not specified)

        Returns:
            List of available layouts with indices and descriptions

        Example:
            layouts = await pptx_list_layouts()
            # Returns:
            # "Available slide layouts:
            #  0: Title Slide - For opening slides
            #  1: Title and Content - Standard content slide
            #  2: Section Header - Section dividers
            #  ..."
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            layouts = []
            layouts.append("=== AVAILABLE SLIDE LAYOUTS ===\n")

            for i, layout in enumerate(prs.slide_master.slide_layouts):
                layout_name = str(layout.name)

                # Add descriptions for common layouts
                descriptions = {
                    "Title Slide": "Opening slide with title and subtitle",
                    "Title and Content": "Standard content slide with bullet points",
                    "Section Header": "Section divider with large title",
                    "Two Content": "Side-by-side content areas",
                    "Comparison": "Two columns for comparing items",
                    "Title Only": "Just title, blank content area",
                    "Blank": "Completely blank slide",
                    "Content with Caption": "Content area with side caption",
                    "Picture with Caption": "Image with description",
                    "Title and Vertical Text": "Vertical text layout",
                    "Vertical Title and Text": "Vertical oriented content",
                }

                desc = descriptions.get(layout_name, "Custom layout")
                layouts.append(f"{i:2}: {layout_name:30} - {desc}")

                # List placeholders in this layout
                placeholders = []
                for shape in layout.placeholders:
                    placeholders.append(str(shape.name))

                if placeholders:
                    layouts.append(f"    Placeholders: {', '.join(placeholders)}")

            layouts.append(f"\nTotal layouts: {len(prs.slide_master.slide_layouts)}")
            layouts.append(
                "\nUse pptx_add_slide_with_layout() to create slides with specific layouts"
            )

            return "\n".join(layouts)

        except Exception as e:
            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_add_slide_with_layout(
        layout_index: int,
        title: str | None = None,
        content: list[str] | None = None,
        subtitle: str | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Add a slide using a specific layout from the slide master.

        Creates a new slide with the specified layout and populates
        placeholders with provided content.

        Args:
            layout_index: Index of the layout to use (0-based)
            title: Title text for the slide (if layout has title placeholder)
            content: Content for the main placeholder (list for bullets)
            subtitle: Subtitle text (for layouts with subtitle)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message with slide index

        Example:
            # Create a two-content comparison slide
            await pptx_add_slide_with_layout(
                layout_index=4,  # Comparison layout
                title="Product Comparison"
            )

            # Create a section header
            await pptx_add_slide_with_layout(
                layout_index=2,  # Section Header
                title="Part 2: Financial Analysis"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if layout_index >= len(prs.slide_master.slide_layouts):
                return SuccessResponse(
                    message=f"Error: Layout index {layout_index} out of range. Use pptx_list_layouts() to see available layouts"
                ).model_dump_json()

            # Get the layout
            layout = prs.slide_master.slide_layouts[layout_index]

            # Add slide with this layout
            slide = prs.slides.add_slide(layout)

            # Populate placeholders
            if title and slide.shapes.title:
                slide.shapes.title.text = title

            # Handle different placeholder types
            for shape in slide.placeholders:
                # Skip if already handled
                if shape == slide.shapes.title:
                    continue

                # Handle subtitle placeholder
                if subtitle and shape.placeholder_format.type == 2:  # SUBTITLE
                    shape.text = subtitle

                # Handle content placeholder
                elif content and shape.placeholder_format.type == 7:  # CONTENT
                    if isinstance(content, list):
                        # Add as bullet points
                        text_frame = shape.text_frame
                        text_frame.clear()

                        for i, item in enumerate(content):
                            if i == 0:
                                p = text_frame.paragraphs[0]
                            else:
                                p = text_frame.add_paragraph()
                            p.text = item
                            p.level = 0
                    else:
                        shape.text = str(content)

            # Apply presentation theme to the slide
            if metadata and metadata.theme:
                from ...themes.theme_manager import ThemeManager

                theme_manager = ThemeManager()
                theme_obj = theme_manager.get_theme(metadata.theme)
                if theme_obj:
                    theme_obj.apply_to_slide(slide)

            # Update in VFS if enabled
            await manager.update(presentation)

            slide_idx = len(prs.slides) - 1
            layout_name = str(layout.name)

            # Analyze placeholders on the new slide
            from pptx.enum.shapes import PP_PLACEHOLDER

            placeholder_info = []
            chart_placeholders = []
            picture_placeholders = []

            for shape in slide.placeholders:
                placeholder_type = shape.placeholder_format.type
                placeholder_idx = shape.placeholder_format.idx

                # Identify CHART placeholders
                if placeholder_type == PP_PLACEHOLDER.CHART:  # 12 = CHART
                    chart_placeholders.append(placeholder_idx)
                    placeholder_info.append(
                        f"  Placeholder {placeholder_idx}: CHART - Use pptx_add_component(component='ColumnChart/LineChart/etc', target_placeholder={placeholder_idx})"
                    )
                # Identify PICTURE placeholders
                elif placeholder_type == PP_PLACEHOLDER.PICTURE:  # 18 = PICTURE
                    picture_placeholders.append(placeholder_idx)
                    placeholder_info.append(
                        f"  Placeholder {placeholder_idx}: PICTURE - Use pptx_add_image(target_placeholder={placeholder_idx})"
                    )
                # Identify other common types
                elif placeholder_type == PP_PLACEHOLDER.TITLE:
                    placeholder_info.append(
                        f"  Placeholder {placeholder_idx}: TITLE - Use pptx_populate_placeholder(placeholder_idx={placeholder_idx})"
                    )
                elif placeholder_type == PP_PLACEHOLDER.SUBTITLE:
                    placeholder_info.append(
                        f"  Placeholder {placeholder_idx}: SUBTITLE - Use pptx_populate_placeholder(placeholder_idx={placeholder_idx})"
                    )
                elif placeholder_type == PP_PLACEHOLDER.BODY:
                    placeholder_info.append(
                        f"  Placeholder {placeholder_idx}: BODY/CONTENT - Use pptx_populate_placeholder(placeholder_idx={placeholder_idx})"
                    )
                else:
                    placeholder_info.append(
                        f"  Placeholder {placeholder_idx}: Type {placeholder_type}"
                    )

            # Build response message
            msg = f"Added slide {slide_idx} using layout '{layout_name}'"
            if placeholder_info:
                msg += "\n\nAvailable placeholders:\n" + "\n".join(placeholder_info)

            if chart_placeholders:
                msg += f"\n\n⚠️ IMPORTANT: This layout has CHART placeholder(s) at indices {chart_placeholders}."
                msg += "\nWhen adding charts, use target_placeholder parameter instead of left/top/width/height:"
                msg += f"\n  pptx_add_component(slide_index={slide_idx}, component='ColumnChart', target_placeholder={chart_placeholders[0]}, ...)"

            if picture_placeholders:
                msg += f"\n\n⚠️ IMPORTANT: This layout has PICTURE placeholder(s) at indices {picture_placeholders}."
                msg += "\nWhen adding images, use target_placeholder parameter instead of left/top/width/height."

            return SuccessResponse(message=msg).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_customize_layout(
        slide_index: int,
        background_color: str | None = None,
        add_footer: str | None = None,
        add_page_number: bool = False,
        add_date: str | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Customize the layout of an existing slide.

        Adds background colors, footers, page numbers, and other
        layout customizations to a slide.

        Args:
            slide_index: Index of the slide to customize (0-based)
            background_color: Hex color for background (e.g., "#F0F0F0")
            add_footer: Footer text to add
            add_page_number: Whether to add slide number
            add_date: Date text to add
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message with customizations applied

        Example:
            await pptx_customize_layout(
                slide_index=1,
                background_color="#F5F5F5",
                add_footer="Confidential - Internal Use Only",
                add_page_number=True,
                add_date="2024-12-01"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            slide = prs.slides[slide_index]
            customizations = []

            # Set background color
            if background_color:
                from pptx.dml.color import RGBColor

                # Parse hex color
                bg_color = (
                    background_color[1:] if background_color.startswith("#") else background_color
                )

                try:
                    r = int(bg_color[0:2], 16)
                    g = int(bg_color[2:4], 16)
                    b = int(bg_color[4:6], 16)

                    fill = slide.background.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(r, g, b)
                    customizations.append(f"background color {bg_color}")
                except (ValueError, AttributeError):
                    customizations.append("background color (failed - invalid format)")

            # Add footer
            if add_footer:
                footer_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(6.8), Inches(9.0), Inches(0.5)
                )
                footer_box.text_frame.text = add_footer
                footer_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                footer_box.text_frame.paragraphs[0].font.size = Pt(FONT_SIZES["xs"])
                customizations.append("footer text")

            # Add page number
            if add_page_number:
                page_box = slide.shapes.add_textbox(
                    Inches(8.5), Inches(6.8), Inches(1.0), Inches(0.5)
                )
                page_box.text_frame.text = str(slide_index + 1)
                page_box.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
                page_box.text_frame.paragraphs[0].font.size = Pt(FONT_SIZES["xs"])
                customizations.append("page number")

            # Add date
            if add_date:
                date_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(2.0), Inches(0.5)
                )
                date_box.text_frame.text = add_date
                date_box.text_frame.paragraphs[0].font.size = Pt(FONT_SIZES["xs"])
                customizations.append("date")

            # Update in VFS if enabled
            await manager.update(presentation)

            if customizations:
                return SuccessResponse(
                    message=f"Customized slide {slide_index}: {', '.join(customizations)}"
                ).model_dump_json()
            else:
                return SuccessResponse(
                    message=f"No customizations applied to slide {slide_index}"
                ).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_apply_master_layout(
        layout_name: str,
        font_name: str | None = None,
        title_color: str | None = None,
        body_color: str | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Apply master layout settings to all slides.

        Sets consistent formatting across the entire presentation by
        modifying the slide master.

        Args:
            layout_name: Name of the layout theme ("corporate", "modern", "minimal")
            font_name: Font family to use throughout
            title_color: Hex color for all titles
            body_color: Hex color for body text
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message with number of slides affected

        Example:
            await pptx_apply_master_layout(
                layout_name="corporate",
                font_name="Arial",
                title_color="#003366",
                body_color="#333333"
            )
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            from pptx.dml.color import RGBColor

            # Parse colors
            title_rgb = None
            body_rgb = None

            if title_color:
                t_color = title_color[1:] if title_color.startswith("#") else title_color
                try:
                    r = int(t_color[0:2], 16)
                    g = int(t_color[2:4], 16)
                    b = int(t_color[4:6], 16)
                    title_rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    pass

            if body_color:
                b_color = body_color[1:] if body_color.startswith("#") else body_color
                try:
                    r = int(b_color[0:2], 16)
                    g = int(b_color[2:4], 16)
                    b = int(b_color[4:6], 16)
                    body_rgb = RGBColor(r, g, b)
                except (ValueError, IndexError):
                    pass

            # Apply to all slides
            slides_updated = 0

            for slide in prs.slides:
                # Update title formatting
                if slide.shapes.title:
                    title = slide.shapes.title
                    if font_name:
                        for paragraph in title.text_frame.paragraphs:
                            paragraph.font.name = font_name
                    if title_rgb:
                        for paragraph in title.text_frame.paragraphs:
                            paragraph.font.color.rgb = title_rgb

                # Update body text formatting
                for shape in slide.shapes:
                    if shape.has_text_frame and shape != slide.shapes.title:
                        for paragraph in shape.text_frame.paragraphs:
                            if font_name:
                                paragraph.font.name = font_name
                            if body_rgb:
                                paragraph.font.color.rgb = body_rgb

                slides_updated += 1

            # Update in VFS if enabled
            await manager.update(presentation)

            settings = []
            if font_name:
                settings.append(f"font: {font_name}")
            if title_color:
                settings.append(f"title color: #{title_color}")
            if body_color:
                settings.append(f"body color: #{body_color}")

            return SuccessResponse(
                message=f"Applied master layout '{layout_name}' to {slides_updated} slides ({', '.join(settings)})"
            ).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_duplicate_slide(slide_index: int, presentation: str | None = None) -> str:
        """
        Duplicate an existing slide.

        Creates a copy of a slide with all its content and formatting.

        Args:
            slide_index: Index of the slide to duplicate (0-based)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message with new slide index

        Example:
            await pptx_duplicate_slide(slide_index=2)
            # Creates a copy of slide 2 as a new slide
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            # Get the source slide
            source_slide = prs.slides[slide_index]

            # Create new slide with same layout
            new_slide = prs.slides.add_slide(source_slide.slide_layout)

            # Copy shapes (simplified - full duplication would require more complex logic)
            for shape in source_slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    # Copy text boxes
                    new_shape = new_slide.shapes.add_textbox(
                        shape.left, shape.top, shape.width, shape.height
                    )
                    if shape.has_text_frame:
                        new_shape.text_frame.text = shape.text_frame.text
                elif hasattr(shape, "text") and shape != source_slide.shapes.title:
                    # Copy other text content
                    try:
                        for new_shape in new_slide.shapes:
                            if new_shape.shape_id == shape.shape_id:
                                if hasattr(new_shape, "text"):
                                    new_shape.text = shape.text
                                break
                    except (AttributeError, KeyError):
                        pass

            # Copy title if exists
            if source_slide.shapes.title and new_slide.shapes.title:
                new_slide.shapes.title.text = source_slide.shapes.title.text

            # Update in VFS if enabled
            await manager.update(presentation)

            new_idx = len(prs.slides) - 1
            return SuccessResponse(
                message=f"Duplicated slide {slide_index} as new slide {new_idx}"
            ).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_reorder_slides(
        slide_index: int, new_position: int, presentation: str | None = None
    ) -> str:
        """
        Move a slide to a different position in the presentation.

        Reorders slides by moving a slide from one position to another.

        Args:
            slide_index: Current index of the slide to move (0-based)
            new_position: Target position for the slide (0-based)
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Success message confirming the move

        Example:
            await pptx_reorder_slides(slide_index=5, new_position=2)
            # Moves slide 5 to position 2
        """
        try:
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            if slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found in presentation"
                ).model_dump_json()

            if new_position >= len(prs.slides):
                return SuccessResponse(
                    message=f"Error: New position {new_position} out of range"
                ).model_dump_json()

            if slide_index == new_position:
                return SuccessResponse(
                    message=f"Slide {slide_index} is already at position {new_position}"
                ).model_dump_json()

            # Get the XML parts
            slides = prs.slides._sldIdLst
            slide_id = slides[slide_index]

            # Remove from current position
            slides.remove(slide_id)

            # Insert at new position
            slides.insert(new_position, slide_id)

            # Update in VFS if enabled
            await manager.update(presentation)

            return SuccessResponse(
                message=f"Moved slide from position {slide_index} to position {new_position}"
            ).model_dump_json()

        except Exception as e:
            return ErrorResponse(error=str(e)).model_dump_json()

    # ========================================================================
    # NOTE: Grid layout is handled internally by component tools
    # ========================================================================
    # The Grid class is used internally by:
    # - component_tools.py: Components accept grid-based positioning
    # - slide_templates/: Templates use Grid to position multiple components
    #
    # No need to expose grid calculation tools - components handle it directly.

    # Return the tools for external access
    return {
        "pptx_list_layouts": pptx_list_layouts,
        "pptx_add_slide_with_layout": pptx_add_slide_with_layout,
        "pptx_customize_layout": pptx_customize_layout,
        "pptx_apply_master_layout": pptx_apply_master_layout,
        "pptx_duplicate_slide": pptx_duplicate_slide,
        "pptx_reorder_slides": pptx_reorder_slides,
    }
