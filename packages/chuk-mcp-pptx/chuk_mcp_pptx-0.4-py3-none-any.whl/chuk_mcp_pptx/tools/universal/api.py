"""
Universal Component API

Single unified tool for adding components anywhere:
- Free-form positioning
- Into placeholders (template-aware)
- Into other components (composition)
- Into layout regions (grid/flex)

All with automatic design system resolution.
"""

import logging
from typing import Any
from pptx.util import Inches

logger = logging.getLogger(__name__)


def register_universal_component_api(mcp, manager):
    """Register the universal component API tools."""

    @mcp.tool
    async def pptx_list_slide_components(
        slide_index: int,
        presentation: str | None = None,
    ) -> str:
        """
        List all components on a slide with validation for placeholders and images.

        This tool is CRITICAL for verifying slide layouts, especially when working with
        templates. It shows all shapes, text boxes, images, charts, and other components
        on a slide, PLUS validates that:
        - All placeholders have been populated with content
        - All images loaded successfully without errors
        - Content matches the layout's purpose

        ⚠️ MANDATORY AFTER POPULATING SLIDES:
        You MUST call this tool after populating placeholders to verify:
        1. No empty placeholders remain (they show "Click to add" in final presentation)
        2. All images loaded correctly (no broken image icons)
        3. Layout integrity is maintained

        Use this tool to:
        - Verify template layouts after adding slides
        - Check placeholder population results
        - Validate image loading
        - Identify existing components before adding new ones
        - Understand component positioning and relationships
        - Debug layout issues

        Args:
            slide_index: Index of the slide (0-based)
            presentation: Presentation name (uses current if not specified)

        Returns:
            JSON string with ComponentListResponse containing:
            - components: List of all components with:
                - component_id: Unique identifier
                - component_type: Type (TextBox, Shape, Image, Chart, Table, etc.)
                - position: X, Y coordinates and dimensions
                - target: Composition target info (parent, container relationships)
                - content_preview: First 50 chars of text content if applicable
            - placeholders: Status of all placeholders:
                - idx: Placeholder index
                - type: Placeholder type (TITLE, BODY, CHART, etc.)
                - name: Placeholder name
                - is_empty: Whether placeholder needs content
                - has_text: Whether placeholder has text
                - has_image: Whether placeholder has image
                - content_preview: Preview of content
            - images: Status of all images:
                - component_id: Image component ID
                - placeholder_idx: Placeholder index if in placeholder
                - loaded_successfully: Whether image loaded without errors
                - error_message: Error details if image failed
                - source: Image source (URL or path)
            - warnings: List of validation warnings:
                - type: Warning type (empty_placeholder, missing_image, layout_mismatch)
                - message: Detailed warning message
                - placeholder_idx: Related placeholder if applicable
                - component_id: Related component if applicable
            - validation_passed: Boolean - true if no warnings, false if issues found
            - component_count: Total number of components
            - slide_index: The slide that was queried

        BEST PRACTICE:
            Always call this after populating placeholders. Check the warnings array
            and validation_passed flag. If validation_passed is false, you MUST fix
            the issues before moving to the next slide.

        Example - Verify template slide with validation:
            # After adding and populating a slide from template
            result = await pptx_list_slide_components(slide_index=5)
            # Response shows:
            # {
            #   "slide_index": 5,
            #   "component_count": 3,
            #   "components": [...],
            #   "placeholders": [
            #     {"idx": 0, "type": "TITLE", "name": "Title 1", "is_empty": false,
            #      "has_text": true, "content_preview": "Our Product"},
            #     {"idx": 1, "type": "BODY", "name": "Content Placeholder 2", "is_empty": false,
            #      "has_text": true, "content_preview": "Key features include..."},
            #     {"idx": 10, "type": "PICTURE", "name": "Picture Placeholder 3", "is_empty": true,
            #      "has_text": false, "has_image": false}
            #   ],
            #   "images": [],
            #   "warnings": [
            #     {"type": "empty_placeholder", "message": "Placeholder 10 (PICTURE: Picture Placeholder 3) is empty and needs content",
            #      "placeholder_idx": 10}
            #   ],
            #   "validation_passed": false
            # }
            #
            # ⚠️ validation_passed is false! Must populate placeholder 10 before moving on

        Example - All validation passed:
            result = await pptx_list_slide_components(slide_index=3)
            # Response shows:
            # {
            #   "validation_passed": true,
            #   "warnings": [],
            #   "placeholders": [
            #     {"idx": 0, "is_empty": false, "has_text": true, ...},
            #     {"idx": 1, "is_empty": false, "has_text": true, ...}
            #   ],
            #   "images": [
            #     {"component_id": "img_0", "loaded_successfully": true, ...}
            #   ]
            # }
            # ✅ All good! Safe to move to next slide
        """
        try:
            from ...models import (
                ErrorResponse,
                ComponentListResponse,
                ComponentInfo,
                ComponentPosition,
                ComponentTarget,
                TargetType,
                PlaceholderStatus,
                ImageStatus,
                ValidationWarning,
            )
            from ...components.tracking import component_tracker
            from ...constants import ErrorMessages
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            # Get presentation
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            # Validate slide index
            if slide_index < 0 or slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found. Presentation has {len(prs.slides)} slides."
                ).model_dump_json()

            slide = prs.slides[slide_index]

            # Get all components on slide
            components = component_tracker.list_on_slide(metadata.name, slide_index)

            # Build component list
            component_list = []
            for comp in components:
                # Convert target_type string to TargetType enum
                target_type = (
                    TargetType(comp.target_type) if comp.target_type else TargetType.FREE_FORM
                )

                component_info = ComponentInfo(
                    id=comp.component_id,
                    type=comp.component_type,
                    position=ComponentPosition(
                        left=comp.left or 0.0,
                        top=comp.top or 0.0,
                        width=comp.width or 0.0,
                        height=comp.height or 0.0,
                    ),
                    target=ComponentTarget(
                        type=target_type,
                        id=comp.target_id,
                    ),
                    parent_id=comp.parent_id,
                    children=comp.children_ids,
                    params=comp.params,
                )
                component_list.append(component_info)

            # VALIDATION: Check all placeholders
            placeholder_statuses = []
            warnings = []

            for shape in slide.placeholders:
                try:
                    idx = shape.placeholder_format.idx
                    ph_type = shape.placeholder_format.type
                    type_name = ph_type.name if hasattr(ph_type, "name") else str(ph_type)

                    is_empty = True
                    has_text = False
                    has_image = False
                    content_preview = None

                    # Check for text content
                    if hasattr(shape, "text_frame") and hasattr(shape.text_frame, "text"):
                        text_content = shape.text_frame.text.strip()
                        if text_content:
                            has_text = True
                            is_empty = False
                            content_preview = (
                                text_content[:50] + "..."
                                if len(text_content) > 50
                                else text_content
                            )
                    elif hasattr(shape, "text"):
                        text_content = shape.text.strip()
                        if text_content:
                            has_text = True
                            is_empty = False
                            content_preview = (
                                text_content[:50] + "..."
                                if len(text_content) > 50
                                else text_content
                            )

                    # Check for image content
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        # Check if placeholder contains an image
                        if hasattr(shape, "image"):
                            has_image = True
                            is_empty = False
                        # Check for picture in placeholder
                        elif len(list(shape.element.iter())) > 1:
                            # Placeholder has child elements, likely not empty
                            for child in shape.element.iter():
                                if (
                                    "pic" in str(child.tag).lower()
                                    or "blip" in str(child.tag).lower()
                                ):
                                    has_image = True
                                    is_empty = False
                                    break

                    # Create status
                    status = PlaceholderStatus(
                        idx=idx,
                        type=type_name,
                        name=shape.name,
                        is_empty=is_empty,
                        has_text=has_text,
                        has_image=has_image,
                        content_preview=content_preview,
                    )
                    placeholder_statuses.append(status)

                    # Add warning if empty
                    if is_empty:
                        warnings.append(
                            ValidationWarning(
                                type="empty_placeholder",
                                message=f"Placeholder {idx} ({type_name}: {shape.name}) is empty and needs content",
                                placeholder_idx=idx,
                                component_id=None,
                            )
                        )

                except Exception as e:
                    logger.warning(f"Could not analyze placeholder: {e}")

            # VALIDATION: Check all images
            image_statuses = []

            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # Find corresponding component
                        comp_id = None
                        placeholder_idx = None
                        for comp in components:
                            if comp.component_type == "Image":
                                comp_id = comp.component_id
                                if comp.target_type == "placeholder":
                                    placeholder_idx = comp.target_id
                                break

                        # Check if image loaded successfully
                        loaded_successfully = True
                        error_message = None
                        source = None

                        # Try to get image info
                        try:
                            if hasattr(shape, "image"):
                                # Image loaded successfully
                                loaded_successfully = True
                            else:
                                # No image data
                                loaded_successfully = False
                                error_message = "Image shape has no image data"
                        except Exception as img_err:
                            loaded_successfully = False
                            error_message = str(img_err)

                        # Try to get source from component params
                        if comp_id:
                            for comp in components:
                                if comp.component_id == comp_id and "image_source" in comp.params:
                                    source = comp.params["image_source"]
                                    break

                        image_status = ImageStatus(
                            component_id=comp_id or f"image_{shape.shape_id}",
                            placeholder_idx=placeholder_idx,
                            loaded_successfully=loaded_successfully,
                            error_message=error_message,
                            source=source,
                        )
                        image_statuses.append(image_status)

                        # Add warning if image failed
                        if not loaded_successfully:
                            warnings.append(
                                ValidationWarning(
                                    type="missing_image",
                                    message=f"Image failed to load: {error_message or 'Unknown error'}",
                                    placeholder_idx=None,
                                    component_id=comp_id,
                                )
                            )

                    except Exception as e:
                        logger.warning(f"Could not analyze image: {e}")

            # Determine if validation passed
            validation_passed = len(warnings) == 0

            response = ComponentListResponse(
                slide_index=slide_index,
                component_count=len(components),
                components=component_list,
                placeholders=placeholder_statuses,
                images=image_statuses,
                warnings=warnings,
                validation_passed=validation_passed,
            )
            return response.model_dump_json()

        except Exception as e:
            logger.error(f"Failed to list components: {e}")
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_add_component(
        slide_index: int,
        component: str,
        target_placeholder: int | None = None,
        target_component: str | None = None,
        target_layout: str | None = None,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
        params: dict[str, Any] | None = None,
        component_id: str | None = None,
        theme: str | None = None,
        presentation: str | None = None,
        placeholder_as_layout: bool = False,
    ) -> str:
        """
        Add a component to a slide with universal targeting and design system awareness.

        This is the universal component insertion API that works everywhere:
        - Free-form positioning on slides
        - Into template placeholders (respects design system)
        - Into other components (composition/nesting)
        - Into layout regions (grid/flex positioning)

        Design System Priority (lowest to highest):
        1. Template design system (automatic)
        2. Placeholder styles (if target_placeholder)
        3. Presentation theme (fallback)
        4. Explicit theme parameter
        5. Individual property overrides in params

        ⚠️ CRITICAL: Content Placeholders (CHART, TABLE, PICTURE, OBJECT)
        When adding slides from templates, ALWAYS check the pptx_add_slide_from_template() response
        for content placeholders. The response explicitly identifies these:

            "🚨 CRITICAL: This layout has CHART placeholder(s) at indices [2]"
            "📊 CRITICAL: This layout has TABLE placeholder(s) at indices [14]"
            "📦 OBJECT placeholder(s) at indices [13] can hold tables, charts, or other content"

        For ALL content placeholders (CHART, TABLE, PICTURE, OBJECT):
        - MUST use target_placeholder parameter (e.g., target_placeholder=2)
        - DO NOT use left/top/width/height (this creates overlays with visible "Click to add" text!)
        - Component automatically inherits placeholder's position and size
        - OBJECT placeholders can accept any component type (Table, Chart, Image)

        WRONG (creates overlay - placeholder remains visible underneath):
            pptx_add_component(slide_index=3, component='Table',
                              left=0.6, top=2.0, width=8.4, height=2.4, ...)
            pptx_add_component(slide_index=5, component='ColumnChart',
                              left=1.0, top=2.0, width=8.0, height=3.5, ...)

        CORRECT (replaces placeholder - clean template design):
            pptx_add_component(slide_index=3, component='Table',
                              target_placeholder=14, params={...})
            pptx_add_component(slide_index=5, component='ColumnChart',
                              target_placeholder=2, params={...})

        BEST PRACTICE - Verify Placeholder Population:
        After adding components to placeholders, ALWAYS verify with:
            components = await pptx_list_slide_components(slide_index=X)

        This ensures:
        - Placeholders are properly populated (not overlaid)
        - Components match template layout expectations
        - Design system colors are correctly applied

        For template-based presentations, use this workflow:
        1. pptx_add_slide_from_template(layout_index=X) - Add slide with layout (READ the response!)
        2. Check response for CHART/PICTURE placeholder warnings
        3. pptx_add_component() - Populate placeholders using target_placeholder
        4. pptx_list_slide_components() - Verify layout matches template

        Args:
            slide_index: Index of the slide (0-based)
            component: Component type (e.g., "metric_card", "badge", "button", "code_block")

            Target modes (pick one):
            target_placeholder: Placeholder index to insert into (template-aware)
            target_component: Component ID to insert into (composition)
            target_layout: Layout region name (grid/flex positioning)

            Positioning (required for free-form, optional for targets):
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches

            Customization:
            params: Component-specific parameters (see pptx_get_component_schema)
            component_id: Unique ID for this component (for targeting later)
            theme: Theme name to override design system
            presentation: Presentation name (uses current if not specified)
            placeholder_as_layout: If True, create a Stack layout at placeholder bounds instead of direct population (allows multi-component placeholders)

        Returns:
            JSON string with component details and placement info

        Examples:
            # 1. Free-form (respects template design system)
            await pptx_add_component(
                slide_index=0,
                component="metric_card",
                left=2.0, top=3.0, width=2.5, height=1.5,
                params={"value": "$150K", "label": "Revenue"}
            )

            # 2. Into placeholder (template-aware, auto-sized)
            await pptx_add_component(
                slide_index=0,
                component="metric_card",
                target_placeholder=1,
                params={"value": "$150K", "label": "Revenue"}
            )

            # 3. Composition (badge in card)
            await pptx_add_component(
                slide_index=0,
                component="card",
                left=2.0, top=2.0, width=4.0, height=3.0,
                component_id="main_card",
                params={"title": "Dashboard"}
            )
            await pptx_add_component(
                slide_index=0,
                component="badge",
                target_component="main_card",
                params={"text": "New", "variant": "success"}
            )

            # 4. Theme override
            await pptx_add_component(
                slide_index=0,
                component="button",
                left=5.0, top=4.0,
                theme="dark-violet",  # Override design system
                params={"text": "Click Me"}
            )

            # 5. Multi-component placeholder (layout approach - RECOMMENDED)
            # First, create a Stack layout at the placeholder
            await pptx_add_component(
                slide_index=0,
                component="stack",
                target_placeholder=1,
                component_id="content_stack",
                params={"direction": "vertical", "gap": "md"}
            )
            # Then add multiple components into the stack
            await pptx_add_component(
                slide_index=0,
                component="text",
                target_component="content_stack",
                params={"text": "Title", "font_size": 24, "bold": True}
            )
            await pptx_add_component(
                slide_index=0,
                component="image",
                target_component="content_stack",
                params={"image_source": "chart.png"}
            )
        """
        try:
            import json
            from ...models import ErrorResponse, ComponentResponse
            from ...themes.design_system import resolve_design_system
            from ...components.registry import get_component_class
            from ...constants import ErrorMessages

            # Handle params - could be dict, JSON string, or None
            if params is None:
                params = {}
            elif isinstance(params, str):
                try:
                    params = json.loads(params)
                except json.JSONDecodeError as e:
                    return ErrorResponse(
                        error=f"Invalid JSON in params: {str(e)}"
                    ).model_dump_json()
            elif not isinstance(params, dict):
                return ErrorResponse(
                    error=f"params must be a dict or JSON string, got {type(params).__name__}"
                ).model_dump_json()

            # Get presentation
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            # Validate slide index
            if slide_index < 0 or slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found. Presentation has {len(prs.slides)} slides."
                ).model_dump_json()

            slide = prs.slides[slide_index]

            # Determine target and positioning
            target_placeholder_obj = None
            final_left, final_top, final_width, final_height = left, top, width, height

            # VALIDATION: Check if using free-form positioning for content that should use placeholders
            if target_placeholder is None and target_component is None and left is not None:
                # This is free-form positioning - check if slide has appropriate placeholders
                from pptx.enum.shapes import PP_PLACEHOLDER

                # Map component types to placeholder types they should use
                component_to_placeholder_map = {
                    "Table": [PP_PLACEHOLDER.TABLE, PP_PLACEHOLDER.OBJECT],
                    "ColumnChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "BarChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "LineChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "AreaChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "PieChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "DoughnutChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "ScatterChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "BubbleChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "WaterfallChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "SparklineChart": [PP_PLACEHOLDER.CHART, PP_PLACEHOLDER.OBJECT],
                    "Image": [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.OBJECT],
                }

                if component in component_to_placeholder_map:
                    # Check if slide has matching placeholders
                    available_placeholders = []
                    placeholder_types = component_to_placeholder_map[component]

                    for shape in slide.placeholders:
                        try:
                            ph_type = shape.placeholder_format.type
                            if ph_type in placeholder_types:
                                available_placeholders.append(
                                    {
                                        "idx": shape.placeholder_format.idx,
                                        "type": ph_type.name
                                        if hasattr(ph_type, "name")
                                        else str(ph_type),
                                        "name": shape.name,
                                    }
                                )
                        except Exception:
                            continue

                    if available_placeholders:
                        # ERROR: They're using free-form positioning when placeholders exist!
                        ph_list = ", ".join(
                            [f"idx={ph['idx']} ({ph['type']})" for ph in available_placeholders]
                        )
                        return ErrorResponse(
                            error=f"❌ PLACEHOLDER REQUIRED: Cannot add {component} with free-form positioning (left/top/width/height) "
                            f"because this slide has {len(available_placeholders)} suitable placeholder(s): {ph_list}. "
                            f"\n\n✅ USE THIS INSTEAD: pptx_add_component(slide_index={slide_index}, component='{component}', "
                            f"target_placeholder={available_placeholders[0]['idx']}, params={{...}})"
                            f"\n\n⚠️ Using free-form positioning creates overlays where the placeholder 'Click to add' text "
                            f"remains visible underneath your content, breaking the template design."
                        ).model_dump_json()

            # MODE 1: Target placeholder
            if target_placeholder is not None:
                # Find placeholder
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == target_placeholder:
                        target_placeholder_obj = shape
                        break

                if not target_placeholder_obj:
                    return ErrorResponse(
                        error=f"Placeholder {target_placeholder} not found on slide {slide_index}"
                    ).model_dump_json()

                # Use placeholder bounds
                final_left = (
                    target_placeholder_obj.left.inches
                    if hasattr(target_placeholder_obj.left, "inches")
                    else target_placeholder_obj.left / 914400
                )
                final_top = (
                    target_placeholder_obj.top.inches
                    if hasattr(target_placeholder_obj.top, "inches")
                    else target_placeholder_obj.top / 914400
                )
                final_width = (
                    target_placeholder_obj.width.inches
                    if hasattr(target_placeholder_obj.width, "inches")
                    else target_placeholder_obj.width / 914400
                )
                final_height = (
                    target_placeholder_obj.height.inches
                    if hasattr(target_placeholder_obj.height, "inches")
                    else target_placeholder_obj.height / 914400
                )

                logger.info(
                    f"Targeting placeholder {target_placeholder}: bounds=({final_left}, {final_top}, {final_width}, {final_height})"
                )

            # MODE 2: Target component (composition)
            elif target_component is not None:
                from ...components.tracking import component_tracker

                # Get parent component tracking info
                parent_tracked = component_tracker.get(
                    presentation=metadata.name,
                    slide_index=slide_index,
                    component_id=target_component,
                )

                if not parent_tracked:
                    return ErrorResponse(
                        error=f"Component '{target_component}' not found on slide {slide_index}. "
                        f"Use component_id parameter when adding components to reference them later."
                    ).model_dump_json()

                # Get the actual component instance
                parent_instance = parent_tracked.instance

                # Calculate position relative to parent
                # For composition, child is positioned within parent's bounds
                parent_bounds = component_tracker.get_bounds(
                    metadata.name, slide_index, target_component
                )
                if not parent_bounds:
                    return ErrorResponse(
                        error=f"Could not get bounds for component '{target_component}'"
                    ).model_dump_json()

                parent_left, parent_top, parent_width, parent_height = parent_bounds

                # Check if parent is a Stack - if so, use its distribute logic
                from ...components.core.stack import Stack

                if parent_instance and isinstance(parent_instance, Stack):
                    # Get current child count for this stack
                    children = component_tracker.get_children(
                        metadata.name, slide_index, target_component
                    )
                    child_count = len(children)

                    # Get next position from stack's distribute
                    positions = parent_instance.distribute(
                        num_items=child_count + 1,  # Include the new child
                        left=parent_left,
                        top=parent_top,
                        container_width=parent_width,
                        container_height=parent_height,
                    )

                    # Use the last position (for the new child)
                    pos = positions[-1]
                    final_left = pos["left"]
                    final_top = pos["top"]
                    final_width = width if width is not None else pos.get("width", parent_width)
                    final_height = height if height is not None else pos.get("height", 1.0)

                # Otherwise, use default composition (center or offset)
                elif left is not None and top is not None:
                    # Treat as offsets from parent's top-left
                    final_left = parent_left + left
                    final_top = parent_top + top
                    final_width = width if width is not None else parent_width * 0.5
                    final_height = height if height is not None else parent_height * 0.3
                else:
                    # Default: center within parent
                    child_width = width if width is not None else parent_width * 0.5
                    child_height = height if height is not None else parent_height * 0.3
                    final_left = parent_left + (parent_width - child_width) / 2
                    final_top = parent_top + (parent_height - child_height) / 2
                    final_width = child_width
                    final_height = child_height

                logger.info(
                    f"Targeting component '{target_component}': "
                    f"parent bounds=({parent_left}, {parent_top}, {parent_width}, {parent_height}), "
                    f"child position=({final_left}, {final_top}, {final_width}, {final_height})"
                )

            # MODE 3: Target layout (grid/flex positioning)
            elif target_layout is not None:
                from ...components.tracking import component_tracker
                from ...models import LayoutType, TargetType

                # Get existing components in this layout on this slide
                all_components = component_tracker.list_on_slide(metadata.name, slide_index)
                layout_components = [
                    c
                    for c in all_components
                    if c.target_type == TargetType.LAYOUT.value and c.target_id == target_layout
                ]

                # Validate layout type
                try:
                    layout_enum = LayoutType(target_layout)
                except ValueError:
                    valid_layouts = [lt.value for lt in LayoutType]
                    return ErrorResponse(
                        error=f"Unknown layout type: '{target_layout}'. "
                        f"Supported: {', '.join(valid_layouts)}"
                    ).model_dump_json()

                # Simple grid layout logic
                if layout_enum == LayoutType.GRID:
                    # Default 2-column grid with gap
                    cols = 2
                    gap = 0.2  # inches
                    margin_left = 1.0
                    margin_top = 2.0
                    col_width = (10.0 - margin_left * 2 - gap * (cols - 1)) / cols

                    # Calculate position based on number of existing items
                    item_index = len(layout_components)
                    row = item_index // cols
                    col = item_index % cols

                    final_left = margin_left + col * (col_width + gap)
                    final_top = margin_top + row * (2.0 + gap)  # 2.0 = row height
                    final_width = width if width is not None else col_width
                    final_height = height if height is not None else 1.5

                elif layout_enum == LayoutType.FLEX_ROW:
                    # Horizontal flex layout
                    gap = 0.2
                    margin_left = 1.0
                    margin_top = 2.0

                    # Calculate X position based on existing components
                    offset_x = sum((c.width or 0) + gap for c in layout_components)

                    final_left = margin_left + offset_x
                    final_top = top if top is not None else margin_top
                    final_width = width if width is not None else 2.0
                    final_height = height if height is not None else 1.5

                elif layout_enum == LayoutType.FLEX_COLUMN:
                    # Vertical flex layout
                    gap = 0.2
                    margin_left = 1.0
                    margin_top = 2.0

                    # Calculate Y position based on existing components
                    offset_y = sum((c.height or 0) + gap for c in layout_components)

                    final_left = left if left is not None else margin_left
                    final_top = margin_top + offset_y
                    final_width = width if width is not None else 3.0
                    final_height = height if height is not None else 1.0

                logger.info(
                    f"Targeting layout '{target_layout}': "
                    f"position=({final_left}, {final_top}, {final_width}, {final_height}), "
                    f"existing items={len(layout_components)}"
                )

            # MODE 4: Free-form (requires explicit positioning)
            else:
                if left is None or top is None:
                    return ErrorResponse(
                        error="Free-form positioning requires 'left' and 'top' parameters"
                    ).model_dump_json()

                # Use defaults for width/height if not specified
                final_width = width if width is not None else 3.0
                final_height = height if height is not None else 2.0

            # Resolve design system with priority hierarchy
            design_system = resolve_design_system(
                slide=slide, placeholder=target_placeholder_obj, theme=theme, params=params
            )

            # Get component class from registry
            component_class = get_component_class(component)
            if not component_class:
                return ErrorResponse(
                    error=f"Unknown component type: '{component}'. Use pptx_list_components to see available components."
                ).model_dump_json()

            # Build comprehensive theme object with all design tokens
            theme_obj = {
                "colors": {
                    "primary": design_system.primary_color,
                    "secondary": design_system.secondary_color,
                    "background": design_system.background_color,
                    "text": design_system.text_color,
                    "border": design_system.border_color,
                },
                "typography": {
                    "font_family": design_system.font_family,
                    "font_size": design_system.font_size,
                    "font_bold": design_system.font_bold,
                    "font_italic": design_system.font_italic,
                },
                "spacing": {
                    "padding": design_system.padding,
                    "margin": design_system.margin,
                    "gap": design_system.gap,
                },
                "borders": {
                    "radius": design_system.border_radius,
                    "width": design_system.border_width,
                },
                # Also expose as flat keys for easy access via get_theme_attr()
                "font_family": design_system.font_family,
                "font_size": design_system.font_size,
                "padding": design_system.padding,
                "margin": design_system.margin,
                "gap": design_system.gap,
                "border_radius": design_system.border_radius,
                "border_width": design_system.border_width,
            }

            # Merge user params with theme object
            component_params = {
                **params,  # User params take priority
                "theme": theme_obj,  # Pass theme object for components that support it
            }

            # Get component's __init__ signature to filter params
            import inspect

            sig = inspect.signature(component_class.__init__)
            accepted_params = set(sig.parameters.keys()) - {"self"}

            # Filter component_params to only include accepted parameters
            filtered_params = {k: v for k, v in component_params.items() if k in accepted_params}

            # Create and render component
            component_instance = component_class(**filtered_params)

            # Check if render is async
            import asyncio

            # Validate and adjust position to fit within slide bounds
            # This prevents overlapping with title areas and slide boundaries
            if (
                final_left is not None
                and final_top is not None
                and final_width is not None
                and final_height is not None
            ):
                from ...layout.helpers import validate_position

                final_left, final_top, final_width, final_height = validate_position(
                    final_left, final_top, final_width, final_height
                )

            # Prepare render kwargs
            render_kwargs = {
                "left": final_left,  # Pass as float - components handle Inches() conversion
                "top": final_top,
                "width": final_width,
                "height": final_height,
            }

            # If targeting a placeholder, pass it to the render method (if component supports it)
            if target_placeholder_obj is not None:
                # Check if render method accepts 'placeholder' parameter
                sig = inspect.signature(component_instance.render)
                if "placeholder" in sig.parameters:
                    render_kwargs["placeholder"] = target_placeholder_obj

            render_result = component_instance.render(slide, **render_kwargs)

            # Await if async
            if asyncio.iscoroutine(render_result):
                await render_result

            # Track component in registry if component_id provided
            if component_id:
                from ...components.tracking import component_tracker

                # Get shape index (last shape added)
                shape_index = len(slide.shapes) - 1 if slide.shapes else None

                # Determine parent_id for composition
                parent_id = None
                if target_component:
                    parent_id = target_component

                # Determine target type using enum
                from ...models import TargetType

                target_id_value: str | int | None = None
                if target_placeholder is not None:
                    target_type_value = TargetType.PLACEHOLDER.value
                    target_id_value = target_placeholder
                elif target_component:
                    target_type_value = TargetType.COMPONENT.value
                    target_id_value = target_component
                elif target_layout:
                    target_type_value = TargetType.LAYOUT.value
                    target_id_value = target_layout
                else:
                    target_type_value = TargetType.FREE_FORM.value

                # Register component instance
                component_tracker.register(
                    presentation=metadata.name,
                    slide_index=slide_index,
                    component_id=component_id,
                    component_type=component,
                    left=final_left,
                    top=final_top,
                    width=final_width,
                    height=final_height,
                    target_type=target_type_value,
                    target_id=target_id_value,
                    parent_id=parent_id,
                    params=params,
                    theme=theme,
                    shape_index=shape_index,
                    instance=component_instance,
                )

                logger.info(f"Component registered with ID: {component_id}")

            # Update metadata and save
            await manager.update_slide_metadata(slide_index)
            await manager.update(presentation)

            pres_name = presentation or manager.get_current_name() or "presentation"

            # Build response message
            target_info = ""
            if target_placeholder is not None:
                target_info = f" in placeholder {target_placeholder}"
            elif target_component:
                target_info = f" in component '{target_component}'"
            elif target_layout:
                target_info = f" in layout '{target_layout}'"

            message = f"Added {component}{target_info} to slide {slide_index}"
            if design_system.overrides:
                message += f" (overrides: {', '.join(design_system.overrides.keys())})"

            return ComponentResponse(
                presentation=pres_name,
                slide_index=slide_index,
                component=component,
                message=message,
                variant=params.get("variant"),
            ).model_dump_json()

        except Exception as e:
            logger.error(f"Failed to add component: {e}", exc_info=True)
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    @mcp.tool
    async def pptx_update_component(
        slide_index: int,
        component_id: str,
        params: dict[str, Any] | None = None,
        left: float | None = None,
        top: float | None = None,
        width: float | None = None,
        height: float | None = None,
        presentation: str | None = None,
    ) -> str:
        """
        Update an existing component's properties or position.

        Allows updating component parameters (like text, colors, data) or repositioning
        without recreating the component.

        Args:
            slide_index: Index of the slide containing the component (0-based)
            component_id: ID of the component to update
            params: New component parameters to update (merged with existing)
            left: New left position in inches (optional)
            top: New top position in inches (optional)
            width: New width in inches (optional)
            height: New height in inches (optional)
            presentation: Presentation name (uses current if not specified)

        Returns:
            JSON string with update confirmation

        Examples:
            # Update component data
            await pptx_update_component(
                slide_index=0,
                component_id="chart1",
                params={
                    "data": {
                        "categories": ["Q1", "Q2", "Q3"],
                        "series": {"Revenue": [100, 120, 140]}
                    }
                }
            )

            # Reposition component
            await pptx_update_component(
                slide_index=0,
                component_id="badge1",
                left=4.0,
                top=3.0
            )

            # Update component params and position
            await pptx_update_component(
                slide_index=0,
                component_id="metric_card",
                params={"value": "$200K", "label": "New Revenue"},
                left=1.5,
                top=2.5
            )
        """
        try:
            import json
            from ...models import ErrorResponse, ComponentResponse
            from ...components.tracking import component_tracker
            from ...components.registry import get_component_class
            from ...constants import ErrorMessages

            # Handle params - could be dict, JSON string, or None
            if params is None:
                params = {}
            elif isinstance(params, str):
                try:
                    params = json.loads(params)
                except json.JSONDecodeError as e:
                    return ErrorResponse(
                        error=f"Invalid JSON in params: {str(e)}"
                    ).model_dump_json()
            elif not isinstance(params, dict):
                return ErrorResponse(
                    error=f"params must be a dict or JSON string, got {type(params).__name__}"
                ).model_dump_json()

            # Get presentation
            result = await manager.get(presentation)
            if not result:
                return ErrorResponse(error=ErrorMessages.NO_PRESENTATION).model_dump_json()

            prs, metadata = result

            # Validate slide index
            if slide_index < 0 or slide_index >= len(prs.slides):
                return ErrorResponse(
                    error=f"Slide index {slide_index} not found. Presentation has {len(prs.slides)} slides."
                ).model_dump_json()

            slide = prs.slides[slide_index]

            # Get existing component
            component_instance = component_tracker.get(
                presentation=metadata.name, slide_index=slide_index, component_id=component_id
            )

            if not component_instance:
                return ErrorResponse(
                    error=f"Component '{component_id}' not found on slide {slide_index}"
                ).model_dump_json()

            # Get component bounds
            bounds = component_tracker.get_bounds(metadata.name, slide_index, component_id)
            if not bounds:
                return ErrorResponse(
                    error=f"Could not get bounds for component '{component_id}'"
                ).model_dump_json()

            old_left, old_top, old_width, old_height = bounds

            # Determine new position/size
            new_left = left if left is not None else old_left
            new_top = top if top is not None else old_top
            new_width = width if width is not None else old_width
            new_height = height if height is not None else old_height

            # Merge params
            merged_params = {**component_instance.params, **(params or {})}

            # Get component class
            component_class = get_component_class(component_instance.component_type)
            if not component_class:
                return ErrorResponse(
                    error=f"Unknown component type: '{component_instance.component_type}'"
                ).model_dump_json()

            # Remove old shape from slide
            if component_instance.shape_index is not None:
                try:
                    shape_to_remove = slide.shapes[component_instance.shape_index]
                    slide.shapes._spTree.remove(shape_to_remove.element)
                except (IndexError, AttributeError):
                    pass  # Shape may have been moved/deleted

            # Recreate component with new params/position
            from ...themes.design_system import resolve_design_system

            design_system = resolve_design_system(
                slide=slide, placeholder=None, theme=component_instance.theme, params=merged_params
            )

            component_params = {
                "bg_color": design_system.background_color,
                "text_color": design_system.text_color,
                "border_color": design_system.border_color,
                "font_family": design_system.font_family,
                "font_size": design_system.font_size,
                **merged_params,
                "theme": {
                    "colors": {
                        "primary": design_system.primary_color,
                        "secondary": design_system.secondary_color,
                        "background": design_system.background_color,
                        "text": design_system.text_color,
                    },
                    "typography": {
                        "font_family": design_system.font_family,
                        "font_size": design_system.font_size,
                    },
                    "spacing": {
                        "padding": design_system.padding,
                        "margin": design_system.margin,
                    },
                },
            }

            new_component = component_class(**component_params)
            new_component.render(
                slide,
                left=Inches(new_left),
                top=Inches(new_top),
                width=Inches(new_width),
                height=Inches(new_height),
            )

            # Update tracker with new values
            new_shape_index = len(slide.shapes) - 1 if slide.shapes else None

            component_tracker.update(
                presentation=metadata.name,
                slide_index=slide_index,
                component_id=component_id,
                left=new_left,
                top=new_top,
                width=new_width,
                height=new_height,
                params=merged_params,
                shape_index=new_shape_index,
            )

            # Update metadata and save
            await manager.update_slide_metadata(slide_index)
            await manager.update(presentation)

            pres_name = presentation or manager.get_current_name() or "presentation"

            # Build response message
            changes = []
            if params:
                changes.append("params")
            if left is not None or top is not None:
                changes.append("position")
            if width is not None or height is not None:
                changes.append("size")

            message = f"Updated {component_id} on slide {slide_index} ({', '.join(changes)})"

            return ComponentResponse(
                presentation=pres_name,
                slide_index=slide_index,
                component=component_instance.component_type,
                message=message,
                variant=merged_params.get("variant"),
            ).model_dump_json()

        except Exception as e:
            logger.error(f"Failed to update component: {e}", exc_info=True)
            from ...models import ErrorResponse

            return ErrorResponse(error=str(e)).model_dump_json()

    return {
        "pptx_add_component": pptx_add_component,
        "pptx_update_component": pptx_update_component,
        "pptx_list_slide_components": pptx_list_slide_components,
    }
