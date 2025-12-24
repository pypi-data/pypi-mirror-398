"""
Inspection and Layout Tools for PowerPoint MCP Server

Provides tools for inspecting slide contents, analyzing layouts,
and making intelligent adjustments to fix positioning issues.
"""

from __future__ import annotations

from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def register_inspection_tools(mcp, manager):
    """Register slide inspection and layout adjustment tools."""

    from ...layout.helpers import (
        validate_position,
        get_safe_content_area,
        SLIDE_WIDTH,
        SLIDE_HEIGHT,
        MARGIN_TOP,
    )

    @mcp.tool
    async def pptx_inspect_slide(
        slide_index: int,
        include_measurements: bool = True,
        check_overlaps: bool = True,
        presentation: str | None = None,
    ) -> str:
        """
        Inspect and describe the contents of a slide.

        Provides detailed information about all elements on a slide including:
        - Shape types and positions
        - Text content
        - Images and their locations
        - Charts and their data
        - Overlapping elements
        - Layout issues

        Args:
            slide_index: Index of the slide to inspect (0-based)
            include_measurements: Include detailed position/size measurements
            check_overlaps: Check for overlapping elements
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Detailed description of slide contents and any layout issues

        Example:
            result = await pptx_inspect_slide(
                slide_index=1,
                include_measurements=True,
                check_overlaps=True
            )
            # Returns detailed analysis like:
            # "Slide 1: 'Quarterly Results'
            #  - Title placeholder at (0.5, 0.5, 9.0, 1.0)
            #  - Chart 'Revenue Chart' at (1.0, 2.0, 8.0, 4.0)
            #  - Image at (7.0, 1.5, 2.0, 1.5)
            #  WARNING: Image overlaps with chart
            #  WARNING: Chart extends beyond safe content area"
        """

        async def _inspect_slide():
            prs = await manager.get_presentation(presentation)
            if not prs:
                return "Error: No presentation found"

            # Ensure slide_index is an integer
            idx = int(slide_index) if isinstance(slide_index, str) else slide_index

            if idx >= len(prs.slides):
                return f"Error: Slide index {idx} out of range"

            slide = prs.slides[idx]

            # Build description
            description = []
            description.append(f"=== SLIDE {idx} INSPECTION ===\n")

            # Get slide title if exists
            if slide.shapes.title:
                description.append(f"Title: '{slide.shapes.title.text}'")
            else:
                description.append("Title: (No title)")

            # Get slide layout name
            layout_name = slide.slide_layout.name
            description.append(f"Layout: {layout_name}\n")

            # Categorize shapes
            placeholders = []
            text_boxes = []
            images = []
            charts = []
            tables = []
            other_shapes = []

            for shape in slide.shapes:
                shape_info = _analyze_shape(shape, include_measurements)

                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    placeholders.append(shape_info)
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    text_boxes.append(shape_info)
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(shape_info)
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    charts.append(shape_info)
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    tables.append(shape_info)
                else:
                    other_shapes.append(shape_info)

            # Report each category
            if placeholders:
                description.append("PLACEHOLDERS:")
                for p in placeholders:
                    description.append(f"  • {p}")

            if text_boxes:
                description.append("\nTEXT BOXES:")
                for t in text_boxes:
                    description.append(f"  • {t}")

            if images:
                description.append("\nIMAGES:")
                for i in images:
                    description.append(f"  • {i}")

            if charts:
                description.append("\nCHARTS:")
                for c in charts:
                    description.append(f"  • {c}")

            if tables:
                description.append("\nTABLES:")
                for t in tables:
                    description.append(f"  • {t}")

            if other_shapes:
                description.append("\nOTHER SHAPES:")
                for s in other_shapes:
                    description.append(f"  • {s}")

            # Check for layout issues
            issues = []

            # Check overlaps if requested
            if check_overlaps:
                overlaps = _check_overlaps(slide.shapes)
                if overlaps:
                    issues.append("\nOVERLAPPING ELEMENTS:")
                    for overlap in overlaps:
                        issues.append(f"  ⚠️ {overlap}")

            # Check bounds
            out_of_bounds = _check_bounds(slide.shapes)
            if out_of_bounds:
                issues.append("\nOUT OF BOUNDS:")
                for oob in out_of_bounds:
                    issues.append(f"  ⚠️ {oob}")

            # Check spacing issues
            spacing_issues = _check_spacing(slide.shapes)
            if spacing_issues:
                issues.append("\nSPACING ISSUES:")
                for issue in spacing_issues:
                    issues.append(f"  ⚠️ {issue}")

            if issues:
                description.append("\n=== LAYOUT ISSUES DETECTED ===")
                description.extend(issues)
                description.append(
                    "\nUse pptx_fix_slide_layout() to automatically fix these issues"
                )
            else:
                description.append("\n✅ No layout issues detected")

            # Add summary
            description.append("\n=== SUMMARY ===")
            description.append(f"Total elements: {len(slide.shapes)}")
            description.append(f"Layout issues: {len(issues) if issues else 0}")

            return "\n".join(description)

        def _analyze_shape(shape, include_measurements):
            """Analyze a single shape and return description."""
            info = []

            # Basic type
            shape_type_name = _get_shape_type_name(shape.shape_type)
            info.append(shape_type_name)

            # Name if available
            if hasattr(shape, "name") and shape.name:
                info.append(f"'{shape.name}'")

            # Position and size
            if include_measurements and hasattr(shape, "left"):
                left = shape.left.inches if hasattr(shape.left, "inches") else 0
                top = shape.top.inches if hasattr(shape.top, "inches") else 0
                width = shape.width.inches if hasattr(shape.width, "inches") else 0
                height = shape.height.inches if hasattr(shape.height, "inches") else 0
                info.append(f"at ({left:.1f}, {top:.1f}) size ({width:.1f} x {height:.1f})")

            # Text content
            if hasattr(shape, "text_frame") and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    preview = text[:50] + "..." if len(text) > 50 else text
                    info.append(f'text: "{preview}"')

            # Chart details
            if shape.shape_type == MSO_SHAPE_TYPE.CHART:
                if hasattr(shape, "chart"):
                    chart_type = shape.chart.chart_type
                    info.append(f"type: {chart_type}")
                    if shape.chart.has_title and shape.chart.chart_title:
                        info.append(f'title: "{shape.chart.chart_title.text_frame.text}"')

            return " ".join(info)

        def _get_shape_type_name(shape_type):
            """Get readable name for shape type."""
            type_names = {
                MSO_SHAPE_TYPE.PLACEHOLDER: "Placeholder",
                MSO_SHAPE_TYPE.TEXT_BOX: "TextBox",
                MSO_SHAPE_TYPE.PICTURE: "Image",
                MSO_SHAPE_TYPE.CHART: "Chart",
                MSO_SHAPE_TYPE.TABLE: "Table",
                MSO_SHAPE_TYPE.GROUP: "Group",
                MSO_SHAPE_TYPE.AUTO_SHAPE: "Shape",
            }
            return type_names.get(shape_type, f"Shape({shape_type})")

        def _check_overlaps(shapes):
            """Check for overlapping shapes."""
            overlaps = []
            shape_list = list(shapes)

            for i, shape1 in enumerate(shape_list):
                # Skip title placeholders
                if _is_title_placeholder(shape1):
                    continue

                for j, shape2 in enumerate(shape_list[i + 1 :], i + 1):
                    # Skip title placeholders
                    if _is_title_placeholder(shape2):
                        continue

                    if _shapes_overlap(shape1, shape2):
                        type1 = _get_shape_type_name(shape1.shape_type)
                        type2 = _get_shape_type_name(shape2.shape_type)
                        overlaps.append(f"{type1} overlaps with {type2}")

            return overlaps

        def _is_title_placeholder(shape):
            """Check if shape is a title placeholder."""
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "placeholder_format"):
                    return shape.placeholder_format.type in [
                        PP_PLACEHOLDER.TITLE,
                        PP_PLACEHOLDER.CENTER_TITLE,
                    ]
            return False

        def _shapes_overlap(shape1, shape2):
            """Check if two shapes overlap."""
            if not (hasattr(shape1, "left") and hasattr(shape2, "left")):
                return False

            # Get bounds for shape1
            left1 = shape1.left.inches if hasattr(shape1.left, "inches") else 0
            top1 = shape1.top.inches if hasattr(shape1.top, "inches") else 0
            right1 = left1 + (shape1.width.inches if hasattr(shape1.width, "inches") else 0)
            bottom1 = top1 + (shape1.height.inches if hasattr(shape1.height, "inches") else 0)

            # Get bounds for shape2
            left2 = shape2.left.inches if hasattr(shape2.left, "inches") else 0
            top2 = shape2.top.inches if hasattr(shape2.top, "inches") else 0
            right2 = left2 + (shape2.width.inches if hasattr(shape2.width, "inches") else 0)
            bottom2 = top2 + (shape2.height.inches if hasattr(shape2.height, "inches") else 0)

            # Check if they overlap
            return not (right1 <= left2 or right2 <= left1 or bottom1 <= top2 or bottom2 <= top1)

        def _check_bounds(shapes):
            """Check for shapes outside slide bounds."""
            issues = []

            for shape in shapes:
                if not hasattr(shape, "left"):
                    continue

                left = shape.left.inches if hasattr(shape.left, "inches") else 0
                top = shape.top.inches if hasattr(shape.top, "inches") else 0
                right = left + (shape.width.inches if hasattr(shape.width, "inches") else 0)
                bottom = top + (shape.height.inches if hasattr(shape.height, "inches") else 0)

                shape_type = _get_shape_type_name(shape.shape_type)

                if left < 0:
                    issues.append(f"{shape_type} extends beyond left edge")
                if top < 0:
                    issues.append(f"{shape_type} extends beyond top edge")
                if right > SLIDE_WIDTH:
                    issues.append(
                        f"{shape_type} extends beyond right edge ({right:.1f} > {SLIDE_WIDTH})"
                    )
                if bottom > SLIDE_HEIGHT:
                    issues.append(
                        f"{shape_type} extends beyond bottom edge ({bottom:.1f} > {SLIDE_HEIGHT})"
                    )

            return issues

        def _check_spacing(shapes):
            """Check for spacing issues."""
            issues = []

            # Check for elements too close to edges (except backgrounds)
            MIN_EDGE_MARGIN = 0.2  # inches

            for shape in shapes:
                # Skip backgrounds and title placeholders
                if _is_title_placeholder(shape):
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Check if it's a background (covers most of slide)
                    if hasattr(shape, "width") and hasattr(shape, "height"):
                        w = shape.width.inches if hasattr(shape.width, "inches") else 0
                        h = shape.height.inches if hasattr(shape.height, "inches") else 0
                        if w > SLIDE_WIDTH * 0.9 and h > SLIDE_HEIGHT * 0.9:
                            continue

                if not hasattr(shape, "left"):
                    continue

                left = shape.left.inches if hasattr(shape.left, "inches") else 0
                top = shape.top.inches if hasattr(shape.top, "inches") else 0

                shape_type = _get_shape_type_name(shape.shape_type)

                if 0 < left < MIN_EDGE_MARGIN:
                    issues.append(f'{shape_type} too close to left edge ({left:.2f}")')
                if 0 < top < MIN_EDGE_MARGIN and top < MARGIN_TOP * 0.5:
                    issues.append(f'{shape_type} too close to top edge ({top:.2f}")')

            return issues

        return await _inspect_slide()

    @mcp.tool
    async def pptx_fix_slide_layout(
        slide_index: int,
        fix_overlaps: bool = True,
        fix_bounds: bool = True,
        fix_spacing: bool = True,
        maintain_relative_positions: bool = True,
        presentation: str | None = None,
    ) -> str:
        """
        Automatically fix layout issues on a slide.

        Analyzes the slide and makes intelligent adjustments to:
        - Resolve overlapping elements
        - Bring out-of-bounds elements back into view
        - Improve spacing and alignment
        - Maintain visual hierarchy

        Args:
            slide_index: Index of the slide to fix (0-based)
            fix_overlaps: Resolve overlapping elements
            fix_bounds: Bring out-of-bounds elements into slide
            fix_spacing: Improve element spacing
            maintain_relative_positions: Try to keep relative positioning
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Summary of fixes applied

        Example:
            result = await pptx_fix_slide_layout(
                slide_index=1,
                fix_overlaps=True,
                fix_bounds=True
            )
            # Returns: "Fixed 3 overlapping elements, adjusted 2 out-of-bounds items"
        """

        async def _fix_layout():
            prs = await manager.get_presentation(presentation)
            if not prs:
                return "Error: No presentation found"

            # Ensure slide_index is an integer
            idx = int(slide_index) if isinstance(slide_index, str) else slide_index

            if idx >= len(prs.slides):
                return f"Error: Slide index {idx} out of range"

            slide = prs.slides[idx]
            fixes_applied = []

            # Get safe content area
            safe_area = get_safe_content_area(has_title=bool(slide.shapes.title))

            # First pass: Fix out of bounds
            if fix_bounds:
                bounds_fixed = _fix_out_of_bounds(slide.shapes, safe_area)
                if bounds_fixed:
                    fixes_applied.append(f"Fixed {bounds_fixed} out-of-bounds elements")

            # Second pass: Fix overlaps
            if fix_overlaps:
                overlaps_fixed = _fix_overlapping_elements(slide.shapes, safe_area)
                if overlaps_fixed:
                    fixes_applied.append(f"Resolved {overlaps_fixed} overlapping elements")

            # Third pass: Improve spacing
            if fix_spacing:
                spacing_improved = _improve_spacing(slide.shapes, safe_area)
                if spacing_improved:
                    fixes_applied.append(f"Improved spacing for {spacing_improved} elements")

            # Note: Changes are persisted in memory; use pptx_save to persist to file

            if fixes_applied:
                return "Layout fixes applied:\n" + "\n".join(f"  • {fix}" for fix in fixes_applied)
            else:
                return "No layout issues found - slide layout is already optimal"

        def _fix_out_of_bounds(shapes, safe_area):
            """Fix shapes that extend beyond slide bounds."""
            fixed_count = 0

            for shape in shapes:
                if not hasattr(shape, "left"):
                    continue

                # Skip title placeholders
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type in [
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        ]:
                            continue

                left = shape.left.inches
                top = shape.top.inches
                width = shape.width.inches
                height = shape.height.inches

                new_left, new_top, new_width, new_height = validate_position(
                    left, top, width, height
                )

                # Apply fixes if needed
                if new_left != left or new_top != top or new_width != width or new_height != height:
                    shape.left = Inches(new_left)
                    shape.top = Inches(new_top)
                    shape.width = Inches(new_width)
                    shape.height = Inches(new_height)
                    fixed_count += 1

            return fixed_count

        def _fix_overlapping_elements(shapes, safe_area):
            """Resolve overlapping shapes by smart repositioning."""
            fixed_count = 0
            shape_list = list(shapes)

            # Sort shapes by area (larger first) to prioritize important elements
            sortable_shapes = []
            for shape in shape_list:
                if hasattr(shape, "width") and hasattr(shape, "height"):
                    area = shape.width.inches * shape.height.inches
                    sortable_shapes.append((area, shape))

            sortable_shapes.sort(reverse=True, key=lambda x: x[0])

            # Track occupied regions
            occupied_regions = []

            for _, shape in sortable_shapes:
                # Skip title placeholders
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type in [
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        ]:
                            continue

                if not hasattr(shape, "left"):
                    continue

                left = shape.left.inches
                top = shape.top.inches
                width = shape.width.inches
                height = shape.height.inches

                # Check if current position overlaps with occupied regions
                needs_move = False
                for region in occupied_regions:
                    if _regions_overlap((left, top, width, height), region):
                        needs_move = True
                        break

                if needs_move:
                    # Find new position
                    new_left, new_top = _find_free_position(
                        width, height, occupied_regions, safe_area
                    )

                    if new_left is not None and new_top is not None:
                        shape.left = Inches(new_left)
                        shape.top = Inches(new_top)
                        fixed_count += 1
                        occupied_regions.append((new_left, new_top, width, height))
                else:
                    occupied_regions.append((left, top, width, height))

            return fixed_count

        def _regions_overlap(region1, region2):
            """Check if two regions overlap."""
            l1, t1, w1, h1 = region1
            l2, t2, w2, h2 = region2

            return not (l1 + w1 <= l2 or l2 + w2 <= l1 or t1 + h1 <= t2 or t2 + h2 <= t1)

        def _find_free_position(width, height, occupied_regions, safe_area):
            """Find a free position for an element."""
            # Try positions in a grid pattern
            x_positions = [
                safe_area["left"] + i * 0.5 for i in range(int(safe_area["width"] / 0.5))
            ]
            y_positions = [
                safe_area["top"] + i * 0.5 for i in range(int(safe_area["height"] / 0.5))
            ]

            for y in y_positions:
                for x in x_positions:
                    # Check if this position works
                    if x + width > safe_area["left"] + safe_area["width"]:
                        continue
                    if y + height > safe_area["top"] + safe_area["height"]:
                        continue

                    # Check overlaps
                    overlaps = False
                    for region in occupied_regions:
                        if _regions_overlap((x, y, width, height), region):
                            overlaps = True
                            break

                    if not overlaps:
                        return x, y

            return None, None

        def _improve_spacing(shapes, safe_area):
            """Improve spacing between elements."""
            improved_count = 0

            # Group shapes by type for better organization
            images = []
            charts = []
            text_boxes = []

            for shape in shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    images.append(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                    charts.append(shape)
                elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    text_boxes.append(shape)

            # Align similar elements if multiple exist
            if len(images) > 1:
                improved_count += _align_elements(images, "horizontal")

            if len(text_boxes) > 1:
                improved_count += _align_elements(text_boxes, "vertical")

            return improved_count

        def _align_elements(shapes, direction):
            """Align multiple similar elements."""
            if not shapes:
                return 0

            aligned = 0

            if direction == "horizontal":
                # Align tops
                tops = [s.top.inches for s in shapes if hasattr(s, "top")]
                if tops:
                    avg_top = sum(tops) / len(tops)
                    for shape in shapes:
                        if hasattr(shape, "top"):
                            if abs(shape.top.inches - avg_top) > 0.1:
                                shape.top = Inches(avg_top)
                                aligned += 1

            elif direction == "vertical":
                # Align lefts
                lefts = [s.left.inches for s in shapes if hasattr(s, "left")]
                if lefts:
                    avg_left = sum(lefts) / len(lefts)
                    for shape in shapes:
                        if hasattr(shape, "left"):
                            if abs(shape.left.inches - avg_left) > 0.1:
                                shape.left = Inches(avg_left)
                                aligned += 1

            return aligned

        return await _fix_layout()

    @mcp.tool
    async def pptx_analyze_presentation_layout(presentation: str | None = None) -> str:
        """
        Analyze the entire presentation for layout consistency and issues.

        Provides a comprehensive report on:
        - Layout usage patterns
        - Common positioning issues
        - Consistency across slides
        - Recommendations for improvement

        Args:
            presentation: Name of presentation (uses current if not specified)

        Returns:
            Comprehensive layout analysis report
        """

        async def _analyze_presentation():
            prs = await manager.get_presentation(presentation)
            if not prs:
                return "Error: No presentation found"

            report = []
            report.append("=== PRESENTATION LAYOUT ANALYSIS ===\n")
            report.append(f"Total slides: {len(prs.slides)}")

            # Analyze each slide
            issues_by_slide = {}
            layout_usage = {}
            element_stats = {"images": 0, "charts": 0, "tables": 0, "text_boxes": 0}

            for i, slide in enumerate(prs.slides):
                layout_name = slide.slide_layout.name
                layout_usage[layout_name] = layout_usage.get(layout_name, 0) + 1

                slide_issues = []

                # Count elements
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        element_stats["images"] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        element_stats["charts"] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        element_stats["tables"] += 1
                    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                        element_stats["text_boxes"] += 1

                # Check for issues
                overlaps = _count_overlaps(slide.shapes)
                if overlaps > 0:
                    slide_issues.append(f"{overlaps} overlapping elements")

                oob = _count_out_of_bounds(slide.shapes)
                if oob > 0:
                    slide_issues.append(f"{oob} out-of-bounds elements")

                if slide_issues:
                    issues_by_slide[i] = slide_issues

            # Report findings
            report.append("\n=== LAYOUT USAGE ===")
            for layout, count in layout_usage.items():
                report.append(f"  • {layout}: {count} slides")

            report.append("\n=== ELEMENT STATISTICS ===")
            report.append(f"  • Images: {element_stats['images']}")
            report.append(f"  • Charts: {element_stats['charts']}")
            report.append(f"  • Tables: {element_stats['tables']}")
            report.append(f"  • Text boxes: {element_stats['text_boxes']}")

            if issues_by_slide:
                report.append("\n=== SLIDES WITH ISSUES ===")
                for slide_idx, issues in issues_by_slide.items():
                    report.append(f"  Slide {slide_idx}: {', '.join(issues)}")
                report.append(f"\nTotal slides with issues: {len(issues_by_slide)}")
                report.append("Use pptx_fix_slide_layout() on affected slides")
            else:
                report.append("\n✅ No layout issues detected in presentation")

            # Recommendations
            report.append("\n=== RECOMMENDATIONS ===")
            if len(layout_usage) > 3:
                report.append("  • Consider using fewer layout variations for consistency")
            if element_stats["images"] > len(prs.slides) * 3:
                report.append("  • High image density - consider reducing for clarity")
            if any(count == 0 for count in element_stats.values()):
                missing = [k for k, v in element_stats.items() if v == 0]
                report.append(f"  • No {', '.join(missing)} found - consider adding for variety")

            return "\n".join(report)

        def _count_overlaps(shapes):
            """Count overlapping shapes."""
            count = 0
            shape_list = list(shapes)

            for i, shape1 in enumerate(shape_list):
                for shape2 in shape_list[i + 1 :]:
                    if _shapes_overlap_simple(shape1, shape2):
                        count += 1
            return count

        def _count_out_of_bounds(shapes):
            """Count out of bounds shapes."""
            count = 0
            for shape in shapes:
                if hasattr(shape, "left"):
                    left = shape.left.inches if hasattr(shape.left, "inches") else 0
                    top = shape.top.inches if hasattr(shape.top, "inches") else 0
                    right = left + (shape.width.inches if hasattr(shape.width, "inches") else 0)
                    bottom = top + (shape.height.inches if hasattr(shape.height, "inches") else 0)

                    if left < 0 or top < 0 or right > SLIDE_WIDTH or bottom > SLIDE_HEIGHT:
                        count += 1
            return count

        def _shapes_overlap_simple(shape1, shape2):
            """Simple overlap check."""
            if not (hasattr(shape1, "left") and hasattr(shape2, "left")):
                return False

            # Skip title placeholders
            for shape in [shape1, shape2]:
                if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format.type in [
                            PP_PLACEHOLDER.TITLE,
                            PP_PLACEHOLDER.CENTER_TITLE,
                        ]:
                            return False

            l1 = shape1.left.inches if hasattr(shape1.left, "inches") else 0
            t1 = shape1.top.inches if hasattr(shape1.top, "inches") else 0
            r1 = l1 + (shape1.width.inches if hasattr(shape1.width, "inches") else 0)
            b1 = t1 + (shape1.height.inches if hasattr(shape1.height, "inches") else 0)

            l2 = shape2.left.inches if hasattr(shape2.left, "inches") else 0
            t2 = shape2.top.inches if hasattr(shape2.top, "inches") else 0
            r2 = l2 + (shape2.width.inches if hasattr(shape2.width, "inches") else 0)
            b2 = t2 + (shape2.height.inches if hasattr(shape2.height, "inches") else 0)

            return not (r1 <= l2 or r2 <= l1 or b1 <= t2 or b2 <= t1)

        return await _analyze_presentation()

    # Return the tools for external access
    return {
        "pptx_inspect_slide": pptx_inspect_slide,
        "pptx_fix_slide_layout": pptx_fix_slide_layout,
        "pptx_analyze_presentation_layout": pptx_analyze_presentation_layout,
    }
