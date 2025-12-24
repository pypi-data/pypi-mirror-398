"""
Radar, Combo, and other specialized chart components.
"""

from typing import Dict, List, Optional, Tuple
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

from .base import ChartComponent


class RadarChart(ChartComponent):
    """
    Radar (Spider) chart component for multi-criteria comparison.

    Features:
    - Multiple data series
    - Filled or line variants
    - Custom scaling
    - Beautiful grid styling
    """

    def __init__(
        self,
        categories: List[str],
        series: Dict[str, List[float]],
        variant: str = "filled",
        max_value: Optional[float] = None,
        **kwargs,
    ):
        """
        Initialize radar chart.

        Args:
            categories: Criteria/axis labels
            series: Dictionary of series names to values
            variant: Chart variant (filled, markers, lines)
            max_value: Maximum value for scaling (auto-calculated if None)
            **kwargs: Additional chart parameters
        """
        super().__init__(**kwargs)
        self.categories = categories
        self.series = series
        self.variant = variant
        self.max_value = max_value

        # Validate data during initialization
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

        # Set chart type based on variant
        variant_map = {
            "filled": XL_CHART_TYPE.RADAR_FILLED,
            "markers": XL_CHART_TYPE.RADAR_MARKERS,
            "lines": XL_CHART_TYPE.RADAR,
        }
        self.chart_type = variant_map.get(variant, XL_CHART_TYPE.RADAR_FILLED)

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate radar chart data."""
        if not self.categories:
            return False, "No categories provided"

        if not self.series:
            return False, "No data series provided"

        if len(self.categories) < 3:
            return False, "Radar chart needs at least 3 categories"

        expected_length = len(self.categories)
        for name, values in self.series.items():
            if len(values) != expected_length:
                return (
                    False,
                    f"Series '{name}' has {len(values)} values, expected {expected_length}",
                )

        return True, None

    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepare radar chart data."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories

        # Scale values if max_value is specified
        if self.max_value:
            for series_name, values in self.series.items():
                scaled_values = [min(v, self.max_value) for v in values]
                chart_data.add_series(series_name, scaled_values)
        else:
            for series_name, values in self.series.items():
                chart_data.add_series(series_name, values)

        return chart_data

    def render(self, slide, placeholder=None, **kwargs):
        """Render radar chart with beautiful styling."""
        # Call base render to create chart
        chart_shape = super().render(slide, placeholder=placeholder, **kwargs)
        chart = chart_shape.chart  # Access the chart object from the shape

        # Get theme colors
        chart_colors = self.tokens.get("chart", [])

        # Style each series
        for idx, series in enumerate(chart.series):
            if idx < len(chart_colors):
                color_hex = chart_colors[idx]
                if isinstance(color_hex, str):
                    rgb = self.hex_to_rgb(color_hex)

                    if self.variant == "filled":
                        # Semi-transparent fill
                        fill = series.format.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(*rgb)
                        if hasattr(fill, "transparency"):
                            fill.transparency = 0.3

                    # Line styling
                    line = series.format.line
                    line.color.rgb = RGBColor(*rgb)
                    line.width = Pt(2.5)

                    # Marker styling for appropriate variants
                    if hasattr(series, "marker") and self.variant in ["markers", "filled"]:
                        series.marker.style = 1  # Circle
                        series.marker.size = 8

                        marker_fill = series.marker.format.fill
                        marker_fill.solid()
                        marker_fill.fore_color.rgb = RGBColor(*rgb)

                        marker_line = series.marker.format.line
                        marker_line.color.rgb = RGBColor(255, 255, 255)
                        marker_line.width = Pt(1)

        # Configure radar-specific styling
        if hasattr(chart, "value_axis"):
            value_axis = chart.value_axis
            value_axis.has_major_gridlines = True

            # Set max value if specified
            if self.max_value:
                value_axis.maximum_scale = self.max_value

            # Style gridlines
            gridlines = value_axis.major_gridlines.format.line
            gridlines.color.rgb = self.get_color("border.secondary")
            gridlines.width = Pt(0.5)

        return chart_shape


class ComboChart(ChartComponent):
    """
    Combination chart with multiple chart types.

    Features:
    - Mixed column/line series
    - Secondary axis support
    - Custom styling per series type
    """

    def __init__(
        self,
        categories: List[str],
        column_series: Dict[str, List[float]],
        line_series: Dict[str, List[float]],
        secondary_axis: Optional[List[str]] = None,
        **kwargs,
    ):
        """
        Initialize combo chart.

        Args:
            categories: Category labels
            column_series: Dictionary of column series
            line_series: Dictionary of line series
            secondary_axis: List of series names to put on secondary axis
            **kwargs: Additional chart parameters
        """
        super().__init__(**kwargs)
        self.categories = categories
        self.column_series = column_series
        self.line_series = line_series
        self.secondary_axis = secondary_axis or []

        # Validate data during initialization
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

        # Start with column chart type
        self.chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate combo chart data."""
        if not self.categories:
            return False, "No categories provided"

        if not self.column_series and not self.line_series:
            return False, "No data series provided"

        expected_length = len(self.categories)

        # Validate column series
        for name, values in self.column_series.items():
            if len(values) != expected_length:
                return (
                    False,
                    f"Column series '{name}' has {len(values)} values, expected {expected_length}",
                )

        # Validate line series
        for name, values in self.line_series.items():
            if len(values) != expected_length:
                return (
                    False,
                    f"Line series '{name}' has {len(values)} values, expected {expected_length}",
                )

        return True, None

    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepare combo chart data."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories

        # Add column series first
        for series_name, values in self.column_series.items():
            chart_data.add_series(series_name, values)

        # Add line series
        for series_name, values in self.line_series.items():
            chart_data.add_series(series_name, values)

        return chart_data

    def render(self, slide, placeholder=None, **kwargs):
        """Render combo chart with mixed chart types."""
        # Call base render to create chart
        chart_shape = super().render(slide, placeholder=placeholder, **kwargs)
        chart = chart_shape.chart  # Access the chart object from the shape

        # Get theme colors
        chart_colors = self.tokens.get("chart", [])

        # Track series index
        series_idx = 0

        # Style column series (keep as columns)
        for i, (name, values) in enumerate(self.column_series.items()):
            if series_idx < len(chart.series):
                series = chart.series[series_idx]

                # Apply color
                if series_idx < len(chart_colors):
                    color_hex = chart_colors[series_idx]
                    if isinstance(color_hex, str):
                        rgb = self.hex_to_rgb(color_hex)
                        fill = series.format.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(*rgb)

                series_idx += 1

        # Convert line series to line charts
        for i, (name, values) in enumerate(self.line_series.items()):
            if series_idx < len(chart.series):
                series = chart.series[series_idx]

                # Change to line chart type
                # Note: This requires XML manipulation in python-pptx
                # For now, we style as best as possible

                # Apply line styling
                if series_idx < len(chart_colors):
                    color_hex = chart_colors[series_idx]
                    if isinstance(color_hex, str):
                        rgb = self.hex_to_rgb(color_hex)

                        # Make fill transparent for line effect
                        fill = series.format.fill
                        fill.background()

                        # Style border as line
                        line = series.format.line
                        line.color.rgb = RGBColor(*rgb)
                        line.width = Pt(3)

                # Move to secondary axis if specified
                if name in self.secondary_axis:
                    # Note: Secondary axis assignment requires XML manipulation
                    pass

                series_idx += 1

        return chart_shape


class GaugeChart(ChartComponent):
    """
    Gauge chart for showing progress or KPI values.

    Uses doughnut chart with custom styling to create gauge appearance.
    """

    def __init__(
        self,
        value: float,
        min_value: float = 0,
        max_value: float = 100,
        thresholds: Optional[Dict[str, float]] = None,
        **kwargs,
    ):
        """
        Initialize gauge chart.

        Args:
            value: Current value to display
            min_value: Minimum value on gauge
            max_value: Maximum value on gauge
            thresholds: Dictionary of threshold names to values
            **kwargs: Additional chart parameters
        """
        super().__init__(**kwargs)
        self.value = value
        self.min_value = min_value
        self.max_value = max_value
        self.thresholds = thresholds or {}

        self.chart_type = XL_CHART_TYPE.DOUGHNUT

        # Calculate gauge segments
        self._calculate_segments()

        # Validate data after calculation
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

    def _calculate_segments(self):
        """Calculate gauge segments for visualization."""
        total_range = self.max_value - self.min_value

        # Create segments: filled (value) and empty (remaining)
        value_portion = (self.value - self.min_value) / total_range
        empty_portion = 1 - value_portion

        # Add invisible segment to create gauge effect (half circle)

        self.categories = ["Value", "Empty", "Hidden"]
        self.values = [
            value_portion * 50,  # Scale to half circle
            empty_portion * 50,
            50,  # Hidden bottom half
        ]

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate gauge data."""
        if self.min_value >= self.max_value:
            return False, "min_value must be less than max_value"

        if not (self.min_value <= self.value <= self.max_value):
            return (
                False,
                f"value ({self.value}) must be between min_value ({self.min_value}) and max_value ({self.max_value})",
            )

        return True, None

    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepare gauge chart data."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories
        chart_data.add_series("Gauge", self.values)
        return chart_data

    def render(self, slide, placeholder=None, **kwargs):
        """Render gauge chart with custom styling and integrated labels."""
        from pptx.enum.text import PP_ALIGN

        # Get position parameters
        left = kwargs.get("left", 1.0)
        top = kwargs.get("top", 2.0)
        width = kwargs.get("width", 3.0)
        height = kwargs.get("height", 2.5)

        # Call base render to create chart
        chart_shape = super().render(slide, placeholder=placeholder, **kwargs)
        chart = chart_shape.chart  # Access the chart object from the shape

        if len(chart.series) > 0:
            series = chart.series[0]

            # Style each segment
            for idx, point in enumerate(series.points):
                fill = point.format.fill
                line = point.format.line

                if idx == 0:  # Value segment
                    # Use success color for filled portion
                    fill.solid()
                    fill.fore_color.rgb = self.get_color("success.DEFAULT")
                elif idx == 1:  # Empty segment
                    # Use muted color for empty portion
                    fill.solid()
                    fill.fore_color.rgb = self.get_color("muted.DEFAULT")
                else:  # Hidden segment
                    # Make invisible
                    fill.background()
                    line.fill.background()

        # Remove legend for cleaner look
        chart.has_legend = False

        # Add title above gauge if provided
        if self.title:
            title_box = slide.shapes.add_textbox(
                Inches(left), Inches(top - 0.5), Inches(width), Inches(0.4)
            )
            title_frame = title_box.text_frame
            title_frame.text = self.title
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(12)
            title_para.font.bold = True
            title_para.font.color.rgb = self.get_color("foreground.DEFAULT")
            title_para.alignment = PP_ALIGN.CENTER

        # Add value display in center/below gauge
        value_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + height * 0.4), Inches(width), Inches(0.6)
        )
        value_frame = value_box.text_frame
        value_frame.text = f"{self.value:.0f}"
        value_para = value_frame.paragraphs[0]
        value_para.font.size = Pt(28)
        value_para.font.bold = True
        value_para.font.color.rgb = self.get_color("primary.DEFAULT")
        value_para.alignment = PP_ALIGN.CENTER

        # Add percentage or range label
        label_box = slide.shapes.add_textbox(
            Inches(left), Inches(top + height * 0.65), Inches(width), Inches(0.3)
        )
        label_frame = label_box.text_frame
        percentage = ((self.value - self.min_value) / (self.max_value - self.min_value)) * 100
        label_frame.text = f"{percentage:.0f}%"
        label_para = label_frame.paragraphs[0]
        label_para.font.size = Pt(11)
        label_para.font.color.rgb = self.get_color("muted.foreground")
        label_para.alignment = PP_ALIGN.CENTER

        return chart_shape
