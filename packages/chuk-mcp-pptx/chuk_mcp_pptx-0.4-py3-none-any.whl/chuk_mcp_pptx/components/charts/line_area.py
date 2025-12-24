"""
Line and Area chart components with variants and registry integration.
"""

from typing import Dict, List, Optional, Tuple
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.text import MSO_ANCHOR
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

from .base import ChartComponent
from ..variants import LINE_CHART_VARIANTS
from ..registry import component, ComponentCategory, prop, example


@component(
    name="LineChart",
    category=ComponentCategory.CHART,
    description="Line chart component for showing trends and changes over time",
    props=[
        prop(
            "categories",
            "array",
            "Category labels for x-axis",
            required=True,
            example=["Jan", "Feb", "Mar", "Apr"],
        ),
        prop(
            "series",
            "object",
            "Dictionary of series names to values",
            required=True,
            example={"Revenue": [100, 120, 115, 130], "Costs": [80, 85, 90, 95]},
        ),
        prop(
            "variant",
            "string",
            "Chart variant",
            options=["line", "smooth", "area", "smooth_area"],
            default="line",
            example="line",
        ),
        prop(
            "style",
            "string",
            "Visual style preset",
            options=["default", "minimal", "detailed"],
            default="default",
            example="default",
        ),
        prop("title", "string", "Chart title", example="Monthly Trends"),
        prop(
            "legend",
            "string",
            "Legend position",
            options=["right", "bottom", "top", "none"],
            default="right",
            example="right",
        ),
    ],
    variants={
        "variant": ["line", "smooth", "area", "smooth_area"],
        "style": ["default", "minimal", "detailed"],
    },
    examples=[
        example(
            "Basic line chart",
            """
chart = LineChart(
    categories=["Jan", "Feb", "Mar", "Apr"],
    series={"Sales": [100, 120, 115, 130]},
    title="Monthly Sales Trend",
    variant="line"
)
chart.render(slide, left=1, top=2)
            """,
            categories=["Jan", "Feb", "Mar", "Apr"],
            series={"Sales": [100, 120, 115, 130]},
            variant="line",
        ),
        example(
            "Smooth line chart",
            """
chart = LineChart(
    categories=["Q1", "Q2", "Q3", "Q4"],
    series={"Revenue": [100, 150, 125, 175], "Target": [110, 140, 130, 180]},
    title="Revenue vs Target",
    variant="smooth",
    style="detailed"
)
chart.render(slide)
            """,
            categories=["Q1", "Q2", "Q3", "Q4"],
            series={"Revenue": [100, 150, 125, 175]},
            variant="smooth",
            style="detailed",
        ),
    ],
    tags=["chart", "line", "trend", "time-series"],
)
class LineChart(ChartComponent):
    """
    Line chart component for showing trends over time.

    Features:
    - Multiple variants (line, smooth, area, smooth_area)
    - Theme-aware styling
    - Data validation
    - Customizable markers and lines

    Variants:
    - line: Standard line chart with markers
    - smooth: Smoothed curves with markers
    - area: Line chart with filled area below
    - smooth_area: Smoothed line with filled area

    Styles:
    - default: Standard appearance with grid
    - minimal: Clean look without grid or markers
    - detailed: Shows values on data points
    """

    def __init__(
        self,
        categories: List[str],
        series: Dict[str, List[float]],
        variant: str = "line",
        style: str = "default",
        **kwargs,
    ):
        """
        Initialize line chart.

        Args:
            categories: Category labels (x-axis)
            series: Dictionary of series names to values
            variant: Chart variant (line, smooth, area, smooth_area)
            style: Visual style (default, minimal, detailed)
            **kwargs: Additional chart parameters (title, theme, legend, etc.)
        """
        super().__init__(style=style, **kwargs)
        self.categories = categories
        self.series = series
        self.variant = variant

        # Validate data during initialization
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

        # Update variant props to include line-specific variants
        self.variant_props = LINE_CHART_VARIANTS.build(variant=variant, style=style)

        # Set chart type based on variant
        if "area" in variant:
            self.chart_type = XL_CHART_TYPE.AREA
        else:
            self.chart_type = XL_CHART_TYPE.LINE_MARKERS

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate line chart data."""
        if not self.categories:
            return False, "No categories provided"

        if not self.series:
            return False, "No data series provided"

        expected_length = len(self.categories)
        for name, values in self.series.items():
            if len(values) != expected_length:
                return (
                    False,
                    f"Series '{name}' has {len(values)} values, expected {expected_length}",
                )

        return True, None

    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepare line chart data."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories

        for series_name, values in self.series.items():
            chart_data.add_series(series_name, values)

        return chart_data

    def render(self, slide, placeholder=None, **kwargs):
        """Render line chart with theme styling."""
        chart_shape = super().render(slide, placeholder=placeholder, **kwargs)
        chart = chart_shape.chart  # Access the chart object from the shape

        # Get theme colors
        chart_colors = self.tokens.get("chart", [])

        # Style each series
        for idx, series in enumerate(chart.series):
            # Apply colors
            if idx < len(chart_colors):
                color_hex = chart_colors[idx]
                rgb = self.hex_to_rgb(color_hex)
                series.format.line.color.rgb = RGBColor(*rgb)
                series.format.line.width = Pt(2)

            # Apply smooth if requested
            smooth = self.variant_props.get("smooth", False)
            if hasattr(series, "smooth"):
                series.smooth = smooth

            # Configure markers
            show_markers = self.variant_props.get("show_markers", True)
            if hasattr(series, "marker") and show_markers:
                series.marker.style = XL_MARKER_STYLE.CIRCLE
                series.marker.size = 6

        # Hide grid for minimal style
        if not self.variant_props.get("show_grid", True):
            if hasattr(chart, "category_axis"):
                chart.category_axis.has_major_gridlines = False
            if hasattr(chart, "value_axis"):
                chart.value_axis.has_major_gridlines = False

        return chart_shape


@component(
    name="AreaChart",
    category=ComponentCategory.CHART,
    description="Area chart component for showing cumulative trends with filled areas",
    props=[
        prop(
            "categories",
            "array",
            "Category labels for x-axis",
            required=True,
            example=["Q1", "Q2", "Q3", "Q4"],
        ),
        prop(
            "series",
            "object",
            "Dictionary of series names to values",
            required=True,
            example={"Sales": [100, 120, 130, 150]},
        ),
        prop(
            "variant",
            "string",
            "Chart variant",
            options=["area", "stacked", "stacked100"],
            default="area",
            example="area",
        ),
        prop(
            "style",
            "string",
            "Visual style preset",
            options=["default", "minimal", "detailed"],
            default="default",
            example="default",
        ),
        prop("title", "string", "Chart title", example="Sales Trends"),
    ],
    variants={
        "variant": ["area", "stacked", "stacked100"],
        "style": ["default", "minimal", "detailed"],
    },
    examples=[
        example(
            "Basic area chart",
            """
chart = AreaChart(
    categories=["Jan", "Feb", "Mar"],
    series={"Revenue": [100, 120, 140]},
    title="Revenue Growth"
)
chart.render(slide)
            """,
            categories=["Jan", "Feb", "Mar"],
            series={"Revenue": [100, 120, 140]},
        )
    ],
    tags=["chart", "area", "cumulative", "trend"],
)
class AreaChart(ChartComponent):
    """
    Area chart component for showing cumulative trends.

    Inherits from ChartComponent and uses area chart types.
    Perfect for showing magnitude changes over time.
    """

    def __init__(
        self,
        categories: List[str],
        series: Dict[str, List[float]],
        variant: str = "area",
        style: str = "default",
        **kwargs,
    ):
        """
        Initialize area chart.

        Args:
            categories: Category labels
            series: Dictionary of series names to values
            variant: Chart variant (area, stacked, stacked100)
            style: Visual style (default, minimal, detailed)
            **kwargs: Additional chart parameters
        """
        super().__init__(style=style, **kwargs)
        self.categories = categories
        self.series = series
        self.variant = variant

        # Validate data
        is_valid, error = self.validate_data()
        if not is_valid:
            raise ValueError(f"Invalid chart data: {error}")

        # Set chart type based on variant
        variant_map = {
            "area": XL_CHART_TYPE.AREA,
            "stacked": XL_CHART_TYPE.AREA_STACKED,
            "stacked100": XL_CHART_TYPE.AREA_STACKED_100,
            "3d": XL_CHART_TYPE.THREE_D_AREA,
        }
        self.chart_type = variant_map.get(variant, XL_CHART_TYPE.AREA)

    def validate_data(self) -> Tuple[bool, Optional[str]]:
        """Validate area chart data."""
        if not self.categories:
            return False, "No categories provided"

        if not self.series:
            return False, "No data series provided"

        expected_length = len(self.categories)
        for name, values in self.series.items():
            if len(values) != expected_length:
                return (
                    False,
                    f"Series '{name}' has {len(values)} values, expected {expected_length}",
                )

        return True, None

    def _prepare_chart_data(self) -> CategoryChartData:
        """Prepare area chart data."""
        chart_data = CategoryChartData()
        chart_data.categories = self.categories

        for series_name, values in self.series.items():
            chart_data.add_series(series_name, values)

        return chart_data


@component(
    name="SparklineChart",
    category=ComponentCategory.CHART,
    description="Sparkline chart component for compact inline trend visualization",
    props=[
        prop("values", "array", "Data values", required=True, example=[10, 15, 13, 17, 14, 20]),
        prop(
            "categories",
            "array",
            "Optional category labels",
            example=["Jan", "Feb", "Mar", "Apr", "May", "Jun"],
        ),
        prop(
            "style",
            "string",
            "Visual style",
            options=["minimal", "default"],
            default="minimal",
            example="minimal",
        ),
    ],
    examples=[
        example(
            "Minimal sparkline",
            """
chart = SparklineChart(
    values=[10, 15, 13, 17, 14, 20],
    style="minimal"
)
chart.render(slide, left=1, top=1, width=2, height=0.5)
            """,
            values=[10, 15, 13, 17, 14, 20],
            style="minimal",
        )
    ],
    tags=["chart", "sparkline", "micro", "inline"],
)
class SparklineChart(LineChart):
    """
    Sparkline chart component for compact trend visualization.

    Simplified line chart designed for small spaces.
    Perfect for dashboards and inline metrics.
    """

    def __init__(
        self,
        values: List[float],
        categories: Optional[List[str]] = None,
        label: Optional[str] = None,
        show_value: bool = True,
        **kwargs,
    ):
        """
        Initialize sparkline chart.

        Args:
            values: Data values
            categories: Optional category labels (auto-generated if not provided)
            label: Optional label to show on left
            show_value: Whether to show current value on right
            **kwargs: Additional chart parameters
        """
        # Auto-generate categories if not provided
        if categories is None:
            categories = [str(i) for i in range(len(values))]

        # Store sparkline-specific properties
        self.values = values
        self.label = label
        self.show_value = show_value

        # Force minimal styling for sparklines
        kwargs["style"] = kwargs.get("style", "minimal")
        kwargs["legend"] = "none"

        super().__init__(categories=categories, series={"Value": values}, variant="line", **kwargs)

    def render(self, slide, placeholder=None, **kwargs):
        """Render sparkline with minimal styling and integrated labels."""
        from pptx.enum.text import PP_ALIGN

        # Get position parameters
        left = kwargs.get("left", 1.0)
        top = kwargs.get("top", 2.0)
        width = kwargs.get("width", 2.0)
        height = kwargs.get("height", 0.5)

        # Reserve space for label and value if needed
        chart_left = left
        chart_width = width

        if self.label:
            # Add 2 inches for label on left
            chart_left += 2.0
            chart_width -= 2.0

        if self.show_value:
            # Reserve 0.8 inches for value on right
            chart_width -= 0.8

        # Update kwargs with adjusted position
        kwargs["left"] = chart_left
        kwargs["width"] = chart_width
        kwargs["height"] = height

        chart_shape = super().render(slide, placeholder=placeholder, **kwargs)
        chart = chart_shape.chart  # Access the chart object from the shape

        # Remove axes for true sparkline effect
        if hasattr(chart, "category_axis"):
            chart.category_axis.visible = False
        if hasattr(chart, "value_axis"):
            chart.value_axis.visible = False

        # Add label on left if provided
        if self.label:
            label_box = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(1.8), Inches(height)
            )
            label_frame = label_box.text_frame
            label_frame.text = self.label
            label_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            label_para = label_frame.paragraphs[0]
            label_para.font.size = Pt(11)
            label_para.font.bold = True
            label_para.font.color.rgb = self.get_color("foreground.DEFAULT")
            label_para.alignment = PP_ALIGN.LEFT

        # Add current value on right if enabled
        if self.show_value:
            current_value = self.values[-1]
            value_box = slide.shapes.add_textbox(
                Inches(chart_left + chart_width + 0.1), Inches(top), Inches(0.7), Inches(height)
            )
            value_frame = value_box.text_frame
            value_frame.text = f"{current_value:.0f}"
            value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            value_para = value_frame.paragraphs[0]
            value_para.font.size = Pt(10)
            value_para.font.bold = True
            value_para.font.color.rgb = self.get_color("primary.DEFAULT")
            value_para.alignment = PP_ALIGN.RIGHT

        return chart_shape
