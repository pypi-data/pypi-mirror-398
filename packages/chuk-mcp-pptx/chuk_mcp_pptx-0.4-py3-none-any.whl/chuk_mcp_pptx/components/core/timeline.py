# src/chuk_mcp_pptx/components/core/timeline.py
"""
Timeline components for PowerPoint presentations.

Provides horizontal timeline visualization for events, milestones,
project phases, and roadmaps.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import get_text_style, FONT_SIZES


class Timeline(Component):
    """
    Timeline component - horizontal event/milestone visualization.

    Perfect for project timelines, product roadmaps, company history,
    and chronological data presentation.

    Variants:
        - default: Standard timeline with primary color
        - minimal: Clean minimal style
        - highlighted: Emphasize key milestones

    Styles:
        - line: Simple line with markers (default)
        - arrow: Line with arrow endpoint
        - segmented: Divided segments between events

    Examples:
        # Simple timeline
        events = [
            {"date": "Jan 2024", "title": "Project Start"},
            {"date": "Mar 2024", "title": "Beta Release"},
            {"date": "Jun 2024", "title": "Launch"}
        ]
        timeline = Timeline(events, theme=theme)
        timeline.render(slide, left=1, top=3, width=8)

        # Detailed timeline with descriptions
        events = [
            {
                "date": "Q1 2024",
                "title": "Planning",
                "description": "Initial planning and research"
            },
            {
                "date": "Q2 2024",
                "title": "Development",
                "description": "Build core features"
            },
            {
                "date": "Q3 2024",
                "title": "Launch",
                "description": "Public release",
                "highlight": True
            }
        ]
        timeline = Timeline(events, variant="highlighted", theme=theme)
        timeline.render(slide, left=1, top=3, width=8)
    """

    def __init__(
        self,
        events: List[Dict[str, Any]],
        variant: str = "default",
        style: str = "line",
        show_descriptions: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize timeline.

        Args:
            events: List of event dicts with keys:
                   - date: Event date/time (required)
                   - title: Event title (required)
                   - description: Event description (optional)
                   - highlight: Whether to highlight (optional)
            variant: Visual variant
            style: Timeline style (line, arrow, segmented)
            show_descriptions: Show event descriptions
            theme: Optional theme
        """
        super().__init__(theme)
        self.events = events
        self.variant = variant
        self.style = style
        self.show_descriptions = show_descriptions

    def _get_line_color(self) -> RGBColor:
        """Get timeline line color."""
        return self.get_color("border.DEFAULT")

    def _get_marker_color(self, is_highlighted: bool = False) -> RGBColor:
        """Get event marker color."""
        if is_highlighted or self.variant == "highlighted":
            return self.get_color("primary.DEFAULT")
        return self.get_color("muted.foreground")

    def _get_highlight_color(self) -> RGBColor:
        """Get highlight color for special events."""
        return self.get_color("primary.DEFAULT")

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self, slide, left: float, top: float, width: float = 8.0, placeholder: Optional[Any] = None
    ) -> list:
        """
        Render timeline to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Timeline width in inches
            placeholder: Optional placeholder to replace

        Returns:
            List of rendered shapes
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        shapes = []

        if not self.events:
            return shapes

        # Calculate positions
        num_events = len(self.events)
        if num_events == 1:
            event_spacing = 0
        else:
            event_spacing = width / (num_events - 1)

        # Render timeline line
        line_top = top + 0.6  # Position line below dates
        shapes.extend(self._render_line(slide, left, line_top, width))

        # Render events
        for i, event in enumerate(self.events):
            event_left = left + (i * event_spacing)
            is_highlighted = event.get("highlight", False)

            # Render marker
            marker_shape = self._render_marker(slide, event_left, line_top, is_highlighted)
            shapes.append(marker_shape)

            # Render date
            date_shape = self._render_date(slide, event_left, top, event.get("date", ""))
            shapes.append(date_shape)

            # Render title
            title_shape = self._render_title(
                slide, event_left, line_top + 0.25, event.get("title", ""), is_highlighted
            )
            shapes.append(title_shape)

            # Render description if enabled
            if self.show_descriptions and "description" in event:
                desc_shape = self._render_description(
                    slide, event_left, line_top + 0.55, event["description"]
                )
                shapes.append(desc_shape)

        return shapes

    def _render_line(self, slide, left: float, top: float, width: float) -> list:
        """Render timeline line."""
        shapes = []

        if self.style == "arrow":
            # Line with arrow
            line = slide.shapes.add_connector(
                1,  # Straight connector
                Inches(left),
                Inches(top),
                Inches(left + width),
                Inches(top),
            )
            line.line.color.rgb = self._get_line_color()
            line.line.width = Pt(2)
            # Add arrowhead
            line.line.end_arrow_type = 2  # Arrow
            shapes.append(line)
        else:
            # Simple line
            line = slide.shapes.add_connector(
                1,  # Straight connector
                Inches(left),
                Inches(top),
                Inches(left + width),
                Inches(top),
            )
            line.line.color.rgb = self._get_line_color()
            line.line.width = Pt(2)
            shapes.append(line)

        return shapes

    def _render_marker(self, slide, left: float, top: float, is_highlighted: bool) -> Any:
        """Render event marker."""
        marker_size = 0.2 if is_highlighted else 0.15

        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(left - marker_size / 2),
            Inches(top - marker_size / 2),
            Inches(marker_size),
            Inches(marker_size),
        )

        marker.fill.solid()
        marker.fill.fore_color.rgb = self._get_marker_color(is_highlighted)
        marker.line.color.rgb = self._get_marker_color(is_highlighted)

        if is_highlighted:
            marker.line.width = Pt(2)

        return marker

    def _render_date(self, slide, left: float, top: float, date_text: str) -> Any:
        """Render event date."""
        font_family = self._get_font_family()
        date_box = slide.shapes.add_textbox(
            Inches(left - 0.75), Inches(top), Inches(1.5), Inches(0.3)
        )

        text_frame = date_box.text_frame
        text_frame.text = date_text
        text_frame.margin_top = 0
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style = get_text_style("small")
        p.font.name = font_family
        p.font.size = Pt(style["font_size"])
        p.font.color.rgb = self.get_color("muted.foreground")
        p.font.bold = True

        return date_box

    def _render_title(
        self, slide, left: float, top: float, title_text: str, is_highlighted: bool
    ) -> Any:
        """Render event title."""
        font_family = self._get_font_family()
        title_box = slide.shapes.add_textbox(
            Inches(left - 1.0), Inches(top), Inches(2.0), Inches(0.3)
        )

        text_frame = title_box.text_frame
        text_frame.text = title_text
        text_frame.margin_top = 0
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        style = get_text_style("small")
        p.font.name = font_family
        p.font.size = Pt(style["font_size"])

        if is_highlighted:
            p.font.color.rgb = self._get_highlight_color()
            p.font.bold = True
        else:
            p.font.color.rgb = self.get_color("foreground.DEFAULT")

        return title_box

    def _render_description(self, slide, left: float, top: float, description_text: str) -> Any:
        """Render event description."""
        font_family = self._get_font_family()
        desc_box = slide.shapes.add_textbox(
            Inches(left - 1.0), Inches(top), Inches(2.0), Inches(0.5)
        )

        text_frame = desc_box.text_frame
        text_frame.text = description_text
        text_frame.margin_top = 0
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        get_text_style("small")
        p.font.name = font_family
        p.font.size = Pt(FONT_SIZES["xs"])  # Smaller for descriptions
        p.font.color.rgb = self.get_color("muted.foreground")

        return desc_box


# TODO: Register component when registry is implemented
# Component metadata for documentation:
# - Variants: default, minimal, highlighted
# - Styles: line, arrow, segmented
# - Props: events, variant, style, show_descriptions
