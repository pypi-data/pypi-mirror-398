# src/chuk_mcp_pptx/components/core/divider.py
"""
Divider component for visual separation.
"""

from typing import Optional, Dict, Any, Literal
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

from ..base import Component
from ...layout.helpers import CONTENT_WIDTH, CONTENT_HEIGHT
from ..registry import component, ComponentCategory, prop, example


@component(
    name="Divider",
    category=ComponentCategory.LAYOUT,
    description="Visual divider line for separating content",
    props=[
        prop(
            "orientation",
            "string",
            "Divider orientation",
            options=["horizontal", "vertical"],
            default="horizontal",
        ),
        prop("thickness", "number", "Line thickness in points", default=1),
        prop("color", "string", "Line color token", default="border.DEFAULT"),
    ],
    examples=[
        example(
            "Horizontal divider",
            """
divider = Divider(orientation="horizontal", thickness=1)
divider.render(slide, left=0.5, top=3.0, width=9.0)
            """,
            orientation="horizontal",
        )
    ],
    tags=["layout", "divider", "separator"],
)
class Divider(Component):
    """
    Divider line for visual separation.

    Usage:
        # Horizontal divider
        divider = Divider(orientation="horizontal")
        divider.render(slide, left=0.5, top=3.0, width=9.0)

        # Vertical divider
        divider = Divider(orientation="vertical")
        divider.render(slide, left=5.0, top=1.5, height=4.0)
    """

    def __init__(
        self,
        orientation: Literal["horizontal", "vertical"] = "horizontal",
        thickness: float = 1,
        color: str = "border.DEFAULT",
        theme: Optional[Dict[str, Any]] = None,
    ):
        super().__init__(theme)
        self.orientation = orientation
        self.thickness = thickness
        self.color = color

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ):
        """Render divider line."""
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        if self.orientation == "horizontal":
            # Horizontal line
            line_width = width or CONTENT_WIDTH
            line_height = Pt(self.thickness)

            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(line_width), line_height
            )
        else:
            # Vertical line
            line_width = Pt(self.thickness)
            line_height = height or CONTENT_HEIGHT

            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), line_width, Inches(line_height)
            )

        # Style line
        line.fill.solid()
        line.fill.fore_color.rgb = self.get_color(self.color)
        line.line.fill.background()

        return line
