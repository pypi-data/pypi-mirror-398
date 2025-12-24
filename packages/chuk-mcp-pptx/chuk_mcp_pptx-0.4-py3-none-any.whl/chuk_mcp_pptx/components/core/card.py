# src/chuk_mcp_pptx/components/core/card.py
"""
Enhanced Card component with variants and composition support.
Uses the new variant system and compositional API.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from ..composition import (
    ComposableComponent,
    CardHeader,
    CardContent,
    CardFooter,
    CardTitle,
    CardDescription,
)
from ..variants import CARD_VARIANTS
from ..registry import component, ComponentCategory, prop, example
from ...tokens.typography import FONT_SIZES, PARAGRAPH_SPACING


@component(
    name="Card",
    category=ComponentCategory.CONTAINER,
    description="Versatile container component with support for variants and composition",
    props=[
        prop(
            "variant",
            "string",
            "Visual style variant",
            options=["default", "outlined", "elevated", "ghost"],
            default="default",
            example="outlined",
        ),
        prop(
            "padding",
            "string",
            "Padding size",
            options=["none", "sm", "md", "lg", "xl"],
            default="md",
            example="lg",
        ),
        prop("left", "number", "Left position in inches", required=True, example=1.0),
        prop("top", "number", "Top position in inches", required=True, example=1.0),
        prop("width", "number", "Width in inches", default=3.0, example=4.0),
        prop("height", "number", "Height in inches", default=2.0, example=3.0),
    ],
    variants={
        "variant": ["default", "outlined", "elevated", "ghost"],
        "padding": ["none", "sm", "md", "lg", "xl"],
    },
    composition={
        "supports": ["CardHeader", "CardTitle", "CardDescription", "CardContent", "CardFooter"],
        "builder": "CompositionBuilder",
    },
    examples=[
        example(
            "Basic card with header and content",
            """
card = Card(variant="outlined", padding="lg")
card.add_child(CardHeader("Welcome", "Get started with our platform"))
card.add_child(CardContent("This is the main content area"))
card.render(slide, left=1, top=1)
            """,
            variant="outlined",
            padding="lg",
        ),
        example(
            "Metric card with composition builder",
            """
builder = CompositionBuilder(theme)
children = builder.title("Revenue").description("$1.2M").badge("↑ 12%", "success").build()

card = Card(variant="elevated")
for child in children:
    card.add_child(child)
card.render(slide, left=2, top=2)
            """,
            variant="elevated",
        ),
    ],
    tags=["container", "card", "layout", "composition"],
)
class Card(ComposableComponent):
    """
    Enhanced Card component with variant and composition support.

    Features:
    - Multiple visual variants (default, outlined, elevated, ghost)
    - Flexible padding options
    - Composition support via subcomponents
    - Theme-aware styling

    Usage:
        # Simple card
        card = Card(variant="outlined", padding="md")
        card.render(slide, left=1, top=1, width=3, height=2)

        # Composed card
        card = Card(variant="elevated")
        card.add_child(CardHeader("Title", "Subtitle"))
        card.add_child(CardContent("Main content"))
        card.add_child(CardFooter("Footer text"))
        card.render(slide, left=1, top=1)
    """

    # Expose composition components as class attributes (shadcn style)
    Header = CardHeader
    Title = CardTitle
    Description = CardDescription
    Content = CardContent
    Footer = CardFooter

    def __init__(
        self, variant: str = "default", padding: str = "md", theme: Optional[Dict[str, Any]] = None
    ):
        """
        Initialize enhanced card.

        Args:
            variant: Visual variant (default, outlined, elevated, ghost)
            padding: Padding size (none, sm, md, lg, xl)
            theme: Optional theme override
        """
        super().__init__(theme)
        self.variant = variant
        self.padding = padding

        # Get variant props
        self.variant_props = CARD_VARIANTS.build(variant=variant, padding=padding)

    def _calculate_min_width(self) -> float:
        """Calculate minimum width needed for content."""
        if not self._children:
            return 3.0

        # Estimate based on title length if present
        for child in self._children:
            if hasattr(child, "text"):
                # Rough estimate: 0.08 inches per character for titles
                text_len = len(child.text)
                min_width = (text_len * 0.08) + 1.0  # Add padding
                return max(2.5, min(min_width, 6.0))  # Cap between 2.5 and 6 inches

        return 3.0

    def _calculate_min_height(self) -> float:
        """
        Calculate minimum height needed for content.

        Uses actual font sizes and padding to ensure text fits properly.
        """
        if not self._children:
            return 1.5

        # Base height for padding (margins) - must match actual render padding
        padding = self.variant_props.get("padding", 0.5)
        total_height = padding * 2  # Top and bottom margins

        # Estimate height per child component based on actual font sizes
        for child in self._children:
            if hasattr(child, "text") and hasattr(child, "__class__"):
                class_name = child.__class__.__name__
                if "Title" in class_name:
                    # h5 font is 16pt = ~0.22" + space before (6pt = ~0.08") + line spacing
                    total_height += 0.35
                elif "Description" in class_name:
                    # body font is 14pt with 6pt space before
                    # More accurate: account for character width based on card width
                    text_len = len(child.text)
                    # Assume ~70 chars per line for typical card widths
                    lines = max(1, (text_len + 69) // 70)  # Ceiling division
                    # 14pt = ~0.19" per line, plus 6pt = ~0.08" space before
                    total_height += 0.08 + (lines * 0.22)  # More generous spacing
                else:
                    total_height += 0.3  # Default content height

        # Add a small buffer for rendering
        total_height += 0.1

        return max(1.5, min(total_height, 4.5))  # Cap between 1.5 and 4.5 inches

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render card to slide or replace a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Card width in inches (optional, auto-calculated if None)
            height: Card height in inches (optional, auto-calculated if None)
            placeholder: Optional placeholder shape to replace

        Returns:
            Shape object representing the card

        Note:
            For cards with text content, omit height parameter to use auto-calculation
            which ensures all content fits properly.
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            import logging

            logger = logging.getLogger(__name__)
            logger.info(
                f"Card targeting placeholder - using bounds: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
            )

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Use calculated width and height if not provided
        card_width = width if width is not None else self._calculate_min_width()

        # Calculate minimum height needed
        calculated_height = self._calculate_min_height()

        # Use the larger of provided height or calculated minimum
        # This ensures content always fits
        if height is not None:
            card_height = max(height, calculated_height)
        else:
            card_height = calculated_height

        # Create card shape
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(card_width),
            Inches(card_height),
        )

        # Apply variant styles
        self._apply_variant_styles(card)

        # Setup text frame
        text_frame = card.text_frame
        text_frame.clear()
        text_frame.word_wrap = True  # Enable wrapping for descriptions
        text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Anchor text to top of shape

        # Apply padding from design system (theme takes priority, then variant, then default)
        padding_inches = self.get_padding("md")  # Uses theme if available
        variant_padding = self.variant_props.get("padding")
        if variant_padding is not None:
            padding_inches = variant_padding
        text_frame.margin_left = Inches(padding_inches)
        text_frame.margin_right = Inches(padding_inches)
        text_frame.margin_top = Inches(padding_inches)
        text_frame.margin_bottom = Inches(padding_inches)

        # Render children if any
        if self._children:
            for child in self._children:
                child.render_into(text_frame, self.theme)

        return card

    def _apply_variant_styles(self, shape):
        """Apply variant-based styling to shape using design tokens."""
        props = self.variant_props

        # Background color
        if props.get("bg_color") and props["bg_color"] != "transparent":
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.get_color(props["bg_color"])
        else:
            shape.fill.background()

        # Border - use design token for width if not specified in variant
        border_width = props.get("border_width")
        if border_width is None:
            border_width = self.get_border_width("2")  # Default from design tokens
        if border_width > 0:
            border_color = props.get("border_color", "border.DEFAULT")
            shape.line.color.rgb = self.get_color(border_color)
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()

        # Border radius - apply from design tokens
        # Note: PowerPoint rounded rectangles have built-in radius, but we can adjust
        try:
            # Attempt to set corner radius if shape supports it
            if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
                radius = self.get_border_radius("md")
                # Convert points to adjustment value (0-1 scale)
                shape.adjustments[0] = min(0.5, radius / 100)
        except Exception:
            pass  # Some shapes don't support adjustments

        # Shadow (for elevated variant)
        if props.get("shadow"):
            shape.shadow.visible = True
            shape.shadow.blur_radius = Pt(self.get_font_size("xs"))
            shape.shadow.distance = Pt(4)
            shape.shadow.angle = 90
            shape.shadow.transparency = 0.3


@component(
    name="MetricCard",
    category=ComponentCategory.DATA,
    description="Specialized card for displaying KPIs and metrics with trend indicators",
    props=[
        prop("label", "string", "Metric label", required=True, example="Revenue"),
        prop("value", "string", "Metric value", required=True, example="$1.2M"),
        prop("change", "string", "Change indicator", example="+12%"),
        prop("trend", "string", "Trend direction", options=["up", "down", "neutral"], example="up"),
        prop(
            "variant",
            "string",
            "Visual variant",
            options=["default", "outlined", "elevated"],
            default="outlined",
        ),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("width", "number", "Width in inches", default=2.5),
        prop("height", "number", "Height in inches", default=1.5),
    ],
    examples=[
        example(
            "Revenue metric with upward trend",
            """
metric = MetricCard(
    label="Revenue",
    value="$1.2M",
    change="+12%",
    trend="up",
    variant="elevated"
)
metric.render(slide, left=1, top=1)
            """,
            label="Revenue",
            value="$1.2M",
            change="+12%",
            trend="up",
        )
    ],
    tags=["metric", "kpi", "data", "card"],
)
class MetricCard(Card):
    """
    Enhanced metric card with variant support.
    Displays KPIs with optional trend indicators.
    """

    def __init__(
        self,
        label: str,
        value: str,
        change: Optional[str] = None,
        trend: Optional[str] = None,
        variant: str = "outlined",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize metric card.

        Args:
            label: Metric label
            value: Metric value
            change: Change indicator (e.g., "+12%")
            trend: Trend direction (up/down/neutral)
            variant: Visual variant
            theme: Optional theme
        """
        super().__init__(variant=variant, padding="md", theme=theme)
        self.label = label
        self.value = value
        self.change = change
        self.trend = trend

    def _calculate_min_width(self) -> float:
        """Calculate minimum width based on value length."""
        # Estimate width needed for the value (which is the largest text)
        # 20pt font is roughly 0.10 inches per character
        value_width = len(str(self.value)) * 0.10
        min_width = value_width + 1.2  # Add padding for margins
        return max(2.0, min(min_width, 3.5))  # Cap between 2.0 and 3.5 inches

    def get_trend_color(self):
        """Get color based on trend."""
        if self.trend == "up":
            return self.get_color("success.DEFAULT")
        elif self.trend == "down":
            return self.get_color("destructive.DEFAULT")
        else:
            return self.get_color("muted.foreground")

    def get_trend_symbol(self) -> str:
        """Get trend symbol."""
        symbols = {"up": "↑", "down": "↓", "neutral": "→"}
        return symbols.get(self.trend, "")

    def _calculate_min_height(self) -> float:
        """Calculate minimum height based on metric content."""
        # Base padding (top + bottom)
        padding = self.variant_props.get("padding", 0.5)
        total_height = padding * 2

        # Label: 12pt = ~0.2"
        total_height += 0.2

        # Value: 20pt = ~0.3" + 4pt space before
        total_height += 0.35

        # Change (if present): 12pt = ~0.2" + 4pt space before
        if self.change:
            total_height += 0.25

        return max(1.2, min(total_height, 2.5))

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ):
        """Render metric card."""
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds
            import logging

            logger = logging.getLogger(__name__)
            logger.info(
                f"MetricCard targeting placeholder - using bounds: ({left:.2f}, {top:.2f}, {width:.2f}, {height:.2f})"
            )

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Use calculated width and height if not provided
        card_width = width if width is not None else self._calculate_min_width()
        card_height = height if height is not None else self._calculate_min_height()

        # Create base card (don't pass placeholder since we already extracted bounds)
        card_shape = super().render(slide, left, top, card_width, card_height)

        # Clear and rebuild content
        text_frame = card_shape.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_ANCHOR.TOP  # Anchor text to top

        padding = self.variant_props.get("padding", 0.5)
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        # Label
        p = text_frame.paragraphs[0]
        p.text = self.label
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(FONT_SIZES["sm"])
        p.font.color.rgb = self.get_color("muted.foreground")
        p.font.name = self.get_theme_attr("font_family", "Inter")

        # Value
        p = text_frame.add_paragraph()
        p.text = self.value
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(PARAGRAPH_SPACING["xs"])
        p.font.size = Pt(FONT_SIZES["2xl"])  # 22pt (was 20pt, using closest design token)
        p.font.bold = True
        p.font.color.rgb = self.get_color(self.variant_props.get("fg_color", "card.foreground"))
        p.font.name = self.get_theme_attr("font_family", "Inter")

        # Change with trend
        if self.change:
            p = text_frame.add_paragraph()
            symbol = self.get_trend_symbol()
            p.text = f"{symbol} {self.change}" if symbol else self.change
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(PARAGRAPH_SPACING["xs"])
            p.font.size = Pt(FONT_SIZES["sm"])
            p.font.color.rgb = self.get_trend_color()
            p.font.name = self.get_theme_attr("font_family", "Inter")

        return card_shape
