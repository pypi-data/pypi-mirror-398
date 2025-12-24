# src/chuk_mcp_pptx/components/core/button.py
"""
Enhanced Button component with variants and composition support.
Uses the variant system and component registry.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

from ..composition import ComposableComponent
from ..variants import BUTTON_VARIANTS
from ..registry import component, ComponentCategory, prop, example
from ...constants import ComponentSizing


@component(
    name="Button",
    category=ComponentCategory.UI,
    description="Interactive button component with multiple variants and sizes",
    props=[
        prop("text", "string", "Button label text", required=True, example="Click me"),
        prop(
            "variant",
            "string",
            "Visual variant",
            options=["default", "secondary", "outline", "ghost", "destructive"],
            default="default",
            example="default",
        ),
        prop(
            "size", "string", "Button size", options=["sm", "md", "lg"], default="md", example="md"
        ),
        prop("left", "number", "Left position in inches", required=True, example=1.0),
        prop("top", "number", "Top position in inches", required=True, example=1.0),
        prop("width", "number", "Width in inches (optional)", example=2.0),
        prop("height", "number", "Height in inches (optional)", example=0.5),
    ],
    variants={
        "variant": ["default", "secondary", "outline", "ghost", "destructive"],
        "size": ["sm", "md", "lg"],
    },
    examples=[
        example(
            "Primary action button",
            """
button = Button(text="Submit", variant="default", size="md")
button.render(slide, left=1, top=1)
            """,
            text="Submit",
            variant="default",
            size="md",
        ),
        example(
            "Destructive action",
            """
button = Button(text="Delete", variant="destructive", size="sm")
button.render(slide, left=1, top=2)
            """,
            text="Delete",
            variant="destructive",
            size="sm",
        ),
        example(
            "Ghost button",
            """
button = Button(text="Cancel", variant="ghost", size="lg")
button.render(slide, left=1, top=3)
            """,
            text="Cancel",
            variant="ghost",
            size="lg",
        ),
    ],
    tags=["button", "ui", "interactive", "action"],
)
class Button(ComposableComponent):
    """
    Enhanced button component with variant support.

    Features:
    - Multiple visual variants (default, secondary, outline, ghost, destructive)
    - Three sizes (sm, md, lg)
    - Theme-aware coloring
    - Composition support for icon buttons

    Usage:
        # Simple button
        button = Button(text="Click me", variant="default", size="md")
        button.render(slide, left=1, top=1)

        # With theme
        button = Button(text="Submit", variant="default", theme=theme_dict)
        button.render(slide, left=2, top=2, width=3, height=0.6)
    """

    def __init__(
        self,
        text: str,
        variant: str = "default",
        size: str = "md",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize button component.

        Args:
            text: Button label text
            variant: Visual variant (default, secondary, outline, ghost, destructive)
            size: Button size (sm, md, lg)
            theme: Optional theme override
        """
        super().__init__(theme)
        self.text = text
        self.variant = variant
        self.size = size

        # Get variant props
        self.variant_props = BUTTON_VARIANTS.build(variant=variant, size=size)

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render button to slide.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Optional width override (uses size default if None)
            height: Optional height override (uses size default if None)
            placeholder: Optional placeholder to replace

        Returns:
            Shape object representing the button
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Get dimensions from variant props or use provided
        btn_width = width if width is not None else self._get_default_width()
        btn_height = height if height is not None else self.variant_props.get("height", 0.5)

        # Create button shape
        button = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(btn_width),
            Inches(btn_height),
        )

        # Apply variant styling
        self._apply_variant_styles(button)

        # Add text
        self._add_text(button)

        return button

    def _get_default_width(self) -> float:
        """Get default width based on text length and size."""
        # Rough estimation: base width + character width
        base_widths = {
            "sm": ComponentSizing.BUTTON_BASE_WIDTH_SM,
            "md": ComponentSizing.BUTTON_BASE_WIDTH_MD,
            "lg": ComponentSizing.BUTTON_BASE_WIDTH_LG,
        }
        char_widths = {
            "sm": ComponentSizing.CHAR_WIDTH_SM,
            "md": ComponentSizing.CHAR_WIDTH_MD,
            "lg": ComponentSizing.CHAR_WIDTH_LG,
        }

        base = base_widths.get(self.size, 2.0)
        char_width = char_widths.get(self.size, 0.07)

        return max(base, len(self.text) * char_width + 0.5)

    def _apply_variant_styles(self, shape):
        """Apply variant-based styling to button shape."""
        props = self.variant_props

        # Background color
        bg_color = props.get("bg_color")
        if bg_color and bg_color != "transparent":
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.get_color(bg_color)
        else:
            shape.fill.background()  # Transparent

        # Border
        border_width = props.get("border_width", 0)
        if border_width > 0:
            border_color = props.get("border_color", "border.DEFAULT")
            shape.line.color.rgb = self.get_color(border_color)
            shape.line.width = Pt(border_width)
        else:
            shape.line.fill.background()  # No border

    def _add_text(self, shape):
        """Add text to button with proper styling."""
        props = self.variant_props

        text_frame = shape.text_frame
        text_frame.clear()

        # Disable word wrap for single-line buttons
        text_frame.word_wrap = False

        # Set padding
        padding = props.get("padding", 0.3)
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding / 2)
        text_frame.margin_bottom = Inches(padding / 2)

        # Add text
        paragraph = text_frame.paragraphs[0]
        paragraph.text = self.text
        paragraph.alignment = PP_ALIGN.CENTER

        # Font styling
        paragraph.font.size = Pt(props.get("font_size", 14))
        # Handle both Theme objects and dict themes
        if hasattr(self.theme, "typography"):
            font_family = self.theme.typography.get("font_family", "Inter")
        else:
            font_family = (
                self.theme.get("font_family", "Inter") if isinstance(self.theme, dict) else "Inter"
            )
        paragraph.font.name = font_family

        # Text color
        fg_color = props.get("fg_color", "foreground.DEFAULT")
        paragraph.font.color.rgb = self.get_color(fg_color)

        # Font weight (bold for emphasis)
        if props.get("font_weight", 400) >= 600:
            paragraph.font.bold = True


@component(
    name="IconButton",
    category=ComponentCategory.UI,
    description="Button component with icon instead of text, perfect for actions",
    props=[
        prop("icon", "string", "Icon name or Unicode character", required=True, example="play"),
        prop(
            "variant",
            "string",
            "Visual variant",
            options=["default", "secondary", "outline", "ghost", "destructive"],
            default="ghost",
            example="ghost",
        ),
        prop(
            "size", "string", "Button size", options=["sm", "md", "lg"], default="md", example="md"
        ),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
    ],
    variants={
        "variant": ["default", "secondary", "outline", "ghost", "destructive"],
        "size": ["sm", "md", "lg"],
    },
    examples=[
        example(
            "Play icon button",
            """
button = IconButton(icon="play", variant="ghost")
button.render(slide, left=1, top=1)
            """,
            icon="play",
            variant="ghost",
        )
    ],
    tags=["button", "icon", "ui", "action"],
)
class IconButton(Button):
    """
    Icon button component - button with an icon instead of text.
    Uses Unicode symbols for icons.

    Common icons available:
    - play, pause, stop, next, previous
    - plus, minus, close, check
    - star, heart, settings, menu
    - search, download, upload, refresh
    - share, edit, delete
    """

    # Icon mappings using Unicode
    ICONS = {
        "play": "▶",
        "pause": "⏸",
        "stop": "⏹",
        "next": "⏭",
        "previous": "⏮",
        "plus": "+",
        "minus": "-",
        "close": "✕",
        "check": "✓",
        "star": "★",
        "heart": "♥",
        "settings": "⚙",
        "menu": "☰",
        "search": "🔍",
        "download": "⬇",
        "upload": "⬆",
        "refresh": "↻",
        "share": "⤴",
        "edit": "✎",
        "delete": "🗑",
        "info": "ℹ",
        "warning": "⚠",
        "error": "✖",
        "success": "✔",
    }

    def __init__(
        self,
        icon: str,
        variant: str = "ghost",
        size: str = "md",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize icon button.

        Args:
            icon: Icon name from ICONS dict or Unicode character
            variant: Visual variant
            size: Button size
            theme: Optional theme override
        """
        # Get icon character or use provided Unicode
        icon_char = self.ICONS.get(icon, icon)
        super().__init__(icon_char, variant, size, theme)

    def _get_default_width(self) -> float:
        """Icon buttons are square."""
        sizes = {"sm": 0.5, "md": 0.6, "lg": 0.8}
        return sizes.get(self.size, 0.6)

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render icon button (square by default).

        Args:
            slide: PowerPoint slide
            left: Left position
            top: Top position
            width: Optional width (defaults to square)
            height: Optional height (defaults to square)
            placeholder: Optional placeholder to replace
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Icon buttons are square by default
        size = width or self._get_default_width()
        return super().render(slide, left, top, size, height or size)


@component(
    name="ButtonGroup",
    category=ComponentCategory.UI,
    description="Group of buttons displayed together",
    props=[
        prop("buttons", "array", "List of button configurations", required=True),
        prop(
            "orientation",
            "string",
            "Layout orientation",
            options=["horizontal", "vertical"],
            default="horizontal",
            example="horizontal",
        ),
        prop("spacing", "number", "Space between buttons in inches", default=0.1, example=0.1),
    ],
    examples=[
        example(
            "Horizontal button group",
            """
group = ButtonGroup(
    buttons=[
        {"text": "Save", "variant": "default"},
        {"text": "Cancel", "variant": "ghost"}
    ],
    orientation="horizontal"
)
group.render(slide, left=1, top=1)
            """,
            orientation="horizontal",
        )
    ],
    tags=["button", "group", "ui", "layout"],
)
class ButtonGroup(ComposableComponent):
    """
    Group of buttons displayed together.
    Handles layout and spacing automatically.
    """

    def __init__(
        self,
        buttons: list,
        orientation: str = "horizontal",
        spacing: float = 0.1,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize button group.

        Args:
            buttons: List of button configs [{"text": "...", "variant": "...", ...}]
            orientation: Layout direction (horizontal/vertical)
            spacing: Space between buttons in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.buttons = buttons
        self.orientation = orientation
        self.spacing = spacing

    def render(self, slide, left: float, top: float, placeholder: Optional[Any] = None) -> list:
        """
        Render button group.

        Args:
            slide: PowerPoint slide
            left: Starting left position
            top: Starting top position
            placeholder: Optional placeholder to replace

        Returns:
            List of button shapes
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        shapes = []
        current_left = left
        current_top = top

        for btn_config in self.buttons:
            # Create button
            btn_text = btn_config.get("text", "Button")
            btn_variant = btn_config.get("variant", "default")
            btn_size = btn_config.get("size", "md")

            button = Button(btn_text, btn_variant, btn_size, self.theme)
            shape = button.render(slide, current_left, current_top)
            shapes.append(shape)

            # Update position for next button
            if self.orientation == "horizontal":
                current_left += shape.width.inches + self.spacing
            else:  # vertical
                current_top += shape.height.inches + self.spacing

        return shapes
