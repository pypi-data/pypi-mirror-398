# src/chuk_mcp_pptx/components/core/image.py
"""
Image component for PowerPoint presentations.

Provides image placement with effects and transformations.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pathlib import Path
import base64
import io
import asyncio
from PIL import Image as PILImage, ImageFilter, ImageEnhance

from ..base import Component
from ..registry import component, ComponentCategory, prop, example


@component(
    name="Image",
    category=ComponentCategory.UI,
    description="Image component with effects, transformations, and filters",
    props=[
        prop("image_source", "string", "Image path or base64 data", required=True),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("width", "number", "Width in inches (optional, maintains ratio)"),
        prop("height", "number", "Height in inches (optional, maintains ratio)"),
        prop("shadow", "boolean", "Add shadow effect", default=False),
        prop("glow", "boolean", "Add glow effect", default=False),
        prop("reflection", "boolean", "Add reflection effect", default=False),
        prop("blur_radius", "number", "Blur radius (0 = no blur)", default=0),
        prop("grayscale", "boolean", "Convert to grayscale", default=False),
        prop("sepia", "boolean", "Apply sepia tone effect", default=False),
        prop(
            "brightness",
            "number",
            "Brightness adjustment (1.0 = normal, <1 = darker, >1 = brighter)",
            default=1.0,
        ),
        prop(
            "contrast",
            "number",
            "Contrast adjustment (1.0 = normal, <1 = less contrast, >1 = more contrast)",
            default=1.0,
        ),
        prop(
            "saturation",
            "number",
            "Saturation adjustment (1.0 = normal, 0 = grayscale, >1 = more vibrant)",
            default=1.0,
        ),
        prop("sharpen", "boolean", "Apply sharpening filter", default=False),
        prop("invert", "boolean", "Invert colors (negative)", default=False),
    ],
    examples=[
        example(
            "Simple image",
            """
image = Image(image_source="path/to/image.png")
image.render(slide, left=2, top=2, width=4)
            """,
            image_source="logo.png",
        ),
        example(
            "Image with shadow",
            """
image = Image(
    image_source="photo.jpg",
    shadow=True
)
image.render(slide, left=1, top=1, width=5, height=3)
            """,
            image_source="photo.jpg",
            shadow=True,
        ),
        example(
            "Image with filters",
            """
image = Image(
    image_source="photo.jpg",
    blur_radius=5,
    brightness=1.2,
    contrast=1.3
)
image.render(slide, left=1, top=1, width=5, height=3)
            """,
            image_source="photo.jpg",
            blur_radius=5,
            brightness=1.2,
        ),
    ],
    tags=["image", "picture", "photo", "media", "filter", "blur", "grayscale"],
)
class Image(Component):
    """
    Image component for adding pictures to slides.

    Features:
    - File path or base64 data support
    - Automatic aspect ratio maintenance
    - Visual effects (shadow, glow, reflection)
    - Flexible sizing

    Usage:
        # Simple image from file
        image = Image(image_source="logo.png")
        image.render(slide, left=2, top=2, width=4)

        # Image with effects
        image = Image(
            image_source="photo.jpg",
            shadow=True,
            glow=True
        )
        image.render(slide, left=1, top=1, width=5, height=3)

        # Base64 image
        image = Image(image_source="data:image/png;base64,iVBOR...")
        image.render(slide, left=3, top=3)
    """

    def __init__(
        self,
        image_source: str,
        shadow: bool = False,
        glow: bool = False,
        reflection: bool = False,
        blur_radius: float = 0,
        grayscale: bool = False,
        sepia: bool = False,
        brightness: float = 1.0,
        contrast: float = 1.0,
        saturation: float = 1.0,
        sharpen: bool = False,
        invert: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize image component.

        Args:
            image_source: Path to image file or base64 data URL
            shadow: Whether to add shadow effect
            glow: Whether to add glow effect
            reflection: Whether to add reflection effect
            blur_radius: Blur radius (0 = no blur)
            grayscale: Convert to grayscale
            sepia: Apply sepia tone effect
            brightness: Brightness adjustment (1.0 = normal)
            contrast: Contrast adjustment (1.0 = normal)
            saturation: Saturation adjustment (1.0 = normal)
            sharpen: Apply sharpening filter
            invert: Invert colors (negative)
            theme: Optional theme override
        """
        super().__init__(theme)

        # Validate required parameters
        if not image_source:
            raise ValueError("Image requires 'image_source' (URL or file path)")
        if not isinstance(image_source, str):
            raise TypeError(
                f"Image 'image_source' must be a string, got {type(image_source).__name__}"
            )

        self.image_source = image_source
        self.shadow = shadow
        self.glow = glow
        self.reflection = reflection
        self.blur_radius = blur_radius
        self.grayscale = grayscale
        self.sepia = sepia
        self.brightness = brightness
        self.contrast = contrast
        self.saturation = saturation
        self.sharpen = sharpen
        self.invert = invert

    async def render(
        self,
        slide,
        left: float,
        top: float,
        width: Optional[float] = None,
        height: Optional[float] = None,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render image to slide or into a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches (optional)
            height: Height in inches (optional)
            placeholder: Optional placeholder shape to populate

        Returns:
            Picture shape object
        """
        # Check if any filters need to be applied
        needs_processing = (
            self.blur_radius > 0
            or self.grayscale
            or self.sepia
            or self.brightness != 1.0
            or self.contrast != 1.0
            or self.saturation != 1.0
            or self.sharpen
            or self.invert
        )

        # Get image source (file path, stream, or URL)
        image_source_for_insertion = None

        if needs_processing:
            # Load image for processing
            pil_image = await self._load_image()

            # Apply filters
            pil_image = self._apply_filters(pil_image)

            # Convert to stream for PowerPoint (wrap blocking I/O)
            image_stream = io.BytesIO()
            await asyncio.to_thread(pil_image.save, image_stream, format="PNG")
            image_stream.seek(0)
            image_source_for_insertion = image_stream
        else:
            # No processing needed, use original
            # Handle HTTP/HTTPS URLs
            if self.image_source.startswith(("http://", "https://")):
                import urllib.request

                try:
                    # Download image from URL - scheme already validated above (http/https only)
                    with urllib.request.urlopen(self.image_source) as response:  # nosec B310
                        image_data = response.read()
                    image_stream = io.BytesIO(image_data)
                    image_source_for_insertion = image_stream
                except Exception as e:
                    raise FileNotFoundError(
                        f"Could not download image from URL: {self.image_source}. Error: {str(e)}"
                    )

            # Handle base64 image data
            elif self.image_source.startswith("data:image/"):
                header, encoded = self.image_source.split(",", 1)
                image_data = base64.b64decode(encoded)
                image_stream = io.BytesIO(image_data)
                image_source_for_insertion = image_stream

            # Handle file path
            elif Path(self.image_source).exists():
                image_source_for_insertion = self.image_source

            else:
                raise FileNotFoundError(f"Image not found: {self.image_source}")

        # Insert into placeholder if provided, otherwise add to slide
        if placeholder is not None:
            # Use placeholder's insert_picture method
            try:
                pic = placeholder.insert_picture(image_source_for_insertion)
            except AttributeError:
                # Fallback if placeholder doesn't support insert_picture
                pic = self._add_picture(slide, image_source_for_insertion, left, top, width, height)
        else:
            pic = self._add_picture(slide, image_source_for_insertion, left, top, width, height)

        # Apply PowerPoint effects
        if self.shadow:
            self._apply_shadow(pic)

        return pic

    def _add_picture(self, slide, image_source, left, top, width, height):
        """Add picture to slide with specified dimensions."""
        if width and height:
            return slide.shapes.add_picture(
                image_source, Inches(left), Inches(top), width=Inches(width), height=Inches(height)
            )
        elif width:
            return slide.shapes.add_picture(
                image_source, Inches(left), Inches(top), width=Inches(width)
            )
        elif height:
            return slide.shapes.add_picture(
                image_source, Inches(left), Inches(top), height=Inches(height)
            )
        else:
            return slide.shapes.add_picture(image_source, Inches(left), Inches(top))

    async def _load_image(self) -> PILImage.Image:
        """Load image from source (file path or base64)."""
        if self.image_source.startswith("data:image/"):
            # Handle base64 data
            header, encoded = self.image_source.split(",", 1)
            image_data = base64.b64decode(encoded)
            image_stream = io.BytesIO(image_data)
            # Wrap blocking I/O in asyncio.to_thread
            return await asyncio.to_thread(PILImage.open, image_stream)
        elif Path(self.image_source).exists():
            # Handle file path - wrap blocking I/O in asyncio.to_thread
            return await asyncio.to_thread(PILImage.open, self.image_source)
        else:
            raise FileNotFoundError(f"Image not found: {self.image_source}")

    def _apply_filters(self, pil_image: PILImage.Image) -> PILImage.Image:
        """Apply PIL filters to image."""
        # Convert to RGB if needed
        if pil_image.mode != "RGB":
            pil_image = pil_image.convert("RGB")

        # Apply blur
        if self.blur_radius > 0:
            pil_image = pil_image.filter(ImageFilter.GaussianBlur(radius=self.blur_radius))

        # Apply sharpen
        if self.sharpen:
            pil_image = pil_image.filter(ImageFilter.SHARPEN)

        # Apply brightness
        if self.brightness != 1.0:
            enhancer = ImageEnhance.Brightness(pil_image)
            pil_image = enhancer.enhance(self.brightness)

        # Apply contrast
        if self.contrast != 1.0:
            enhancer = ImageEnhance.Contrast(pil_image)
            pil_image = enhancer.enhance(self.contrast)

        # Apply saturation
        if self.saturation != 1.0:
            enhancer = ImageEnhance.Color(pil_image)
            pil_image = enhancer.enhance(self.saturation)

        # Apply grayscale (overrides saturation)
        if self.grayscale:
            pil_image = pil_image.convert("L").convert("RGB")

        # Apply sepia
        if self.sepia:
            pil_image = self._apply_sepia(pil_image)

        # Apply invert
        if self.invert:
            from PIL import ImageOps

            pil_image = ImageOps.invert(pil_image)

        return pil_image

    def _apply_sepia(self, pil_image: PILImage.Image) -> PILImage.Image:
        """Apply sepia tone filter."""
        # Sepia matrix transformation
        width, height = pil_image.size
        pixels = pil_image.load()

        for y in range(height):
            for x in range(width):
                r, g, b = pixels[x, y]

                # Sepia formula
                tr = int(0.393 * r + 0.769 * g + 0.189 * b)
                tg = int(0.349 * r + 0.686 * g + 0.168 * b)
                tb = int(0.272 * r + 0.534 * g + 0.131 * b)

                # Clamp values
                pixels[x, y] = (min(255, tr), min(255, tg), min(255, tb))

        return pil_image

    def _apply_shadow(self, picture_shape):
        """Apply shadow effect to picture."""
        shadow_format = picture_shape.shadow
        shadow_format.inherit = False
        shadow_format.visible = True
        shadow_format.distance = Pt(4)
        shadow_format.blur_radius = Pt(4)
        shadow_format.transparency = 0.5
        shadow_format.angle = 45
