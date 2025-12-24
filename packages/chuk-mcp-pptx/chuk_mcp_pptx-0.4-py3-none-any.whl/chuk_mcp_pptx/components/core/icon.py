# src/chuk_mcp_pptx/components/core/icon.py
"""
Icon components for PowerPoint presentations.

Provides standardized icon symbols for visual communication.
Uses Unicode symbols and shapes for common business icons.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from ..base import Component


# Icon mappings using Unicode symbols
ICON_SYMBOLS = {
    # Status & Actions
    "check": "✓",
    "checkmark": "✓",
    "x": "✕",
    "cross": "✕",
    "close": "✕",
    "plus": "+",
    "minus": "−",
    # Arrows & Direction
    "arrow-up": "↑",
    "arrow-down": "↓",
    "arrow-left": "←",
    "arrow-right": "→",
    "chevron-up": "⌃",
    "chevron-down": "⌄",
    "chevron-left": "‹",
    "chevron-right": "›",
    # Common Symbols
    "star": "★",
    "star-outline": "☆",
    "heart": "❤",
    "warning": "⚠",
    "info": "ℹ",
    "question": "?",
    "exclamation": "!",
    "bell": "🔔",
    "calendar": "📅",
    "clock": "🕐",
    "mail": "✉",
    "phone": "📞",
    "location": "📍",
    "pin": "📌",
    "tag": "🏷",
    "bookmark": "🔖",
    "flag": "⚑",
    # Business & Office
    "document": "📄",
    "folder": "📁",
    "chart": "📊",
    "graph": "📈",
    "trending-up": "📈",
    "trending-down": "📉",
    "briefcase": "💼",
    "dollar": "$",
    "euro": "€",
    "pound": "£",
    # People & Social
    "user": "👤",
    "users": "👥",
    "team": "👥",
    "crown": "👑",
    "trophy": "🏆",
    # Tech & UI
    "gear": "⚙",
    "settings": "⚙",
    "search": "🔍",
    "filter": "⊗",
    "menu": "☰",
    "home": "🏠",
    "link": "🔗",
    "lock": "🔒",
    "unlock": "🔓",
    "download": "⬇",
    "upload": "⬆",
    "refresh": "⟳",
    # Misc
    "lightbulb": "💡",
    "rocket": "🚀",
    "target": "🎯",
    "key": "🔑",
    "fire": "🔥",
    "zap": "⚡",
    "circle": "●",
    "square": "■",
    "diamond": "◆",
    "dot": "•",
}


class Icon(Component):
    """
    Icon component - standardized icon symbols.

    Provides common business and UI icons using Unicode symbols.
    Perfect for bullet points, status indicators, and visual markers.

    Variants:
        - default: Foreground color
        - primary: Primary color
        - success: Green
        - warning: Yellow/orange
        - error: Red
        - muted: Muted color

    Sizes:
        - sm: 12pt (0.17")
        - md: 16pt (0.22") - default
        - lg: 20pt (0.28")
        - xl: 24pt (0.33")

    Examples:
        # Simple checkmark
        icon = Icon("check", variant="success", theme=theme)
        icon.render(slide, left=1, top=2)

        # Large star
        icon = Icon("star", size="xl", variant="warning", theme=theme)
        icon.render(slide, left=1, top=2)

        # Custom symbol
        icon = Icon("→", variant="primary", theme=theme)
        icon.render(slide, left=1, top=2)
    """

    # Size mapping in points
    SIZE_MAP = {
        "sm": 12,
        "md": 16,
        "lg": 20,
        "xl": 24,
        "2xl": 32,
    }

    def __init__(
        self,
        icon: str,
        variant: str = "default",
        size: str = "md",
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize icon.

        Args:
            icon: Icon name or Unicode symbol
            variant: Color variant
            size: Icon size (sm, md, lg, xl, 2xl)
            theme: Optional theme
        """
        super().__init__(theme)

        # Get symbol - either from map or use as-is
        self.symbol = ICON_SYMBOLS.get(icon, icon)
        self.variant = variant
        self.size = size

    def _get_icon_color(self) -> RGBColor:
        """Get color based on variant."""
        color_map = {
            "default": "foreground.DEFAULT",
            "primary": "primary.DEFAULT",
            "success": "success.DEFAULT",
            "warning": "warning.DEFAULT",
            "error": "destructive.DEFAULT",
            "muted": "muted.foreground",
        }
        color_path = color_map.get(self.variant, "foreground.DEFAULT")
        return self.get_color(color_path)

    def _get_size_inches(self) -> float:
        """Get icon size in inches for box dimensions."""
        pt_size = self.SIZE_MAP.get(self.size, 16)
        # Convert pt to inches with some padding
        return (pt_size / 72) + 0.1

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(self, slide, left: float, top: float, placeholder: Optional[Any] = None) -> Any:
        """
        Render icon to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            placeholder: Optional placeholder to replace

        Returns:
            Shape containing the icon
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        from pptx.util import Inches

        font_family = self._get_font_family()
        size_inches = self._get_size_inches()
        pt_size = self.SIZE_MAP.get(self.size, 16)

        # Create text box for icon
        icon_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(size_inches), Inches(size_inches)
        )

        text_frame = icon_box.text_frame
        text_frame.text = self.symbol
        text_frame.margin_top = 0
        text_frame.margin_bottom = 0
        text_frame.margin_left = 0
        text_frame.margin_right = 0
        text_frame.vertical_anchor = 1  # Middle

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.name = font_family
        p.font.size = Pt(pt_size)
        p.font.color.rgb = self._get_icon_color()

        return icon_box


class IconList(Component):
    """
    Icon list component - list with icon bullets.

    Creates a vertical list where each item has an icon prefix.
    Perfect for feature lists, requirements, checklists.

    Examples:
        # Feature list with checkmarks
        items = [
            ("check", "Fast performance"),
            ("check", "Easy to use"),
            ("check", "Fully customizable")
        ]
        icon_list = IconList(items, variant="success", theme=theme)
        icon_list.render(slide, left=1, top=2, width=6)

        # Mixed icons
        items = [
            ("rocket", "Fast deployment"),
            ("lock", "Secure by default"),
            ("zap", "High performance")
        ]
        icon_list = IconList(items, theme=theme)
        icon_list.render(slide, left=1, top=2, width=6)
    """

    def __init__(
        self,
        items: list,
        variant: str = "default",
        icon_size: str = "md",
        spacing: float = 0.35,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize icon list.

        Args:
            items: List of (icon, text) tuples
            variant: Color variant for icons
            icon_size: Icon size
            spacing: Spacing between items in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.items = items
        self.variant = variant
        self.icon_size = icon_size
        self.spacing = spacing

    def _get_font_family(self) -> str:
        """Get font family from theme."""
        return self.get_theme_attr("font_family", "Calibri")

    def render(
        self, slide, left: float, top: float, width: float = 6.0, placeholder: Optional[Any] = None
    ) -> list:
        """
        Render icon list to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Total width in inches
            placeholder: Optional placeholder to replace

        Returns:
            List of rendered shapes
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        from pptx.util import Inches
        from ...tokens.typography import get_text_style

        font_family = self._get_font_family()
        shapes = []
        current_top = top
        icon_width = 0.4  # Space for icon

        for icon_name, text in self.items:
            # Render icon
            icon = Icon(icon_name, variant=self.variant, size=self.icon_size, theme=self.theme)
            icon_shape = icon.render(slide, left=left, top=current_top)
            shapes.append(icon_shape)

            # Render text
            text_box = slide.shapes.add_textbox(
                Inches(left + icon_width),
                Inches(current_top),
                Inches(width - icon_width),
                Inches(0.3),
            )
            text_frame = text_box.text_frame
            text_frame.text = text
            text_frame.margin_top = 0
            text_frame.word_wrap = True

            p = text_frame.paragraphs[0]
            style = get_text_style("body")
            p.font.name = font_family
            p.font.size = Pt(style["font_size"])
            p.font.color.rgb = self.get_color("foreground.DEFAULT")

            shapes.append(text_box)
            current_top += self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# Icon - Variants: default, primary, success, warning, error, muted
# Icon - Sizes: sm, md, lg, xl, 2xl
# IconList - Props: items, variant, icon_size, spacing
