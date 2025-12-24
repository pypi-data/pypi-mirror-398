# src/chuk_mcp_pptx/components/core/shape.py
"""
Shape components for PowerPoint presentations.

Provides basic geometric shapes with design system integration.
"""

from typing import Optional, Dict, Any
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from ..base import Component
from ..registry import component, ComponentCategory, prop, example


# Shape type mapping
SHAPE_TYPES = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,  # Alias
    "diamond": MSO_SHAPE.DIAMOND,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
    "arrow_right": MSO_SHAPE.RIGHT_ARROW,
    "arrow_left": MSO_SHAPE.LEFT_ARROW,
    "arrow_up": MSO_SHAPE.UP_ARROW,
    "arrow_down": MSO_SHAPE.DOWN_ARROW,
    "star": MSO_SHAPE.STAR_5_POINT,
    "star_5": MSO_SHAPE.STAR_5_POINT,
    "star_6": MSO_SHAPE.STAR_6_POINT,
    "star_7": MSO_SHAPE.STAR_7_POINT,
    "hexagon": MSO_SHAPE.HEXAGON,
    "chevron": MSO_SHAPE.CHEVRON,
    "plus": MSO_SHAPE.MATH_PLUS,
    "minus": MSO_SHAPE.MATH_MINUS,
    "multiply": MSO_SHAPE.MATH_MULTIPLY,
    "divide": MSO_SHAPE.MATH_DIVIDE,
    "callout": MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
    "cloud": MSO_SHAPE.CLOUD,
    "heart": MSO_SHAPE.HEART,
    "lightning": MSO_SHAPE.LIGHTNING_BOLT,
}


@component(
    name="Shape",
    category=ComponentCategory.UI,
    description="Geometric shape component with theme integration",
    props=[
        prop(
            "shape_type",
            "string",
            "Type of shape",
            options=list(SHAPE_TYPES.keys()),
            default="rectangle",
            example="star",
        ),
        prop("text", "string", "Text content", example="Click me"),
        prop("fill_color", "string", "Fill color (hex or semantic)", example="#FF5733"),
        prop("line_color", "string", "Border color", example="primary.DEFAULT"),
        prop("line_width", "number", "Border width in points", default=1.0),
        prop("left", "number", "Left position in inches", required=True),
        prop("top", "number", "Top position in inches", required=True),
        prop("width", "number", "Width in inches", required=True),
        prop("height", "number", "Height in inches", required=True),
    ],
    examples=[
        example(
            "Star shape",
            """
shape = Shape(shape_type="star", fill_color="warning.DEFAULT")
shape.render(slide, left=2, top=2, width=2, height=2)
            """,
            shape_type="star",
        ),
        example(
            "Rounded rectangle with text",
            """
shape = Shape(
    shape_type="rounded_rectangle",
    text="Important",
    fill_color="primary.DEFAULT"
)
shape.render(slide, left=1, top=1, width=3, height=1.5)
            """,
            shape_type="rounded_rectangle",
            text="Important",
        ),
    ],
    tags=["shape", "geometry", "basic", "ui"],
)
class Shape(Component):
    """
    Basic geometric shape component.

    Features:
    - 25+ shape types (rectangles, circles, arrows, stars, etc.)
    - Theme-aware colors
    - Text content support
    - Customizable borders

    Usage:
        # Simple shape
        shape = Shape(shape_type="oval", fill_color="primary.DEFAULT")
        shape.render(slide, left=2, top=2, width=3, height=2)

        # Shape with text
        shape = Shape(
            shape_type="rounded_rectangle",
            text="Success!",
            fill_color="success.DEFAULT"
        )
        shape.render(slide, left=1, top=1, width=4, height=2)
    """

    def __init__(
        self,
        shape_type: str = "rectangle",
        text: Optional[str] = None,
        fill_color: Optional[str] = None,
        line_color: Optional[str] = None,
        line_width: float = 1.0,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize shape component.

        Args:
            shape_type: Type of shape (see SHAPE_TYPES)
            text: Optional text content
            fill_color: Fill color (hex or semantic color path)
            line_color: Border color (hex or semantic color path)
            line_width: Border width in points
            theme: Optional theme override
        """
        super().__init__(theme)
        self.shape_type = shape_type
        self.text = text
        self.fill_color = fill_color
        self.line_color = line_color
        self.line_width = line_width

    def render(
        self,
        slide,
        left: float,
        top: float,
        width: float,
        height: float,
        placeholder: Optional[Any] = None,
    ) -> Any:
        """
        Render shape to slide or replace a placeholder.

        Args:
            slide: PowerPoint slide object
            left: Left position in inches
            top: Top position in inches
            width: Width in inches
            height: Height in inches
            placeholder: Optional placeholder shape to replace

        Returns:
            Shape object
        """
        # If placeholder provided, extract bounds and delete it
        bounds = self._extract_placeholder_bounds(placeholder)
        if bounds is not None:
            left, top, width, height = bounds

        # Delete placeholder after extracting bounds
        self._delete_placeholder_if_needed(placeholder)

        # Get MSO shape type
        mso_shape = SHAPE_TYPES.get(self.shape_type.lower(), MSO_SHAPE.RECTANGLE)

        # Create shape
        shape = slide.shapes.add_shape(
            mso_shape, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Apply fill color
        if self.fill_color:
            color_rgb = self._parse_color(self.fill_color)
            shape.fill.solid()
            shape.fill.fore_color.rgb = color_rgb
        else:
            # Use theme default
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.get_color("accent.DEFAULT")

        # Apply line color and width
        if self.line_color:
            color_rgb = self._parse_color(self.line_color)
            shape.line.color.rgb = color_rgb
        else:
            shape.line.color.rgb = self.get_color("border.DEFAULT")

        shape.line.width = Pt(self.line_width)

        # Add text if provided
        if self.text and shape.has_text_frame:
            self._add_text(shape)

        return shape

    def _parse_color(self, color_str: str) -> RGBColor:
        """Parse color string (hex or semantic path)."""
        if color_str.startswith("#"):
            # Parse hex color
            color_str = color_str[1:]
            return RGBColor(
                int(color_str[0:2], 16), int(color_str[2:4], 16), int(color_str[4:6], 16)
            )
        else:
            # Use semantic color from theme
            return self.get_color(color_str)

    def _add_text(self, shape):
        """Add text content to shape."""
        text_frame = shape.text_frame
        text_frame.text = self.text
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.1)
        text_frame.margin_right = Inches(0.1)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)

        # Center text
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Set text color from theme
        paragraph.font.color.rgb = self.get_color("foreground.DEFAULT")
