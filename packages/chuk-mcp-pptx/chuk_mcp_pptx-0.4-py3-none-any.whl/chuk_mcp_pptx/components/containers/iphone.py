"""
iPhone device container component for PowerPoint presentations.

Provides an authentic iPhone mockup frame for displaying chat conversations
and other mobile content.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES


class iPhoneContainer(Component):
    """
    iPhone device container component.

    Creates a realistic iPhone frame with:
    - Rounded corners matching iPhone design
    - Notch at the top (modern iPhone style)
    - Status bar (time, signal, battery)
    - Home indicator at bottom
    - Dark or light theme support
    - Content area for rendering conversations

    Examples:
        # iPhone with iMessage conversation
        iphone = iPhoneContainer(
            title="iMessage",
            show_notch=True,
            theme=theme
        )

        # Render the container and get content area bounds
        content_area = iphone.render(slide, left=1, top=1)

        # Render conversation inside the content area
        conversation.render(
            slide,
            left=content_area['left'],
            top=content_area['top'],
            width=content_area['width']
        )
    """

    def __init__(
        self,
        title: str = "iPhone",
        show_notch: bool = True,
        variant: str = "pro",  # pro, pro-max, standard
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize iPhone container.

        Args:
            title: Title shown in the status bar area
            show_notch: Whether to show the notch (True for modern iPhones)
            variant: iPhone variant (pro, pro-max, standard)
            theme: Optional theme (respects dark/light mode)
        """
        super().__init__(theme)
        self.title = title
        self.show_notch = show_notch
        self.variant = variant

    def _get_device_color(self) -> RGBColor:
        """Get device frame color based on theme."""
        # Check if theme is dark mode
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                # If background is dark (low RGB values), use light frame
                if sum(bg) < 384:  # 128 * 3
                    return RGBColor(45, 45, 48)  # Space gray/graphite
                else:
                    return RGBColor(255, 255, 255)  # Silver/white

        # Default to space gray
        return RGBColor(45, 45, 48)

    def _get_screen_bg_color(self) -> RGBColor:
        """Get screen background color based on theme."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return RGBColor(bg[0], bg[1], bg[2])

        # Default to white
        return RGBColor(255, 255, 255)

    def _get_text_color(self) -> RGBColor:
        """Get text color based on theme."""
        if self.theme and isinstance(self.theme, dict):
            fg = self.theme.get("colors", {}).get("foreground", {}).get("DEFAULT")
            if fg and isinstance(fg, (list, tuple)) and len(fg) >= 3:
                return RGBColor(fg[0], fg[1], fg[2])

        # Default to black
        return RGBColor(0, 0, 0)

    def _is_dark_mode(self) -> bool:
        """Check if theme is dark mode."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return sum(bg) < 384  # 128 * 3
        return False

    def render(
        self, slide, left: float, top: float, width: float = 3.5, height: float = 6.5
    ) -> Dict[str, float]:
        """
        Render iPhone container to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Device width in inches (default 3.5)
            height: Device height in inches (default 6.5)

        Returns:
            Dict with content area bounds: {left, top, width, height}
        """
        shapes = []

        # Device frame (outer rounded rectangle)
        device_frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        device_frame.fill.solid()
        device_frame.fill.fore_color.rgb = self._get_device_color()
        device_frame.line.color.rgb = self._get_device_color()
        device_frame.line.width = Pt(2)
        shapes.append(device_frame)

        # Screen area (slightly inset)
        screen_margin = 0.08
        screen_left = left + screen_margin
        screen_top = top + screen_margin
        screen_width = width - (screen_margin * 2)
        screen_height = height - (screen_margin * 2)

        screen = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(screen_left),
            Inches(screen_top),
            Inches(screen_width),
            Inches(screen_height),
        )
        screen.fill.solid()
        screen.fill.fore_color.rgb = self._get_screen_bg_color()
        screen.line.fill.background()
        shapes.append(screen)

        # Notch (if enabled)
        notch_height = 0.15
        status_bar_top = screen_top

        if self.show_notch:
            notch_width = screen_width * 0.35
            notch_left = screen_left + (screen_width - notch_width) / 2

            notch = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(notch_left),
                Inches(screen_top),
                Inches(notch_width),
                Inches(notch_height),
            )
            notch.fill.solid()
            notch.fill.fore_color.rgb = self._get_device_color()
            notch.line.fill.background()
            shapes.append(notch)

            status_bar_top += notch_height

        # Status bar
        status_bar_height = 0.25
        status_bar = slide.shapes.add_textbox(
            Inches(screen_left + 0.1),
            Inches(status_bar_top),
            Inches(screen_width - 0.2),
            Inches(status_bar_height),
        )
        status_frame = status_bar.text_frame
        status_frame.text = "9:41"  # Classic iPhone time
        status_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        status_p = status_frame.paragraphs[0]
        status_p.alignment = PP_ALIGN.CENTER
        status_p.font.size = Pt(FONT_SIZES["xs"])
        status_p.font.bold = True
        status_p.font.color.rgb = self._get_text_color()
        shapes.append(status_bar)

        # Home indicator (bottom)
        home_indicator_height = 0.08
        home_indicator_width = screen_width * 0.3
        home_indicator_left = screen_left + (screen_width - home_indicator_width) / 2
        home_indicator_top = screen_top + screen_height - 0.15

        home_indicator = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(home_indicator_left),
            Inches(home_indicator_top),
            Inches(home_indicator_width),
            Inches(home_indicator_height),
        )
        home_indicator.fill.solid()

        # Home indicator color based on theme
        if self._is_dark_mode():
            home_indicator.fill.fore_color.rgb = RGBColor(120, 120, 128)
        else:
            home_indicator.fill.fore_color.rgb = RGBColor(60, 60, 67)

        home_indicator.line.fill.background()
        shapes.append(home_indicator)

        # Calculate content area (between status bar and home indicator)
        content_top = status_bar_top + status_bar_height + 0.05
        content_bottom = home_indicator_top - 0.05
        content_height = content_bottom - content_top

        # Add padding for realistic spacing
        padding_h = 0.15  # iPhone has tighter horizontal padding
        padding_v = 0.15

        content_left = screen_left + padding_h
        content_width = screen_width - (padding_h * 2)

        return {
            "left": content_left,
            "top": content_top + padding_v,
            "width": content_width,
            "height": content_height - (padding_v * 2),
        }


# TODO: Register component when registry is implemented
# Component metadata for documentation:
# iPhoneContainer - Device mockup with notch, status bar, home indicator
# Supports dark/light themes
