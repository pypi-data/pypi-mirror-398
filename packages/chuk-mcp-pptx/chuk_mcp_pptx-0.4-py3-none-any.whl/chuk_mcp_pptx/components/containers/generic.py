"""
Generic container components for PowerPoint presentations.

Provides simple, flexible containers for displaying chat conversations
and other content without specific device chrome.
"""

from typing import Optional, Dict, Any
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.platform_colors import get_container_ui_color
from ...constants import Platform, ColorKey, Theme


class ChatContainer(Component):
    """
    Generic chat container component.

    Creates a simple, clean container for chat conversations with:
    - Optional header with title
    - Optional border
    - Rounded corners
    - Theme-aware colors
    - Flexible sizing

    Perfect for when you want to show a conversation without specific
    device or platform chrome.

    Examples:
        # Simple container
        container = ChatContainer(
            title="Conversation",
            show_header=True,
            theme=theme
        )

        content_area = container.render(slide, left=1, top=1, width=6, height=5)

        # Render conversation inside
        conversation.render(
            slide,
            left=content_area['left'],
            top=content_area['top'],
            width=content_area['width']
        )
    """

    def __init__(
        self,
        title: Optional[str] = None,
        show_header: bool = False,
        show_border: bool = True,
        variant: str = "default",  # default, outlined, filled
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize chat container.

        Args:
            title: Optional title shown in header
            show_header: Whether to show header bar
            show_border: Whether to show border
            variant: Container variant (default, outlined, filled)
            theme: Optional theme (respects dark/light mode)
        """
        super().__init__(theme)
        self.title = title
        self.show_header = show_header
        self.show_border = show_border
        self.variant = variant

    def _is_dark_mode(self) -> bool:
        """Check if theme is dark mode."""
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return sum(bg) < 384
        return False

    def _get_container_bg_color(self) -> RGBColor:
        """Get container background color."""
        if self.variant == "filled":
            if self.theme and isinstance(self.theme, dict):
                card = self.theme.get("colors", {}).get("card", {}).get("DEFAULT")
                if card and isinstance(card, (list, tuple)) and len(card) >= 3:
                    return RGBColor(card[0], card[1], card[2])

        # Default to theme background
        if self.theme and isinstance(self.theme, dict):
            bg = self.theme.get("colors", {}).get("background", {}).get("DEFAULT")
            if bg and isinstance(bg, (list, tuple)) and len(bg) >= 3:
                return RGBColor(bg[0], bg[1], bg[2])

        hex_color = get_container_ui_color(Platform.GENERIC, ColorKey.CONTENT_BG, Theme.LIGHT)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_border_color(self) -> RGBColor:
        """Get border color."""
        if self.theme and isinstance(self.theme, dict):
            border = self.theme.get("colors", {}).get("border", {}).get("DEFAULT")
            if border and isinstance(border, (list, tuple)) and len(border) >= 3:
                return RGBColor(border[0], border[1], border[2])

        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.GENERIC, ColorKey.BORDER, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_header_bg_color(self) -> RGBColor:
        """Get header background color."""
        if self.theme and isinstance(self.theme, dict):
            muted = self.theme.get("colors", {}).get("muted", {}).get("DEFAULT")
            if muted and isinstance(muted, (list, tuple)) and len(muted) >= 3:
                return RGBColor(muted[0], muted[1], muted[2])

        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.GENERIC, ColorKey.HEADER_BG, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color."""
        if self.theme and isinstance(self.theme, dict):
            fg = self.theme.get("colors", {}).get("foreground", {}).get("DEFAULT")
            if fg and isinstance(fg, (list, tuple)) and len(fg) >= 3:
                return RGBColor(fg[0], fg[1], fg[2])

        theme_mode = Theme.DARK if self._is_dark_mode() else Theme.LIGHT
        hex_color = get_container_ui_color(Platform.GENERIC, ColorKey.HEADER_TEXT, theme_mode)
        return RGBColor(*self.hex_to_rgb(hex_color))

    def render(
        self, slide, left: float, top: float, width: float = 6.0, height: float = 5.0
    ) -> Dict[str, float]:
        """
        Render chat container to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Container width in inches
            height: Container height in inches

        Returns:
            Dict with content area bounds
        """
        shapes = []

        # Container frame
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
        )
        container.fill.solid()
        container.fill.fore_color.rgb = self._get_container_bg_color()

        # Border
        if self.show_border:
            container.line.color.rgb = self._get_border_color()
            container.line.width = Pt(1)
        else:
            container.line.fill.background()

        shapes.append(container)

        current_y = top
        header_height = 0.0

        # Header (if enabled)
        if self.show_header and self.title:
            header_height = 0.4

            header = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(header_height),
            )
            header.fill.solid()
            header.fill.fore_color.rgb = self._get_header_bg_color()
            header.line.fill.background()
            shapes.append(header)

            # Header text
            header_text = slide.shapes.add_textbox(
                Inches(left + 0.2),
                Inches(top + 0.05),
                Inches(width - 0.4),
                Inches(header_height - 0.1),
            )
            header_frame = header_text.text_frame
            header_frame.text = self.title
            header_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            header_p = header_frame.paragraphs[0]
            header_p.alignment = PP_ALIGN.CENTER
            header_p.font.size = Pt(14)
            header_p.font.bold = True
            header_p.font.color.rgb = self._get_text_color()
            shapes.append(header_text)

            current_y += header_height

        # Calculate content area with padding
        padding_h = 0.25  # Horizontal padding
        padding_v = 0.25  # Vertical padding

        content_top = current_y + padding_v
        content_height = height - (content_top - top) - padding_v

        return {
            "left": left + padding_h,
            "top": content_top,
            "width": width - (padding_h * 2),
            "height": content_height,
        }


# TODO: Register component when registry is implemented
# Future: Add social media containers (Twitter/X, Instagram, LinkedIn, YouTube)
