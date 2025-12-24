"""
ChatGPT chat components for PowerPoint presentations.

Provides realistic ChatGPT-style AI chat interfaces.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, FONT_FAMILIES
from ...tokens.platform_colors import CHAT_COLORS
from ...constants import Platform, ColorKey


class ChatGPTMessage(Component):
    """
    ChatGPT-style message component.

    Replicates the OpenAI ChatGPT interface with:
    - Clean, minimal design
    - User messages (right-aligned, subtle background)
    - Assistant messages (left-aligned with avatar)
    - Monospace code blocks support
    - Clean typography

    Examples:
        # User message
        msg = ChatGPTMessage(
            text="Explain recursion in simple terms",
            variant="user",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=8)

        # Assistant message
        msg = ChatGPTMessage(
            text="Recursion is when a function calls itself to solve a problem by breaking it into smaller parts.",
            variant="assistant",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=8)
    """

    def __init__(
        self, text: str, variant: str = "assistant", theme: Optional[Dict[str, Any]] = None
    ):
        """
        Initialize ChatGPT message.

        Args:
            text: Message text
            variant: Message variant (user, assistant, system)
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.variant = variant

    def _get_bg_color(self) -> Optional[RGBColor]:
        """Get background color."""
        if self.variant == ColorKey.USER:
            hex_color = CHAT_COLORS[Platform.CHATGPT][ColorKey.USER]
        elif self.variant == ColorKey.SYSTEM:
            hex_color = CHAT_COLORS[Platform.CHATGPT][ColorKey.SYSTEM]
        else:
            hex_color = CHAT_COLORS[Platform.CHATGPT][ColorKey.ASSISTANT]
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _get_text_color(self) -> RGBColor:
        """Get text color."""
        hex_color = CHAT_COLORS[Platform.CHATGPT][ColorKey.TEXT]
        return RGBColor(*self.hex_to_rgb(hex_color))

    def _calculate_content_height(self, width: float) -> float:
        """Estimate content height."""
        chars_per_line = int(width * 8)  # More conservative for better wrapping
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        line_height = 0.2
        padding = 0.4 if self.variant == "user" else 0.5
        return (lines * line_height) + padding

    def render(self, slide, left: float, top: float, width: float = 8.0) -> list:
        """Render ChatGPT message."""
        shapes = []

        content_height = self._calculate_content_height(width)
        avatar_size = 0.4

        # Create background container
        container = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(content_height),
        )

        container.fill.solid()
        container.fill.fore_color.rgb = self._get_bg_color()
        container.line.fill.background()

        shapes.append(container)

        # Avatar (for assistant messages)
        if self.variant == "assistant":
            avatar_left = left + 0.35
            avatar_top = top + 0.2

            avatar_shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(avatar_left),
                Inches(avatar_top),
                Inches(avatar_size),
                Inches(avatar_size),
            )

            # ChatGPT green
            avatar_shape.fill.solid()
            avatar_shape.fill.fore_color.rgb = RGBColor(
                *self.hex_to_rgb(CHAT_COLORS[Platform.CHATGPT][ColorKey.AVATAR])
            )
            avatar_shape.line.fill.background()

            # Add "AI" text in avatar
            av_text = avatar_shape.text_frame
            av_text.text = "AI"
            av_p = av_text.paragraphs[0]
            av_p.alignment = PP_ALIGN.CENTER
            av_text.vertical_anchor = MSO_ANCHOR.MIDDLE
            av_p.font.size = Pt(FONT_SIZES["sm"])
            av_p.font.bold = True
            av_p.font.color.rgb = self.get_color("primary.foreground")

            shapes.append(avatar_shape)

            # Text starts after avatar
            text_left = left + 0.35 + avatar_size + 0.2
            text_width = width - (0.35 + avatar_size + 0.4)
        elif self.variant == "user":
            # User messages: simple padding
            text_left = left + 0.5
            text_width = width - 1.0
        else:
            # System messages: centered
            text_left = left + 0.5
            text_width = width - 1.0

        # Message text
        text_box = slide.shapes.add_textbox(
            Inches(text_left), Inches(top + 0.2), Inches(text_width), Inches(content_height - 0.4)
        )

        text_frame = text_box.text_frame
        text_frame.text = self.text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_left = 0
        text_frame.margin_right = 0

        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.font.size = Pt(FONT_SIZES["base"])
        p.font.name = FONT_FAMILIES["sans"][0]
        p.font.color.rgb = self._get_text_color()
        p.line_spacing = 1.35

        shapes.append(text_box)

        return shapes


class ChatGPTConversation(Component):
    """
    ChatGPT-style conversation component.

    Displays a complete ChatGPT conversation with alternating user/assistant messages.

    Examples:
        messages = [
            {"text": "What is Python?", "variant": "user"},
            {"text": "Python is a high-level programming language...", "variant": "assistant"},
            {"text": "Can you show an example?", "variant": "user"},
            {"text": "Here's a simple Python example: print('Hello!')", "variant": "assistant"}
        ]

        conversation = ChatGPTConversation(messages, theme=theme)
        conversation.render(slide, left=0.5, top=1.5, width=9)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.0,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize ChatGPT conversation.

        Args:
            messages: List of message dicts with ChatGPTMessage params
            spacing: Spacing between messages (usually 0 for ChatGPT style)
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 8.0) -> list:
        """Render ChatGPT conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = ChatGPTMessage(
                text=msg_data.get("text", ""),
                variant=msg_data.get("variant", "assistant"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            msg_height = message._calculate_content_height(width)
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# ChatGPTMessage - Variants: user, assistant, system
# ChatGPTConversation - Full conversation flow
