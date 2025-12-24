"""
Slack chat components for PowerPoint presentations.

Provides realistic Slack-style chat interfaces for workplace communication demos.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES
from ...tokens.platform_colors import CHAT_COLORS
from ...constants import Platform, ColorKey


class SlackMessage(Component):
    """
    Slack-style message component.

    Replicates Slack's workplace messaging interface with:
    - Left-aligned messages with avatar and sender name
    - Timestamp on the same line as sender
    - Threaded conversation support
    - Reactions/emoji support
    - Clean, professional workspace design

    Examples:
        # Simple message
        msg = SlackMessage(
            text="Hey team, quick update on the project",
            sender="John Doe",
            timestamp="9:41 AM",
            avatar_text="JD",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Message with reaction
        msg = SlackMessage(
            text="Great work everyone! 🎉",
            sender="Sarah Smith",
            timestamp="9:45 AM",
            avatar_text="SS",
            reactions=["👍 3", "🎉 5"],
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(
        self,
        text: str,
        sender: str,
        timestamp: str,
        avatar_text: Optional[str] = None,
        reactions: Optional[List[str]] = None,
        thread_count: Optional[int] = None,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Slack message.

        Args:
            text: Message text
            sender: Sender name (required for Slack)
            timestamp: Message timestamp
            avatar_text: Avatar initials
            reactions: List of reactions (e.g., ["👍 3", "❤️ 2"])
            thread_count: Number of thread replies (optional)
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.sender = sender
        self.timestamp = timestamp
        self.avatar_text = avatar_text or sender[0]
        self.reactions = reactions or []
        self.thread_count = thread_count

    def _calculate_message_height(self, width: float) -> float:
        """Estimate message height."""
        # Account for sender name
        sender_height = 0.2

        # Account for timestamp
        timestamp_height = 0.2

        # Calculate text height
        text_width = width - 0.65  # Account for avatar and margins (like Teams)
        chars_per_line = int(text_width * 13)
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        text_height = lines * 0.19

        # Account for reactions
        reactions_height = 0.28 if self.reactions else 0

        # Account for thread indicator
        thread_height = 0.25 if self.thread_count else 0

        padding = 0.15
        return (
            sender_height
            + timestamp_height
            + text_height
            + reactions_height
            + thread_height
            + padding
        )

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render Slack message."""
        shapes = []

        avatar_size = 0.4  # Reduced to match Teams
        avatar_left = left
        content_left = left + avatar_size + 0.2  # More spacing like Teams
        content_width = width - avatar_size - 0.25

        # Avatar (circular like Teams)
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,  # Circular instead of rounded rectangle
            Inches(avatar_left),
            Inches(top),
            Inches(avatar_size),
            Inches(avatar_size),
        )
        avatar.fill.solid()
        avatar.fill.fore_color.rgb = RGBColor(
            *self.hex_to_rgb(CHAT_COLORS[Platform.SLACK][ColorKey.AVATAR])
        )
        avatar.line.fill.background()

        # Avatar text
        av_frame = avatar.text_frame
        av_frame.text = self.avatar_text
        av_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        av_p = av_frame.paragraphs[0]
        av_p.alignment = PP_ALIGN.CENTER
        av_p.font.size = Pt(FONT_SIZES["sm"])
        av_p.font.bold = True
        av_p.font.color.rgb = self.get_color("primary.foreground")
        shapes.append(avatar)

        current_top = top

        # Sender name (separate from timestamp like Teams)
        sender_box = slide.shapes.add_textbox(
            Inches(content_left), Inches(current_top), Inches(content_width), Inches(0.2)
        )
        sender_frame = sender_box.text_frame
        sender_frame.text = self.sender
        sender_frame.word_wrap = False
        sender_frame.vertical_anchor = MSO_ANCHOR.TOP
        sender_p = sender_frame.paragraphs[0]
        sender_p.alignment = PP_ALIGN.LEFT
        sender_p.font.size = Pt(FONT_SIZES["base"])
        sender_p.font.bold = True
        sender_p.font.color.rgb = RGBColor(
            *self.hex_to_rgb(CHAT_COLORS[Platform.SLACK][ColorKey.TEXT])
        )
        shapes.append(sender_box)
        current_top += 0.18

        # Timestamp (on separate line like Teams)
        timestamp_box = slide.shapes.add_textbox(
            Inches(content_left), Inches(current_top), Inches(content_width), Inches(0.18)
        )
        timestamp_frame = timestamp_box.text_frame
        timestamp_frame.text = self.timestamp
        timestamp_frame.word_wrap = False
        timestamp_p = timestamp_frame.paragraphs[0]
        timestamp_p.alignment = PP_ALIGN.LEFT
        timestamp_p.font.size = Pt(FONT_SIZES["sm"])
        timestamp_p.font.color.rgb = RGBColor(
            *self.hex_to_rgb(CHAT_COLORS[Platform.SLACK][ColorKey.SECONDARY_TEXT])
        )
        shapes.append(timestamp_box)
        current_top += 0.18

        # Message text
        text_box = slide.shapes.add_textbox(
            Inches(content_left), Inches(current_top), Inches(content_width), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = self.text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.LEFT
        text_p.font.size = Pt(FONT_SIZES["sm"])
        text_p.font.color.rgb = RGBColor(
            *self.hex_to_rgb(CHAT_COLORS[Platform.SLACK][ColorKey.TEXT])
        )
        shapes.append(text_box)

        # Calculate actual text height
        chars_per_line = int(content_width * 11)
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        text_height = lines * 0.18
        current_top += text_height + 0.05

        # Reactions
        if self.reactions:
            reactions_text = "   ".join(self.reactions)
            reactions_box = slide.shapes.add_textbox(
                Inches(content_left), Inches(current_top), Inches(content_width), Inches(0.25)
            )
            reactions_frame = reactions_box.text_frame
            reactions_frame.text = reactions_text
            reactions_p = reactions_frame.paragraphs[0]
            reactions_p.alignment = PP_ALIGN.LEFT
            reactions_p.font.size = Pt(FONT_SIZES["sm"])
            reactions_p.font.color.rgb = RGBColor(
                *self.hex_to_rgb(CHAT_COLORS[Platform.SLACK][ColorKey.SECONDARY_TEXT])
            )
            shapes.append(reactions_box)
            current_top += 0.25

        # Thread indicator
        if self.thread_count:
            thread_box = slide.shapes.add_textbox(
                Inches(content_left), Inches(current_top), Inches(content_width), Inches(0.2)
            )
            thread_frame = thread_box.text_frame
            thread_frame.text = (
                f"💬 {self.thread_count} {'reply' if self.thread_count == 1 else 'replies'}"
            )
            thread_p = thread_frame.paragraphs[0]
            thread_p.alignment = PP_ALIGN.LEFT
            thread_p.font.size = Pt(FONT_SIZES["sm"])
            thread_p.font.color.rgb = RGBColor(
                *self.hex_to_rgb(CHAT_COLORS[Platform.SLACK][ColorKey.LINK])
            )
            shapes.append(thread_box)

        return shapes


class SlackConversation(Component):
    """
    Slack-style conversation component.

    Displays a complete Slack channel conversation with automatic spacing.

    Examples:
        messages = [
            {
                "text": "Quick standup in 5 minutes!",
                "sender": "Sarah",
                "timestamp": "9:00 AM",
                "avatar_text": "SM"
            },
            {
                "text": "On my way",
                "sender": "John",
                "timestamp": "9:01 AM",
                "avatar_text": "JD",
                "reactions": ["👍 2"]
            }
        ]
        conversation = SlackConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=8)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.15,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize Slack conversation.

        Args:
            messages: List of message dicts with SlackMessage params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render Slack conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = SlackMessage(
                text=msg_data.get("text", ""),
                sender=msg_data.get("sender", "User"),
                timestamp=msg_data.get("timestamp", ""),
                avatar_text=msg_data.get("avatar_text"),
                reactions=msg_data.get("reactions"),
                thread_count=msg_data.get("thread_count"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            msg_height = message._calculate_message_height(width)
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# SlackMessage - Workplace messaging with reactions and threads
# SlackConversation - Full conversation flow
