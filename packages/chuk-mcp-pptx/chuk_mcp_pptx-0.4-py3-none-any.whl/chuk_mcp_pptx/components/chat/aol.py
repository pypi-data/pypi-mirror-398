"""
AOL Instant Messenger chat components for PowerPoint presentations.

Provides nostalgic AOL Instant Messenger (AIM)-style chat interfaces.
Recreates the classic 2000s messaging experience.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES, FONT_FAMILIES


class AIMBubble(Component):
    """
    AOL Instant Messenger-style chat bubble component.

    Recreates the classic AIM interface from the 2000s with:
    - No bubbles - just text with screen name
    - Different colors for sent vs received
    - Screen name followed by colon
    - Classic AIM yellow/blue color scheme
    - Simple, nostalgic design

    Examples:
        # Sent message
        msg = AIMBubble(
            text="Hey! What's up?",
            screen_name="xXCoolDude2003Xx",
            variant="sent",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=7)

        # Received message
        msg = AIMBubble(
            text="Not much, you?",
            screen_name="sk8rgrl2004",
            variant="received",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=7)
    """

    def __init__(
        self,
        text: str,
        screen_name: str,
        variant: str = "received",
        timestamp: Optional[str] = None,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize AIM bubble.

        Args:
            text: Message text
            screen_name: AIM screen name
            variant: Message variant (sent, received)
            timestamp: Optional timestamp
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.screen_name = screen_name
        self.variant = variant
        self.timestamp = timestamp

    def _get_screen_name_color(self) -> RGBColor:
        """Get screen name color."""
        if self.variant == "sent":
            # Classic AIM red for sent
            return self.get_color("destructive.DEFAULT")
        else:
            # Classic AIM blue for received
            return self.get_color("primary.DEFAULT")

    def _get_text_color(self) -> RGBColor:
        """Get message text color."""
        # Classic black text
        return self.get_color("foreground.DEFAULT")

    def _calculate_message_height(self, width: float) -> float:
        """Estimate message height."""
        # Screen name line
        header_height = 0.22  # Increased from 0.18

        # Calculate text height
        chars_per_line = int(width * 11)  # More generous
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        text_height = lines * 0.19  # Increased from 0.16

        # Timestamp
        timestamp_height = 0.18 if self.timestamp else 0  # Increased from 0.15

        return header_height + text_height + timestamp_height + 0.08  # More padding

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render AIM message."""
        shapes = []
        current_top = top

        # Screen name (e.g., "xXCoolDude2003Xx:")
        sn_box = slide.shapes.add_textbox(
            Inches(left),
            Inches(current_top),
            Inches(width),
            Inches(0.22),  # Increased from 0.18
        )
        sn_frame = sn_box.text_frame
        sn_frame.text = f"{self.screen_name}:"
        sn_frame.word_wrap = False
        sn_p = sn_frame.paragraphs[0]
        sn_p.alignment = PP_ALIGN.LEFT
        sn_p.font.size = Pt(FONT_SIZES["sm"])  # Larger
        sn_p.font.bold = True
        sn_p.font.name = FONT_FAMILIES["sans"][0]  # Classic AIM font
        sn_p.font.color.rgb = self._get_screen_name_color()
        shapes.append(sn_box)
        current_top += 0.20  # More spacing

        # Message text
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(current_top), Inches(width), Inches(0.5)
        )
        text_frame = text_box.text_frame
        text_frame.text = self.text
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_p = text_frame.paragraphs[0]
        text_p.alignment = PP_ALIGN.LEFT
        text_p.font.size = Pt(FONT_SIZES["sm"])  # Larger
        text_p.font.name = FONT_FAMILIES["sans"][0]
        text_p.font.color.rgb = self._get_text_color()
        text_p.line_spacing = 1.3  # Better line spacing
        shapes.append(text_box)

        # Calculate text height
        chars_per_line = int(width * 11)
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)
        text_height = lines * 0.19
        current_top += text_height

        # Timestamp (small, gray)
        if self.timestamp:
            ts_box = slide.shapes.add_textbox(
                Inches(left), Inches(current_top + 0.02), Inches(width), Inches(0.15)
            )
            ts_frame = ts_box.text_frame
            ts_frame.text = self.timestamp
            ts_p = ts_frame.paragraphs[0]
            ts_p.alignment = PP_ALIGN.LEFT
            ts_p.font.size = Pt(FONT_SIZES["xs"])
            ts_p.font.color.rgb = self.get_color("muted.foreground")  # Gray
            shapes.append(ts_box)

        return shapes


class AIMConversation(Component):
    """
    AOL Instant Messenger-style conversation component.

    Displays a complete AIM conversation with classic styling.

    Examples:
        messages = [
            {
                "text": "Hey! Want to hang out later?",
                "screen_name": "sk8rgrl2004",
                "variant": "received",
                "timestamp": "5:30 PM"
            },
            {
                "text": "Yeah! Let's go to the mall",
                "screen_name": "xXCoolDude2003Xx",
                "variant": "sent",
                "timestamp": "5:31 PM"
            }
        ]
        conversation = AIMConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=7)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.15,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize AIM conversation.

        Args:
            messages: List of message dicts with AIMBubble params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 7.0) -> list:
        """Render AIM conversation."""
        shapes = []
        current_top = top

        for msg_data in self.messages:
            message = AIMBubble(
                text=msg_data.get("text", ""),
                screen_name=msg_data.get("screen_name", "User"),
                variant=msg_data.get("variant", "received"),
                timestamp=msg_data.get("timestamp"),
                theme=self.theme,
            )

            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height
            msg_height = message._calculate_message_height(width)
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register components when registry is implemented
# Component metadata for documentation:
# AIMBubble - Classic AOL Instant Messenger style (no bubbles, just text)
# AIMConversation - Full conversation flow
