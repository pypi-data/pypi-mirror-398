"""
Chat message components for PowerPoint presentations.

Provides chat/messaging UI components for showing conversations,
customer support examples, messaging apps, and communication flows.
"""

from typing import Optional, Dict, Any, List
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from ..base import Component
from ...tokens.typography import FONT_SIZES
from ...constants import ComponentSizing
from ..core.avatar import Avatar
from ...tokens.typography import get_text_style


class ChatMessage(Component):
    """
    Chat message component - single message bubble.

    Perfect for showing conversations, chat interfaces, customer support
    examples, and messaging flows in presentations.

    Variants:
        - sent: Message sent by user (aligned right, primary color)
        - received: Message received from other (aligned left, muted color)
        - system: System message (centered, minimal styling)

    Examples:
        # Simple received message
        msg = ChatMessage(
            text="Hello! How can I help you?",
            variant="received",
            theme=theme
        )
        msg.render(slide, left=1, top=2, width=6)

        # Sent message with sender info
        msg = ChatMessage(
            text="I need help with my account",
            sender="John Doe",
            timestamp="10:30 AM",
            variant="sent",
            theme=theme
        )
        msg.render(slide, left=1, top=3, width=6)

        # Message with avatar
        msg = ChatMessage(
            text="Sure, I can help with that!",
            sender="Support Agent",
            avatar_text="SA",
            variant="received",
            show_avatar=True,
            theme=theme
        )
        msg.render(slide, left=1, top=4, width=6)
    """

    def __init__(
        self,
        text: str,
        sender: Optional[str] = None,
        timestamp: Optional[str] = None,
        avatar_text: Optional[str] = None,
        avatar_icon: Optional[str] = None,
        variant: str = "received",
        show_avatar: bool = False,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize chat message.

        Args:
            text: Message text content
            sender: Sender name (optional)
            timestamp: Message timestamp (optional)
            avatar_text: Avatar initials (optional)
            avatar_icon: Avatar icon (optional)
            variant: Message variant (sent, received, system)
            show_avatar: Whether to show avatar
            theme: Optional theme
        """
        super().__init__(theme)
        self.text = text
        self.sender = sender
        self.timestamp = timestamp
        self.avatar_text = avatar_text
        self.avatar_icon = avatar_icon
        self.variant = variant
        self.show_avatar = show_avatar

    def _get_bubble_color(self) -> RGBColor:
        """Get message bubble background color."""
        if self.variant == "sent":
            return self.get_color("primary.DEFAULT")
        elif self.variant == "system":
            return self.get_color("muted.DEFAULT")
        else:  # received
            return self.get_color("card.DEFAULT")

    def _get_text_color(self) -> RGBColor:
        """Get message text color."""
        if self.variant == "sent":
            return self.get_color("primary.foreground")
        else:
            return self.get_color("foreground.DEFAULT")

    def _get_meta_color(self) -> RGBColor:
        """Get metadata (sender/timestamp) color."""
        if self.variant == "sent":
            return self.get_color("primary.foreground")
        else:
            return self.get_color("muted.foreground")

    def _calculate_bubble_height(self, width: float) -> float:
        """Estimate bubble height based on text length."""
        # Rough estimate: ~70 characters per line for typical widths
        chars_per_line = int(width * 11)  # Approximate characters per inch
        lines = max(1, (len(self.text) + chars_per_line - 1) // chars_per_line)

        # Base height for text
        line_height = 0.2  # ~14pt line height in inches
        text_height = lines * line_height

        # Add padding and metadata
        padding = 0.3  # Top and bottom padding
        meta_height = 0.2 if (self.sender or self.timestamp) else 0

        return text_height + padding + meta_height

    def render(self, slide, left: float, top: float, width: float = 4.0) -> list:
        """
        Render chat message to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Maximum message width in inches

        Returns:
            List of rendered shapes [bubble, avatar if shown]
        """
        shapes = []

        # Avatar dimensions
        avatar_size = 0.6
        avatar_gap = 0.15

        # Calculate positions based on variant
        if self.variant == "sent":
            # Align right
            bubble_width = min(width * 0.7, width - 0.5)
            bubble_left = left + width - bubble_width
            show_avatar_right = self.show_avatar
            show_avatar_left = False
        elif self.variant == "system":
            # Center
            bubble_width = min(width * 0.8, width - 0.5)
            bubble_left = left + (width - bubble_width) / 2
            show_avatar_right = False
            show_avatar_left = False
        else:  # received
            # Align left
            bubble_width = min(width * 0.7, width - 0.5)
            bubble_left = left
            if self.show_avatar:
                bubble_left += avatar_size + avatar_gap
                bubble_width -= avatar_size + avatar_gap
            show_avatar_right = False
            show_avatar_left = self.show_avatar

        # Calculate bubble height
        bubble_height = self._calculate_bubble_height(bubble_width)

        # Render avatar (left side for received)
        if show_avatar_left:
            avatar = Avatar(
                text=self.avatar_text,
                icon=self.avatar_icon or "user",
                variant="filled",
                size="sm",
                color_variant="default",
                theme=self.theme,
            )
            avatar_shape = avatar.render(slide, left=left, top=top, diameter=avatar_size)
            shapes.append(avatar_shape)

        # Create message bubble
        bubble = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bubble_left),
            Inches(top),
            Inches(bubble_width),
            Inches(bubble_height),
        )

        # Style bubble
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = self._get_bubble_color()

        # Border for received messages
        if self.variant == "received":
            bubble.line.color.rgb = self.get_color("border.DEFAULT")
            bubble.line.width = Pt(ComponentSizing.BORDER_WIDTH_THIN)
        else:
            bubble.line.fill.background()

        # Setup text frame
        text_frame = bubble.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        padding = 0.15
        text_frame.margin_left = Inches(padding)
        text_frame.margin_right = Inches(padding)
        text_frame.margin_top = Inches(padding)
        text_frame.margin_bottom = Inches(padding)

        # Sender name (if provided)
        current_p = text_frame.paragraphs[0]
        if self.sender:
            current_p.text = self.sender
            if self.variant == "system":
                current_p.alignment = PP_ALIGN.CENTER
            else:
                # Both sent and received should be left-aligned within bubble
                current_p.alignment = PP_ALIGN.LEFT

            current_p.font.size = Pt(FONT_SIZES["xs"])
            current_p.font.bold = True
            current_p.font.color.rgb = self._get_meta_color()

            # Add new paragraph for message text
            current_p = text_frame.add_paragraph()
            current_p.space_before = Pt(ComponentSizing.SPACE_MD)

        # Message text
        current_p.text = self.text
        if self.variant == "system":
            current_p.alignment = PP_ALIGN.CENTER
        else:
            # All message text within bubbles should be left-aligned
            current_p.alignment = PP_ALIGN.LEFT

        style = get_text_style("body")
        current_p.font.size = Pt(style["font_size"])
        current_p.font.color.rgb = self._get_text_color()

        # Timestamp (if provided)
        if self.timestamp:
            p = text_frame.add_paragraph()
            p.text = self.timestamp
            p.space_before = Pt(ComponentSizing.SPACE_LG)
            if self.variant == "system":
                p.alignment = PP_ALIGN.CENTER
            elif self.variant == "sent":
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.RIGHT  # Timestamp typically right-aligned

            p.font.size = Pt(FONT_SIZES["xs"])
            p.font.color.rgb = self._get_meta_color()

        shapes.append(bubble)

        # Render avatar (right side for sent)
        if show_avatar_right:
            avatar = Avatar(
                text=self.avatar_text,
                icon=self.avatar_icon or "user",
                variant="filled",
                size="sm",
                color_variant="primary",
                theme=self.theme,
            )
            avatar_left = bubble_left + bubble_width + avatar_gap
            avatar_shape = avatar.render(slide, left=avatar_left, top=top, diameter=avatar_size)
            shapes.append(avatar_shape)

        return shapes


class ChatConversation(Component):
    """
    Chat conversation component - series of chat messages.

    Displays a complete conversation thread with multiple messages.
    Perfect for showing customer support flows, chat examples, and
    messaging scenarios.

    Examples:
        # Simple conversation
        messages = [
            {
                "text": "Hi! I need help with my order",
                "variant": "sent",
                "timestamp": "10:30 AM"
            },
            {
                "text": "Hello! I'd be happy to help. What's your order number?",
                "sender": "Support",
                "avatar_text": "SA",
                "variant": "received",
                "timestamp": "10:31 AM"
            },
            {
                "text": "It's #12345",
                "variant": "sent",
                "timestamp": "10:32 AM"
            }
        ]

        conversation = ChatConversation(messages, theme=theme)
        conversation.render(slide, left=1, top=2, width=7)

        # Conversation with avatars
        messages = [
            {
                "text": "Welcome to our support chat!",
                "sender": "Bot",
                "avatar_icon": "user",
                "variant": "system"
            },
            {
                "text": "Hello!",
                "sender": "John",
                "avatar_text": "JD",
                "variant": "sent",
                "show_avatar": True
            },
            {
                "text": "How can I assist you today?",
                "sender": "Sarah",
                "avatar_text": "SM",
                "variant": "received",
                "show_avatar": True
            }
        ]

        conversation = ChatConversation(messages, spacing=0.3, theme=theme)
        conversation.render(slide, left=1, top=2, width=7)
    """

    def __init__(
        self,
        messages: List[Dict[str, Any]],
        spacing: float = 0.25,
        theme: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize chat conversation.

        Args:
            messages: List of message dicts with ChatMessage params
            spacing: Spacing between messages in inches
            theme: Optional theme
        """
        super().__init__(theme)
        self.messages = messages
        self.spacing = spacing

    def render(self, slide, left: float, top: float, width: float = 6.0) -> list:
        """
        Render chat conversation to slide.

        Args:
            slide: PowerPoint slide
            left: Left position in inches
            top: Top position in inches
            width: Conversation width in inches

        Returns:
            List of all rendered shapes
        """
        shapes = []
        current_top = top

        for msg_data in self.messages:
            # Create message
            message = ChatMessage(
                text=msg_data.get("text", ""),
                sender=msg_data.get("sender"),
                timestamp=msg_data.get("timestamp"),
                avatar_text=msg_data.get("avatar_text"),
                avatar_icon=msg_data.get("avatar_icon"),
                variant=msg_data.get("variant", "received"),
                show_avatar=msg_data.get("show_avatar", False),
                theme=self.theme,
            )

            # Render message
            msg_shapes = message.render(slide, left, current_top, width)
            shapes.extend(msg_shapes)

            # Calculate height of rendered message (approximate)
            bubble_width = width * 0.7 if msg_data.get("variant") != "system" else width * 0.8
            msg_height = message._calculate_bubble_height(bubble_width)

            # Update top position for next message
            current_top += msg_height + self.spacing

        return shapes


# TODO: Register component when registry is implemented
# Component metadata for documentation:
# ChatMessage - Variants: sent, received, system
# ChatMessage - Props: text, sender, timestamp, avatar_text, avatar_icon, show_avatar
# ChatConversation - Props: messages (list of ChatMessage configs), spacing
