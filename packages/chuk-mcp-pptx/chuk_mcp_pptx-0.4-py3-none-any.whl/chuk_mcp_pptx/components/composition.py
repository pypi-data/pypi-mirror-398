"""
Composition patterns for PowerPoint components.
Inspired by shadcn/ui's compositional API (Card.Header, Card.Content, etc.)
"""

from __future__ import annotations


from typing import Any, TypeVar, Optional
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from abc import ABC, abstractmethod

from .base import Component


T = TypeVar("T", bound="ComposableComponent")


class ComposableComponent(Component):
    """
    Base class for components that support composition.
    Allows nested subcomponents like Card.Header, Card.Content, etc.
    """

    def __init__(self, theme: Optional[dict[str, Any | None]] = None):
        super().__init__(theme)
        self._children: list["SubComponent"] = []

    def add_child(self, child: "SubComponent"):
        """Add a child subcomponent."""
        child.parent = self
        child.theme = self.theme  # Inherit theme
        self._children.append(child)
        return self

    def get_children(self) -> list["SubComponent"]:
        """Get all child components."""
        return self._children

    def clear_children(self):
        """Remove all children."""
        self._children = []


class SubComponent(Component, ABC):
    """
    Base class for subcomponents (like Card.Header).
    Must be used within a parent ComposableComponent.
    """

    def __init__(self, theme: Optional[dict[str, Any | None]] = None):
        super().__init__(theme)
        self.parent: ComposableComponent | None = None

    @abstractmethod
    def render_into(self, parent_shape, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """
        Render this subcomponent into the parent's shape.

        Args:
            parent_shape: The parent shape to render into
            theme: Optional theme override

        Returns:
            The rendered element
        """
        pass


# Card composition components


class CardHeader(SubComponent):
    """Card header subcomponent."""

    def __init__(
        self,
        title: str,
        description: str | None = None,
        theme: Optional[dict[str, Any | None]] = None,
    ):
        super().__init__(theme)
        self.title = title
        self.description = description

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render header into text frame."""
        theme = theme or self.theme

        # Clear existing paragraphs if this is the first child
        if not text_frame.paragraphs[0].text:
            text_frame.paragraphs[0].text = ""

        # Title
        p = (
            text_frame.paragraphs[0]
            if not text_frame.paragraphs[0].text
            else text_frame.add_paragraph()
        )
        p.text = self.title
        style = self.get_text_style("h4")
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])
        p.font.bold = True
        p.font.color.rgb = self.get_color("card.foreground")

        # Description
        if self.description:
            p = text_frame.add_paragraph()
            p.text = self.description
            p.space_before = Pt(4)
            style = self.get_text_style("body-sm")
            p.font.name = style["font_family"]
            p.font.size = Pt(style["font_size"])
            p.font.color.rgb = self.get_color("muted.foreground")

        return p


class CardContent(SubComponent):
    """Card content subcomponent."""

    def __init__(self, content: str, theme: Optional[dict[str, Any | None]] = None):
        super().__init__(theme)
        self.content = content

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render content into text frame."""
        theme = theme or self.theme

        p = text_frame.add_paragraph()
        p.text = self.content
        p.space_before = Pt(8)
        style = self.get_text_style("body")
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])
        p.font.color.rgb = self.get_color("foreground.DEFAULT")

        return p


class CardFooter(SubComponent):
    """Card footer subcomponent."""

    def __init__(
        self, text: str, align: str = "left", theme: Optional[dict[str, Any | None]] = None
    ):
        super().__init__(theme)
        self.text = text
        self.align = align

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render footer into text frame."""
        theme = theme or self.theme

        p = text_frame.add_paragraph()
        p.text = self.text
        p.space_before = Pt(12)
        style = self.get_text_style("caption")
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])
        p.font.color.rgb = self.get_color("muted.foreground")

        if self.align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif self.align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        return p


class CardTitle(SubComponent):
    """Standalone card title subcomponent."""

    def __init__(self, text: str, theme: Optional[dict[str, Any | None]] = None):
        super().__init__(theme)
        self.text = text

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render title into text frame."""
        theme = theme or self.theme

        p = (
            text_frame.paragraphs[0]
            if not text_frame.paragraphs[0].text
            else text_frame.add_paragraph()
        )
        p.text = self.text
        p.alignment = PP_ALIGN.CENTER
        style = self.get_text_style("h5")  # Use h5 for smaller, more compact titles
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])
        p.font.bold = True
        p.font.color.rgb = self.get_color("card.foreground")

        return p


class CardDescription(SubComponent):
    """Standalone card description subcomponent."""

    def __init__(self, text: str, theme: Optional[dict[str, Any | None]] = None):
        super().__init__(theme)
        self.text = text

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render description into text frame."""
        theme = theme or self.theme

        p = text_frame.add_paragraph()
        p.text = self.text
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)
        style = self.get_text_style("body")
        p.font.name = style["font_family"]
        p.font.size = Pt(style["font_size"])
        p.font.color.rgb = self.get_color("muted.foreground")

        return p


# Layout composition components


class Stack(SubComponent):
    """Vertical stack layout."""

    def __init__(
        self,
        children: list[SubComponent],
        spacing: float = 0.1,
        theme: Optional[dict[str, Any | None]] = None,
    ):
        super().__init__(theme)
        self.children = children
        self.spacing = spacing

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render stacked children."""
        theme = theme or self.theme

        for i, child in enumerate(self.children):
            if i > 0:
                # Add spacing between items
                p = text_frame.add_paragraph()
                p.space_before = Pt(self.spacing * 72)  # Convert inches to points

            child.render_into(text_frame, theme)

        return text_frame


class Separator(SubComponent):
    """Visual separator line."""

    def __init__(self, theme: Optional[dict[str, Any | None]] = None):
        super().__init__(theme)

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render separator as text."""
        theme = theme or self.theme

        p = text_frame.add_paragraph()
        p.text = "─" * 50  # Horizontal line using box drawing character
        p.space_before = Pt(8)
        p.space_after = Pt(8)
        p.font.size = Pt(8)
        p.font.color.rgb = self.get_color("border.DEFAULT")

        return p


class Badge(SubComponent):
    """Badge/label component."""

    def __init__(
        self, text: str, variant: str = "default", theme: Optional[dict[str, Any | None]] = None
    ):
        super().__init__(theme)
        self.text = text
        self.variant = variant

    def render_into(self, text_frame, theme: Optional[dict[str, Any | None]] = None) -> Any:
        """Render badge inline."""
        theme = theme or self.theme

        p = text_frame.add_paragraph()
        p.text = f"[{self.text}]"
        p.font.size = Pt(10)
        p.font.bold = True

        # Color based on variant
        color_map = {
            "default": "primary.DEFAULT",
            "success": "success.DEFAULT",
            "warning": "warning.DEFAULT",
            "destructive": "destructive.DEFAULT",
            "secondary": "secondary.DEFAULT",
        }
        p.font.color.rgb = self.get_color(color_map.get(self.variant, "primary.DEFAULT"))

        return p


# Composition helper functions


def compose(*children: SubComponent) -> list[SubComponent]:
    """
    Helper function to compose multiple subcomponents.

    Example:
        card = ComposableCard()
        card.children = compose(
            CardHeader("Title", "Description"),
            CardContent("Main content here"),
            CardFooter("Footer text")
        )
    """
    return list(children)


def with_separator(*children: SubComponent) -> list[SubComponent]:
    """
    Compose children with separators between them.

    Example:
        compose(
            CardHeader("Section 1"),
            Separator(),
            CardContent("Content"),
            Separator(),
            CardFooter("Footer")
        )
    """
    result = []
    for i, child in enumerate(children):
        if i > 0:
            result.append(Separator())
        result.append(child)
    return result


class CompositionBuilder:
    """
    Fluent builder for composing components.

    Example:
        composition = (CompositionBuilder()
            .header("Title", "Subtitle")
            .content("Main content")
            .footer("Footer text")
            .build())
    """

    def __init__(self, theme: Optional[dict[str, Any | None]] = None):
        self.theme = theme
        self._children: list[SubComponent] = []

    def header(self, title: str, description: str | None = None) -> "CompositionBuilder":
        """Add header."""
        self._children.append(CardHeader(title, description, self.theme))
        return self

    def title(self, text: str) -> "CompositionBuilder":
        """Add title."""
        self._children.append(CardTitle(text, self.theme))
        return self

    def description(self, text: str) -> "CompositionBuilder":
        """Add description."""
        self._children.append(CardDescription(text, self.theme))
        return self

    def content(self, text: str) -> "CompositionBuilder":
        """Add content."""
        self._children.append(CardContent(text, self.theme))
        return self

    def footer(self, text: str, align: str = "left") -> "CompositionBuilder":
        """Add footer."""
        self._children.append(CardFooter(text, align, self.theme))
        return self

    def badge(self, text: str, variant: str = "default") -> "CompositionBuilder":
        """Add badge."""
        self._children.append(Badge(text, variant, self.theme))
        return self

    def separator(self) -> "CompositionBuilder":
        """Add separator."""
        self._children.append(Separator(self.theme))
        return self

    def custom(self, component: SubComponent) -> "CompositionBuilder":
        """Add custom component."""
        self._children.append(component)
        return self

    def build(self) -> list[SubComponent]:
        """Build the composition."""
        return self._children
