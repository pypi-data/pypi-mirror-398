"""
Presentation Manager for PowerPoint MCP Server

Manages PowerPoint presentations with support for chuk-artifacts integration.
Each presentation is stored as a BLOB namespace for persistence and multi-server access.

Uses chuk-mcp-server's built-in artifact store context for storage.
Uses Pydantic models throughout for type safety and validation.

MULTI-INSTANCE SUPPORT:
- Uses artifact store as the source of truth for all presentation data
- Instance-level dictionaries are used only as a short-lived cache (TTL: 60s)
- All operations query the artifact store to find presentations by session
- No shared state required between server instances
"""

from __future__ import annotations

import asyncio
import base64
import io
import logging
import time
from datetime import datetime
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

from ..models import (
    PresentationMetadata,
    SlideMetadata,
    PresentationInfo,
    ListPresentationsResponse,
)

logger = logging.getLogger(__name__)


class PresentationManager:
    """
    Manages PowerPoint presentations with chuk-artifacts integration.

    Uses chuk-mcp-server's built-in artifact store context for flexible storage
    (memory, filesystem, sqlite, s3). Each presentation is stored as a BLOB
    namespace with automatic session management.
    Presentations and metadata are Pydantic models for type safety.

    MULTI-INSTANCE SAFE:
    - Artifact store is the source of truth
    - Local caches have 60s TTL and are refreshed from artifact store
    - Supports multiple concurrent server instances accessing same session
    """

    # MIME type for PowerPoint presentations
    PPTX_MIME_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    # Cache TTL in seconds (how long to trust local cache before refreshing)
    CACHE_TTL = 60

    def __init__(self, base_path: str = "presentations") -> None:
        """
        Initialize the presentation manager.

        Args:
            base_path: Base path prefix for presentation names
        """
        self.base_path = base_path
        # Local caches with TTL - NOT the source of truth
        self._presentations: dict[str, PresentationType] = {}
        self._metadata: dict[str, PresentationMetadata] = {}
        self._namespace_ids: dict[str, str] = {}  # name -> namespace_id mapping
        self._cache_timestamps: dict[str, float] = {}  # name -> last_loaded_time
        self._current_presentation: str | None = None
        logger.info(
            f"PresentationManager initialized (multi-instance safe), base path: {base_path}"
        )

    def _get_store(self):
        """Get the artifact store from context."""
        from chuk_mcp_server import get_artifact_store, has_artifact_store

        if has_artifact_store():
            return get_artifact_store()
        return None

    def _sanitize_name(self, name: str) -> str:
        """Sanitize presentation name to prevent directory traversal."""
        safe_name = "".join(c for c in name if c.isalnum() or c in ("-", "_"))
        if not safe_name:
            safe_name = "presentation"
        return safe_name

    def _is_cache_valid(self, name: str) -> bool:
        """Check if cached presentation is still valid (within TTL)."""
        if name not in self._cache_timestamps:
            return False
        age = time.time() - self._cache_timestamps[name]
        return age < self.CACHE_TTL

    def _update_cache_timestamp(self, name: str) -> None:
        """Update the cache timestamp for a presentation."""
        self._cache_timestamps[name] = time.time()

    async def _find_namespace_in_store(self, name: str) -> str | None:
        """
        Find a presentation's namespace_id in the artifact store by querying the session.

        This is the key method for multi-instance support - it always queries the
        artifact store to find presentations, regardless of local cache state.

        Args:
            name: Presentation name to search for

        Returns:
            namespace_id if found, None otherwise
        """
        store = self._get_store()
        if not store:
            return None

        try:
            safe_name = self._sanitize_name(name)
            expected_name = f"{self.base_path}/{safe_name}"

            # Always query the artifact store for the current session state
            namespaces = await store.list_namespaces()
            logger.debug(f"Searching for '{expected_name}' in {len(namespaces)} namespaces")

            for ns_info in namespaces:
                if ns_info.name == expected_name:
                    logger.debug(
                        f"Found namespace {ns_info.namespace_id} for presentation '{name}'"
                    )
                    return ns_info.namespace_id

            logger.debug(f"Presentation '{name}' not found in artifact store")
            return None
        except Exception as e:
            logger.error(f"Failed to search for presentation in artifact store: {e}")
            return None

    def get_namespace_id(self, name: str) -> str | None:
        """
        Get the namespace ID for a presentation by name (from cache only).

        Note: This returns cached value. For multi-instance safety, use
        _find_namespace_in_store() to query the artifact store directly.
        """
        return self._namespace_ids.get(name)

    def get_artifact_uri(self, name: str) -> str | None:
        """
        Get the artifact URI for a presentation.

        Args:
            name: Presentation name

        Returns:
            Artifact URI string or None if not found
        """
        namespace_id = self._namespace_ids.get(name)
        if namespace_id:
            return f"artifact://chuk-mcp-pptx/{self.base_path}/{name}"
        return None

    async def _save_to_store(self, name: str, prs: PresentationType) -> bool:
        """
        Save presentation to artifact store.

        Args:
            name: Presentation name
            prs: Presentation object

        Returns:
            True if successful, False otherwise
        """
        store = self._get_store()
        if not store:
            logger.debug("No artifact store available, skipping persistence")
            return False

        from chuk_mcp_server import NamespaceType, StorageScope

        try:
            # Convert presentation to bytes (wrap blocking I/O)
            buffer = io.BytesIO()
            await asyncio.to_thread(prs.save, buffer)
            buffer.seek(0)
            data = buffer.read()

            # Check if namespace already exists
            namespace_id = self._namespace_ids.get(name)

            if namespace_id:
                # Update existing namespace
                await store.write_namespace(namespace_id, data=data)
                logger.info(f"Updated presentation in artifact store: {name} ({namespace_id})")
            else:
                # Create new BLOB namespace
                safe_name = self._sanitize_name(name)
                namespace_info = await store.create_namespace(
                    type=NamespaceType.BLOB,
                    scope=StorageScope.SESSION,
                    name=f"{self.base_path}/{safe_name}",
                    metadata={
                        "mime_type": self.PPTX_MIME_TYPE,
                        "presentation_name": name,
                        "file_extension": ".pptx",
                    },
                )
                self._namespace_ids[name] = namespace_info.namespace_id

                # Write the presentation data
                await store.write_namespace(namespace_info.namespace_id, data=data)
                logger.info(
                    f"Saved presentation to artifact store: {name} ({namespace_info.namespace_id})"
                )

            return True
        except Exception as e:
            logger.error(f"Failed to save to artifact store: {e}")
            return False

    async def _load_from_store(
        self, name: str, force_refresh: bool = False
    ) -> PresentationType | None:
        """
        Load presentation from artifact store by name.

        Uses _find_namespace_in_store() to query the session for the presentation.
        This ensures multi-instance compatibility by always checking the artifact store.

        Args:
            name: Presentation name
            force_refresh: If True, bypass cache even if valid

        Returns:
            Presentation object or None if not found
        """
        store = self._get_store()
        if not store:
            logger.debug("No artifact store available")
            return None

        try:
            # Check cache first (if not forcing refresh)
            if not force_refresh and self._is_cache_valid(name) and name in self._presentations:
                logger.debug(
                    f"Using cached presentation '{name}' (age: {time.time() - self._cache_timestamps[name]:.1f}s)"
                )
                return self._presentations[name]

            # Find namespace in artifact store (multi-instance safe)
            namespace_id = await self._find_namespace_in_store(name)
            if not namespace_id:
                logger.debug(f"Presentation '{name}' not found in artifact store")
                return None

            # Cache the namespace_id
            self._namespace_ids[name] = namespace_id

            # Read from artifact store
            data = await store.read_namespace(namespace_id)
            if data is None:
                logger.debug(f"No data found for namespace: {namespace_id}")
                return None

            buffer = io.BytesIO(data)
            prs = Presentation(buffer)

            # Update cache
            self._presentations[name] = prs
            self._update_cache_timestamp(name)

            logger.info(f"Loaded presentation '{name}' from artifact store ({namespace_id})")
            return prs
        except Exception as e:
            logger.error(f"Failed to load presentation '{name}' from artifact store: {e}")
            return None

    async def _delete_from_store(self, name: str) -> bool:
        """
        Delete presentation from artifact store.

        Args:
            name: Presentation name

        Returns:
            True if successful, False otherwise
        """
        store = self._get_store()
        if not store:
            logger.debug("No artifact store available")
            return False

        try:
            namespace_id = self._namespace_ids.get(name)
            if not namespace_id:
                logger.warning(f"Presentation not found in namespace mapping for deletion: {name}")
                return False

            await store.destroy_namespace(namespace_id)
            del self._namespace_ids[name]
            logger.info(f"Deleted presentation from artifact store: {name} ({namespace_id})")
            return True
        except Exception as e:
            logger.error(f"Failed to delete from artifact store: {e}")  # nosec B608
            return False

    async def _load_template_from_store(self, template_name: str) -> bytes | None:
        """
        Load a template from artifact store by name.

        Args:
            template_name: Name of the template presentation in artifact store

        Returns:
            Template data as bytes or None if not found
        """
        store = self._get_store()
        if not store:
            logger.error("No artifact store available for template loading")
            return None

        try:
            # Check if template exists in our namespace mapping
            namespace_id = self._namespace_ids.get(template_name)
            if namespace_id:
                data = await store.read_namespace(namespace_id)
                if data:
                    logger.info(f"Loaded template from artifact store: {template_name}")
                    return data

            # Try to find template by searching for it
            logger.warning(f"Template not found in namespace mapping: {template_name}")
            return None
        except Exception as e:
            logger.error(f"Failed to load template from artifact store: {e}")
            return None

    async def create(
        self,
        name: str,
        theme: str | None = None,
        template_name: str | None = None,
    ) -> PresentationMetadata:
        """
        Create a new presentation, optionally from a template.

        Args:
            name: Presentation name
            theme: Optional theme to apply
            template_name: Optional name of a template presentation in artifact store to use as base

        Returns:
            PresentationMetadata for the new presentation
        """
        if template_name:
            # First check builtin templates via TemplateManager
            from ..templates import TemplateManager

            template_manager = TemplateManager()
            template_data = await template_manager.get_template_data(template_name)

            if template_data:
                # Create presentation from builtin template
                buffer = io.BytesIO(template_data)
                prs = await asyncio.to_thread(Presentation, buffer)
                logger.info(f"Created presentation from builtin template: {template_name}")

                # Remove all existing slides from the template - we only want the layouts/master
                # Users will add slides using pptx_add_slide_from_template
                slide_ids_to_remove = list(range(len(prs.slides) - 1, -1, -1))
                for slide_idx in slide_ids_to_remove:
                    rId = prs.slides._sldIdLst[slide_idx].rId
                    prs.part.drop_rel(rId)
                    del prs.slides._sldIdLst[slide_idx]
                logger.info(
                    f"Removed {len(slide_ids_to_remove)} template slides, keeping only layouts"
                )
            else:
                # Fallback to artifact store template
                template_data = await self._load_template_from_store(template_name)
                if template_data:
                    # Create presentation from template bytes
                    buffer = io.BytesIO(template_data)
                    prs = await asyncio.to_thread(Presentation, buffer)
                    logger.info(
                        f"Created presentation from artifact store template: {template_name}"
                    )

                    # Remove all existing slides from the template - we only want the layouts/master
                    slide_ids_to_remove = list(range(len(prs.slides) - 1, -1, -1))
                    for slide_idx in slide_ids_to_remove:
                        rId = prs.slides._sldIdLst[slide_idx].rId
                        prs.part.drop_rel(rId)
                        del prs.slides._sldIdLst[slide_idx]
                    logger.info(
                        f"Removed {len(slide_ids_to_remove)} template slides, keeping only layouts"
                    )
                else:
                    # Fallback to blank presentation if template not found
                    logger.warning(
                        f"Template {template_name} not found in builtin or artifact store, creating blank presentation"
                    )
                    prs = Presentation()
        else:
            # Create blank presentation
            prs = Presentation()
            logger.info("Created blank presentation")

        # If presentation already exists, delete the old one first to ensure clean overwrite
        if name in self._presentations:
            logger.info(f"Presentation '{name}' already exists, overwriting...")
            # Delete old namespace if it exists
            old_namespace_id = self._namespace_ids.get(name)
            if old_namespace_id:
                store = self._get_store()
                if store:
                    try:
                        await store.delete_namespace(old_namespace_id)
                        logger.info(
                            f"Deleted old artifact for presentation '{name}' ({old_namespace_id})"
                        )
                    except Exception as e:
                        logger.warning(f"Could not delete old artifact: {e}")
                # Remove from namespace tracking
                del self._namespace_ids[name]

        self._presentations[name] = prs
        self._current_presentation = name
        self._update_cache_timestamp(name)  # Mark as freshly cached

        # Auto-save to artifact store (creates namespace)
        saved = await self._save_to_store(name, prs)

        # Create metadata
        metadata = PresentationMetadata(
            name=name,
            slide_count=len(prs.slides),
            theme=theme,
            vfs_path=self.get_artifact_uri(name),
            namespace_id=self.get_namespace_id(name),
            is_saved=saved,
            template_path=template_name,
        )
        self._metadata[name] = metadata

        return metadata

    async def get(
        self, name: str | None = None, force_refresh: bool = False
    ) -> tuple[PresentationType, PresentationMetadata] | None:
        """
        Get a presentation and its metadata by name.

        MULTI-INSTANCE SAFE: Always checks artifact store if cache is stale or forced.

        Args:
            name: Presentation name (uses current if not specified)
            force_refresh: If True, bypass cache and reload from artifact store

        Returns:
            Tuple of (Presentation, PresentationMetadata) or None if not found
        """
        pres_name = name or self._current_presentation
        if not pres_name:
            return None

        # Check if cache is valid (TTL-based)
        cache_valid = self._is_cache_valid(pres_name) and not force_refresh

        # If cache is valid, use it
        if cache_valid and pres_name in self._presentations:
            prs = self._presentations[pres_name]
            metadata = self._metadata.get(pres_name)

            # Create metadata if missing (for backwards compatibility)
            if not metadata:
                metadata = PresentationMetadata(
                    name=pres_name,
                    slide_count=len(prs.slides),
                    vfs_path=self.get_artifact_uri(pres_name),
                    namespace_id=self.get_namespace_id(pres_name),
                    is_saved=True,
                )
                self._metadata[pres_name] = metadata

            logger.debug(f"Using cached presentation '{pres_name}'")
            return (prs, metadata)

        # Cache is stale or missing - load from artifact store
        logger.debug(f"Cache miss or stale for '{pres_name}', loading from artifact store")
        loaded_prs = await self._load_from_store(pres_name, force_refresh=force_refresh)
        if loaded_prs is not None:
            # _load_from_store already updated cache
            # Create/update metadata
            metadata = PresentationMetadata(
                name=pres_name,
                slide_count=len(loaded_prs.slides),
                vfs_path=self.get_artifact_uri(pres_name),
                namespace_id=self.get_namespace_id(pres_name),
                is_saved=True,
            )
            self._metadata[pres_name] = metadata

            return (loaded_prs, metadata)

        return None

    async def get_presentation(self, name: str | None = None) -> PresentationType | None:
        """
        Get just the presentation object (async).

        This method attempts to load from artifact store if not found in memory.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Presentation object or None if not found
        """
        result = await self.get(name)
        if result:
            prs, _ = result
            return prs
        return None

    async def get_metadata(self, name: str | None = None) -> PresentationMetadata | None:
        """
        Get just the metadata (async).

        This method attempts to load from artifact store if not found in memory.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            PresentationMetadata or None if not found
        """
        result = await self.get(name)
        if result:
            _, metadata = result
            return metadata
        return None

    async def save(self, name: str | None = None) -> bool:
        """
        Save presentation to artifact store.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            True if successful, False otherwise
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return False

        return await self._save_to_store(pres_name, self._presentations[pres_name])

    async def update(self, name: str | None = None) -> bool:
        """
        Update presentation in artifact store after modifications.

        This should be called after any modification to ensure persistence.
        Also updates metadata (slide count, modified time, etc.).
        Updates cache timestamp to keep cache fresh.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            True if successful, False otherwise
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return False

        # Update metadata before saving
        prs = self._presentations[pres_name]
        metadata = self._metadata.get(pres_name)
        if metadata:
            metadata.slide_count = len(prs.slides)
            metadata.modified_at = datetime.now()

        # Save to artifact store
        result = await self._save_to_store(pres_name, prs)

        # Update cache timestamp after successful save
        if result:
            self._update_cache_timestamp(pres_name)

        return result

    async def delete(self, name: str) -> bool:
        """
        Delete a presentation from memory and artifact store.

        Args:
            name: Presentation name

        Returns:
            True if successful, False otherwise
        """
        if name not in self._presentations:
            return False

        # Clean up all cache data
        del self._presentations[name]
        if name in self._metadata:
            del self._metadata[name]
        if name in self._cache_timestamps:
            del self._cache_timestamps[name]
        if name in self._namespace_ids:
            del self._namespace_ids[name]

        # Update current if we deleted it
        if self._current_presentation == name:
            self._current_presentation = (
                next(iter(self._presentations), None) if self._presentations else None
            )

        # Delete from artifact store
        await self._delete_from_store(name)

        return True

    async def list_presentations(self) -> ListPresentationsResponse:
        """
        List all presentations with metadata.

        MULTI-INSTANCE SAFE: Queries artifact store to list all presentations in the session.

        Returns:
            ListPresentationsResponse with presentation info
        """
        presentations: list[PresentationInfo] = []
        store = self._get_store()

        if store:
            # Query artifact store for all presentations in this session
            try:
                namespaces = await store.list_namespaces()
                logger.debug(f"Found {len(namespaces)} namespaces in artifact store")

                for ns_info in namespaces:
                    # Filter for presentations (not templates)
                    if (
                        ns_info.name.startswith(f"{self.base_path}/")
                        and "/templates/" not in ns_info.name
                    ):
                        # Extract presentation name from namespace
                        name_parts = ns_info.name.split("/")
                        if len(name_parts) >= 2:
                            raw_name = name_parts[-1]
                            # Check if we have it in cache
                            prs = None
                            if raw_name in self._presentations and self._is_cache_valid(raw_name):
                                prs = self._presentations[raw_name]
                            else:
                                # Load to get accurate slide count
                                prs = await self._load_from_store(raw_name)

                            slide_count = len(prs.slides) if prs else 0

                            presentations.append(
                                PresentationInfo(
                                    name=raw_name,
                                    slide_count=slide_count,
                                    is_current=(raw_name == self._current_presentation),
                                    file_path=f"artifact://chuk-mcp-pptx/{ns_info.name}",
                                    namespace_id=ns_info.namespace_id,
                                )
                            )
                            logger.debug(
                                f"Listed presentation: {raw_name} ({ns_info.namespace_id})"
                            )
            except Exception as e:
                logger.error(f"Failed to list presentations from artifact store: {e}")
                # Fallback to memory-only listing
                for name, prs in self._presentations.items():
                    metadata = self._metadata.get(name)
                    presentations.append(
                        PresentationInfo(
                            name=name,
                            slide_count=len(prs.slides),
                            is_current=(name == self._current_presentation),
                            file_path=(
                                self.get_artifact_uri(name)
                                if metadata and metadata.is_saved
                                else None
                            ),
                            namespace_id=self.get_namespace_id(name),
                        )
                    )
        else:
            # No artifact store - list from memory only
            logger.debug("No artifact store available, listing from memory only")
            for name, prs in self._presentations.items():
                metadata = self._metadata.get(name)
                presentations.append(
                    PresentationInfo(
                        name=name,
                        slide_count=len(prs.slides),
                        is_current=(name == self._current_presentation),
                        file_path=(
                            self.get_artifact_uri(name) if metadata and metadata.is_saved else None
                        ),
                        namespace_id=self.get_namespace_id(name),
                    )
                )

        return ListPresentationsResponse(
            presentations=presentations,
            total=len(presentations),
            current=self._current_presentation,
        )

    async def set_current(self, name: str) -> bool:
        """
        Set the current presentation.

        Args:
            name: Presentation name

        Returns:
            True if successful, False if presentation not found
        """
        if name not in self._presentations:
            # Try loading from artifact store
            prs = await self._load_from_store(name)
            if prs:
                self._presentations[name] = prs
            else:
                return False

        self._current_presentation = name
        return True

    def get_current_name(self) -> str | None:
        """Get the name of the current presentation."""
        return self._current_presentation

    async def update_slide_metadata(self, slide_index: int) -> None:
        """
        Update metadata for a slide after modifications (async).

        Args:
            slide_index: Index of the slide to update
        """
        if not self._current_presentation:
            return

        # Use async get to ensure we have the presentation loaded
        result = await self.get(self._current_presentation)
        if not result:
            return

        prs, metadata = result

        # Ensure we have enough slide metadata entries
        while len(metadata.slides) <= slide_index:
            metadata.slides.append(SlideMetadata(index=len(metadata.slides), layout="Blank"))

        # Get the slide
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            slide_meta = metadata.slides[slide_index]

            # Update metadata from slide
            slide_meta.shape_count = len(slide.shapes)
            slide_meta.has_title = slide.shapes.title is not None
            if slide_meta.has_title and slide.shapes.title:
                slide_meta.title_text = slide.shapes.title.text

            # Check for charts, tables, images
            for shape in slide.shapes:
                # Check for charts - has_chart raises exception on non-chart shapes
                try:
                    if hasattr(shape, "has_chart") and shape.has_chart:
                        slide_meta.has_chart = True
                except Exception:
                    pass

                # Check for tables - has_table raises exception on non-table shapes
                try:
                    if hasattr(shape, "has_table") and shape.has_table:
                        slide_meta.has_table = True
                except Exception:
                    pass

                # Check for images by shape type
                if shape.shape_type == 13:  # Picture type
                    slide_meta.has_images = True

            metadata.update_modified()

    async def export_base64(self, name: str | None = None) -> str | None:
        """
        Export presentation as base64.

        Args:
            name: Presentation name (uses current if not specified)

        Returns:
            Base64-encoded presentation data or None
        """
        pres_name = name or self._current_presentation
        if not pres_name or pres_name not in self._presentations:
            return None

        prs = self._presentations[pres_name]
        try:
            buffer = io.BytesIO()
            await asyncio.to_thread(prs.save, buffer)
            buffer.seek(0)
            return base64.b64encode(buffer.read()).decode("utf-8")
        except Exception as e:
            logger.error(f"Failed to export as base64: {e}")
            return None

    async def import_base64(self, data: str, name: str, as_template: bool = False) -> bool:
        """
        Import presentation from base64.

        Args:
            data: Base64-encoded presentation data
            name: Name for the imported presentation
            as_template: If True, saves as a template (doesn't set as current)

        Returns:
            True if successful, False otherwise
        """
        try:
            buffer = io.BytesIO(base64.b64decode(data))
            prs = await asyncio.to_thread(Presentation, buffer)
            self._presentations[name] = prs
            self._update_cache_timestamp(name)  # Mark as freshly cached

            if not as_template:
                self._current_presentation = name

            # Auto-save to artifact store
            await self._save_to_store(name, prs)

            # Create metadata
            metadata = PresentationMetadata(
                name=name,
                slide_count=len(prs.slides),
                vfs_path=self.get_artifact_uri(name),
                namespace_id=self.get_namespace_id(name),
                is_saved=True,
            )
            self._metadata[name] = metadata

            logger.info(
                f"Imported presentation: {name}" + (" (as template)" if as_template else "")
            )
            return True
        except Exception as e:
            logger.error(f"Failed to import from base64: {e}")
            return False

    async def import_template(self, file_path: str, template_name: str) -> bool:
        """
        Import a PowerPoint file as a template into the artifact store.

        Args:
            file_path: Path to the PowerPoint file
            template_name: Name to save the template as

        Returns:
            True if successful, False otherwise
        """
        store = self._get_store()
        if not store:
            logger.error("No artifact store available for template import")
            return False

        from chuk_mcp_server import NamespaceType, StorageScope

        try:
            # Read the PowerPoint file
            with open(file_path, "rb") as f:
                data = await asyncio.to_thread(f.read)

            # Verify it's a valid presentation
            buffer = io.BytesIO(data)
            prs = await asyncio.to_thread(Presentation, buffer)
            logger.info(f"Validated template file: {file_path} ({len(prs.slides)} slides)")

            # Create namespace for template
            safe_name = self._sanitize_name(template_name)
            namespace_info = await store.create_namespace(
                type=NamespaceType.BLOB,
                scope=StorageScope.SESSION,
                name=f"{self.base_path}/templates/{safe_name}",
                metadata={
                    "mime_type": self.PPTX_MIME_TYPE,
                    "template_name": template_name,
                    "file_extension": ".pptx",
                    "is_template": True,
                    "slide_count": len(prs.slides),
                },
            )
            self._namespace_ids[template_name] = namespace_info.namespace_id

            # Write template data
            await store.write_namespace(namespace_info.namespace_id, data=data)

            # Store in memory for immediate use
            self._presentations[template_name] = prs
            self._update_cache_timestamp(template_name)  # Mark as freshly cached

            # Create metadata
            metadata = PresentationMetadata(
                name=template_name,
                slide_count=len(prs.slides),
                vfs_path=self.get_artifact_uri(template_name),
                namespace_id=namespace_info.namespace_id,
                is_saved=True,
            )
            self._metadata[template_name] = metadata

            logger.info(f"Imported template: {template_name} ({namespace_info.namespace_id})")
            return True
        except Exception as e:
            logger.error(f"Failed to import template: {e}")
            return False

    def clear_all(self) -> None:
        """Clear all presentations from memory cache."""
        self._presentations.clear()
        self._metadata.clear()
        self._namespace_ids.clear()
        self._cache_timestamps.clear()
        self._current_presentation = None

        # Note: This doesn't delete from artifact store, only clears memory cache
        # Artifact store presentations persist based on session/scope
        # They will be reloaded on next access
